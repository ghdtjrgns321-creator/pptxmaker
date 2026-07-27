"""차트·다이어그램 렌더 모듈 — pptx-visuals 스킬의 단일 출처.

build_pptx.py가 import해서 사용한다. 차트는 네이티브 차트(PPT에서 데이터 수정 가능),
다이어그램은 도형+커넥터 조합(timeline·matrix_2x2·venn 등 9종)으로 렌더한다.
색·폰트·크기는 인자로 받은 brand dict(brand-kit.yaml)에서만 나온다.
"""

from pathlib import Path

from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


def _rgb(hex6):
    return RGBColor.from_string(hex6)


# --- 아이콘 (Lucide MIT — make_icons.py가 brand 색 PNG로 사전 변환) ---
ICON_DIR = Path(__file__).resolve().parents[2] / "pptx-build" / "assets" / "icons"


def add_icon(slide, name, x, y, size, color="primary"):
    """브랜드 색 아이콘 PNG 삽입. color: primary | accent | white.
    없는 이름이면 명시적 실패(조용한 누락 금지) — 목록은 make_icons.py ICONS."""
    path = ICON_DIR / f"{name}_{color}.png"
    if not path.exists():
        raise ValueError(f"unknown icon: {name} ({path})")
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(size), Inches(size))


CHART_TYPES = {
    "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,  # 항목 간 크기 비교(카테고리명 짧을 때)
    "hbar": XL_CHART_TYPE.BAR_CLUSTERED,  # 순위·비중 비교(카테고리명 길 때 — 레이블 잘림 방지)
    "line": XL_CHART_TYPE.LINE,  # 시계열 추세
    "pie": XL_CHART_TYPE.PIE,  # 구성비(합계 100%·항목 ≤5)
    "doughnut": XL_CHART_TYPE.DOUGHNUT,  # 구성비 + 중앙 여백 활용
    "stacked_bar": XL_CHART_TYPE.COLUMN_STACKED,  # 구성 변화 비교
}


# --- 차트 ---
def add_chart(slide, spec, brand, x, y, w, h):
    """네이티브 차트 삽입 + brand 팔레트 시리즈 색 적용. GraphicFrame 반환."""
    ct = CHART_TYPES.get(spec.get("chart_type", "bar"), XL_CHART_TYPE.COLUMN_CLUSTERED)
    data = CategoryChartData()
    data.categories = spec["categories"]
    for name, values in spec["series"].items():
        data.add_series(name, values)
    gf = slide.shapes.add_chart(ct, Inches(x), Inches(y), Inches(w), Inches(h), data)
    ch = gf.chart
    ch.has_legend = len(spec["series"]) > 1 or ct in (
        XL_CHART_TYPE.PIE,
        XL_CHART_TYPE.DOUGHNUT,
    )
    if ch.has_legend:
        ch.legend.position = XL_LEGEND_POSITION.BOTTOM
        ch.legend.include_in_layout = False
    _apply_series_colors(ch, spec.get("chart_type", "bar"), brand, spec.get("emphasis"))
    _apply_chart_style(ch, spec.get("chart_type", "bar"), brand)
    return gf


def _apply_chart_style(chart, kind, brand):
    """PPT 기본 차트 티 제거 — 자동 제목 삭제, 글자 축소·브랜드 폰트, 데이터 레이블."""
    c, f, s = brand["colors"], brand["fonts"], brand["sizes"]
    chart.has_title = False  # 슬라이드 제목과 중복되는 차트 자동 제목 제거
    chart.font.size = Pt(s.get("caption", 9) + 1)
    chart.font.name = f["body"]
    chart.font.color.rgb = _rgb(c["text"])
    if chart.has_legend:
        chart.legend.font.size = Pt(s.get("caption", 9) + 1)
    for axis_name in ("category_axis", "value_axis"):
        try:
            axis = getattr(chart, axis_name)
            axis.tick_labels.font.size = Pt(s.get("caption", 9) + 1)
            axis.tick_labels.font.name = f["body"]
            axis.format.line.color.rgb = _rgb(c["muted"])
            if axis.has_major_gridlines:
                axis.major_gridlines.format.line.color.rgb = _rgb(c["bg_alt"])
        except ValueError:  # pie 등 축 없는 차트
            continue
    if kind in ("bar", "hbar", "line", "stacked_bar"):
        plot = chart.plots[0]
        plot.has_data_labels = True
        plot.data_labels.font.size = Pt(s.get("caption", 9) + 1)
        plot.data_labels.font.name = f["body"]
        plot.data_labels.font.color.rgb = _rgb(c["text"])


GRAY_BASE = "C7CBD1"  # BCG 원칙: 데이터는 전부 회색, 강조만 accent 1색 (mpl_exhibits와 동일값)


def _apply_series_colors(chart, kind, brand, emphasis=None):
    """Office 기본 테마색 제거. emphasis(카테고리/시리즈명 목록)가 있으면
    '전부 회색 + 강조만 accent' BCG 규칙, 없으면 기존 팔레트 순환."""
    c = brand["colors"]
    palette = [c["accent"], c["primary"], c["muted"]]
    emph = set(emphasis or [])

    def pick(i, name):
        if emph:
            return _rgb(c["accent"]) if name in emph else _rgb(GRAY_BASE)
        return _rgb(palette[i % len(palette)])

    if kind in ("pie", "doughnut"):
        ser = chart.plots[0].series[0]
        cats = list(chart.plots[0].categories)
        for i, pt in enumerate(ser.points):
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = pick(i, cats[i] if i < len(cats) else "")
        return
    single = len(list(chart.series)) == 1
    for i, series in enumerate(chart.series):
        if emph and single and kind in ("bar", "hbar"):
            # 단일 시리즈 막대: 포인트 단위 강조(카테고리명 매칭)
            cats = list(chart.plots[0].categories)
            for j, pt in enumerate(series.points):
                pt.format.fill.solid()
                pt.format.fill.fore_color.rgb = pick(j, cats[j] if j < len(cats) else "")
            continue
        color = pick(i, series.name)
        if kind == "line":
            series.format.line.color.rgb = color
        else:
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = color


def add_callouts(slide, annotations, brand, x, y, w, h):
    """네이티브 차트 위 BCG식 콜아웃 — accent 테두리 박스(+지시선)를 도형으로 얹는다.
    이미지가 아니라 pptx 도형이므로 PPT에서 문구·위치 편집 가능.
    annotations: [{"text": "...", "at": [fx, fy], "point_to": [fx, fy](선택)}]
    fx/fy는 개체 영역(x,y,w,h) 안의 분수 좌표(0~1)."""
    c, f, s = brand["colors"], brand["fonts"], brand["sizes"]
    for ann in annotations:
        fx, fy = ann.get("at", [0.66, 0.08])
        bw, bh = 3.0, 0.55
        bx, by = x + fx * w, y + fy * h
        bx = min(bx, x + w - bw)  # 영역 밖 탈출 방지
        if ann.get("point_to"):
            px, py = ann["point_to"]
            ln = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(bx + bw / 2),
                Inches(by + bh),
                Inches(x + px * w),
                Inches(y + py * h),
            )
            ln.line.color.rgb = _rgb(c["accent"])
            ln.line.width = Pt(1.0)
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(bx), Inches(by), Inches(bw), Inches(bh)
        )
        try:
            box.adjustments[0] = 0.12
        except (IndexError, ValueError):
            pass
        box.fill.solid()
        box.fill.fore_color.rgb = _rgb(c["bg"])
        box.line.color.rgb = _rgb(c["accent"])
        box.line.width = Pt(1.2)
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = ann["text"]
        run.font.size = Pt(s.get("caption", 9) + 1)
        run.font.bold = True
        run.font.name = f["body"]
        run.font.color.rgb = _rgb(c["text"])


# --- 다이어그램 (도형+화살표 DSL) ---
def add_diagram(slide, spec, brand, x, y, w, h):
    """diagram 슬라이드 본문 렌더. **layout은 필수** — 기본값을 두지 않는다.

    남은 어휘 9종: matrix_2x2 · spectrum · harvey_table · check_matrix · venn · timeline ·
    icon_rows · stat_split · split_detail. layout 누락·오타는 아래 else의 ValueError로
    **빌드 시점에 시끄럽게** 죽는다(조용한 폴백 금지 — 그게 어휘 수렴의 기계적 원인이었다).

    **박스+화살표 계열 전멸(2026-07-25 사용자 확정 폐기 9종)**: `flow`·`process_band`·`layers`·
    `branch`·`cards`·`from_to`·`band_table`·`pro_con`·`contrast_split`. 사유는 하나다 — 팔레트
    전수 렌더에서 **전면 반려**("그냥 도형 나열 / 고등학생 ppt"). 노드 문법을 dense로 재작업한
    개선판까지 반려됐고, 실전 사용은 0건이었다. 이 물성들의 정본은 (A) 골든 도해다:
    분기·라우팅=G01·G04·G09·G11 · 집합 대응=G05 · 관계 전수=G08 · 단계=G10·G14 · 병렬 카드=
    `dense.hero_card` · 경계=G19(crosswalk: archetype-catalog).
    (B)에 남는 **정당한** 쓸모는 (A)에 어휘가 없는 물성뿐이다 — 시간축·좌표계·다축 충족도·교집합."""
    layout = spec.get("layout")
    if layout == "matrix_2x2":
        _matrix_2x2(slide, spec, brand, x, y, w, h)
    elif layout == "spectrum":
        _spectrum(slide, spec, brand, x, y, w, h)
    elif layout in ("harvey_table", "check_matrix"):
        _score_grid(slide, spec, brand, x, y, w, h, mode=layout)
    elif layout == "venn":
        _venn(slide, spec, brand, x, y, w, h)
    elif layout == "icon_rows":
        _icon_rows(slide, spec.get("rows", []), brand, x, y, w, h, tag_head=spec.get("tag_head"))
    elif layout == "stat_split":
        _stat_split(slide, spec, brand, x, y, w, h)
    elif layout == "split_detail":
        _split_detail(slide, spec, brand, x, y, w, h)
    elif layout == "timeline":
        _timeline(slide, spec["nodes"], brand, x, y, w, h)
    else:
        raise ValueError(f"unknown diagram layout: {layout}")


def _icon_rows(slide, rows, brand, x, y, w, h, tag_head=None, dark=False):
    """아이콘 행 리스트(McKinsey 실측 골격) — [아이콘 배지 | 볼드 리드 | 설명(+선택 tag 열)].
    서사형 내용의 기본 그릇(2026-07-25 이후 (B)의 정성 나열 정본 — 표로 파편화하던 band_table은 폐기).
    rows: [{"icon"?, "badge"?, "lead", "desc", "tag"?}] 3~6개. tag_head: 3열 머리(예: '실측 근거').

    `badge`(2026-07-25): 배지 자리에 **수량**을 넣는 옵션. 골든 문법 G18(수 배지 원인 귀속 행 —
    "떨어진 것이 각각 왜인가", 합이 분모와 맞는 전수 귀속)을 새 L-ID 없이 이 행이 흡수하기 위한
    최소 옵션이다(사용자 확정 2026-07-25). 없으면 종전대로 순번(i+1) — 순번은 귀속 수량을 말할
    수 없어서, 폴백만으로는 그 물성이 안 된다."""
    c, f, s = brand["colors"], brand["fonts"], brand["sizes"]
    ink = "FFFFFF" if dark else c["text"]
    lead_c = "FFFFFF" if dark else c["primary"]
    tag_w = 2.4 if any(r.get("tag") for r in rows) else 0.0
    head_h = 0.0
    if tag_head and tag_w:
        head_h = 0.32
        tb = slide.shapes.add_textbox(Inches(x + w - tag_w), Inches(y), Inches(tag_w), Inches(0.3))
        p = tb.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = tag_head
        r.font.size = Pt(s["caption"] + 1)
        r.font.bold = True
        r.font.name = f["head"]
        r.font.color.rgb = _rgb(c["muted"])
    n = max(len(rows), 1)
    rh = (h - head_h) / n
    badge_d = 0.5
    for i, row in enumerate(rows):
        ry = y + head_h + i * rh
        if i:  # 행 사이 헤어라인
            ln = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(ry), Inches(x + w), Inches(ry)
            )
            ln.line.color.rgb = _rgb("3A3F47" if dark else c["bg_alt"])
            ln.line.width = Pt(0.75)
        by = ry + rh / 2 - badge_d / 2
        badge = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x), Inches(by), Inches(badge_d), Inches(badge_d)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = _rgb("2A2E35" if dark else c["bg_alt"])
        badge.line.fill.background()
        icon = row.get("icon")
        if icon and (ICON_DIR / f"{icon}_accent.png").exists():
            add_icon(
                slide, icon, x + badge_d / 2 - 0.15, by + badge_d / 2 - 0.15, 0.3, color="accent"
            )
        else:  # 아이콘 없으면 수 배지(`badge`) · 없으면 순번 배지
            tf = badge.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = str(row.get("badge", i + 1))
            r.font.size = Pt(11)
            r.font.bold = True
            r.font.name = f["head"]
            r.font.color.rgb = _rgb(c["accent"])
        lead_w = 1.9
        tb = slide.shapes.add_textbox(
            Inches(x + badge_d + 0.18), Inches(ry + 0.08), Inches(lead_w), Inches(rh - 0.14)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = row["lead"]
        r.font.size = Pt(s["body"])
        r.font.bold = True
        r.font.name = f["head"]
        r.font.color.rgb = _rgb(lead_c)
        dx = x + badge_d + 0.18 + lead_w + 0.15
        dw = w - (dx - x) - tag_w - (0.2 if tag_w else 0)
        tb = slide.shapes.add_textbox(Inches(dx), Inches(ry + 0.08), Inches(dw), Inches(rh - 0.14))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        for seg, emph in _parse_emph_runs(row.get("desc", "")):
            r = p.add_run()
            r.text = seg
            r.font.size = Pt(s["body"] - 1)
            r.font.bold = emph
            r.font.name = f["body"]
            r.font.color.rgb = _rgb(c["accent"] if emph else ink)
        if row.get("tag"):
            tb = slide.shapes.add_textbox(
                Inches(x + w - tag_w), Inches(ry + 0.08), Inches(tag_w), Inches(rh - 0.14)
            )
            tf = tb.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = row["tag"]
            r.font.size = Pt(s["body"] - 2)
            r.font.name = f["body"]
            r.font.color.rgb = _rgb(c["muted"])


def _parse_emph_runs(text):
    """`**강조**` 세그먼트 분해(빌더 _parse_emph와 동일 문법)."""
    import re

    parts, i = [], 0
    for mt in re.finditer(r"\*\*(.+?)\*\*", text):
        if mt.start() > i:
            parts.append((text[i : mt.start()], False))
        parts.append((mt.group(1), True))
        i = mt.end()
    if i < len(text):
        parts.append((text[i:], False))
    return parts or [(text, False)]


def _stat_split(slide, spec, brand, x, y, w, h):
    """좌 빅넘버 존 + 세로 룰 + 우 아이콘 행(McKinsey '$1.2T' 페이지 골격).
    spec: {"stat": {"value", "label", "desc"?}, "rows": [icon_rows 행]}"""
    c, f, s = brand["colors"], brand["fonts"], brand["sizes"]
    stat = spec["stat"]
    lw = min(3.2, w * 0.28)
    tb = slide.shapes.add_textbox(Inches(x), Inches(y + h * 0.18), Inches(lw - 0.3), Inches(1.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = str(stat["value"])
    r.font.size = Pt(40)
    r.font.bold = True
    r.font.name = f["head"]
    r.font.color.rgb = _rgb(c["accent"])
    tb = slide.shapes.add_textbox(
        Inches(x), Inches(y + h * 0.18 + 1.0), Inches(lw - 0.3), Inches(h - 1.4)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = stat.get("label", "")
    r.font.size = Pt(s["body"])
    r.font.bold = True
    r.font.name = f["head"]
    r.font.color.rgb = _rgb(c["primary"])
    if stat.get("desc"):
        p2 = tf.add_paragraph()
        p2.space_before = Pt(4)
        r2 = p2.add_run()
        r2.text = stat["desc"]
        r2.font.size = Pt(s["body"] - 1)
        r2.font.name = f["body"]
        r2.font.color.rgb = _rgb(c["muted"])
    ln = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x + lw), Inches(y + 0.1), Inches(x + lw), Inches(y + h - 0.1)
    )
    ln.line.color.rgb = _rgb(c["muted"])
    ln.line.width = Pt(0.75)
    _icon_rows(
        slide,
        spec.get("rows", []),
        brand,
        x + lw + 0.35,
        y,
        w - lw - 0.35,
        h,
        tag_head=spec.get("tag_head"),
    )


def _split_detail(slide, spec, brand, x, y, w, h):
    """좌 시각(중첩 다이어그램/차트 아님 — 도형 다이어그램만) + 우 아이콘 행
    (McKinsey 동심원+설명 리스트 골격). spec: {"visual": <diagram spec>, "rows": [...]}"""
    vw = w * 0.42
    add_diagram(slide, spec["visual"], brand, x, y, vw, h)
    _icon_rows(
        slide,
        spec.get("rows", []),
        brand,
        x + vw + 0.4,
        y,
        w - vw - 0.4,
        h,
        tag_head=spec.get("tag_head"),
    )


def _venn(slide, spec, brand, x, y, w, h):
    """벤 다이어그램(BCG p12·McKinsey p117 실측) — 겹치는 원 2~3개, 교집합 라벨.
    투명 채움 대신 굵은 외곽선 원(편집 가능·브랜드 룰 유지).
    spec: {"sets": [{"label", "sub"?}] 2~3개, "overlap": "교집합 라벨(선택)"}"""
    c, f, s = brand["colors"], brand["fonts"], brand["sizes"]
    sets = spec["sets"][:3]
    n = len(sets)
    d = min(h * 0.82, w / (1 + 0.62 * (n - 1)))
    total_w = d * (1 + 0.62 * (n - 1))
    ox = x + (w - total_w) / 2
    oy = y + (h - d) / 2 - (0.15 if n == 3 else 0)
    colors = [c["accent"], c["muted"], c["primary"]]
    for i, st in enumerate(sets):
        cx = ox + i * d * 0.62
        cy = oy + (d * 0.28 if (n == 3 and i == 1) else 0)
        ring = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx), Inches(cy), Inches(d), Inches(d))
        ring.fill.background()
        ring.line.color.rgb = _rgb(colors[i % 3])
        ring.line.width = Pt(2.5)
        # 라벨은 원의 바깥쪽 반원 중심(교집합 회피)
        lx = cx + (0.1 * d if i == 0 else (0.55 * d if i == n - 1 else 0.3 * d))
        ly = cy + (d - 0.75 if (n == 3 and i == 1) else 0.18)
        tb = slide.shapes.add_textbox(Inches(lx), Inches(ly), Inches(d * 0.36), Inches(0.7))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = st["label"]
        r.font.size = Pt(s["body"])
        r.font.bold = True
        r.font.name = f["head"]
        r.font.color.rgb = _rgb(colors[i % 3])
        if st.get("sub"):
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            r2 = p2.add_run()
            r2.text = st["sub"]
            r2.font.size = Pt(s["body"] - 2)
            r2.font.name = f["body"]
            r2.font.color.rgb = _rgb(c["muted"])
    if spec.get("overlap"):
        mid_x = ox + total_w / 2 - 1.1
        mid_y = oy + d / 2 - 0.25 + (d * 0.1 if n == 3 else 0)
        tb = slide.shapes.add_textbox(Inches(mid_x), Inches(mid_y), Inches(2.2), Inches(0.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = spec["overlap"]
        r.font.size = Pt(s["body"] - 1)
        r.font.bold = True
        r.font.name = f["head"]
        r.font.color.rgb = _rgb(c["text"])


def _matrix_2x2(slide, spec, brand, x, y, w, h):
    """2×2 매트릭스(정성 판단의 시각 구조화) — 사분면 배경 + 축 라벨 + 항목 점 배치.
    spec: {"x_axis": ["좌", "우"], "y_axis": ["하", "상"], "quadrants": [4개 라벨(선택,
    좌상→우상→좌하→우하)], "items": [{"label", "x": 0~1, "y": 0~1, "emphasis"?}]}"""
    c, f, s = brand["colors"], brand["fonts"], brand["sizes"]
    ml, mb = 0.55, 0.4  # 축 라벨 여백
    px, py, pw, ph = x + ml, y, w - ml - 0.2, h - mb
    for qi in range(4):  # 사분면 배경(체커)
        qx = px + (qi % 2) * pw / 2
        qy = py + (qi // 2) * ph / 2
        q = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(qx), Inches(qy), Inches(pw / 2), Inches(ph / 2)
        )
        q.fill.solid()
        q.fill.fore_color.rgb = _rgb(c["bg_alt"] if qi in (0, 3) else c["bg"])
        q.line.color.rgb = _rgb(c["muted"])
        q.line.width = Pt(0.5)
    for qi, qlabel in enumerate(spec.get("quadrants", [])[:4]):
        qx = px + (qi % 2) * pw / 2
        qy = py + (qi // 2) * ph / 2
        tb = slide.shapes.add_textbox(
            Inches(qx + 0.08), Inches(qy + 0.05), Inches(pw / 2 - 0.16), Inches(0.3)
        )
        p = tb.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = qlabel
        r.font.size = Pt(s["caption"])
        r.font.bold = True
        r.font.name = f["head"]
        r.font.color.rgb = _rgb(c["muted"])
    # 축 라벨: x축 좌/우(하단), y축 하/상(좌측 세로)
    xa = spec.get("x_axis", ["", ""])
    ya = spec.get("y_axis", ["", ""])
    for i, (tx, align) in enumerate(zip(xa, (PP_ALIGN.LEFT, PP_ALIGN.RIGHT))):
        tb = slide.shapes.add_textbox(
            Inches(px + i * pw / 2), Inches(py + ph + 0.05), Inches(pw / 2), Inches(0.3)
        )
        p = tb.text_frame.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = tx
        r.font.size = Pt(s["caption"] + 1)
        r.font.bold = True
        r.font.name = f["head"]
        r.font.color.rgb = _rgb(c["primary"])
    for i, ty in enumerate(ya):  # [하, 상]
        tb = slide.shapes.add_textbox(
            Inches(x - 0.05),
            Inches(py + (1 - i) * (ph / 2) + ph / 4 - 0.15),
            Inches(ml),
            Inches(0.3),
        )
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = ty
        r.font.size = Pt(s["caption"] + 1)
        r.font.bold = True
        r.font.name = f["head"]
        r.font.color.rgb = _rgb(c["primary"])
    for it in spec.get("items", []):
        emph = it.get("emphasis")
        d = 0.3 if emph else 0.22
        ix = px + it["x"] * pw - d / 2
        iy = py + (1 - it["y"]) * ph - d / 2
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(ix), Inches(iy), Inches(d), Inches(d))
        dot.fill.solid()
        dot.fill.fore_color.rgb = _rgb(c["accent"] if emph else c["muted"])
        dot.line.color.rgb = _rgb(c["bg"])
        dot.line.width = Pt(1.0)
        tb = slide.shapes.add_textbox(
            Inches(ix - 0.6), Inches(iy + d + 0.02), Inches(d + 1.2), Inches(0.28)
        )
        tb.text_frame.word_wrap = False
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = it["label"]
        r.font.size = Pt(s["body"] - 2)
        r.font.bold = bool(emph)
        r.font.name = f["body"]
        r.font.color.rgb = _rgb(c["primary"] if emph else c["text"])


def _spectrum(slide, spec, brand, x, y, w, h):
    """성숙도 스펙트럼 — 단계 밴드(점진 음영) + 현위치 마커(accent) + 단계별 설명.
    spec: {"stages": [{"label", "sub"?}], "marker": <0-base 단계 인덱스>}"""
    c, f, s = brand["colors"], brand["fonts"], brand["sizes"]
    stages = spec["stages"]
    n = len(stages)
    marker = spec.get("marker")
    band_y = y + h * 0.32
    band_h = 0.7
    gap = 0.08
    sw = (w - gap * (n - 1)) / n
    shades = ["F2F3F5", "DDE0E4", "C7CBD1", "A9AFB8", "8A9099", "5C636D"]
    for i, st in enumerate(stages):
        sx = x + i * (sw + gap)
        seg = slide.shapes.add_shape(
            MSO_SHAPE.PENTAGON if i < n - 1 else MSO_SHAPE.RECTANGLE,
            Inches(sx),
            Inches(band_y),
            Inches(sw),
            Inches(band_h),
        )
        cur = marker is not None and i == marker
        seg.fill.solid()
        seg.fill.fore_color.rgb = _rgb(c["accent"] if cur else shades[min(i, len(shades) - 1)])
        seg.line.fill.background()
        tf = seg.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = st["label"]
        r.font.size = Pt(s["body"] - 1)
        r.font.bold = True
        r.font.name = f["head"]
        r.font.color.rgb = _rgb("FFFFFF" if (cur or i >= 3) else c["primary"])
        if st.get("sub"):
            tb = slide.shapes.add_textbox(
                Inches(sx), Inches(band_y + band_h + 0.12), Inches(sw), Inches(h - band_h - 0.6)
            )
            tb.text_frame.word_wrap = True
            p2 = tb.text_frame.paragraphs[0]
            p2.alignment = PP_ALIGN.CENTER
            r2 = p2.add_run()
            r2.text = st["sub"]
            r2.font.size = Pt(s["body"] - 2)
            r2.font.name = f["body"]
            r2.font.color.rgb = _rgb(
                c["text"] if (marker is not None and i == marker) else c["muted"]
            )
        if marker is not None and i == marker:  # 현위치 캡 마커
            tri = slide.shapes.add_shape(
                MSO_SHAPE.ISOSCELES_TRIANGLE,
                Inches(sx + sw / 2 - 0.12),
                Inches(band_y - 0.3),
                Inches(0.24),
                Inches(0.2),
            )
            tri.rotation = 180
            tri.fill.solid()
            tri.fill.fore_color.rgb = _rgb(c["accent"])
            tri.line.fill.background()


def _score_grid(slide, spec, brand, x, y, w, h, mode="harvey_table"):
    """하비볼/체크 매트릭스 — 행 라벨 + 열별 점수 기호. 정성 비교의 시각 구조화.
    spec: {"cols": ["기존", "본 시스템"], "rows": [{"label", "scores": [0~4 | "y"/"n"/"-"]}],
           "emphasis_col": <0-base 열 인덱스(선택)>}  harvey: 0=빈원~4=꽉찬원."""
    c, f, s = brand["colors"], brand["fonts"], brand["sizes"]
    cols, rows = spec["cols"], spec["rows"]
    label_w = min(3.4, w * 0.34)
    cw = (w - label_w) / len(cols)
    rh = min(0.62, (h - 0.5) / (len(rows) + 1))
    emph_col = spec.get("emphasis_col")
    for j, col in enumerate(cols):  # 열 머리
        tb = slide.shapes.add_textbox(
            Inches(x + label_w + j * cw), Inches(y), Inches(cw), Inches(0.35)
        )
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = col
        r.font.size = Pt(s["body"] - 1)
        r.font.bold = True
        r.font.name = f["head"]
        r.font.color.rgb = _rgb(c["accent"] if j == emph_col else c["primary"])
    _hline_y = y + 0.4
    ln = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(_hline_y), Inches(x + w), Inches(_hline_y)
    )
    ln.line.color.rgb = _rgb(c["primary"])
    ln.line.width = Pt(1.2)
    for i, row in enumerate(rows):
        ry = y + 0.5 + i * rh
        if i % 2 == 0:  # 줄무늬
            band = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(x), Inches(ry - 0.04), Inches(w), Inches(rh - 0.04)
            )
            band.fill.solid()
            band.fill.fore_color.rgb = _rgb(c["bg_alt"])
            band.line.fill.background()
        tb = slide.shapes.add_textbox(
            Inches(x + 0.08), Inches(ry), Inches(label_w - 0.1), Inches(rh)
        )
        tb.text_frame.word_wrap = True
        p = tb.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = row["label"]
        r.font.size = Pt(s["body"] - 1)
        r.font.name = f["body"]
        r.font.color.rgb = _rgb(c["text"])
        for j, score in enumerate(row["scores"]):
            cx = x + label_w + j * cw + cw / 2
            color = c["accent"] if j == emph_col else c["muted"]
            if mode == "check_matrix":
                mark = {"y": "✓", "n": "✗", "-": "–"}.get(str(score).lower(), str(score))
                tb2 = slide.shapes.add_textbox(
                    Inches(cx - 0.3), Inches(ry - 0.02), Inches(0.6), Inches(rh)
                )
                p2 = tb2.text_frame.paragraphs[0]
                p2.alignment = PP_ALIGN.CENTER
                r2 = p2.add_run()
                r2.text = mark
                r2.font.size = Pt(s["body"] + 2)
                r2.font.bold = True
                r2.font.name = f["body"]
                r2.font.color.rgb = _rgb(
                    color if mark == "✓" else (c["muted"] if mark == "–" else c["primary"])
                )
                continue
            d = 0.26  # harvey ball
            base = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(cx - d / 2),
                Inches(ry + rh / 2 - d / 2 - 0.04),
                Inches(d),
                Inches(d),
            )
            sc = max(0, min(4, int(score)))
            if sc == 4:
                base.fill.solid()
                base.fill.fore_color.rgb = _rgb(color)
            else:
                base.fill.solid()
                base.fill.fore_color.rgb = _rgb(c["bg"])
            base.line.color.rgb = _rgb(color)
            base.line.width = Pt(1.2)
            if 0 < sc < 4:  # 부분 채움: PIE 도형 오버레이
                pie = slide.shapes.add_shape(
                    MSO_SHAPE.PIE,
                    Inches(cx - d / 2),
                    Inches(ry + rh / 2 - d / 2 - 0.04),
                    Inches(d),
                    Inches(d),
                )
                try:
                    pie.adjustments[0] = -90.0
                    pie.adjustments[1] = -90.0 + 90.0 * sc
                except (IndexError, ValueError):
                    pass
                pie.fill.solid()
                pie.fill.fore_color.rgb = _rgb(color)
                pie.line.fill.background()


def _arrowhead(connector):
    """커넥터 끝에 삼각 화살촉(python-pptx 미노출 → XML 직접 부착)."""
    ln = connector.line._get_or_add_ln()
    tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
    ln.append(tail)


def _mix(a, b, t):
    """브랜드 색 두 개의 중간색(t=0 → a, t=1 → b). 헤어라인·진행 램프 파생용.

    `goldenfab.kit.mix`와 같은 산식이지만 저쪽은 RGBColor 도메인, 여기는 brand dict의
    **hex 문자열** 도메인이다 — 스킬 간 import + 변환 글루가 산식(3줄)보다 길어지므로
    산식만 국소로 둔다. 색 **값**은 여전히 brand-kit 단일 출처(hex 리터럴 0, 통합D 래칫).
    """
    return "".join(
        f"{round(int(a[i : i + 2], 16) + (int(b[i : i + 2], 16) - int(a[i : i + 2], 16)) * t):02X}"
        for i in (0, 2, 4)
    )


def _tbox(slide, x, y, w, h, text, pt, font, color, *, bold=False, align=PP_ALIGN.LEFT):
    """좌측 정렬 기본 텍스트 상자(마진 0 · 한글 어절 줄바꿈) — 노드 내부 텍스트 단일 경로."""
    from pptx.enum.lang import MSO_LANGUAGE_ID

    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(max(w, 0.2)), Inches(max(h, 0.16)))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p._p.get_or_add_pPr().set("eaLnBrk", "0")  # §4 어절 단위 줄바꿈(런 언어 태그와 쌍)
    p.alignment = align
    p.line_spacing = 1.05
    r = p.add_run()
    r.text = text
    r.font.size = Pt(pt)
    r.font.bold = bold
    r.font.name = font
    r.font.color.rgb = _rgb(color)
    r.font.language_id = MSO_LANGUAGE_ID.KOREAN
    return tb


def _timeline(slide, nodes, brand, x, y, w, h):
    """수평 타임라인 — 축선 + 넘버 마일스톤 원 + 상단 라벨/하단 설명 (로드맵용)."""
    c, f, s = brand["colors"], brand["fonts"], brand["sizes"]
    n = len(nodes)
    line_y = y + h * 0.45
    axis = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x + 0.3), Inches(line_y), Inches(x + w - 0.3), Inches(line_y)
    )
    axis.line.color.rgb = _rgb(c["primary"])
    axis.line.width = Pt(2.0)
    _arrowhead(axis)
    step = (w - 1.4) / max(n - 1, 1)
    for i, nd in enumerate(nodes):
        cx = x + 0.7 + i * step
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(cx - 0.19), Inches(line_y - 0.19), Inches(0.38), Inches(0.38)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = _rgb(c["accent"])
        dot.line.color.rgb = _rgb(c["bg"])
        dot.line.width = Pt(1.5)
        tf = dot.text_frame
        tf.word_wrap = False
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(i + 1)
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.name = f["head"]
        run.font.color.rgb = _rgb("FFFFFF")
        # 상단 라벨(굵게) / 하단 설명(작게) — 박스를 본문 영역 안으로 클램프(양끝 잘림 방지)
        bx = max(x, cx - step / 2 + 0.1)
        bw = min(x + w, cx + step / 2 - 0.1) - bx
        lb = slide.shapes.add_textbox(Inches(bx), Inches(line_y - 0.85), Inches(bw), Inches(0.5))
        p1 = lb.text_frame.paragraphs[0]
        lb.text_frame.word_wrap = True
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = nd["label"]
        r1.font.size = Pt(s["body"])
        r1.font.bold = True
        r1.font.name = f["head"]
        r1.font.color.rgb = _rgb(c["primary"])
        if nd.get("sub"):
            sb = slide.shapes.add_textbox(
                Inches(bx), Inches(line_y + 0.32), Inches(bw), Inches(0.8)
            )
            sb.text_frame.word_wrap = True
            p2 = sb.text_frame.paragraphs[0]
            p2.alignment = PP_ALIGN.CENTER
            r2 = p2.add_run()
            r2.text = nd["sub"]
            r2.font.size = Pt(s["body"] - 2)
            r2.font.name = f["body"]
            r2.font.color.rgb = _rgb(c["muted"])
