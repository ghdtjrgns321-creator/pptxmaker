"""골든 레이아웃 기계 오딧 러너 — 2초. 제3자 채점 **전에** 이걸 통과시킨다.

사용:
    uv run python .claude/skills/pptx-build/scripts/audit_golden.py            # 검사
    uv run python .claude/skills/pptx-build/scripts/audit_golden.py --selftest # 규칙 민감도 검증

규칙 본체는 `goldenfab/audit.py`. 여기는 **레이아웃별 기대값**만 선언한다.
새 레이아웃을 오딧에 넣으려면 SPECS에 한 항목을 추가한다.

왜 이 파일이 있나: S4 한 장에 제3자 채점 2회 = 30분이 들었는데, 두 채점자의 FAIL 9건 중
**7건이 순수 산수**였다. 채점자는 규칙을 **발견**할 뿐 적용은 코드가 한다 — 한 번 캔 규칙을
여기 내려두면 다음 장부터 2초에 잡힌다. 상세: goldenfab/audit.py 모듈 독스트링.
"""

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

HERE = Path(__file__).resolve().parent
sys.path.insert(0, str(HERE))

from goldenfab import audit as A  # noqa: E402
from goldenfab import grid as G  # noqa: E402  — 항목 수 파생 자 정본(pitch·track)
from goldenfab.kit import load_kit  # noqa: E402
from goldenfab.registry import LAYOUTS  # noqa: E402

K = load_kit()
C = K["rgb"]


def render(name):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    LAYOUTS[name](prs, None)
    return list(prs.slides)[-1]


def air_pairs_problem_grid():
    """S4 공기 검사 — 좌표는 장(배치) + 부품(LAYOUT)에서 파생한다.

    2026-07-26까지 이 검사는 폐기된 sparse(`_variant_k`) 좌표를 재고 있었다. 골든이 dense로
    승격된 뒤로는 **렌더되지 않는 좌표를 검사하는 hollow 검사**였다 — 통과해도 아무것도
    보증하지 못한다. 부품화하면서 실제로 그려지는 자리에서 파생하도록 옮겼다.
    """
    import goldenfab.dense as D
    import goldenfab.s04_dense as S4
    from goldenfab.figures.fan_in import LAYOUT as FAN
    from goldenfab.figures.gate_branch import LAYOUT as GATE

    rule = D.BAND_RULE_DY + 0.012  # 소제목 top → 룰 bottom
    gy = S4.GATE_BOX.y
    row1 = gy + GATE["row1_dy"]
    row_pitch = GATE["row_h"] + GATE["cap_dy"] + GATE["cap_h"] + GATE["row_air"]
    fy, n = S4.FAN_BOX.y, len(S4.DEFAULT["limits"])
    return [
        ("밴드1 룰 → 첫 행", S4.BAND1_HEAD_Y + rule, gy + GATE["ok_dy"]),
        (
            "2종 캡션 → 1종 행",
            row1 + GATE["row_h"] + GATE["cap_dy"] + GATE["cap_h"],
            row1 + row_pitch,
        ),
        ("1종 행 → 복귀 호 라벨", row1 + row_pitch + GATE["row_h"], gy + GATE["back_label_dy"]),
        ("복귀 호 → 밴드2 라벨", gy + GATE["back_dy"], S4.BAND2_HEAD_Y),
        ("밴드2 룰 → 한계 첫 칩", S4.BAND2_HEAD_Y + rule, fy),
        (
            "한계 마지막 칩 → 하한",
            fy + (n - 1) * FAN["item_pitch"] + FAN["item_h"],
            A.CONTENT_BOTTOM + A.AIR_MIN,
        ),
    ]


def air_pairs_tech_evidence():
    """S8 이분 매핑 — 소제목 룰 → 첫 칩, 마지막 칩 → 캡션. 좌표는 장(배치) 상수에서 파생.

    2026-07-26까지 이 검사는 폐기된 sparse(`s08_variants`) 좌표를 재고 있었다(S4와 같은 형태의
    hollow 검사 — 렌더되지 않는 좌표는 통과해도 아무것도 보증하지 못한다). 부품화하면서 실제로
    그려지는 자리에서 파생하도록 옮겼다.
    """
    import goldenfab.s08_dense as S8

    n_concepts = len({c for _t, _g, cs in S8.DEFAULT_C["terms"] for c in cs})
    pitch = G.pitch(n_concepts, S8.MAP_TOP, S8.MAP_BOTTOM, S8.MAP_H)
    last_bottom = S8.MAP_TOP + (n_concepts - 1) * pitch + S8.MAP_H
    return [
        ("매핑 룰 → 첫 칩", S8.MAP_HEAD_Y + 0.34 + 0.012, S8.MAP_TOP),  # _subhead 룰 dy
        ("마지막 칩 → 캡션", last_bottom, 3.36),  # build의 map_caption y
    ]


def air_pairs_tech_tree():
    """S9 카탈로그 — 소제목 → 첫 행, 마지막 행 → 하단 카드. 좌표는 장(배치) + 부품에서 파생.

    2026-07-26까지 폐기된 sparse(`s09_variants`) 좌표를 재던 hollow 검사였다(S4·S8과 같은 형태).
    """
    import goldenfab.s09_dense as S9
    from goldenfab.figures.relation_catalog import LAYOUT as CAT

    rows = S9.DEFAULT["cat_rows"]
    box = S9.CAT_BOX
    h = CAT["row_h"]
    pitch = G.pitch(len(rows), box.y, box.y + box.h - h, h, what="카탈로그 행")
    last_bottom = box.y + (len(rows) - 1) * pitch + h
    return [
        ("카탈로그 소제목 → 첫 행", 0.80 + 0.24, box.y),
        ("마지막 행 → 하단 카드", last_bottom, 3.62),  # build의 hero_card y
    ]


def air_pairs_mirror_matrix():
    """S17 그룹 미러 — 칼럼 헤드 → 첫 행, 그룹 경계 헤어라인 앞뒤, 마지막 행 → 하한.

    좌표는 전부 s17 모듈 상수·행 수에서 파생 — 렌더와 같은 산식이라 어긋나면 렌더가 틀린 것.
    """
    import goldenfab.s17_variants as S17

    groups = S17.VARIANT_C_DEFAULTS["groups"]
    n_rows = sum(len(axes) for _no, _q, _cl, axes in groups)
    eff_bottom = S17.LAST_BOTTOM - (len(groups) - 1) * S17.GGAP
    pitch = G.pitch(n_rows, S17.ROW0, eff_bottom, S17.CARD_H, cap=S17.PITCH_CAP)
    pairs = [("칼럼 헤드 → 첫 행", S17.COLHEAD_Y + 0.26, S17.ROW0)]
    i = 0
    for gi, (_no, _q, _cl, axes) in enumerate(groups):
        g_top = S17.ROW0 + i * pitch + gi * S17.GGAP
        if gi > 0:
            sep_y = g_top - ((pitch - S17.CARD_H) + S17.GGAP) / 2
            prev_bottom = S17.ROW0 + (i - 1) * pitch + (gi - 1) * S17.GGAP + S17.CARD_H
            pairs.append((f"그룹{gi} 마지막 → 헤어라인", prev_bottom, sep_y))
            pairs.append((f"헤어라인 → 그룹{gi + 1} 첫 행", sep_y, g_top))
        i += len(axes)
    last_bottom = S17.ROW0 + (n_rows - 1) * pitch + (len(groups) - 1) * S17.GGAP + S17.CARD_H
    pairs.append(("마지막 행 → 하한", last_bottom, A.CONTENT_BOTTOM + A.AIR_MIN))
    return pairs


def air_pairs_boundary():
    """S16 경계(dense) — 캔버스 → 다리, 다리 → 카드, 카드 → 하한. 전부 모듈 상수 파생.

    2026-07-26 dense 승격 전에는 **s18_variants(sparse) 상수**를 재고 있었다. 이 검사는 도형이
    아니라 모듈 상수끼리의 산술이라, 렌더러가 s16_dense로 바뀌어도 조용히 통과했다 —
    검사가 다른 장을 재는 hollow였다. 렌더하는 모듈과 재는 모듈은 같아야 한다.
    """
    import goldenfab.s16_dense as S16

    return [
        ("캔버스 → 다리", S16.CANVAS_Y + S16.CANVAS_H, S16.BRIDGE_Y),
        ("다리 → 카드", S16.BRIDGE_Y, S16.CARD_Y),
        ("카드 → 하한", S16.CARD_Y + S16.CARD_H, A.CONTENT_BOTTOM_DENSE + A.AIR_MIN),
    ]


# 레이아웃별 기대값. progress=허용 셰브런 수, ink_allow=primary 채움 허용 수(결론 바 1개).
SPECS = {
    "problem_grid": {
        "progress": 0,  # XOR을 셰브런으로 그리면 거짓 인과 (반려 4회)
        "ink_allow": 1,  # primary 채움 = 결론 바 하나뿐
        # 인과도 구획만 행 노드로 취급. dense 승격(2026-07-26)으로 세로가 압축되면서 하단
        # 카드 패널이 y=4.30까지 올라왔다 — 4.6이면 그 패널(h=2.36)이 계열에 섞여 "높이 2종"이
        # 된다. 인과도 마지막 행이 y=2.93이므로 4.0이 두 구획을 가르는 자리다.
        "node_top_max": 4.0,
        "air": air_pairs_problem_grid,
    },
    # 그림이 있는 장 — 침범 감시가 목적(2026-07-15 s10 버그: 비율 가정이 틀려 텍스트 1.74" 침범)
    #
    # `known`: 이 오딧을 만들기 **전에** 사용자 승인으로 확정된 장의 기존 위반. 규칙을 끄지 않고
    # 기준선으로 박는다 — **더 나빠지면 잡히고, 지금 상태는 통과**(래칫). 규칙을 빼면 hollow가 되고,
    # 그대로 빨간불이면 게이트가 상시 적색이라 무시당한다(둘 다 이 프로젝트가 이미 겪은 실패다).
    # 해소는 해당 장 재설계 시 — 재설계는 사용자 결정이라 여기서 임의로 하지 않는다.
    "screenshot": {
        "progress": 0,
        "ink_allow": 0,
        "node_top_max": 0,
        "air": lambda: [],
        "known": {
            # 관찰 번호 01~04가 accent 런 4 + 캡션 바 1. "사용처"로는 2(번호 장치·바)지만
            # 계수기는 런 단위라 5로 센다 — 계수 입도 문제이지 설계 결함이 아니다.
            "§2 accent 상한": 5,
            # ("P2③④ 경계": 1 — 2026-07-15 **해소**. 리터럴 피치 1.05를 항목 수에서 파생시켜
            #  마지막 본문 bottom 6.47 → 6.35. 3·4·5·6개 전부 하한 내. 기준선 삭제 = 빚 청산.)
        },
    },
    "tech_capture": {
        "progress": 0,
        "ink_allow": 0,
        "node_top_max": 0,
        "air": lambda: [],
        "known": {"§2 accent 상한": 6, "P4⑩ 검정 충돌": 2},
    },
    # 기술 설명 장 2종 (2026-07-15 표 → 관계 시각화 재설계). **known 없음** — 새로 그렸으므로
    # 빚 없이 시작한다. 여기서 기준선을 깔면 재설계의 의미가 없다.
    # ink_allow=2: 결론 바 + 승인된 다크 면 1개. S4의 "primary 채움 = 결론 바 하나뿐"은 그 장의
    # 규칙이다 — 거기선 검정이 '파국'과 '우리 원칙' 두 뜻을 져서 충돌했다. 여기선 두 번째 검정이
    # §E가 규정한 **다크 JSON 카드**(S8)와 **트리 루트 `기준서 1115`**(S9)로, 결론 바와 뜻이
    # 겹치지 않는다(코드 실물 / 기준서 최상위). 사용자도 코드 카드를 명시 승인했다("코드를
    # 보여주는 시각화는 좋은데").
    "tech_evidence": {  # S8 이분 매핑 — 실무어 → 기준서 개념
        "progress": 0,
        "ink_allow": 2,  # 결론 바 + JSON 카드(§E 다크 라운드 rect)
        "node_top_max": 0,
        "air": air_pairs_tech_evidence,
    },
    "tech_tree": {  # S9 간선 7종 카탈로그 + 위계 트리
        "progress": 0,
        "ink_allow": 2,  # 결론 바 + 트리 루트(기준서 1115)
        "node_top_max": 0,
        # 같은 노드의 재등장이 메시지다 — `변동대가` ×4(트리 개념층 + 카탈로그 hier·cp 출발 +
        # term 도착), `문단 50` ×3(트리 + cp·case 도착). 좌 트리와 우 카탈로그가 **같은 실물을
        # 두 각도로** 본다는 게 이 장의 구성이라, 재등장이 곧 두 구획이 맞물린다는 증거다.
        #
        # 5 → 8은 **악화가 아니라 드러남**이다. 이전 5는 좌 트리(bg_alt)와 우 카탈로그(흰+검정
        # 테두리)의 서명이 **갈려 있어서** 같은 이름이 다른 노드로 세어졌던 것이고, 그 갈림이
        # 바로 제3자가 잡은 클래스 충돌(#10)이었다. NODE_STYLE로 통일하니 진짜 수가 보인다 —
        # 규칙을 느슨하게 한 게 아니라 **측정이 정확해졌다**. 실측 8을 박아 늘면 FAIL.
        "dup_allow": 8,
        "air": air_pairs_tech_tree,
    },
    # 차별점 장 2종 (2026-07-16 밀도 재깎기). 이전까지 SPECS 미등록이라 채움률 규칙이 한 번도
    # 안 돌았고, S17이 채움률 14~29% FAIL 9건·판정 대비 FAIL 1건을 그대로 안고 있었다 —
    # 규칙이 있어도 러너에 안 물리면 없는 규칙이다. **known 없음** — 새로 깎았으므로 빚 0.
    "mirror_matrix": {  # S17 그룹 미러 — 좌 질문 칼럼 + 중앙 축 기둥 + 좌우 날개
        # 골든 덱에서 제외된 장(2026-07-16) — registry 창고용으로 sparse variant 렌더러를 유지하므로
        # 프로파일도 sparse다. dense로 채점하면 sparse 출처선(7.45)이 경계 위반으로 잡힌다.
        "profile": "sparse",
        "progress": 0,
        # 결론 바 + 축 스파인 칩 7 — 칩은 전부 같은 뜻(비교 축 라벨)이라 잉크 충돌이 아니다.
        # 계수기는 도형 단위라 8로 박는다(늘면 FAIL).
        "ink_allow": 8,
        "node_top_max": 0,  # 날개 높이 균일은 모듈 audit(wings=14, h 단일)이 이미 강제
        "air": air_pairs_mirror_matrix,
    },
    "boundary": {  # S18 경계 — 안(응답 4유형 계약) / 밖(차단·유보) + 하단 잔여 한계
        "progress": 0,
        "ink_allow": 3,  # 결론 바 + 차단 바 2(경계 진입을 끊는 짧은 바 — §5 차단 어휘)
        "node_top_max": 0,
        "air": air_pairs_boundary,
    },
}


def audit(name):
    """장별 오딧 = 전역 묶음(`generic_checks` 단일 출처) + 이 장에만 있는 오라클.

    2026-07-26 dense 승격 전에는 전역 검사를 **여기서 다시 나열**했다. 그래서 임계가 sparse에
    박제됐고, 골든이 dense가 된 순간 accent·검정충돌·경계 3종이 한꺼번에 적색이 됐다
    (registry 주석이 예고한 그 자리다). 규칙은 한 곳, 임계만 프로파일로 가른다.
    """
    spec = SPECS[name]
    sl = render(name)
    shapes = list(sl.shapes)
    print(f"── {name}: 도형 {len(shapes)}")
    res = [("§5 진행형 도형", *A.check_progress_shapes(shapes, spec["progress"]))]
    res += A.generic_checks(
        shapes,
        str(C["accent"]),
        dup_allow=spec.get("dup_allow", 0),
        screenshot=(name == "screenshot"),
        profile=spec.get("profile", "dense"),
    )
    res += [
        ("P4⑩ 검정 충돌", *A.check_ink_collision(shapes, str(C["primary"]), spec["ink_allow"])),
        ("P4⑧ 노드 높이", *A.check_node_heights(shapes, spec["node_top_max"])),
        ("§6 공기", *A.check_air(spec["air"]())),
    ]
    return A.report(res, known=spec.get("known"))


def selftest():
    """규칙 민감도 — 일부러 어긴 슬라이드를 잡는가. 안 잡으면 죽은 코드다."""
    from pptx.enum.shapes import MSO_SHAPE

    ok = []
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    sl = prs.slides.add_slide(prs.slide_layouts[6])

    # 셰브런 1개 심기 → §5가 잡아야
    sl.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(1), Inches(1), Inches(1), Inches(0.5))
    hit, msg, n = A.check_progress_shapes(list(sl.shapes), 0)
    ok.append(("§5 셰브런 검출", not hit, msg, n))

    # 흐린 판정 단어 심기 → P4⑤가 잡아야
    tb = sl.shapes.add_textbox(Inches(1), Inches(2), Inches(2), Inches(0.3))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = "✕ 불가역"
    r.font.color.rgb = C["muted"]  # 3.2:1
    hit, msg, n = A.check_verdict_contrast(list(sl.shapes), None)
    ok.append(("P4⑤ 흐린 판정 검출", not hit, msg, n))

    # 산문은 muted가 정상 — 판정으로 오탐하면 안 된다(2026-07-15 규칙 정밀화의 회귀 방지).
    # 이 케이스가 없으면 "없다"를 뺀 것이 규칙을 죽인 건지 고친 건지 구분할 수 없다.
    prs_p = Presentation()
    prs_p.slide_width, prs_p.slide_height = Inches(13.333), Inches(7.5)
    slp = prs_p.slides.add_slide(prs_p.slide_layouts[6])
    tbp = slp.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(0.3))
    rp = tbp.text_frame.paragraphs[0].add_run()
    rp.text = "문단 뭉치에는 순서도 관계도 없다"  # 17자 서사 본문 — 판정 아님
    rp.font.color.rgb = C["muted"]
    hit_p, msg_p, n_p = A.check_verdict_contrast(list(slp.shapes), None)
    ok.append(("P4⑤ 산문 오탐 0", hit_p, msg_p, n_p))

    # 그림 침범 — 구 s10 방식(height만 지정 → 비율이 폭을 정함)을 재현하면 잡아야
    import goldenfab.s10_dense as S10

    prs2 = Presentation()
    prs2.slide_width, prs2.slide_height = Inches(13.333), Inches(7.5)
    sl2 = prs2.slides.add_slide(prs2.slide_layouts[6])
    sl2.shapes.add_picture(S10.DEFAULT["img"], Inches(0.6), Inches(1.8), height=Inches(4.75))
    tb2 = sl2.shapes.add_textbox(Inches(8.6), Inches(1.8), Inches(4.1), Inches(0.3))
    tb2.text_frame.paragraphs[0].add_run().text = "이 화면에서 보이는 것"
    hit, msg, n = A.check_picture_overlap(list(sl2.shapes))
    ok.append(("§F 그림 침범 검출", not hit, msg, n))

    # 정상(fit_picture)은 통과해야 — 오탐 0
    hit3, msg3, n3 = A.check_picture_overlap(list(render("screenshot").shapes))
    ok.append(("§F 정상 오탐 0", hit3, msg3, n3))

    # 노드 클래스 — 같은 이름을 다른 스타일로 그리면 잡아야(제3자가 s9에서 찾은 최대 결함)
    from pptx.enum.shapes import MSO_SHAPE as _MS

    prs_n = Presentation()
    prs_n.slide_width, prs_n.slide_height = Inches(13.333), Inches(7.5)
    sln = prs_n.slides.add_slide(prs_n.slide_layouts[6])
    for i, (fill, line) in enumerate(((C["bg_alt"], None), (C["bg"], C["primary"]))):
        sh = sln.shapes.add_shape(
            _MS.ROUNDED_RECTANGLE, Inches(1 + i * 2), Inches(1), Inches(1.5), Inches(0.3)
        )
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
        if line is None:
            sh.line.fill.background()
        else:
            sh.line.color.rgb = line
        sh.text_frame.paragraphs[0].add_run().text = "변동대가"  # 같은 이름, 다른 스타일
    hit_n, msg_n, n_n = A.check_node_class(list(sln.shapes))
    ok.append(("P4⑩ 클래스 충돌 검출", not hit_n, msg_n, n_n))

    # 정상(단일 출처 NODE_STYLE)은 통과해야 — 오탐 0
    hit_t, msg_t, n_t = A.check_node_class(list(render("tech_tree").shapes))
    ok.append(("P4⑩ 정상 오탐 0", hit_t, msg_t, n_t))

    # S4의 두 ◇ `근거 충분?`(실선 vs 점선 = 의도된 대조)을 오탐하면 안 된다 — 초판이 테두리를
    # 서명에 넣어 이걸 잡았고, 그럼 duplicate 규칙("구별 장치를 둬라")과 서로 반대를 요구한다.
    hit_k, msg_k, n_k = A.check_node_class(list(render("problem_grid").shapes))
    ok.append(("P4⑩ 의도된 대조 오탐 0", hit_k, msg_k, n_k))

    # 밀도 밴드 — 희소 슬라이드(도형 4개)는 잡고, 골든 본문 장은 전수 통과해야 (2026-07-16
    # 변형 출발점 배선: 스냅샷 회귀가 물러난 자리의 밀집화 게이트. 임계는 스냅샷 파생 — 리터럴 0)
    band = A.density_band()
    prs_d = Presentation()
    prs_d.slide_width, prs_d.slide_height = Inches(13.333), Inches(7.5)
    sld = prs_d.slides.add_slide(prs_d.slide_layouts[6])
    for i in range(4):
        tb = sld.shapes.add_textbox(Inches(1), Inches(1 + i * 0.5), Inches(3), Inches(0.3))
        tb.text_frame.paragraphs[0].add_run().text = f"성긴 항목 {i + 1}"
    hit_d, msg_d, n_d = A.check_density(list(sld.shapes), band)
    ok.append(("§6-D 희소 검출", not hit_d, msg_d, n_d))

    body_keys = [k for k in SPECS]  # 오딧 등록 본문 장 = 골든 본문 표본
    dens_fails = []
    for k in body_keys:
        # 스크린샷 장(§F)은 밀도 예외 — 캡처가 내용을 진다
        okk, _m, _n = A.check_density(
            list(render(k).shapes), band, screenshot=(k in A.DENSITY_EXEMPT)
        )
        if not okk:
            dens_fails.append(k)
    ok.append(
        (
            "§6-D 골든 오탐 0",
            not dens_fails,
            f"골든 본문 {len(body_keys) - len(dens_fails)}/{len(body_keys)} 통과 {dens_fails}",
            len(dens_fails),
        )
    )

    # ── 텍스트 겹침(check_text_collision) 민감도 — 배경 없는 같은 폭 두 문장이 세로로 포개지면
    # 잡고, 좁은 라벨·fill 노드·정상 간격·골든 전수는 오탐 0 (2026-07-24 카드 목업 3회 반려 하강).
    def _txt(sl, x, y, w, h, s):
        tbx = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tbx.text_frame.word_wrap = True
        tbx.text_frame.paragraphs[0].add_run().text = s
        return tbx

    def _blank():
        p = Presentation()
        p.slide_width, p.slide_height = Inches(13.333), Inches(7.5)
        return p.slides.add_slide(p.slide_layouts[6])

    # (a) 검출 — evidence_card 함의줄↔출처줄 유형(fill 없는 같은 폭 문장이 세로로 겹침)
    slo = _blank()
    _txt(slo, 0.6, 2.60, 3.5, 0.18, "출처 · 7_JOURNEY.md §7.3")
    _txt(slo, 0.6, 2.73, 3.5, 0.24, "유사도 필터가 정답을 거른다")
    hit_c, msg_c, n_c = A.check_text_collision(list(slo.shapes))
    ok.append(("텍스트 겹침 검출", not hit_c, msg_c, n_c))

    # (b) 정상 간격 — 오탐 0(같은 두 문장을 안 겹치게 배치)
    slg = _blank()
    _txt(slg, 0.6, 2.60, 3.5, 0.18, "출처 · 7_JOURNEY.md §7.3")
    _txt(slg, 0.6, 2.88, 3.5, 0.24, "유사도 필터가 정답을 거른다")
    hit_g, msg_g, n_g = A.check_text_collision(list(slg.shapes))
    ok.append(("텍스트 겹침 정상 오탐 0", hit_g, msg_g, n_g))

    # (c) 경계 — 좁은 라벨이 넓은 캡션 밑(폭 크게 다름 = 다이어그램 라벨) → 오탐 0(width_sim)
    slw = _blank()
    _txt(slw, 0.6, 6.30, 12.1, 0.20, "지식그래프 — 기준서의 구조를 그대로 담는다")
    _txt(slw, 3.0, 6.42, 1.0, 0.18, "hier 79")
    hit_w, msg_w, n_w = A.check_text_collision(list(slw.shapes))
    ok.append(("텍스트 겹침 좁은라벨 오탐 0", hit_w, msg_w, n_w))

    # (d) 경계 — fill 있는 노드끼리 겹침 → 오탐 0(배경 있으면 판독 구분·의도, check_adhoc 담당)
    slf = _blank()
    for i, xx in enumerate((2.0, 2.1)):
        shf = slf.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(xx), Inches(2.0 + i * 0.15), Inches(2.0), Inches(0.4)
        )
        shf.fill.solid()
        shf.fill.fore_color.rgb = C["bg_alt"]
        shf.text_frame.paragraphs[0].add_run().text = f"노드 {i}"
    hit_f, msg_f, n_f = A.check_text_collision(list(slf.shapes))
    ok.append(("텍스트 겹침 fill노드 오탐 0", hit_f, msg_f, n_f))

    # (e) 통과작 전수 — SPECS 골든 레이아웃 전부 오탐 0(§6-D 골든 오탐 0과 같은 회귀 방어)
    tc_fails = [k for k in SPECS if not A.check_text_collision(list(render(k).shapes))[0]]
    ok.append(
        (
            "텍스트 겹침 골든 오탐 0",
            not tc_fails,
            f"골든 {len(SPECS) - len(tc_fails)}/{len(SPECS)} 통과 {tc_fails}",
            len(tc_fails),
        )
    )

    # 대비 산수 자체
    cr_muted = A.contrast(C["muted"])
    cr_primary = A.contrast(C["primary"])
    ok.append(
        (
            "대비 계산",
            cr_muted < 4.5 < cr_primary,
            f"muted {cr_muted} < 4.5 < primary {cr_primary}",
            0,
        )
    )

    print("── 민감도 자체검증 (규칙이 위반을 실제로 잡는가)")
    return A.report(ok)


if __name__ == "__main__":
    if "--selftest" in sys.argv:
        sys.exit(1 if selftest() else 0)
    fails = []
    for nm in SPECS:
        fails += audit(nm)
    sys.exit(1 if fails else 0)
