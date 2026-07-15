"""compare_golden — 골든 레퍼런스 덱의 도형 전수 회귀 게이트 (스냅샷 대조).

비교 속성: 도형 수 / shape_type / left·top·width·height(±0.005") / 전체 텍스트 /
솔리드 채움색 / 런 색 / 표 셀 / 첫 텍스트 런의 폰트 크기·볼드.

## 2026-07-15 단일화 — 무엇이 바뀌었나

전에는 `golden/`(동결 원본)과 `goldenfab/`(공장)을 **각각 빌드해서 서로 비교**했다. 두 트리에
같은 레이아웃 코드가 두 벌 있었기 때문이다(Phase 2 이식 비계). 그래서 장 하나 고치려면 두
파일을 똑같이 고쳐야 했고, 한쪽만 고치면 이 게이트가 깨졌다 — 사본끼리의 비교라 독립 검증력은
없으면서 동기화 비용만 들었다.

지금은 goldenfab이 유일 소스이고, **결과물(도형 서명)을 스냅샷으로 박아** 대조한다(골든 파일
테스트의 정석). 레퍼런스 덱 정의는 `goldenfab/reference.py`.

사용:
    uv run python .../compare_golden.py [리포트경로]      # 검사 (종료코드 0 = PASS)
    uv run python .../compare_golden.py --update-snapshot  # 의도적 변경 후 기준선 재생성
                                                          # → git diff로 반드시 리뷰할 것
"""

import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.util import Inches

HERE = Path(__file__).resolve().parent
ROOT = HERE.parents[3]  # scripts -> pptx-build -> skills -> .claude -> pptmaker
sys.path.insert(0, str(HERE))

SNAPSHOT = HERE.parent / "assets" / "golden-snapshot.json"

EMU = 914400
TOL = 0.005


def _color_str(color_fmt):
    """ColorFormat -> 비교용 문자열. 테마색·미지정도 안전하게 처리."""
    try:
        if color_fmt.type is None:
            return None
        return str(color_fmt.rgb)
    except Exception:
        try:
            return str(color_fmt.theme_color)
        except Exception:
            return None


def _fill_str(fill_fmt):
    try:
        if fill_fmt.type == MSO_FILL_TYPE.SOLID:
            return _color_str(fill_fmt.fore_color)
    except Exception:
        pass
    return None


def _run_colors(sh):
    """모든 런의 색 — 런 단위 강조색(01 accent 등)을 잡는다. 첫 런만 보면 놓친다."""
    if not sh.has_text_frame:
        return None
    out = []
    for p in sh.text_frame.paragraphs:
        for r in p.runs:
            out.append(_color_str(r.font.color))
    return tuple(out) if out else None


def _cells(sh):
    """표 셀의 (텍스트, 채움) — 표는 GraphicFrame 1덩어리라 text_frame/fill로는 안 잡힌다."""
    if not getattr(sh, "has_table", False):
        return None
    out = []
    for row in sh.table.rows:
        for cell in row.cells:
            out.append((cell.text, _fill_str(cell.fill)))
    return tuple(out)


def shape_sig(sh):
    sig = {
        "type": str(sh.shape_type),
        "l": sh.left / EMU if sh.left is not None else None,
        "t": sh.top / EMU if sh.top is not None else None,
        "w": sh.width / EMU if sh.width is not None else None,
        "h": sh.height / EMU if sh.height is not None else None,
        "text": sh.text_frame.text if sh.has_text_frame else "",
        "fill": None,
        "run": None,
        "run_colors": _run_colors(sh),
        "cells": _cells(sh),
    }
    # str(fill.type) 비교 금지 — 실제 repr은 "<MSO_FILL_TYPE.SOLID: 1>"이라 문자열 대조는
    # 영원히 거짓이 되고 게이트가 색을 못 본 채 통과한다(2026-07-15 실측: 225개 중 0개 비교).
    try:
        sig["fill"] = _fill_str(sh.fill)
    except Exception:
        pass
    if sh.has_text_frame:
        for p in sh.text_frame.paragraphs:
            if p.runs:
                r = p.runs[0]
                sig["run"] = (r.font.size.pt if r.font.size else None, r.font.bold)
                break
    return sig


def diff_shapes(a, b):
    """도형 서명 비교 — 불일치 필드 목록 반환(빈 리스트 = 일치)."""
    bad = []
    for k in ("type", "text", "fill", "run", "run_colors", "cells"):
        if a[k] != b[k]:
            bad.append(f"{k}: {a[k]!r} != {b[k]!r}")
    for k in ("l", "t", "w", "h"):
        av, bv = a[k], b[k]
        if (av is None) != (bv is None) or (av is not None and abs(av - bv) > TOL):
            bad.append(f"{k}: {av} != {bv}")
    return bad


def jsonable(sig):
    """JSON 왕복 정규화 — 스냅샷은 tuple을 list로 저장하므로 살아있는 쪽도 맞춰야 비교가 성립한다."""
    return json.loads(json.dumps(sig))


def deck_signature(prs):
    """덱 전체 도형 서명 — 스냅샷 저장/대조의 단일 형식."""
    return [[jsonable(shape_sig(sh)) for sh in sl.shapes] for sl in prs.slides]


def compare_slide(golden_sigs, fab_slide, name):
    """기준(스냅샷) 서명 vs 지금 빌드한 슬라이드."""
    gs = golden_sigs
    fs = [jsonable(shape_sig(s)) for s in fab_slide.shapes]
    rows, mism = [], 0
    if len(gs) != len(fs):
        rows.append(f"| {name} | 도형 수 | {len(gs)} != {len(fs)} | FAIL |")
        mism += 1
    n = min(len(gs), len(fs))
    ok = 0
    for i in range(n):
        bad = diff_shapes(gs[i], fs[i])
        if bad:
            mism += 1
            rows.append(f"| {name} | 도형#{i} | {'; '.join(bad)[:120]} | FAIL |")
        else:
            ok += 1
    rows.insert(0, f"| {name} | 전체 | {ok}/{len(gs)} 일치 | {'PASS' if mism == 0 else 'FAIL'} |")
    return rows, mism, ok, len(gs)


def build_current():
    """지금 코드(goldenfab)로 레퍼런스 덱을 조립. 기준선은 스냅샷 파일이지 다른 코드 트리가 아니다."""
    from goldenfab.reference import build_reference, slide_names

    return build_reference(), slide_names()


def load_snapshot():
    if not SNAPSHOT.exists():
        raise SystemExit(
            f"기준선 스냅샷이 없다: {SNAPSHOT}\n"
            "  최초 생성: uv run python .../compare_golden.py --update-snapshot"
        )
    data = json.loads(SNAPSHOT.read_text(encoding="utf-8"))
    return data["slides"], data.get("names", [])


def save_snapshot(prs, names):
    sigs = deck_signature(prs)
    SNAPSHOT.parent.mkdir(parents=True, exist_ok=True)
    SNAPSHOT.write_text(
        json.dumps(
            {
                "_comment": (
                    "골든 레퍼런스 덱 도형 서명 기준선 — compare_golden 회귀 게이트의 정답지. "
                    "goldenfab/reference.py가 만드는 19장을 박제한 것. 손으로 고치지 말 것 — "
                    "의도적 레이아웃 변경 후 `compare_golden.py --update-snapshot`으로 재생성하고 "
                    "git diff로 무엇이 바뀌었는지 리뷰한다."
                ),
                "slides": sigs,
                "names": names,
            },
            ensure_ascii=False,
            indent=1,
        ),
        encoding="utf-8",
    )
    return sum(len(s) for s in sigs)


def param_efficacy():
    """파라미터 유효성 — 상이 입력이 출력 텍스트에 반영되는지 2케이스."""
    from goldenfab.registry import LAYOUTS as FAB

    results = []
    f = Presentation()
    f.slide_width, f.slide_height = Inches(13.333), Inches(7.5)
    FAB["closing"](
        f,
        {
            "value_prop": "테스트 가치제안 문장 A" + chr(10) + "둘째 줄 B",
            "kicker": "CLOSING — 테스트덱",
        },
    )
    texts = " ".join(sh.text_frame.text for sh in f.slides[0].shapes if sh.has_text_frame)
    results.append(
        ("closing 커스텀", "테스트 가치제안 문장 A" in texts and "CLOSING — 테스트덱" in texts)
    )

    f2 = Presentation()
    f2.slide_width, f2.slide_height = Inches(13.333), Inches(7.5)
    FAB["screenshot"](f2, {"headline": "커스텀 헤드라인 검사용", "caption": "커스텀 캡션 — 검사"})
    texts2 = " ".join(sh.text_frame.text for sh in f2.slides[0].shapes if sh.has_text_frame)
    results.append(
        ("screenshot 커스텀", "커스텀 헤드라인 검사용" in texts2 and "커스텀 캡션 — 검사" in texts2)
    )

    # Phase B 파라미터화 10종 — headline override가 출력에 반영되는지(내장 아님) 검증
    for t in (
        "problem_grid",
        "exec_graph",
        "tech_evidence",
        "tech_tree",
        "tech_mechanism",
        "tech_capture",
        "ab_simulation",
        "validation",
        "mirror_matrix",
        "boundary",
    ):
        ft = Presentation()
        ft.slide_width, ft.slide_height = Inches(13.333), Inches(7.5)
        sent = "SENT_" + t.upper()
        FAB[t](ft, {"headline": sent})
        tt = " ".join(sh.text_frame.text for sh in ft.slides[0].shapes if sh.has_text_frame)
        results.append((f"{t} override", sent in tt))
    return results


def main():
    args = [a for a in sys.argv[1:] if not a.startswith("--")]
    prs, names = build_current()

    if "--update-snapshot" in sys.argv:
        n = save_snapshot(prs, names)
        print(f"기준선 재생성: {SNAPSHOT}\n  슬라이드 {len(names)} · 도형 {n}")
        print("  git diff로 무엇이 바뀌었는지 반드시 리뷰할 것.")
        return

    out_path = Path(args[0]) if args else ROOT / "golden" / "variants" / "compare_full.md"
    snap_slides, snap_names = load_snapshot()
    lines = [
        "# compare_golden — 골든 레퍼런스 덱 vs 기준선 스냅샷 대조",
        "",
        f'허용오차 ±{TOL}" · 비교 속성: 수/type/text/fill/run(pt·bold)/run_colors/cells/l·t·w·h',
        f"기준선: `{SNAPSHOT.relative_to(ROOT)}` · 소스: `goldenfab/reference.py`",
        "",
        "| 슬라이드 | 대상 | 내용 | 판정 |",
        "|---|---|---|---|",
    ]
    total_mism = total_ok = total_n = 0
    if len(snap_slides) != len(names):
        lines.append(
            f"| 전체 | 슬라이드 수 | 기준선 {len(snap_slides)} != 현재 {len(names)} | FAIL |"
        )
        total_mism += 1
    if snap_names and snap_names != names:
        # 장 순서·구성이 바뀌었는데 도형 수가 우연히 맞으면 도형 대조만으론 못 잡는다
        lines.append("| 전체 | 장 구성 | 기준선 != 현재 (SLIDE_ORDER 변경) | FAIL |")
        total_mism += 1
    for gs, fs2, name in zip(snap_slides, prs.slides, names):
        rows, mism, ok, n = compare_slide(gs, fs2, name)
        lines.extend(rows)
        total_mism += mism
        total_ok += ok
        total_n += n
    lines.append("")
    lines.append(
        f"총계: 슬라이드 {len(names)}/{len(snap_slides)}, 도형 일치 {total_ok}/{total_n}, "
        f"불일치 {total_mism}건"
    )
    lines.append("")
    lines.append("## 파라미터 유효성 (상이 입력 반영)")
    par = param_efficacy()
    for name, ok in par:
        lines.append(f"- {name}: {'PASS' if ok else 'FAIL'}")
    out_path.parent.mkdir(parents=True, exist_ok=True)
    out_path.write_text(chr(10).join(lines) + chr(10), encoding="utf-8")
    print(chr(10).join(lines[-30:]))
    assert total_mism == 0, f"COMPARE FAIL {total_mism}건"
    assert all(ok for _, ok in par), "PARAM FAIL"
    print(f"COMPARE PASS -> {out_path}")


if __name__ == "__main__":
    main()
