#!/usr/bin/env python
"""deck-spec(JSON) + brand-kit(YAML) -> 네이티브 .pptx 빌드.

표는 네이티브 표(GraphicFrame.table), 차트는 네이티브 차트(GraphicFrame.chart)로
생성한다 — 이미지가 아니므로 파워포인트에서 데이터만 고치면 그대로 반영된다.

브랜드 일관성은 brand-kit.yaml 단일 출처에서만 나온다. 레이아웃 상수도 이 파일에
박제되어 있어, 같은 spec은 항상 같은 결과를 낸다("공장 일관성").

usage: python build_pptx.py <deck-spec.json> <out.pptx> [--brand brand-kit.yaml]
"""

import argparse
import json
import re
import sys
from pathlib import Path

import yaml
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# 차트·다이어그램 렌더는 pptx-visuals 스킬이 단일 출처
sys.path.insert(0, str(Path(__file__).resolve().parents[2] / "pptx-visuals" / "scripts"))
import mpl_exhibits  # noqa: E402
import visuals  # noqa: E402

# 16:9 고정 캔버스 (inch)
EMU_W, EMU_H = 13.333, 7.5

# 마스터 틀(컨설팅 표준) 고정 좌표 — 모든 슬라이드에 동일 위치로 스탬프
FOOT_Y = EMU_H - 0.42  # 푸터 텍스트 y (회사명·페이지번호)
HAIR_Y = EMU_H - 0.5  # 푸터 헤어라인 y
HEADER_HAIR_Y = 1.42  # 제목(+부제) 아래 헤더 구분선 y (파트 없는 구형 레이아웃)
BODY_TOP = 1.72  # 본문 시작 y (파트 없는 구형 레이아웃)
# 내비게이션 헤더(우석진 실측 + v4.5 헤드라인 2줄 허용)
TAB_Y = 0.2  # 섹션 탭 y
KICKER_Y = 0.62  # 소분류 라벨 y
TITLE_Y = 0.82  # 계층번호 헤드라인 y (28pt 2줄 = 0.95in 확보)
SUB_Y = 1.80  # 리드 문단 y (14pt 2줄 = 0.55in 확보)
NAV_RULE_Y = 2.42  # 헤더 하단 헤어라인 y
NAV_BODY_TOP = 2.56  # 파트 체계 사용 시 본문 시작 y
ROMAN = ["Ⅰ", "Ⅱ", "Ⅲ", "Ⅳ", "Ⅴ", "Ⅵ"]
SRC_Y = EMU_H - 0.85  # 출처선 y (각주 위 — 겹침 방지)
FOOTNOTE_Y = EMU_H - 0.62  # 각주 y (푸터 헤어라인 위)
PAGE_W = 2.0  # 페이지번호 박스 폭 (우측 정렬)
BODY_BOTTOM = EMU_H - 0.95  # 본문(개체) 하한 — 출처·각주 영역 침범 금지
COMMENT_W = 4.4  # 해설 칼럼 폭 (개체 좌측)


def load_brand(path):
    b = yaml.safe_load(Path(path).read_text(encoding="utf-8"))
    b.setdefault("brand", {}).setdefault("name", "My Company")
    return b


def rgb(hex6):
    return RGBColor.from_string(hex6)


class Deck:
    """brand-kit을 들고 슬라이드 타입별 렌더러를 제공하는 빌더."""

    def __init__(self, brand):
        self.b = brand
        self.c = brand["colors"]
        self.f = brand["fonts"]
        self.s = brand["sizes"]
        self.margin = brand["layout"]["margin"]
        self.logo = brand["layout"].get("logo") or ""
        self.prs = Presentation()
        self.prs.slide_width = Inches(EMU_W)
        self.prs.slide_height = Inches(EMU_H)
        self.blank = self.prs.slide_layouts[6]
        self.parts = []  # [{no,title,items:[제목…]}] — part 타입 pre-pass로 채움
        self.body_top = BODY_TOP  # 파트 체계 사용 시 NAV_BODY_TOP으로 전환
        self._page_no = 0
        self._total = 0
        self.meta = {}
        self._comment_side = "left"  # 해설 칼럼 좌/우 교차 토글(형식 반복 방지)
        self.render_dir = Path("_workspace/render")  # mpl 익스히빗 PNG 출력처(main이 재설정)
        # 장 스크립트(adapted/novel) 상대경로 기준 — main이 spec 위치로 설정
        self._spec_dir: Path | None = None

    # --- 저수준 헬퍼 ---
    def _slide(self, bg=None):
        s = self.prs.slides.add_slide(self.blank)
        fill = s.background.fill
        fill.solid()
        fill.fore_color.rgb = rgb(bg or self.c["bg"])
        return s

    def _text(
        self,
        slide,
        x,
        y,
        w,
        h,
        text,
        size,
        color,
        *,
        font=None,
        bold=False,
        align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP,
    ):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = font or self.f["body"]
        run.font.color.rgb = rgb(color)
        return tb

    @staticmethod
    def _parse_emph(text):
        """`**x**` 구간을 강조(True)로 분리. BCG식 인라인 데이터 강조용."""
        parts, i = [], 0
        for mt in re.finditer(r"\*\*(.+?)\*\*", text):
            if mt.start() > i:
                parts.append((text[i : mt.start()], False))
            parts.append((mt.group(1), True))
            i = mt.end()
        if i < len(text):
            parts.append((text[i:], False))
        return parts or [(text, False)]

    def _rich_text(self, slide, x, y, w, h, text, size, color, *, font=None, align=PP_ALIGN.LEFT):
        """본문 텍스트 — `**핵심**`은 그 run만 bold + accent색으로. 나머지는 color."""
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        for seg, emph in self._parse_emph(text):
            run = p.add_run()
            run.text = seg
            run.font.size = Pt(size)
            run.font.name = font or self.f["body"]
            run.font.bold = bool(emph)
            run.font.color.rgb = rgb(self.c["accent"] if emph else color)
        return tb

    def _accent_bar(self, slide, x, y, w=2.2, h=0.09, color=None):
        from pptx.enum.shapes import MSO_SHAPE

        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = rgb(color or self.c["accent"])
        bar.line.fill.background()
        return bar

    def _hline(self, slide, x1, x2, y, color, weight=0.75):
        """가는 수평 구분선(헤어라인)."""
        from pptx.enum.shapes import MSO_CONNECTOR

        ln = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y), Inches(x2), Inches(y)
        )
        ln.line.color.rgb = rgb(color)
        ln.line.width = Pt(weight)
        return ln

    def _title_block(self, slide, title, subtitle=None, spec=None):
        """헤더. 파트 체계 사용 시(우석진 실측): 섹션 탭 → 소분류 kicker → 계층번호 제목
        → accent 틱 거버닝 메시지 → 헤어라인. 파트 없으면 구형(제목+이중룰) 유지."""
        m = self.margin
        if self.parts:
            self._nav_tabs(slide, (spec or {}).get("_part_no"))
            part_no = (spec or {}).get("_part_no")
            if part_no:
                part = self.parts[part_no - 1]
                self._text(
                    slide,
                    m,
                    KICKER_Y,
                    EMU_W - 2 * m,
                    0.28,
                    f"{part_no}. {part['title']}",
                    10,
                    self.c["muted"],
                )
            num = (spec or {}).get("_num")
            title_text = f"{num}. {title}" if num else title
            self._text(
                slide,
                m,
                TITLE_Y,
                EMU_W - 2 * m,
                0.95,
                title_text,
                self.s["section"],
                self.c["primary"],
                font=self.f["head"],
                bold=True,
            )
            if subtitle:
                # 거버닝 메시지 좌측 accent 틱 — 전 장 반복되는 시각 모티프
                self._accent_bar(slide, m, SUB_Y + 0.05, w=0.05, h=0.2)
                self._text(
                    slide,
                    m + 0.16,
                    SUB_Y,
                    EMU_W - 2 * m - 0.16,
                    0.58,
                    subtitle,
                    self.s.get("sub", 15),
                    self.c["text"],
                    font=self.f["body"],
                )
            self._hline(slide, m, EMU_W - m, NAV_RULE_Y, self.c["muted"], weight=0.75)
            return
        self._text(
            slide,
            m,
            0.32,
            EMU_W - 2 * m,
            # 0.52: 제목 상자 바닥(0.84)이 부제 top(0.86) **아래**여야 한다. 0.6이면 0.06"
            # 겹쳐 check_text_collision이 레거시 전 장에서 상시 적색이었다(2026-07-25 실측:
            # 팔레트 18/18). 텍스트는 TOP 앵커라 상자만 줄이는 건 렌더상 무변화.
            0.52,
            title,
            self.s["section"],
            self.c["primary"],
            font=self.f["head"],
            bold=True,
        )
        if subtitle:
            self._text(
                slide,
                m,
                0.86,
                EMU_W - 2 * m,
                0.45,
                subtitle,
                self.s.get("sub", 15),
                self.c["muted"],
                font=self.f["body"],
            )
        self._hline(slide, m, EMU_W - m, HEADER_HAIR_Y, self.c["primary"], weight=2.5)
        self._hline(slide, m, EMU_W - m, HEADER_HAIR_Y + 0.05, self.c["muted"], weight=0.75)

    def _nav_tabs(self, slide, cur_part_no):
        """상단 섹션 탭(현재 위치 하이라이트) + 우측 문서명·페이지 + 전폭 룰."""
        from pptx.enum.shapes import MSO_SHAPE

        m = self.margin
        info_w = 3.8
        tab_w = min(2.3, (EMU_W - 2 * m - info_w) / max(len(self.parts), 1))
        for p in self.parts:
            x = m + (p["no"] - 1) * tab_w
            label = f"{ROMAN[p['no'] - 1]}. {p['title']}"
            if p["no"] == cur_part_no:
                chip = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Inches(x), Inches(TAB_Y), Inches(tab_w), Inches(0.3)
                )
                chip.fill.solid()
                chip.fill.fore_color.rgb = rgb(self.c["bg_alt"])
                chip.line.fill.background()
                self._text(
                    slide,
                    x,
                    TAB_Y + 0.02,
                    tab_w,
                    0.26,
                    label,
                    10,
                    self.c["primary"],
                    bold=True,
                    align=PP_ALIGN.CENTER,
                )
                self._hline(slide, x, x + tab_w, TAB_Y + 0.31, self.c["accent"], weight=2.0)
            else:
                self._text(
                    slide,
                    x,
                    TAB_Y + 0.02,
                    tab_w,
                    0.26,
                    label,
                    10,
                    self.c["muted"],
                    align=PP_ALIGN.CENTER,
                )
        self._text(
            slide,
            EMU_W - m - info_w,
            TAB_Y + 0.02,
            info_w,
            0.26,
            f"{self.meta.get('project', '')}   {self._page_no:02d} / {self._total:02d}",
            self.s["caption"],
            self.c["muted"],
            align=PP_ALIGN.RIGHT,
        )
        self._hline(slide, m, EMU_W - m, TAB_Y + 0.34, self.c["primary"], weight=1.2)

    def _num_circle(self, slide, x, y, num, *, d=0.44, dark_bg=False):
        """넘버 서클 — 간지·목차·불릿에 반복되는 시각 모티프(accent 테두리 원 + 번호)."""
        from pptx.enum.shapes import MSO_SHAPE

        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
        c.fill.background()
        c.line.color.rgb = rgb(self.c["accent"])
        c.line.width = Pt(1.5)
        tf = c.text_frame
        tf.word_wrap = False
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = num
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.name = self.f["head"]
        run.font.color.rgb = rgb("FFFFFF" if dark_bg else self.c["primary"])
        return c

    def _footnotes(self, slide, notes):
        """각주(7~8pt) — 참조·단서를 푸터 헤어라인 위에. BCG식 하단 각주."""
        txt = (
            "   ".join(f"{i}. {n}" for i, n in enumerate(notes, 1))
            if len(notes) > 1
            else str(notes[0])
        )
        self._text(
            slide,
            self.margin,
            FOOTNOTE_Y,
            EMU_W - 2 * self.margin,
            0.3,
            txt,
            self.s.get("foot", 8),
            self.c["muted"],
            font=self.f["body"],
        )

    def _next_side(self):
        """개체 슬라이드 해설 칼럼의 좌/우를 결정적으로 교차 — 연속 동일 배치 방지."""
        side = self._comment_side
        self._comment_side = "right" if side == "left" else "left"
        return side

    def _intro(self, slide, spec):
        """intro 서술 문단(2~4문장, 전폭) — 컨설팅 보고서식 리드. 본문 y 오프셋 반환."""
        text = spec.get("intro")
        if not text:
            return 0.0
        tb = slide.shapes.add_textbox(
            Inches(self.margin),
            Inches(self.body_top + 0.05),
            Inches(EMU_W - 2 * self.margin),
            Inches(0.85),
        )
        tf = tb.text_frame
        tf.word_wrap = True
        self._emph_runs(tf.paragraphs[0], text, self.s["body"], self.c["text"])
        return 1.0

    def _commentary(self, slide, spec, *, below=None, side="left", y_off=0.0):
        """해설 칼럼/블록 — 개체(차트·표·다이어그램)에 분석 텍스트를 붙여 문서 밀도를 만든다.
        below=None이면 side 방향 칼럼(개체는 반대편), below=y이면 그 아래 전폭 블록.
        한 텍스트박스 안에서 문단으로 흘려 줄바꿈이 겹치지 않게 한다(자동 플로우)."""
        items = spec.get("commentary", [])
        if not items:
            return False
        m = self.margin
        if below is None:
            x = m if side == "left" else EMU_W - m - (COMMENT_W - 0.3)
            y = self.body_top + y_off + 0.12
            w = COMMENT_W - 0.3
            h = BODY_BOTTOM - y - 0.08
        else:
            x, y, w, h = m, below, EMU_W - 2 * m, BODY_BOTTOM - below
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        size = self.s["body"]
        first = True
        for it in items:
            text = it if isinstance(it, str) else it.get("text", "")
            subs = [] if isinstance(it, str) else it.get("sub", [])
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_before = Pt(0 if len(tf.paragraphs) == 1 else 8)
            p.space_after = Pt(2)
            self._emph_runs(p, f"• {text}", size, self.c["text"])
            for sub in subs:
                p2 = tf.add_paragraph()
                p2.space_after = Pt(1)
                self._emph_runs(p2, f"   – {sub}", size - 1, self.c["text"])
        return True

    def _emph_runs(self, para, text, size, color):
        """문단에 `**강조**` 지원 run들을 채운다(_rich_text의 문단 버전)."""
        for seg, emph in self._parse_emph(text):
            run = para.add_run()
            run.text = seg
            run.font.size = Pt(size)
            run.font.name = self.f["body"]
            run.font.bold = bool(emph)
            run.font.color.rgb = rgb(self.c["accent"] if emph else color)

    def _flatten(self, slide):
        """전 도형 그림자 제거 — PowerPoint 기본 그림자가 2010년대 PPT 티의 주범."""
        for sh in slide.shapes:
            try:
                sh.shadow.inherit = False
            except (AttributeError, ValueError, NotImplementedError):
                continue  # GraphicFrame(표·차트)은 미지원 — 자체 그림자 없음

    def _frame(self, slide, page_no, total):
        """마스터 틀(기관 문서형 푸터): 헤어라인 + 좌 문서명 + 중앙 브랜드 워드마크
        + 우하단 페이지번호 + 우측 세로 저작권선. 상수 좌표라 전 장 동일 위치 스탬프."""
        m = self.margin
        self._hline(slide, m, EMU_W - m, HAIR_Y, self.c["muted"], weight=0.5)
        # 좌: 문서/프로젝트명
        self._text(
            slide,
            m,
            FOOT_Y,
            4.0,
            0.32,
            self.meta.get("project", ""),
            self.s["caption"],
            self.c["muted"],
        )
        # 중앙: 브랜드 워드마크
        self._text(
            slide,
            EMU_W / 2 - 3,
            FOOT_Y,
            6.0,
            0.32,
            self.b["brand"]["name"],
            self.s["caption"],
            self.c["muted"],
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        # 우: 페이지번호
        self._text(
            slide,
            EMU_W - m - PAGE_W,
            FOOT_Y,
            PAGE_W,
            0.32,
            f"{page_no:02d} / {total:02d}",
            self.s["caption"],
            self.c["muted"],
            align=PP_ALIGN.RIGHT,
        )
        self._vcopyright(slide)

    def _vcopyright(self, slide):
        """우측 여백에 세로로 흐르는 저작권선 — 실제 기관 덱의 대표 시그니처."""
        yr = self.meta.get("year", "")
        txt = f"© {yr} {self.b['brand']['name']}. All rights reserved.".replace("©  ", "© ")
        tb = slide.shapes.add_textbox(Inches(EMU_W - 2.6), Inches(3.55), Inches(4.6), Inches(0.28))
        tb.rotation = 270
        tf = tb.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(7)
        run.font.name = self.f["body"]
        run.font.color.rgb = rgb(self.c["muted"])

    def _source(self, slide, text):
        """출처선: 정량 슬라이드의 좌하단(푸터 위). 셀링 신뢰의 근거 표기."""
        self._text(
            slide,
            self.margin,
            SRC_Y,
            EMU_W - 2 * self.margin,
            0.28,
            f"출처: {text}",
            self.s["caption"],
            self.c["muted"],
        )

    def _logo_mark(self, slide, color=None):
        txt = self.b["brand"]["name"]
        self._text(
            slide,
            self.margin,
            EMU_H - 0.55,
            5,
            0.4,
            txt,
            self.s["caption"],
            color or self.c["muted"],
            bold=True,
        )

    # --- 슬라이드 타입별 렌더러 ---

    def section(self, spec):
        s = self._slide(bg=self.c["bg_alt"])
        m = self.margin
        self._accent_bar(s, m, 3.15, w=1.2, h=0.1)
        self._text(
            s,
            m,
            3.35,
            EMU_W - 2 * m,
            1.2,
            spec["title"],
            self.s["title"],
            self.c["primary"],
            font=self.f["head"],
            bold=True,
        )
        return s

    def bullets(self, spec):
        s = self._slide()
        self._title_block(s, spec["title"], spec.get("subtitle"), spec)
        m = self.margin
        y = self.body_top + self._intro(s, spec) + 0.08
        for i, bt in enumerate(spec.get("bullets", []), 1):
            text = bt if isinstance(bt, str) else bt.get("text", "")
            subs = [] if isinstance(bt, str) else bt.get("sub", [])
            icon = None if isinstance(bt, str) else bt.get("icon")
            # 리드 마커: 아이콘 지정 시 픽토그램, 아니면 넘버 서클 모티프
            if icon:
                visuals.add_icon(s, icon, m, y - 0.04, 0.36, color="accent")
            else:
                self._num_circle(s, m, y - 0.04, f"{i:02d}", d=0.36)
            self._rich_text(
                s,
                m + 0.52,
                y,
                EMU_W - 2 * m - 0.52,
                0.5,
                text,
                self.s["body"],
                self.c["text"],
            )
            y += 0.46
            for sub in subs:  # 2단계 중첩 (BCG식 – 서브불릿)
                self._text(s, m + 0.72, y, 0.3, 0.3, "–", self.s["body"] - 1, self.c["muted"])
                self._rich_text(
                    s,
                    m + 1.0,
                    y,
                    EMU_W - 2 * m - 1.0,
                    0.4,
                    sub,
                    self.s["body"] - 1,
                    self.c["text"],
                )
                y += 0.36
            y += 0.1
        return s

    def two_column(self, spec):
        s = self._slide()
        self._title_block(s, spec["title"], spec.get("subtitle"), spec)
        m = self.margin
        off = self._intro(s, spec)
        top = self.body_top + off
        colw = (EMU_W - 2 * m - 0.6) / 2
        for idx, key in enumerate(("left", "right")):
            col = spec.get(key, {})
            x = m + idx * (colw + 0.6)
            self._text(
                s,
                x,
                top,
                colw,
                0.4,
                col.get("heading", ""),
                self.s["head"],
                self.c["primary"],
                font=self.f["head"],
                bold=True,
            )
            self._hline(s, x, x + colw, top + 0.44, self.c["primary"], weight=1.5)
            y = top + 0.62
            for it in col.get("items", []):
                text = it if isinstance(it, str) else it.get("text", "")
                subs = [] if isinstance(it, str) else it.get("sub", [])
                self._text(s, x, y, 0.3, 0.35, "•", self.s["body"], self.c["muted"])
                self._rich_text(
                    s,
                    x + 0.28,
                    y,
                    colw - 0.28,
                    0.5,
                    text,
                    self.s["body"],
                    self.c["text"],
                )
                y += 0.44
                for sub in subs:
                    self._text(
                        s,
                        x + 0.5,
                        y,
                        0.3,
                        0.3,
                        "–",
                        self.s["body"] - 1,
                        self.c["muted"],
                    )
                    self._rich_text(
                        s,
                        x + 0.8,
                        y,
                        colw - 0.8,
                        0.4,
                        sub,
                        self.s["body"] - 1,
                        self.c["text"],
                    )
                    y += 0.34
                y += 0.06
        return s

    def metrics(self, spec):
        s = self._slide()
        self._title_block(s, spec["title"], spec.get("subtitle"), spec)
        items = spec.get("items", [])[:4]
        m = self.margin
        n = max(len(items), 1)
        colw = (EMU_W - 2 * m) / n
        # 해설이 있으면 카드를 상단으로 올리고 아래 전폭 해설 블록
        top = self.body_top + (0.25 if spec.get("commentary") else 0.75)
        icon_off = 0.6 if any(isinstance(it, dict) and it.get("icon") for it in items) else 0
        for i, it in enumerate(items):
            x = m + i * colw
            if it.get("icon"):
                visuals.add_icon(
                    s, it["icon"], x + colw / 2 - 0.25, top - 0.08, 0.5, color="accent"
                )
            # 숫자는 primary(잉크) — hero 오렌지 카드 대신 절제된 문서형
            self._text(
                s,
                x,
                top + icon_off,
                colw,
                1.1,
                str(it.get("value", "")),
                self.s["title"],
                self.c["primary"],
                font=self.f["head"],
                bold=True,
                align=PP_ALIGN.CENTER,
            )
            # 값 아래 짧은 accent 밑줄 — 포인트는 여기 극소량만
            self._hline(
                s,
                x + colw * 0.32,
                x + colw * 0.68,
                top + icon_off + 1.18,
                self.c["accent"],
                weight=2.0,
            )
            self._text(
                s,
                x,
                top + icon_off + 1.3,
                colw,
                0.7,
                it.get("label", ""),
                self.s["body"],
                self.c["muted"],
                align=PP_ALIGN.CENTER,
            )
        self._commentary(s, spec, below=top + icon_off + 2.25)
        return s

    def table(self, spec):
        s = self._slide()
        self._title_block(s, spec["title"], spec.get("subtitle"), spec)
        headers = spec["headers"]
        rows = spec["rows"]
        matrix = spec.get("style") == "matrix"  # 비교 매트릭스: 마지막 열(제안) 강조
        off = self._intro(s, spec)
        # 해설 칼럼이 있으면 표는 반대편으로(좌/우 교차) — 남는 공백을 분석 텍스트로 채운다
        x, w = self._exhibit_zone(s, spec, off)
        nrows, ncols = len(rows) + 1, len(headers)
        # 본문 영역을 채우도록 행 높이를 늘린다(상단 고정·하단 공백 방지)
        avail = BODY_BOTTOM - (self.body_top + off + 0.12)
        row_h = min(0.68, max(0.42, avail / nrows))
        ty = self.body_top + off + 0.12
        gf = s.shapes.add_table(
            nrows,
            ncols,
            Inches(x),
            Inches(ty),
            Inches(w),
            Inches(row_h * nrows),
        )
        tbl = gf.table
        for r in range(nrows):
            tbl.rows[r].height = Inches(row_h)
        # v4.5 헤어라인 스타일(McKinsey 실측) — 검정 헤더 채움·얼룩말 줄무늬 제거,
        # 위계는 상단 강한 룰 + 헤더 아래·행 사이 헤어라인으로만.
        for j, h in enumerate(headers):
            cell = tbl.cell(0, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(self.c["bg"])
            last = matrix and j == ncols - 1
            self._cell_text(
                cell,
                h,
                self.c["accent"] if last else self.c["primary"],
                bold=True,
                align=PP_ALIGN.CENTER if (matrix and j) else PP_ALIGN.LEFT,
            )
        for i, row in enumerate(rows, 1):
            for j, val in enumerate(row):
                cell = tbl.cell(i, j)
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb(self.c["bg"])
                last = matrix and j == ncols - 1
                self._cell_text(
                    cell,
                    str(val),
                    self.c["primary"] if last else self.c["text"],
                    bold=last,
                    align=PP_ALIGN.CENTER if (matrix and j) else PP_ALIGN.LEFT,
                )
        self._hline(s, x, x + w, ty, self.c["primary"], weight=1.5)
        self._hline(s, x, x + w, ty + row_h, self.c["primary"], weight=0.75)
        for i in range(2, nrows):
            self._hline(s, x, x + w, ty + row_h * i, self.c["bg_alt"], weight=0.75)
        self._hline(s, x, x + w, ty + row_h * nrows, self.c["muted"], weight=0.75)
        return s

    def _cell_text(self, cell, text, color, bold=False, align=PP_ALIGN.LEFT):
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(self.s["body"] - 3)
        run.font.bold = bold
        run.font.name = self.f["body"]
        run.font.color.rgb = rgb(color)

    def _exhibit_zone(self, slide, spec, off):
        """해설 칼럼(좌/우 교차) 렌더 후 개체가 쓸 (x, w)를 돌려준다."""
        m = self.margin
        side = self._next_side() if spec.get("commentary") else None
        has_comment = self._commentary(slide, spec, side=side or "left", y_off=off)
        if has_comment:
            x = m + COMMENT_W if side == "left" else m
            return x, EMU_W - 2 * m - COMMENT_W
        return m, EMU_W - 2 * m

    def _banner(self, slide, text):
        """하단 결론 배너(BCG p28 실측) — 전폭 primary 스트립. source·footnotes와 병용 금지."""
        from pptx.enum.shapes import MSO_SHAPE

        m = self.margin
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(m),
            Inches(EMU_H - 0.92),
            Inches(EMU_W - 2 * m),
            Inches(0.36),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = rgb(self.c["primary"])
        bar.line.fill.background()
        tf = bar.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        for seg, emph in self._parse_emph(text):
            run = p.add_run()
            run.text = seg
            run.font.size = Pt(self.s["body"])
            run.font.bold = True
            run.font.name = self.f["head"]
            run.font.color.rgb = rgb("FFFFFF" if not emph else self.c["accent"])

    def chart(self, spec):
        s = self._slide()
        self._title_block(s, spec["title"], spec.get("subtitle"), spec)
        off = self._intro(s, spec)
        # 멀티패널(BCG p8 실측): 미니 차트 2~3개 나란히, 각자 소제목+밑줄
        if spec.get("panels"):
            m = self.margin
            panels = spec["panels"][:3]
            n = len(panels)
            pw = (EMU_W - 2 * m - 0.5 * (n - 1)) / n
            top = self.body_top + off + 0.1
            sub = spec.get("sub_table")
            sub_h = min(1.3, 0.34 * (len(sub["rows"]) + 1)) if sub else 0.0
            panel_bottom = BODY_BOTTOM - (sub_h + 0.15 if sub else 0)
            for i, pn in enumerate(panels):
                px = m + i * (pw + 0.5)
                self._text(
                    s,
                    px,
                    top,
                    pw,
                    0.55,
                    pn["title"],
                    self.s.get("sub", 15) - 1,
                    self.c["primary"],
                    font=self.f["head"],
                    bold=True,
                    align=PP_ALIGN.CENTER,
                )
                self._hline(
                    s, px + pw * 0.2, px + pw * 0.8, top + 0.6, self.c["primary"], weight=1.5
                )
                visuals.add_chart(s, pn, self.b, px, top + 0.75, pw, panel_bottom - top - 0.85)
            if sub:
                self._mini_table(s, sub, m, panel_bottom + 0.15, EMU_W - 2 * m, sub_h)
            return s
        x, w = self._exhibit_zone(s, spec, off)
        cy = self.body_top + off + 0.12
        ch = BODY_BOTTOM - self.body_top - off - 0.12
        sub = spec.get("sub_table")  # 차트 아래 부속 데이터 행(BCG p18) — 한 장 안 조합
        sub_h = 0.0
        if sub:
            sub_h = min(1.5, 0.34 * (len(sub["rows"]) + 1))
            ch -= sub_h + 0.15
        if self._is_image_exhibit(spec):
            # 하이브리드 이미지 경로: 네이티브가 못 만드는 유형은 mpl PNG(수치편집불가).
            # 콜아웃·각주는 mpl이 이미지 안에 굽는다 — 도형 오버레이와 이중 표기 금지.
            png = self.render_dir / f"exhibit_{self._page_no:02d}_{spec['chart_type']}.png"
            mpl_exhibits.render(spec, self.b, png, w_in=w, h_in=ch)
            # 종횡비 보존 삽입 — tight crop된 PNG를 존에 강제 스트레치하면 원이 타원이 된다
            from PIL import Image

            iw, ih = Image.open(png).size
            if iw / ih >= w / ch:
                pic_h = w * ih / iw
                s.shapes.add_picture(
                    str(png), Inches(x), Inches(cy + (ch - pic_h) / 2), Inches(w), Inches(pic_h)
                )
            else:
                pic_w = ch * iw / ih
                s.shapes.add_picture(
                    str(png), Inches(x + (w - pic_w) / 2), Inches(cy), Inches(pic_w), Inches(ch)
                )
        else:
            visuals.add_chart(s, spec, self.b, x, cy, w, ch)
            if spec.get("annotations"):
                # 네이티브 경로 콜아웃은 pptx 도형 — PPT에서 문구·위치 편집 가능
                visuals.add_callouts(s, spec["annotations"], self.b, x, cy, w, ch)
        if sub:
            self._mini_table(s, sub, x, cy + ch + 0.15, w, sub_h)
        return s

    def _mini_table(self, slide, sub, x, y, w, h):
        """차트 부속 데이터 행 — 네이티브 미니 표(수치편집가능). headers+rows 1~3행."""
        headers, rows = sub["headers"], sub["rows"]
        nrows, ncols = len(rows) + 1, len(headers)
        gf = slide.shapes.add_table(nrows, ncols, Inches(x), Inches(y), Inches(w), Inches(h))
        tbl = gf.table
        row_h = h / nrows
        for j, htxt in enumerate(headers):
            cell = tbl.cell(0, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(self.c["bg"])
            self._cell_text(cell, str(htxt), self.c["primary"], bold=True)
        for i, row in enumerate(rows, 1):
            for j, val in enumerate(row):
                cell = tbl.cell(i, j)
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb(self.c["bg"])
                self._cell_text(cell, str(val), self.c["text"])
        self._hline(slide, x, x + w, y, self.c["primary"], weight=1.0)
        self._hline(slide, x, x + w, y + row_h, self.c["muted"], weight=0.75)
        self._hline(slide, x, x + w, y + h, self.c["muted"], weight=0.75)

    def _is_image_exhibit(self, spec):
        """mpl 이미지 경로 판별 — 확장 유형이거나 render:image 명시."""
        return spec.get("chart_type") in mpl_exhibits.MPL_TYPES or spec.get("render") == "image"

    def diagram(self, spec):
        s = self._slide()
        self._title_block(s, spec["title"], spec.get("subtitle"), spec)
        m = self.margin
        off = self._intro(s, spec)
        # 2026-07-25: 좌/우 교차 배치 분기(layers·branch·cards)는 그 어휘가 폐기돼 삭제했다 —
        # 남은 9종은 컴포지트(본문 전폭 사용) 또는 전폭 가로 둘로 갈린다.
        if spec.get("layout") in ("icon_rows", "stat_split", "split_detail"):
            # v4.5 컴포지트 조판 — 자체 밀도가 높으므로 본문 영역 전체 사용(하단 공백 금지)
            avail = BODY_BOTTOM - (self.body_top + off + 0.12)
            visuals.add_diagram(
                s, spec, self.b, m, self.body_top + off + 0.12, EMU_W - 2 * m, avail
            )
        else:
            # timeline·matrix_2x2·spectrum·venn·harvey_table·check_matrix는 전폭 가로 —
            # 다이어그램 상단, 해설은 아래 전폭 블록
            avail = BODY_BOTTOM - (self.body_top + off + 0.12)
            diag_h = 2.35 if spec.get("commentary") else min(4.2, avail)
            visuals.add_diagram(
                s, spec, self.b, m, self.body_top + off + 0.12, EMU_W - 2 * m, diag_h
            )
            self._commentary(s, spec, below=self.body_top + off + 0.12 + diag_h + 0.25)
        return s

    def cta(self, spec):
        s = self._slide(bg=self.c["primary"])
        m = self.margin
        self._accent_bar(s, m, 2.7, w=1.6, h=0.12)
        self._text(
            s,
            m,
            2.95,
            EMU_W - 2 * m,
            1.2,
            spec["title"],
            self.s["title"],
            "FFFFFF",
            font=self.f["head"],
            bold=True,
        )
        contact = spec.get("contact", {})
        lines = [f"{k}: {v}" for k, v in contact.items() if v]
        self._text(
            s,
            m,
            4.4,
            EMU_W - 2 * m,
            1.5,
            "\n".join(lines),
            self.s["body"],
            self.c["bg_alt"],
        )
        self._logo_mark(s, color="FFFFFF")
        return s

    # 2026-07-25: legacy `cover`·`toc`·`part`를 폐기했다 — 골든 프레임
    # (`golden.cover`/`golden.toc`/`golden.part`, goldenfab/layouts.py)이 정본이고 이쪽은 v3 시절
    # 사본이라 "어느 쪽이 표준인가"가 갈렸다(구세대 견본 spec이 이 타입을 쓰고 있었다).
    # 이 타입을 쓰면 아래 unknown slide type ValueError로 시끄럽게 죽는다.
    RENDERERS = {
        "section": section,
        "bullets": bullets,
        "two_column": two_column,
        "metrics": metrics,
        "table": table,
        "chart": chart,
        "diagram": diagram,
        "cta": cta,
    }

    def _render_golden(self, sl, slide_no="?"):
        """type "golden.<layout>" — goldenfab registry 호출.

        공장 문. 골든 기본값은 회귀 게이트 전용이므로 여기로는 못 나간다 —
        content가 골든 DEFAULT를 전부 덮지 않으면 빌드 중단(content_contract).
        """
        import inspect

        from goldenfab.content_contract import assert_content
        from goldenfab.registry import LAYOUTS as GOLDEN_LAYOUTS

        key = sl["type"].split(".", 1)[1]
        fn = GOLDEN_LAYOUTS.get(key)
        if fn is None:
            raise ValueError(f"unknown golden layout: {key} (등록: {sorted(GOLDEN_LAYOUTS)})")
        params = inspect.signature(fn).parameters
        if len(params) >= 2:
            assert_content(key, fn, sl.get("content"), slide_no)
            result = fn(self.prs, sl.get("content"))
        else:
            if sl.get("content"):
                raise ValueError(
                    f"golden.{key}는 콘텐츠 파라미터 미지원(골든 내장) — content 제거 필요"
                )
            result = fn(self.prs)
        slide = result[0] if isinstance(result, tuple) else result
        if slide is None:  # 내부 side-effect 함수(반환 없음) — 방금 추가된 슬라이드 사용
            slide = self.prs.slides[-1]
        self._flatten(slide)
        return slide

    def _render_scripted(self, sl, slide_no="?"):
        """type "adapted.<골든타입>" / "novel" — 장 스크립트 디스패치 (2026-07-16 변형 출발점).

        골든은 복제 템플릿이 아니라 변형 가능한 출발점이다(design-rules 골든셋 원칙 개정).
        내용의 카디널리티·구획이 골든과 다르면 goldenfab 부품(kit·grid)을 재사용한 장
        스크립트로 변형(adapted)하고, 골든에 없는 물성이면 신규 설계(novel)한다.
        스냅샷 회귀는 goldenfab 보호 전용이므로 여기 장은 대상이 아니다 — 품질은
        audit_deck.py(전역 오딧+밀도 밴드)와 병렬 채점(design-rules P4)이 진다.

        **content 계약(2026-07-25 배선):** 이 문에는 계약 검사가 **없었다** — 골든 출발 스크립트를
        그대로 쓰면 K-IFRS 글이 남아도 전 게이트 green이었다(`_render_golden`에만 assert_content가
        있었다). 이제 `assert_scripted_content`가 "스크립트가 실제 읽는 골든 기본값 키"를 요구한다.
        """
        import importlib.util

        script = sl.get("script")
        if not script:
            raise ValueError(
                f"S{slide_no} {sl['type']}: 'script' 경로 필수 — 장 스크립트가 레이아웃 본체다"
            )
        p = Path(script)
        if not p.is_absolute():
            p = (self._spec_dir / p) if self._spec_dir else p
        if not p.exists():
            raise ValueError(f"S{slide_no} {sl['type']}: 장 스크립트 없음: {p}")
        mod_spec = importlib.util.spec_from_file_location(f"chapter_s{slide_no}", p)
        if mod_spec is None or mod_spec.loader is None:
            raise ValueError(f"S{slide_no} {sl['type']}: 장 스크립트 로드 불가: {p}")
        mod = importlib.util.module_from_spec(mod_spec)
        mod_spec.loader.exec_module(mod)
        if not hasattr(mod, "build"):
            raise ValueError(f"S{slide_no} {sl['type']}: 장 스크립트에 build(prs, content) 없음")
        from goldenfab.content_contract import assert_scripted_content

        assert_scripted_content(
            sl["type"], mod, p.read_text(encoding="utf-8"), sl.get("content"), slide_no
        )
        result = mod.build(self.prs, sl.get("content"))
        slide = result[0] if isinstance(result, tuple) else result
        if slide is None:
            slide = self.prs.slides[-1]
        self._flatten(slide)
        return slide

    def build(self, spec):
        self.meta = spec.get("meta", {})
        body_types = {
            "section",
            "bullets",
            "two_column",
            "table",
            "chart",
            "diagram",
            "metrics",
        }
        # PART pre-pass: 간지 수집 + 본문 계층 번호(N-M) 부여 (내비게이션 단일 출처)
        # **경계(2026-07-25)**: legacy `part` 렌더러가 폐기돼 이 분기는 legacy 경로에서 더는
        # 발화하지 않는다(골든 덱은 `golden.part`로 자체 간지·탭을 그린다). 그래서 아래 내비
        # 기계(`_nav_tabs`·계층번호 `_num`·NAV_BODY_TOP)는 현재 **도달 불가**다. 지울지
        # 살릴지는 "legacy (B) 경로를 어디까지 유지하나"라는 별개 결정이라 사용자 몫으로 남긴다 —
        # 여기 적어 두는 이유는, 안 적으면 다음 사람이 살아 있는 기계로 착각하고 손대기 때문이다.
        seq = 0
        for sl in spec["slides"]:
            if sl["type"] == "part":
                seq = 0
                sl["_part_no"] = len(self.parts) + 1
                self.parts.append({"no": sl["_part_no"], "title": sl["title"], "items": []})
            elif sl["type"] in body_types and self.parts:
                seq += 1
                sl["_part_no"] = self.parts[-1]["no"]
                sl["_num"] = f"{sl['_part_no']}-{seq}"
                self.parts[-1]["items"].append(sl.get("title", ""))
        self.body_top = NAV_BODY_TOP if self.parts else BODY_TOP
        total = len(spec["slides"])
        self._total = total
        for i, sl in enumerate(spec["slides"], start=1):
            if sl["type"].startswith("golden."):
                # 골든 엔진(goldenfab) 디스패치 — 골든 장은 헤더·결론바·출처를 자체 포함하므로
                # 레거시 푸터 스탬프·배너를 적용하지 않는다. compare_golden.py가 회귀 게이트.
                self._render_golden(sl, i)
                continue
            if sl["type"].startswith("adapted.") or sl["type"] == "novel":
                # 장 스크립트 디스패치 — 골든을 출발점으로 변형(adapted) / 신규 물성(novel).
                # 골든 장과 같이 자체 완결이므로 푸터 스탬프·배너 미적용.
                self._render_scripted(sl, i)
                continue
            renderer = self.RENDERERS.get(sl["type"])
            if renderer is None:
                raise ValueError(f"unknown slide type: {sl['type']}")
            self._page_no = i
            slide = renderer(self, sl)
            self._flatten(slide)  # v4.5: PowerPoint 기본 그림자 전면 제거(플랫 룩)
            # 마스터 틀: 표지·간지·CTA를 제외한 전 장에 푸터 스탬프
            if sl["type"] not in ("cover", "part", "cta"):
                self._frame(slide, i, total)
            if sl.get("banner") and sl["type"] not in ("cover", "toc", "part", "cta"):
                # 배너는 하단 스트립을 점유 — source·footnotes와 병용 금지(겹침)
                self._banner(slide, sl["banner"])
            else:
                if sl.get("source"):
                    self._source(slide, sl["source"])
                if sl.get("footnotes"):
                    self._footnotes(slide, sl["footnotes"])
        self.gate_fails = self._gate_density(spec)
        return self.prs

    def _gate_density(self, spec):
        """빌드 시 기계 게이트(2026-07-20 신설 · 2026-07-25 legacy 확장) — 빌드가 **스스로**
        `goldenfab.audit.generic_checks`로 채점한다. 오케스트레이터가 consistency-qa를 안 돌려도
        성김·겹침·넘침이 잡힌다. 규칙 단일 출처는 audit.py, 러너 참조는 audit_deck.py.

        프로파일 배정: golden.*(정형 4종 제외)=sparse · adapted./novel=dense ·
        **legacy 렌더러 장(section·bullets·two_column·metrics·table·chart·diagram·cta)=legacy**.
        legacy는 2026-07-25까지 이 게이트의 대상이 아니었다 — 남은 legacy 어휘(차트 6·mpl 9·표·
        문서형·다이어그램 9종)가 무검증으로 나갔다. 밀도 밴드는 골든 파생이라 legacy에 안 건다
        (legacy 밀도의 정본은 consistency-qa audit_pptx의 본문 60단어 하한).
        반환: 실패 문자열 목록(빈 목록=통과). 빌드는 막지 않는다 — main()이 종료코드로 신호한다.
        """
        import sys as _sys

        gf = str(Path(__file__).resolve().parent)
        if gf not in _sys.path:
            _sys.path.insert(0, gf)
        from goldenfab import audit as A
        from goldenfab.kit import load_kit

        fixed = {"cover", "toc", "part", "closing"}
        accent = str(load_kit()["rgb"]["accent"])
        band = A.density_band()
        slides = list(self.prs.slides)
        fails = []
        for i, sl in enumerate(spec["slides"], start=1):
            t = sl.get("type", "")
            base = t.split(".", 1)[1] if t.startswith(("golden.", "adapted.")) else t
            opts = sl.get("audit", {})
            if t.startswith("golden.") and base not in fixed:
                profile = "dense" if opts.get("dense", False) else "sparse"
            elif t.startswith("adapted.") or t == "novel":
                profile = "dense"
            elif t in self.RENDERERS:
                profile = "legacy"  # legacy 렌더러 장 — 2026-07-25 배선
            else:
                continue
            if i > len(slides):
                continue
            screenshot = base == "screenshot" or opts.get("screenshot", False)
            shapes = list(slides[i - 1].shapes)
            res = A.generic_checks(
                shapes,
                accent,
                # 밀도 밴드는 골든 스냅샷 파생이라 legacy 프레임엔 안 건다(93/111 상시 적색)
                band=None if profile == "legacy" else band,
                dup_allow=opts.get("dup_allow", 0),
                screenshot=screenshot,
                profile=profile,
            )
            fails += [f"S{i} {t}: {x}" for x in A.report(res, known=opts.get("known"))]
        return fails


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("spec")
    ap.add_argument("out")
    ap.add_argument(
        "--brand",
        default=str(Path(__file__).parent.parent / "assets" / "brand-kit.yaml"),
    )
    args = ap.parse_args()

    brand = load_brand(args.brand)
    spec = json.loads(Path(args.spec).read_text(encoding="utf-8"))
    deck = Deck(brand)
    deck.render_dir = Path(args.out).resolve().parent / "render"
    deck._spec_dir = Path(args.spec).resolve().parent
    prs = deck.build(spec)
    prs.save(args.out)
    print(f"OK: {args.out} ({len(prs.slides)} slides)")
    # 빌드 시 밀도 게이트 결과 — pptx는 이미 저장(검사 가능), 실패면 종료코드로 신호(기계 강제)
    fails = getattr(deck, "gate_fails", [])
    if fails:
        print(
            f"\n밀도 게이트 FAIL {len(fails)}건 — 성김/탑헤비(글자수 미달). 카드로 메우지 말고 근거 보강:"
        )
        for f in fails:
            print(f"  - {f}")
        return 1
    return 0


if __name__ == "__main__":
    sys.exit(main())
