"""deckkit — 토큰·원시도구·페이지 골격. 형태는 정하지 않는다.

## 왜 이 모듈이 생겼나 (2026-08-03)

하네스를 못 보는 새 컨텍스트에 계약만 주고 같은 덱을 만들게 한 실험에서, 하네스로 만든 것보다
조밀하고(도형 70 대 50/장) 수치도 정확한 산출물이 나왔다. 사용자가 그쪽을 택했고, 골든덱은
기준선에서 내렸다. 이 모듈은 그 산출물에서 **토큰과 골격만** 뽑은 것이다.

## 무엇을 소유하고 무엇을 소유하지 않나

소유한다 — 색·글꼴 위계·여백(`brand.yaml`), 텍스트/사각/화살표/막대/이미지, 표지·간지·헤더·
푸터·패널. 즉 **일관성이 필요한 것**.

소유하지 않는다 — 장을 어떻게 구성할지. 그건 매번 내용에 맞게 정한다. 부품이 장 구성을 정하면
전 장이 같은 형식이 된다(2026-07-29 실측: 36장이 동일한 4카드).
"""

import math
import os
import struct

import yaml
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

_B = yaml.safe_load((__import__("pathlib").Path(__file__).parent / "brand.yaml").read_text("utf-8"))

COLORS = _B["colors"]
SIZES = _B["sizes"]
FONT = _B["font"]
FR = _B["frame"]
LIM = _B["limits"]

INK = COLORS["ink"]
# 구조는 근흑이 지고 오렌지는 포인트로만 — 무채색 + 포인트 하나(Muted Ember).
# NAVY/BLUE/AMBER 이름은 호출부 호환을 위해 남겼다. 값은 브랜드 토큰이 단일 출처다.
NAVY, NAVY_LT = COLORS["dark"], COLORS["dark_lt"]
BLUE, BLUE_LT = COLORS["struct"], COLORS["struct_lt"]
AMBER, AMBER_LT = COLORS["accent"], COLORS["accent_lt"]
GRAY, LINE, PANEL, WHITE = COLORS["muted"], COLORS["line"], COLORS["panel"], COLORS["white"]
DARK, DARK_LT = NAVY, NAVY_LT  # 뜻이 드러나는 이름 — 새 코드는 이쪽을 쓴다
STRUCT, STRUCT_LT = BLUE, BLUE_LT
ACCENT, ACCENT_LT = AMBER, AMBER_LT

SW, SH = FR["slide_w"], FR["slide_h"]
M, TOP, BOT = FR["margin"], FR["top"], FR["bottom"]
RULE = FR["rule"]
CW = SW - 2 * M
RIGHT = SW - M

_ALIGN = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}
_ANCH = {"t": MSO_ANCHOR.TOP, "m": MSO_ANCHOR.MIDDLE, "b": MSO_ANCHOR.BOTTOM}

# 조판 중 쌓이는 경고 — 장 스크립트가 끝나고 읽는다
WARN: list[str] = []
CUR = {"n": 0, "title": ""}
IMG_DIR = ""


def set_context(n: int, title: str = "", img_dir: str = ""):
    CUR["n"], CUR["title"] = n, title
    if img_dir:
        global IMG_DIR
        IMG_DIR = img_dir


def C(h: str) -> RGBColor:
    return RGBColor.from_string(h)


def In(v: float) -> Inches:
    return Inches(v)


def _em(txt: str) -> float:
    """문자열의 대략적 폭(em). 전각 1.0 · 영숫자 0.55 · 공백 0.35."""
    w = 0.0
    for ch in txt:
        w += 0.35 if ch == " " else (1.0 if ord(ch) > 0x2000 else 0.55)
    return w


def _font(run, size, bold, color, italic=False):
    """라틴·동아시아·복합 자형을 모두 같은 글꼴로 — 하나라도 빠지면 한글이 다른 꼴로 나온다."""
    f = run.font
    f.size, f.bold, f.italic = Pt(size), bold, italic
    f.color.rgb = C(color)
    f.name = FONT
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(f"{{http://schemas.openxmlformats.org/drawingml/2006/main}}{tag[2:]}")
        if el is None:
            el = rPr.makeelement(
                f"{{http://schemas.openxmlformats.org/drawingml/2006/main}}{tag[2:]}", {}
            )
            rPr.append(el)
        el.set("typeface", FONT)


def T(slide, x, y, w, h, lines, valign="t", pad=0.03, wrap=True, name=""):
    """텍스트 블록. lines = [{t, s, b, c, a, sb, sa, ls}] · t는 문자열 또는 런 목록.

    한글은 **어절 단위로만 줄바꿈**한다(`eaLnBrk=0`). 이걸 안 걸면 PowerPoint가 글자 단위로
    끊어 고아 글자가 생기고, 무하네스 실험에서 그걸 손으로 7곳 고쳐야 했다.
    """
    tb = slide.shapes.add_textbox(In(x), In(y), In(w), In(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = In(pad)
    tf.margin_top = tf.margin_bottom = In(0.01)
    tf.vertical_anchor = _ANCH[valign]
    used, inner = 0.0, max(w - 2 * pad, 0.2)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p._p.get_or_add_pPr().set("eaLnBrk", "0")
        p.alignment = _ALIGN[ln.get("a", "l")]
        sb, sa = ln.get("sb", 0), ln.get("sa", 0)
        p.space_before, p.space_after = Pt(sb), Pt(sa)
        ls = ln.get("ls", 1.0)
        p.line_spacing = ls
        t = ln["t"]
        runs = (
            t
            if isinstance(t, list)
            else [(t, ln.get("s", SIZES["body"]), ln.get("b", False), ln.get("c", INK))]
        )
        maxs, em = 0.0, 0.0
        for rt, rs, rb, rc in runs:
            for j, seg in enumerate(str(rt).split("\n")):
                if j:
                    em = math.ceil(em / inner - 1e-6) * inner
                r = p.add_run()
                r.text = ("\n" if j else "") + seg
                _font(r, rs, rb, rc)
                em += _em(seg) * rs / 72.0
            maxs = max(maxs, rs)
        used += max(1, math.ceil(em / inner - 1e-6)) * maxs * 1.22 * ls / 72.0 + (sb + sa) / 72.0
    if used > h + 0.03:
        WARN.append(f'S{CUR["n"]:02d} {name or "텍스트"} 넘침 추정 {used:.2f}" > {h:.2f}"')
    return tb


def R(slide, x, y, w, h, fill=None, line=None, lw=0.75, rounded=0.0):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, In(x), In(y), In(w), In(h)
    )
    if rounded:
        shp.adjustments[0] = rounded
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = C(fill)
    else:
        shp.fill.background()
    if line:
        shp.line.color.rgb = C(line)
        shp.line.width = Pt(lw)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    shp.text_frame.word_wrap = True
    if x + w > RIGHT + LIM["right_tolerance"]:
        WARN.append(f'S{CUR["n"]:02d} 우변 초과 {x + w - RIGHT:.2f}" (도형 x={x:.2f} w={w:.2f})')
    return shp


def ARROW(slide, x, y, w, h, color=BLUE):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, In(x), In(y), In(w), In(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = C(color)
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _arrowhead(shp, head=True, tail=False):
    ln = shp.line._get_or_add_ln()
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    for tag, on in (("headEnd", tail), ("tailEnd", head)):
        if not on:
            continue
        el = ln.find(f"{ns}{tag}")
        if el is None:
            el = ln.makeelement(f"{ns}{tag}", {})
            ln.append(el)
        el.set("type", "triangle")
        el.set("w", "med")
        el.set("len", "med")


def CONNECT(slide, x1, y1, x2, y2, *, elbow=False, head=True, color=None, w=1.0, dash=False):
    """점과 점을 잇는다 — 직선 또는 꺾은선(elbow), 화살촉 선택.

    **도해 부품을 만들지 않는 이유.** 흐름도·관계도를 `flow()` `graph()` 같은 블록으로 굳히면
    안 맞는 내용에도 그 형식으로 그리게 된다(2026-07-29 실측: 전 장이 같은 형식). 그래서
    잉크만 준다 — 사각(`R`)과 이 커넥터면 머메이드·아스키로 적을 수 있는 노드·간선 도해는
    무엇이든 그릴 수 있고, 배치는 매번 내용이 정한다.

    언제 그릴지는 재료가 말한다 — `scan_material.py`가 절마다 화살표 수·순서 유무를 센다.
    """
    kind = MSO_CONNECTOR.ELBOW if elbow else MSO_CONNECTOR.STRAIGHT
    c = slide.shapes.add_connector(kind, In(x1), In(y1), In(x2), In(y2))
    c.line.color.rgb = C(color or LINE)
    c.line.width = Pt(w)
    if dash:
        c.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    if head:
        _arrowhead(c)
    return c


def _edge(shp, side):
    x, y = Emu(shp.left).inches, Emu(shp.top).inches
    w, h = Emu(shp.width).inches, Emu(shp.height).inches
    return {
        "l": (x, y + h / 2),
        "r": (x + w, y + h / 2),
        "t": (x + w / 2, y),
        "b": (x + w / 2, y + h),
        "c": (x + w / 2, y + h / 2),
    }[side]


def LINK(slide, a, b, *, sides=None, elbow=False, head=True, color=None, w=1.0, dash=False):
    """도형 a → 도형 b를 잇는다. sides=("r","l") 처럼 접점을 지정, 없으면 상대 위치로 고른다."""
    ax, ay = (
        Emu(a.left).inches + Emu(a.width).inches / 2,
        Emu(a.top).inches + Emu(a.height).inches / 2,
    )
    bx, by = (
        Emu(b.left).inches + Emu(b.width).inches / 2,
        Emu(b.top).inches + Emu(b.height).inches / 2,
    )
    if sides is None:
        sides = (
            ("r", "l")
            if abs(bx - ax) >= abs(by - ay) and bx >= ax
            else (
                ("l", "r")
                if abs(bx - ax) >= abs(by - ay)
                else (("b", "t") if by >= ay else ("t", "b"))
            )
        )
    p1, p2 = _edge(a, sides[0]), _edge(b, sides[1])
    return CONNECT(slide, *p1, *p2, elbow=elbow, head=head, color=color, w=w, dash=dash)


def BAR(slide, x, y, w, h, frac, color, track=None):
    """값을 길이로 인코딩. track을 주면 전체 대비가 보인다."""
    if track:
        R(slide, x, y, w, h, fill=track)
    frac = max(0.0, min(1.0, frac))
    if frac > 0:
        R(slide, x, y, max(w * frac, 0.02), h, fill=color)


def _png_size(path):
    with open(path, "rb") as f:
        head = f.read(33)
    return struct.unpack(">II", head[16:24])


def image_ar(fname) -> float:
    """그림의 가로/세로 비율. 칸을 그림에 맞출 때 쓴다."""
    path = fname if os.path.isabs(str(fname)) else os.path.join(IMG_DIR, str(fname))
    iw, ih = _png_size(path)
    return iw / ih


def image_h(fname, w, caption=False, cap_h=0.24) -> float:
    """**폭 w를 다 쓸 때 필요한 칸 높이.**

    칸을 먼저 정하고 그림을 우겨넣으면 비율이 안 맞는 만큼 빈다(2026-08-03 지적: "칸이
    안 맞아서 보기 싫다"). 순서를 뒤집는다 — 폭을 정하고 높이를 **그림에게 물어본다.**
    """
    return w / image_ar(fname) + (cap_h if caption else 0)


def IMAGE(slide, fname, x, y, w, h, caption=None, cap_h=0.24):
    """비율 유지로 칸에 맞춘다. 결과 폭이 판독 하한 미만이면 경고 — 개수를 줄이거나 크롭할 것."""
    path = fname if os.path.isabs(str(fname)) else os.path.join(IMG_DIR, str(fname))
    iw, ih = _png_size(path)
    ar = iw / ih
    box_h = h - (cap_h if caption else 0)
    pw, ph = (box_h * ar, box_h) if w / box_h > ar else (w, w / ar)
    px, py = x + (w - pw) / 2, y + (box_h - ph) / 2
    pic = slide.shapes.add_picture(path, In(px), In(py), In(pw), In(ph))
    pic.line.color.rgb = C(LINE)
    pic.line.width = Pt(0.75)
    if pw < LIM["image_min_w"]:
        WARN.append(
            f'S{CUR["n"]:02d} 이미지 폭 {pw:.2f}" < 판독 하한 {LIM["image_min_w"]}"'
            f" ({os.path.basename(path)}) — 개수를 줄이거나 크롭할 것"
        )
    fill = (pw * ph) / (w * box_h) if w * box_h > 0 else 1.0
    if fill < LIM["image_fill_min"]:
        WARN.append(
            f"S{CUR['n']:02d} 그림이 칸의 {fill:.0%}만 채운다 ({os.path.basename(path)})"
            f' — 칸 {w:.2f}×{box_h:.2f}", 그림 {pw:.2f}×{ph:.2f}".'
            f" `image_h(파일, 폭)`으로 칸 높이를 그림에 맞출 것"
        )
    if caption:
        T(
            slide,
            x,
            y + box_h + 0.02,
            w,
            cap_h,
            [{"t": caption, "s": SIZES["tiny"], "c": GRAY, "a": "c", "ls": 1.1}],
            name="caption",
        )
    return pic


# ── 페이지 골격 ──────────────────────────────────────────────────────────────
def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ── 부(部) 내비게이션 ────────────────────────────────────────────────────────
# 제목 옆 우측은 **덱 내내 같은 메뉴**다. 장마다 다른 문장을 넣던 자리였는데(2026-08-03 이전),
# 그러면 "지금 어느 부에 있나"를 볼 곳이 간지밖에 없다 — 간지를 넘긴 순간부터 위치를 잃는다.
# 목록은 조판이 정하지 않는다. `set_parts(outline.parts(계약))`로 계약에서 읽어 온다.
PARTS: list[tuple[str, str]] = []
# 탭 줄은 헤더 밑줄(`FR["rule"]`)에 앉는다 — 현재 부의 강조 바가 그 밑줄의 일부가 되어
# "지금 열린 탭"으로 읽힌다. 좌측 1.15" 근흑 마디와 같은 선이라 요소를 새로 만들지 않는다.
# x는 **고정**이다 — 장마다 탭이 움직이면 눈이 매번 찾아야 하고, 그러면 내비가 아니다.
# 대신 항목 사이 간격을 남는 폭으로 나눠 우변까지 꽉 채운다.
NAV = {"x": 6.00, "y": 0.60, "h": 0.30, "min_gap": 0.16, "bar": 0.030, "pad": 0.26}


def set_parts(parts):
    """부 목록을 덱 단위로 고정한다. parts = [(로마자, 제목, ...), ...] — 리드는 무시한다."""
    PARTS[:] = [(str(p[0]), str(p[1])) for p in parts]


def nav_label(roman, title) -> str:
    return f"{roman} {title}"


def _nav_widths() -> list[float]:
    """항목마다 제 글자만큼 — 균등분할하면 짧은 부에 빈칸이 생기고 긴 부가 눌린다."""
    return [_em(nav_label(r, t)) * SIZES["small"] / 72.0 + 0.06 for r, t in PARTS]


def nav_gap() -> float:
    """`NAV["x"]`부터 우변까지를 항목이 나눠 쓰고 남은 것이 간격이다."""
    ws = _nav_widths()
    if len(ws) < 2:
        return 0.0
    return (RIGHT - NAV["x"] - sum(ws)) / (len(ws) - 1)


def part_nav(s, part):
    """제목 옆을 가로로 채우는 부 탭 줄. 현재 부만 accent + 밑줄 위 강조 바."""
    if not PARTS:
        WARN.append(f"S{CUR['n']:02d} 부 목록이 비었다 — set_parts()를 먼저 부른다")
        return
    ws, gap = _nav_widths(), nav_gap()
    if gap < NAV["min_gap"]:
        WARN.append(
            f'S{CUR["n"]:02d} 부 탭 간격 {gap:.2f}" < {NAV["min_gap"]}"'
            " — 부 이름이 길거나 부가 많다. 계약의 부 제목을 줄인다"
        )
    x, hit = NAV["x"], False
    for (roman, title), w in zip(PARTS, ws):
        cur = roman == part
        hit = hit or cur
        tb = T(
            s,
            x,
            NAV["y"],
            w,
            NAV["h"],
            [
                {
                    "t": nav_label(roman, title),
                    "s": SIZES["small"],
                    "b": cur,
                    "c": ACCENT if cur else GRAY,
                    "a": "c",
                }
            ],
            valign="m",
            name=f"부내비:{roman}",
        )
        tb.name = f"부내비:{roman}"  # 게이트가 이 이름으로 찾는다
        if cur:
            R(s, x, RULE - 0.006, w, NAV["bar"], fill=ACCENT)
        x += w + gap
    if part and not hit:
        WARN.append(f"S{CUR['n']:02d} 부 「{part}」가 계약 부 목록에 없다")


def page_head(s, n, title, footer_text="", width=None):
    """제목 + 밑줄 + 푸터. 내비 없는 골격 — 목차처럼 부에 속하지 않는 장이 쓴다."""
    set_context(n, title)
    T(
        s,
        M - 0.02,
        0.28,
        width or CW,
        0.56,
        [{"t": title, "s": SIZES["title"], "b": True, "c": INK}],
        valign="m",
        name="제목",
    )
    R(s, M, RULE, CW, 0.018, fill=LINE)
    R(s, M, RULE, 1.15, 0.018, fill=BLUE)
    footer(s, n, footer_text)


def header(s, n, part, title, footer_text=""):
    """본문 장 상단. part = 로마자("Ⅲ") — 이 장이 속한 부, 계약 `장 목록`의 부 열이 정한다."""
    page_head(s, n, title, footer_text, width=NAV["x"] - M - NAV["pad"])
    part_nav(s, part)


def footer(s, n, text=""):
    R(s, M, 7.00, CW, 0.012, fill=LINE)
    if text:
        T(s, M, 7.06, 8.0, 0.26, [{"t": text, "s": SIZES["tiny"], "c": GRAY}], name="꼬리")
    T(
        s,
        11.0,
        7.06,
        1.71,
        0.26,
        [{"t": f"{n:02d}", "s": SIZES["tiny"], "c": GRAY, "a": "r"}],
        name="쪽번호",
    )


def dark_page(prs, n):
    """어두운 전면 장 — 표지·간지의 바탕.

    배경은 **도형이 아니라 슬라이드 배경**으로 칠한다. 도형이면 폭 13.33"가 우변 경고를 낸다.
    """
    s = blank(prs)
    set_context(n)
    bg = s.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = C(NAVY)
    T(
        s,
        11.0,
        7.02,
        1.71,
        0.26,
        [{"t": f"{n:02d}", "s": SIZES["tiny"], "c": GRAY, "a": "r"}],
        name="쪽번호",
    )
    return s


def cover(prs, kicker, title, subtitle=None, lead=None, facts=()):
    """표지. title은 줄바꿈(\\n)으로 행을 나눈다. facts = [(라벨, 값, 부연), ...]

    정형 장이라 덱마다 달라지면 안 된다 — 2026-08-03에 이게 템플릿에 없어서 매번 새로
    발명되고 있었다(그 과정에서 토큰 밖 색도 생겼다).
    """
    s = dark_page(prs, 1)
    R(s, 0, 0, 0.14, SH, fill=ACCENT)
    T(
        s,
        M,
        1.52,
        8.0,
        0.30,
        [{"t": kicker, "s": SIZES["small"], "b": True, "c": ACCENT}],
        name="킥커",
    )
    T(
        s,
        M,
        1.94,
        11.4,
        1.66,
        [
            {"t": ln, "s": SIZES["cover"], "b": True, "c": WHITE, "ls": 1.16}
            for ln in title.split("\n")
        ],
        name="표제",
    )
    y = 3.88
    if subtitle:
        T(
            s,
            M,
            y,
            9.0,
            0.56,
            [{"t": subtitle, "s": SIZES["div"], "b": True, "c": ACCENT}],
            name="부제",
        )
        y += 0.64
    if lead:
        T(
            s,
            M,
            y,
            10.6,
            0.40,
            [{"t": lead, "s": SIZES["lead"], "c": LINE, "ls": 1.25}],
            name="리드",
        )
    if facts:
        n = len(facts)
        gap = 0.24
        fw = (CW - gap * (n - 1)) / n
        fy = 5.42
        for i, (lab, val, sub) in enumerate(facts):
            fx = M + i * (fw + gap)
            R(s, fx, fy, fw, 0.86, fill=NAVY_LT)
            T(
                s,
                fx + 0.20,
                fy + 0.10,
                fw - 0.40,
                0.66,
                [
                    {"t": lab, "s": SIZES["tiny"], "c": GRAY},
                    {"t": val, "s": SIZES["h"], "b": True, "c": WHITE, "sb": 2},
                    {"t": sub, "s": SIZES["tiny"], "c": GRAY, "sb": 2},
                ],
                name=f"표지칩:{lab}",
            )
    return s


def divider(prs, n, roman, title, lead=None, chapters=()):
    """간지. chapters = ["S4 문제 인식", ...] — 그 부에 무엇이 들었는지 미리 보인다."""
    s = dark_page(prs, n)
    R(s, 0, 2.05, 2.30, 0.05, fill=ACCENT)
    T(
        s,
        M,
        2.30,
        3.0,
        0.62,
        [{"t": roman, "s": SIZES["div"], "b": True, "c": ACCENT}],
        name="로마자",
    )
    T(
        s,
        M,
        3.00,
        10.0,
        0.70,
        [{"t": title, "s": SIZES["cover"], "b": True, "c": WHITE}],
        name="간지제목",
    )
    if lead:
        T(
            s,
            M,
            3.86,
            8.4,
            0.44,
            [{"t": lead, "s": SIZES["lead"], "c": LINE, "ls": 1.25}],
            name="간지리드",
        )
    if chapters:
        x = M
        for ch in chapters:
            w = 0.30 + len(ch) * 0.115
            R(s, x, 4.62, w, 0.34, fill=NAVY_LT)
            T(
                s,
                x,
                4.62,
                w,
                0.34,
                [{"t": ch, "s": SIZES["tiny"], "c": LINE, "a": "c"}],
                valign="m",
                name="간지칩",
            )
            x += w + 0.14
    return s


def toc(prs, n, parts, footer_text=""):
    """목차. parts = [(로마자, 부 제목, "S4 · S5", 요약), ...] — 부마다 한 행."""
    s = blank(prs)
    page_head(s, n, "목차", footer_text, width=CW)
    rh = (BOT - TOP) / len(parts)
    for i, (roman, title, span, summary) in enumerate(parts):
        y = TOP + i * rh
        R(s, M, y + rh - 0.02, CW, 0.012, fill=LINE)
        R(s, M, y + 0.10, 0.52, 0.52, fill=BLUE)
        T(
            s,
            M,
            y + 0.10,
            0.52,
            0.52,
            [{"t": roman, "s": SIZES["h"], "b": True, "c": WHITE, "a": "c"}],
            valign="m",
            name="부번호",
        )
        T(
            s,
            M + 0.74,
            y,
            4.4,
            rh - 0.06,
            [{"t": title, "s": SIZES["title"] - 6, "b": True, "c": INK}],
            valign="m",
            name=f"부:{title}",
        )
        T(
            s,
            M + 5.30,
            y,
            5.6,
            rh - 0.06,
            [{"t": summary, "s": SIZES["small"], "c": GRAY, "ls": 1.2}],
            valign="m",
            name="부요약",
        )
        T(
            s,
            M + CW - 1.5,
            y,
            1.5,
            rh - 0.06,
            [{"t": span, "s": SIZES["small"], "b": True, "c": BLUE, "a": "r"}],
            valign="m",
            name="쪽범위",
        )
    return s


# ── 여기 `Panel`이 있었다 (2026-08-03 삭제) ──────────────────────────────────
#
# 소제목으로 가둔 구역을 그리고 `row(h)`로 세로 슬롯을 나눠 주던 클래스다. 만든 이유는
# "쌓은 블록의 좌우가 어긋나는 사고를 규칙이 아니라 구조로 막는다"였는데, **실전 덱 19장이
# 단 한 번도 부르지 않았다.** 정렬은 조판이 좌표로 직접 잡아서 해결됐다.
#
# 대신 계약만 받은 새 컨텍스트 6회는 **6/6이 불렀다.** 기본값이 `fill=PANEL, line=LINE`이라
# 부르는 순간 회색 면이 딸려 왔고, 그때마다 본문 면적의 **83~94%가 회색 덩어리**가 됐다
# (완성본 19장의 같은 수치는 0%). 2026-07-29의 `figures/*.LAYOUT` 사고와 같은 구조다 —
# **기본값이 장 구성을 정한다.**
#
# 면을 뺀 "자리 배분기"로 남기지도 않았다. 배분기도 결국 칸이고, 칸이 있으면 재료를 칸에
# 맞춰 휘게 된다. 자리는 매번 내용이 정한다.
