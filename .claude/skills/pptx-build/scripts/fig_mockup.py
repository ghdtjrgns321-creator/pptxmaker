#!/usr/bin/env python
"""부품 목업 러너 — 도구함 부품 1종을 데이터 여러 벌로 찍어 사람이 검증한다 (2026-07-26).

부품화의 검증 스텝(재설계 계획 1단계 ②③)이 이 스크립트다. 부품은 `draw(slide, box, data, kit)`
계약을 지키므로 **임의 데이터로 부를 수 있다** — 장 함수였을 때는 불가능했던 것이고, 그래서
목업이 범용 폴백(B)으로 샜다.

각 슬라이드가 보여주는 것:
  · 샘플 라벨(데이터가 무엇인지) + 부품 ID·물성
  · box 경계 가이드(옅은 점선) — 부품이 자리 밖으로 새지 않는지 **눈으로** 확인
  · 항목 수를 바꾼 벌들 — 최소 · 골든 실측 · 밀도 · 극단

수용 한계를 넘는 데이터는 `grid.pitch`/`track`이 시끄럽게 죽인다. 러너는 그 예외를 잡아
"몇 개에서 죽는지"를 마지막 장에 적는다(조용한 절단이 없다는 증거).

usage: uv run python .claude/skills/pptx-build/scripts/fig_mockup.py <부품이름> [--no-render]
       uv run python .claude/skills/pptx-build/scripts/fig_mockup.py --all
"""

import importlib
import subprocess
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[4]
sys.path.insert(0, str(Path(__file__).parent))

from pptx import Presentation  # noqa: E402
from pptx.util import Inches  # noqa: E402

from goldenfab import dense  # noqa: E402
from goldenfab.kit import SLIDE_H, SLIDE_W, add_box, add_text, load_kit  # noqa: E402

FIG_DIR = Path(__file__).parent / "goldenfab" / "figures"
OUT_PPTX = ROOT / "results" / "검토"
OUT_PNG = ROOT / "_render" / "부품"

# 목업 기본 자리 = **콘텐츠 영역 전체**. 부품은 자리에 맞춰 늘어나야 하므로 목업도 그렇게 찍는다
# (자리를 골든 장 크기로 고정해두면 "그 장에서만 어떻게 보이나"가 되고, 디자인 판단이 안 된다).
# 헤더 룰 0.66 아래 ~ 출처선 7.14 위 · 좌우 마진은 grid와 동일.
DEFAULT_BOX = (0.6, 1.0, 12.133, 5.9)


def _guide(slide, box, kit):
    """box 경계 가이드 — 부품이 자리 안에만 그리는지 눈으로 보이게. 목업 전용."""
    g = add_box(slide, box[0], box[1], box[2], box[3], fill=None, line=kit["rgb"]["bg_alt"])
    g.line.width = Inches(0.008)
    ln = g.line._get_or_add_ln()
    ln.append(ln.makeelement(_qn("a:prstDash"), {"val": "sysDot"}))
    return g


def _qn(tag):
    from pptx.oxml.ns import qn

    return qn(tag)


def _slide(prs, kit):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(s, 0, 0, SLIDE_W, SLIDE_H, fill=kit["rgb"]["bg"])
    return s


def build(mod, kit, box=DEFAULT_BOX, *, guide=False):
    """부품 모듈 → 목업 Presentation. 반환: (prs, 실패 로그)

    사람이 판정하는 것은 **디자인**이다 — 그래서 한 장에 케이스 하나만, 페이지를 채워 크게 찍는다.
    자리 가이드(점선)는 내부 검증용이라 `--guide`로만 켠다.
    """
    meta = mod.META
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(SLIDE_W), Inches(SLIDE_H)
    fails = []

    for sample in mod.SAMPLES:
        label, data, ov = (sample + (None,))[:3] if len(sample) < 3 else sample
        s = _slide(prs, kit)
        dense.compact_header(s, f"{meta['name']} — {meta['물성']}", label)
        if guide:
            _guide(s, box, kit)
        try:
            mod.draw(s, box, data, kit, ov)
        except Exception as e:  # 수용 한계 초과는 시끄럽게 죽는 게 정상 — 목업에 기록한다
            fails.append((label, str(e)))
            add_text(
                s,
                box[0],
                box[1] + box[3] / 2,
                box[2],
                0.4,
                f"[한계 초과] {e}",
                kit["sizes"]["caption"],
                kit["fonts"]["body"],
                kit["rgb"]["accent"],
            )
    return prs, fails


def render_png(pptx_path, png_dir):
    """전장 PNG — 렌더 정본 `render_deck.ps1`에 위임한다. 눈검증용(실패해도 pptx는 남는다).

    COM을 직접 부르지 않는다: 정본은 PowerPoint 단일 인스턴스를 전역 뮤텍스로 직렬화하고
    `sNN.png` 이름으로 내보낸다. 직접 부르면 동시 렌더가 서로의 Quit에 죽고, 파일명이
    `슬라이드NN.PNG`로 나와 **Stop 게이트의 렌더 감지(`s[0-9][0-9].png`)가 못 본다**
    — 렌더를 하고도 "렌더 없음"으로 차단됐다(2026-07-26 실측).
    """
    png_dir.mkdir(parents=True, exist_ok=True)
    r = subprocess.run(
        [
            "pwsh",
            "-File",
            str(Path(__file__).parent / "render_deck.ps1"),
            "-Pptx",
            str(pptx_path),
            "-OutDir",
            str(png_dir),
        ],
        capture_output=True,
        text=True,
    )
    return r.returncode == 0, (r.stderr or "")[-300:]


def run(name, do_render=True, guide=False):
    mod = importlib.import_module(f"goldenfab.figures.{name}")
    kit = load_kit()
    prs, fails = build(mod, kit, guide=guide)
    OUT_PPTX.mkdir(parents=True, exist_ok=True)
    out = OUT_PPTX / f"부품_{name}.pptx"
    try:
        prs.save(str(out))
    except PermissionError:
        # 사용자가 목업을 열어둔 채 재빌드하는 건 정상 흐름이다 — 죽이지 말고 옆에 저장하고 알린다.
        out = out.with_name(out.stem + "_new.pptx")
        prs.save(str(out))
        print("  (원본이 PowerPoint에 열려 있어 _new 로 저장했다 — 닫고 다시 돌리면 원본을 덮는다)")
    print(f"built: {out}  ({len(prs.slides)}장)")
    for label, err in fails:
        print(f"  [한계] {label}: {err.splitlines()[0][:110]}")
    if do_render:
        ok, err = render_png(out, OUT_PNG / name)
        print(f"  render: {'OK ' + str(OUT_PNG / name) if ok else 'FAIL ' + err}")
    return out


def main():
    args = [a for a in sys.argv[1:] if not a.startswith("--")]
    do_render = "--no-render" not in sys.argv
    if "--all" in sys.argv:
        # 부품 판별은 **META 유무**다 — 파일명 규칙으로 거르면 1층 원소(elements.py)까지
        # 부품으로 찍으려다 죽는다. 계약을 가진 것만 부품이다.
        names = sorted(
            p.stem
            for p in FIG_DIR.glob("*.py")
            if not p.stem.startswith("_") and "META" in p.read_text(encoding="utf-8")
        )
    elif args:
        names = args
    else:
        print(__doc__)
        sys.exit(2)
    for n in names:
        run(n, do_render, guide="--guide" in sys.argv)


if __name__ == "__main__":
    main()
