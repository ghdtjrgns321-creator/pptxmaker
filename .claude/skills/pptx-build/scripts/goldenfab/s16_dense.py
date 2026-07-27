"""S16 정직한 한계(경계) — dense v6(2026-07-21). hero_card 재건 + §8 게이트 통과.

v5 실패(사용자 반려): 하단 한계 3장을 add_box+텍스트로 손으로 그려 밋밋했다 — dense.py의
'카드=hero_card 단일'(design-rules §8)을 어긴 것. 이번엔 정본 컴포넌트 `hero_card`(배지+태그+
배너+히어로 아이콘+2단 불릿+칩)로 재건 → 통과작 기술카드급 리치. `check_adhoc_card`(§8 기계
게이트)가 이제 밋밋 카드를 막는다.

P0 북극성:
- #1 형태는 탭의 일 — 경계 = 안↔밖 도해(상단, 밖의 질문이 경계에 부딪혀 차단) + 정직한 한계
  3종(하단, hero_card). 카드로 도해를 뭉개지 않는다.
- #2 밀도 — 안 상자는 내용 크기에 맞춤(스탯 욱여넣기·빈칸 금지). 카드는 hero_card 풀 스펙.
- #3 시각화 — 도해(화살표·✕·차단벽) + 실측 수치(밖 4건·2건, 안 78/92·0). 다리 문장이 s15 14건과 연결.

콘텐츠 파라미터화(DEFAULT 노출 — 승격 대비). 헬퍼는 s16 자립.
"""

from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from . import dense as D
from . import grid as G
from .kit import SLIDE_H, SLIDE_W, add_box, add_text, load_kit

K = load_kit()
C, S, F = K["rgb"], K["sizes"], K["fonts"]

DEFAULT = {
    "kicker": "4. 문제와 해결",
    "headline": "정직한 한계 — 이 시스템이 서 있는 경계",
    "sublead": "안이면 결정적으로 답하고, 밖이면 거절·유보 — 못 하는 것까지 실측으로 세어 두었다",
    "outside_label": "경계 밖",
    "inside_title": "K-IFRS 1115 — 경계 안",
    "inside_policy": "등재된 기준서·용어 안에서만 결정적으로 답한다 — 아는 척이 구조적으로 불가능한 설계.",
    "inside_evidence": [
        ("홀드아웃 ", {}),
        ("78 / 92", {"bold": True, "color": C["accent"]}),
        (" 재현 · 실행 에러 ", {}),
        ("0", {"bold": True, "color": C["accent"]}),
        ("건 · 격리 0/92", {}),
    ],
    "out_left_head": "타 기준서 질문",
    "out_left_sub": "IAS 38 · 1002 · 1008 · 1037",
    "out_left_verdict": "강제 OUT",
    "out_left_cnt": "4건",
    "out_right_head": "등재되지 않은 표현",
    "out_right_sub": "임베딩식 유사 확장 없음",
    "out_right_verdict": "유보",
    "out_right_cnt": "2건",
    "bridge": "미재현 14건 중 6건이 이 경계에서 나온다 — 타 기준서 4 · 진입 누락 2. 나머지는 경계 안의 결함이 아니라 생성·검증 계층의 문제였다.",
    "limits": [
        {
            "num": "1",
            "title": "결론 미확정",
            "tag": "생성 계층의 한계",
            "ic": "scale",
            "banner": "헤지 2건 · 결론 유보",
            "items": [
                ("근거는 확보", " — 문단 다 찾고도 미확정"),
                ("정답 부재", " — 만들어낼 답 자체가 없음"),
            ],
            "chips": [
                ("circle-x", "", "타 기준서·범위 밖 — 검색 대상 아님"),
                ("target", "", "실제 개선은 생성 계층에서만"),
            ],
        },
        {
            "num": "2",
            "title": "인용 완전성",
            "tag": "결론과 별개 축",
            "ic": "file-text",
            "banner": "하드 인용 재현 59.1%",
            "items": [
                ("결론 맞아도", " — 근거 전부 인용은 못 함"),
                ("정확도와 독립", " — 인용률은 별개 지표"),
            ],
            "chips": [
                ("triangle-alert", "", "결론 정확도와는 별개 축이다"),
                ("search", "", "인용률은 별도 지표로 추적"),
            ],
        },
        {
            "num": "3",
            "title": "검증 자체의 한계",
            "tag": "측정 신뢰도",
            "ic": "gauge",
            "banner": "92건 전수 열람 · 비홀드아웃",
            "items": [
                ("순수 홀드아웃 아님", " — 성능 과대 여지"),
                ("재검증 필요", " — 안 본 질문으로 확인"),
            ],
            "chips": [
                ("search", "", "순수 홀드아웃이 아니라 과대 여지"),
                ("circle-check-big", "", "안 본 질문으로 최종 재검증"),
            ],
        },
    ],
    "source": "출처: 6_TEST-DECISIONS.md §6.3 · 4_SEARCH-PIPELINE.md (라우팅 강제 OUT) · 00_factsheet.md §D·§H",
}

# ── 경계 도해 좌표(dense, 압축 — 하단 hero_card 3장에 세로 양보) ──
CANVAS_Y, CANVAS_H = 1.12, 1.9
IN_W, IN_H = 4.5, 1.38
IN_X = (SLIDE_W - IN_W) / 2
IN_Y = 1.42
LINE_Y = IN_Y + IN_H / 2
CHIP_Y, CHIP_H = LINE_Y - 0.33, 0.66
BRIDGE_Y = 3.14  # 캔버스 하단(3.02)에서 AIR_MIN 0.12 확보 — 3.12면 0.10으로 §6 미달
CARD_Y, CARD_H = 3.46, 3.0


def _dash_arrow(slide, x1, y1, x2, y2):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = C["muted"]
    conn.line.width = Pt(1.5)
    ln = conn.line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "sm", "len": "sm"}))


def _chip(slide, x, w, head, sub):
    chip = add_box(
        slide, x, CHIP_Y, w, CHIP_H, fill=C["bg"], line=C["muted"], line_w=1.0, shape="round"
    )
    tf = chip.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.12)
    for i, (txt, col) in enumerate(((head, C["primary"]), (sub, C["muted"]))):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p._p.get_or_add_pPr().set("eaLnBrk", "0")
        r = p.add_run()
        r.text = txt
        r.font.name = F["head"] if i == 0 else F["body"]
        r.font.bold = i == 0
        r.font.size = Pt(S["caption"])
        r.font.color.rgb = col
        r.font.language_id = MSO_LANGUAGE_ID.KOREAN


def _block_mark(slide, bar_x, x_char, verdict, cnt, label_x):
    """차단 표식 — 경계 벽(primary 바) + ✕(충돌점) + 판정·실측(accent)."""
    add_box(slide, bar_x, LINE_Y - 0.22, 0.05, 0.44, fill=C["primary"])
    add_text(
        slide, x_char, LINE_Y - 0.14, 0.3, 0.28, "✕", S["body"], F["head"], C["primary"], bold=True
    )
    add_text(
        slide,
        label_x,
        LINE_Y - 0.52,
        0.95,
        0.22,
        verdict,
        S["foot"],
        F["head"],
        C["accent"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        label_x,
        LINE_Y - 0.3,
        0.95,
        0.2,
        "실측 " + cnt,
        S["foot"],
        F["body"],
        C["muted"],
        align=PP_ALIGN.CENTER,
    )


def build(prs, c=None):
    c = {**DEFAULT, **(c or {})}
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    D.compact_header(slide, c["kicker"], c["headline"])
    add_text(
        slide, G.MARGIN_L, 0.76, D.FULL_W, 0.24, c["sublead"], S["caption"], F["body"], C["muted"]
    )

    # ══ 상단: 경계 도해 ══
    add_box(slide, G.MARGIN_L, CANVAS_Y, D.FULL_W, CANVAS_H, fill=C["bg_alt"])
    add_text(
        slide,
        G.MARGIN_L + 0.2,
        CANVAS_Y + 0.1,
        3.0,
        0.24,
        c["outside_label"],
        S["caption"],
        F["head"],
        C["muted"],
        bold=True,
    )
    region = add_box(
        slide, IN_X, IN_Y, IN_W, IN_H, fill=C["bg"], line=C["primary"], line_w=2.5, shape="round"
    )
    region.adjustments[0] = 0.06
    add_text(
        slide,
        IN_X,
        IN_Y + 0.14,
        IN_W,
        0.3,
        c["inside_title"],
        S["head"],
        F["head"],
        C["primary"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        IN_X + 0.28,
        IN_Y + 0.48,
        IN_W - 0.56,
        0.44,
        c["inside_policy"],
        S["caption"],
        F["body"],
        C["muted"],
        align=PP_ALIGN.CENTER,
        line_spacing=1.2,
    )
    D.hrule(slide, IN_X + 0.5, IN_Y + 0.98, IN_W - 1.0, color=D.line_mid, weight=0.75)
    add_text(
        slide,
        IN_X,
        IN_Y + 1.06,
        IN_W,
        0.28,
        [c["inside_evidence"]],
        S["caption"],
        F["body"],
        C["primary"],
        align=PP_ALIGN.CENTER,
    )

    # 밖 좌: 타 기준서 → 강제 OUT
    _chip(slide, 0.85, 2.5, c["out_left_head"], c["out_left_sub"])
    _dash_arrow(slide, 3.42, LINE_Y, IN_X - 0.42, LINE_Y)
    _block_mark(slide, IN_X - 0.05, IN_X - 0.40, c["out_left_verdict"], c["out_left_cnt"], 3.45)

    # 밖 우: 미등재 표현 → 유보
    rcx = G.RIGHT_EDGE - 2.5
    _chip(slide, rcx, 2.5, c["out_right_head"], c["out_right_sub"])
    _dash_arrow(slide, rcx, LINE_Y, IN_X + IN_W + 0.47, LINE_Y)
    _block_mark(
        slide,
        IN_X + IN_W,
        IN_X + IN_W + 0.13,
        c["out_right_verdict"],
        c["out_right_cnt"],
        rcx - 0.95,
    )

    # ══ 다리: 경계 ↔ 떨어진 14건 연결 ══
    add_text(
        slide,
        G.MARGIN_L,
        BRIDGE_Y,
        D.FULL_W,
        0.24,
        [
            [
                ("경계의 비용은 실측돼 있다 — ", {"bold": True, "color": C["primary"]}),
                (c["bridge"], {}),
            ]
        ],
        S["foot"],
        F["body"],
        C["muted"],
    )

    # ══ 하단: 정직한 한계 N종 — hero_card(정본 컴포넌트). 폭은 한계 수에서 파생(§F 가로 짝) ══
    gap = 0.27
    sec_w = G.track(len(c["limits"]), G.MARGIN_L, G.RIGHT_EDGE, gap, 2.2, what="한계 카드")
    for i, lim in enumerate(c["limits"]):
        sx = G.MARGIN_L + i * (sec_w + gap)
        D.hero_card(
            slide,
            sx,
            CARD_Y,
            sec_w,
            CARD_H,
            lim["num"],
            lim["title"],
            lim["tag"],
            lim["ic"],
            lim["banner"],
            lim["items"],
            lim["chips"],
        )

    D.source_line(slide, c["source"])
    return slide
