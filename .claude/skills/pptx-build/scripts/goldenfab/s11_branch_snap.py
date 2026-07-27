"""S11 주입 뒤 분기 섹션 단독 덤프 — 1:N flat tree 정렬 스냅 (사용자 스펙 2026-07-11).

구조: [주입된 질문] ─┬─ 사실관계 부족  되물음…
                     ├─ 산술 필요     계산을…
                     └─ 그 외         확정…
실행: 이 타입은 goldenfab 레지스트리 경유로만 렌더된다 — 골든 19장 확인은
      `uv run python golden/build_golden.py`(2026-07-15 단일화로 시안 개별 실행 경로 폐지).
"""

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches, Pt

from . import grid as G
from .kit import SLIDE_H, SLIDE_W, add_box, add_text, load_kit, set_shape_text

K = load_kit()
C, S, F = K["rgb"], K["sizes"], K["fonts"]

# ── 스펙 상수 (행 y값 = 각 갈래의 세로중심, 커넥터가 닿는 선) ──
GROUP_TOP = 4.55
ROW_PITCH = 0.42
ROW_H = 0.32

BOX_X, BOX_W, BOX_H = 0.6, 1.35, 0.46
CONN_X = BOX_X + BOX_W + 0.15  # 2.1
LABEL_X = CONN_X + 0.25  # 2.35
LABEL_W = 1.35
DESC_X = LABEL_X + LABEL_W  # 3.7

BRANCHES = [
    ("사실관계 부족", "되물음 — 부족한 사실을 묻는다"),
    ("산술 필요", "계산을 확인한 뒤 답변"),
    ("그 외", "확정 또는 조건부 답변"),
]


def _line(slide, x1, y1, x2, y2):
    ln = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    ln.line.color.rgb = C["muted"]
    ln.line.width = Pt(1.25)
    return ln


def build_branch(slide, dy=0.0, branches=None, bottom=None):
    """1:N flat tree. dy: 본판 통합용 y 오프셋(상대구조 불변).

    2026-07-25 항목 수 파생: 갈래 수를 `range(3)`으로 박아두면 갈래가 4개인 content에서
    **4번째가 조용히 사라졌다**(`zip(centers, BRANCHES)`). 이제 갈래 수에서 중심을 파생하고,
    루트 박스는 블록의 세로중심(짝수 갈래면 행 사이)에 붙는다 — 부채꼴의 축은 항상 블록 중심이다.
    블록이 콘텐츠 하한을 넘으면 시끄럽게 죽인다(§F: 조용히 넘치는 값을 만들지 않는다).
    """
    branches = BRANCHES if branches is None else branches
    n = len(branches)
    if n < 1:
        raise ValueError("갈래 0개 — 그릴 것이 없다.")
    centers = [GROUP_TOP + i * ROW_PITCH + dy for i in range(n)]
    limit = G.CONTENT_BOTTOM if bottom is None else bottom
    block_bottom = centers[-1] + ROW_H / 2
    if block_bottom > limit + 1e-9:
        room = int((limit - ROW_H / 2 - GROUP_TOP - dy) / ROW_PITCH) + 1
        raise ValueError(
            f'갈래 {n}개는 이 자리에 안 들어간다 — 블록 bottom {block_bottom:.3f}" > 하한 '
            f'{limit:.3f}". 최대 {room}개. content를 줄이거나 레이아웃을 다시 설계할 것.'
        )
    mid = (centers[0] + centers[-1]) / 2
    root = add_box(slide, BOX_X, mid - BOX_H / 2, BOX_W, BOX_H, fill=C["primary"], shape="round")
    set_shape_text(root, "주입된 질문", S["body"], F["head"], C["bg"], bold=True)
    _line(slide, BOX_X + BOX_W, mid, CONN_X, mid)  # 박스 → 세로줄
    _line(slide, CONN_X, centers[0], CONN_X, centers[-1])  # 세로줄 — 첫 갈래 ~ 마지막 갈래
    shapes = {"root": root, "chips": [], "hlines": []}
    chip_w, chip_h = 5.4, 0.36
    from pptx.enum.lang import MSO_LANGUAGE_ID
    from pptx.enum.text import PP_ALIGN

    for cy, (cond, desc) in zip(centers, branches):
        shapes["hlines"].append(_line(slide, CONN_X, cy, LABEL_X - 0.02, cy))
        chip = add_box(
            slide,
            LABEL_X,
            cy - chip_h / 2,
            chip_w,
            chip_h,
            fill=C["bg_alt"],
            shape="round",
        )
        tf = chip.text_frame
        tf.word_wrap = False
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = Inches(0.15)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p._p.get_or_add_pPr().set("eaLnBrk", "0")
        r1 = p.add_run()
        r1.text = cond
        r1.font.name, r1.font.bold = F["head"], True
        r1.font.size = Pt(S["body"])
        r1.font.color.rgb = C["primary"]
        r2 = p.add_run()
        r2.text = f"    {desc}"
        r2.font.name = F["body"]
        r2.font.size = Pt(S["body"])
        r2.font.color.rgb = C["muted"]
        for r in (r1, r2):
            r.font.language_id = MSO_LANGUAGE_ID.KOREAN
        shapes["chips"].append(chip)
    return shapes, centers


def audit(shapes, centers):
    EMU = 914400
    rows = []
    root = shapes["root"]
    root_c = (root.top + root.height / 2) / EMU
    mid = (centers[0] + centers[-1]) / 2  # 부채꼴 축 = 블록 세로중심(갈래 수에서 파생)
    rows.append(("주입박스 세로중심", f"{root_c:.3f}", f"{mid:.3f}", root_c - mid))
    chip_ls = [sh.left / EMU for sh in shapes["chips"]]
    rows.append(
        ("칩 left 편차", f"{min(chip_ls):.3f}", f"{max(chip_ls):.3f}", max(chip_ls) - min(chip_ls))
    )
    for i, (sh, cy) in enumerate(zip(shapes["chips"], centers)):
        c = (sh.top + sh.height / 2) / EMU
        rows.append((f"칩{i + 1} 세로중심", f"{c:.3f}", f"{cy:.3f}", c - cy))
    for i, (ln, cy) in enumerate(zip(shapes["hlines"], centers)):
        y1 = ln.begin_y / EMU
        y2 = ln.end_y / EMU
        rows.append(
            (
                f"가로줄{i + 1} y=갈래중심",
                f"{y1:.3f}/{y2:.3f}",
                f"{cy:.3f}",
                max(abs(y1 - cy), abs(y2 - cy)),
            )
        )
    print(f"| {'항목':<18} | {'실측':>12} | {'기준':>8} | {'편차':>6} |")
    fails = 0
    for name, got, want, dev in rows:
        ok = abs(dev) < 0.005
        fails += 0 if ok else 1
        print(f"| {name:<18} | {got:>12} | {want:>8} | {dev:+.3f} {'OK' if ok else 'FAIL'} |")
    assert fails == 0, f"SNAP AUDIT FAIL {fails}건"
    print("snap audit pass")


def main():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    shapes, centers = build_branch(slide)
    audit(shapes, centers)
    out = Path(__file__).parent / "variants" / "s11_branch_snap.pptx"
    prs.save(str(out))
    print(f"saved {out}")


if __name__ == "__main__":
    main()
