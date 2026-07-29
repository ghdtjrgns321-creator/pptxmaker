"""goldenfab.layouts — 골든 확정 레이아웃 cover·toc·part.

**여기가 이 3종의 확정 구현이자 유일 소스다**(2026-07-15 단일화).
좌표·서식 임의 변경 금지(design-rules). 콘텐츠만 인자(dict)로 받는다.
회귀 검증기(compare_golden)는 2026-07-29 아카이브 — 골든이 고정 기준선이 되어 대조 대상이 없다.
"""

from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

from . import grid as G  # noqa: F401 (후속 타입에서 사용)
from .kit import SLIDE_H, SLIDE_W, add_box, add_text, load_kit, mix, set_shape_text

K = load_kit()
C, S, F = K["rgb"], K["sizes"], K["fonts"]
MARGIN = K["layout"]["margin"]

# ── 표지 GRID (사용자 확정값 2026-07-10 — 임의 좌표 금지) ──
COVER_L = 0.65
COVER_R = SLIDE_W - COVER_L
COVER_RULE_Y = 0.92
COVER_ACCENT_TOP = 2.75
COVER_TITLE_TOP = 3.15
COVER_SUB_TOP = 4.25
COVER_YEAR_TOP = 6.35
COVER_BAND_TOP = 6.95


def cover(prs, c):
    """표지. c: {kicker, title, value_prop, year, brand_name}."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rule_w = COVER_R - COVER_L
    add_text(
        slide, COVER_L, 0.55, 9.0, 0.3, c["kicker"], S["caption"], F["head"], C["muted"], bold=True
    )
    add_box(slide, COVER_L, COVER_RULE_Y, rule_w, 0.014, fill=C["muted"])
    add_box(slide, COVER_L, COVER_ACCENT_TOP, 0.3, 0.3, fill=C["accent"])
    add_text(
        slide,
        COVER_L,
        COVER_TITLE_TOP,
        rule_w,
        0.75,
        c["title"],
        S["display"],
        F["head"],
        C["primary"],
        bold=True,
    )
    add_text(
        slide, COVER_L, COVER_SUB_TOP, 11.0, 0.35, c["value_prop"], S["sub"], F["body"], C["text"]
    )
    add_text(slide, COVER_L, COVER_YEAR_TOP, 4.0, 0.28, c["year"], S["head"], F["body"], C["muted"])
    band_h = SLIDE_H - COVER_BAND_TOP
    add_box(slide, 0, COVER_BAND_TOP, SLIDE_W, band_h, fill=C["primary"])
    add_text(
        slide,
        COVER_L,
        COVER_BAND_TOP + band_h / 2 - 0.15,
        6.0,
        0.3,
        c["brand_name"],
        S["head"],
        F["head"],
        C["bg"],
        bold=True,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    return slide


# ── 목차 GRID (vC5 확정값) ──
TOC_TITLE_TOP, TOC_RULE_Y = 0.5, 1.35
TOC_PANEL_X = 9.6
TOC_LIST_R = TOC_PANEL_X - 0.5
TOC_NUM_X, TOC_TEXT_X = MARGIN, 1.35
TOC_ROW1_Y, TOC_ROW_PITCH = 2.0, 0.82


def toc(prs, c):
    """목차. c: {title, items: [(제목, 서브, 페이지)...], deck_title, deck_sub}."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    add_box(slide, TOC_PANEL_X, 0, SLIDE_W - TOC_PANEL_X, SLIDE_H, fill=C["primary"])
    ph_c = mix(C["primary"], C["muted"], 0.45)
    ph_x, ph_y, ph_w, ph_h = TOC_PANEL_X + 0.9, 0.9, (SLIDE_W - TOC_PANEL_X) - 1.8, 0.9
    add_box(slide, ph_x, ph_y, ph_w, ph_h, fill=None, line=ph_c, line_w=1.0)
    add_text(
        slide,
        ph_x,
        ph_y + ph_h / 2 - 0.15,
        ph_w,
        0.3,
        "회사 로고",
        S["caption"],
        F["body"],
        ph_c,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        TOC_PANEL_X + 0.35,
        SLIDE_H - 1.15,
        SLIDE_W - TOC_PANEL_X - 0.7,
        0.7,
        [[(c["deck_title"], {"bold": True, "color": C["bg"]})], [(c["deck_sub"], {})]],
        S["caption"],
        F["body"],
        C["bg_alt"],
        line_spacing=1.4,
    )
    add_text(
        slide,
        MARGIN,
        TOC_TITLE_TOP,
        4.0,
        0.55,
        c["title"],
        S["title"],
        F["head"],
        C["primary"],
        bold=True,
    )
    add_box(slide, MARGIN, TOC_RULE_Y, TOC_LIST_R - MARGIN, 0.014, fill=C["muted"])
    n = len(c["items"])
    for i, (title, subs, page) in enumerate(c["items"]):
        y = TOC_ROW1_Y + i * TOC_ROW_PITCH
        add_text(
            slide, TOC_NUM_X, y, 0.6, 0.3, f"0{i + 1}", S["head"], F["head"], C["accent"], bold=True
        )
        add_text(
            slide,
            TOC_TEXT_X,
            y - 0.02,
            6.0,
            0.32,
            title,
            S["head"],
            F["head"],
            C["primary"],
            bold=True,
        )
        add_text(
            slide,
            TOC_TEXT_X,
            y + 0.3,
            TOC_LIST_R - TOC_TEXT_X - 1.0,
            0.26,
            subs,
            S["caption"],
            F["body"],
            C["muted"],
        )
        add_text(
            slide,
            TOC_LIST_R - 0.9,
            y,
            0.9,
            0.3,
            f"p.{page}",
            S["body"],
            F["body"],
            C["muted"],
            align=PP_ALIGN.RIGHT,
        )
        if i < n - 1:
            add_box(
                slide,
                TOC_TEXT_X,
                y + TOC_ROW_PITCH - 0.17,
                TOC_LIST_R - TOC_TEXT_X,
                0.01,
                fill=C["bg_alt"],
            )
    return slide


# ── 간지 GRID (확정값) ──
PART_L = 0.65
PART_TAG_TOP, PART_TAG_W, PART_TAG_H = 2.0, 1.3, 0.42
PART_TITLE_TOP = 2.55
PART_SUB_TOP = 4.4
PART_SUB_H = 0.35
PART_SUB_X = PART_L + 0.05 + 0.2
PART_DOTS_TOP = 5.75
PART_DOT_D, PART_DOT_PITCH = 0.13, 0.35
PART_WM_RIGHT = 12.0


def part(prs, c):
    """간지. c: {no(1~), title, lead, total(기본 6)}."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    no, total = c["no"], c.get("total", 6)
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["primary"])
    add_text(
        slide,
        PART_WM_RIGHT - 6.0,
        1.9,
        6.0,
        2.6,
        f"0{no}",
        S["display"] * 4,
        F["head"],
        mix(C["primary"], C["muted"], 0.22),
        bold=True,
        align=PP_ALIGN.RIGHT,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    pill = add_box(slide, PART_L, PART_TAG_TOP, PART_TAG_W, PART_TAG_H, fill=C["accent"])
    set_shape_text(pill, f"PART 0{no}", S["caption"], F["head"], C["bg"], bold=True)
    add_text(
        slide,
        PART_L,
        PART_TITLE_TOP,
        7.5,
        0.75,
        c["title"],
        S["display"],
        F["head"],
        C["bg"],
        bold=True,
    )
    add_box(slide, PART_L, PART_SUB_TOP, 0.05, PART_SUB_H, fill=C["accent"])
    add_text(
        slide,
        PART_SUB_X,
        PART_SUB_TOP,
        9.5,
        PART_SUB_H,
        c["lead"],
        S["sub"],
        F["body"],
        C["bg_alt"],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    for i in range(total):
        dot_fill = C["accent"] if i == no - 1 else C["muted"]
        add_box(
            slide,
            PART_L + i * PART_DOT_PITCH,
            PART_DOTS_TOP,
            PART_DOT_D,
            PART_DOT_D,
            fill=dot_fill,
            shape="oval",
        )
    return slide


LAYOUTS = {"cover": cover, "toc": toc, "part": part}
