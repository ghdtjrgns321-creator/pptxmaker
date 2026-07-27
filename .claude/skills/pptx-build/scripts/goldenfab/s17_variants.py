"""S17 시안 — 차별점: 일반 임베딩 RAG 대비 7축 (아키타입: 경쟁 매트릭스 + 다크 요지 사이드바).

모방: golden/ref/s14_matrix.png (McKinsey p26 — 본문 매트릭스 + 우측 다크 패널).
콘텐츠 실물: 00_factsheet.md §G — 7축 전수 + 정직 조항(1115 한정·타 기준서 4건·유연성 없음).
실행: 이 타입은 goldenfab 레지스트리 경유로만 렌더된다 — 골든 19장 확인은
      `uv run python golden/build_golden.py`(2026-07-15 단일화로 시안 개별 실행 경로 폐지).
"""

from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from . import grid as G
from .kit import SLIDE_H, SLIDE_W, add_box, add_text, load_kit, mix

K = load_kit()
C, S, F = K["rgb"], K["sizes"], K["fonts"]

KICKER = "5. 차별점"
SOURCE = "출처: 00_factsheet.md §G (6_TEST-DECISIONS.md·4_SEARCH-PIPELINE.md)"

AXES = [  # §G 7축 전수
    ("검색 진입", "임베딩 유사도 — 점수 경쟁", "용어사전 문자 매칭 — 결정적"),
    ("근거 설명", "유사도 점수 (근거 불명)", "그래프 경로 — 어떤 간선을 지났는지"),
    ("판단 절차", "없음 — 생성이 즉석 구성", "원문 앵커 판단트리 41개 사전 조립"),
    ("AI 개입 범위", "검색 전 구간", "용어 색인 1개로 한정"),
    ("오류 방향", "허위 확정 (2종 오류)", "못 찾으면 유보 — 1종 안전"),
    ("데이터 검증", "불가 — 벡터는 검산이 없다", "회계 항등식 전수 검증"),
    ("재현성", "비결정 — 같은 질문, 다른 답", "결정적 — 같은 질문, 같은 경로"),
]


def header(slide, headline, kicker=KICKER):
    add_text(
        slide, G.MARGIN_L, 0.42, 8.0, 0.28, kicker, S["caption"], F["head"], C["muted"], bold=True
    )
    add_text(
        slide,
        G.MARGIN_L,
        0.72,
        G.RIGHT_EDGE - G.MARGIN_L,
        0.55,
        headline,
        S["section"],
        F["head"],
        C["primary"],
        bold=True,
    )
    add_box(slide, G.MARGIN_L, G.RULE_Y, G.RIGHT_EDGE - G.MARGIN_L, 0.014, fill=C["muted"])


def variant_a(prs):
    """A — 좌 7축 네이티브 매트릭스 + 우 다크 요지 사이드바(McKinsey p26 구도)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    header(slide, "일반 임베딩 RAG과 무엇이 다른가 — 7개 축 전부 같은 방향")
    from pptx.enum.lang import MSO_LANGUAGE_ID

    # ── 좌: 매트릭스 ──
    tbl_x, tbl_y, tbl_w = G.MARGIN_L, 2.0, 8.6
    n_rows = len(AXES) + 1
    gf = slide.shapes.add_table(
        n_rows, 3, Inches(tbl_x), Inches(tbl_y), Inches(tbl_w), Inches(0.5 * n_rows)
    )
    tbl = gf.table
    for ci, cw in enumerate((1.65, 3.3, 3.65)):
        tbl.columns[ci].width = Inches(cw)
    data = [("비교 축", "일반 임베딩 RAG", "이 시스템")] + list(AXES)
    for ri, row in enumerate(data):
        tbl.rows[ri].height = Inches(0.42 if ri == 0 else 0.5)
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            if ri == 0:
                cell.fill.fore_color.rgb = C["primary"]
            elif ci == 2:
                cell.fill.fore_color.rgb = mix(C["bg"], C["accent"], 0.08)
            else:
                cell.fill.fore_color.rgb = C["bg_alt"] if ri % 2 == 0 else C["bg"]
            cell.margin_left = cell.margin_right = Inches(0.1)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p._p.get_or_add_pPr().set("eaLnBrk", "0")
            r = p.add_run()
            r.text = val
            r.font.name = F["head"] if ri == 0 or ci == 0 else F["body"]
            r.font.size = Pt(S["caption"])
            r.font.bold = ri == 0 or ci == 0 or ci == 2
            if ri == 0:
                r.font.color.rgb = C["bg"]
            elif ci == 1:
                r.font.color.rgb = C["muted"]
            else:
                r.font.color.rgb = C["primary"]
            r.font.language_id = MSO_LANGUAGE_ID.KOREAN
    # ── 우: 다크 요지 사이드바 ──
    px = 9.55
    pw = G.RIGHT_EDGE - px  # 3.18
    add_box(slide, px, 1.75, pw, G.CONTENT_BOTTOM - 1.75, fill=C["primary"])
    tx, tw = px + 0.3, pw - 0.6
    add_text(
        slide,
        tx,
        2.15,
        tw,
        1.3,
        [
            [("차이는 성능이 아니라 구조다", {"bold": True, "color": C["bg"]})],
            [
                (
                    "일곱 축 모두 같은 문장으로 요약된다 — 확률 신호가 하던 일을 기준서의 구조가 대신한다.",
                    {},
                )
            ],
        ],
        S["body"],
        F["body"],
        C["bg_alt"],
        line_spacing=1.35,
    )
    add_box(slide, tx, 3.85, tw, 0.014, fill=mix(C["primary"], C["muted"], 0.5))
    add_text(
        slide,
        tx,
        4.1,
        tw,
        1.9,
        [
            [("정직한 한계", {"bold": True, "color": C["bg"]})],
            [
                (
                    "코퍼스는 1115호 한정 — 타 기준서 질문 4건은 재현하지 못했다. "
                    "임베딩식의 유연한 유사 확장도 없다. 대신 그 경계 안에서는 결정적으로 동작한다.",
                    {},
                )
            ],
        ],
        S["caption"],
        F["body"],
        C["bg_alt"],
        line_spacing=1.35,
    )
    # 결론 바 + 출처
    bar = add_box(slide, G.MARGIN_L, G.BAR_Y, G.RIGHT_EDGE - G.MARGIN_L, G.BAR_H, fill=C["primary"])
    tf = bar.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "일반 RAG의 확률 신호 자리마다 기준서의 구조가 들어가 있다 — 7축 전부, 예외 없이"
    run.font.name, run.font.bold = F["head"], True
    run.font.size = Pt(S["body"])
    run.font.color.rgb = C["bg"]
    add_text(
        slide,
        G.MARGIN_L,
        G.SOURCE_Y,
        G.RIGHT_EDGE - G.MARGIN_L,
        0.3,
        SOURCE,
        S["foot"],
        F["body"],
        C["muted"],
    )
    return slide


MIRROR = [  # (축, RAG 키워드, 시스템 키워드) — §G 7축, 키워드 우선
    ("검색 진입", "임베딩 유사도\n점수 경쟁", "용어사전 문자 매칭\n결정적 진입"),
    ("근거 설명", "유사도 점수\n이유를 설명 못 함", "그래프 경로\n지나온 간선이 곧 근거"),
    ("판단 절차", "없음\n생성이 즉석에서 구성", "판단트리 41개\n원문 앵커로 사전 조립"),
    ("AI 개입", "검색 전 구간", "용어 색인 1개로 한정"),
    ("오류 방향", "허위 확정\n2종 오류", "못 찾으면 유보\n1종 안전"),
    ("데이터 검증", "불가\n벡터는 검산이 없다", "회계 항등식\n전수 검증"),
    ("재현성", "같은 질문, 다른 답", "같은 질문, 같은 경로"),
]


def _wing_card(slide, x, y, w, h, text, dashed, align_right, head_pt=None):
    from pptx.enum.lang import MSO_LANGUAGE_ID
    from pptx.oxml.ns import qn

    card = add_box(slide, x, y, w, h, fill=C["bg"] if not dashed else C["bg_alt"])
    ln_el = card.line
    ln_el.color.rgb = C["muted"] if dashed else C["primary"]
    ln_el.width = Pt(1.0 if dashed else 1.25)
    if dashed:
        _l = ln_el._get_or_add_ln()
        _l.append(_l.makeelement(qn("a:prstDash"), {"val": "dash"}))
    tf = card.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.12)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    lines = text.split("\n")
    for li, seg in enumerate(lines):
        p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.RIGHT if align_right else PP_ALIGN.LEFT
        p._p.get_or_add_pPr().set("eaLnBrk", "0")
        if li == 0 and not dashed:
            rc = p.add_run()
            rc.text = "✓  "
            rc.font.name, rc.font.bold = F["head"], True
            rc.font.size = Pt(head_pt or S["caption"])
            rc.font.color.rgb = C["primary"]
        r = p.add_run()
        r.text = seg
        r.font.name = F["head"] if li == 0 else F["body"]
        r.font.bold = li == 0
        r.font.size = Pt((head_pt or S["caption"]) if li == 0 else S["caption"])
        r.font.color.rgb = (C["muted"] if dashed else C["primary"]) if li == 0 else C["muted"]
        r.font.language_id = MSO_LANGUAGE_ID.KOREAN
    return card


def variant_b(prs):
    """B — 미러 매트릭스: 중앙 축 기둥, 좌 점선 날개(RAG) vs 우 실선 날개(이 시스템)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    header(slide, "일반 임베딩 RAG과 무엇이 다른가 — 7개 축 전부 같은 방향")
    from .kit import set_shape_text

    wing_l_x, wing_w = G.MARGIN_L, 3.15
    spine_x, spine_w = 3.95, 1.3
    wing_r_x = 5.45
    wing_r_w = 3.45
    row0, pitch, card_h = 2.42, 0.55, 0.47
    add_text(
        slide,
        wing_l_x,
        2.02,
        wing_w,
        0.3,
        "일반 임베딩 RAG",
        S["body"],
        F["head"],
        C["muted"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        wing_r_x,
        2.02,
        wing_r_w,
        0.3,
        "이 시스템",
        S["body"],
        F["head"],
        C["primary"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    for i, (axis, rag, ours) in enumerate(MIRROR):
        y = row0 + i * pitch
        _wing_card(slide, wing_l_x, y, wing_w, card_h, rag, dashed=True, align_right=True)
        pill = add_box(
            slide, spine_x, y + (card_h - 0.34) / 2, spine_w, 0.34, fill=C["primary"], shape="round"
        )
        set_shape_text(pill, axis, S["caption"], F["head"], C["bg"], bold=True)
        _wing_card(slide, wing_r_x, y, wing_r_w, card_h, ours, dashed=False, align_right=False)
    # ── 우: 다크 요지 사이드바 ──
    px = 9.2
    pw = G.RIGHT_EDGE - px
    add_box(slide, px, 1.75, pw, G.CONTENT_BOTTOM - 1.75, fill=C["primary"])
    tx, tw = px + 0.3, pw - 0.6
    add_text(
        slide,
        tx,
        2.15,
        tw,
        1.5,
        [
            [("차이는 성능이 아니라\n구조다", {"bold": True, "color": C["bg"]})],
            [
                (
                    "일곱 축 모두 같은 문장으로 요약된다 — 확률 신호가 하던 일을 기준서의 구조가 대신한다.",
                    {},
                )
            ],
        ],
        S["body"],
        F["body"],
        C["bg_alt"],
        line_spacing=1.35,
    )
    add_box(slide, tx, 4.0, tw, 0.014, fill=mix(C["primary"], C["muted"], 0.5))
    add_text(
        slide,
        tx,
        4.25,
        tw,
        1.9,
        [
            [("정직한 한계", {"bold": True, "color": C["bg"]})],
            [
                (
                    "코퍼스는 1115호 한정 — 타 기준서 질문 4건은 재현하지 못했다. "
                    "임베딩식의 유연한 유사 확장도 없다. 대신 그 경계 안에서는 결정적으로 동작한다.",
                    {},
                )
            ],
        ],
        S["caption"],
        F["body"],
        C["bg_alt"],
        line_spacing=1.35,
    )
    bar = add_box(slide, G.MARGIN_L, G.BAR_Y, G.RIGHT_EDGE - G.MARGIN_L, G.BAR_H, fill=C["primary"])
    tf = bar.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "일반 RAG의 확률 신호 자리마다 기준서의 구조가 들어가 있다 — 7축 전부, 예외 없이"
    run.font.name, run.font.bold = F["head"], True
    run.font.size = Pt(S["body"])
    run.font.color.rgb = C["bg"]
    add_text(
        slide,
        G.MARGIN_L,
        G.SOURCE_Y,
        G.RIGHT_EDGE - G.MARGIN_L,
        0.3,
        SOURCE,
        S["foot"],
        F["body"],
        C["muted"],
    )
    return slide


# ── variant_c 재설계 좌표(2026-07-16 밀도 재깎기) — GRID에서 파생, 아래 파생식 외 리터럴 금지 ──
GRP_X = G.MARGIN_L  # 좌측 독자-질문 칼럼
GRP_W = 1.55
WING_GAP = 0.22  # 날개 ↔ 스파인 — 치환 커넥터(점선→칩→실선 화살표)가 지나갈 폭
WING_L_X = GRP_X + GRP_W + 0.20  # 2.35
SPINE_W = 1.35
WING_W = (G.RIGHT_EDGE - WING_L_X - SPINE_W - 2 * WING_GAP) / 2  # 4.37 — 좌우 동일(대결 대칭)
SPINE_X = WING_L_X + WING_W + WING_GAP
WING_R_X = SPINE_X + SPINE_W + WING_GAP
COLHEAD_Y = 2.02
ROW0 = 2.42
PITCH_CAP = 0.54  # 행 피치 상한 — 실제 피치는 행 수에서 파생(grid.pitch)
CARD_H = 0.42
GGAP = 0.12  # 그룹 경계 추가 공기(헤어라인 자리) — 경계 총 공기 = (피치-카드) + GGAP
LAST_BOTTOM = 6.32  # 마지막 카드 bottom 한계 (하한 6.35 - 공기 여유)

VARIANT_C_DEFAULTS = {  # 텍스트 내용만 — 좌표·색·폰트·도형종류는 코드에 고정
    "kicker": KICKER,
    # 헤드라인 "같은 방향"은 본문 주장("전부 다르다")과 첫 독해가 충돌한다(채점 지적) — 배너와 정렬
    "headline": "일반 임베딩 RAG과 무엇이 다른가 — 7개 축, 예외 없이",
    "left_head": "일반 임베딩 RAG",
    "spine_head": "비교 축",
    "right_head": "이 시스템",
    # (번호, 독자 질문, 요지, [(축, RAG 2줄, 이 시스템 2줄), ...]) — §G 7축을 독자 질문 3개로 묶음.
    # 질문은 15pt 1줄(≤9자 권장·마지막 그룹만 2줄 여유), 카드 줄은 키워드+기전(§H 키워드 1~2줄).
    "groups": [
        (
            "01",
            "어떻게 찾나",
            "점수 경쟁 없이 사전과 그래프가 경로를 정한다",
            [
                (
                    "검색 진입",
                    "임베딩 유사도\n점수 경쟁으로 문서를 고른다",
                    "용어사전 문자 매칭\n결정적 진입 — 점수 경쟁이 없다",
                ),
                (
                    "근거 설명",
                    "유사도 점수 하나\n왜 그 답인지 설명하지 못한다",
                    "그래프 경로\n지나온 간선이 곧 근거다",
                ),
            ],
        ),
        (
            "02",
            "누가 판단하나",
            "판단은 사전 조립 — AI는 색인 하나",
            [
                (
                    "판단 절차",
                    "사전 절차 없음\n생성이 즉석에서 판단을 구성",
                    "판단트리 41개\n원문 앵커로 사전 조립",
                ),
                (
                    "AI 개입",
                    "검색 전 구간\n진입부터 생성까지 AI 판단",
                    "용어 색인 1개로 한정\n그래프·간선은 기계 생성",
                ),
            ],
        ),
        (
            "03",
            "틀리면 어떻게 되나",
            "치명 방향(허위 확정)을 구조로 차단",
            [
                (
                    "오류 방향",
                    "허위 확정 — 2종 오류\n틀린 답을 확신 어조로 낸다",
                    "못 찾으면 유보 — 1종 안전\n놓침은 추가 검토로 회복",
                ),
                (
                    "데이터 검증",
                    "검산 불능\n벡터는 합을 물을 수 없다",
                    "회계 항등식 전수 검증\n언급 = 해소 + 외부 · 미해소 0",
                ),
                (
                    "재현성",
                    "비결정\n같은 질문에 답이 매번 다르다",
                    "결정적\n같은 질문, 같은 경로, 같은 답",
                ),
            ],
        ),
    ],
    "bar": "일반 RAG의 확률 신호 자리마다 기준서의 구조가 들어가 있다 — 7축 전부, 예외 없이",
    "source": SOURCE,
}


def _q_lines(question):
    """15pt 질문이 GRP_W 안에서 몇 줄로 접히는지 — 글자폭 근사(오딧과 같은 자: 한글 1자 ≈ 0.8pt)."""
    per_line = int(GRP_W / (S["head"] * 0.8 / 72))
    return 1 if len(question) <= per_line else 2


def variant_c(prs, c=None):
    """C — 그룹 미러 매트릭스: 좌측 독자-질문 칼럼(3그룹) + 중앙 축 기둥 + 좌우 날개 대결.

    2026-07-16 재깎기: 이전 판은 날개 5.17" 카드에 키워드 5~6자가 떠 있어 채움률 14~29%
    (오딧 FAIL 9건), 15pt 런 0개로 28pt→9pt 직락 위계였다. 7축을 독자 질문 3개(어떻게 찾나 /
    누가 판단하나 / 틀리면 어떻게 되나)로 묶어 구획을 만들고, 카드를 두 줄(키워드 12pt bold +
    기전 9pt)로 밀도를 올렸다. c=텍스트 override(None=골든).
    """
    c = {**VARIANT_C_DEFAULTS, **(c or {})}
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    header(slide, c["headline"], c["kicker"])
    from .kit import set_shape_text

    groups = c["groups"]
    n_rows = sum(len(axes) for _no, _q, _cl, axes in groups)
    # 피치는 행·그룹 수에서 파생 — 축이 늘면 좁아지고, 수용 한계를 넘으면 시끄럽게 죽는다(§F)
    eff_bottom = LAST_BOTTOM - (len(groups) - 1) * GGAP
    pitch = G.pitch(n_rows, ROW0, eff_bottom, CARD_H, cap=PITCH_CAP, what="비교 축")
    add_text(
        slide,
        WING_L_X,
        COLHEAD_Y,
        WING_W,
        0.26,
        c["left_head"],
        S["body"],
        F["head"],
        C["muted"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        SPINE_X,
        COLHEAD_Y,
        SPINE_W,
        0.26,
        c["spine_head"],
        S["caption"],
        F["head"],
        C["muted"],
        bold=True,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.BOTTOM,
    )
    add_text(
        slide,
        WING_R_X,
        COLHEAD_Y,
        WING_W,
        0.26,
        c["right_head"],
        S["body"],
        F["head"],
        C["primary"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    def _swap_arrow(x1, x2, yc, dashed):
        """치환 커넥터 — 확률 신호(점선)가 들어와 구조(실선 화살표)로 나간다. 좌→우 한 방향."""
        from pptx.enum.shapes import MSO_CONNECTOR
        from pptx.oxml.ns import qn

        conn = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(yc), Inches(x2), Inches(yc)
        )
        conn.line.color.rgb = C["muted"] if dashed else C["primary"]
        conn.line.width = Pt(1.25)
        ln = conn.line._get_or_add_ln()
        if dashed:
            ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))
        ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "sm", "len": "sm"}))
        return conn

    i = 0  # 전역 행 인덱스
    for gi, (no, question, claim, axes) in enumerate(groups):
        g_top = ROW0 + i * pitch + gi * GGAP
        if gi > 0:  # 그룹 경계 헤어라인 — 경계 총 공기의 정중앙
            sep_y = g_top - ((pitch - CARD_H) + GGAP) / 2
            add_box(slide, G.MARGIN_L, sep_y, G.RIGHT_EDGE - G.MARGIN_L, 0.012, fill=C["bg_alt"])
        # 좌측 질문 블록: 번호(accent) → 질문(15pt) → 요지(caption)
        add_text(
            slide,
            GRP_X,
            g_top - 0.02,
            GRP_W,
            0.2,
            no,
            S["caption"],
            F["head"],
            C["accent"],
            bold=True,
        )
        ql = _q_lines(question)
        add_text(
            slide,
            GRP_X,
            g_top + 0.18,
            GRP_W,
            0.28 * ql,
            question,
            S["head"],
            F["head"],
            C["primary"],
            bold=True,
        )
        add_text(
            slide,
            GRP_X,
            g_top + 0.18 + 0.28 * ql + 0.04,
            GRP_W,
            0.34,
            claim,
            S["caption"],
            F["body"],
            C["muted"],
            line_spacing=1.2,
        )
        for axis, rag, ours in axes:
            y = ROW0 + i * pitch + gi * GGAP
            yc = y + CARD_H / 2
            _wing_card(
                slide,
                WING_L_X,
                y,
                WING_W,
                CARD_H,
                rag,
                dashed=True,
                align_right=True,
                head_pt=S["body"],
            )
            # 치환 커넥터: 확률 신호(점선)가 축으로 들어와 구조(실선 ▶)로 우측에 배달된다
            _swap_arrow(WING_L_X + WING_W + 0.02, SPINE_X - 0.02, yc, dashed=True)
            _swap_arrow(SPINE_X + SPINE_W + 0.02, WING_R_X - 0.02, yc, dashed=False)
            pill = add_box(
                slide,
                SPINE_X,
                y + (CARD_H - 0.34) / 2,
                SPINE_W,
                0.34,
                fill=C["primary"],
                shape="round",
            )
            set_shape_text(pill, axis, S["caption"], F["head"], C["bg"], bold=True)
            _wing_card(
                slide,
                WING_R_X,
                y,
                WING_W,
                CARD_H,
                ours,
                dashed=False,
                align_right=False,
                head_pt=S["body"],
            )
            i += 1
    bar = add_box(slide, G.MARGIN_L, G.BAR_Y, G.RIGHT_EDGE - G.MARGIN_L, G.BAR_H, fill=C["primary"])
    tf = bar.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = c["bar"]
    run.font.name, run.font.bold = F["head"], True
    run.font.size = Pt(S["body"])
    run.font.color.rgb = C["bg"]
    add_text(
        slide,
        G.MARGIN_L,
        G.SOURCE_Y,
        G.RIGHT_EDGE - G.MARGIN_L,
        0.3,
        c["source"],
        S["foot"],
        F["body"],
        C["muted"],
    )
    return slide


def audit(prs):
    from pptx.oxml.ns import qn

    EMU = 914400
    fails = []
    for si, sl in enumerate(prs.slides):
        accents, tables = 0, 0
        for sh in sl.shapes:
            t, le = sh.top / 914400, sh.left / 914400
            b, r = t + sh.height / 914400, le + sh.width / 914400
            full_bleed = sh.width / 914400 >= SLIDE_W - 0.05
            if t < 6.4 and b > 6.37 and not full_bleed:
                fails.append((si, "bottom", round(b, 2), sh.shape_id))
            if r > 12.75 and not full_bleed:
                fails.append((si, "right", round(r, 2), sh.shape_id))
            if sh.has_table:
                tables += 1
                if len(sh.table.rows) != 8:
                    fails.append((si, "rows≠8", len(sh.table.rows)))
            for el in sh._element.iter(qn("a:srgbClr")):
                if el.get("val") == "D66E3A":
                    if el.getparent().tag.endswith("}solidFill"):
                        accents += 1
        if accents > 4:
            fails.append((si, "accent>4", accents))
        if si == 0 and tables != 1:
            fails.append((si, "table≠1", tables))
        if si in (1, 2):
            wing_h = 0.47 if si == 1 else CARD_H  # variant_c는 재깎기 카드 높이
            wings = sum(
                1
                for sh in sl.shapes
                if abs(sh.height / EMU - wing_h) < 0.015 and sh.width / EMU > 3.0
            )
            pills = sum(
                1
                for sh in sl.shapes
                if abs(sh.height / EMU - 0.34) < 0.005 and 1.25 < sh.width / EMU < 1.45
            )
            if wings != 14 or pills != 7:
                fails.append((si, "mirror", wings, pills))
        if (
            si == 2
        ):  # variant_c 재깎기 구조: 그룹 3(15pt 질문·accent 번호) + 경계 헤어라인 2 + 채움률·판정 대비
            from . import audit as GA

            n_groups = len(VARIANT_C_DEFAULTS["groups"])
            heads15 = sum(
                1
                for sh in sl.shapes
                for text, _col, pt, bold in GA.runs_of(sh)
                if pt == S["head"] and bold and text.strip()
            )
            if heads15 != n_groups:
                fails.append((si, f"15pt 질문≠{n_groups}", heads15))
            seps = sum(
                1
                for sh in sl.shapes
                if abs(sh.height / EMU - 0.012) < 0.003
                and sh.width / EMU > 11
                and sh.top / EMU > 2.0
            )
            if seps != n_groups - 1:
                fails.append((si, f"그룹 헤어라인≠{n_groups - 1}", seps))
            from pptx.enum.shapes import MSO_SHAPE_TYPE

            n_axes = sum(len(axes) for _n2, _q2, _c2, axes in VARIANT_C_DEFAULTS["groups"])
            conns = sum(1 for sh in sl.shapes if sh.shape_type == MSO_SHAPE_TYPE.LINE)
            if conns != n_axes * 2:  # 치환 커넥터 — 행마다 점선 진입 + 실선 배달
                fails.append((si, f"치환 커넥터≠{n_axes * 2}", conns))
            shapes_c = list(sl.shapes)
            for rule_name, (ok, msg, _n) in (
                ("채움률", GA.check_fill_ratio(shapes_c)),
                ("판정 대비", GA.check_verdict_contrast(shapes_c, None)),
            ):
                if not ok:
                    fails.append((si, rule_name, msg))
    assert not fails, f"AUDIT FAIL {fails}"
    print("audit pass")


def main():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    variant_a(prs)
    variant_b(prs)
    variant_c(prs)
    audit(prs)
    out = Path(__file__).parent / "variants" / "s17_variants.pptx"
    prs.save(str(out))
    print(f"saved {out}")


if __name__ == "__main__":
    main()
