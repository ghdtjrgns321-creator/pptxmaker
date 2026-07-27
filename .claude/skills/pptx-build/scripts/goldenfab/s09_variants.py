"""S9 시안 — 기술 2: 지식그래프 (§8d 기술 장 문법: 서사 3칼럼 + 실물 증거. UI 패스).

재료: 00_factsheet.md §C·§D (원문 2_DATA-TAXONOMY·3_KNOWLEDGE-GRAPH).
실물 증거: 좌 = 위계 트리(한 개념에서 뻗는 fan) / 우 = 간선 7종 카탈로그(종류별 실물 1건).
실행: 이 타입은 goldenfab 레지스트리 경유로만 렌더된다 — 골든 19장 확인은
      `uv run python golden/build_golden.py`(2026-07-15 단일화로 시안 개별 실행 경로 폐지).

2026-07-15 사용자 반려 "무엇을어떻게만들었나 시각화는 이해가 안되는 시각화임" → 우측 표를
카탈로그로 교체. 이 docstring이 폐기 전까지 "3D 캡처 + 노드·간선 등록부 표"라고 적혀 있었는데
3D 캡처는 **한 번도 렌더된 적이 없었다**(IMG_3D 상수는 선언만 되고 미사용) — 표로 갈아치우며
주석을 안 고친 거짓 기록이라 함께 제거했다.
"""

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from . import grid as G
from .kit import SLIDE_H, SLIDE_W, add_box, add_text, load_kit, set_shape_text
from .s08_variants import (
    chip_w,
)  # 칩 폭 파생식 단일 출처 — 두 장이 같은 자를 쓴다(피치는 grid.pitch)

K = load_kit()
C, S, F = K["rgb"], K["sizes"], K["fonts"]

# 콘텐츠 기본값(골든 내용) — c=None이면 이 값, override 시 텍스트만 교체(좌표·색 고정)
DEFAULT = {
    "kicker": "3. 기술 설명 — TECH 02 · Retrieve",
    "headline": "지식그래프 — 기준서의 구조를 그대로 옮긴 결정적 지도",
    "narratives": [
        (
            "왜 필요한가",
            "문단 뭉치에는 순서도 관계도 없다. 기준서는 계층·상호참조·사례가 얽힌 구조 — 그 구조를 보존해야 근거가 근거로 이어진다.",
        ),
        (
            "파이프라인에서의 역할",
            "Retrieve의 지도. 진입 개념에서 관계를 한 단계만 따라 문단·사례·근거를 수집한다 — 수집 범위 자체가 그래프로 고정된다.",
        ),
        (
            "어떻게 만들었나",
            "기준서 공식 소제목 80개를 그대로 개념 노드로, 뼈대 간선은 기준서에서 기계 생성 — AI 판단 0. 고립 노드 0을 감사로 확인했다.",
        ),
    ],
    "struct_head": "구조 — 기준서의 위계 그대로",
    "root": "기준서 1115",
    "concepts": ["변동대가", "보증", "⋯ 80"],
    "paras": ["문단 50", "문단 56", "B33"],
    # 거터 라벨에 **단위**를 박는다 — 2026-07-15 제3자 채점 FAIL(#5): `문단 250`(개수)과
    # `문단 50`(문단 번호)이 같은 행에 93px 간격으로 나란해 어느 쪽이 개수고 어느 쪽이 ID인지
    # 글자만으로 못 갈랐다. 이 장은 "명사+숫자" 꼴이 개수·문서번호·문단ID 세 뜻을 겸한다
    # (기준서 1115=문서번호 · 개념 80=개수 · 문단 50=ID). 오독을 유발한 게 아니라 강제했다.
    "layer_tag1": "개념층 80개",
    "layer_tag2": "문단층 250개",
    "term_chip": "“볼륨디스카운트”",
    "term_cap": "용어 423 등재 · 400 진입",  # 등재는 423, 그래프 노드는 간선 보유 400 — 층이 다르다
    "case_chip": "사례 188",
    "cat_title": (
        "간선 7종 — 무엇이 무엇을 잇고, 어디서 왔나",
        "    7종 2,629 + BC 근거층 65 = 2,694 · 고립 0",
    ),
    # (종류, 수, 출발 노드, 종류 라벨, 도착 노드, 출처 규칙, 출발 종류, 도착 종류)
    # 전부 data/ontology 실측 1건 — 창작 0. 노드 종류: concept|para|term|case
    # 1~5행이 전부 **변동대가** 하나에 달려 있다(좌측 트리와 같은 개념) — 7종이 따로 노는 게
    # 아니라 한 개념에 실제로 다 붙는다는 걸 실물로 보인다. 6·7행은 `계약을 식별함`으로 연결.
    "cat_rows": [
        (
            "hier",
            "79",
            "변동대가",
            "속한다",
            "거래가격을 산정함",
            "기준서 목차 level · 기계 생성",
            "concept",
            "concept",
        ),
        ("cp", "250", "변동대가", "관할", "문단 50", "documentId · 기계 생성", "concept", "para"),
        (
            "e3",
            "244",
            "문단 56",
            "인용",
            "문단 53",
            "원문에 “문단 53” · 회계 항등식",
            "para",
            "para",
        ),
        (
            "case",
            "1,220",
            "QNA 201909A",
            "인용",
            "문단 50",
            "감사본 related_paragraphs",
            "case",
            "para",
        ),
        (
            "term",
            "755",
            "“볼륨디스카운트”",
            "진입",
            "변동대가",
            "사람 1차자료 전수검수 (TECH 01)",
            "term",
            "concept",
        ),
        (
            "example",
            "74",
            "IE 사례 1",
            "공식예시",
            "계약을 식별함",
            "IE 목차 매핑 · 사람 확정",
            "case",
            "concept",
        ),
        (
            "e2",
            "7",
            "계약을 식별함",
            "먼저판단",
            "수행의무를 식별함",
            "5단계 모형 1→2 · 사람 확정",
            "concept",
            "concept",
        ),
    ],
    # 노드 모양 범례는 두지 않는다 — 칩 안의 실제 텍스트("문단 50"·"변동대가")가 종류를 이미
    # 말한다. 범례를 얹으면 해독표를 하나 더 만드는 셈(§6 "축 매트릭스보다 스토리 구도").
    "node_caption": "노드 929 = 개념 80 + 문단 250 + 사례 188 + BC 그룹 11 + 용어 400(423 중 간선 보유)",
    "bar": "임베딩이 놓치는 '법적 이웃'을 관계가 잡는다 — 텍스트가 아니라 구조로 검색한다",
    "source": "출처: 2_DATA-TAXONOMY.md · 3_KNOWLEDGE-GRAPH.md (00_factsheet.md §C·§D)",
}


def _arrow(slide, x1, y1, x2, y2, *, width=1.5, size="med"):
    """방향 있는 커넥터(§5 — 커넥터는 전부 화살표). 트리·카탈로그 공용.

    variant_a 안의 클로저였던 것을 모듈 레벨로 올렸다 — build_catalog가 같은 화살표 문법을
    써야 하는데 함수 안에 갇혀 있으면 복붙본이 하나 더 생긴다(골든이 두 벌이었던 그 병).
    """
    from pptx.oxml.ns import qn

    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = C["muted"]
    conn.line.width = Pt(width)
    ln = conn.line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": size, "len": size}))
    return conn


def header(slide, headline, kicker):
    add_text(
        slide, G.MARGIN_L, 0.42, 8.0, 0.28, kicker, S["caption"], F["head"], C["muted"], bold=True
    )
    add_text(
        slide,
        G.MARGIN_L,
        0.72,
        G.RIGHT_EDGE - G.MARGIN_L,
        0.55,
        headline,
        S["section"],
        F["head"],
        C["primary"],
        bold=True,
    )
    add_box(slide, G.MARGIN_L, G.RULE_Y, G.RIGHT_EDGE - G.MARGIN_L, 0.014, fill=C["muted"])


def narrative_row(slide, narratives):
    nar_w = (G.RIGHT_EDGE - G.MARGIN_L - 2 * 0.4) / 3
    for i, (head, body) in enumerate(narratives):
        nx = G.MARGIN_L + i * (nar_w + 0.4)
        add_text(
            slide,
            nx,
            G.CONTENT_TOP,
            nar_w,
            0.28,
            [(f"0{i + 1}", {"color": C["accent"]}), (f"  {head}", {})],
            S["head"],
            F["head"],
            C["primary"],
            bold=True,
        )
        add_text(
            slide,
            nx,
            G.CONTENT_TOP + 0.36,
            nar_w,
            0.95,
            body,
            S["caption"],
            F["body"],
            C["muted"],
            line_spacing=1.25,
        )
        if i < 2:
            add_box(slide, nx + nar_w + 0.2, G.CONTENT_TOP, 0.012, 1.2, fill=C["bg_alt"])


# ── 간선 카탈로그 GRID (표를 걷어낸 자리 — 2026-07-15 사용자 반려 "이해가 안 가는 시각화") ──
#
# 물성 선언(design-rules P1): **실물 대조(종류별 1건)**. 이 장이 말하려는 건 "무엇이 무엇을
#   잇는가"인데, 표는 그걸 40자 문장으로 서술하고 `수` 열에 숫자를 세웠다. 게다가 그 열에
#   **노드와 간선이 섞여**(80은 노드, 79·1,220·244는 간선) 더할 수 없는 숫자가 한 줄에 섰다.
#   합이 941인데 제목은 929였고, 표라서 아무도 검산하지 않았다(간선도 폐기값 2,697을 썼다).
#   → 각 행이 **실제 노드 도해**(`[A] ─종류─▶ [B]`). 글자가 아니라 도형과 화살표다.
#   노드 종류는 **모양**으로 구분한다(개념=round 진한 테두리 / 문단=rect / 사례=rect bg_alt /
#   용어=round 흰 칩). §6 "숫자 스탯 단독 나열 금지 — 관계 인코딩과 결합".
#
# 어휘 분리(§6 장 내 중복 금지): 좌측 트리 = **fan**(한 개념에서 뻗음) / 우측 카탈로그 =
#   **행 대조**(종류별 1건 병렬). 같은 노드를 두 각도로 본다 — 1~5행이 전부 좌측 트리의
#   `변동대가`에 달려 있어 두 구획이 실물로 맞물린다.
#
# 색(P4⑩): accent 예산 **1**(3칼럼 번호 런이 3을 이미 씀, 상한 4). 카탈로그 행에는 accent를
#   쓰지 않는다 — 7행 중 3행(term·example·e2)이 "사람 확정"이라 색으로 구분하고 싶지만 예산이
#   없고, 무엇보다 이 장의 주장은 "사람 개입이 **적다**"이다. 출처 열의 글자가 그 일을 한다.
CAT_X = 5.00
CAT_W = G.RIGHT_EDGE - CAT_X  # 7.733
CAT_HEAD_Y = 3.30
CAT_TOP = 3.66  # 소제목 bottom(3.54)에서 공기 0.12 — 3.62는 0.08로 모자랐다(§6)
CAT_BOTTOM = G.CONTENT_BOTTOM  # 6.35
CAT_H = 0.26  # 행 노드 높이 — 계열 균일(P4⑧)
KIND_W = 1.05  # 종류·수 앵커 열
# 노드 칩 폭은 **글자에서 파생**(s08.chip_w 재사용) — 균일 1.72"로 그렸더니 "변동대가" 4자가
# 21%, "문단 50"이 26%로 오딧 채움률 7건이 걸렸다. 정렬은 폭이 아니라 **화살표 구간**으로 잡는다:
# A는 우변(A_RIGHT) 고정, B는 좌변(B_LEFT) 고정 → 화살표가 모든 행에서 같은 자리·같은 길이.
A_RIGHT = CAT_X + KIND_W + 1.72  # 7.77
ARROW_W = 1.05
B_LEFT = A_RIGHT + ARROW_W  # 8.82
B_MAX_W = 1.85
SRC_X = B_LEFT + B_MAX_W + 0.30  # 10.97 — 출처 열
SRC_W = G.RIGHT_EDGE - SRC_X  # 1.763


def _cat_pitch(n):
    """행 피치를 **행 수에서 파생**(design-rules §F) + 겹침 가드. 자는 `grid.pitch` 정본."""
    return G.pitch(n, CAT_TOP, CAT_BOTTOM, CAT_H, what="카탈로그 행")


# ── 노드 클래스 인코딩 — **이 장의 단일 출처** (좌 트리 · 우 카탈로그 공용) ──
#
# 2026-07-15 제3자 채점 FAIL(s9 2/10, 지적 #10): 좌 트리는 `bg_alt` 채움을 **내부 노드**
# (개념·문단)에 쓰고, 우 카탈로그는 같은 `bg_alt`를 **사례**에 썼다. 독자가 왼쪽에서 배운
# "회색 = 그래프 내부 노드"가 오른쪽에서 "회색 = 사례"로 뒤집힌다. `변동대가`는 한 장 안에서
# 회색 덩어리(좌)이자 흰 바탕 검정 테두리 볼드(우)였다 — **노드 클래스가 이 장의 주제인데
# 그 인코딩이 자기모순**이었다.
#
# 원인은 구조적이다: 트리의 `_node`(fill 기본 bg_alt)와 카탈로그의 `_node_chip`이 **별개
# 함수**라 각자 색을 정했다. 골든이 두 벌이어서 갈라졌던 것과 같은 병 — 그래서 dict 하나로
# 합친다. 노드를 그리는 길은 이 함수 하나뿐이다.
NODE_STYLE = {
    "root": dict(fill=C["primary"], line=None, line_w=None, shape="round"),  # 기준서 최상위
    "concept": dict(
        fill=C["bg"], line=C["primary"], line_w=1.0, shape="round"
    ),  # 기준서가 세운 개념
    "para": dict(fill=C["bg"], line=C["muted"], line_w=0.75, shape="rect"),  # 본문 문단
    "case": dict(fill=C["bg_alt"], line=None, line_w=None, shape="rect"),  # 외부 사례(QNA·IE·감리)
    "term": dict(fill=C["bg"], line=C["muted"], line_w=0.75, shape="round"),  # 사람이 등재한 진입어
}
BOLD_KINDS = ("root", "concept")


def _node_chip(slide, x, y, w, text, kind, *, h=None, size=None):
    """노드를 그리는 **유일한 길**. 종류 = 모양 + 채움(NODE_STYLE). 표에서는 전부 같은 글자였다."""
    b = add_box(slide, x, y, w, h or CAT_H, **NODE_STYLE[kind])
    set_shape_text(
        b,
        text,
        size or S["foot"],
        F["head"] if kind in BOLD_KINDS else F["body"],
        C["bg"] if kind == "root" else C["primary"],
        bold=(kind in BOLD_KINDS),
    )
    return b


def build_catalog(slide, rows):
    """간선 7종 실물 카탈로그. 각 행 = 실제 노드 도해 + 출처 규칙 1줄. 반환: 행 수(오딧용)."""
    pitch = _cat_pitch(len(rows))
    for i, (kind, n, a, label, b, src, ka, kb) in enumerate(rows):
        y = CAT_TOP + i * pitch
        cy = y + CAT_H / 2
        add_text(  # 종류 + 수 — 좌측 앵커(수치는 관계에 종속된 보조다)
            slide,
            CAT_X,
            y + 0.01,
            KIND_W,
            0.24,
            [(f"{kind} ", {"color": C["primary"], "bold": True}), (n, {"color": C["muted"]})],
            S["foot"],
            F["body"],
            C["muted"],
            anchor=MSO_ANCHOR.MIDDLE,
        )
        wa = chip_w(a, S["foot"])
        wb = min(chip_w(b, S["foot"]), B_MAX_W)
        _node_chip(slide, A_RIGHT - wa, y, wa, a, ka)  # A는 우변 고정
        _node_chip(slide, B_LEFT, y, wb, b, kb)  # B는 좌변 고정
        _arrow(slide, A_RIGHT, cy, B_LEFT, cy)
        add_text(  # 관계 라벨은 화살표 **위**에 — 선 위에 얹으면 글자가 선에 잘린다(P4⑥)
            slide,
            A_RIGHT,
            cy - 0.20,
            ARROW_W,
            0.17,
            label,
            S["foot"] - 1,
            F["body"],
            C["muted"],
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            SRC_X,
            y + 0.01,
            SRC_W,
            0.24,
            src,
            S["foot"] - 1,
            F["body"],
            C["muted"],
            anchor=MSO_ANCHOR.MIDDLE,
            line_spacing=1.05,
        )
    return len(rows)


def bar_and_source(slide, text, source):
    bar = add_box(slide, G.MARGIN_L, G.BAR_Y, G.RIGHT_EDGE - G.MARGIN_L, G.BAR_H, fill=C["primary"])
    tf = bar.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name, run.font.bold = F["head"], True
    run.font.size = Pt(S["body"])
    run.font.color.rgb = C["bg"]
    add_text(
        slide,
        G.MARGIN_L,
        G.SOURCE_Y,
        G.RIGHT_EDGE - G.MARGIN_L,
        0.3,
        source,
        S["foot"],
        F["body"],
        C["muted"],
    )


def variant_a(prs, c=None):
    """A — 서사 3칼럼 + 실물 증거(좌 위계 트리 + 우 간선 7종 카탈로그). c: 텍스트 override(None=골든 기본)."""
    c = {**DEFAULT, **(c or {})}
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    header(slide, c["headline"], c["kicker"])
    narrative_row(slide, c["narratives"])

    # ── 하단 좌: 위계 도해 (기준서의 질서 그대로) ──
    add_text(
        slide,
        G.MARGIN_L,
        3.3,
        4.0,
        0.24,
        c["struct_head"],
        S["caption"],
        F["head"],
        C["primary"],
        bold=True,
    )

    # 4층 트리 — 루트 → 개념(fan) → 문단(fan), 용어·사례 칩 부착. 예시 라벨은 실데이터.
    # 노드는 전부 `_node_chip`(NODE_STYLE 단일 출처)을 거친다 — 여기 별도 `_node`를 두었더니
    # 카탈로그와 색이 갈려 같은 회색이 좌우에서 반대 뜻이 됐다(제3자 채점 FAIL #10).
    TREE_H = 0.34

    def _node(x, y, w, name, kind):
        return _node_chip(slide, x, y, w, name, kind, h=TREE_H, size=S["caption"])

    root_x, root_w = 2.35, 1.6
    _node(root_x, 3.6, root_w, c["root"], "root")
    # 개념층 (y 4.45): 변동대가 · 보증 · ⋯ (80) — 좌표 고정, 텍스트만 c에서
    concept_pos = [(2.05, 1.0), (3.2, 0.7), (4.05, 0.65)]  # (cx, cw)
    for name, (cx, cw) in zip(c["concepts"], concept_pos):
        _node(cx, 4.45, cw, name, "concept")
        _arrow(slide, root_x + root_w / 2, 3.94, cx + cw / 2, 4.45)
    # 문단층 (y 5.35): 변동대가→50·56, 보증→B33 (상호참조 화살표 공간 확보용 간격)
    para_pos = [(1.75, 0.8, 2.55), (2.95, 0.8, 2.55), (3.95, 0.6, 3.55)]  # (px, pw, src_cx)
    for name, (px, pw, src_cx) in zip(c["paras"], para_pos):
        _node(px, 5.35, pw, name, "para")
        _arrow(slide, src_cx, 4.79, px + pw / 2, 5.35)
    # 층 태그 (좌측 열)
    add_text(
        slide,
        G.MARGIN_L,
        4.5,
        1.15,
        0.22,
        c["layer_tag1"],
        S["caption"],
        F["head"],
        C["primary"],
        bold=True,
    )
    add_text(
        slide,
        G.MARGIN_L,
        5.4,
        1.15,
        0.22,
        c["layer_tag2"],
        S["caption"],
        F["head"],
        C["primary"],
        bold=True,
    )
    # 용어 칩 → 변동대가 (aliases 실물: 볼륨디스카운트 → 변동대가)
    _node(G.MARGIN_L, 3.6, 1.45, c["term_chip"], "term")
    add_text(slide, G.MARGIN_L, 3.96, 1.5, 0.18, c["term_cap"], S["foot"], F["body"], C["muted"])
    _arrow(slide, G.MARGIN_L + 1.45, 3.77, 2.05, 4.62)  # 용어는 개념 좌변으로 진입(루트 fan과 분리)
    # 사례 칩 → 문단층 (인용 문단으로 연결)
    _node(G.MARGIN_L, 5.85, 1.1, c["case_chip"], "case")
    _arrow(slide, G.MARGIN_L + 1.1, 6.02, 1.75, 5.69)
    # 상호참조 양방향 (문단 50 ↔ 문단 56)
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(2.55), Inches(5.52), Inches(2.95), Inches(5.52)
    )
    conn.line.color.rgb = C["muted"]
    conn.line.width = Pt(1.25)
    _ln = conn.line._get_or_add_ln()
    from pptx.oxml.ns import qn as _qn

    _ln.append(_ln.makeelement(_qn("a:headEnd"), {"type": "triangle", "w": "sm", "len": "sm"}))
    _ln.append(_ln.makeelement(_qn("a:tailEnd"), {"type": "triangle", "w": "sm", "len": "sm"}))
    add_text(  # 노드 검산 — 표의 `수` 열이 노드·간선을 섞어 941≠929를 만들던 자리
        slide,
        G.MARGIN_L,
        6.21,  # `사례 188` 칩 bottom 6.19 아래 — 6.14는 칩과 0.05 겹쳤다
        4.3,
        0.14,
        c["node_caption"],
        S["foot"] - 1,
        F["body"],
        C["muted"],
    )
    # ── 하단 우: 간선 7종 실물 카탈로그 (무엇이 무엇을 잇고, 어디서 왔나) ──
    add_text(
        slide,
        CAT_X,
        CAT_HEAD_Y,
        CAT_W,
        0.24,
        [
            [
                (c["cat_title"][0], {"bold": True, "color": C["primary"], "font": F["head"]}),
                (c["cat_title"][1], {}),
            ]
        ],
        S["caption"],
        F["body"],
        C["muted"],
    )
    build_catalog(slide, c["cat_rows"])
    bar_and_source(slide, c["bar"], c["source"])
    return slide


def audit(prs):
    EMU = 914400
    fails = []
    for si, sl in enumerate(prs.slides):
        for sh in sl.shapes:
            L, T, W, H = sh.left / EMU, sh.top / EMU, sh.width / EMU, sh.height / EMU
            if L == 0 and T == 0 and W > 13:
                continue
            if T + H > G.CONTENT_BOTTOM + 0.02 and T < G.BAR_Y - 0.02:
                fails.append((si, "bottom>6.35", round(T + H, 3)))
            if 11.5 < W < 13 and abs((L + W) - G.RIGHT_EDGE) > 0.02:
                fails.append((si, "풀폭 right≠12.733", round(L + W, 3)))
    assert not fails, f"AUDIT FAIL {fails}"
    print("audit pass")


def main():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    variant_a(prs)
    audit(prs)
    out = Path(__file__).parent / "variants" / "s09_variants.pptx"
    prs.save(str(out))
    print(f"saved {out}")


if __name__ == "__main__":
    main()
