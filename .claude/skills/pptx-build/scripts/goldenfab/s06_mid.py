"""S06 중간본 — 골든 s06 ↔ 파일럿 s06 사이 (2026-07-20). 골든 미편입.

파일럿(밀도 최대: 헤더+플로우+4카드+하단 3패널)과 골든(공기 많음: 큰 헤더+리드 문단+
텍스트 4컬럼+검은 재서술 바) 사이의 중간지점. 사용자 확정 구성(2026-07-20):

- 헤더: 파일럿처럼 압축 (키커 + 14pt 한 줄 + 룰). 골든의 28pt·리드 문단 폐지.
- 플로우 도해: 파일럿 원형 **그대로** (줄이지 않음). 이 슬라이드의 주인공.
- 4 아이콘 카드: 파일럿의 기호·아이콘 다용 유지. 3패널이 빠진 지면만큼 아래로 키워
  내부(배너·불릿·칩·아이콘)를 넉넉하게 — 빽빽함 완화, 스트레치 없이 지면을 채움.
- 제거: 골든의 검은 전폭 재서술 바 + 파일럿 하단 3패널(AS-IS/IMPACT/5레이어).

콘텐츠는 s06_pilot과 동일 재고(units.md v4). 창작 0.

실행:
    uv run python -c "
    import sys; sys.path.insert(0, '.claude/skills/pptx-build/scripts')
    from goldenfab.reference import new_presentation
    from goldenfab.s06_mid import build
    prs = new_presentation(); build(prs); prs.save('golden/_pilot_s06/mid.pptx')"
"""

from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from . import grid as G
from .kit import ROOT, SLIDE_H, SLIDE_W, add_box, add_text, load_kit, mix, set_shape_text

K = load_kit()
C, S, F = K["rgb"], K["sizes"], K["fonts"]
ICONS = ROOT / "assets/icons"

# ── 콘텐츠 (units.md v4 재고 — s06_pilot과 동일, 이 파일 밖 문구 금지) ──────
KICKER = "2. 파이프라인"
HEADLINE = "질문에서 응답까지 — 분기까지 전부 결정적인 실행 그래프"
SOURCE = "출처: 4_SEARCH-PIPELINE.md (00_factsheet.md §C)"

# 플로우 시각화 — 골든 s06(s06_variants.variant_c)에서 그대로 이식. 텍스트·도형·색 수정 금지.
NODE_NAMES = ["질문", "Analyze", "판정", "Retrieve", "Generate", "Format", "응답"]
NODE_TAGS = [
    "용어사전 매칭",
    "그래프 1홉 탐색",
    "판단트리 주입",
    "경고·꼬리질문",
]  # 노드 인덱스 1·3·4·5
IN_LABEL = "IN"
OUT_LABEL = "OUT"
REJECT_BOX = "거절 메시지"
REJECT_DESC = "범위 밖 질문은 본선에 진입하지 못하고 즉시 거절된다"

# 카드 4 = 골든 s06 '핵심 기술 4' (s06_variants.DEFAULT_C["details"] + NODES desc) 그대로.
# 플로우가 4노드를 보이므로 카드는 4기술 — 중복 없는 골든 원본 축.
# 불릿 = (볼드 리드, 설명) 2단 — 밀도는 올리되 서술형 금지.
NODE_CARDS = [  # (배지, 기술명, 태그, 히어로 아이콘, 배너=리드, 불릿4[(리드,설명)], 칩3[(아이콘,강조,라벨)])
    (
        "1",
        "용어사전",
        "후보 진입점",
        "book-open",
        "등재 423 · AI 신규 창작 0",
        [
            ("자료 3종", " 질의매핑·사례제목·부록A"),
            ("AI 초안 + 전수 검수", " 로 등재"),
            ("확정 라우터 아님", " — 후보 진입점"),
            ("substring 진입", " — 임베딩 아님"),
        ],
        [
            ("book-open", "288", " 질의 매핑"),
            ("file-text", "123", " 사례 제목"),
            ("target", "9", " 부록A 정의"),
        ],
    ),
    (
        "2",
        "지식그래프",
        "기계 생성 뼈대",
        "workflow",
        "노드 929 · 간선 2,697",
        [
            ("공식 소제목 →", " 개념 80 · 문단 250 배정"),
            ("계층·참조 간선", " 기준서에서 기계 생성"),
            ("AI 생성물", " 용어 색인 하나뿐"),
            ("탐색 hops=1", " 상호참조 E3 1홉"),
        ],
        [
            ("git-branch", "개념 80", ""),
            ("layers", "문단 250", ""),
            ("circle-x", "AI 창작 0", " 간선"),
        ],
    ),
    (
        "3",
        "판단트리",
        "판단 순서의 사전 조립",
        "git-branch",
        "조건-분기 골격 41 · 원문 앵커",
        [
            ("흩어진 판단 절차", " 미리 조립"),
            ("진입·트리거 개념", " 최다 겹침 주입"),
            ("주제 직속 트리", " 전부 프롬프트 주입"),
            ("via_topic 다중 주입", " 원문 앵커 포함"),
        ],
        [
            ("git-branch", "트리 41", " 조건-분기"),
            ("link", "원문 앵커", " 포함"),
            ("circle-check-big", "창작 0", ""),
        ],
    ),
    (
        "4",
        "구조화 출력",
        "코드 수준 강제",
        "lock",
        "PydanticAI 스키마 강제",
        [
            ("답변 형식", " 코드로 강제"),
            ("위반 시", " result_validator 자동 재시도"),
            ("출력 스키마", " 자유 서술 차단"),
            ("감리 경고·꼬리질문", " 마지막 보조 장치"),
        ],
        [
            ("lock", "스키마", " 강제"),
            ("settings", "재시도", " 자동"),
            ("file-text", "감리", " 보조"),
        ],
    ),
]
SIDE_NOTE = "초기 설계의 유사도 재정렬(rerank) 단계는 제거 — 그래프가 이미 결정적으로 선별하므로 재정렬할 것이 없다"

# ── GRID ──────────────────────────────────────────────────────────────────
FULL_W = G.RIGHT_EDGE - G.MARGIN_L
RULE_Y = 0.66
LANE_Y, FNODE_H = 1.05, 0.60  # 골든 플로우 레인 — variant_c 원본 지오메트리 (압축 헤더 밑에 배치)
CARD_Y, CARD_H = 2.85, 3.80  # 카드: 골든 플로우 밑, 히어로 아이콘으로 속 채운 축소 카드
HERO_ICON = 0.48  # 카드 중앙 히어로 아이콘 크기 (불릿 밀도 확보 위해 축소)
SIDE_NOTE_Y = 6.74
SOURCE_Y = 7.14

line_soft = mix(C["primary"], C["bg"], 0.80)
line_mid = mix(C["primary"], C["bg"], 0.55)
ink_body = mix(C["muted"], C["primary"], 0.45)


def _icon(slide, name, variant, x, y, size):
    p = ICONS / f"{name}_{variant}.png"
    return slide.shapes.add_picture(str(p), Inches(x), Inches(y), Inches(size), Inches(size))


def _arrow(slide, x1, y1, x2, y2, color, w=1.5):
    """골든 플로우 화살표 — 실선 + 삼각 머리."""
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = color
    conn.line.width = Pt(w)
    ln = conn.line._get_or_add_ln()
    tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
    ln.append(tail)
    return conn


def _hrule(slide, x, y, w, color=None, weight=0.75):
    return add_box(slide, x, y, w, 0.001, fill=None, line=color or line_soft, line_w=weight)


def _chip_row(slide, x, y, w, icon, em, label, *, h=0.30):
    """아이콘 칩 — 중간본은 파일럿보다 크게 (h·아이콘·폰트 상향)."""
    chip = add_box(slide, x, y, w, h, fill=C["bg_alt"])
    tf = chip.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.30)
    p = tf.paragraphs[0]
    p._p.get_or_add_pPr().set("eaLnBrk", "0")
    p.alignment = PP_ALIGN.LEFT
    for txt, color, bold in ((em, C["accent"], True), (label, C["muted"], False)):
        if not txt:
            continue
        r = p.add_run()
        r.text = txt
        r.font.name = F["head"]
        r.font.size = Pt(S["caption"])
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.language_id = MSO_LANGUAGE_ID.KOREAN
    _icon(slide, icon, "primary", x + 0.08, y + (h - 0.17) / 2, 0.17)


def build(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])

    # ── 헤더 — 파일럿과 동일 압축 (키커 + 14pt 한 줄 + 룰) ──
    add_text(
        slide, G.MARGIN_L, 0.14, 4.0, 0.16, KICKER, S["foot"], F["head"], C["muted"], bold=True
    )
    add_text(
        slide,
        G.MARGIN_L,
        0.32,
        FULL_W,
        0.28,
        HEADLINE,
        S["sub"],
        F["head"],
        C["primary"],
        bold=True,
    )
    _hrule(slide, G.MARGIN_L, RULE_Y, FULL_W, color=C["primary"], weight=1.4)

    # ══ 플로우 시각화 — 골든 s06(s06_variants.variant_c) 그대로 이식 ══════════
    #    좌표만 압축 헤더 밑으로 당김(LANE_Y). 노드명·태그·도형·색·폭은 골든 원본.
    mid = LANE_Y + FNODE_H / 2
    widths = [1.0, 1.45, 1.15, 1.45, 1.45, 1.45, 1.0]
    gap = (G.RIGHT_EDGE - G.MARGIN_L - sum(widths)) / 6
    xs, x = [], G.MARGIN_L
    for w in widths:
        xs.append(x)
        x += w + gap
    node_style = [
        ("round", C["bg"], C["muted"], C["primary"]),
        ("round", C["bg_alt"], None, C["primary"]),
        ("diamond", C["bg"], C["accent"], C["primary"]),
        ("round", C["bg_alt"], None, C["primary"]),
        ("round", C["bg_alt"], None, C["primary"]),
        ("round", C["bg_alt"], None, C["primary"]),
        ("round", C["primary"], None, C["bg"]),
    ]
    for (kind, fill, line, ink), name, nx, w in zip(node_style, NODE_NAMES, xs, widths):
        h = FNODE_H + 0.25 if kind == "diamond" else FNODE_H
        y = mid - h / 2
        shp = add_box(
            slide, nx, y, w, h, fill=fill, line=line, line_w=1.25 if line else None, shape=kind
        )
        set_shape_text(
            shp, name, S["caption" if kind == "diamond" else "body"], F["head"], ink, bold=True
        )
    for i in range(6):
        _arrow(slide, xs[i] + widths[i], mid, xs[i + 1], mid, C["muted"], w=1.5)
    for i, tag in zip([1, 3, 4, 5], NODE_TAGS):
        add_text(
            slide,
            xs[i] - 0.35,
            LANE_Y + FNODE_H + 0.1,
            widths[i] + 0.7,
            0.24,
            tag,
            S["caption"],
            F["head"],
            C["primary"],
            bold=True,
            align=PP_ALIGN.CENTER,
        )
    add_text(
        slide,
        xs[2] + widths[2] + 0.03,
        mid - 0.32,
        0.5,
        0.22,
        IN_LABEL,
        S["caption"],
        F["head"],
        C["accent"],
        bold=True,
    )
    # ── OUT 분기 — 다이아 아래 수직 낙하, 거절 박스 + 우측 인라인 설명 ──
    dia_cx = xs[2] + widths[2] / 2
    dia_bot = mid + (FNODE_H + 0.25) / 2
    rej_w, rej_h = 2.0, 0.45
    rej_x, rej_y = dia_cx - rej_w / 2, dia_bot + 0.48
    rej = add_box(slide, rej_x, rej_y, rej_w, rej_h, fill=C["bg_alt"], shape="round")
    set_shape_text(rej, REJECT_BOX, S["caption"], F["head"], C["primary"], bold=True)
    _arrow(slide, dia_cx, dia_bot, dia_cx, rej_y, C["muted"], w=1.5)
    add_text(
        slide,
        dia_cx + 0.12,
        dia_bot + 0.12,
        0.6,
        0.22,
        OUT_LABEL,
        S["caption"],
        F["head"],
        C["accent"],
        bold=True,
    )
    add_text(
        slide,
        rej_x + rej_w + 0.2,
        rej_y + 0.1,
        4.2,
        0.26,
        REJECT_DESC,
        S["caption"],
        F["body"],
        C["muted"],
    )

    # ══ 동형 카드 4 — 히어로 아이콘판 (배지·제목→태그→배너→히어로→불릿3→칩3) ══
    card_w = (FULL_W - 3 * 0.18) / 4
    for i, (num, name, tag, icon, banner, procs, chips) in enumerate(NODE_CARDS):
        cx = G.MARGIN_L + i * (card_w + 0.18)
        cy = CARD_Y
        add_box(slide, cx, cy, card_w, CARD_H, fill=C["bg"], line=line_soft, line_w=1.0)
        # 헤더 행: 배지 · 제목
        badge = add_box(slide, cx + 0.14, cy + 0.13, 0.28, 0.28, fill=C["primary"], shape="oval")
        set_shape_text(badge, num, S["foot"], F["head"], C["bg"], bold=True)
        add_text(
            slide,
            cx + 0.52,
            cy + 0.11,
            card_w - 0.66,
            0.32,
            name,
            S["head"],
            F["head"],
            C["primary"],
            bold=True,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        # 태그 (accent)
        add_text(
            slide,
            cx + 0.14,
            cy + 0.50,
            card_w - 0.28,
            0.20,
            tag,
            S["caption"],
            F["head"],
            C["accent"],
            bold=True,
        )
        # 라이트 배너
        ban = add_box(slide, cx + 0.14, cy + 0.76, card_w - 0.28, 0.30, fill=C["bg_alt"])
        set_shape_text(ban, banner, S["caption"], F["head"], C["primary"], bold=True)
        ban.text_frame.word_wrap = False
        # 히어로 아이콘 — accent 심볼 (중앙 정렬, 불릿 밀도 확보 위해 축소)
        _icon(slide, icon, "accent", cx + (card_w - HERO_ICON) / 2, cy + 1.08, HERO_ICON)
        # 불릿 4 — (볼드 리드 + 설명) 2단, 촘촘히
        for j, (lead, detail) in enumerate(procs):
            add_text(
                slide,
                cx + 0.15,
                cy + 1.64 + j * 0.29,
                card_w - 0.26,
                0.28,
                [
                    [
                        ("· ", {"bold": True, "color": C["accent"]}),
                        (lead, {"bold": True, "color": C["primary"]}),
                        (detail, {}),
                    ]
                ],
                S["caption"],
                F["body"],
                ink_body,
                anchor=MSO_ANCHOR.MIDDLE,
                line_spacing=1.0,
            )
        # 구분 룰 + 아이콘 칩 3
        _hrule(slide, cx + 0.14, cy + 2.90, card_w - 0.28)
        for k, (cicon, em, label) in enumerate(chips):
            _chip_row(
                slide, cx + 0.14, cy + 2.99 + k * 0.24, card_w - 0.28, cicon, em, label, h=0.23
            )

    add_text(
        slide, G.MARGIN_L, SIDE_NOTE_Y, FULL_W, 0.16, SIDE_NOTE, S["foot"], F["body"], C["muted"]
    )
    add_text(slide, G.MARGIN_L, SOURCE_Y, FULL_W, 0.2, SOURCE, S["foot"], F["body"], C["muted"])
    return slide
