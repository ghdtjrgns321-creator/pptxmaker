"""S11 판단트리 — 고밀도 개편본 (s06 공식, 2026-07-20). 골든 편입 대상.

원본 s11_variants.variant_d(28pt 헤더 + 3서사 + 3단계 다중주입·세갈래 분기(좌) + 적용
판단트리(우) + 검은 바)를 S8·S9 승리 공식으로 재배치:
- 압축 헤더(dense.compact_header)
- 상단 시그니처 밴드: 적용 실물 판단트리(주인공, 좌) + 주제기반 3단계 다중주입(우)
- 하단: s06 표준 hero_card 4장(흩어진 조건 / 답변 세 갈래 / 원문 앵커 41 / 순서 창작 안 함)
- 검은 전폭 바 폐지 → 결론은 카드에 흡수 (하단 한줄평 바 금지, design-rules §8)

콘텐츠는 s11_variants.DEFAULT·STEPS 원문 그대로. 도해는 top-band 좌표로 새로 그린다.
"""

from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt

from . import dense as D
from . import grid as G
from .kit import SLIDE_H, SLIDE_W, add_box, add_text, load_kit, set_shape_text
from .figures import Box
from .figures import n_branch as NBRANCH
from .figures import numbered_steps as STEPS

K = load_kit()

KICKER = "3. 기술 설명 — TECH 03 · Generate"

SOURCE = (
    "출처: data/ontology/judgment_trees.json(트리 41 · 사용자 검수 완료) · "
    "3_KNOWLEDGE-GRAPH.md §3.3 · 4_SEARCH-PIPELINE.md L72·78"
)

NARRATIVES = [  # §8d 서사 3칼럼 — 두괄식 볼드 헤드 + 완전 문장
    (
        "왜 필요한가",
        "하나의 회계 판단이 문단 32·35·36·37과 B9~B13에 흩어져 있다. "
        "문단을 조각으로 검색해서는 판단의 순서가 복원되지 않는다.",
    ),
    (
        "파이프라인에서의 역할",
        "Generate 직전에 주입되는 답변의 골격. 질문에 걸린 트리가 "
        "판단 순서를 제공하므로 AI가 순서를 창작하지 않는다.",
    ),
    (
        "어떻게 만들었나",
        "기준서 본문의 조건-분기만 원문 앵커(문단 번호)와 함께 옮겨 "
        "41개를 미리 조립 — AI 창작 아님, 전 트리 사용자 검수 완료.",
    ),
]

# content_contract 계약 키 — golden.tech_mechanism의 필수 키 출처.
# 2026-07-26 `s11_variants`(sparse 시안)에서 이관 + 부품화로 나온 글(steps·branches)을 추가.
DEFAULT = {  # 텍스트 내용만(좌표·색·크기·도형은 코드 고정). c=None이면 골든 기본값.
    "kicker": KICKER,
    "headline": "판단트리 41개 — 흩어진 조건-분기를 판단 순서로 미리 조립",
    "narratives": NARRATIVES,
    "steps_head": "트리는 이렇게 걸린다 — 주제기반 다중주입",
    "branch_label": "주입 뒤 — 답변은 질문 성격에 따라 세 갈래로 갈린다",
    "flow_head": "적용 실물 — 질문에서 트리 분기까지",
    "question": "“볼륨디스카운트 조항이 있는 계약은…”",
    "topic_label": "주제 지목 — 변동대가",
    "tree_chip": "「변동대가 (추정)」 걸림",
    "diamond": "변동금액 포함?",
    "diamond_anchor": "문단 51",
    "method1_head": "기댓값",
    "method1_sub": "확률 가중 합",
    "method2_head": "가능성 최고 금액",
    "method2_sub": "단일 결과치",
    "flow_foot": "문단 53 — 더 잘 예측하는 방법 하나를 일관 적용(문단 54)",
    "bar": "질문이 여러 판단에 걸치면 걸린 절차를 모두 넣는다 — 트리 41개, 전부 원문 앵커",
    "source": SOURCE,
    "steps": [
        {"head": "주제 지목", "body": "35개 주제에서 최대 3개 지목 — 개념 80개를 직접 고르지 않는다"},
        {"head": "직속 개념 매칭", "body": "직속 개념이 트리 발동 개념과 겹치면 걸린다 — 배경 개념 제외"},
        {"head": "걸린 트리 전부 주입", "body": "1개만 고르지 않고 걸린 절차 모두 주입 — 없으면 진입 개념 폴백"},
    ],
    "branch_source": "주입된 질문",
    "branches": [
        {"label": "사실관계 부족", "desc": " 되물음 — 부족한 사실을 묻는다"},
        {"label": "산술 필요", "desc": " 계산을 확인한 뒤 답변"},
        {"label": "그 외", "desc": " 확정 또는 조건부 답변"},
    ],
}
C, S, F = K["rgb"], K["sizes"], K["fonts"]

# 카드 주제 = 판단트리의 고유 성질 (제네릭 왜/역할/구축 대신)
CARDS = [
    (
        "1",
        "흩어진 조건 조립",
        "판단이 문단에 흩어짐",
        "layers",
        "조각 검색은 순서 복원 불가",
        [
            ("한 회계 판단", " 문단 32·35·36·37·B9~13"),
            ("조각 검색", " 판단 순서 안 복원"),
            ("트리", " 순서를 미리 조립"),
        ],
        [
            ("layers", "41", " 트리"),
            ("git-branch", "조건-분기", ""),
            ("circle-check-big", "순서", " 복원"),
        ],
    ),
    (
        "2",
        "걸린 절차 전부 주입",
        "1개만 고르지 않음",
        "git-branch",
        "여러 판단 걸치면 전부",
        [
            ("1개만 안 고름", " 걸린 절차 모두 주입"),
            ("직속 개념 없으면", " 진입 개념 전체 폴백"),
            ("전부 원문 앵커", " 근거 추적 가능"),
        ],
        [
            ("git-branch", "전부", " 주입"),
            ("circle-check-big", "원문 앵커", ""),
            ("layers", "41", " 트리"),
        ],
    ),
    (
        "3",
        "원문 앵커 41",
        "AI 창작 아님",
        "file-text",
        "조건-분기만 옮겨 41개",
        [
            ("조건-분기만", " 원문 앵커=문단번호"),
            ("41개", " 미리 조립"),
            ("전 트리", " 사용자 검수 완료"),
        ],
        [
            ("file-text", "41", " 트리"),
            ("circle-x", "AI 창작", " 0"),
            ("circle-check-big", "전수", " 검수"),
        ],
    ),
    (
        "4",
        "순서 창작 안 함",
        "Generate 직전 골격",
        "cpu",
        "트리가 판단 순서 제공",
        [
            ("Generate 직전", " 답변 골격 주입"),
            ("걸린 트리", " 판단 순서 제공"),
            ("AI", " 순서 창작 없이 채우기"),
        ],
        [
            ("git-branch", "골격", " 주입"),
            ("circle-check-big", "순서", " 트리"),
            ("circle-x", "AI 순서", " 0"),
        ],
    ),
]


def _subhead(slide, x, y, w, text):
    add_text(slide, x, y, w, 0.28, text, S["caption"], F["head"], C["primary"], bold=True)
    D.hrule(slide, x, y + 0.30, w, color=D.line_mid, weight=1.0)


def _decision_tree(slide, c, x0, x1):
    """적용 실물 판단트리 — 질문 → 트리 걸림 → 다이아 분기 → 두 방법 (세로)."""
    cx = (x0 + x1) / 2
    tw = x1 - x0
    _subhead(slide, x0, 0.82, tw, c["flow_head"])
    q = add_box(
        slide,
        x0 + 0.3,
        1.28,
        tw - 0.6,
        0.36,
        fill=C["bg"],
        line=C["muted"],
        line_w=1.0,
        shape="round",
    )
    set_shape_text(q, c["question"], S["foot"], F["body"], C["primary"])
    D.arrow(slide, cx, 1.64, cx, 1.86, C["muted"])
    add_text(
        slide,
        cx + 0.12,
        1.63,
        tw / 2 - 0.12,
        0.18,
        c["topic_label"],
        S["foot"],
        F["body"],
        C["muted"],
    )
    chip = add_box(slide, x0 + 0.8, 1.88, tw - 1.6, 0.34, fill=C["primary"], shape="round")
    set_shape_text(chip, c["tree_chip"], S["foot"], F["head"], C["bg"], bold=True)
    D.arrow(slide, cx, 2.22, cx, 2.40, C["muted"])
    dia_w, dia_h = 1.5, 0.58
    dia = add_box(
        slide,
        cx - dia_w / 2,
        2.40,
        dia_w,
        dia_h,
        fill=C["bg"],
        line=C["accent"],
        line_w=1.75,
        shape="diamond",
    )
    set_shape_text(dia, c["diamond"], S["foot"], F["head"], C["primary"], bold=True)
    add_text(
        slide,
        cx + dia_w / 2 + 0.06,
        2.58,
        1.0,
        0.2,
        c["diamond_anchor"],
        S["foot"],
        F["body"],
        C["muted"],
    )
    my = 3.02
    D.arrow(slide, cx - 0.32, 2.98, x0 + 1.0, my, C["muted"])
    D.arrow(slide, cx + 0.32, 2.98, x1 - 1.0, my, C["muted"])
    for mx, head, sub in (
        (x0 + 0.08, c["method1_head"], c["method1_sub"]),
        (x1 - 1.83, c["method2_head"], c["method2_sub"]),
    ):
        b = add_box(
            slide, mx, my, 1.75, 0.44, fill=C["bg"], line=C["muted"], line_w=1.0, shape="round"
        )
        tf = b.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p._p.get_or_add_pPr().set("eaLnBrk", "0")
        for txt, bold, color, sz in (
            (head, True, C["primary"], S["foot"]),
            ("  " + sub, False, C["muted"], S["foot"] - 1),
        ):
            r = p.add_run()
            r.text = txt
            r.font.name, r.font.bold, r.font.size, r.font.color.rgb = F["head"], bold, Pt(sz), color
            r.font.language_id = MSO_LANGUAGE_ID.KOREAN


# 3단계 — 번호 배지 + 텍스트 흐름(박스 없음 → 빈 여백 자체가 안 생긴다). 문구는 원본 사실 보존.
# ── 주입 흐름·갈래 자리 (골든 실측) ────────────────────────────────────────────
# 도해는 `figures.numbered_steps`·`figures.n_branch`가 그린다. 이 장은 자리만 정한다.
# 글(단계·갈래 문구)은 2026-07-26 부품화로 content로 나갔다 — 모듈 상수로 두면 실전 덱에
# 골든 글이 그대로 새 나간다(계약 ③).
STEPS_DY = 1.28  # 소제목 밑 첫 배지 top
BR_BAND_TOP, BR_BAND_BOTTOM = 2.46, 3.46  # 갈래가 쓸 수 있는 세로 밴드


def _inject_and_branch(slide, c, x0, x1):
    """주제기반 다중주입(번호 흐름) + 주입 뒤 N갈래 분기 — 둘 다 개수 파생 부품이다."""
    tw = x1 - x0
    _subhead(slide, x0, 0.82, tw, c["steps_head"])
    STEPS.draw(slide, Box(x0, STEPS_DY, tw, 1.00), {"steps": c["steps"]}, K)
    add_text(
        slide, x0, 2.26, tw, 0.2, c["branch_label"], S["foot"], F["head"], C["primary"], bold=True
    )
    NBRANCH.draw(
        slide,
        Box(x0, BR_BAND_TOP, tw, BR_BAND_BOTTOM - BR_BAND_TOP),
        {"source": c["branch_source"], "branches": c["branches"]},
        K,
    )


def build(prs, c=None):
    c = {**DEFAULT, **(c or {})}
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    D.compact_header(slide, c["kicker"], c["headline"])

    # ── 상단 시그니처 밴드: 적용 판단트리(좌) + 3단계 다중주입(우) ──
    _decision_tree(slide, c, G.MARGIN_L, 6.0)
    add_box(slide, 6.3, 0.86, 0.012, 2.55, fill=C["bg_alt"])  # 세로 구분선
    _inject_and_branch(slide, c, 6.6, G.RIGHT_EDGE)

    # ── 하단: s06 표준 hero_card 4장 ──
    card_w = (D.FULL_W - 3 * 0.18) / 4
    for i, (num, title, tag, ic, banner, items, chips) in enumerate(CARDS):
        cx = G.MARGIN_L + i * (card_w + 0.18)
        D.hero_card(
            slide, cx, 3.62, card_w, 3.32, num, title, tag, ic, banner, items, chips, hero=0.44
        )
    # 결론("걸린 절차를 모두 주입 · 전부 원문 앵커")은 카드1·3에 흡수 — 하단 한줄평 바 금지(§8).
    D.source_line(slide, c["source"])
    return slide
