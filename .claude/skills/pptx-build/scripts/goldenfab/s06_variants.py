"""S6 실행그래프 — 골든 정본 (00_factsheet.md §C).

구성: 헤더 + 실행 레인(부품) + 노드별 핵심 기술 4칼럼 + 결론 바.
도해는 `figures.routing_lane`이 그린다 — 이 파일은 자리와 노드 종류만 정한다.

**시안 2종을 2026-07-26에 지웠다.** 파일명이 `_variants`인 이유가 그것이었는데, 레지스트리는
`variant_c` 하나만 물고 있었고 A·B는 이 파일 `main()`만 부르는 죽은 코드였다(약 290줄).
부품화로 레인이 밖으로 나가면서 두 시안이 공유하던 헬퍼(`macro_band`·`micro_flow`·`header`·
`_arrow`)도 함께 갈 곳을 잃었다. 골든 확인은 `uv run python golden/build_golden.py`.
"""

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from . import grid as G
from .figures import Box
from .figures import routing_lane as ROUTE
from .kit import SLIDE_H, SLIDE_W, add_box, add_text, load_kit

K = load_kit()
C, S, F = K["rgb"], K["sizes"], K["fonts"]


def _subhead(slide, x, y, w, text):
    add_text(slide, x, y, w, 0.32, text, S["head"], F["head"], C["primary"], bold=True)
    add_box(slide, x, y + 0.35, w, 0.012, fill=C["muted"])


# ── 실행 레인 배치 (2026-07-26 부품화) ────────────────────────────────────────────
# 이 장이 정하는 것은 **자리와 노드 종류**뿐이고 그림은 `figures.routing_lane`이 그린다.
# 여기 있던 `_arrow`는 골든 표준 화살표와 완전히 같은 사본이었다(같은 화살표가 14벌 흩어져
# 있었다는 실측의 실물 하나) — `elements.arrow` 단일 출처로 흡수하고 지웠다.
LANE_BOX = Box(G.MARGIN_L, 2.225, G.RIGHT_EDGE - G.MARGIN_L, 1.775)
LANE_KINDS = ["start", "process", "gate", "process", "process", "process", "end"]
_TAGGED = (1, 3, 4, 5)  # 태그가 붙는 노드 자리 — 단자·판정에는 안 붙는다


def _lane_tags(tags):
    """노드 자리별 태그(없으면 None). 태그 목록은 태그 붙는 노드 수와 같아야 한다."""
    if len(tags) != len(_TAGGED):
        raise ValueError(f"s06: node_tags {len(tags)}개 — 태그 노드는 {len(_TAGGED)}자리다")
    at = dict(zip(_TAGGED, tags, strict=True))
    return [at.get(i) for i in range(len(LANE_KINDS))]


DEFAULT_C = {
    "kicker": "2. 파이프라인",
    "headline": "질문에서 응답까지 — 분기까지 전부 결정적인 실행 그래프",
    "subtitle": "진입부 임베딩 유사도 0 — LLM이 고르는 것은 35토픽 목록뿐, 경로는 그래프가 결정한다.",
    "node_names": ["질문", "Analyze", "판정", "Retrieve", "Generate", "Format", "응답"],
    "node_tags": ["용어사전 매칭", "그래프 1홉 탐색", "판단트리 주입", "경고·꼬리질문"],
    "in_label": "IN",
    "reject_box": "거절 메시지",
    "out_label": "OUT",
    "reject_desc": "범위 밖 질문은 본선에 진입하지 못하고 즉시 거절된다",
    "detail_head": "핵심 기술 — 결정성을 어떻게 만들었나",
    "details": [
        (
            "01  용어사전 — 후보 진입점",
            "등재 423 · AI 신규 창작 0.",
            " 사람이 만든 자료 3종(질의 매핑 288 · 사례 제목 123 · 부록A 정의 9)에서 AI 초안 + 사람 전수 검수로 등재. 확정 라우터가 아니라 후보 진입점이다.",
        ),
        (
            "02  지식그래프 — 기계 생성 뼈대",
            "노드 929 · 간선 2,697.",
            " 기준서 공식 소제목이 그대로 개념 80, 문단 250 배정. 계층·참조 간선은 기준서에서 기계 생성 — AI가 만드는 것은 용어 색인 하나뿐이다.",
        ),
        (
            "03  판단트리 — 판단 순서의 사전 조립",
            "조건-분기 골격 41개, 원문 앵커 포함.",
            " 기준서에 흩어진 판단 절차를 미리 조립해 두고, 진입 개념과 트리거 개념의 최다 겹침 + 주제 직속 트리를 전부 프롬프트에 주입한다.",
        ),
        (
            "04  구조화 출력 — 코드 수준 강제",
            "PydanticAI 스키마 강제.",
            " 답변 형식을 코드로 강제하고, 위반하면 result_validator가 자동 재시도한다. 감리 경고·꼬리질문은 마지막에 부가되는 보조 장치다.",
        ),
    ],
    "rerank_note": "초기 설계에 있던 유사도 재정렬(rerank) 단계는 제거 — 그래프가 이미 결정적으로 선별하므로 재정렬할 것이 없다",
    "bar": "질문에서 답까지 4개 노드 전부가 결정적으로 동작한다",
    "source": "출처: 4_SEARCH-PIPELINE.md (00_factsheet.md §C)",
}


def variant_c(prs, c=None):
    """C — 원문(4_SEARCH-PIPELINE.md L8~19) 아스키 구조의 도형화: 분기 포함 실행 그래프.

    c: 텍스트 내용 override(None=골든 기본값). 좌표·색·도형 종류는 고정."""
    c = {**DEFAULT_C, **(c or {})}
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    # 헤더 — header() 인라인(kicker·headline만 c에서, 출력은 header 호출과 동일)
    add_text(
        slide,
        G.MARGIN_L,
        0.42,
        6.0,
        0.28,
        c["kicker"],
        S["caption"],
        F["head"],
        C["muted"],
        bold=True,
    )
    add_text(
        slide,
        G.MARGIN_L,
        0.72,
        G.RIGHT_EDGE - G.MARGIN_L,
        0.55,
        c["headline"],
        S["section"],
        F["head"],
        C["primary"],
        bold=True,
    )
    add_box(slide, G.MARGIN_L, G.RULE_Y, G.RIGHT_EDGE - G.MARGIN_L, 0.014, fill=C["muted"])
    add_text(
        slide,
        G.MARGIN_L,
        G.CONTENT_TOP,
        G.RIGHT_EDGE - G.MARGIN_L,
        0.3,
        c["subtitle"],
        S["sub"],
        F["body"],
        C["text"],
    )
    # ── 실행 레인 — 부품 호출(figures.routing_lane) ──────────────────────────────
    # 2026-07-26 부품화: 레인·판정·거절 분기를 그리던 코드가 여기 있었다. 자리(LANE_BOX)만
    # 이 장이 정하고 그림은 부품이 그린다 — 좌표는 골든 실측 그대로여서 픽셀 동일이다.
    # `LANE_KINDS`가 노드 종류를 선언한다(폭·모양이 여기서 나온다). `strict=True`인 이유:
    # content가 노드 수가 다른 이름 목록을 주입하면 조용히 잘리지 않고 시끄럽게 죽어야 한다.
    ROUTE.draw(
        slide,
        LANE_BOX,
        {
            "nodes": [
                {"name": name, "kind": kind, **({"tag": tag} if tag else {})}
                for name, kind, tag in zip(
                    c["node_names"], LANE_KINDS, _lane_tags(c["node_tags"]), strict=True
                )
            ],
            "in_label": c["in_label"],
            "branch": {
                "box": c["reject_box"],
                "label": c["out_label"],
                "desc": c["reject_desc"],
            },
        },
        K,
    )
    # ── 하단: 노드별 핵심 기술 — 평어 상세 4칼럼 (원문 §C·L47·L53·L63~65) ──
    _subhead(slide, G.MARGIN_L, 4.35, G.RIGHT_EDGE - G.MARGIN_L, c["detail_head"])
    det_w = (G.RIGHT_EDGE - G.MARGIN_L - 3 * 0.35) / 4  # 2.771
    for i, (head, lead, body) in enumerate(c["details"]):
        dx = G.MARGIN_L + i * (det_w + 0.35)
        add_text(
            slide, dx, 4.85, det_w, 0.28, head, S["caption"], F["head"], C["primary"], bold=True
        )
        add_text(
            slide,
            dx,
            5.18,
            det_w,
            1.05,
            [[(lead, {"bold": True, "color": C["primary"]}), (body, {})]],
            S["caption"],
            F["body"],
            C["muted"],
            line_spacing=1.25,
        )
        if i < 3:
            add_box(slide, dx + det_w + 0.175, 4.85, 0.012, 1.05, fill=C["bg_alt"])
    add_text(
        slide,
        G.MARGIN_L,
        6.05,
        G.RIGHT_EDGE - G.MARGIN_L,
        0.3,
        c["rerank_note"],
        S["caption"],
        F["body"],
        C["muted"],
    )
    bar = add_box(slide, G.MARGIN_L, G.BAR_Y, G.RIGHT_EDGE - G.MARGIN_L, G.BAR_H, fill=C["primary"])
    tf = bar.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = c["bar"]
    r.font.name, r.font.bold = F["head"], True
    r.font.size = Pt(S["body"])
    r.font.color.rgb = C["bg"]
    add_text(
        slide,
        G.MARGIN_L,
        G.SOURCE_Y,
        G.RIGHT_EDGE - G.MARGIN_L,
        0.3,
        c["source"],
        S["foot"],
        F["body"],
        C["muted"],
    )
    return slide


def audit(prs):
    """오딧 — bottom>6.35 / 밴드·풀폭 right / 룰-채움 겹침 / accent ≤4."""
    EMU = 914400
    accent_hex = str(C["accent"])
    fails = []
    for si, sl in enumerate(prs.slides):
        rules, solids, acc = [], [], 0
        for sh in sl.shapes:
            L, T, W, H = sh.left / EMU, sh.top / EMU, sh.width / EMU, sh.height / EMU
            if (L == 0 and T == 0 and W > 13) or (W < 4 and H > 6):  # 배경·사이드바 제외
                continue
            if H <= 0.02 and W > 1:
                rules.append((L, T, W))
            elif H > 0.3 and W > 0.5 and sh.shape_type != 17:
                solids.append((L, T, W, H))
            if T + H > G.CONTENT_BOTTOM + 0.02 and T < G.BAR_Y - 0.02:
                fails.append((si, "bottom>6.35", round(T + H, 3)))
            if W > 11.5 and W < 13 and abs((L + W) - G.RIGHT_EDGE) > 0.02:
                fails.append((si, "풀폭 right≠12.733", round(L + W, 3)))
            try:
                if str(sh.fill.fore_color.rgb) == accent_hex:
                    acc += 1
            except Exception:
                pass
            if sh.has_text_frame:
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        if (
                            r.font.color
                            and r.font.color.rgb is not None
                            and str(r.font.color.rgb) == accent_hex
                        ):
                            acc += 1
        for rl, rt, rw in rules:
            for sl_, st, sw, shh in solids:
                if rl < sl_ + sw and rl + rw > sl_ and st - 0.005 <= rt <= st + shh - 0.005:
                    fails.append((si, "룰-채움 겹침", round(rt, 3)))
        if acc > 4:
            fails.append((si, "accent 초과", acc))
        print(f"slide{si}: accent {acc}곳")
    assert not fails, f"AUDIT FAIL {fails}"
    print("audit pass")


def main():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    variant_c(prs)
    audit(prs)
    from pathlib import Path

    out = Path(__file__).parent / "variants" / "s06_variants.pptx"
    prs.save(str(out))
    print(f"saved {out}")


if __name__ == "__main__":
    main()
