"""S8 시안 2종 — 기술 1: 용어사전 (구축→적용→실증 3구획 휴리스틱, UI 패스 — 문구 러프 허용).

재료: 00_factsheet.md §C 구축·작동 상세 + §D (원문 2_DATA-TAXONOMY §2.6).
시안 A(좌앵커형): 좌측 앵커(대형 423) + 우측 3구획 세로 스택(구축 플로우·적용 예시·분해 바).
시안 B(풀폭 플로우형): 상단 구축 셰브런 밴드 + 중단 적용 예시·다중 목적지 + 하단 분해 바.
실행: 이 타입은 goldenfab 레지스트리 경유로만 렌더된다 — 골든 19장 확인은
      `uv run python golden/build_golden.py`(2026-07-15 단일화로 시안 개별 실행 경로 폐지).
"""

from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from . import grid as G
from ._variant_h import _shape_step
from .kit import SLIDE_H, SLIDE_W, add_box, add_text, load_kit, mix, set_shape_text

K = load_kit()
C, S, F = K["rgb"], K["sizes"], K["fonts"]

KICKER = "3. 기술 설명 — TECH 01 · Analyze"
SOURCE = "출처: 2_DATA-TAXONOMY.md §2.6 · 4_SEARCH-PIPELINE.md (00_factsheet.md §C·§D)"
OVERLAP = G.FLOW_H * 0.55 / 2

BREAKDOWN = [  # 등재 423 분해 (합 423 검산은 audit)
    ("자동", 316, C["primary"]),
    ("위임판단", 86, mix(C["primary"], C["muted"], 0.5)),
    ("사용자검토", 18, C["muted"]),
    ("사용자확정", 1, C["accent"]),
    ("제외", 2, mix(C["muted"], C["bg"], 0.5)),
]


def header(slide, headline, kicker=KICKER):
    add_text(
        slide, G.MARGIN_L, 0.42, 8.0, 0.28, kicker, S["caption"], F["head"], C["muted"], bold=True
    )
    add_text(
        slide,
        G.MARGIN_L,
        0.72,
        G.RIGHT_EDGE - G.MARGIN_L,
        0.55,
        headline,
        S["section"],
        F["head"],
        C["primary"],
        bold=True,
    )
    add_box(slide, G.MARGIN_L, G.RULE_Y, G.RIGHT_EDGE - G.MARGIN_L, 0.014, fill=C["muted"])


def _subhead(slide, x, y, w, text):
    add_text(slide, x, y, w, 0.32, text, S["head"], F["head"], C["primary"], bold=True)
    add_box(slide, x, y + 0.35, w, 0.012, fill=C["muted"])


def _arrow(slide, x1, y1, x2, y2):
    from pptx.oxml.ns import qn

    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = C["muted"]
    conn.line.width = Pt(1.5)
    ln = conn.line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))


def bar_and_source(slide, text, right=None, source=SOURCE):
    r = right or G.RIGHT_EDGE
    bar = add_box(slide, G.MARGIN_L, G.BAR_Y, r - G.MARGIN_L, G.BAR_H, fill=C["primary"])
    tf = bar.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name, run.font.bold = F["head"], True
    run.font.size = Pt(S["body"])
    run.font.color.rgb = C["bg"]
    add_text(
        slide, G.MARGIN_L, G.SOURCE_Y, r - G.MARGIN_L, 0.3, source, S["foot"], F["body"], C["muted"]
    )


def breakdown_bar(slide, x, y, w, h):
    """등재 423 비례 스택 바 + 범례. 미세 구간은 최소 가시폭(주석으로 정직 표기)."""
    total = sum(n for _, n, _ in BREAKDOWN)
    min_w = 0.06
    seg_x = x
    for name, n, color in BREAKDOWN:
        seg_w = max(w * n / total, min_w)
        add_box(slide, seg_x, y, seg_w, h, fill=color)
        seg_x += seg_w
    leg_w = w / len(BREAKDOWN)
    for i, (name, n, color) in enumerate(BREAKDOWN):
        leg_x = x + i * leg_w
        add_box(slide, leg_x, y + h + 0.1, 0.13, 0.13, fill=color)
        add_text(
            slide,
            leg_x + 0.19,
            y + h + 0.06,
            leg_w - 0.2,
            0.22,
            f"{name} {n}",
            S["caption"],
            F["body"],
            C["muted"],
        )


def _apply_example(slide, x, y):
    """적용 예시 — 질문 속 문자가 개념 후보로 이어지는 그림."""
    q = add_box(slide, x, y, 3.4, 0.6, fill=C["bg"], line=C["muted"], line_w=1.0, shape="round")
    tf = q.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r1 = p.add_run()
    r1.text = "“리베이트 조항이 있는 계약은…”"
    r1.font.name, r1.font.size, r1.font.color.rgb = F["body"], Pt(S["caption"]), C["primary"]
    _arrow(slide, x + 3.4, y + 0.3, x + 4.0, y + 0.3)
    chip = add_box(slide, x + 4.0, y + 0.12, 1.5, 0.36, fill=C["primary"], shape="round")
    set_shape_text(chip, "변동대가", S["caption"], F["head"], C["bg"], bold=True)
    add_text(
        slide,
        x + 5.7,
        y + 0.16,
        2.6,
        0.3,
        "후보 진입점 — 확정 라우터 아님",
        S["caption"],
        F["body"],
        C["muted"],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        x,
        y + 0.75,
        8.2,
        0.26,
        "질문 문장에 등재 용어가 문자 그대로 나타나면 잡는다(substring) — 유사도 점수 없음. 최종 선택은 질문 문맥이 담당.",
        S["caption"],
        F["body"],
        C["muted"],
    )


def variant_a(prs):
    """A — 좌앵커(대형 423) + 우측 3구획 스택(구축 → 적용 → 실증)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    header(slide, "용어사전 — 사람이 검수한 진입 색인, AI 신규 창작은 0")
    lx, lw = G.COL_L_X, G.COL_L_W
    add_text(
        slide,
        lx,
        G.SUBHEAD_Y,
        lw,
        0.26,
        "진입 색인",
        S["caption"],
        F["head"],
        C["muted"],
        bold=True,
    )
    add_text(slide, lx, 2.08, lw, 0.7, "423", S["display"], F["head"], C["accent"], bold=True)
    add_text(slide, lx, 2.78, lw, 0.48, "등재 용어", S["title"], F["head"], C["primary"], bold=True)
    add_text(
        slide,
        lx,
        3.42,
        lw,
        1.5,
        "원천은 사람이 만든 1차 자료 3종뿐이다. AI는 초안만 내고 사람이 전수 검수했다 — 색인이 틀리면 '못 찾음'으로 안전하게 드러난다.",
        S["sub"],
        F["body"],
        C["text"],
        line_spacing=1.3,
    )
    add_text(
        slide,
        lx,
        5.2,
        lw,
        1.0,
        "뼈대(개념·간선)는 기준서에서 기계 생성 — AI 개입을 색인 하나로 한정한 신빙성 배치.",
        S["caption"],
        F["body"],
        C["muted"],
        line_spacing=1.25,
    )
    add_box(
        slide, G.V_RULE_X, G.CONTENT_TOP, 0.014, G.CONTENT_BOTTOM - G.CONTENT_TOP, fill=C["bg_alt"]
    )

    rx, rw = G.COL_R_X, G.COL_R_W
    _subhead(slide, rx, G.SUBHEAD_Y, rw, "구축 — 3원천에서 사람 전수 검수로")
    srcs = [("질의 매핑 288", 0), ("사례 제목 123", 1), ("부록A 정의 9", 2)]
    for name, i in srcs:
        b = add_box(slide, rx, 2.28 + i * 0.42, 1.9, 0.34, fill=C["bg_alt"], shape="round")
        set_shape_text(b, name, S["caption"], F["body"], C["primary"])
    steps = [("AI 초안", 6.6), ("사람 전수 검수", 8.6)]
    for name, sx in steps:
        b = add_box(slide, sx, 2.7, 1.7, 0.5, fill=C["bg_alt"], shape="round")
        set_shape_text(b, name, S["caption"], F["head"], C["primary"], bold=True)
    fin = add_box(slide, 10.9, 2.7, 1.83, 0.5, fill=C["primary"], shape="round")
    set_shape_text(fin, "등재 423", S["body"], F["head"], C["bg"], bold=True)
    for i in range(3):
        _arrow(slide, rx + 1.9, 2.45 + i * 0.42, 6.6, 2.95)
    _arrow(slide, 8.3, 2.95, 8.6, 2.95)
    _arrow(slide, 10.3, 2.95, 10.9, 2.95)
    _subhead(slide, rx, 3.75, rw, "적용 — 질문 속 문자를 그대로 잡는다")
    _apply_example(slide, rx, 4.3)
    _subhead(slide, rx, G.BOX_Y + 0.05, rw, "실증 — 등재 423의 결정 분해")
    breakdown_bar(slide, rx, G.BOX_Y + 0.45, rw, 0.24)
    bar_and_source(
        slide, "AI가 만드는 것은 색인 하나 — 틀려도 1종(놓침)으로 드러나는 자리에만 둔다"
    )
    return slide


def variant_b(prs):
    """B — 상단 풀폭 구축 셰브런 밴드 + 중단 적용·다중 목적지 + 하단 분해 바."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    header(slide, "용어사전 — 3원천, 사람 전수 검수, 신규 창작 0")
    add_text(
        slide,
        G.MARGIN_L,
        G.CONTENT_TOP,
        G.RIGHT_EDGE - G.MARGIN_L,
        0.3,
        "질문이 그래프에 진입하는 첫 관문 — 사람이 만든 자료만으로 등재했다.",
        S["sub"],
        F["body"],
        C["text"],
    )
    band_y = 2.35
    step_w = (G.RIGHT_EDGE - G.MARGIN_L + 3 * OVERLAP) / 4
    band = [
        ("1차 자료 3종", "사람 원천"),
        ("AI 초안", "제안만"),
        ("사람 전수 검수", "애매 3건 판정"),
        ("등재 423", "창작 0"),
    ]
    for i, (head, desc) in enumerate(band):
        _shape_step(
            slide,
            G.MARGIN_L + i * (step_w - OVERLAP),
            band_y,
            step_w,
            G.FLOW_H,
            "pentagon" if i == 0 else "chevron",
            C["bg_alt"] if i < 3 else C["accent"],
            head,
            desc,
            C["primary"] if i < 3 else C["bg"],
            mix(C["primary"], C["bg"], 0.35) if i < 3 else C["bg_alt"],
        )
    caps = [
        "질의 매핑 288 · 사례 제목 123 · 부록A 정의 9",
        "후보 제안까지만 — 확정 권한 없음",
        "전건 결정 로그(누가·왜) 기록",
        "자동 316 · 위임 86 · 검토 18 · 확정 1 · 제외 2",
    ]
    for i, cap in enumerate(caps):
        add_text(
            slide,
            G.MARGIN_L + i * (step_w - OVERLAP) + 0.1,
            band_y + G.FLOW_H + 0.12,
            step_w - 0.5,
            0.55,
            cap,
            S["caption"],
            F["body"],
            C["muted"],
            line_spacing=1.2,
        )
    _subhead(slide, G.MARGIN_L, 4.15, 8.0, "적용 — substring 진입")
    _apply_example(slide, G.MARGIN_L, 4.7)
    _subhead(slide, 9.0, 4.15, G.RIGHT_EDGE - 9.0, "다중 목적지")
    term = add_box(slide, 9.0, 4.75, 1.4, 0.4, fill=C["bg_alt"], shape="round")
    set_shape_text(term, "“계약”", S["caption"], F["head"], C["primary"], bold=True)
    for i, concept in enumerate(["계약 식별", "계약변경", "계약원가"]):
        cy = 4.55 + i * 0.5
        _arrow(slide, 10.4, 4.95, 10.95, cy + 0.2)
        cb = add_box(
            slide, 10.95, cy, 1.75, 0.4, fill=C["bg"], line=C["muted"], line_w=0.75, shape="round"
        )
        set_shape_text(cb, concept, S["caption"], F["body"], C["primary"])
    add_text(
        slide,
        9.0,
        6.1,
        G.RIGHT_EDGE - 9.0,
        0.24,
        "한 용어 → 여러 개념 후보",
        S["caption"],
        F["body"],
        C["muted"],
    )
    bar_and_source(
        slide, "AI가 만드는 것은 색인 하나 — 틀려도 1종(놓침)으로 드러나는 자리에만 둔다"
    )
    return slide


# 콘텐츠 기본값(골든 내용) — c=None이면 이 값, override 시 텍스트만 교체(좌표·색·폰트·도형 고정)
DEFAULT_C = {
    "headline": "용어사전 — 실무 언어를 기준서 개념에 잇는 진입 색인",
    "kicker": KICKER,
    # 상단 3칼럼 서사 (head, body) — 번호 접두는 코드에서 위치로 생성
    "narratives": [
        (
            "왜 필요한가",
            "실무는 '리베이트'라 말하고 기준서는 '고객에게 지급할 대가'라 쓴다. 이 언어 간극을 잇지 않으면 검색이 시작조차 되지 않는다.",
        ),
        (
            "파이프라인에서의 역할",
            "Analyze 단계의 진입 관문. 질문에 등재 용어가 문자 그대로 나타나면 연결 개념으로 그래프에 진입한다 — 유사도 점수 없이, 사람이 승인한 경로로만.",
        ),
        (
            "어떻게 만들었나",
            "사람 1차 자료 3종(질의 매핑 288 · 사례 제목 123 · 부록A 정의 9)에서 AI가 초안을 내고 사람이 전수 검수해 423개 등재 — 신규 창작 0건.",
        ),
    ],
    # 소제목은 **짧게** — 15pt로 31자를 쓰면 5.2"라 범례(4.75 시작)를 덮어 글자가 뭉갠다
    # (2026-07-15 렌더 실측). 발췌 건수·등급 분해는 아래 캡션이 진다.
    "map_head": "실무가 쓰는 말 → 기준서가 쓰는 말",
    # 실물 엔트리 (data/ontology/aliases.json 실측 — 축약만, 창작 0).
    # (용어, 등급, [연결 개념]) — **우열 개념은 코드가 이 목록에서 파생**한다(등장순 dedup).
    # 그래서 N:1 수렴(`고객에게 지급할 대가` ← 리베이트+상품권)이 손으로 적는 게 아니라
    # 데이터에서 나온다. 등급 문자열도 aliases.json 원문 그대로("위임판단" 아님).
    "terms": [
        ("리베이트", "자동", ["고객에게 지급할 대가"]),
        ("밀어내기", "자동", ["위탁약정"]),
        ("볼륨디스카운트", "자동", ["변동대가", "변동대가 추정치를 제약함"]),
        ("상품권", "자동(위임판단)", ["고객에게 지급할 대가", "고객이 행사하지 아니한 권리"]),
        ("반품의 회계처리", "자동", ["반품권이 있는 판매", "본인 대 대리인의 고려사항"]),
    ],
    # 범례는 **형태 언어만**, 수치는 캡션만. 2026-07-15 제3자 채점 FAIL(#3): "자동 316"·
    # "자동(위임판단) 86"이 범례와 캡션에 글자 그대로 두 번 있었다.
    "map_legend": "실선 = 자동    점선 = 자동(위임판단) — AI가 판단, 사람이 승인",
    "map_caption": "등재 423 중 발췌 5건 — 423 = 자동 316 + 자동(위임판단) 86 + 검토 18 + 사용자 확정 1 + 제외 2",
    "json_title": "aliases.json — 실제 엔트리",
    "json_lines": [  # 상품권 엔트리 실물 축약
        '{ "term": "상품권",',
        '  "sources": ["query-mapping"],',
        '  "grade": "자동(위임판단)",',
        '  "concepts": [',
        '    "고객에게 지급할 대가",',
        '    "고객이 행사하지 아니한 권리" ],',
        '  "decision": {',
        '    "by": "AI 위임 판단",',
        '    "reason": "미행사 상품권 =',
        '        B44~47 정면 조항" } }',
    ],
    "json_caption": "모든 엔트리가 결정 로그(누가·왜)를 갖는다 — 전건 추적.",
    "bar": "AI가 만드는 것은 색인 하나 — 틀려도 1종(놓침)으로 드러나는 자리에만 둔다",
    "source": SOURCE,
}


# ── 이분 매핑 GRID (표를 걷어낸 자리 — 2026-07-15 사용자 반려 "이해가 안 가는 시각화") ──
#
# 물성 선언(design-rules P1): **잇기(N:M)**. 용어사전이 하는 일은 실무 언어를 기준서 언어에
#   잇는 것이다. 표는 그 관계를 4열 나열로 바꿔 쉼표 하나로 뭉갰다 — 한 용어가 여러 개념으로
#   갈라지는 것(1:N)도, 여러 용어가 한 개념으로 모이는 것(N:1)도 안 보였다. 그게 본질인데.
#   게다가 죽은 열이 절반이었다: `원천` 5행 중 4행이 같은 값, `등급` 4행이 "자동".
#   → 좌열(실무어) ─선─▶ 우열(기준서 개념). 갈래와 수렴이 **선의 모양 그 자체**로 보인다.
#   §6 "숫자·항목 단독 나열 금지 — 관계 인코딩과 결합" / §6 "나열은 도형에 가두기".
#
# 색(P4⑩): accent 예산은 **1**뿐이다 — 3칼럼 번호(01·02·03) 런이 이미 3을 쓴다(오딧 상한 4).
#   그 1을 `자동(위임판단)` 칩 테두리에 준다. 이 장의 결론 바가 "AI가 만드는 것은 색인 하나"이므로
#   **AI가 판단한 자리**가 accent를 가져가는 게 맞다. 선은 muted — 선까지 accent면 예산 초과다.
#   §5 점선 = 불확실 → 위임판단 매핑은 점선.
MAP_X = G.MARGIN_L  # 0.6
MAP_HEAD_Y = 3.35
MAP_RULE_DY = 0.35
MAP_L_RIGHT = 2.45  # 좌열 **우변 고정** — 선 출발점이 일정해야 매핑이 기둥으로 읽힌다
MAP_R_X = 4.75  # 우열 **좌변 고정** — 선 도착점 일정. 칩 폭은 글자에서 파생(아래)
MAP_R_MAX = 8.20  # 우열 우한 (JSON 카드 8.7 앞)
CHIP_INSET = 0.44  # 글자 좌우 여백
CHIP_FLOOR = 0.95  # 최소 폭 — 3자짜리도 칩으로 보이게
MAP_TOP = 3.88
MAP_BOTTOM = 5.97  # 캡션(6.09)에서 공기 0.12 확보 — 6.05로 두면 0.04로 붙는다(§6)
MAP_H = 0.24  # 칩 높이 — 계열 균일(P4⑧)
DELEGATED = "자동(위임판단)"  # aliases.json grade 원문. 이 값만 점선+accent


def chip_w(text, pt=None, inset=CHIP_INSET, floor=CHIP_FLOOR):
    """칩 폭을 **글자에서 파생**한다 — §6 "칩·박스는 풀폭 금지, 내용 폭 + 인셋으로".

    균일 폭(좌 1.85 / 우 3.45)으로 그렸더니 오딧 채움률이 6건 걸렸다: "위탁약정" 4자가
    3.45" 칩에서 12%, "상품권" 3자가 16%. 정렬된 열이라 균일 폭이 자연스러워 보였지만,
    죽은 회색은 죽은 회색이다(S4 한계 칩 3.30"→1.60"과 같은 병).
    정렬은 **폭이 아니라 한쪽 변**으로 잡는다 — 좌열은 우변, 우열은 좌변 고정.
    글자폭 근사는 오딧과 같은 식(한국어 1자 ≈ 0.8×pt)이라 규칙과 설계가 같은 자를 쓴다.
    """
    pt = pt or S["caption"]
    return max(len(text.strip()) * pt * 0.8 / 72 + inset, floor)


def _map_line(slide, x1, y1, x2, y2, *, dashed=False):
    """잇기 선. 화살촉은 목적지(기준서 개념) 쪽 — 방향이 곧 "진입"이다."""
    from pptx.oxml.ns import qn

    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = C["muted"]
    conn.line.width = Pt(1.0)
    ln = conn.line._get_or_add_ln()
    if dashed:
        ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "sm", "len": "sm"}))
    return conn


def bipartite_map(slide, terms):
    """실무어 → 기준서 개념 이분 매핑. 우열은 terms에서 **파생**(등장순 dedup).

    우열을 손으로 적지 않는 이유: 적는 순간 데이터와 어긋날 수 있고, N:1 수렴이 "내가 그렇게
    그렸다"가 되어버린다. 파생하면 `고객에게 지급할 대가`에 두 선이 모이는 건 aliases.json이
    그렇게 생겼기 때문이다. 반환: (좌 y맵, 우 y맵, 선 수) — 오딧·자기검증용.
    """
    concepts = []  # 등장순 dedup — dict.fromkeys는 순서 보존
    for _t, _g, cs in terms:
        for c_ in cs:
            if c_ not in concepts:
                concepts.append(c_)

    r_pitch = G.pitch(len(concepts), MAP_TOP, MAP_BOTTOM, MAP_H, what="연결 개념")
    r_y = {c_: MAP_TOP + i * r_pitch for i, c_ in enumerate(concepts)}
    # 좌열은 우열 **범위 중앙**에 맞춘다 — 5개를 7개 옆에 위에서부터 쌓으면 아래가 빈다(P4④)
    r_span = (len(concepts) - 1) * r_pitch
    l_pitch = G.pitch(len(terms), MAP_TOP, MAP_BOTTOM, MAP_H, what="용어")
    l_span = (len(terms) - 1) * l_pitch
    l_top = MAP_TOP + (r_span - l_span) / 2
    l_y = {t: l_top + i * l_pitch for i, (t, _g, _cs) in enumerate(terms)}

    # 여러 용어가 한 개념으로 모이면(N:1) 도착점을 **분산**한다 — 2026-07-15 제3자 채점
    # FAIL(#6): `고객에게 지급할 대가` 좌변 한 점에 리베이트의 실선 화살촉과 상품권의 점선
    # 화살촉이 겹쳐 하나로 뭉갰다(18배 확대 확인). 점선 종단이 실선 화살촉에 삼켜져,
    # 범례가 구분하라는 "위임판단"이 **도착점에서 실선으로 보였다** — N:1 수렴이 이 장의
    # 요점인데 바로 그 수렴 지점에서 형태 언어가 무너진 것이다.
    fan_in = {}
    for _t, _g, cs in terms:
        for c_ in cs:
            fan_in[c_] = fan_in.get(c_, 0) + 1
    slot = dict.fromkeys(fan_in, 0)

    # 선 먼저 → 칩이 화살촉을 덮는다(칩 변에서 깔끔히 끝남)
    n_lines = 0
    for t, g, cs in terms:
        for c_ in cs:
            n_in = fan_in[c_]
            # 도착 y를 칩 높이 안에서 n_in등분 — 1개면 정중앙, 2개면 위·아래로 갈린다
            off = (slot[c_] + 1) / (n_in + 1) if n_in > 1 else 0.5
            slot[c_] += 1
            _map_line(
                slide,
                MAP_L_RIGHT,
                l_y[t] + MAP_H / 2,
                MAP_R_X,
                r_y[c_] + MAP_H * off,
                dashed=(g == DELEGATED),
            )
            n_lines += 1

    for t, g, _cs in terms:
        deleg = g == DELEGATED
        w = chip_w(t)
        b = add_box(
            slide,
            MAP_L_RIGHT - w,  # 우변 고정
            l_y[t],
            w,
            MAP_H,
            fill=C["bg"],
            line=C["accent"] if deleg else C["muted"],
            line_w=1.25 if deleg else 0.75,
            shape="round",
        )
        set_shape_text(b, t, S["caption"], F["head"], C["primary"], bold=True)
    for c_ in concepts:
        b = add_box(
            slide,
            MAP_R_X,  # 좌변 고정
            r_y[c_],
            min(chip_w(c_), MAP_R_MAX - MAP_R_X),
            MAP_H,
            fill=C["bg_alt"],
            shape="round",
        )
        set_shape_text(b, c_, S["caption"], F["body"], C["primary"])
    return l_y, r_y, n_lines


def variant_c(prs, c=None):
    """C3 — 서사 3칼럼(왜/역할/구축) + 실물 증거(이분 매핑 + JSON 카드). c: 텍스트 override(None=골든 기본값). 좌표·색·폰트 고정."""
    c = {**DEFAULT_C, **(c or {})}
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    header(slide, c["headline"], c["kicker"])
    # ── 상단: 왜 / 역할 / 구축 3칼럼 (두괄식 볼드 헤드 + 완전 문장) ──
    nar_w = (G.RIGHT_EDGE - G.MARGIN_L - 2 * 0.4) / 3
    for i, (head, body) in enumerate(c["narratives"]):
        nx = G.MARGIN_L + i * (nar_w + 0.4)
        add_text(
            slide,
            nx,
            G.CONTENT_TOP,
            nar_w,
            0.28,
            [(f"0{i + 1}", {"color": C["accent"]}), (f"  {head}", {})],
            S["head"],
            F["head"],
            C["primary"],
            bold=True,
        )
        add_text(
            slide,
            nx,
            G.CONTENT_TOP + 0.36,
            nar_w,
            0.95,
            body,
            S["caption"],
            F["body"],
            C["muted"],
            line_spacing=1.25,
        )
        if i < 2:
            add_box(slide, nx + nar_w + 0.2, G.CONTENT_TOP, 0.012, 1.2, fill=C["bg_alt"])
    # ── 하단 좌: 이분 매핑 (실물 5건 → 개념 7개, 선 9개) ──
    map_w = MAP_R_MAX - MAP_X  # 7.60
    _subhead(slide, MAP_X, MAP_HEAD_Y, map_w, c["map_head"])
    add_text(  # 범례는 소제목과 같은 줄 우측 — 룰 아래에 두면 매핑 첫 행(3.88)을 침범한다
        slide,
        MAP_R_X,
        MAP_HEAD_Y + 0.06,
        MAP_R_MAX - MAP_R_X,
        0.22,
        c["map_legend"],
        S["foot"],
        F["body"],
        C["muted"],
        align=PP_ALIGN.RIGHT,
    )
    bipartite_map(slide, c["terms"])
    add_text(
        slide,
        MAP_X,
        6.09,
        map_w,
        0.26,
        c["map_caption"],
        S["caption"],
        F["body"],
        C["muted"],
    )
    # ── 하단 우: 실물 JSON 카드 (모노스페이스, 다크) ──
    jx, jy, jw, jh = 8.7, 3.35, G.RIGHT_EDGE - 8.7, 2.55
    card = add_box(slide, jx, jy, jw, jh, fill=C["primary"], shape="round")
    try:
        card.adjustments[0] = 0.06  # 라운드 과대 방지
    except (IndexError, ValueError):
        pass
    add_text(
        slide,
        jx + 0.25,
        jy + 0.15,
        jw - 0.5,
        0.24,
        c["json_title"],
        S["caption"],
        F["head"],
        C["bg_alt"],
        bold=True,
    )
    add_text(
        slide,
        jx + 0.25,
        jy + 0.5,
        jw - 0.5,
        jh - 0.7,
        [[(ln, {})] for ln in c["json_lines"]],
        S["foot"],
        "Consolas",
        C["bg"],
        line_spacing=1.25,
    )
    add_text(
        slide,
        jx,
        6.09,  # 매핑 캡션과 같은 baseline — 두 구획의 캡션이 어긋나면 바닥이 흔들려 보인다
        jw,
        0.26,
        c["json_caption"],
        S["caption"],
        F["body"],
        C["muted"],
    )
    bar_and_source(slide, c["bar"], source=c["source"])
    return slide


def audit(prs):
    EMU = 914400
    accent_hex = str(C["accent"])
    assert sum(n for _, n, _ in BREAKDOWN) == 423, "분해 합 검산 실패"
    fails = []
    for si, sl in enumerate(prs.slides):
        rules, solids, acc = [], [], 0
        for sh in sl.shapes:
            L, T, W, H = sh.left / EMU, sh.top / EMU, sh.width / EMU, sh.height / EMU
            if L == 0 and T == 0 and W > 13:
                continue
            if H <= 0.02 and W > 1:
                rules.append((L, T, W))
            elif H > 0.3 and W > 0.5 and sh.shape_type != 17:
                solids.append((L, T, W, H))
            if T + H > G.CONTENT_BOTTOM + 0.02 and T < G.BAR_Y - 0.02:
                fails.append((si, "bottom>6.35", round(T + H, 3)))
            if 11.5 < W < 13 and abs((L + W) - G.RIGHT_EDGE) > 0.02:
                fails.append((si, "풀폭 right≠12.733", round(L + W, 3)))
            try:
                if str(sh.fill.fore_color.rgb) == accent_hex:
                    acc += 1
            except Exception:
                pass
            if sh.has_text_frame:
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        if (
                            r.font.color
                            and r.font.color.rgb is not None
                            and str(r.font.color.rgb) == accent_hex
                        ):
                            acc += 1
        for rl, rt, rw in rules:
            for sl_, st, sw, shh in solids:
                if rl < sl_ + sw and rl + rw > sl_ and st - 0.005 <= rt <= st + shh - 0.005:
                    fails.append((si, "룰-채움 겹침", round(rt, 3)))
        if acc > 4:
            fails.append((si, "accent 초과", acc))
        print(f"slide{si}: accent {acc}곳")
    assert not fails, f"AUDIT FAIL {fails}"
    print("audit pass")


def main():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    variant_a(prs)
    variant_b(prs)
    variant_c(prs)
    audit(prs)
    from pathlib import Path

    out = Path(__file__).parent / "variants" / "s08_variants.pptx"
    prs.save(str(out))
    print(f"saved {out}")


if __name__ == "__main__":
    main()
