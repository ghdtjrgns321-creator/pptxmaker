"""물성 assert — design-rules P1을 산문에서 코드 게이트로 승격.

왜 필요한가(2026-07-14 실증): 스킬을 로드하고 design-rules를 통독하고 P1/P4를 인용한 직후에도
§5 위반(분기 내용에 셰브런 밴드)과 P4① 위반(문장 박스 나열)을 냈다. **규칙을 읽는 것과 적용하는
것은 다르다** — 내용의 물성을 오분류하면 규칙은 무력하다. 그래서 물성을 데이터로 선언하게 하고
실제 생성 도형과 기계가 대조한다. 선언과 형태가 어긋나면 빌드가 죽는다.

사용:
    from .shape_kind import assert_shape_kind
    SHAPE_KIND = "branch"          # 레이아웃 상단에 선언
    ...
    assert_shape_kind(slide, SHAPE_KIND)   # 도형 다 그린 뒤 마지막에

물성(kind):
    sequence — 분기 없는 단순 진행(셰브런 밴드 허용)
    branch   — 갈림길/상호배타(XOR). 셰브런 밴드 금지, 뿌리→커넥터→갈래로.
    contrast — 대비. 동일 지오메트리 + 채움 명도 차등 쌍이 있어야 함(문장 박스 나열 금지).
"""

from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.shapes import MSO_SHAPE

EMU = 914400
_PROGRESS_SHAPES = (MSO_SHAPE.CHEVRON, MSO_SHAPE.PENTAGON)
KINDS = ("sequence", "branch", "contrast")


def _autoshape_types(slide):
    out = []
    for sh in slide.shapes:
        try:
            out.append(sh.auto_shape_type)
        except Exception:
            pass
    return out


def _solid_boxes(slide):
    """솔리드 채움 도형의 (w, h, fill) — 대비 쌍 탐지용."""
    out = []
    for sh in slide.shapes:
        try:
            if sh.fill.type == MSO_FILL_TYPE.SOLID:  # str 비교 금지 — enum 직접 대조
                out.append(
                    (
                        round(sh.width / EMU, 2),
                        round(sh.height / EMU, 2),
                        str(sh.fill.fore_color.rgb),
                    )
                )
        except Exception:
            pass
    return out


def _contrast_pairs(slide, min_w=1.0, min_h=0.5, tol=0.02):
    """동일 지오메트리 + 채움 상이 쌍 수 — 대비가 형태로 인코딩됐는지."""
    boxes = [b for b in _solid_boxes(slide) if b[0] >= min_w and b[1] >= min_h]
    n = 0
    for i in range(len(boxes)):
        for j in range(i + 1, len(boxes)):
            same_geom = (
                abs(boxes[i][0] - boxes[j][0]) < tol and abs(boxes[i][1] - boxes[j][1]) < tol
            )
            if same_geom and boxes[i][2] != boxes[j][2]:
                n += 1
    return n


def assert_shape_kind(slide, kind):
    """선언 물성 ↔ 실제 도형 대조. 어긋나면 ValueError(빌드 중단)."""
    if kind not in KINDS:
        raise ValueError(f"알 수 없는 물성 kind={kind!r} — 허용: {KINDS}")

    progress = [t for t in _autoshape_types(slide) if t in _PROGRESS_SHAPES]

    if kind == "branch" and progress:
        raise ValueError(
            f"물성 위반(P1 · design-rules §5): kind='branch'인데 진행형 도형"
            f"(셰브런/펜타곤) {len(progress)}개 사용. "
            "§5 '셰브런 밴드는 분기 없는 단순 진행에만' — 갈림길(XOR)을 셰브런으로 이으면 "
            "'A한 뒤 B한다'는 거짓 인과가 된다. 분기는 뿌리→커넥터→갈래(판정 diamond)로 그릴 것."
        )

    if kind == "contrast" and _contrast_pairs(slide) == 0:
        raise ValueError(
            "물성 위반(P1 · design-rules §5): kind='contrast'인데 대비 인코딩 0 — "
            "동일 지오메트리·채움 상이 도형 쌍이 없다. 문장을 박스에 넣은 나열은 대비가 아니다"
            "(P4 ①텍스트 덤프). 명도 차등 쌍·비례 바·경로 종착 대비 등으로 형태에 인코딩할 것."
        )

    return True
