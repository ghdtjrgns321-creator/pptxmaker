"""S9 지식그래프 — 고밀도 개편본 (s06 공식, 2026-07-20). 골든 편입 대상.

원본 s09_variants.variant_a(28pt 헤더 + 3서사 + 위계 트리(좌) + 간선 카탈로그(우) + 검은 바)를
S8과 동일한 승리 공식으로 재배치:
- 압축 헤더(dense.compact_header)
- 상단 시그니처 밴드: 위계 트리(주인공, 좌) + 간선 7종 카탈로그(주인공, 우) — 검증된 자산 재사용
- 하단: s06 표준 hero_card 4장(왜/역할/구축/규모검증) — 카드 컴포넌트 단일(dense.hero_card)
- 검은 전폭 바 폐지 → 결론은 카드1 불릿에 흡수 (하단 한줄평 바 금지, design-rules §8)

트리는 원본 좌표를 top-band로 선형 압축(ty/TH), 카탈로그는 s09_variants.build_catalog를
상단 밴드 좌표로 재사용(모듈 상수 오버라이드). 로직·색·노드 클래스 인코딩은 원본 그대로.
"""

from pptx.enum.shapes import MSO_CONNECTOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from . import dense as D
from . import grid as G
from .figures import Box
from .figures import card_row as CARDS_ROW
from .figures import elements as E
from .figures import relation_catalog as RELCAT
from .kit import SLIDE_H, SLIDE_W, add_box, add_text, load_kit

K = load_kit()

# content_contract 계약 키 — registry가 이 장을 golden.tech_tree로 부를 때 필수 키 출처.
# 2026-07-26 `s09_variants`(sparse 시안)에서 이관. 골든 전용 글이다.
DEFAULT = {
    "kicker": "3. 기술 설명 — TECH 02 · Retrieve",
    "headline": "지식그래프 — 기준서의 구조를 그대로 옮긴 결정적 지도",
    "narratives": [
        (
            "왜 필요한가",
            "문단 뭉치에는 순서도 관계도 없다. 기준서는 계층·상호참조·사례가 얽힌 구조 — 그 구조를 보존해야 근거가 근거로 이어진다.",
        ),
        (
            "파이프라인에서의 역할",
            "Retrieve의 지도. 진입 개념에서 관계를 한 단계만 따라 문단·사례·근거를 수집한다 — 수집 범위 자체가 그래프로 고정된다.",
        ),
        (
            "어떻게 만들었나",
            "기준서 공식 소제목 80개를 그대로 개념 노드로, 뼈대 간선은 기준서에서 기계 생성 — AI 판단 0. 고립 노드 0을 감사로 확인했다.",
        ),
    ],
    "struct_head": "구조 — 기준서의 위계 그대로",
    "root": "기준서 1115",
    "concepts": ["변동대가", "보증", "⋯ 80"],
    "paras": ["문단 50", "문단 56", "B33"],
    # 거터 라벨에 **단위**를 박는다 — 2026-07-15 제3자 채점 FAIL(#5): `문단 250`(개수)과
    # `문단 50`(문단 번호)이 같은 행에 93px 간격으로 나란해 어느 쪽이 개수고 어느 쪽이 ID인지
    # 글자만으로 못 갈랐다. 이 장은 "명사+숫자" 꼴이 개수·문서번호·문단ID 세 뜻을 겸한다
    # (기준서 1115=문서번호 · 개념 80=개수 · 문단 50=ID). 오독을 유발한 게 아니라 강제했다.
    "layer_tag1": "개념층 80개",
    "layer_tag2": "문단층 250개",
    "term_chip": "“볼륨디스카운트”",
    "term_cap": "용어 423 등재 · 400 진입",  # 등재는 423, 그래프 노드는 간선 보유 400 — 층이 다르다
    "case_chip": "사례 188",
    "cat_title": (
        "간선 7종 — 무엇이 무엇을 잇고, 어디서 왔나",
        "    7종 2,629 + BC 근거층 65 = 2,694 · 고립 0",
    ),
    # (종류, 수, 출발 노드, 종류 라벨, 도착 노드, 출처 규칙, 출발 종류, 도착 종류)
    # 전부 data/ontology 실측 1건 — 창작 0. 노드 종류: concept|para|term|case
    # 1~5행이 전부 **변동대가** 하나에 달려 있다(좌측 트리와 같은 개념) — 7종이 따로 노는 게
    # 아니라 한 개념에 실제로 다 붙는다는 걸 실물로 보인다. 6·7행은 `계약을 식별함`으로 연결.
    "cat_rows": [
        (
            "hier",
            "79",
            "변동대가",
            "속한다",
            "거래가격을 산정함",
            "기준서 목차 level · 기계 생성",
            "concept",
            "concept",
        ),
        ("cp", "250", "변동대가", "관할", "문단 50", "documentId · 기계 생성", "concept", "para"),
        (
            "e3",
            "244",
            "문단 56",
            "인용",
            "문단 53",
            "원문에 “문단 53” · 회계 항등식",
            "para",
            "para",
        ),
        (
            "case",
            "1,220",
            "QNA 201909A",
            "인용",
            "문단 50",
            "감사본 related_paragraphs",
            "case",
            "para",
        ),
        (
            "term",
            "755",
            "“볼륨디스카운트”",
            "진입",
            "변동대가",
            "사람 1차자료 전수검수 (TECH 01)",
            "term",
            "concept",
        ),
        (
            "example",
            "74",
            "IE 사례 1",
            "공식예시",
            "계약을 식별함",
            "IE 목차 매핑 · 사람 확정",
            "case",
            "concept",
        ),
        (
            "e2",
            "7",
            "계약을 식별함",
            "먼저판단",
            "수행의무를 식별함",
            "5단계 모형 1→2 · 사람 확정",
            "concept",
            "concept",
        ),
    ],
    # 노드 모양 범례는 두지 않는다 — 칩 안의 실제 텍스트("문단 50"·"변동대가")가 종류를 이미
    # 말한다. 범례를 얹으면 해독표를 하나 더 만드는 셈(§6 "축 매트릭스보다 스토리 구도").
    "node_caption": "노드 929 = 개념 80 + 문단 250 + 사례 188 + BC 그룹 11 + 용어 400(423 중 간선 보유)",
    "bar": "임베딩이 놓치는 '법적 이웃'을 관계가 잡는다 — 텍스트가 아니라 구조로 검색한다",
    "source": "출처: 2_DATA-TAXONOMY.md · 3_KNOWLEDGE-GRAPH.md (00_factsheet.md §C·§D)",
}


def _arrow(slide, x1, y1, x2, y2, *, width=1.5, size="med"):
    """방향 있는 커넥터(§5 — 커넥터는 전부 화살표). 트리·카탈로그 공용.

    variant_a 안의 클로저였던 것을 모듈 레벨로 올렸다 — build_catalog가 같은 화살표 문법을
    써야 하는데 함수 안에 갇혀 있으면 복붙본이 하나 더 생긴다(골든이 두 벌이었던 그 병).
    """
    from pptx.oxml.ns import qn

    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = C["muted"]
    conn.line.width = Pt(width)
    ln = conn.line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": size, "len": size}))
    return conn


def header(slide, headline, kicker):
    add_text(
        slide, G.MARGIN_L, 0.42, 8.0, 0.28, kicker, S["caption"], F["head"], C["muted"], bold=True
    )
    add_text(
        slide,
        G.MARGIN_L,
        0.72,
        G.RIGHT_EDGE - G.MARGIN_L,
        0.55,
        headline,
        S["section"],
        F["head"],
        C["primary"],
        bold=True,
    )
    add_box(slide, G.MARGIN_L, G.RULE_Y, G.RIGHT_EDGE - G.MARGIN_L, 0.014, fill=C["muted"])


def narrative_row(slide, narratives):
    nar_w = (G.RIGHT_EDGE - G.MARGIN_L - 2 * 0.4) / 3
    for i, (head, body) in enumerate(narratives):
        nx = G.MARGIN_L + i * (nar_w + 0.4)
        add_text(
            slide,
            nx,
            G.CONTENT_TOP,
            nar_w,
            0.28,
            [(f"0{i + 1}", {"color": C["accent"]}), (f"  {head}", {})],
            S["head"],
            F["head"],
            C["primary"],
            bold=True,
        )
        add_text(
            slide,
            nx,
            G.CONTENT_TOP + 0.36,
            nar_w,
            0.95,
            body,
            S["caption"],
            F["body"],
            C["muted"],
            line_spacing=1.25,
        )
        if i < 2:
            add_box(slide, nx + nar_w + 0.2, G.CONTENT_TOP, 0.012, 1.2, fill=C["bg_alt"])


# ── 간선 카탈로그 GRID (표를 걷어낸 자리 — 2026-07-15 사용자 반려 "이해가 안 가는 시각화") ──
#
# 물성 선언(design-rules P1): **실물 대조(종류별 1건)**. 이 장이 말하려는 건 "무엇이 무엇을
#   잇는가"인데, 표는 그걸 40자 문장으로 서술하고 `수` 열에 숫자를 세웠다. 게다가 그 열에
#   **노드와 간선이 섞여**(80은 노드, 79·1,220·244는 간선) 더할 수 없는 숫자가 한 줄에 섰다.
#   합이 941인데 제목은 929였고, 표라서 아무도 검산하지 않았다(간선도 폐기값 2,697을 썼다).
#   → 각 행이 **실제 노드 도해**(`[A] ─종류─▶ [B]`). 글자가 아니라 도형과 화살표다.
#   노드 종류는 **모양**으로 구분한다(개념=round 진한 테두리 / 문단=rect / 사례=rect bg_alt /
#   용어=round 흰 칩). §6 "숫자 스탯 단독 나열 금지 — 관계 인코딩과 결합".
#
# 어휘 분리(§6 장 내 중복 금지): 좌측 트리 = **fan**(한 개념에서 뻗음) / 우측 카탈로그 =
#   **행 대조**(종류별 1건 병렬). 같은 노드를 두 각도로 본다 — 1~5행이 전부 좌측 트리의
#   `변동대가`에 달려 있어 두 구획이 실물로 맞물린다.
#
# 색(P4⑩): accent 예산 **1**(3칼럼 번호 런이 3을 이미 씀, 상한 4). 카탈로그 행에는 accent를
#   쓰지 않는다 — 7행 중 3행(term·example·e2)이 "사람 확정"이라 색으로 구분하고 싶지만 예산이
#   없고, 무엇보다 이 장의 주장은 "사람 개입이 **적다**"이다. 출처 열의 글자가 그 일을 한다.
C, S, F = K["rgb"], K["sizes"], K["fonts"]

# content_contract 계약 키 — registry가 이 dense를 golden.tech_tree로 부를 때 필수 키 출처.

# 하단 카드 4 = s06 표준 hero_card (배지·태그·중앙 배너·히어로 아이콘·불릿·칩)
# (번호, 제목, 태그, 히어로 아이콘, 배너, 불릿3[(리드,설명)], 칩3[(아이콘,강조,라벨)])
# 카드 주제 = 지식그래프의 고유 성질 (제네릭 왜/역할/구축 대신, 각 장 내용에 맞춤)
CARD_KEYS = ("num", "title", "tag", "ic", "banner", "items", "chips")

CARDS = [
    (
        "1",
        "구조로 검색",
        "텍스트 아닌 관계",
        "workflow",
        "임베딩이 놓치는 법적 이웃",
        [
            ("문단 뭉치", " 순서·관계 없음"),
            ("기준서", " 계층·참조·사례 구조"),
            ("법적 이웃", " 관계가 잡는다"),
        ],
        [
            ("git-branch", "관계", " 로 연결"),
            ("search", "구조", " 검색"),
            ("circle-x", "임베딩", " 놓침"),
        ],
    ),
    (
        "2",
        "1홉 수집 지도",
        "Retrieve의 지도",
        "search",
        "수집 범위가 그래프로 고정",
        [
            ("진입 개념에서", " 관계 1단계만"),
            ("문단·사례·근거", " 결정적 수집"),
            ("수집 범위 자체", " 그래프로 고정"),
        ],
        [
            ("git-branch", "1홉", " 탐색"),
            ("layers", "문단·사례", " 수집"),
            ("circle-check-big", "결정적", " 범위"),
        ],
    ),
    (
        "3",
        "기계 생성 뼈대",
        "AI 판단 0",
        "git-branch",
        "공식 소제목 → 개념 80",
        [
            ("공식 소제목 80", " 그대로 개념 노드"),
            ("뼈대 간선", " 기준서에서 기계 생성"),
            ("AI 판단 0", " 사람 개입 최소"),
        ],
        [
            ("git-branch", "개념 80", ""),
            ("circle-x", "AI 판단", " 0"),
            ("book-open", "문단 250", ""),
        ],
    ),
    (
        "4",
        "규모·검증",
        "노드 929 · 간선 2,694",
        "layers",
        "고립 0 · 전수 감사 확인",
        [
            ("노드 929", " 개념+문단+사례+용어"),
            ("간선 2,694", " 7종 + BC 근거층"),
            ("고립 노드 0", " 감사로 확인"),
        ],
        [
            ("layers", "929", " 노드"),
            ("git-branch", "2,694", " 간선"),
            ("circle-check-big", "0", " 고립"),
        ],
    ),
]

# ── 트리 선형 압축 — 원본 y(3.4~6.2) → 상단 밴드(0.95~) ──
BAND_A, BAND_B = 0.80, 0.95  # new_y = (y-3.4)*0.80 + 0.95


def ty(y):
    return (y - 3.4) * BAND_A + BAND_B


TH = 0.34 * BAND_A  # 압축된 노드 높이

# ── 간선 카탈로그 자리 (골든 실측) ──────────────────────────────────────────────
# 좌열(종류·수) 좌변이 CAT_X, 마지막 행 top이 CAT_BOTTOM에 오도록. 폭은 프레임 우변까지.
CAT_BOX = Box(5.00, 1.22, G.RIGHT_EDGE - 5.00, 3.28 - 1.22 + 0.26)


def _cat_row(r):
    """골든 cat_rows(8튜플) → 부품 계약. 튜플 자리를 이름으로 바꾼다."""
    kind, n, a, label, b, src, ka, kb = r
    return {
        "kind": kind, "count": n, "a": a, "a_kind": ka,
        "label": label, "b": b, "b_kind": kb, "src": src,
    }



def _tree(slide, c):
    """원본 위계 트리(variant_a)를 상단 밴드로 선형 압축 이식. 노드는 NODE_STYLE 단일 출처."""
    add_text(
        slide,
        G.MARGIN_L,
        0.80,
        4.0,
        0.24,
        c["struct_head"],
        S["caption"],
        F["head"],
        C["primary"],
        bold=True,
    )

    def _node(x, y, w, name, kind):
        return E.node_chip(slide, K, x, ty(y), w, TH, name, kind, size=S["foot"])

    root_x, root_w = 2.35, 1.6
    _node(root_x, 3.6, root_w, c["root"], "root")
    concept_pos = [(2.05, 1.0), (3.2, 0.7), (4.05, 0.65)]
    for name, (cx, cw) in zip(c["concepts"], concept_pos):
        _node(cx, 4.45, cw, name, "concept")
        E.arrow(slide, K, root_x + root_w / 2, ty(3.94), cx + cw / 2, ty(4.45))
    para_pos = [(1.75, 0.8, 2.55), (2.95, 0.8, 2.55), (3.95, 0.6, 3.55)]
    for name, (px, pw, src_cx) in zip(c["paras"], para_pos):
        _node(px, 5.35, pw, name, "para")
        E.arrow(slide, K, src_cx, ty(4.79), px + pw / 2, ty(5.35))
    add_text(
        slide,
        G.MARGIN_L,
        ty(4.5),
        1.15,
        0.20,
        c["layer_tag1"],
        S["foot"],
        F["head"],
        C["primary"],
        bold=True,
    )
    add_text(
        slide,
        G.MARGIN_L,
        ty(5.4),
        1.15,
        0.20,
        c["layer_tag2"],
        S["foot"],
        F["head"],
        C["primary"],
        bold=True,
    )
    _node(G.MARGIN_L, 3.6, 1.45, c["term_chip"], "term")
    add_text(
        slide, G.MARGIN_L, ty(3.98), 1.5, 0.16, c["term_cap"], S["foot"], F["body"], C["muted"]
    )
    E.arrow(slide, K, G.MARGIN_L + 1.45, ty(3.77), 2.05, ty(4.62))
    _node(G.MARGIN_L, 5.85, 1.1, c["case_chip"], "case")
    E.arrow(slide, K, G.MARGIN_L + 1.1, ty(6.02), 1.75, ty(5.69))
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(2.55), Inches(ty(5.52)), Inches(2.95), Inches(ty(5.52))
    )
    conn.line.color.rgb = C["muted"]
    conn.line.width = Pt(1.25)
    ln = conn.line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:headEnd"), {"type": "triangle", "w": "sm", "len": "sm"}))
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "sm", "len": "sm"}))
    add_text(
        slide,
        G.MARGIN_L,
        ty(6.24),
        4.3,
        0.14,
        c["node_caption"],
        S["foot"] - 1,
        F["body"],
        C["muted"],
    )


def build(prs, c=None):
    c = {**DEFAULT, **(c or {})}
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    D.compact_header(slide, c["kicker"], c["headline"])

    # ── 상단 시그니처 밴드 좌: 위계 트리 (주인공) ──
    _tree(slide, c)

    # ── 상단 시그니처 밴드 우: 간선 카탈로그 (주인공) ──
    # 2026-07-26 부품화: 전엔 여기서 `S9.CAT_TOP` 같은 **다른 모듈의 전역을 덮어썼다**.
    # 그러면 그 값이 build를 부른 뒤에만 맞아서, 그 상수를 읽는 오딧이 실행 순서에 따라
    # 다른 것을 재게 된다. 이제 자리는 인자(box)로만 전달한다.
    add_text(
        slide,
        CAT_BOX.x,
        0.80,
        CAT_BOX.w,
        0.24,
        [
            [
                (c["cat_title"][0], {"bold": True, "color": C["primary"], "font": F["head"]}),
                (c["cat_title"][1], {}),
            ]
        ],
        S["caption"],
        F["body"],
        C["muted"],
    )
    RELCAT.draw(slide, CAT_BOX, {"rows": [_cat_row(r) for r in c["cat_rows"]]}, K)

    # ── 하단: s06 표준 hero_card 4장 ──
    # 카드 띠 — 부품 호출(figures.card_row). 장은 자리만 정하고 장수 파생은 부품이 한다.
    CARDS_ROW.draw(
        slide,
        Box(G.MARGIN_L, 3.62, D.FULL_W, 3.32),
        {"cards": [dict(zip(CARD_KEYS, c, strict=True)) for c in CARDS]},
        K,
    )
    # 결론("법적 이웃을 구조로 검색")은 카드1에 흡수 — 하단 전폭 한줄평 바 금지(§8).
    D.source_line(slide, c["source"])
    return slide
