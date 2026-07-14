"""S15 시안 — 실제 결과: 골든테스트 92건 전/후 (신규 어휘: 덤벨 차트. UI 패스).

콘텐츠 실물: 00_factsheet.md §F(전/후 5지표)·§H(홀드아웃 성과) — 6_TEST-DECISIONS.md 실측.
실행: uv run python golden/s15_variants.py → golden/variants/s15_variants.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from . import grid as G
from .kit import SLIDE_H, SLIDE_W, add_box, add_text, load_kit, mix

K = load_kit()
C, S, F = K["rgb"], K["sizes"], K["fonts"]

KICKER = "4. 문제와 해결"
SOURCE = "출처: 6_TEST-DECISIONS.md (00_factsheet.md §F·§H) — 사람 작성 골든테스트 92건 홀드아웃"

# (라벨, 전 표기, 후 표기, 전 %, 후 %) — %는 위치 계산용 (§F 실측에서 파생: 72/92·78/92)
DUMBBELLS = [
    ("결론 재현 (92건)", "72건", "78건", 72 / 92 * 100, 78 / 92 * 100),
    ("산술 필요 판정", "33%", "100%", 33, 100),
    ("꼬리질문 일관성", "60%", "100%", 60, 100),
]

QUAL_CHIPS = [  # 수치 축이 없는 개선 2건 (§F)
    ("판단트리 오선택 14/27", "주제기반 다중주입으로 구조 해소"),
    ("생성 비용 — 월 상한(429) 초과", "문서 상한 복원으로 건당 약 14원"),
]

AX_L, AX_R = 2.75, 7.3  # 덤벨 축


def header(slide, headline):
    add_text(
        slide, G.MARGIN_L, 0.42, 8.0, 0.28, KICKER, S["caption"], F["head"], C["muted"], bold=True
    )
    add_text(
        slide,
        G.MARGIN_L,
        0.72,
        G.RIGHT_EDGE - G.MARGIN_L,
        0.55,
        headline,
        S["section"],
        F["head"],
        C["primary"],
        bold=True,
    )
    add_box(slide, G.MARGIN_L, G.RULE_Y, G.RIGHT_EDGE - G.MARGIN_L, 0.014, fill=C["muted"])


def _subhead(slide, x, y, w, text):
    add_text(slide, x, y, w, 0.3, text, S["head"], F["head"], C["primary"], bold=True)
    add_box(slide, x, y + 0.34, w, 0.012, fill=C["muted"])


def bar_and_source(slide, text):
    bar = add_box(slide, G.MARGIN_L, G.BAR_Y, G.RIGHT_EDGE - G.MARGIN_L, G.BAR_H, fill=C["primary"])
    tf = bar.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name, run.font.bold = F["head"], True
    run.font.size = Pt(S["body"])
    run.font.color.rgb = C["bg"]
    add_text(
        slide,
        G.MARGIN_L,
        G.SOURCE_Y,
        G.RIGHT_EDGE - G.MARGIN_L,
        0.3,
        SOURCE,
        S["foot"],
        F["body"],
        C["muted"],
    )


def _pos(v):
    return AX_L + (AX_R - AX_L) * v / 100


def variant_a(prs):
    """A — 좌 덤벨 3행(전 muted → 후 accent) + 우 대형 숫자·정성 개선 칩."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    header(slide, "골든테스트 92건 — 전환 후 무엇이 달라졌나")
    # ── 좌: 덤벨 차트 ──
    _subhead(slide, G.MARGIN_L, G.CONTENT_TOP, 7.35, "전 → 후 실측 (같은 92건, 같은 채점 기준)")
    # 눈금: 0 / 50 / 100
    grid_top, grid_bottom = 2.5, 5.35
    for v in (0, 50, 100):
        gx = _pos(v)
        add_box(slide, gx, grid_top, 0.008, grid_bottom - grid_top, fill=C["bg_alt"])
        add_text(
            slide,
            gx - 0.3,
            grid_bottom + 0.05,
            0.6,
            0.2,
            f"{v}%",
            S["caption"],
            F["body"],
            C["muted"],
            align=PP_ALIGN.CENTER,
        )
    dot = 0.17
    for i, (label, before_t, after_t, bv, av) in enumerate(DUMBBELLS):
        cy = 2.85 + i * 0.95
        add_text(
            slide,
            G.MARGIN_L,
            cy - 0.13,
            2.05,
            0.3,
            label,
            S["body"],
            F["head"],
            C["primary"],
            bold=True,
        )
        bx, ax_ = _pos(bv), _pos(av)
        add_box(slide, bx, cy - 0.025, ax_ - bx, 0.05, fill=mix(C["bg"], C["muted"], 0.45))
        add_box(slide, bx - dot / 2, cy - dot / 2, dot, dot, fill=C["muted"], shape="oval")
        add_box(slide, ax_ - dot / 2, cy - dot / 2, dot, dot, fill=C["accent"], shape="oval")
        add_text(
            slide,
            bx - 0.75,
            cy - 0.12,
            0.6,
            0.24,
            before_t,
            S["caption"],
            F["body"],
            C["muted"],
            align=PP_ALIGN.RIGHT,
        )
        add_text(
            slide,
            ax_ + 0.15,
            cy - 0.13,
            0.75,
            0.26,
            after_t,
            S["body"],
            F["head"],
            C["primary"],
            bold=True,
        )
    add_text(
        slide,
        G.MARGIN_L,
        5.75,
        7.35,
        0.55,
        "결론 재현은 92건 분모의 비율 위치로 표시(72→78). 산술 판정·꼬리질문은 전환 후 조정(AI 판단 위임·thinking 조정)의 실측이다.",
        S["caption"],
        F["body"],
        C["muted"],
        line_spacing=1.25,
    )
    # ── 우: 대형 숫자 + 정성 개선 칩 ──
    rx = 8.15
    rw = G.RIGHT_EDGE - rx
    add_box(slide, 7.9, G.CONTENT_TOP, 0.012, G.CONTENT_BOTTOM - G.CONTENT_TOP, fill=C["bg_alt"])
    _subhead(slide, rx, G.CONTENT_TOP, rw, "홀드아웃 최종")
    add_text(slide, rx, 2.42, rw, 0.75, "78 / 92", S["display"], F["head"], C["accent"], bold=True)
    add_text(
        slide,
        rx,
        3.2,
        rw,
        0.55,
        "결론 재현 — 실행 에러 0건, 정답 문서(QNA) 전면 차단 상태의 수치다 (격리 증명 0/92).",
        S["caption"],
        F["body"],
        C["text"],
        line_spacing=1.25,
    )
    add_text(
        slide,
        rx,
        3.85,
        rw,
        0.28,
        [
            [
                ("평균 총점 ", {}),
                ("4.79 → 4.88", {"bold": True, "color": C["primary"]}),
                ("  (5점 만점)", {}),
            ]
        ],
        S["body"],
        F["body"],
        C["muted"],
    )
    from pptx.enum.lang import MSO_LANGUAGE_ID

    for i, (head, body) in enumerate(QUAL_CHIPS):
        chip = add_box(slide, rx, 4.4 + i * 0.85, rw, 0.72, fill=C["bg_alt"], shape="round")
        tf = chip.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = Inches(0.15)
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.LEFT
        p1._p.get_or_add_pPr().set("eaLnBrk", "0")
        r1 = p1.add_run()
        r1.text = head
        r1.font.name, r1.font.bold = F["head"], True
        r1.font.size = Pt(S["caption"])
        r1.font.color.rgb = C["primary"]
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        p2._p.get_or_add_pPr().set("eaLnBrk", "0")
        r2 = p2.add_run()
        r2.text = f"→ {body}"
        r2.font.name = F["body"]
        r2.font.size = Pt(S["caption"])
        r2.font.color.rgb = C["muted"]
        for r in (r1, r2):
            r.font.language_id = MSO_LANGUAGE_ID.KOREAN
    bar_and_source(
        slide, "같은 92건, 같은 기준으로 잰 전과 후 — 개선은 조정이 아니라 구조에서 왔다"
    )
    return slide


MISS_ROWS = [  # §H L124 — 미재현 14건 케이스 렌즈 전수
    ("소프트축 채점 차이", 6),
    ("질문 프레임 상이", 4),
    ("진입 누락", 2),
    ("코퍼스 범위 밖", 2),
]


def variant_b(prs):
    """B — 92건 워플(재현 78/미재현 14) + 미재현 전수 원인 분해 바."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    header(slide, "골든테스트 92건 — 78건 재현, 못 맞힌 14건은 원인까지 전수 분해")
    from pptx.enum.lang import MSO_LANGUAGE_ID

    # ── 좌: 92건 워플 (23×4) ──
    lx, lw = G.MARGIN_L, 6.1
    _subhead(slide, lx, G.CONTENT_TOP, lw, "92건은 어떻게 끝났나")
    cell, pitch = 0.2, 0.25
    cols = 23
    wx, wy = lx, 2.5
    hit_color = C["primary"]
    miss_color = mix(C["bg"], C["muted"], 0.35)
    for n in range(92):
        r, c = divmod(n, cols)
        add_box(
            slide,
            wx + c * pitch,
            wy + r * pitch,
            cell,
            cell,
            fill=hit_color if n < 78 else miss_color,
        )
    waffle_h = 4 * pitch  # 1.04
    ly = wy + waffle_h + 0.25
    add_text(slide, lx, ly, 1.2, 0.6, "78", S["display"], F["head"], C["primary"], bold=True)
    add_text(
        slide,
        lx + 1.25,
        ly + 0.28,
        1.6,
        0.28,
        "결론 재현",
        S["body"],
        F["head"],
        C["primary"],
        bold=True,
    )
    add_text(slide, lx + 3.3, ly, 0.95, 0.6, "14", S["display"], F["head"], C["muted"], bold=True)
    add_text(
        slide, lx + 4.3, ly + 0.28, 1.7, 0.28, "미재현", S["body"], F["head"], C["muted"], bold=True
    )
    add_text(
        slide,
        lx,
        ly + 0.85,
        lw,
        0.75,
        "정답 문서(QNA)를 전면 차단하고 잰 성적이다 — 격리 증명 0/92 · 실행 에러 0건 · "
        "소프트 채점 92/92. 전환 직후 72에서 조정 후 78로.",
        S["caption"],
        F["body"],
        C["muted"],
        line_spacing=1.3,
    )
    # ── 우: 미재현 14건 전수 분해 바 ──
    rx = 6.75
    rw = G.RIGHT_EDGE - rx
    add_box(slide, 6.5, G.CONTENT_TOP, 0.012, G.CONTENT_BOTTOM - G.CONTENT_TOP, fill=C["bg_alt"])
    _subhead(slide, rx, G.CONTENT_TOP, rw, "미재현 14건 — 어디서 놓쳤나 (전수)")
    bar_x0 = rx + 1.9
    bar_max = 3.2  # 6건 기준 최대 폭
    for i, (label, n) in enumerate(MISS_ROWS):
        by = 2.55 + i * 0.62
        add_text(
            slide,
            rx,
            by + 0.02,
            1.85,
            0.26,
            label,
            S["caption"],
            F["head"],
            C["primary"],
            bold=True,
        )
        add_box(slide, bar_x0, by, bar_max * n / 6, 0.3, fill=C["primary"])
        add_text(
            slide,
            bar_x0 + bar_max * n / 6 + 0.1,
            by + 0.01,
            0.6,
            0.28,
            f"{n}건",
            S["body"],
            F["head"],
            C["primary"],
            bold=True,
        )
    add_text(
        slide,
        rx,
        5.1,
        rw,
        0.26,
        "6 + 4 + 2 + 2 = 14 — 잔여 없음, 전수 귀속",
        S["caption"],
        F["body"],
        C["muted"],
    )
    note = add_box(slide, rx, 5.5, rw, 0.8, fill=C["bg_alt"], shape="round")
    tf = note.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.15)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p._p.get_or_add_pPr().set("eaLnBrk", "0")
    r = p.add_run()
    r.text = (
        "문단 렌즈로 다시 갈라도 결론은 같다 — 상한 밀림 45% · 검색 실패 33% 등 "
        "검색 단계가 78%, 생성 단계 문제가 아니다."
    )
    r.font.name = F["body"]
    r.font.size = Pt(S["caption"])
    r.font.color.rgb = C["primary"]
    r.font.language_id = MSO_LANGUAGE_ID.KOREAN
    bar_and_source(
        slide, "78/92 — 못 맞힌 14건의 원인까지 전수로 안다, 이것이 이 수치를 믿을 수 있는 이유다"
    )
    return slide


MISS_BUCKETS = [  # 6_TEST §6.3 표 — 버킷 4 전수 (합 14)
    (
        "6",
        "기타 소프트축",
        "타 기준서 4건(1115 코퍼스 밖 — 인용 불가) + 헤지 2건(근거 확보, 결론 미확정)",
    ),
    (
        "4",
        "질문 프레임",
        "정답 규칙은 판단트리가 이미 주입 — 기준서 개념 교차연결, 고칠 대상이 아님",
    ),
    ("2", "진입 누락", "정답 개념 진입 실패 — 검색 개선으로 움직일 실익은 1건뿐"),
    ("2", "범위 밖 판정", "라우팅이 코퍼스 범위 밖으로 판정, 정답 회수 0"),
]

STATS = [  # 6_TEST §6.2 표 실물
    ("소프트 채점", "92 / 92"),
    ("실행 에러", "0건"),
    ("개선 추적", "재현 72 → 78"),
    ("평균 총점", "4.79 → 4.88 (5점)"),
]


def _kv_row(slide, x, y, w, k, v):
    from pptx.enum.lang import MSO_LANGUAGE_ID

    box = add_text(slide, x, y, w, 0.28, "", S["body"], F["body"], C["muted"])
    tf = box.text_frame
    p = tf.paragraphs[0]
    p._p.get_or_add_pPr().set("eaLnBrk", "0")
    r1 = p.add_run()
    r1.text = k
    r1.font.name = F["body"]
    r1.font.size = Pt(S["caption"])
    r1.font.color.rgb = C["muted"]
    r2 = p.add_run()
    r2.text = f"   {v}"
    r2.font.name, r2.font.bold = F["head"], True
    r2.font.size = Pt(S["body"])
    r2.font.color.rgb = C["primary"]
    for r in (r1, r2):
        r.font.language_id = MSO_LANGUAGE_ID.KOREAN


VARIANT_C_DEFAULTS = {  # 텍스트 내용만 — 좌표·색·크기·도형은 함수 코드에 고정, KICKER/SOURCE/STATS/MISS_BUCKETS 흡수
    "kicker": KICKER,
    "headline": "골든테스트 92건 — 정답 문서를 차단한 채 사람의 결론을 재현하는가",
    "col1_head": "① 어떻게 만든 시험인가",
    "doc_caption": "실제 질의회신 92건 — 질의자·회계기준원·해석위원회 작성",
    "doc_q_head": "질문",
    "doc_q_desc": "실무 사실관계",
    "doc_r_head": "회신",
    "doc_r_desc": "전문가의 결론",
    "dests": [
        "질문만 시스템에 입력",
        "회신은 채점 정답지로\n3축 대조 — 결론·문단·분기",
        "문서 자체는 검색 차단 ✕\n격리 증명 0/92",
    ],
    "col1_note": "자작 시나리오를 버리고 사람이 쓴 실제 문답으로 — AI가 만든 문제를 AI가 푸는 자가순환과, 자기 답을 자기가 인용하는 순환을 모두 차단했다. 재현 판정 = 결론축 ≥ 1 그리고 3축 합 ≥ 4.",
    "col2_head": "② 최종 기록",
    "final_big": "78 / 92",
    "final_label": "사람 전문가 결론 재현",
    "stats": STATS,
    "col2_note": "하드 인용 재현율 59.1%는 별도 보고 — 인용률만으로 품질을 재지 않는다.",
    "col3_head": "③ 미재현 14건 — 각각 왜",
    "buckets": MISS_BUCKETS,
    "col3_note": "합 6+4+2+2 = 14, 전수 귀속 — 검색 재설계로 실제 움직일 수 있는 것은 1건뿐이다.",
    "bar": "정답이 적힌 문서를 못 보게 한 상태에서 78/92 — 실패 14건까지 전수 해부된 수치다",
    "source": SOURCE,
}


def variant_c(prs, c=None):
    """C — 3구역: ①시험 구성(순환 차단 2겹) ②최종 기록 ③미재현 14건 각 설명. c=텍스트 override(None=골든)."""
    c = {**VARIANT_C_DEFAULTS, **(c or {})}
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    # 헤더 인라인 (kicker·headline만 c에서 — 좌표·색·크기는 header()와 동일하게 고정)
    add_text(
        slide,
        G.MARGIN_L,
        0.42,
        8.0,
        0.28,
        c["kicker"],
        S["caption"],
        F["head"],
        C["muted"],
        bold=True,
    )
    add_text(
        slide,
        G.MARGIN_L,
        0.72,
        G.RIGHT_EDGE - G.MARGIN_L,
        0.55,
        c["headline"],
        S["section"],
        F["head"],
        C["primary"],
        bold=True,
    )
    add_box(slide, G.MARGIN_L, G.RULE_Y, G.RIGHT_EDGE - G.MARGIN_L, 0.014, fill=C["muted"])
    from pptx.enum.lang import MSO_LANGUAGE_ID

    # ── ① 좌: 문서 분해도 — 실제 질의회신 한 건을 세 갈래로 쪼갠다 ──
    ax, aw = G.MARGIN_L, 3.8
    _subhead(slide, ax, G.CONTENT_TOP, aw, c["col1_head"])
    from pptx.enum.shapes import MSO_CONNECTOR
    from pptx.oxml.ns import qn

    def _fan_arrow(x1, y1, x2, y2, dashed=False):
        conn = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
        )
        conn.line.color.rgb = C["muted"] if dashed else C["primary"]
        conn.line.width = Pt(1.5)
        ln = conn.line._get_or_add_ln()
        if dashed:
            ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))
        ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "sm", "len": "sm"}))
        return conn

    # 문서 카드: 질문부 + 회신부 2밴드
    doc_x, doc_y, doc_w, doc_h = ax, 2.5, 1.5, 1.7
    add_text(
        slide,
        ax,
        2.22,
        aw,
        0.24,
        c["doc_caption"],
        S["caption"],
        F["body"],
        C["muted"],
    )
    add_box(slide, doc_x, doc_y, doc_w, doc_h, fill=C["bg"], line=C["primary"], line_w=1.5)
    add_box(slide, doc_x, doc_y + doc_h / 2, doc_w, 0.014, fill=C["primary"])
    add_text(
        slide,
        doc_x,
        doc_y + 0.14,
        doc_w,
        0.55,
        [[(c["doc_q_head"], {"bold": True, "color": C["primary"]})], [(c["doc_q_desc"], {})]],
        S["caption"],
        F["body"],
        C["muted"],
        align=PP_ALIGN.CENTER,
        line_spacing=1.2,
    )
    add_text(
        slide,
        doc_x,
        doc_y + doc_h / 2 + 0.14,
        doc_w,
        0.55,
        [[(c["doc_r_head"], {"bold": True, "color": C["primary"]})], [(c["doc_r_desc"], {})]],
        S["caption"],
        F["body"],
        C["muted"],
        align=PP_ALIGN.CENTER,
        line_spacing=1.2,
    )
    # 분해 3갈래 → 목적지 칩
    dest_x, dest_w = ax + 2.1, aw - 2.1
    dest_style = [  # (dy, fill, color, border, dashed, src_y) — 좌표·색 고정, 텍스트는 c["dests"]에서
        (2.5, C["bg"], C["primary"], C["muted"], False, doc_y + 0.42),
        (3.45, C["bg"], C["primary"], C["accent"], False, doc_y + doc_h - 0.42),
        (4.5, C["bg_alt"], C["muted"], None, True, doc_y + doc_h / 2),
    ]

    for (dy, fill, color, border, dashed, src_y), text in zip(dest_style, c["dests"]):
        chip_h = 0.62
        _fan_arrow(doc_x + doc_w, src_y, dest_x - 0.05, dy + chip_h / 2, dashed=dashed)
        chip = add_box(
            slide, dest_x, dy, dest_w, chip_h, fill=fill, line=border, line_w=1.5, shape="round"
        )
        tf = chip.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = Inches(0.1)
        for li, seg in enumerate(text.split("\n")):
            pp = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
            pp.alignment = PP_ALIGN.LEFT
            pp._p.get_or_add_pPr().set("eaLnBrk", "0")
            r = pp.add_run()
            r.text = seg
            r.font.name = F["head"] if li == 0 else F["body"]
            r.font.bold = li == 0
            r.font.size = Pt(S["caption"])
            r.font.color.rgb = color
            r.font.language_id = MSO_LANGUAGE_ID.KOREAN
    add_text(
        slide,
        ax,
        5.4,
        aw,
        0.9,
        c["col1_note"],
        S["caption"],
        F["body"],
        C["muted"],
        line_spacing=1.3,
    )
    # ── ② 중: 홀드아웃 최종 기록 ──
    bx = 4.85
    bw = 2.9
    add_box(slide, 4.6, G.CONTENT_TOP, 0.012, G.CONTENT_BOTTOM - G.CONTENT_TOP, fill=C["bg_alt"])
    _subhead(slide, bx, G.CONTENT_TOP, bw, c["col2_head"])
    add_text(
        slide, bx, 2.4, bw, 0.75, c["final_big"], S["display"], F["head"], C["accent"], bold=True
    )
    add_text(
        slide,
        bx,
        3.15,
        bw,
        0.3,
        c["final_label"],
        S["body"],
        F["head"],
        C["primary"],
        bold=True,
    )
    for i, (k, v) in enumerate(c["stats"]):
        _kv_row(slide, bx, 3.65 + i * 0.42, bw, k, v)
    add_text(
        slide,
        bx,
        5.5,
        bw,
        0.75,
        c["col2_note"],
        S["caption"],
        F["body"],
        C["muted"],
        line_spacing=1.25,
    )
    # ── ③ 우: 미재현 14건 각 설명 ──
    cx = 8.1
    cw = G.RIGHT_EDGE - cx
    add_box(slide, 7.85, G.CONTENT_TOP, 0.012, G.CONTENT_BOTTOM - G.CONTENT_TOP, fill=C["bg_alt"])
    _subhead(slide, cx, G.CONTENT_TOP, cw, c["col3_head"])
    for i, (n, head, desc) in enumerate(c["buckets"]):
        my = 2.35 + i * 0.86
        badge = add_box(slide, cx, my, 0.44, 0.44, fill=C["primary"])
        tf = badge.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = n
        r.font.name, r.font.bold = F["head"], True
        r.font.size = Pt(S["body"])
        r.font.color.rgb = C["bg"]
        add_text(
            slide,
            cx + 0.6,
            my - 0.02,
            cw - 0.6,
            0.26,
            head,
            S["body"],
            F["head"],
            C["primary"],
            bold=True,
        )
        add_text(
            slide,
            cx + 0.6,
            my + 0.26,
            cw - 0.6,
            0.5,
            desc,
            S["caption"],
            F["body"],
            C["muted"],
            line_spacing=1.2,
        )
    add_text(
        slide,
        cx,
        5.85,
        cw,
        0.45,
        c["col3_note"],
        S["caption"],
        F["head"],
        C["primary"],
        bold=True,
        line_spacing=1.2,
    )
    bar = add_box(slide, G.MARGIN_L, G.BAR_Y, G.RIGHT_EDGE - G.MARGIN_L, G.BAR_H, fill=C["primary"])
    tfb = bar.text_frame
    tfb.vertical_anchor = MSO_ANCHOR.MIDDLE
    pb = tfb.paragraphs[0]
    pb.alignment = PP_ALIGN.CENTER
    rb = pb.add_run()
    rb.text = c["bar"]
    rb.font.name, rb.font.bold = F["head"], True
    rb.font.size = Pt(S["body"])
    rb.font.color.rgb = C["bg"]
    add_text(
        slide,
        G.MARGIN_L,
        G.SOURCE_Y,
        G.RIGHT_EDGE - G.MARGIN_L,
        0.3,
        c["source"],
        S["foot"],
        F["body"],
        C["muted"],
    )
    return slide


def audit(prs):
    from pptx.oxml.ns import qn

    fails = []
    for si, sl in enumerate(prs.slides):
        accents, ovals = 0, 0
        for sh in sl.shapes:
            t, le = sh.top / 914400, sh.left / 914400
            b, r = t + sh.height / 914400, le + sh.width / 914400
            full_bleed = sh.width / 914400 >= SLIDE_W - 0.05
            if t < 6.4 and b > 6.37 and not full_bleed:
                fails.append((si, "bottom", round(b, 2), sh.shape_id))
            if r > 12.75 and not full_bleed:
                fails.append((si, "right", round(r, 2), sh.shape_id))
            for el in sh._element.iter(qn("a:prstGeom")):
                if el.get("prst") == "ellipse":
                    ovals += 1
            for el in sh._element.iter(qn("a:srgbClr")):
                if el.get("val") == "D66E3A":
                    if el.getparent().tag.endswith("}solidFill"):
                        accents += 1
        if accents > 4:
            fails.append((si, "accent>4", accents))
        if si == 0 and ovals != 6:  # vA 덤벨 점 3쌍
            fails.append((si, "dots≠6", ovals))
        if si == 1:  # vB 워플 92칸 (정사각 cell 0.2)
            waffle = sum(
                1
                for sh in sl.shapes
                if abs(sh.width / 914400 - 0.2) < 0.005 and abs(sh.height / 914400 - 0.2) < 0.005
            )
            if waffle != 92:
                fails.append((si, "waffle≠92", waffle))
    assert not fails, f"AUDIT FAIL {fails}"
    print("audit pass")


def main():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    variant_a(prs)
    variant_b(prs)
    variant_c(prs)
    audit(prs)
    out = Path(__file__).parent / "variants" / "s15_variants.pptx"
    prs.save(str(out))
    print(f"saved {out}")


if __name__ == "__main__":
    main()
