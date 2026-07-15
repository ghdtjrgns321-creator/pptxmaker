# 시안 K — S4 문제정의. 사용자 확정 구조(2026-07-15): 상단 게이트 인과도 + 하단 한계 4가지 수렴.
#
# 물성 선언(design-rules P1):
#   상단 = 분기(XOR) — 두 시스템의 차이는 "근거 부족을 인지하는 판정(◇)의 유무"다.
#     공통 뿌리(질문) → ◇ → 충분/부족, 부족에서 다시 두 갈래(확정 강행 / 유보).
#     **셰브런·펜타곤 0개** — 진행형 도형으로 XOR을 그리면 "확정한 다음 유보한다"는 거짓 인과가 된다
#     (2026-07-15 사용자 반려: "저 비교는 파이프라인으로 할 게 아니다" — 둘 다 4칸 밴드로 그려서
#      구조가 같아 보였다. 말하려는 것과 정반대였다).
#   상단 종착 = 대비(회복 가능성) — 2종은 경로가 재무제표까지 **끝까지 가서 박힘**(✕ 불가역),
#     1종은 경로가 **중간에 멈추고 되돌아옴**(↩ 시간만). 경로 도달 거리로 인코딩한다.
#     문장 박스 나열 금지(P4 ① 텍스트 덤프).
#   하단 = 수렴(fan-in) — 한계 4가지가 전부 2종 하나로 모인다. 나열이 아니라 "전부 한 방향으로
#     표출된다"가 형태. 상단=fan-out / 하단=fan-in으로 어휘 분리(§6 장 내 시각 어휘 중복 금지).
#
# 색 의미(P4 ⑩ — 제3자 채점 2회로 확정):
#   v1: 검정이 위=안전 / 아래=치명으로 뒤집혀 있었다.
#   v4: 그걸 고쳤더니 **재무제표(파국)와 결론 바(우리 원칙)가 똑같은 검정+샤프+흰볼드**가 됐다
#       — 같은 죄의 반복. 훑는 사람은 검은 덩어리 둘을 같은 종류로 분류한다.
#   확정 규칙: **accent는 면(面)으로만 위험을 말한다** — 2종 칩 · 재무제표 2곳뿐(accent 예산 2/4).
#     판정 단어(✕ 불가역 · ✕ 세울 수 없다 · ↩)는 전부 **primary 잉크**. accent 글자는 흰 바탕에서
#     3.41:1이라 판정을 지기엔 너무 흐리다(제3자 실측). **판정을 말하는 단어는 그것을 설명하는
#     산문보다 흐릴 수 없다.**
#   **primary 채움 = 결론 바 하나뿐**(우리 원칙). 나머지 면은 전부 무채색.
#   muted(#8A9099)는 이 장이 스스로 "무시해도 됨"으로 규정한 색이다(눈썹·푸터·괘선·버린 분기
#     "충분"·못 세운 ◇). **그 잉크로 제품의 강점을 쓰지 않는다** — v5에서 "↩ 되돌아옴"이 그랬다.
#   ⚠ accent를 셀 땐 채움만이 아니라 **테두리·커넥터 선**도 센다(v4에서 선을 안 세서 7개를 4로 통과시켰다).
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from . import grid as G
from .kit import SLIDE_H, SLIDE_W, add_box, add_text, load_kit, mix

K = load_kit()
C, S, F = K["rgb"], K["sizes"], K["fonts"]

# 논증을 지는 주석 잉크. muted(#8A9099)는 3.2:1이라 이 장의 **추론 전부**가 가장 안 읽히는
# 잉크로 쓰였다(제3자 실측) — 단순 명찰이 17.9:1을 가져갔다. 신규 hex 금지(§2)라 kit.mix로 파생.
INK_NOTE = mix(C["muted"], C["primary"], 0.55)

# ── 밴드 GRID (GRID 상수 파생 — 임의 좌표 금지) ──
B1_HEAD_Y = G.CONTENT_TOP  # 1.75 상단 밴드 소제목
BAND_RULE_DY = 0.36

# 인과도 3행 (상단부터: 충분 / 부족·일반LLM / 부족·본시스템)
# 두 행 모두 하위 캡션 1줄을 갖는다 — 캡션을 행 오른쪽에 세로로 쌓으면 넘쳐서 아래 밴드를
# 침범한다(v2 렌더 실측: 폭 1.51"에 4줄). 폭을 주려면 행 아래 풀폭 1줄이 유일한 자리.
ROW_A_Y = 2.26  # 충분 → 확정 (밴드1 룰 bottom 2.122에서 공기 0.138)
ROW_BC_H = 0.48
ROW_A_H = ROW_BC_H  # 같은 계열(행 노드)은 높이 균일 — v4는 0.42/0.48/0.58 3종이 섞여 있었다(P4⑧)
ROW_B_Y = 2.90  # 부족 → 일반 LLM (2종). bottom 3.38, 캡션 3.42~3.64
ROW_C_Y = 3.78  # 부족 → 본 시스템 (1종). B행 캡션에서 공기 0.14. bottom 4.26
CAP_DY = 0.04  # 행 bottom → 하위 캡션 top
CAP_H = 0.22
ROW_A_CY = ROW_A_Y + ROW_A_H / 2  # 2.44
ROW_B_CY = ROW_B_Y + ROW_BC_H / 2  # 3.12
ROW_C_CY = ROW_C_Y + ROW_BC_H / 2  # 3.96

# 복귀 호 — 1종은 경로가 **질문으로 되돌아온다**. v5까지 이걸 글자로만 주장했다:
# 반대편(2종)은 진짜 화살표로 프레임 끝까지 그려지는데 되돌아옴은 회색 캡션 하나였다.
# 제3자 2차 채점: "비교의 두 항이 서로 다른 층위(시각 vs 텍스트)에 있다" → FAIL 5건 중 4건이
# 이 요소 하나로 수렴했다. 실제 선으로 고리를 닫으면 좌하단 공백도 함께 메워진다.
BACK_Y = 4.62  # 복귀 호의 수평 구간
BACK_LABEL_Y = 4.38  # 호 위 라벨 (row C bottom 4.26에서 공기 0.12, 호에서 0.02 위)
BACK_LABEL_X = 2.40
BACK_LABEL_W = 3.50

# 뿌리 → ◇ → (충분 / 부족)
Q_W, Q_H = 1.55, 1.05
DIA_X, DIA_W, DIA_H = 2.45, 1.35, 0.95
LACK_CY = (ROW_B_CY + ROW_C_CY) / 2  # 3.83 — 부족 갈래 중심
DIA_CY = (ROW_A_CY + LACK_CY) / 2  # 3.225 — ◇ 중심 = 두 출력의 수직 중앙(§5 정렬 린터)
DIA_RIGHT = DIA_X + DIA_W  # 3.80

L2_X, L2_W = 4.30, 1.25  # 2단계 칩(확정 / 근거 부족)
L2_RIGHT = L2_X + L2_W  # 5.55
ROW_X = 5.95  # 3단계 행 시작
ACT_W = 2.35  # [시스템 · 행동]
ERR_X = ROW_X + ACT_W + 0.40  # 8.70
ERR_W = 1.90
TERM_X = ERR_X + ERR_W + 0.40  # 11.00
TERM_W = G.RIGHT_EDGE - TERM_X  # 1.733 — 2종만 풀폭 도달

# 하단 수렴 — 복귀 호(4.60)에서 공기 0.16
B2_HEAD_Y = 4.76  # 룰 5.12
# 칩 폭은 글자에 맞춘다 — 3.30"는 채움률 14~22%로 바마다 죽은 회색이 372~409px였다(제3자 실측).
# 2.30"도 26~35%로 부족(오딧 P4④). 가장 짧은 라벨("환각 위험" 5자) 기준으로 1.60".
LIM_X, LIM_W = G.MARGIN_L, 1.60  # 0.6~2.2
LIM_TOP = 5.28  # 룰 bottom 5.132에서 공기 0.148
LIM_PITCH = 0.25
LIM_H = 0.24  # 4항목 5.28~6.27
CONV_PT_X = 4.95  # 4갈래가 **한 점**으로 모이는 지점(진짜 fan-in)
CONV_CY = LIM_TOP + (LIM_PITCH * 3 + LIM_H) / 2  # 5.75 — 4항목 세로 중앙
# 수렴 종착 = 2종 칩이 아니라 **못 세운 ◇**. 소제목이 "왜 판정을 세울 수 없나"인데 2종으로
# 수렴시키면 상단이 이미 배달한 결론에 다시 착지한다(제3자: 재탕 FAIL). 상단의 **서 있는 ◇**와
# 같은 크기로 두어 "같은 판정인데 하나는 서고 하나는 못 선다"가 형태로 읽히게 한다.
CONV_X = 6.05
GHOST_CY = CONV_CY
MARK_X = CONV_X + DIA_W + 0.15  # ✕ 세울 수 없다
MARK_W = 1.55
CONV_NOTE_X = MARK_X + MARK_W + 0.25
CONV_NOTE_W = G.RIGHT_EDGE - CONV_NOTE_X

DEFAULT = {
    "kicker": "1. 문제 정의",
    "headline": "일반 LLM은 틀리는 게 아니라, 틀리는 방향이 문제다",
    "band1_head": "근거가 부족할 때 무엇을 하느냐가 갈린다",
    "question": '"이 계약, 수익을\n지금 인식해도\n됩니까?"',
    "gate": "근거 충분?",
    "branch_ok": "충분",
    "branch_lack": "부족",
    "ok_chip": "확정",
    "ok_note": "근거가 충분하면 두 시스템 모두 확정한다 — 여기서는 차이가 없다. 문제는 전부 아래, 근거가 부족한 구간에서 생긴다.",
    "lack_chip": "근거 부족",
    "rows": [
        {
            "act": "일반 LLM · 확정 강행",
            "err": "2종 · 허위 확정",
            "term": "재무제표 반영",
            "mark": "✕ 불가역",
            "note": "그대로 왜곡으로 직결된다",
        },
        {
            "act": "본 시스템 · 유보 + 근거 제시",
            "err": "1종 · 놓침",
            "term": "",
            "mark": "",
            "note": "",
        },
    ],
    "back_label": "↩ 근거를 보고 다시 묻는다 — 손실은 시간뿐",
    "band2_head": "범용 LLM은 왜 이 판정을 세울 수 없나 — 한계 4가지가 한 곳에서 만난다",
    "limits": ["근거 추적 불가", "재현성 없음", "판단 출처 창작", "환각 위험"],
    "converge_mark": "✕ 세울 수 없다",
    "converge_note": "판정이 서지 않으니 '부족'을 인지할 수단이 없다 — 그래서 늘 위쪽 경로로 나간다. 편향을 프롬프트가 아니라 구조로 걸어야 하는 이유다.",
    "bar": "확실할 때만 확정하고, 애매하면 근거를 보여주며 유보한다",
    "source": "출처: 1_OVERVIEW.md · PROJECT_OVERVIEW.md · 5_INTERFACE.md (00_factsheet.md §A·§B)",
}


def _band_head(slide, y, text):
    """밴드 소제목 + 풀폭 룰 (§6 구획은 여백이 아니라 장치로 나눈다)."""
    w = G.RIGHT_EDGE - G.MARGIN_L
    add_text(slide, G.MARGIN_L, y, w, 0.32, text, S["head"], F["head"], C["primary"], bold=True)
    add_box(slide, G.MARGIN_L, y + BAND_RULE_DY, w, 0.012, fill=C["muted"])


def _center(box, text, size, color, bold=True):
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = tf.margin_bottom = Inches(0)
    tf.margin_left = tf.margin_right = Inches(0.06)
    p = tf.paragraphs[0]
    p._p.get_or_add_pPr().set("eaLnBrk", "0")
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name, r.font.bold = F["head"], bold
    r.font.size = Pt(size)
    r.font.color.rgb = color


def _arrow(slide, x1, y1, x2, y2, *, elbow=False, color=None):
    """방향 화살표 — S6 확정 어휘(tailEnd 삼각형) 재사용. 화살촉은 항상 **끝점**(x2,y2)."""
    from pptx.enum.shapes import MSO_CONNECTOR
    from pptx.oxml.ns import qn

    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.ELBOW if elbow else MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    conn.line.color.rgb = color or C["muted"]
    conn.line.width = Pt(1.5)
    ln = conn.line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    return conn


def _line(slide, x1, y1, x2, y2):
    """화살촉 없는 선 — 수렴 대각선용(4갈래에 화살촉을 다 붙이면 한 점에서 뭉개진다)."""
    from pptx.enum.shapes import MSO_CONNECTOR

    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = C["muted"]
    conn.line.width = Pt(1.0)
    return conn


def _dash(shp):
    """테두리 점선화 — §5 형태 언어: 점선 = 확률·불확실(못 세운 판정), 실선 = 결정."""
    from pptx.enum.dml import MSO_LINE_DASH_STYLE

    shp.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return shp


def _branch_label(slide, x, y, text, ink):
    add_text(slide, x, y, 0.62, 0.22, text, S["caption"], F["head"], ink, bold=True)


def variant_k(prs, c=None):
    """S4 — 상단: 게이트 인과도(분기+종착 대비) · 하단: 한계 4가지 수렴. c: 텍스트 override."""
    c = {**DEFAULT, **(c or {})}
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])

    # ── 헤더 ──
    add_text(
        slide,
        G.MARGIN_L,
        0.42,
        6.0,
        0.28,
        c["kicker"],
        S["caption"],
        F["head"],
        C["muted"],
        bold=True,
    )
    add_text(
        slide,
        G.MARGIN_L,
        0.72,
        G.RIGHT_EDGE - G.MARGIN_L,
        0.55,
        c["headline"],
        S["section"],
        F["head"],
        C["primary"],
        bold=True,
    )
    add_box(slide, G.MARGIN_L, G.RULE_Y, G.RIGHT_EDGE - G.MARGIN_L, 0.014, fill=C["muted"])

    # ── 구획 1: 게이트 인과도 ──
    _band_head(slide, B1_HEAD_Y, c["band1_head"])

    # 뿌리(질문) → ◇ 판정
    q = add_box(
        slide,
        G.MARGIN_L,
        DIA_CY - Q_H / 2,
        Q_W,
        Q_H,
        fill=C["bg_alt"],
        line=C["muted"],
        line_w=0.75,
        shape="round",
    )
    _center(q, c["question"], S["caption"], C["primary"])
    _arrow(slide, G.MARGIN_L + Q_W, DIA_CY, DIA_X, DIA_CY)

    dia = add_box(
        slide,
        DIA_X,
        DIA_CY - DIA_H / 2,
        DIA_W,
        DIA_H,
        fill=C["bg"],
        line=C["primary"],
        line_w=1.25,
        shape="diamond",
    )
    _center(dia, c["gate"], S["caption"], C["primary"])

    # ◇ → 충분(위) / 부족(아래)
    _arrow(slide, DIA_RIGHT, DIA_CY, L2_X, ROW_A_CY, elbow=True)
    _arrow(slide, DIA_RIGHT, DIA_CY, L2_X, LACK_CY, elbow=True)
    # 분기 라벨 — accent 예산은 2종 경로 전용이므로 여기엔 쓰지 않는다(§2 상한 4)
    _branch_label(slide, DIA_RIGHT + 0.06, ROW_A_CY - 0.30, c["branch_ok"], C["muted"])
    _branch_label(slide, DIA_RIGHT + 0.06, LACK_CY + 0.06, c["branch_lack"], C["primary"])

    # 충분 → 확정 (양쪽 공통, 무채색)
    ok = add_box(
        slide,
        L2_X,
        ROW_A_Y,
        L2_W,
        ROW_A_H,
        fill=C["bg_alt"],
        line=C["muted"],
        line_w=0.75,
        shape="round",
    )
    _center(ok, c["ok_chip"], S["body"], C["primary"])
    add_text(
        slide,
        L2_RIGHT + 0.4,
        ROW_A_Y - 0.02,
        G.RIGHT_EDGE - L2_RIGHT - 0.4,
        ROW_A_H,
        c["ok_note"],
        S["caption"],
        F["body"],
        INK_NOTE,
        line_spacing=1.3,
    )

    # 부족 칩 → 두 갈래 (테두리도 accent 예산에 든다 — 여기선 primary로)
    lack = add_box(
        slide,
        L2_X,
        LACK_CY - ROW_A_H / 2,
        L2_W,
        ROW_A_H,
        fill=C["bg"],
        line=C["primary"],
        line_w=1.25,
        shape="round",
    )
    _center(lack, c["lack_chip"], S["body"], C["primary"])
    for cy in (ROW_B_CY, ROW_C_CY):
        _arrow(slide, L2_RIGHT, LACK_CY, ROW_X, cy, elbow=True)

    # 3단계 두 행 — 경로 도달 거리로 회복 가능성 인코딩
    row_style = [
        (ROW_B_Y, C["accent"], C["bg"], True),  # 2종: accent 채움, 끝까지 도달
        (ROW_C_Y, C["bg_alt"], C["primary"], False),  # 1종: 무채색, 중간에 멈춤
    ]
    for (ry, err_fill, err_ink, fatal), row in zip(row_style, c["rows"]):
        act = add_box(
            slide,
            ROW_X,
            ry,
            ACT_W,
            ROW_BC_H,
            fill=C["bg_alt"],
            line=C["muted"],
            line_w=0.75,
            shape="round",
        )
        _center(act, row["act"], S["caption"], C["primary"])
        rcy = ry + ROW_BC_H / 2
        _arrow(slide, ROW_X + ACT_W, rcy, ERR_X, rcy)
        err = add_box(slide, ERR_X, ry, ERR_W, ROW_BC_H, fill=err_fill, shape="round")
        _center(err, row["err"], S["caption"], err_ink)

        cap_y = ry + ROW_BC_H + CAP_DY
        if fatal:
            # 2종 — 경로가 **끝까지 가서** 재무제표에 박힌다. 도달 거리가 곧 불가역성.
            # 종착을 primary(검정)로 두면 결론 바와 같은 옷이 된다 — 파국과 우리 원칙이 한 색
            # (제3자: 바이트 동일). accent = 2종 경로 전용이므로 종착도 accent, 샤프 코너로 종결.
            _arrow(slide, ERR_X + ERR_W, rcy, TERM_X, rcy)
            term = add_box(slide, TERM_X, ry, TERM_W, ROW_BC_H, fill=C["accent"], shape="rect")
            _center(term, row["term"], S["caption"], C["bg"])
            add_text(
                slide,
                TERM_X,
                cap_y,
                TERM_W,
                CAP_H,
                row["mark"],
                S["caption"],
                F["head"],
                C["primary"],
                bold=True,
                align=PP_ALIGN.CENTER,
            )
        else:
            # 1종 — 경로가 **질문으로 되돌아온다**. 실제 복귀 호를 그린다(아래 _back_arc).
            # 짧게 끝나는 것만으로는 "되돌아옴"이 안 보인다 — 반대편은 진짜 선인데 이쪽만
            # 글자면 비교의 두 항이 다른 층위가 된다(제3자 2차 채점 FAIL 4건의 공통 원인).
            pass
        # 2종 캡션만 행 아래. 1종의 설명은 복귀 호 라벨이 진다(호 자체가 주장이므로).
        if fatal:
            add_text(
                slide,
                ERR_X,
                cap_y,
                TERM_X - ERR_X - 0.1,
                CAP_H,
                row["note"],
                S["caption"],
                F["body"],
                INK_NOTE,
            )

    # ── 복귀 호 — 1종은 질문으로 되돌아온다(고리를 닫는다) ──
    # 2종은 프레임 끝까지 나가서 안 돌아오고, 1종은 여기로 돌아온다. 두 경로가 **같은 층위**
    # (둘 다 실제 선)에 있어야 비교가 성립한다.
    back_x = ERR_X + ERR_W / 2
    q_cx = G.MARGIN_L + Q_W / 2
    _line(slide, back_x, ROW_C_Y + ROW_BC_H, back_x, BACK_Y)
    _line(slide, back_x, BACK_Y, q_cx, BACK_Y)
    # 화살촉은 **질문 카드 쪽**(끝점)에 — back=True면 시작점(아래)에 붙어 "질문에서 나간다"로
    # 거꾸로 읽힌다(v6 렌더 실측).
    _arrow(slide, q_cx, BACK_Y, q_cx, DIA_CY + Q_H / 2)
    add_text(
        slide,
        BACK_LABEL_X,
        BACK_LABEL_Y,
        BACK_LABEL_W,
        CAP_H,
        c["back_label"],
        S["caption"],
        F["head"],
        C["primary"],
        bold=True,
    )

    # ── 구획 2: 한계 4가지 수렴(fan-in) ──
    _band_head(slide, B2_HEAD_Y, c["band2_head"])
    for i, lim in enumerate(c["limits"]):
        ly = LIM_TOP + i * LIM_PITCH
        chip = add_box(slide, LIM_X, ly, LIM_W, LIM_H, fill=C["bg_alt"], shape="round")
        _center(chip, lim, S["caption"], C["primary"], bold=False)
        # 각 칩 우변 → **한 점**으로 모이는 대각선. 평행선 4개 + 세로 막대로 그리면 수렴이
        # 형태로 안 보인다(v3 렌더 실측: 그냥 화살표 5개로 읽혔다).
        _line(slide, LIM_X + LIM_W, ly + LIM_H / 2, CONV_PT_X, CONV_CY)
    # 수렴점 → **못 세운 ◇**. 상단 ◇와 같은 크기·같은 글자, 테두리만 점선(§5 점선 = 확률·불확실).
    # 같은 판정인데 하나는 서고 하나는 못 선다 — 그 대비가 이 구획의 전부다.
    _arrow(slide, CONV_PT_X, CONV_CY, CONV_X, CONV_CY)
    ghost = add_box(
        slide,
        CONV_X,
        GHOST_CY - DIA_H / 2,
        DIA_W,
        DIA_H,
        fill=C["bg"],
        line=C["muted"],
        line_w=1.25,
        shape="diamond",
    )
    _dash(ghost)
    _center(ghost, c["gate"], S["caption"], C["muted"])
    add_text(
        slide,
        MARK_X,
        GHOST_CY - CAP_H / 2,
        MARK_W,
        CAP_H,
        c["converge_mark"],
        S["caption"],
        F["head"],
        C["primary"],
        bold=True,
    )
    add_text(
        slide,
        CONV_NOTE_X,
        GHOST_CY - DIA_H / 2,
        CONV_NOTE_W,
        DIA_H,
        c["converge_note"],
        S["caption"],
        F["body"],
        INK_NOTE,
        line_spacing=1.3,
    )

    # ── 결론 바 + 출처 ──
    bar = add_box(slide, G.MARGIN_L, G.BAR_Y, G.RIGHT_EDGE - G.MARGIN_L, G.BAR_H, fill=C["primary"])
    _center(bar, c["bar"], S["body"], C["bg"])
    add_text(
        slide,
        G.MARGIN_L,
        G.SOURCE_Y,
        G.RIGHT_EDGE - G.MARGIN_L,
        0.3,
        c["source"],
        S["foot"],
        F["body"],
        C["muted"],
    )
