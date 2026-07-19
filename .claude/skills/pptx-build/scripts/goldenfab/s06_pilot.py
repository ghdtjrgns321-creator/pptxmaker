"""S06 파일럿 v6 — 골든 BCG 문법 베이스 + 8개 교정 (2026-07-18). 골든 미편입.

v5 실패 진단: 레퍼런스의 WHY/HOW/IMPACT 다크 바 골격을 이식 — 그건 밀도 원리의
예시였지 템플릿이 아니었다. 베이스는 골든 s06(BCG: 흰 바탕·가는 룰·무채+Ember).

v6 = 골든 s06 원형에서 시작해 8개 교정만 더한다:
- 유지: 전폭 플로우 도해(주인공 — 흰 라운드 노드·accent 다이아·화살표·노드 캡션·OUT 드랍),
        가는 룰 구획, 흰 바탕, 무채 + Ember 소량
- 교정 1·5·8: 키커 + 14pt 헤드라인 한 줄 + 룰 (28pt·리드 문단 폐지, 전 텍스트 8~14pt)
- 교정 2: 하단 검은 재서술 바 폐지 → 그 지면에 콘텐츠 패널 3 (WHY·IMPACT·5레이어)
- 교정 3: Lucide 아이콘 ~30 (행·칩·카드)
- 교정 4: units.md v4 명사구 재고 ~101단위
- 교정 6: 전 요소가 패널·카드·칩 안 (BCG 헤어라인 박스)
- 교정 7: 동형 카드 4 (배지→제목→태그→아이콘→배너→불릿3→아이콘 칩3)

실행:
    uv run python -c "
    import sys; sys.path.insert(0, '.claude/skills/pptx-build/scripts')
    from goldenfab.reference import new_presentation
    from goldenfab.s06_pilot import build
    prs = new_presentation(); build(prs); prs.save('golden/_pilot_s06/pilot.pptx')"
"""

from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from . import grid as G
from .kit import ROOT, SLIDE_H, SLIDE_W, add_box, add_text, load_kit, mix, set_shape_text

K = load_kit()
C, S, F = K["rgb"], K["sizes"], K["fonts"]
ICONS = ROOT / "assets/icons"

# ── 콘텐츠 (units.md v4 재고 — 이 파일 밖 문구 금지) ───────────────────────
KICKER = "2. 파이프라인"
HEADLINE = "질문에서 응답까지 — 분기까지 전부 결정적인 실행 그래프"
SOURCE = "출처: 2_PIPELINE.md §2.5·2.6 · 5_SEARCH-RUNTIME.md §5.2~5.13 (단위 재고: golden/_pilot_s06/units.md)"

FLOW = [  # (노드, 노드 캡션) — 골든 s06 원형 어휘
    ("질문", None),
    ("analyze", "용어사전 매칭"),
    ("판정", None),
    ("retrieve", "그래프 1홉 탐색"),
    ("generate", "판단트리 주입"),
    ("format", "감리 넛지"),
    ("응답", None),
]
OUT_BOX = "거절 메시지"
OUT_NOTE = "범위 밖 질문 — 본선 미진입 · 즉시 거절"

NODE_CARDS = [  # (배지, 노드명, 태그, 아이콘, 배너, 불릿 3, 아이콘 칩 3[(아이콘, 강조, 라벨)])
    (
        "1",
        "analyze",
        "범위·주제 판정",
        "search",
        "임베딩 0 — 사전 substring 진입",
        [
            "멀티턴 → standalone 질의 재구성",
            "scope guard — 타 기준서 강제 OUT",
            "topic_hint 우선 — gold 상단 진입",
        ],
        [
            ("circle-check-big", "42/42", " calc 판정"),
            ("book-open", "288", " 쿼리매핑 사전"),
            ("calendar", "최근 3턴", " 히스토리"),
        ],
    ),
    (
        "2",
        "retrieve",
        "결정적 1홉 탐색",
        "database",
        "용어 423 → 관할 문단 → e3 이웃",
        [
            "케이스·IE via_topic 직결 — 플러드 차단",
            "fetch 순서 보존 — $in 재배열",
            "MongoDB 원문 조회 — to_thread",
        ],
        [
            ("target", "69→23", " 케이스 정제"),
            ("gauge", "3.7→0.6s", " retrieve"),
            ("circle-x", "임베딩 0", ""),
        ],
    ),
    (
        "3",
        "generate",
        "듀얼 LLM 생성",
        "cpu",
        "계산 gpt-4.1-mini · 서술 Gemini",
        [
            "판단트리 다중 주입 — 오선택 0/31",
            "인용 교집합 필터 — 실재분만",
            "thinking medium — 꼬리질문 100%",
        ],
        [
            ("circle-check-big", "산술 100%", ""),
            ("git-branch", "트리 41", " 다중 주입"),
            ("settings", "temp 0.0", " 결정성"),
        ],
    ),
    (
        "4",
        "format",
        "감리 넛지 부착",
        "file-text",
        "수집 문서 중 감리사례 1건 선별",
        [
            "감독원 지적사례 넛지 문장 부착",
            "상황형(clarify) 넛지 스킵",
            "summary_matcher 제거 — 그래프가 담아옴",
        ],
        [
            ("file-text", "감리 22건", " 코퍼스"),
            ("link", "넛지 1건", " 부착"),
            ("circle-x", "코사인 0", ""),
        ],
    ),
]

ASIS_LABEL = "AS-IS · 점수 기반 선별 (폐기)"
ASIS_ROWS = [  # (아이콘, 제목, 설명) — 한 줄 행
    ("layers", "임베딩+BM25+RRF", "하이브리드 검색 폐기"),
    ("scale", "Cohere 리랭커", "105건 중 103건 탈락"),
    ("settings", "weight_score", "근거 없는 리터럴"),
    ("target", "슬롯 상한 3·2·3", "매직넘버 핀포인트"),
    ("link", "tree·summary_matcher", "코사인 매칭"),
]
TOBE_LABEL = "TO-BE"
TOBE_CHAIN = [("쿼리매핑 288", "확장"), ("용어 423", "진입"), ("그래프 1홉", None)]
WHY_CAP = "검색 진입부 임베딩 호출 0"

IMPACT_ROWS = [  # (아이콘, 층, 실측)
    ("circle-check-big", "결정성", "같은 질문 → 같은 근거 문단"),
    ("trending-up", "재현율", "하드 49.1→59.1% · 소프트 78/92"),
    ("gauge", "응답 속도", "retrieve 3.7→0.6s · 전체 −22%"),
    ("coins", "비용", "건당 약 14원 · 상한 12건"),
]
IMPACT_CAP = "근거 선정에 LLM 개입 0"

LAYER_ROWS = [  # (아이콘, 제목, 설명)
    ("workflow", "① 판단트리 주입", "절차는 원문 앵커 — 창작 0"),
    ("lock", "② 구조화 출력", "cited 필드 스키마 강제"),
    ("triangle-alert", "③ 빈 인용 차단", "ModelRetry 재생성 ×2"),
    ("shield-check", "④ reasoning_guard", "미확인 확정 → 조건부 강등"),
    ("circle-check-big", "⑤ 교집합 필터", "컨텍스트 실재분만 통과"),
]
LAYER_CAP = "강제하는 것은 근거의 존재 — 없는 인용은 화면 밖"
SIDE_NOTE = "* 곁가지 3 — clarify fast-path · search_id 캐시 스킵 · 범위 밖 즉시 거절"

# ── GRID ──────────────────────────────────────────────────────────────────
FULL_W = G.RIGHT_EDGE - G.MARGIN_L
RULE_Y = 0.66
FLOW_Y, NODE_H = 0.92, 0.46
CARD_Y, CARD_H = 2.06, 2.30
PANEL_Y, PANEL_H = 4.58, 2.38
SOURCE_Y = 7.16

line_soft = mix(C["primary"], C["bg"], 0.80)
line_mid = mix(C["primary"], C["bg"], 0.55)
ink_body = mix(C["muted"], C["primary"], 0.45)


def _icon(slide, name, variant, x, y, size):
    p = ICONS / f"{name}_{variant}.png"
    return slide.shapes.add_picture(str(p), Inches(x), Inches(y), Inches(size), Inches(size))


def _arrow(slide, x1, y1, x2, y2, color, w=1.5):
    """골든 플로우 화살표 — 실선 + 삼각 머리."""
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = color
    conn.line.width = Pt(w)
    ln = conn.line._get_or_add_ln()
    tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
    ln.append(tail)
    return conn


def _hrule(slide, x, y, w, color=None, weight=0.75):
    r = add_box(slide, x, y, w, 0.001, fill=None, line=color or line_soft, line_w=weight)
    return r


def _panel_label(slide, x, y, w, text):
    """골든 섹션 문법 — 볼드 라벨 + 가는 룰."""
    add_text(slide, x, y, w, 0.2, text, S["caption"], F["head"], C["primary"], bold=True)
    _hrule(slide, x, y + 0.24, w, color=line_mid, weight=1.0)


def _chip_row(slide, x, y, w, icon, em, label, *, h=0.22):
    chip = add_box(slide, x, y, w, h, fill=C["bg_alt"])
    tf = chip.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.26)
    p = tf.paragraphs[0]
    p._p.get_or_add_pPr().set("eaLnBrk", "0")
    p.alignment = PP_ALIGN.LEFT
    for txt, color, bold in ((em, C["accent"], True), (label, C["muted"], False)):
        if not txt:
            continue
        r = p.add_run()
        r.text = txt
        r.font.name = F["head"]
        r.font.size = Pt(S["foot"])
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.language_id = MSO_LANGUAGE_ID.KOREAN
    _icon(slide, icon, "primary", x + 0.06, y + (h - 0.14) / 2, 0.14)


def build(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])

    # ── 헤더 — 키커 + 14pt 한 줄 + 룰 (상단 의식 압축: ~0.7") ──
    add_text(
        slide, G.MARGIN_L, 0.14, 4.0, 0.16, KICKER, S["foot"], F["head"], C["muted"], bold=True
    )
    add_text(
        slide,
        G.MARGIN_L,
        0.32,
        FULL_W,
        0.28,
        HEADLINE,
        S["sub"],
        F["head"],
        C["primary"],
        bold=True,
    )
    _hrule(slide, G.MARGIN_L, RULE_Y, FULL_W, color=C["primary"], weight=1.4)

    # ══ 주인공 — 골든 s06 전폭 플로우 (흰 라운드 노드·accent 다이아·화살표) ══
    node_w = [1.00, 1.30, 1.10, 1.30, 1.30, 1.20, 1.00]
    gap = (FULL_W - sum(node_w)) / (len(FLOW) - 1)
    nx = G.MARGIN_L
    ncy = FLOW_Y + NODE_H / 2
    dia_cx = dia_bot = 0.0
    for j, ((name, capt), w) in enumerate(zip(FLOW, node_w)):
        if name == "판정":
            dia_cx = nx + w / 2
            dia_h = NODE_H + 0.26
            dia_bot = ncy + dia_h / 2
            nd = add_box(
                slide,
                nx + (w - 1.10) / 2,
                ncy - dia_h / 2,
                1.10,
                dia_h,
                fill=C["bg"],
                line=C["accent"],
                line_w=1.6,
                shape="diamond",
            )
            set_shape_text(nd, name, S["caption"], F["head"], C["accent"], bold=True)
            nd.text_frame.word_wrap = False
        else:
            first, last = j == 0, j == len(FLOW) - 1
            nd = add_box(
                slide,
                nx,
                FLOW_Y,
                w,
                NODE_H,
                fill=C["primary"] if last else (C["bg"] if first else C["bg_alt"]),
                line=C["primary"] if first else None,
                line_w=1.2,
                shape="round",
            )
            set_shape_text(
                nd, name, S["body"], F["head"], C["bg"] if last else C["primary"], bold=True
            )
            nd.text_frame.word_wrap = False
        if capt:
            add_text(
                slide,
                nx - 0.15,
                FLOW_Y + NODE_H + 0.06,
                w + 0.3,
                0.16,
                capt,
                S["foot"],
                F["body"],
                C["muted"],
                align=PP_ALIGN.CENTER,
            )
        if j < len(FLOW) - 1:
            x1 = nx + w + (0.10 if name == "판정" else 0.04)
            x2 = nx + w + gap - 0.04
            _arrow(slide, x1, ncy, x2, ncy, C["accent"] if name == "판정" else line_mid, w=1.5)
            if name == "판정":
                add_text(
                    slide,
                    nx + w,
                    ncy - 0.26,
                    gap,
                    0.16,
                    "IN",
                    S["foot"],
                    F["head"],
                    C["accent"],
                    bold=True,
                    align=PP_ALIGN.CENTER,
                )
        nx += w + gap

    _arrow(slide, dia_cx, dia_bot, dia_cx, dia_bot + 0.18, C["accent"], w=1.5)
    add_text(
        slide,
        dia_cx + 0.06,
        dia_bot,
        0.5,
        0.16,
        "OUT",
        S["foot"],
        F["head"],
        C["accent"],
        bold=True,
    )
    ob_w = 1.10
    ob = add_box(slide, dia_cx - ob_w / 2, dia_bot + 0.20, ob_w, 0.26, fill=C["bg_alt"])
    set_shape_text(ob, OUT_BOX, S["foot"], F["head"], C["primary"], bold=True)
    add_text(
        slide,
        dia_cx + ob_w / 2 + 0.15,
        dia_bot + 0.20,
        3.2,
        0.26,
        OUT_NOTE,
        S["foot"],
        F["body"],
        C["muted"],
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # ══ 동형 카드 4 — 배지→제목→태그→아이콘→배너→불릿3→아이콘 칩3 ══════════
    card_w = (FULL_W - 3 * 0.18) / 4
    for i, (num, name, tag, icon, banner, procs, chips) in enumerate(NODE_CARDS):
        cx = G.MARGIN_L + i * (card_w + 0.18)
        cy = CARD_Y
        add_box(slide, cx, cy, card_w, CARD_H, fill=C["bg"], line=line_soft, line_w=1.0)
        badge = add_box(slide, cx + 0.10, cy + 0.09, 0.24, 0.24, fill=C["primary"], shape="oval")
        set_shape_text(badge, num, S["foot"], F["head"], C["bg"], bold=True)
        add_text(
            slide,
            cx + 0.42,
            cy + 0.06,
            card_w - 0.80,
            0.28,
            name,
            S["body"],
            F["head"],
            C["primary"],
            bold=True,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        _icon(slide, icon, "accent", cx + card_w - 0.36, cy + 0.09, 0.24)
        add_text(
            slide,
            cx + 0.42,
            cy + 0.35,
            card_w - 0.56,
            0.18,
            tag,
            S["foot"],
            F["head"],
            C["accent"],
            bold=True,
        )
        ban = add_box(slide, cx + 0.10, cy + 0.56, card_w - 0.20, 0.22, fill=C["bg_alt"])
        set_shape_text(ban, banner, S["foot"], F["head"], C["primary"], bold=True)
        ban.text_frame.word_wrap = False
        for j, txt in enumerate(procs):
            add_text(
                slide,
                cx + 0.12,
                cy + 0.84 + j * 0.215,
                card_w - 0.22,
                0.22,
                [[("· ", {"bold": True, "color": C["accent"]}), (txt, {})]],
                S["foot"],
                F["body"],
                ink_body,
                anchor=MSO_ANCHOR.MIDDLE,
                line_spacing=1.0,
            )
        _hrule(slide, cx + 0.10, cy + 1.52, card_w - 0.20)
        for k, (cicon, em, label) in enumerate(chips):
            _chip_row(
                slide, cx + 0.10, cy + 1.58 + k * 0.23, card_w - 0.20, cicon, em, label, h=0.20
            )
    add_text(
        slide,
        G.MARGIN_L,
        CARD_Y + CARD_H + 0.05,
        FULL_W,
        0.16,
        SIDE_NOTE,
        S["foot"],
        F["body"],
        C["muted"],
    )

    # ══ 하단 패널 3 — 재서술 바 대신 콘텐츠 (WHY · IMPACT · 5레이어) ═══════
    panel_w = (FULL_W - 2 * 0.185) / 3

    # 패널 A. WHY — 걷어낸 것 → 남긴 것
    ax = G.MARGIN_L
    add_box(slide, ax, PANEL_Y, panel_w, PANEL_H, fill=C["bg"], line=line_soft, line_w=1.0)
    _panel_label(slide, ax + 0.14, PANEL_Y + 0.10, panel_w - 0.28, ASIS_LABEL)
    for i, (icon, title, sub) in enumerate(ASIS_ROWS):
        ry = PANEL_Y + 0.42 + i * 0.25
        _icon(slide, icon, "primary", ax + 0.14, ry + 0.02, 0.17)
        add_text(
            slide,
            ax + 0.38,
            ry,
            panel_w - 0.52,
            0.2,
            [[(title, {"bold": True, "color": C["primary"]}), ("  " + sub, {})]],
            S["foot"],
            F["body"],
            ink_body,
            anchor=MSO_ANCHOR.MIDDLE,
        )
    ty = PANEL_Y + 1.78
    add_text(
        slide,
        ax + 0.14,
        ty - 0.04,
        0.5,
        0.18,
        TOBE_LABEL,
        S["foot"],
        F["head"],
        C["accent"],
        bold=True,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    chx = ax + 0.62
    ch_w = (panel_w - 0.62 - 0.14 - 2 * 0.24) / 3
    for label, arrow_label in TOBE_CHAIN:
        cb = add_box(slide, chx, ty - 0.06, ch_w, 0.26, fill=C["bg"], line=C["primary"], line_w=1.1)
        set_shape_text(cb, label, S["foot"], F["head"], C["primary"], bold=True)
        cb.text_frame.word_wrap = False
        if arrow_label:
            add_text(
                slide,
                chx + ch_w - 0.02,
                ty - 0.06,
                0.28,
                0.26,
                "→",
                S["caption"],
                F["head"],
                C["accent"],
                bold=True,
                align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE,
            )
        chx += ch_w + 0.24
    _chip_row(
        slide, ax + 0.14, PANEL_Y + PANEL_H - 0.36, panel_w - 0.28, "search", WHY_CAP, "", h=0.24
    )

    # 패널 B. IMPACT — 보장되는 것
    bx = G.MARGIN_L + panel_w + 0.185
    add_box(slide, bx, PANEL_Y, panel_w, PANEL_H, fill=C["bg"], line=line_soft, line_w=1.0)
    _panel_label(slide, bx + 0.14, PANEL_Y + 0.10, panel_w - 0.28, "IMPACT · 보장되는 것")
    for i, (icon, tier, meas) in enumerate(IMPACT_ROWS):
        ry = PANEL_Y + 0.46 + i * 0.375
        _icon(slide, icon, "primary", bx + 0.14, ry + 0.01, 0.20)
        add_text(
            slide,
            bx + 0.42,
            ry - 0.03,
            1.10,
            0.22,
            tier,
            S["caption"],
            F["head"],
            C["primary"],
            bold=True,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        add_text(
            slide,
            bx + 0.42,
            ry + 0.155,
            panel_w - 0.56,
            0.17,
            meas,
            S["foot"],
            F["body"],
            ink_body,
        )
    _chip_row(
        slide,
        bx + 0.14,
        PANEL_Y + PANEL_H - 0.36,
        panel_w - 0.28,
        "shield-check",
        IMPACT_CAP,
        "",
        h=0.24,
    )

    # 패널 C. 환각 방지 5레이어
    cx0 = G.MARGIN_L + 2 * (panel_w + 0.185)
    add_box(slide, cx0, PANEL_Y, panel_w, PANEL_H, fill=C["bg"], line=line_soft, line_w=1.0)
    _panel_label(slide, cx0 + 0.14, PANEL_Y + 0.10, panel_w - 0.28, "환각 방지 5레이어")
    for i, (icon, title, sub) in enumerate(LAYER_ROWS):
        ry = PANEL_Y + 0.44 + i * 0.27
        _icon(slide, icon, "primary", cx0 + 0.14, ry + 0.02, 0.17)
        add_text(
            slide,
            cx0 + 0.38,
            ry,
            panel_w - 0.52,
            0.2,
            [[(title, {"bold": True, "color": C["primary"]}), ("  " + sub, {})]],
            S["foot"],
            F["body"],
            ink_body,
            anchor=MSO_ANCHOR.MIDDLE,
        )
    _chip_row(
        slide,
        cx0 + 0.14,
        PANEL_Y + PANEL_H - 0.36,
        panel_w - 0.28,
        "circle-x",
        LAYER_CAP,
        "",
        h=0.24,
    )

    add_text(slide, G.MARGIN_L, SOURCE_Y, FULL_W, 0.2, SOURCE, S["foot"], F["body"], C["muted"])
    return slide
