"""S4 문제정의 — 부품 조합으로 전환 (2026-07-26). 도해는 `figures/`에서 꺼낸다.

## 이 장이 소유하는 것 / 소유하지 않는 것

장(3층 배치)이 정하는 것은 **어떤 부품을 어디에 얼마나 크게 앉히나**뿐이다. 도해 자체는
2층 부품이 그린다:

    구획1 게이트 인과도  → figures.gate_branch   (판정 하나에서 결과가 갈린다)
    구획2 한계 수렴      → figures.fan_in        (원인 N개가 한 곳에서 만난다)

`box` 좌표는 골든 실측값 그대로여서 부품 호출로 바꿔도 `compare_golden`이 **픽셀 동일**이다
— 그게 부품화에 손실이 없음을 증명하는 유일한 방법이다.

## 구성 근거 (재작업 이력 요약)

초판은 게이트 인과도를 상단 띠로 재드로하고 한계 수렴을 4개 hero_card로 파괴했다(관계를
목록으로 바꾸면 "한 곳에서 만난다"는 주장이 사라진다). v1은 재사용했으나 검은 바 자리를
메우려 도해를 세로로 늘려 밀도가 떨어졌다("비어 보인다"). v2에서 헤더를 압축하고 검은 결론
바를 우하단 해결 패널로 돌려 지금 구성이 됐다.
"""

from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches, Pt

from . import dense as D
from . import grid as G
from .figures import Box
from .figures import fan_in as FAN
from .figures import gate_branch as GATE
from .kit import SLIDE_H, SLIDE_W, add_box, add_text, load_kit

K = load_kit()
C, S, F = K["rgb"], K["sizes"], K["fonts"]

# content_contract 계약 키 — registry가 이 장을 golden.problem_grid로 부를 때 필수 키 출처.
# 이 값들은 **기본값일 뿐 골든 전용 글**이다. 실전 덱은 content로 전량 덮어야 하고, 덮지 않으면
# 공장 문(content_contract)이 빌드를 중단한다 — 안 그러면 K-IFRS 글이 다른 프로젝트 덱으로 샌다.
#
# 2026-07-26 `_variant_k`(sparse 시안)에서 이관: 옛 `limits`는 라벨만 있는 리스트라 이 장이 읽지
# 않는 죽은 키였다(모듈 상수를 대신 그렸다). 부품이 라벨+설명 두 줄을 받으므로 계약도 그 형태다.
DEFAULT = {
    "kicker": "1. 문제 정의",
    "headline": "일반 LLM은 틀리는 게 아니라, 틀리는 방향이 문제다",
    "band1_head": "근거가 부족할 때 무엇을 하느냐가 갈린다",
    "question": '"이 계약, 수익을\n지금 인식해도\n됩니까?"',
    "gate": "근거 충분?",
    "branch_ok": "충분",
    "branch_lack": "부족",
    "ok_chip": "확정",
    "ok_note": "근거가 충분하면 두 시스템 모두 확정한다 — 여기서는 차이가 없다. 문제는 전부 아래, 근거가 부족한 구간에서 생긴다.",
    "lack_chip": "근거 부족",
    "rows": [
        {
            "act": "일반 LLM · 확정 강행",
            "err": "2종 · 허위 확정",
            "term": "재무제표 반영",
            "mark": "✕ 불가역",
            "note": "그대로 왜곡으로 직결된다",
        },
        {
            "act": "본 시스템 · 유보 + 근거 제시",
            "err": "1종 · 놓침",
            "term": "",
            "mark": "",
            "note": "",
        },
    ],
    "back_label": "↩ 근거를 보고 다시 묻는다 — 손실은 시간뿐",
    "band2_head": "범용 LLM은 왜 이 판정을 세울 수 없나 — 한계 4가지가 한 곳에서 만난다",
    "limits": [
        {"label": "근거 추적 불가", "sub": "어느 문단으로 답했는지 못 밝힌다"},
        {"label": "재현성 없음", "sub": "같은 질문에 매번 다른 판정"},
        {"label": "판단 출처 창작", "sub": "없는 조항·문단을 지어낸다"},
        {"label": "환각 위험", "sub": "그럴듯한 허위를 확정으로 낸다"},
    ],
    "converge_mark": "✕ 세울 수 없다",
    "converge_note": "판정이 서지 않으니 '부족'을 인지할 수단이 없다 — 늘 위쪽 경로로 샌다.",
    "bar": "확실할 때만 확정하고, 애매하면 근거를 보여주며 유보한다",
    "source": "출처: 1_OVERVIEW.md · PROJECT_OVERVIEW.md · 5_INTERFACE.md (00_factsheet.md §A·§B)",
}

KICKER = DEFAULT["kicker"]
HEADLINE = DEFAULT["headline"]
SOURCE = DEFAULT["source"]

# ── 부품 자리 (골든 실측) ──────────────────────────────────────────────────────
# 도해 두 개가 세로로 쌓이고, 오른쪽 아래는 해결 패널이 차지한다. 그래서 하단 도해의 폭은
# 프레임 끝이 아니라 패널 좌변에서 끊긴다.
BAND1_HEAD_Y = 0.90
GATE_BOX = Box(G.MARGIN_L, 1.29, G.RIGHT_EDGE - G.MARGIN_L, 2.50)
BAND2_HEAD_Y = 4.02
PANEL_X = 8.55
FAN_BOX = Box(G.MARGIN_L, 4.66, PANEL_X - G.MARGIN_L, 2.12)


def _solution_panel(slide):
    """검은 바로 있던 해결 원리를, 빈 우하단을 채우는 단일 액센트 패널로.

    수렴 도해가 "왜 못 세우나"를 말한 뒤 이 패널이 "그래서 이 시스템은 세운다"로 받는다.
    하단 전폭 바가 아니고 4카드도 아니다 — 형태 자체가 금지된 어휘다(design-rules §8).

    **남은 부채:** 이 패널의 글이 아직 코드에 박혀 있다. 도해와 달리 패널은 배치 틀 쪽
    작업이라 이번 부품화 범위 밖이고, 배치 틀 도구함을 만들 때 함께 밖으로 밀어낸다.
    """
    px = PANEL_X
    pw = G.RIGHT_EDGE - px
    py, ph = 4.30, 2.36
    add_box(slide, px, py, pw, ph, fill=C["bg_alt"])
    add_box(slide, px, py, 0.05, ph, fill=C["accent"])  # 액센트 좌변 바(컨테인먼트)
    D.icon(slide, "circle-check-big", "accent", px + 0.24, py + 0.24, 0.30)
    add_text(
        slide,
        px + 0.64,
        py + 0.22,
        pw - 0.82,
        0.32,
        "그래서 — 편향을 구조로 건다",
        S["head"],
        F["head"],
        C["accent"],
        bold=True,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        px + 0.26,
        py + 0.72,
        pw - 0.5,
        0.9,
        [
            [
                ("확실할 때만 확정", {"bold": True, "color": C["primary"]}),
                ("하고, 애매하면 근거를 보여주며 ", {}),
                ("유보", {"bold": True, "color": C["primary"]}),
                ("한다. 판정을 프롬프트가 아니라 시스템 ", {}),
                ("구조", {"bold": True, "color": C["primary"]}),
                ("로 세운다.", {}),
            ]
        ],
        S["caption"],
        F["body"],
        D.ink_body,
        line_spacing=1.4,
    )
    # 조건 요약 칩 2 — 확정/유보의 트리거(밀도 보강)
    cy = py + ph - 0.52
    cw = (pw - 0.5) / 2 - 0.08
    for i, (em, label) in enumerate((("확정", "근거 충분"), ("유보", "근거 부족"))):
        cx = px + 0.26 + i * (cw + 0.16)
        chip = add_box(slide, cx, cy, cw, 0.36, fill=C["bg"], line=D.line_mid, line_w=0.75)
        tf = chip.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        tf.margin_left = Inches(0.14)
        p = tf.paragraphs[0]
        p._p.get_or_add_pPr().set("eaLnBrk", "0")

        for txt, color, bold in ((em + "  ", C["accent"], True), (label, C["muted"], False)):
            r = p.add_run()
            r.text = txt
            r.font.name, r.font.bold, r.font.color.rgb = F["head"], bold, color
            r.font.size = Pt(S["foot"])
            r.font.language_id = MSO_LANGUAGE_ID.KOREAN


def build(prs, c=None):
    c = {**DEFAULT, **(c or {})}
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])

    # 헤더는 content에서 읽는다 — 상수로 두면 실전 덱에 골든 글이 새 나간다.
    D.compact_header(slide, c.get("kicker") or KICKER, c.get("headline") or HEADLINE)

    # ── 구획 1: 판정 하나에서 결과가 갈린다 ──
    D.band_head(slide, BAND1_HEAD_Y, c["band1_head"])
    GATE.draw(
        slide,
        GATE_BOX,
        {
            "question": c["question"],
            "gate": c["gate"],
            "branch_ok": c["branch_ok"],
            "branch_lack": c["branch_lack"],
            "ok_chip": c["ok_chip"],
            "ok_note": c["ok_note"],
            "lack_chip": c["lack_chip"],
            "rows": c["rows"],
            "back_label": c["back_label"],
        },
        K,
    )

    # ── 구획 2: 한계 N가지가 한 곳에서 만난다 ──
    D.band_head(slide, BAND2_HEAD_Y, c["band2_head"])
    FAN.draw(
        slide,
        FAN_BOX,
        {
            "items": c["limits"],
            "target": c["gate"],  # 상단의 **서 있는** ◇와 같은 판정 — 여기선 못 선다
            "mark": c["converge_mark"],
            "note": c["converge_note"],
            "soft": True,
        },
        K,
    )
    _solution_panel(slide)

    D.source_line(slide, SOURCE)
    return slide
