"""1층 원소 — 도해 부품들이 공유하는 최소 단위 (2026-07-26 신설).

여기 있는 것을 장·부품이 **다시 만들지 않는다.** 재설계 전 실측이 그 비용을 보여준다:
같은 화살표가 14벌·같은 소제목이 10벌·수 배지가 3계열로 흩어져 있었다(부품 비율 4.0%).

원소는 자기 자리를 모른다 — 좌표는 부르는 쪽(2층 도해)이 정한다.
색·크기는 전부 `kit`(brand-kit.yaml)에서 오고 이 모듈에 hex·pt 리터럴을 두지 않는다.
"""

from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

from ..kit import add_box, add_text, mix, set_shape_text

# 글자 폭 근사 — 오딧(`audit.text_need_height`)과 **같은 자**를 쓴다. 규칙과 설계가 어긋나면
# 설계는 통과하는데 검사가 잡거나 그 반대가 된다.
CJK_W = 0.8  # 한국어 1자 ≈ 0.8 × pt


def text_w(text, pt):
    """글자 폭(inch). 칩·라벨 폭 파생의 단일 출처."""
    return len(str(text).strip()) * pt * CJK_W / 72


def chip_w(text, pt, inset, floor=0.0):
    """칩 폭 = 글자 폭 + 좌우 여백, 최소 폭 보장.

    **자리에 맞춰 자르지 않는다.** 폭만 줄이면 글자는 그대로라 칩이 2줄이 되고 글자가 칩
    높이 밖으로 넘친다(2026-07-26 목업 눈검증 — 상자는 자리 안이라 기하 게이트를 통과한다).
    안 들어가면 부르는 쪽이 `measure`로 알아채고 시끄럽게 죽어야 한다.
    """
    return max(text_w(text, pt) + inset, floor)


def chip(slide, x, y, w, h, text, kit, *, style="plain", pt=None, font=None):
    """라벨 칩 한 개. style: plain(회색 배지) · outline(테두리) · soft(강조 테두리=불확실·위임)

    칩은 **한 스타일 계열**이어야 한다 — 같은 회색이 어떤 자리에선 내부 노드, 어떤 자리에선
    사례를 뜻하면 그 도해는 채점에서 떨어진다(2026-07-15 제3자 채점 FAIL의 원인).
    """
    C, S, F = kit["rgb"], kit["sizes"], kit["fonts"]
    pt = pt or S["caption"]
    fills = {"plain": C["bg_alt"], "outline": C["bg"], "soft": C["bg"]}
    lines = {"plain": None, "outline": C["muted"], "soft": C["accent"]}
    widths = {"plain": None, "outline": 0.75, "soft": 1.25}
    b = add_box(
        slide,
        x,
        y,
        w,
        h,
        fill=fills[style],
        line=lines[style],
        line_w=widths[style],
        shape="round",
    )
    set_shape_text(b, text, pt, font or F["head"], C["primary"], bold=(style != "plain"))
    return b


def card_size(title, note, kit, *, pt=None, note_pt=None, inset=0.40, floor=1.15, h1=0.36, h2=0.22):
    """카드 고유 크기 — 폭은 **두 줄 중 긴 쪽**에서, 높이는 줄 수에서. 반환 (w, h).

    칩(`chip`)과 갈리는 지점: 칩은 라벨 한 겹이고 카드는 **라벨 + 보조 한 줄**이다.
    도해가 한 겹만 갖고 있으면 화면이 커질수록 초라해진다 — 노드마다 "무엇"과 "얼마/왜"가
    같이 있어야 밀도가 생긴다(2026-07-26 목업 반려: "그냥 너무 단순한 초등학생 PPT").
    """
    S = kit["sizes"]
    pt = pt or S["body"]
    note_pt = note_pt or S["caption"]
    w = max(text_w(title, pt), text_w(note, note_pt) if note else 0.0) + inset
    return max(w, floor), h1 + (h2 if note else 0.0)


def card(
    slide,
    x,
    y,
    w,
    h,
    title,
    kit,
    *,
    note=None,
    metric=None,
    style="plain",
    pt=None,
    note_pt=None,
    shape="round",
):
    """2~3층 노드 — (수치) + 굵은 라벨 + 회색 보조줄. style: plain|outline|soft|solid(중심·주장)."""
    C, S, F = kit["rgb"], kit["sizes"], kit["fonts"]
    pt = pt or S["body"]
    note_pt = note_pt or S["caption"]
    fills = {"plain": C["bg_alt"], "outline": C["bg"], "soft": C["bg"], "solid": C["primary"]}
    lines = {"plain": None, "outline": C["muted"], "soft": C["accent"], "solid": None}
    widths = {"plain": None, "outline": 0.75, "soft": 1.25, "solid": None}
    ink = C["bg"] if style == "solid" else C["primary"]
    sub = mix(C["bg"], C["primary"], 0.42) if style == "solid" else C["muted"]
    b = add_box(
        slide, x, y, w, h, fill=fills[style], line=lines[style], line_w=widths[style], shape=shape
    )
    paras = []
    if metric:
        paras.append([(metric, {"size": pt * 1.7, "bold": True, "color": ink})])
    paras.append([(title, {"size": pt, "bold": True, "color": ink})])
    if note:
        paras.append([(note, {"size": note_pt, "color": sub})])
    add_text(
        slide,
        x,
        y,
        w,
        h,
        paras,
        pt,
        F["head"],
        ink,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
        line_spacing=1.08,
    )
    return b


def ribbon(slide, kit, x1, y1, x2, y2, t1, t2, *, emph=False, hollow=False):
    """흐름 리본 — (x1,y1)에서 두께 t1로 나가 (x2,y2)에서 두께 t2로 닿는 사다리꼴.

    얇은 실선으로 이으면 **굵기가 곧 양**이라는 문법이 연결부에서 끊긴다(2026-07-26 목업
    눈검증 — 288건 기둥에서 머리카락 선이 나가 196건 띠에 닿았다). 면이 이어져야 양이 흐른다.
    """
    from pptx.util import Inches, Pt

    b = slide.shapes.build_freeform(Inches(x1), Inches(y1))
    b.add_line_segments(
        [
            (Inches(x2), Inches(y2)),
            (Inches(x2), Inches(y2 + t2)),
            (Inches(x1), Inches(y1 + t1)),
        ],
        close=True,
    )
    shp = b.convert_to_shape()
    shp.shadow.inherit = False
    C = kit["rgb"]
    if hollow:  # 빠져나간 양 — 채우면 남은 것과 같은 무게로 보인다
        shp.fill.background()
        shp.line.color.rgb = mix(C["muted"], C["bg"], 0.3)
        shp.line.width = Pt(0.75)
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = (
            mix(C["accent"], C["bg"], 0.55) if emph else mix(C["muted"], C["bg"], 0.68)
        )
        shp.line.fill.background()
    return shp


def legend(slide, kit, x, y, w, entries, *, size=0.15, gap=0.22):
    """범례 한 줄 — entries=[(kind, 설명), ...]. kind: full|half|solid|soft|hollow.

    마크의 뜻을 산문 각주로만 적으면 눈이 표와 각주를 왕복한다. 마크를 **실제로 그려** 옆에
    뜻을 붙인다(같은 잉크로 두 번 말하지 않는다).
    """
    C, S, F = kit["rgb"], kit["sizes"], kit["fonts"]
    pt = S["caption"]
    cur = x
    for kind, label in entries:
        cy = y + 0.10
        if kind in ("full", "half"):
            add_box(
                slide,
                cur,
                cy - size / 2,
                size,
                size,
                fill=C["primary"] if kind == "full" else None,
                line=None if kind == "full" else C["primary"],
                line_w=None if kind == "full" else 1.0,
                shape="oval",
            )
        else:
            sw = {
                "solid": (mix(C["accent"], C["bg"], 0.55), None),
                "soft": (mix(C["muted"], C["bg"], 0.68), None),
                "hollow": (None, mix(C["muted"], C["bg"], 0.3)),
            }[kind]
            add_box(
                slide, cur, cy - size / 2, size * 1.6, size, fill=sw[0], line=sw[1], line_w=0.75
            )
        adv = size * (1.6 if kind not in ("full", "half") else 1.0) + 0.08
        add_text(slide, cur + adv, y, w, 0.20, label, pt, F["body"], C["muted"])
        cur += adv + text_w(label, pt) + gap
    return cur - x


def edge_point(cx, cy, w, h, ang):
    """중심 (cx,cy)·크기 (w,h)인 상자에서 각도 `ang`(rad) 방향의 **변 위 한 점**.

    관계선은 상자 중심이 아니라 변에서 시작·끝나야 한다. 중심점끼리 이으면 선이 상자를
    관통한다(수직·수평 방향에서 특히 — 2026-07-26 방사 목업에서 12시·6시 선이 중심을 뚫었다).
    """
    import math

    dx, dy = math.cos(ang), math.sin(ang)
    tx = (w / 2) / abs(dx) if abs(dx) > 1e-6 else float("inf")
    ty = (h / 2) / abs(dy) if abs(dy) > 1e-6 else float("inf")
    t = min(tx, ty)
    return cx + dx * t, cy + dy * t


def link(slide, kit, x1, y1, x2, y2, *, dashed=False, width=1.0, head=True, color="muted"):
    """잇기 선 — 화살촉은 **목적지 쪽**. 방향이 곧 '무엇이 무엇으로 가나'다.

    점선 = 불확실·위임 판단(§5). 색은 muted가 기본 — 선까지 accent면 강조 예산이 넘친다
    (accent는 "AI가 판단한 자리" 같은 **주장이 있는 곳**에 쓴다).
    """
    from pptx.enum.shapes import MSO_CONNECTOR
    from pptx.oxml.ns import qn
    from pptx.util import Inches, Pt

    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = kit["rgb"][color]
    conn.line.width = Pt(width)
    ln = conn.line._get_or_add_ln()
    if dashed:
        ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))
    if head:
        ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "sm", "len": "sm"}))
    return conn


def elbow(slide, kit, x1, y1, x2, y2, *, width=0.75, dashed=False, dot=True):
    """직교 꺾은선 — 사선 대신 ㄱ자로 잇는다. 시작점에 접점 표시(작은 흰 원).

    실측 근거: `ref/bcggen_p023.png`(BCG 허브앤스포크 실물)는 중심에서 부서로 **사선을
    쏘지 않는다** — 점선 궤도 위 접점에서 나와 수평→수직→수평으로 꺾는다. 사선 방사는
    각도가 제각각이라 정돈돼 보이지 않고, 꺾은선은 배선도처럼 읽혀 "체계"라는 인상을 준다.
    """
    from pptx.enum.shapes import MSO_CONNECTOR
    from pptx.oxml.ns import qn
    from pptx.util import Inches, Pt

    C = kit["rgb"]
    mid = (x1 + x2) / 2
    pts = (
        [(x1, y1), (mid, y1), (mid, y2), (x2, y2)]
        if abs(x2 - x1) >= abs(y2 - y1)
        else [(x1, y1), (x1, (y1 + y2) / 2), (x2, (y1 + y2) / 2), (x2, y2)]
    )
    for (ax, ay), (bx, by) in zip(pts, pts[1:], strict=False):
        if abs(ax - bx) < 1e-6 and abs(ay - by) < 1e-6:
            continue
        conn = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Inches(ax), Inches(ay), Inches(bx), Inches(by)
        )
        conn.line.color.rgb = C["muted"]
        conn.line.width = Pt(width)
        if dashed:
            ln = conn.line._get_or_add_ln()
            ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))
    if dot:
        d = 0.20
        add_box(
            slide,
            x1 - d / 2,
            y1 - d / 2,
            d,
            d,
            fill=C["bg"],
            line=C["muted"],
            line_w=0.75,
            shape="oval",
        )
    return pts


def badge(slide, kit, cx, cy, text, *, d=0.24, on_dark=False):
    """번호 배지 — 도해의 노드와 옆 설명표의 행을 **같은 번호**로 잇는다.

    실측 근거: `ref/mck_p025.png`은 동심원 안 노드와 오른쪽 Description 표를 ①~⑤로 묶는다.
    번호가 없으면 독자가 위치로 대응을 추측해야 하고, 항목이 넷을 넘으면 그 추측이 깨진다.
    """
    C, S, F = kit["rgb"], kit["sizes"], kit["fonts"]
    b = add_box(
        slide,
        cx - d / 2,
        cy - d / 2,
        d,
        d,
        fill=C["bg"] if on_dark else C["primary"],
        shape="oval",
    )
    set_shape_text(
        b, str(text), S["foot"], F["head"], C["primary"] if on_dark else C["bg"], bold=True
    )
    return b


def curve(slide, kit, x1, x2, y, *, bulge=0.5, dashed=False, width=1.0, head=True, down=False):
    """호(arc) — **같은 줄에 있는** 두 노드를 잇는다. 직선으로 이으면 사이 노드를 관통한다.

    `bulge`: 부풀림 높이(inch). 곡선 자체가 "이 둘은 같은 층에 있다"를 말한다.
    `down`: 아래쪽으로 부푼다. 한 줄 위에만 호를 쌓으면 링크가 많을 때 서로 겹쳐 못 읽고
    아래 절반은 백지가 된다 — 위/아래를 갈라 쓰면 같은 자리에 두 배가 들어간다.
    """
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn
    from pptx.util import Inches, Pt

    x, w = min(x1, x2), abs(x2 - x1)
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ARC, Inches(x), Inches(y - bulge), Inches(w), Inches(bulge * 2)
    )
    # PowerPoint ARC 기본은 **1/4원**(12시→3시)이라 그대로 쓰면 호가 반쪽만 그려진다
    # (2026-07-26 목업 눈검증 — 왼쪽 노드에서 출발하지 않는 호가 나왔다).
    # 각도 단위는 1/60000도, 시계방향·0=3시 → 9시(180°)에서 3시(360°)까지가 위쪽 반원,
    # 3시(0°)에서 9시(180°)까지가 아래쪽 반원.
    av = shp._element.spPr.find(qn("a:prstGeom")).find(qn("a:avLst"))
    angles = (0, 10800000) if down else (10800000, 21600000)
    for nm, val in zip(("adj1", "adj2"), angles):
        gd = av.makeelement(qn("a:gd"), {"name": nm, "fmla": f"val {val}"})
        av.append(gd)
    shp.shadow.inherit = False
    shp.fill.background()
    shp.line.color.rgb = kit["rgb"]["muted"]
    shp.line.width = Pt(width)
    ln = shp.line._get_or_add_ln()
    if dashed:
        ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))
    if head:
        ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "sm", "len": "sm"}))
    return shp


def band(slide, kit, x, y, w, h, *, emph=False):
    """흐름 띠 — 높이가 곧 양이다. 라벨은 부르는 쪽이 얹는다."""

    C = kit["rgb"]
    fill = C["accent"] if emph else mix(C["muted"], C["bg"], 0.45)
    return add_box(slide, x, y, w, h, fill=fill)
