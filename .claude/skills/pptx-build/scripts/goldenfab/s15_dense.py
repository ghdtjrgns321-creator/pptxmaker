"""S15 골든테스트(검증) — dense v6(2026-07-21). P0 북극성: 결과 탭 고유 형태(기술 4카드 아님).

P0 적용:
- #1 형태는 탭의 일이 정한다 — s15 = 결과 탭. 기술설명 4카드 패턴 금지(그건 s08·s09·s11·s12 전용).
  결과가 말해야 할 것 = ① 몇/몇 합격 ② 떨어진 것 전수 원인 ③ 정답 차단 방법.
- #3 어울리는 시각화 — 비율 = 링도넛(78/92), 실패 분석 = **못 고침 13 vs 실익 1 비대칭**(이
  비대칭이 펀치라인: "14 떨어졌지만 실제 개선 대상은 1건뿐"), 방법 = 격리 흐름 도해.
- #2 밀도 — 모든 존 정보 · 원인마다 근거·수치 · 실측 박기.

콘텐츠 파라미터화(DEFAULT 노출 — 승격 대비). preflight_dense green 후 제시.
"""

from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from . import dense as D
from . import grid as G
from .kit import SLIDE_H, SLIDE_W, add_box, add_text, load_kit, mix

K = load_kit()
C, S, F = K["rgb"], K["sizes"], K["fonts"]

DEFAULT = {
    "kicker": "4. 문제와 해결",
    "headline": "골든테스트 92건 — 정답 문서를 차단하고 78건 재현, 떨어진 14건은 실제 개선 대상 1건뿐",
    "sublead": "실제 질의회신 92건 · 격리 증명 0/92 · 실행 에러 0 · 소프트 채점 92/92 · 재현 72→78 — 개선은 튜닝이 아니라 온톨로지 구조에서 왔다",
    "hit": 78,
    "total": 92,
    "setup_head": "① 어떻게 시험했나 — 정답을 못 보게 하고",
    "dests": [
        ("질문만 시스템 입력", "circle-check-big", C["primary"]),
        ("회신 = 채점 정답지\n3축: 결론·문단·분기", "circle-check-big", C["accent"]),
        ("문서 자체는 검색 차단\n격리 증명 0/92", "circle-x", C["muted"]),
    ],
    "result_head": "② 결과 — 92건 중",
    "result_stats": [
        ("실행 에러", "0건"),
        ("소프트 채점", "92 / 92"),
        ("재현 추적", "72 → 78"),
        ("평균 총점", "4.79 → 4.88"),
    ],
    "result_note": "정답이 적힌 문서를 검색에서 차단한 채 사람 전문가의 결론을 재현한 수치 — 하드 인용 재현 59.1%는 별도 보고.",
    "miss_head": "③ 떨어진 14건 — 전수 원인 귀속",
    "unfix_head": "못 고치는 유형 · 13건 — 시스템이 손댈 수 없는 구조적 비대상",
    "unfix": [
        (
            "6",
            "코퍼스 밖",
            "타 기준서 4건은 1115 코퍼스 밖이라 인용 자체가 불가 · 범위 밖 2건은 라우팅이 정답 회수 0으로 판정 — 대상 자체가 아니다",
        ),
        (
            "2",
            "결론 미확정",
            "헤지 2건 — 근거는 확보했으나 전문가도 결론을 확정하지 않은 유보 사안, 시스템이 만들어낼 정답 자체가 없다",
        ),
        (
            "4",
            "질문 프레임 상이",
            "정답 규칙은 판단트리가 이미 주입 · 기준서 개념도 교차연결됨 — 시스템 결함이 아니라 질문을 보는 각도가 다른 것뿐",
        ),
        (
            "1",
            "진입 실패(실익 없음)",
            "정답 개념에 검색이 미진입한 케이스 중, 검색을 고쳐도 회수 실익이 없는 1건 — 개선해도 안 잡힌다",
        ),
    ],
    "fix_head": "실제 개선 대상 · 1건",
    "fix_lead": "회수 가능한 건 단 1건 — 나머지 13은 손댈 대상이 아니다.",
    "fix_points": [
        ("이 1건", "진입 개념 후보를 놓친 검색 문제 — 진입점을 넓히면 회수 가능."),
        ("그래서", "실패의 78%가 검색 단계 — 개선축은 생성 아닌 검색뿐."),
    ],
    "source": "출처: 6_TEST-DECISIONS.md (00_factsheet.md §F·§H) — 사람 작성 골든테스트 92건 홀드아웃",
}


def _subhead(slide, x, y, w, text):
    add_text(slide, x, y, w, 0.26, text, S["caption"], F["head"], C["primary"], bold=True)
    D.hrule(slide, x, y + 0.32, w, color=D.line_mid, weight=1.0)


def _fan_arrow(slide, x1, y1, x2, y2, dashed=False):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = C["muted"] if dashed else C["primary"]
    conn.line.width = Pt(1.4)
    ln = conn.line._get_or_add_ln()
    if dashed:
        ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "sm", "len": "sm"}))


def _ring(slide, x, y, sz, hit, total):
    """78/92 링도넛 — 비율을 링으로 인코딩. 가운데 78."""
    cd = CategoryChartData()
    cd.categories = ["재현", "미재현"]
    cd.add_series("s", (hit, total - hit))
    gf = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT, Inches(x), Inches(y), Inches(sz), Inches(sz), cd
    )
    ch = gf.chart
    ch.has_legend = False
    ch.has_title = False
    plot = ch.plots[0]
    plot.has_data_labels = False
    try:
        plot.doughnut_hole_size = 64
    except Exception:  # noqa: BLE001
        pass
    pts = plot.series[0].points
    for pt, col in zip(pts, (C["primary"], mix(C["bg"], C["muted"], 0.28))):
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = col
        pt.format.line.color.rgb = C["bg"]
        pt.format.line.width = Pt(2)
    # 가운데: 78 (홀 안에 맞춤 — 겹침 없이) + /92 작게. '재현'은 링 밖 캡션.
    add_text(
        slide,
        x,
        y + sz / 2 - 0.36,
        sz,
        0.42,
        str(hit),
        S["title"],
        F["head"],
        C["primary"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        x,
        y + sz / 2 + 0.10,
        sz,
        0.2,
        f"/ {total}",
        S["caption"],
        F["head"],
        C["muted"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        x - 0.3,
        y + sz + 0.02,
        sz + 0.6,
        0.2,
        "재현 (홀드아웃)",
        S["foot"],
        F["head"],
        C["primary"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def build(prs, c=None):
    c = {**DEFAULT, **(c or {})}
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    D.compact_header(slide, c["kicker"], c["headline"])
    add_text(
        slide, G.MARGIN_L, 0.76, D.FULL_W, 0.24, c["sublead"], S["caption"], F["body"], C["muted"]
    )

    # ══ 상단 좌: ① 방법(격리 시험 흐름) ══
    ax = G.MARGIN_L
    _subhead(slide, ax, 1.14, 5.4, c["setup_head"])
    doc_x, doc_y, doc_w, doc_h = ax, 1.72, 1.5, 1.32
    add_box(slide, doc_x, doc_y, doc_w, doc_h, fill=C["bg"], line=C["primary"], line_w=1.5)
    add_box(slide, doc_x, doc_y + doc_h / 2, doc_w, 0.014, fill=C["primary"])
    add_text(
        slide,
        doc_x,
        doc_y + 0.10,
        doc_w,
        0.55,
        [[("질문", {"bold": True, "color": C["primary"]})], [("실무 사실관계", {})]],
        S["caption"],
        F["body"],
        C["muted"],
        align=PP_ALIGN.CENTER,
        line_spacing=1.15,
    )
    add_text(
        slide,
        doc_x,
        doc_y + doc_h / 2 + 0.10,
        doc_w,
        0.55,
        [[("회신", {"bold": True, "color": C["primary"]})], [("전문가 결론", {})]],
        S["caption"],
        F["body"],
        C["muted"],
        align=PP_ALIGN.CENTER,
        line_spacing=1.15,
    )
    dest_x, dest_w = ax + 2.3, 3.1
    src_ys = [doc_y + 0.34, doc_y + doc_h - 0.34, doc_y + doc_h / 2]
    dest_ys = [1.68, 2.24, 2.8]
    for (text, ic, color), sy, dyy in zip(c["dests"], src_ys, dest_ys):
        dashed = color == C["muted"]
        _fan_arrow(slide, doc_x + doc_w, sy, dest_x - 0.05, dyy + 0.25, dashed=dashed)
        chip = add_box(
            slide,
            dest_x,
            dyy,
            dest_w,
            0.5,
            fill=C["bg_alt"] if dashed else C["bg"],
            line=color,
            line_w=1.4,
            shape="round",
        )
        D.icon(slide, ic, "accent", dest_x + 0.14, dyy + 0.1, 0.3)
        tf = chip.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.55)
        tf.margin_right = Inches(0.08)
        for li, seg in enumerate(text.split("\n")):
            pp = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
            pp.alignment = PP_ALIGN.LEFT
            pp._p.get_or_add_pPr().set("eaLnBrk", "0")
            rr = pp.add_run()
            rr.text = seg
            rr.font.name = F["head"] if li == 0 else F["body"]
            rr.font.bold = li == 0
            rr.font.size = Pt(S["foot"])
            rr.font.color.rgb = C["primary"] if li == 0 else C["muted"]

    # ══ 상단 우: ② 결과(링도넛 78/92 + 지표) ══
    rx = 6.2
    add_box(slide, rx - 0.35, 1.05, 0.012, 2.35, fill=C["bg_alt"])
    _subhead(slide, rx, 1.14, G.RIGHT_EDGE - rx, c["result_head"])
    _ring(slide, rx, 1.62, 1.66, c["hit"], c["total"])
    # 지표 2×2 그리드 — 넓은 단행은 죽은 회색이라 좁은 칸으로(P4④·§6-D)
    sx = rx + 1.95
    col_w, sgap = 2.2, 0.12
    for i, (k, v) in enumerate(c["result_stats"]):
        col, rowi = i % 2, i // 2
        gx = sx + col * (col_w + sgap)
        yy = 1.66 + rowi * 0.46
        row = add_box(slide, gx, yy, col_w, 0.38, fill=C["bg_alt"])
        tf = row.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.14)
        p = tf.paragraphs[0]
        p._p.get_or_add_pPr().set("eaLnBrk", "0")
        for txt, col_, bold in ((k + "  ", C["muted"], False), (v, C["accent"], True)):
            r = p.add_run()
            r.text = txt
            r.font.name, r.font.bold, r.font.size, r.font.color.rgb = (
                F["head"],
                bold,
                Pt(S["foot"]),
                col_,
            )
    add_text(
        slide,
        sx,
        1.66 + 2 * 0.46 + 0.08,
        G.RIGHT_EDGE - sx,
        0.7,
        c["result_note"],
        S["foot"],
        F["body"],
        C["muted"],
        line_spacing=1.25,
    )

    # ══ 하단: ③ 떨어진 14건 — 못 고침 13(넓게) vs 실익 1(좁게 accent) ══
    _subhead(slide, ax, 3.62, D.FULL_W, c["miss_head"])
    uw = 8.55  # 못 고침 블록 폭(넓게 = 13건)
    ux, uy, uh = ax, 4.04, 2.86
    add_box(slide, ux, uy, uw, uh, fill=C["bg_alt"], shape="round")
    add_text(
        slide,
        ux + 0.24,
        uy + 0.16,
        uw - 0.4,
        0.26,
        c["unfix_head"],
        S["body"],
        F["head"],
        C["muted"],
        bold=True,
    )
    D.hrule(slide, ux + 0.24, uy + 0.5, uw - 0.48, color=D.line_mid, weight=0.75)
    for i, (n, cause, detail) in enumerate(c["unfix"]):
        ry = uy + 0.62 + i * 0.56
        badge = add_box(slide, ux + 0.24, ry + 0.02, 0.46, 0.46, fill=C["muted"])
        bt = badge.text_frame
        bt.vertical_anchor = MSO_ANCHOR.MIDDLE
        bp = bt.paragraphs[0]
        bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run()
        br.text = n
        br.font.name, br.font.bold, br.font.size, br.font.color.rgb = (
            F["head"],
            True,
            Pt(S["body"]),
            C["bg"],
        )
        add_text(
            slide,
            ux + 0.86,
            ry,
            2.3,
            0.5,
            cause,
            S["body"],
            F["head"],
            C["primary"],
            bold=True,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        add_text(
            slide,
            ux + 3.3,
            ry,
            uw - 3.55,
            0.5,
            detail,
            S["caption"],
            F["body"],
            C["muted"],
            line_spacing=1.18,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    fx, fw = ux + uw + 0.2, D.FULL_W - uw - 0.2
    add_box(
        slide,
        fx,
        uy,
        fw,
        uh,
        fill=mix(C["bg"], C["accent"], 0.10),
        line=C["accent"],
        line_w=1.5,
        shape="round",
    )
    add_box(slide, fx, uy, 0.06, uh, fill=C["accent"])
    D.icon(slide, "target", "accent", fx + 0.24, uy + 0.2, 0.4)
    add_text(
        slide,
        fx + 0.74,
        uy + 0.22,
        fw - 0.9,
        0.3,
        c["fix_head"],
        S["body"],
        F["head"],
        C["accent"],
        bold=True,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        fx + 0.24,
        uy + 0.7,
        fw - 0.44,
        0.42,
        c["fix_lead"],
        S["caption"],
        F["body"],
        C["primary"],
        line_spacing=1.22,
    )
    D.hrule(
        slide, fx + 0.24, uy + 1.28, fw - 0.48, color=mix(C["accent"], C["bg"], 0.5), weight=0.75
    )
    for i, (head, desc) in enumerate(c["fix_points"]):
        py = uy + 1.42 + i * 0.74
        add_text(
            slide,
            fx + 0.24,
            py,
            fw - 0.44,
            0.24,
            head,
            S["caption"],
            F["head"],
            C["accent"],
            bold=True,
        )
        add_text(
            slide,
            fx + 0.24,
            py + 0.26,
            fw - 0.44,
            0.4,
            desc,
            S["foot"],
            F["body"],
            C["muted"],
            line_spacing=1.15,
        )

    D.source_line(slide, c["source"])
    return slide
