"""goldenfab.audit — 슬라이드 기계 오딧. **채점자가 캔 규칙을 코드로 제련해 두는 곳.**

## 왜 이 파일이 있나 (2026-07-15)

S4 한 장에 제3자 채점 2회 = **30분**이 들었다. 그런데 두 채점자가 낸 FAIL 9건을 분류하니
**7건이 순수 산수**였다 — 대비비·채움률·높이 종류·텍스트 중복·라벨 정렬·같은 fill 충돌·accent 계수.
15분짜리 모델이 필요한 일이 아니었다. 내 오딧이 그 규칙을 **안 갖고 있어서** 못 잡았을 뿐이다.

즉 **채점자는 검사기가 아니라 광산이다.** 한 번 캔 광석을 코드로 내려두면 다시 캘 필요가 없다.
design-rules의 철학("확정할 때마다 판단을 박제")과 같되, **산문이 아니라 실행되는 코드로** 박제한다
— 산문 규칙은 잊히지만 assert는 건너뛸 수 없다(P1 주석의 자기 인정).

## 운영 규칙

1. 레이아웃을 만들면 먼저 이 오딧을 통과시킨다(2초).
2. 제3자 채점은 **1회만**. 재채점 금지 — 오늘 실증: 6/10 → 5/10으로 **안 수렴**했고 두 채점자가
   **서로 다른 항목**을 찾았다. 주관적 채점에는 고정점이 없다.
3. 채점자가 새 결함을 찾으면 **반드시 여기에 규칙으로 내린다**. 산문으로 남기면 다음 장에서 또 캔다.
4. 회를 거듭할수록 채점자가 찾을 게 줄어든다. 그게 이 파일이 하는 일이다.

## 규칙 출처

각 규칙에 발견 경위를 적는다 — 왜 있는지 모르는 규칙은 다음 사람이 지운다.
"""

import re

from pptx.enum.shapes import MSO_SHAPE

EMU = 914400
AIR_MIN = 0.12  # §6 공기 하한
ACCENT_MAX = 4  # §2 accent 상한
ACCENT_MAX_DENSE = 6  # dense accent 상한(채움·선만, 런 제외) — 압축 헤더 룰 등 여유
CONTENT_BOTTOM = 6.35  # P2③ (sparse 골든 본문 하한 — 결론 바 6.60 위)
CONTENT_BOTTOM_DENSE = 7.1  # dense 본문 하한 — 바 없이 출처선(dense.SOURCE_Y 7.14) 바로 위까지
RIGHT_EDGE = 12.733  # P2④
CONTRAST_MIN = 4.5  # WCAG AA (일반 텍스트)

# ── legacy 프레임 상수 (2026-07-25 실측 · PROFILES "legacy") ──────────────────
# **2026-07-29부터 이 프로파일은 대상이 0이다** — 겨냥하던 렌더러가 아카이브됐다. 아래는
# "왜 프레임마다 임계가 달라야 하는가"의 실측 근거로 남긴다(이하 2026-07-25 당시 기록).
# legacy 렌더러(chart·표·문서형·다이어그램)는 골든과 **다른 프레임**을 썼다:
# 하단이 배너(top 6.58)·출처선(6.65)·각주(6.88)·푸터 헤어라인(7.0)·푸터 텍스트(7.08)로 층층이 차
# 있다. 골든 하한(6.35·7.1)을 그대로 쓰면 111/111 적색이 되고, 그건 이 프로젝트가 이미 겪은
# "상시 적색은 무시당하고 우회된다" 경로다. 그래서 **규칙은 하나, 임계만 프레임별**로 둔다.
#
# 하한을 **푸터 헤어라인(7.0)**으로 잡은 근거(legacy 산출물 111장 전수 실측): 본문 도형의 최대
# bottom이 정확히 7.0이었고(`_stat_split` 캡션 상자), 7.0 아래는 전부 푸터 텍스트 영역이라
# 그 아래로 내려간 본문은 푸터와 실제로 겹친다. 더 위로(6.55) 올리면 spectrum 단계 설명(6.9)·
# stat_split 캡션이 오탐으로 걸린다 — 상자에 여유를 두는 legacy idiom이지 침범이 아니다.
# **경계(정직):** 그래서 이 검사는 legacy에서 "푸터 침범"까지만 본다. 본문이 출처선(6.65)을
# 덮는 경우는 여기서 안 걸리고 `check_text_collision`(글자 실점유 비교)이 진다.
CONTENT_BOTTOM_LEGACY = 7.0
LEGACY_FRAME_TOPS = (6.88, 6.95, 7.0, 7.08)
ACCENT_MAX_LEGACY = 8  # 실측 상한(harvey_table 점수볼 = 채움 4 + 선 4). 런은 legacy도 제외.


# ── 픽셀·색 산수 ──────────────────────────────────────────────


def _lum(hexstr):
    """상대 휘도 (WCAG)."""

    def ch(v):
        v /= 255
        return v / 12.92 if v <= 0.03928 else ((v + 0.055) / 1.055) ** 2.4

    r, g, b = (int(hexstr[i : i + 2], 16) for i in (0, 2, 4))
    return 0.2126 * ch(r) + 0.7152 * ch(g) + 0.0722 * ch(b)


def contrast(fg, bg="FFFFFF"):
    """대비비. 판정 단어가 산문보다 흐린지 재는 데 쓴다."""
    a, b = _lum(str(fg)), _lum(str(bg))
    hi, lo = max(a, b), min(a, b)
    return round((hi + 0.05) / (lo + 0.05), 2)


def shape_kind(sh):
    try:
        return sh.auto_shape_type
    except Exception:
        return None


# ── 색 읽기는 **읽기 전용**이어야 한다 (2026-07-25 버그 수리) ────────────────────
# 왜: 초판은 `sh.line.color.type` · `run.font.color.rgb`로 색을 읽었는데 python-pptx의 ColorFormat은
# **접근하는 순간 XML을 만든다** — `<a:ln><a:solidFill/></a:ln>`(빈 solidFill)이 붙고,
# PowerPoint는 그걸 **검정 헤어라인 테두리**로 그린다. 빌드 시 게이트(`_gate_density`)는 `prs.save()`
# **전에** 도니 그 테두리가 파일에 그대로 저장됐다 — 즉 **오딧이 산출물을 오염시켰다.**
# legacy 배선 후 fixture 렌더 눈검증에서 전 텍스트 상자에 테두리가 생겨 발각됐다(게이트는 전부 green).
# 교훈: 검사기는 대상을 절대 변형하지 않는다. 아래는 lxml로 XML을 **읽기만** 한다.
_NS_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def _solid_srgb(parent):
    """`<parent><a:solidFill><a:srgbClr val="RRGGBB"/>` → "RRGGBB". 없으면 None(생성하지 않음).

    schemeClr(테마색)·gradFill·pattFill은 None — 초판도 `.rgb` 예외로 None이었다(동등).
    """
    if parent is None:
        return None
    sf = parent.find(f"{_NS_A}solidFill")
    if sf is None:
        return None
    clr = sf.find(f"{_NS_A}srgbClr")
    val = clr.get("val") if clr is not None else None
    return val.upper() if val else None


def _spPr(sh):
    el = getattr(sh, "_element", None)
    return getattr(el, "spPr", None) if el is not None else None


def fill_hex(sh):
    return _solid_srgb(_spPr(sh))


def line_hex(sh):
    spPr = _spPr(sh)
    return _solid_srgb(spPr.find(f"{_NS_A}ln") if spPr is not None else None)


def runs_of(sh):
    """(텍스트, hex색, pt, bold) 목록 — **rPr XML을 읽기만** 한다(위 오염 주석 참조).

    `run.font`는 `get_or_add_rPr()`를 호출하고 `font.color`는 빈 `<a:solidFill/>`을 넣는다.
    텍스트 색에 빈 solidFill이 박히면 상속 색이 끊긴다 — 검사기가 색을 바꿀 수는 없다.
    """
    out = []
    if not sh.has_text_frame:
        return out
    for p in sh.text_frame.paragraphs:
        for r in p.runs:
            rPr = r._r.find(f"{_NS_A}rPr")
            sz = rPr.get("sz") if rPr is not None else None
            b = rPr.get("b") if rPr is not None else None
            out.append(
                (
                    r.text,
                    _solid_srgb(rPr),
                    int(sz) / 100 if sz else None,
                    None if b is None else b in ("1", "true"),
                )
            )
    return out


def box(sh):
    """도형의 **화면상** 경계 상자 (x, y, w, h) 인치 — 회전을 반영한다.

    왜 회전을 보나(2026-07-25 legacy 배선 실측): legacy 프레임의 세로 저작권선은
    `add_textbox(10.73, 3.55, 4.6, 0.28)` + `rotation=270`이다. python-pptx의 left/top/width/
    height는 **회전 전** 값이라 그대로 쓰면 폭 4.6"의 가로 띠가 화면 중앙을 지나는 것으로
    보인다 — 그래서 mpl 익스히빗 10장이 `check_picture_overlap`에 "저작권선 침범"으로 걸렸다
    (전부 오탐). 회전 도형이 하나만 있어도 기하 검사 전체가 오염되므로 여기 한 곳에서 바로잡는다.
    90·270도(가로세로 교환)만 정확히 처리하고, 그 외 각도는 회전 전 상자를 그대로 쓴다 —
    임의 각도의 AABB는 더 커져서 오탐을 만들고, 이 하네스는 90도 배수만 쓴다.
    """
    x, y = sh.left / EMU, sh.top / EMU
    w, h = sh.width / EMU, sh.height / EMU
    rot = round(getattr(sh, "rotation", 0) or 0) % 360
    if rot in (90, 270):
        cx, cy = x + w / 2, y + h / 2
        w, h = h, w
        x, y = cx - w / 2, cy - h / 2
    return (x, y, w, h)


# ── 규칙 ──────────────────────────────────────────────────────


def check_accent(shapes, accent, cap=ACCENT_MAX, count_runs=True):
    """§2 accent 상한. **채움만 세면 안 된다 — 테두리·커넥터 선도 accent 예산이다.**

    출처: v4에서 채움·런만 세서 4로 통과시켰는데 제3자가 7군데를 셌다(선 미계수).

    **count_runs(2026-07-20 dense 인지):** dense 장은 2단 불릿의 **강조**를 accent 텍스트 런으로
    쓴다(정상 표현, 장당 28~31런). 그건 '장식적 accent 남용'(이 규칙이 막는 것)이 아니므로
    dense에선 count_runs=False로 채움·선만 센다. sparse 골든은 기존대로 런까지 센다(기본 True).
    """
    f = [s for s in shapes if fill_hex(s) == accent]
    ln = [s for s in shapes if line_hex(s) == accent]
    rn = [r for s in shapes for r in runs_of(s) if r[1] == accent] if count_runs else []
    tot = len(f) + len(ln) + len(rn)
    ok = tot <= cap
    run_note = f" + 런 {len(rn)}" if count_runs else " (런 제외·dense)"
    return ok, f"accent 채움 {len(f)} + 선 {len(ln)}{run_note} = {tot} (상한 {cap})", tot


def check_verdict_contrast(shapes, prose_pt):
    """**판정을 말하는 단어는 그것을 설명하는 산문보다 흐릴 수 없다.**

    출처: v5에서 결론 단어 '↩ 되돌아옴'(3.22) · '✕ 불가역'(3.41)이 그것을 설명하는
    산문(8.47)보다 2.5배 흐렸다. 단순 명찰은 16.16을 가져갔다. 제3자 2회 모두 지적.
    """
    bad = []
    for s in shapes:
        for text, col, pt, _bold in runs_of(s):
            t = text.strip()
            if not col or not t:
                continue
            # 판정 단어는 **짧은 라벨**이다. 기호로 시작하거나(✕ 불가역 · ↩ 되돌아옴), 짧으면서
            # 판정어로 끝난다. 길이 상한이 없으면 산문 속 "없다"까지 걸린다 — 2026-07-15
            # tech_tree 실측: '문단 뭉치에는 순서도 관계도 없다'(17자, 서사 본문)를 판정으로
            # 오탐했다. 서사는 muted가 정상이다(§4 위계 4단). 규칙을 느슨하게 한 게 아니라
            # **판정과 산문을 구분**하도록 정확하게 했다 — 기호 시작 케이스는 그대로 잡힌다.
            is_verdict = bool(re.match(r"^\s*[✕✓↩→]", t)) or (
                len(t) <= 12 and ("불가역" in t or t.endswith("없다"))
            )
            if is_verdict:
                cr = contrast(col)
                if cr < CONTRAST_MIN:
                    bad.append((t[:18], col, cr))
    return not bad, f"판정 단어 대비 < {CONTRAST_MIN}: {len(bad)} {bad[:3]}", len(bad)


def check_fill_ratio(shapes, min_ratio=0.30, min_w=1.5):
    """도형이 글자에 비해 과하게 넓지 않은가(죽은 회색).

    출처: v5의 한계 칩이 폭 3.30"에 채움률 14~22%였다 — 단어와 커넥터 사이 238px 공백.
    글자폭은 pt·글자수로 근사(정밀 측정 불가) — 한국어 1자 ≈ 0.8×pt.
    """
    bad = []
    for s in shapes:
        if not s.has_text_frame or fill_hex(s) is None:
            continue
        _x, _y, w, _h = box(s)
        if w < min_w:
            continue
        rs = runs_of(s)
        if not rs:
            continue
        txt = "".join(r[0] for r in rs)
        pt = max((r[2] or 9) for r in rs)
        est = len(txt.strip()) * pt * 0.8 / 72
        ratio = est / w
        if ratio < min_ratio:
            bad.append((txt.strip()[:14], round(w, 2), f"{ratio:.0%}"))
    return not bad, f"채움률 < {min_ratio:.0%}: {len(bad)} {bad[:3]}", len(bad)


H_VOCAB = 1.5  # 이보다 높은 카드는 좌변 바만으로 통과할 수 없다 — 아래 (2026-08-02) 참고


def check_adhoc_card(shapes, w_min=2.4, h_min=0.9):
    """즉흥 밋밋 카드 금지 — 카드 크기 박스가 **구조 자식**(아이콘·배지·구분선·중첩 박스) 0이면 FAIL.

    출처(2026-07-21): dense.py에 '카드=hero_card 단일 · 장별 즉흥 카드 금지'(design-rules §8)가
    **prose로만** 있어 반복 무시됐다 — s16 한계 카드를 add_box+텍스트로 손으로 밋밋하게 그려
    글자수·채움률·넘침 게이트를 전부 통과시켰다(사용자 반려: "밋밋하고 쓰레기"). hero_card는
    배지(oval)·히어로 아이콘(picture)·배너·구분선(rule)·칩(중첩 박스)을 **항상** 낳는다 —
    그 구조 자식이 하나도 없는 카드 크기 박스 = 즉흥 밋밋 카드. '리치·의미 정합'은 기계화 못
    하지만 '구조 0'은 여기서 막는다(그 위는 눈검증 몫). dense 전용(sparse 골든 카드 idiom은 다름).

    문턱 하향 3.0x1.5 → **2.4x0.9**(2026-07-25 래칫). 아래 실측에 나오는 (B) 어휘
    (flow·layers·branch·cards·from_to)는 **그 후 전부 폐기**됐다 — 문턱을 이 값으로 정한
    근거로만 보존한다(지우면 왜 2.4x0.9인지 다시 실측해야 한다).
    초판 문턱이 **카드 크기**만 봐서
    다이어그램 노드가 전부 눈 밖이었다 — 다이어그램 18종 전면 반려 당시 실측으로 `flow`
    노드(2.66x1.4)·`layers`(12.13x1.0)·`branch` 자식(7.63x1.24)·`from_to`(4.69x1.15)가
    **같은 결함(구조 0)인데 미달로 통과**하고 `cards`(5.92x1.7) 한 장만 걸렸다. 패턴이
    같으면 문턱도 같아야 한다. 오탐 실측(하향 승인 근거): dense 장 모듈 10/10 · 다이어그램
    18장 전수 · 개선 4장 = **0건**. 더 낮추면(2.4x0.6) s14_dense 1 · s16_dense 2가 오탐으로
    걸린다 — 그게 하한이 여기인 이유다(정밀 > 양: 오탐 게이트는 백스톱에 음소거된다).

    **판정 상향 — 좌변 바는 구조로 치지 않는다(2026-08-02).** 초판은 구조 자식을 한 덩어리로
    셌고, `sh < 0.06 or sw < 0.08`가 얇은 바를 "구분선"으로 인정했다. local-ai-assist에서
    카드 14장 전부가 **accent 좌변 바(0.055") 하나로 통과**했고(어휘는 0), FAIL 메시지가
    `hero_card 미사용`이라 적혀 있어 그 통과가 "hero_card 안 써도 괜찮다"로 오독됐다.
    그래서 구조 자식을 **어휘**(아이콘 PICTURE · 배지 OVAL · 중첩 박스)와 **바**로 나누고,
    `ch >= H_VOCAB(1.5)`인 카드는 어휘 0이면 FAIL로 잡는다. 메시지도 실물에 맞췄다.
    `ch < 1.5`에는 초판 판정(구조 자식 전무)을 그대로 남긴다 — 다이어그램 노드가 그 대역이다.

    오탐 실측(상향 승인 근거, 통과작 전수 152장): golden-deck 30 · golden-pilot_s06 11덱 45 ·
    probe 8덱 42 · results-검토 3덱 35 = **0건**. `ch >= 1.5` 단서가 오탐 0의 핵심이다 —
    없으면 골든 s16 다이어그램 노드(4.50x1.38 · 바 1개)가 걸린다. 대가로 local의 2.74x1.13
    카드 1장은 빠져나간다(13/14 검출).
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    def center_in(inner, ox, oy, ow, oh):
        ix, iy, iw, ih = inner
        cx, cy = ix + iw / 2, iy + ih / 2
        return ox <= cx <= ox + ow and oy <= cy <= oy + oh

    bad = []
    for c in shapes:
        if shape_kind(c) not in (MSO_SHAPE.RECTANGLE, MSO_SHAPE.ROUNDED_RECTANGLE):
            continue
        if fill_hex(c) is None and line_hex(c) is None:
            continue
        cx, cy, cw, ch = box(c)
        if cw < w_min or ch < h_min:
            continue
        if cw > 13 and ch > 7:
            continue  # 풀블리드 배경
        fh = fill_hex(c)
        if fh is not None and _lum(fh) < 0.45:
            continue  # 어두운 패널(코드/JSON/스키마/터미널 아티팩트) — 카드 아님(P0 #1 실물)
        c_area = cw * ch
        has_text = bool(c.has_text_frame and c.text_frame.text.strip())
        n_vocab = 0  # 어휘 — 아이콘·배지·중첩 박스(배너·칩)
        n_bar = 0  # 얇은 선·좌변 바 — 구조로 치지 않는다(아래 근거)
        for s in shapes:
            if s._element is c._element:
                continue
            sx, sy, sw, sh = box(s)
            if not center_in((sx, sy, sw, sh), cx, cy, cw, ch):
                continue
            if s.has_text_frame and s.text_frame.text.strip():
                has_text = True
            if s.shape_type == MSO_SHAPE_TYPE.PICTURE:
                n_vocab += 1  # 아이콘
            elif shape_kind(s) == MSO_SHAPE.OVAL:
                n_vocab += 1  # 배지
            elif line_hex(s) is not None or fill_hex(s) is not None:
                if sh < 0.06 or sw < 0.08:
                    n_bar += 1  # 구분선·세로 바
                elif shape_kind(s) in (MSO_SHAPE.RECTANGLE, MSO_SHAPE.ROUNDED_RECTANGLE) and (
                    sw * sh < c_area * 0.7
                ):
                    n_vocab += 1  # 중첩 박스(배너·칩·내부 패널)
        if has_text and (n_vocab + n_bar == 0 or (ch >= H_VOCAB and n_vocab == 0)):
            lbl = c.text_frame.text.strip()[:14] if c.has_text_frame else "(카드)"
            bad.append((lbl, round(cw, 1), round(ch, 1)))
    return (
        not bad,
        f"카드 어휘 부재(아이콘·배지·중첩 박스 0 · 좌변 바만): {len(bad)} {bad[:3]}",
        len(bad),
    )


def check_duplicate_nodes(shapes, allow=0):
    """**구별 장치 없는** 반복 = 재탕(P4③).

    출처: v4에서 '2종 · 허위 확정'이 두 밴드에 같은 글자·같은 fill·같은 테두리로 있었다.

    ⚠ 테두리(색·점선)가 다르면 재탕이 아니라 **대조**다 — v7의 두 ◇ '근거 충분?'은 실선(성립)
    vs 점선(성립 불가)으로 뒤가 앞을 부정한다. 제3자도 "반복이 아니라 대조"로 PASS를 줬는데
    초판 규칙이 (텍스트, 채움)만 봐서 오탐을 냈다. 구별 장치까지 서명에 넣는다.

    `allow`: 같은 노드의 재등장이 **그 자체로 메시지**인 레이아웃을 위한 기대값(기본 0).
    tech_tree 카탈로그가 그렇다 — `변동대가`가 hier·cp의 출발이자 term의 도착으로 3번 나오는
    건 "한 개념에 세 종류의 간선이 붙는다"를 보이는 것이다. 규칙이 잡으려는 건 **정보량 0인**
    반복이고 이건 정보량이 있다. 다만 예외를 여는 게 아니라 **수를 박는다** — 실측보다 늘면
    FAIL이라, 무심코 하나 더 그리면 잡힌다.
    """
    seen = set()
    dup = []
    for s in shapes:
        f = fill_hex(s)
        if not f or not s.has_text_frame:
            continue
        t = s.text_frame.text.strip()
        if len(t) < 3:
            continue
        try:
            dash = str(s.line.dash_style)
        except Exception:
            dash = None
        sig = (t, f, line_hex(s), dash)  # 구별 장치 = 테두리 색 + 점선 여부
        if sig in seen:
            dup.append(t[:20])
        seen.add(sig)
    return len(dup) <= allow, f"구별 장치 없는 반복: {len(dup)} (허용 {allow}) {dup[:3]}", len(dup)


def check_node_class(shapes, min_len=2):
    """같은 이름의 노드가 한 장 안에서 **다른 스타일**로 그려지면 FAIL (P4⑩의 쌍둥이).

    출처: 2026-07-15 제3자 채점이 s9에 2/10을 줬고 최대 결함이 이것이었다 — 좌 트리는
    `bg_alt` 채움을 **내부 노드**(개념·문단)에, 우 카탈로그는 같은 `bg_alt`를 **사례**에 썼다.
    독자가 왼쪽에서 배운 "회색 = 내부 노드"가 오른쪽에서 "회색 = 사례"로 뒤집힌다.
    `변동대가`는 한 장에서 회색 덩어리이자 흰 바탕 검정 테두리 볼드였다.

    ⚠ `check_duplicate_nodes`의 **정확한 반대**다. 저건 "같은 이름 + 같은 서명"의 무의미한
    반복을 잡고, 이건 "같은 이름 + **다른** 서명"의 모순을 잡는다. 둘 다 있어야 한다 —
    duplicate만 있으면 "스타일을 바꾸면 통과"라는 잘못된 탈출구가 열린다.

    ⚠ **채움만 본다.** 초판이 (채움, 테두리)를 서명으로 삼았더니 S4의 두 ◇ `근거 충분?`을
    오탐했다 — 그건 실선(성립) vs 점선(성립 불가)으로 **뒤가 앞을 부정하는 의도된 대조**이고,
    제3자도 "반복이 아니라 대조"로 PASS를 준 설계다. 그대로 뒀으면 duplicate 규칙("구별 장치를
    둬라")과 이 규칙("같은 스타일이어야")이 **서로 반대를 요구**해 S4가 어느 쪽으로도 못 빠져
    나갔다. 경계는 이렇다:
        **채움 = 클래스**(이게 무엇인가 — 개념·문단·사례) → 같은 이름이면 같아야 한다.
        **테두리·점선 = 상태**(성립/불성립·자동/위임) → 달라도 된다, 그게 대조다.
    s9의 실제 결함은 채움이 갈린 것(`bg_alt` 좌 = 내부 노드 / 우 = 사례)이라 이 경계로 잡힌다.

    원인은 늘 구조적이다: 노드를 그리는 함수가 둘이면 색이 갈린다. 단일 출처(NODE_STYLE)로
    합치는 게 근본 수정이고, 이 규칙은 그게 다시 갈라지는지 지킨다.
    """
    sigs = {}
    bad = []
    for s in shapes:
        f = fill_hex(s)
        if not f or not s.has_text_frame:
            continue
        t = s.text_frame.text.strip()
        if len(t) < min_len:
            continue
        if t in sigs and sigs[t] != f:
            bad.append((t[:18], sigs[t], f))
        else:
            sigs.setdefault(t, f)
    return not bad, f"같은 노드 다른 채움: {len(bad)} {[b[0] for b in bad][:3]}", len(bad)


def text_need_height(sh):
    """이 상자의 텍스트를 다 그리는 데 필요한 높이(인치) 근사 — 넘침·겹침 검사의 단일 출처.

    줄당 글자수 ≈ 상자폭 / (pt·CW/72), 필요 줄 수 = ceil(글자수/줄당),
    필요 높이 = Σ(줄수·pt·1.2·행간/72). **근사 한계:** 글자폭은 혼합 텍스트(한글·숫자·공백)
    평균 ≈ 0.62×pt(순 한글 1.0, 숫자·공백 0.4~0.5)라 순 한글 장문은 과소추정된다.
    두 검사가 각자 계산하면 한쪽만 고쳐져 갈라지므로 여기 한 곳에 둔다.
    """
    w = sh.width / EMU
    need = 0.0
    for para in sh.text_frame.paragraphs:
        ptxt = "".join(r.text for r in para.runs)
        if not ptxt:
            need += 0.12  # 빈 문단도 한 줄 차지
            continue
        pts = [r.font.size.pt for r in para.runs if r.font.size]
        pt = max(pts) if pts else 12
        cw = pt * 0.62 / 72  # 혼합 텍스트 평균 글자폭
        per_line = max(1, int((w - 0.04) / cw))
        lines = max(1, -(-len(ptxt) // per_line))  # ceil
        ls = para.line_spacing if isinstance(para.line_spacing, (int, float)) else 1.0
        need += lines * pt * 1.2 * ls / 72
    return need


def check_text_overflow(shapes, tol=1.5):
    """텍스트가 제 상자 높이를 넘겨 흘러나오나(세로 오버플로우) — 근사 추정(2026-07-21 신설).

    add_text는 word_wrap=True라 가로는 감기지만, 감긴 줄 수가 상자 높이를 넘으면 **아래로 흘러**
    이웃 요소와 겹친다(s15 fix 블록 사고). 렌더 없이 잡으려고 `text_need_height`로 근사하고,
    필요 높이가 상자 높이×tol을 넘으면 넘침. tol·CW는 오탐(짧은 라벨)과 미탐(경계) 사이 타협 —
    총체적 겹침이 아니라 **단일 상자 세로 초과**만 잡는 하한 대리지표다(눈검증은 여전히 필요).

    **폭 계수를 못 올리는 이유(2026-07-25 실측 · 통과작 163장 분모).** 한글 미탐(순 한글 42자가
    5줄로 흘러 상자를 넘겼는데 미탐)을 고치려 계수를 올려 봤다:
    0.75→오탐 6장 · 0.85→7장 · 1.00→12장 · 전각/반각 구분 + 어절 그리디(가장 정확)→5장.
    전부 통과작(골든 부제·작은 캡션 상자·1자 마크)이다. 게다가 **오탐과 진짜 결함이 같은 비율**에
    있다 — legacy 28pt 헤드라인 2줄(결함, 0.93/0.52=1.79×)과 골든 부제 2줄(정상, 0.93/0.55=1.69×)은
    비율로 못 가른다. 둘을 가르는 건 "아래에 무엇이 있나"이고 그건 이 규칙이 아니라
    `check_text_collision`의 일이다. → **계수·tol 그대로 둔다.** 올리려면 통과작 오탐 처리가 선행.
    """
    bad = []
    for s in shapes:
        if not s.has_text_frame:
            continue
        tf = s.text_frame
        if not tf.text.strip():
            continue
        w, h = s.width / EMU, s.height / EMU
        if w < 0.3 or h < 0.05:
            continue
        need = text_need_height(s)
        if need > h * tol:
            bad.append((tf.text.strip()[:12], round(need, 2), round(h, 2)))
    return not bad, f"텍스트 상자 넘침(추정) {len(bad)} {bad[:3]}", len(bad)


def check_picture_overlap(shapes):
    """그림이 다른 요소를 침범하지 않는가.

    출처(2026-07-15 사용자 버그 리포트): s10이 이미지 비율을 `3200/2000`(1.60)으로 **가정**하고
    폭 7.6"를 전제로 옆 칼럼을 8.6"에 놨는데, 실제 파일은 1899×926(**2.05**)이라 폭 9.74"로
    렌더돼 텍스트를 **1.74" 침범**했다. 그 `img_w`는 계산만 하고 안 쓰는 죽은 변수였다.
    → `kit.fit_picture`가 근본 수정(박스에 맞춤). 이 규칙은 재발 감시.

    ⚠ python-pptx는 접근할 때마다 **새 프록시**를 만들어 `s is p`로는 자기 자신도 못 거른다 —
    `_element` 동일성으로 비교해야 한다(이 함정에 한 번 빠져 "그림이 자기와 겹침"을 냈다).
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    pics = [s for s in shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    hits = []
    for p in pics:
        x, y, w, h = box(p)
        if w < 1.0 and h < 1.0:
            continue  # 아이콘(콘텐츠 그림 아님) — 라벨과 나란히·겹쳐 놓이는 게 정상(dense 칩)
        for s in shapes:
            if s._element is p._element:
                continue
            sx, sy, sw, sh = box(s)
            if sw > 13 and sh > 7:
                continue  # 풀블리드 배경
            if sx < x + w and x < sx + sw and sy < y + h and y < sy + sh:
                lbl = s.text_frame.text.strip()[:16] if s.has_text_frame else "(무텍스트)"
                hits.append(lbl)
    return not hits, f"PICTURE {len(pics)}개 · 침범 {len(hits)} {hits[:3]}", len(hits)


def check_text_collision(
    shapes, min_depth=0.03, cover=0.6, wsim=0.55, small_frac=0.5, contain_frac=0.7
):
    """**배경 없는 순수 텍스트** 두 상자가 같은 열에서 세로로 포개지나 — 글자가 판독 불가로 겹침.

    출처(2026-07-24 카드 목업 3회 반려): evidence_card 함의줄이 하단 출처줄 위에 겹침, hero_card
    불릿이 세로로 뭉개져 서로 겹침, 출처줄과 인용 블록 겹침. 사용자가 **눈으로 3번** 잡았다.
    기존 오딧의 사각지대였다 — check_text_overflow는 '제 상자 세로 초과'(단일 박스)만,
    check_picture_overlap은 '그림 대 요소'만 본다. **텍스트 요소끼리의 겹침**은 아무도 안 봤다.
    이 규칙은 check_picture_overlap의 형제다: 저건 그림이 침범, 이건 텍스트가 텍스트를 침범.

    역할 분담(단일 출처의 두 반쪽): check_text_overflow(제 상자 초과·경미한 흘러내림·tol 대리지표)
    + 이 규칙(배치 충돌·글자 실점유 세로 포갬). 초판은 "렌더 추정을 안 쓴다 — 순수 기하라 추정
    오탐 0"이었으나 **틀렸다**(2026-07-25 legacy 실측 17장 오탐): 순수 상자 기하는 "상자를 넉넉히
    잡는 idiom"을 겹침으로 오판한다. 이제 아래 (5)로 `text_need_height` 근사를 **상한 클램프**로만
    쓴다 — 추정이 커지는 방향으로는 판정에 못 끼어들어 추정 오탐이 새로 생기지 않는다.

    오탐 경계 — 통과작 전수 실측(2026-07-24)으로 좁힌 네 축(핵심 난제):
    (1) **배경 있는 도형 제외**(fill_hex 있으면 스킵): 판정 마름모 노드·아이콘 칩·배지·bg 패널은
        배경색으로 구분돼 겹쳐도 판독되고 의도된 배치다(다이어그램 노드 근접, 배지 위 라벨).
        실측: s04·problem_grid 라벨 대 노드가 여기서 걸러진다. 결함 3건은 **전부 fill 없는
        add_text**끼리였다 — 배경 없는 글자만 포개면 판독 불가.
    (2) **세로 파일업만**(ox >= cover x 좁은 폭): 두 상자가 같은 열에 겹쳐(가로 폭 대부분 공유)
        세로로 파고들 때만. 마크·번호가 라벨과 **가로로 나란한** 한 줄 인접('01 큰 노드=')은
        ox가 좁아 제외 — 결함은 문장이 **아래로** 흐르는 세로 겹침이다(screenshot 번호 대 설명).
    (3) **같은 폭 컬럼끼리만**(width_sim = 좁폭/넓폭 >= wsim): 결함은 카드 폭 두 문장이 세로로
        포개짐(폭비 ~1). 넓은 캡션·부제(w12·3.78) 밑에 **좁은 간선/판정 라벨**(w1.0~1.8)이
        얹히는 다이어그램 배치(폭비 <=0.47)는 본문 충돌이 아니라 라벨-오버-영역이라 제외한다.
        실측: tech_tree 6건('지식그래프…'x'hier 79' 등)·s16 '✕'x'실측 4건'(0.32)이 전부 여기서.
    (4) 포함(작은 라벨이 큰 텍스트 박스 안): 면적비 < small_frac AND 교집합 >= contain_frac x
        작은 면적이면 제외 — 비슷한 크기 포갬(불릿 뭉갬)은 포함이 아니라 충돌이다.

    (5) **글자가 실제로 차지하는 높이까지만**(h = min(상자 높이, text_need_height)): 상자를 넉넉히
        잡고 텍스트를 위에 붙이는 idiom(legacy `bullets` 상자 0.5"에 피치 0.46" → 0.04" 상자 겹침,
        `_stat_split` 40pt 숫자 상자 1.4"에 실제 0.67")에서 **글자는 안 겹치는데 상자만 겹친다**.
        2026-07-25 legacy 111장 실측: 이 축이 없으면 17장이 오탐(전부 상자 여유 0.04~0.4")이었다.
        상자가 글자로 꽉 찬 경우(need ≥ h)엔 종전과 동일해 검출력이 그대로다(주입 프로브:
        꽉 찬 상자 0.05~0.25" 포갬 전건 검출. 이를 상시 확인하던 `test_wiring` ④.5는
        2026-07-29 아카이브 — 검출력 회귀는 지금 아무도 안 지킨다).
        **미탐 경계(정직):** `text_need_height`는 혼합 텍스트 평균(0.62×pt)이라 **순 한글 장문은
        과소추정**한다 — 상자 여유가 큰 자리에서 그만큼 얕은 겹침을 놓칠 수 있다. 근사를 정확하게
        (전각 1.0/반각 0.5 + 어절 그리디) 키우면 legacy 통과작 1건이 오탐이 된다(2026-07-25 실측,
        163장 분모) → **오탐 0을 택했다**(오탐 게이트는 백스톱에 음소거된다).

    경계 한계(정직): '배경 없는·같은 폭·세로 파일업'까지만 기계로 막는다. 배경 있는 요소 위 텍스트
    겹침·좁은 라벨 대 넓은 본문·가로 충돌은 여기서 안 잡힌다(눈검증 몫) — 정밀(오탐 0) > 양.
    """
    tb = []
    for s in shapes:
        if not s.has_text_frame or not s.text_frame.text.strip():
            continue
        if fill_hex(s) is not None:
            continue  # 배경 있는 도형(노드·칩·배지·패널) — 순수 텍스트 요소 아님(1)
        x, y, w, h = box(s)
        if w > 13 and h > 7:
            continue  # 풀블리드 배경
        h = min(h, max(text_need_height(s), 0.05))  # 글자 실점유까지만(5)
        tb.append((s, x, y, w, h))
    bad = []
    for i in range(len(tb)):
        _s1, ax, ay, aw, ah = tb[i]
        for j in range(i + 1, len(tb)):
            _s2, bx, by, bw, bh = tb[j]
            ox = min(ax + aw, bx + bw) - max(ax, bx)
            oy = min(ay + ah, by + bh) - max(ay, by)
            if ox <= 0 or oy <= 0:
                continue
            if oy < min_depth:
                continue  # 세로 침투 없음(접점·정확 타일링·EMU 반올림)
            if ox < cover * min(aw, bw):
                continue  # 가로로 나란함(마크+라벨·번호+설명) — 세로 파일업 아님(2)
            if min(aw, bw) < wsim * max(aw, bw):
                continue  # 폭이 크게 다름 — 좁은 라벨이 넓은 캡션/영역 밑에 얹힘(3, 다이어그램)
            small, large = min(aw * ah, bw * bh), max(aw * ah, bw * bh)
            if large > 0 and small / large < small_frac and ox * oy >= contain_frac * small:
                continue  # 작은 라벨이 큰 텍스트 박스 안(4 포함)
            t1 = tb[i][0].text_frame.text.strip()[:10]
            t2 = tb[j][0].text_frame.text.strip()[:10]
            bad.append((t1, t2, round(oy, 3)))
    return not bad, f"텍스트 상자 세로 겹침: {len(bad)} {bad[:3]}", len(bad)


def check_bounds(shapes, bottom=CONTENT_BOTTOM, right=RIGHT_EDGE, skip_tops=()):
    """P2③ 콘텐츠 하한 · P2④ 풀폭 우변."""
    over, badr = [], []
    for s in shapes:
        x, y, w, h = box(s)
        if w > 13 and h > 7:
            continue  # 풀블리드 배경
        if any(abs(y - t) < 0.01 for t in skip_tops):
            continue
        if y + h > bottom + 0.01:
            over.append((round(y + h, 2), s.text_frame.text[:14] if s.has_text_frame else ""))
        if abs(x - 0.6) < 0.001 and w > 11 and abs(x + w - right) > 0.005:
            badr.append(round(x + w, 3))
    return (
        not (over or badr),
        f"bottom>{bottom}: {len(over)} {over[:2]} · 풀폭 right 위반: {badr}",
        len(over) + len(badr),
    )


DENSITY_EXEMPT = {"screenshot"}

# 골든 스냅샷 17장의 장 종류 키 — assets/golden-snapshot.json의 slides 배열과 같은 순서.
# 2026-07-29 goldenfab/reference.py(골든 생산 코드)를 아카이브하며 그 SLIDE_ORDER에서 옮겼다.
# 스냅샷과 한 몸이므로 스냅샷을 갈면 이것도 같이 간다.
SNAPSHOT_KEYS = [
    "cover",
    "toc",
    "part",
    "problem_grid",
    "part",
    "exec_graph",
    "part",
    "tech_evidence",
    "tech_tree",
    "screenshot",
    "tech_mechanism",
    "tech_capture",
    "part",
    "ab_simulation",
    "validation",
    "boundary",
    "closing",
]


def density_band(snapshot_path=None):
    """골든 스냅샷에서 본문 장의 밀도 하한을 파생 — **텍스트 글자수가 주 지표**(2026-07-20).

    출처(2026-07-16 배선 재설계): 골든이 "복제 템플릿"에서 "변형 가능한 출발점"이 되면서
    스냅샷 회귀가 실전 덱을 못 지킨다 — 변형·신규 장의 "밀집화"를 재는 기계 게이트가 이것.
    임계는 리터럴이 아니라 **골든의 가장 성긴 본문 장**에서 파생한다(§4 범용 해결) —
    골든이 바뀌어 스냅샷을 재생성하면 밴드도 따라 움직인다.

    **왜 글자수인가(2026-07-20 preflight_dense 통합):** 도형 수로 재면 밀도 부족을 hero_card
    4장으로 메우는 반사(4카드 강제)를 유발한다 — 골든 항목형은 도형이 적어도(예: ab_simulation
    37도형) 글자가 꽉 차 통과작급이다. 그래서 chars를 주 하한으로, 도형 수는 텅 빈 장만 거르는
    느슨한 바닥(shapes_floor). 스크린샷 장은 DENSITY_EXEMPT로 파생·검사에서 뺀다.
    """
    import json
    from pathlib import Path

    fixed = {"cover", "toc", "part", "closing"}  # 정형 장 — 밀도 대상 아님
    if snapshot_path is None:
        snapshot_path = Path(__file__).resolve().parents[2] / "assets/golden-snapshot.json"
    snap = json.loads(Path(snapshot_path).read_text(encoding="utf-8"))
    pairs = [(key, sl) for key, sl in zip(SNAPSHOT_KEYS, snap["slides"]) if key not in fixed]
    body = [sl for _k, sl in pairs]
    dense_body = [sl for k, sl in pairs if k not in DENSITY_EXEMPT]  # 글자수 하한은 비스크린샷에서
    shape_counts = [len(sl) for sl in body]
    frame_counts = [sum(1 for sh in sl if (sh.get("text") or "").strip()) for sl in body]
    char_counts = [sum(len((sh.get("text") or "").strip()) for sh in sl) for sl in dense_body]
    return {
        "shapes_min": min(shape_counts),
        "frames_min": min(frame_counts),
        "chars_min": round(min(char_counts) * 0.85),  # 통과작 최저 텍스트량 × 경계 여유
        "shapes_floor": 20,  # 텅 빈 장만 거르는 느슨한 바닥(도형 수로 카드 강제 금지)
        "n_body": len(body),
    }


def check_density(shapes, band, screenshot=False):
    """골든 밀도 밴드 — **글자수가 주 지표**(2026-07-20). 변형(adapted)·신규(novel) 장이 골든
    본문 최저 텍스트량보다 성기면 FAIL. 도형 수는 텅 빈 장만 거르는 바닥이다 — 도형 수로 재면
    밀도 부족을 카드로 메우는 반사(4카드 강제)를 유발하므로 chars를 주 하한으로 한다.
    스크린샷 장(§F)은 캡처가 내용을 지므로 밀도 예외.
    """
    n_shapes = len(shapes)
    n_chars = sum(
        len(s.text_frame.text.strip())
        for s in shapes
        if s.has_text_frame and s.text_frame.text.strip()
    )
    floor = band.get("shapes_floor", 20)
    if screenshot:
        msg = f"글자 {n_chars} · 도형 {n_shapes} — 스크린샷 장(§F) 밀도 예외"
        ok, deficit = True, 0
    else:
        ok = n_chars >= band["chars_min"] and n_shapes >= floor
        deficit = max(0, band["chars_min"] - n_chars)
        msg = (
            f"글자 {n_chars}(하한 {band['chars_min']}) · 도형 {n_shapes}(바닥 {floor}) "
            f"— 골든 본문 {band['n_body']}장 텍스트량 파생"
        )
    return bool(ok), msg, int(deficit)


PROFILES = ("sparse", "dense", "legacy")


ICON_MAX = 0.8  # 이보다 큰 그림은 스크린샷으로 본다(assets/icons는 0.17~0.48 정사각)


def count_icons(shapes):
    """이 장이 쓴 **아이콘 어휘** 개수 — 작은 정사각 PICTURE만 센다.

    스크린샷과 구분하는 자는 크기다. `dense.icon()`이 넣는 그림은 size×size(0.17~0.48)이고
    `kit.fit_picture()`가 넣는 캡처는 폭·높이가 인치 단위로 크다.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    n = 0
    for s in shapes:
        if s.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        _x, _y, sw, sh = box(s)
        if sw <= ICON_MAX and sh <= ICON_MAX and 0.6 <= (sw / sh if sh else 0) <= 1.6:
            n += 1
    return n


def check_icon_vocab(icon_total, body_slides):
    """덱 단위 — 본문 전체에서 아이콘 어휘가 **통째로** 죽었으면 FAIL(2026-08-02 신설).

    출처: local-ai-assist 2026-08-02. 조판 스크립트 3,547줄에 `icon` 문자열이 0회였고 본문
    19장의 아이콘 PICTURE가 0개였다(골든 17장은 96개). 사용자 반려는 "카드를 왜 안 썼냐 ·
    비어 보인다"였는데, 골든의 리치함은 `hero_card`라는 **함수**가 아니라 배지·아이콘·배너라는
    **잉크**에서 온다. 부품함을 import하고도 어휘를 안 집으면 이 상태가 되고, 장 단위 검사는
    전부 통과한다(각 장은 그 나름 꽉 차 있으므로) — 그래서 덱 단위 자리가 필요했다.

    **하한을 1보다 올리지 않는다.** 골든도 본문 14장 중 9장(64%)만 쓴다. "장별 커버리지 하한"을
    만들면 그게 "카드 4장 써라"가 태어난 승격(한 장의 지적 → 모든 장의 명령)과 같은 경로다.
    잡는 것은 **"어휘가 통째로 죽었다"** 하나뿐이다.

    정직한 한계: 아이콘 하나 붙이고 통과시키는 건 여전히 가능하다. "의미에 맞게 골랐나"는
    눈검증(⑥) 몫이다. 실측 — 골든 96개 PASS / local-ai-assist 0개 FAIL.
    """
    ok = icon_total > 0
    return (
        ok,
        f"본문 {body_slides}장 아이콘 어휘 {icon_total}개 — 통째로 미사용"
        f'(assets/icons {ICON_MAX}" 이하 정사각 PICTURE 기준)',
        0 if ok else 1,
    )


def generic_checks(shapes, accent, band=None, dup_allow=0, screenshot=False, profile="sparse"):
    """레이아웃 무관 전역 오딧 묶음 — **모든 렌더 경로의 단일 출처**.

    레이아웃 특정 규칙(진행형 도형 허용 수·ink_allow·공기 쌍·노드 높이)은 여기 없다 —
    그건 장 스크립트가 자체 assert하거나 채점(P4)이 본다. score_deck.py가 이 묶음의 러너.
    screenshot=True면 밀도 검사를 예외 처리(§F 스크린샷 장).

    **profile — 프레임은 셋, 규칙은 하나.** 임계만 프레임별로 갈리고 규칙 코드는 공유한다
    (스킬별 경쟁 임계 신설 금지). 어느 검사가 어느 프로파일에 도는지는 **오탐 0 실측**으로 정했다:
      · `sparse` — sparse 골든 장(결론 바 6.60·출처 7.15 프레임).
      · `dense`  — dense 장/adapted/novel(2026-07-20): accent 런 제외(2단 불릿 **강조**는 정상)·
        상한 완화, 경계는 dense 하한(압축 헤더로 세로를 더 씀), 칩 폭 보정, §8 즉흥 카드.
      · `legacy` — **2026-07-29부터 대상이 0이다**(이 프로파일이 겨냥하던 build_pptx 렌더러가
        아카이브됐다). 상수와 분기는 남겨 뒀다 — 아래 실측 근거가 "왜 프레임마다 임계가 다른가"를
        말해주기 때문이다. 다시 쓸 일이 없으면 지워도 된다.
        (이하 2026-07-25 당시 기록)
        전엔 이 경로에 generic_checks가 **아예 안 돌았다**(러너 3종이 golden./adapted./novel만
        대상, consistency-qa는 import조차 안 함) — 남은 legacy 어휘가 무검증으로 나갔다.
        **legacy 산출물 111장 전수 실측**으로 오탐 0인 부분집합만 물린다:
          물림 = accent(런 제외·상한 8) · 경계(legacy 프레임 상수) · 노드 재탕 · 노드 클래스 ·
                 그림 침범 · 텍스트 넘침 · 텍스트 겹침 · §8 즉흥 카드   (전부 0/111)
          미배선 3종 = **채움률**(풀폭 배너 12.13"×0.36 다크 스트립이 구조상 26% → 오탐 1) ·
                 **판정 단어 대비**(check_matrix의 회색 ✓ 3.22·강조어 accent 3.41 = 어휘 idiom →
                 오탐 4) · **밀도 밴드**(골든 458자 하한은 legacy 프레임에 부적합 → 93/111 적색.
                 이 프로파일의 밀도 정본이던 외부 검사기도 2026-07-29 함께 아카이브됐다).
        미배선 3종을 억지로 물리면 상시 적색이 되고, 그건 이 프로젝트가 이미 겪은 우회 경로다.
    """
    if profile not in PROFILES:
        raise ValueError(f"generic_checks: 알 수 없는 profile {profile!r} (허용 {PROFILES})")
    dense, legacy = profile == "dense", profile == "legacy"
    if dense:
        accent_res = check_accent(shapes, accent, cap=ACCENT_MAX_DENSE, count_runs=False)
        bounds_res = check_bounds(shapes, bottom=CONTENT_BOTTOM_DENSE, skip_tops=(7.14,))
        # dense 칩(아이콘+짧은 라벨)은 폭 대비 텍스트가 적어도 죽은 회색이 아니다 — min_w를 칩
        # 폭(≤2.98") 위로 올려 진짜 넓은 죽은 박스(≥3.2")만 잡는다.
        fill_res = check_fill_ratio(shapes, min_w=3.2)
    elif legacy:
        accent_res = check_accent(shapes, accent, cap=ACCENT_MAX_LEGACY, count_runs=False)
        bounds_res = check_bounds(shapes, bottom=CONTENT_BOTTOM_LEGACY, skip_tops=LEGACY_FRAME_TOPS)
        fill_res = None  # 미배선(풀폭 배너 idiom) — 위 docstring 참조
    else:
        accent_res = check_accent(shapes, accent)
        bounds_res = check_bounds(shapes, skip_tops=(6.60, 7.15))
        fill_res = check_fill_ratio(shapes)
    res = [("§2 accent 상한", *accent_res)]
    if not legacy:
        res.append(("P4⑤ 판정 단어 대비", *check_verdict_contrast(shapes, None)))
    if fill_res is not None:
        res.append(("P4④ 채움률", *fill_res))
    res += [
        ("P4③ 노드 재탕", *check_duplicate_nodes(shapes, allow=dup_allow)),
        ("P4⑩ 노드 클래스", *check_node_class(shapes)),
        ("§F 그림 침범", *check_picture_overlap(shapes)),
        ("텍스트 넘침", *check_text_overflow(shapes)),
        ("텍스트 겹침", *check_text_collision(shapes)),
        ("P2③④ 경계", *bounds_res),
    ]
    if dense or legacy:
        # §8 즉흥 밋밋 카드 금지 — hero_card 미사용(구조 0)을 기계로 차단.
        # dense 통과작 8장 · legacy 111장 전수 오탐 0. 어두운 코드·아티팩트 패널은 제외.
        # sparse 골든은 제외(카드 idiom이 다르다).
        res.append(("§8 즉흥 카드", *check_adhoc_card(shapes)))
    if band:
        res.append(("§6-D 밀도 밴드", *check_density(shapes, band, screenshot=screenshot)))
    return res


def report(results, known=None):
    """(이름, ok, 상세, 위반수) 목록 → 출력 + 실패 목록.

    **위반 수는 규칙이 직접 반환한다** — 상세 문자열에서 정규식으로 긁으면 규칙마다 형식이
    달라 조용히 None이 되고 래칫이 죽는다(초판이 `채움 2 (허용 0)`을 못 읽어 그랬다).

    known: {규칙명: 기준선 위반 수} — 이 오딧 **이전에** 사용자 승인으로 확정된 장의 기존 빚.
      규칙을 끄지 않고 기준선으로 박는다(래칫): 기준선 이하면 통과, **넘으면 FAIL**.
      · 규칙을 빼면 hollow-PASS가 된다(존재≠검사).
      · 그대로 빨간불이면 게이트가 상시 적색이라 무시당하고, 결국 우회한다 —
        이 프로젝트가 이미 겪은 실패 경로다.
      기준선을 내리는 건 해당 장 재설계 때. 좋아지면 알려서 기준선을 조인다.
    """
    known = known or {}
    fails = []
    for name, ok, detail, n in results:
        base = known.get(name)
        if not ok and base is not None and n is not None and n <= base:
            mark = "DEBT" if n == base else "GOOD"
            note = "기지 빚(기준선)" if n == base else f"개선 {base}→{n} · 기준선을 낮출 것"
            print(f"{mark} {name:22} {detail}  ← {note}")
            continue
        print(f"{'OK  ' if ok else 'FAIL'} {name:22} {detail}")
        if not ok:
            fails.append(name)
    print("\n" + ("AUDIT PASS" if not fails else f"AUDIT FAIL: {fails}"))
    return fails
