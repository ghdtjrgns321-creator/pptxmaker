"""S14 트러블슈팅(리랭커→지식그래프) — dense 개편 v5(2026-07-20). R1: 골든 variant_c 재사용.

반복 실패에서 확정된 교정(사용자 3지적):
1. 좋은 시각화를 부수지 말 것 — v3/v4에서 골든 A/B 시뮬(팬·컷·경로)을 컴팩트 3칩으로 부수고
   hero_card 4장을 우겨넣었다. 되살린다.
2. 카드 강제 금지 — 밀도 부족을 카드로 메우려는 반사가 병. variant_c는 624글자로 하한(660)에서
   단 36 부족 → **카드가 아니라 주석 텍스트로** 채운다(§8 R2·§7 사용자 3지적).
3. 게이트가 도형 수로 밀도를 재던 게 카드 강제의 기계적 원인이었다 → 게이트를 글자수로 재보정
   (preflight_dense). variant_a(36도형이지만 712글자)가 통과작급인 게 그 증거.

그래서 v5 = 골든 `s14_variants.variant_c`(A/B 시뮬)를 **재드로 없이 그대로 재사용**:
- 헤더 → compact_header + 서브리드(검은 바 결론 흡수, §8)
- 검은 전폭 바 제거
- 압축 헤더·바 제거 공간에 2패널을 ty(위치)+sh(높이)로 재배치(§8: 둘 다 스케일해야 안 뜬다)
- 두 패널 note를 나머지 균열(점수 근거없음·라우팅 비결정)·구조(위계·문단번호·상호참조)로
  보강 → 글자수 하한 통과. 카드 0.

형제 대조: S9=그래프 설명 / S14=A/B 대결(고유). preflight_dense green 후 제시.
"""

from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from . import dense as D
from . import grid as G
from .kit import SLIDE_H, SLIDE_W, add_box, add_text, load_kit, mix

K = load_kit()

KICKER = "4. 문제와 해결"
SOURCE = "출처: 7_JOURNEY.md §7.3·§7.4 (재구축 여정) · 00_factsheet.md §A'·§D"
CHAIN_TOP, CHAIN_BOTTOM, CHAIN_H = 3.20, 5.30, 0.46
CHAIN_PITCH_CAP = 0.78  # 여유 있을 때의 골든 리듬(3단 = 3.20·3.98·4.76)

def chain_layout(n):
    """근거 체인 n단의 `(y, fill, color, border)` 리스트.

    2026-07-25 항목 수 파생: 전엔 `step_y = [3.2, 3.98, 4.76]`과 `step_style` 3줄 리스트를
    `zip(c["steps"], …)`으로 물렸다 — 단계가 4개면 **4번째가 조용히 사라지고** 2개면 아래가
    비었다. 이제 y는 `grid.pitch`로, 강조 인코딩은 규칙으로 파생한다:
    **마지막 = accent 테두리(도달점), 중간 = 다크 채움(전환점), 나머지 = 헤어라인.**
    좌표·강조가 개수에서 나오므로 단계가 늘어도 "도달 경로가 근거"라는 물성이 유지된다.
    """
    p = G.pitch(n, CHAIN_TOP, CHAIN_BOTTOM, CHAIN_H, cap=CHAIN_PITCH_CAP, what="근거 단계")
    dark = (n - 1) // 2  # n=3 → 1(가운데). n=1이면 마지막과 겹치므로 accent가 우선한다.
    out = []
    for i in range(n):
        if i == n - 1:
            style = (C["bg"], C["primary"], C["accent"])
        elif i == dark:
            style = (C["primary"], C["bg"], None)
        else:
            style = (C["bg"], C["primary"], C["muted"])
        out.append((CHAIN_TOP + i * p, *style))
    return out

def _solid_arrow(slide, x1, y1, x2, y2):
    from pptx.enum.shapes import MSO_CONNECTOR
    from pptx.oxml.ns import qn

    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = C["primary"]
    conn.line.width = Pt(1.75)
    ln = conn.line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    return conn

# content_contract 계약 키 — golden.ab_simulation의 필수 키 출처.
# 2026-07-26 `s14_variants`(sparse 시안)에서 이관. 골든 전용 글이다.
_GOLDEN_TEXT = {  # 이관해온 골든 원문 — 계약 dict가 아니라 DEFAULT의 재료

    "headline": "같은 질문, 두 시스템 — 점수는 떨어뜨리고, 경로는 도달한다",
    "kicker": KICKER,
    "left_title": "리랭커 — 점수가 근거다, 그런데 점수엔 근거가 없다",
    "question": "“볼륨디스카운트 조항이 있는 계약은…”",
    "rank_labels": ["1위", "2위", "3위"],
    "surface_card": "표면 단어가\n비슷한 문단",
    "cut_label": "상위 컷 — 아래는 버려진다",
    "reject_mark": "✕",
    "reject_card": "전문가 큐레이션 문단 — 정답인데 표면 유사도가 낮다",
    "left_note": "실측 — 전문가 배정 큐레이션 문서가 표면 유사도가 낮아 105건 중 103건 탈락(98%). 이미 bypass로 우회 중이었다 — 유사도 필터가 정답을 거른다는 신호.",
    "right_title": "지식그래프 — 경로가 근거다",
    "steps": [
        ("용어사전 매칭 — 볼륨디스카운트 → 변동대가", "사람이 전수 검수한 색인 (등재 423)"),
        ("개념 노드 — 변동대가", "목차 위계의 그 자리 — 측정 > 거래가격 산정 > 변동대가"),
        ("관할 문단 50~54 — 변동대가 추정", "개념에 배정된 문단으로 — 간선 자체가 근거 ✓"),
    ],
    "right_note": "탈락시킬 점수가 없다 — 연결이 있으면 도달하고, 어떤 간선을 지났는지가 그대로 답변의 근거가 된다.",
    # 구 S17(7축 비교) 삭제 후 그 주장을 이 바가 진다(2026-07-16) — 검색뿐 아니라 전 축이 같은 치환
    "bar": "이 한 축이 아니라 검색·근거·판단·검증·재현 전부 — 확률 신호 자리마다 기준서의 구조가 들어가 있다",
    "source": SOURCE,
}

# ── 이 장이 쓰는 도해 원소 (2026-07-26 `s14_variants`에서 이관) ─────────────────
def _dash_arrow(slide, x1, y1, x2, y2):
    """점선 화살표 — 확률 신호(불확실) 표기용."""
    from pptx.enum.shapes import MSO_CONNECTOR
    from pptx.oxml.ns import qn

    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = C["muted"]
    conn.line.width = Pt(1.25)
    ln = conn.line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "sm", "len": "sm"}))
    return conn

def _chip(slide, x, y, w, h, text, fill, color, border=None, bold=True, size=None):
    from pptx.enum.lang import MSO_LANGUAGE_ID

    box = add_box(slide, x, y, w, h, fill=fill, line=border, line_w=1.25, shape="round")
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.06)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p._p.get_or_add_pPr().set("eaLnBrk", "0")
    r = p.add_run()
    r.text = text
    r.font.name = F["head"] if bold else F["body"]
    r.font.bold = bold
    r.font.size = Pt(size or S["caption"])
    r.font.color.rgb = color
    r.font.language_id = MSO_LANGUAGE_ID.KOREAN
    return box


# 콘텐츠 기본값(골든 내용) — c=None이면 이 값, override 시 텍스트만 교체(좌표·색·도형 고정)
C, S, F = K["rgb"], K["sizes"], K["fonts"]

KICKER = _GOLDEN_TEXT["kicker"]
HEADLINE = _GOLDEN_TEXT["headline"]
SOURCE = _GOLDEN_TEXT["source"]
SUBLEAD = "확률 신호가 있던 자리마다 기준서의 구조가 들어간다 — 임베딩·가중치·리랭커 없이 온톨로지 그래프로 동작한다"

# 두 패널 note 보강 — 나머지 균열/구조를 텍스트로(카드 아님). 원문 근거(7_JOURNEY §7.3).
LEFT_NOTE = (
    "실측 — 전문가 배정 큐레이션 문서가 표면 유사도 낮아 105건 중 103건 탈락(98%). 게다가 점수엔 "
    "근거가 없고(왜 3점인지 설명 불가), 3중 매칭 신호가 서로 압도해 라우팅이 비결정적 — 같은 "
    "질문에 토픽이 바뀌면 답도 바뀌었다. 유사도 필터가 정답을 거른다는 신호."
)
RIGHT_NOTE = (
    "탈락시킬 점수가 없다 — 연결이 있으면 도달하고, 지난 간선이 그대로 답변의 근거다. 위계(계층 "
    "간선)·사례마다 붙은 문단번호·문단 간 상호참조가 이미 기준서에 명시돼 있어, 확률로 추정할 "
    "필요가 없다. 판단트리도 원문 앵커로 재검수했다."
)

# content_contract 계약 키 — registry가 이 dense를 golden.ab_simulation으로 부를 때 필수 키 출처.
# dense가 덧대는 left_note/right_note도 골든 텍스트라 계약 키에 포함(누락 시 타 프로젝트에 K-IFRS 누출).
DEFAULT = {**_GOLDEN_TEXT, "left_note": LEFT_NOTE, "right_note": RIGHT_NOTE}

# ── affine 세로 변환: 위치(ty)+높이(sh) 둘 다 스케일(§8) — 골든 [1.85,6.35] → dense [1.12,6.72] ──
_A = (6.72 - 1.12) / (6.35 - 1.85)
_B = 1.12 - _A * 1.85


def ty(y):
    return _A * y + _B


def sh(h):
    return _A * h


def build(prs, c=None):
    c = {**DEFAULT, **(c or {})}
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    # content에서 읽는다 — 모듈 상수를 그대로 쓰면 실전 덱에 K-IFRS 글이 그대로 나간다.
    # (2026-07-26 dense 승격 때 compare_golden PARAM 오라클이 잡았다: sparse 원본은
    #  `header(slide, c["headline"], c["kicker"])`인데 dense 이식에서 상수로 굳었다.)
    D.compact_header(slide, c.get("kicker") or KICKER, c.get("headline") or HEADLINE)
    add_text(
        slide,
        G.MARGIN_L,
        0.76,
        D.FULL_W,
        0.24,
        c.get("sublead") or SUBLEAD,
        S["caption"],
        F["body"],
        C["muted"],
    )

    # ══ 좌 판: 리랭커 (bg_alt) — 골든 variant_c 그대로, ty+sh ══
    lx, lw = G.MARGIN_L, 5.75
    add_box(slide, lx, ty(1.85), lw, sh(6.35 - 1.85), fill=C["bg_alt"])
    px = lx + 0.25
    add_text(
        slide,
        px,
        ty(2.0),
        lw - 0.5,
        sh(0.28),
        c["left_title"],
        S["head"],
        F["head"],
        C["primary"],
        bold=True,
    )
    q_w = 3.7
    q_x = lx + (lw - q_w) / 2
    _chip(
        slide,
        q_x,
        ty(2.42),
        q_w,
        sh(0.42),
        c["question"],
        C["bg"],
        C["primary"],
        border=C["muted"],
        bold=False,
    )
    cards_y = 3.75
    card_w = 1.62
    card_xs = [px, px + 1.82, px + 3.64]
    qcx = lx + lw / 2
    for cx, lab in zip(card_xs, c["rank_labels"]):
        ccx = cx + card_w / 2
        _dash_arrow(slide, qcx, ty(2.84), ccx, ty(cards_y))
        add_text(
            slide,
            (qcx + ccx) / 2 - 0.25,
            ty((2.84 + cards_y) / 2) - 0.12,
            0.5,
            0.22,
            lab,
            S["caption"],
            F["head"],
            C["primary"],
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        _chip(
            slide,
            cx,
            ty(cards_y),
            card_w,
            sh(0.5),
            c["surface_card"],
            C["bg"],
            C["muted"],
            border=C["muted"],
            bold=False,
        )
    cut_y = 4.62
    cut = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(px),
        Inches(ty(cut_y)),
        Inches(lx + lw - 0.25),
        Inches(ty(cut_y)),
    )
    cut.line.color.rgb = C["primary"]
    cut.line.width = Pt(1.25)
    cut.line._get_or_add_ln().append(
        cut.line._get_or_add_ln().makeelement(qn("a:prstDash"), {"val": "dash"})
    )
    add_text(
        slide,
        lx + lw - 2.3,
        ty(cut_y) - 0.26,
        2.05,
        0.22,
        c["cut_label"],
        S["caption"],
        F["body"],
        C["muted"],
        align=PP_ALIGN.RIGHT,
    )
    add_text(
        slide,
        px + 0.35,
        ty(4.86),
        0.4,
        sh(0.34),
        c["reject_mark"],
        S["title"],
        F["head"],
        C["primary"],
        bold=True,
    )
    _chip(
        slide,
        px + 0.85,
        ty(4.82),
        3.4,
        sh(0.5),
        c["reject_card"],
        mix(C["bg_alt"], C["muted"], 0.25),
        C["primary"],
        bold=False,
    )
    add_text(
        slide,
        px,
        ty(5.5),
        lw - 0.5,
        sh(0.85),
        c["left_note"],
        S["caption"],
        F["body"],
        C["muted"],
        line_spacing=1.3,
    )

    # ══ 우 판: 지식그래프 (백 + accent 상단 룰) ══
    rx = 6.75
    rw = G.RIGHT_EDGE - rx
    add_box(slide, rx, ty(1.85), rw, 0.035, fill=C["accent"])
    add_box(slide, rx, ty(1.885), rw, sh(6.35 - 1.885), fill=C["bg"], line=C["bg_alt"], line_w=1.0)
    qx = rx + 0.25
    add_text(
        slide,
        qx,
        ty(2.0),
        rw - 0.5,
        sh(0.28),
        c["right_title"],
        S["head"],
        F["head"],
        C["primary"],
        bold=True,
    )
    chain_w = 3.55
    ccx = rx + 0.35 + chain_w / 2
    _chip(
        slide,
        rx + 0.35,
        ty(2.42),
        chain_w,
        sh(0.42),
        c["question"],
        C["bg"],
        C["primary"],
        border=C["muted"],
        bold=False,
    )
    # 스텝 y·강조는 골든 정본 chain_layout에 위임(단계 수 파생) — dense는 ty/sh affine만 얹는다
    prev_bottom = 2.84
    for (text, note), (syb, fill, color, border) in zip(
        c["steps"], chain_layout(len(c["steps"]))
    ):
        _solid_arrow(slide, ccx, ty(prev_bottom), ccx, ty(syb))
        _chip(slide, rx + 0.35, ty(syb), chain_w, sh(CHAIN_H), text, fill, color, border=border)
        add_text(
            slide,
            rx + 0.35 + chain_w + 0.15,
            ty(syb) + sh(0.1),
            rw - chain_w - 0.85,
            sh(0.4),
            note,
            S["caption"],
            F["body"],
            C["muted"],
            line_spacing=1.2,
        )
        prev_bottom = syb + CHAIN_H
    add_text(
        slide,
        qx,
        ty(5.5),
        rw - 0.5,
        sh(0.85),
        c["right_note"],
        S["caption"],
        F["body"],
        C["text"],
        line_spacing=1.3,
    )

    D.source_line(slide, SOURCE)
    return slide
