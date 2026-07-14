"""S14 시안 — 트러블슈팅: 리랭커 → 지식그래프 전환 (신규 어휘: 전환 대비 2패널. UI 패스).

콘텐츠 실물: 7_JOURNEY.md §7.3(균열 4·공통 원인·대체표), 00_factsheet.md §A'(103/105)·§D(간선 수).
Part Ⅳ는 전/후 서사의 합법 무대(design-rules §8d 격리 조항).
실행: uv run python golden/s14_variants.py → golden/variants/s14_variants.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from . import grid as G
from .kit import SLIDE_H, SLIDE_W, add_box, add_text, load_kit, mix

K = load_kit()
C, S, F = K["rgb"], K["sizes"], K["fonts"]

KICKER = "4. 문제와 해결"
SOURCE = "출처: 7_JOURNEY.md §7.3·§7.4 (재구축 여정) · 00_factsheet.md §A'·§D"

CRACKS = [  # 7_JOURNEY §7.3 균열 4 — 원문 축약
    (
        "리랭커가 정답을 떨어뜨린다",
        "표면 유사도로 정렬 — 전문가 큐레이션 문서 105건 중 103건 탈락(98%)",
    ),
    (
        "점수에 근거가 없다",
        "“왜 이 문단이 3점인가”를 설명하지 못한다 — 순위를 좌우하는 값이 근거 불명",
    ),
    (
        "라우팅이 비결정적이다",
        "3중 매칭 신호가 서로를 압도 — 같은 질문에 토픽이 바뀌면 답도 바뀐다",
    ),
    ("판단트리가 미검수다", "AI 교차생성 결과물을 기준서 원문 대조 없이 신뢰하고 있었다"),
]

STRUCTURES = [  # 전환 후 — 기준서의 명시적 구조 (§D 실측 수)
    (
        "위계가 이미 있다",
        "목차 계층과 5단계 모형을 그대로 간선으로 — 순위 점수 없이 위상 순서로 정렬 (계층 간선 79)",
    ),
    (
        "사례마다 문단번호가 붙어 있다",
        "QNA·감리·적용사례가 인용한 문단 번호로 그래프에 바로 연결 (사례 간선 1,220)",
    ),
    (
        "상호참조가 명시돼 있다",
        "문단이 문단을 지목(상호참조 간선 244) — 이웃 확장이 확률 아닌 관계로 결정된다",
    ),
]


def header(slide, headline, kicker=KICKER):
    add_text(
        slide, G.MARGIN_L, 0.42, 8.0, 0.28, kicker, S["caption"], F["head"], C["muted"], bold=True
    )
    add_text(
        slide,
        G.MARGIN_L,
        0.72,
        G.RIGHT_EDGE - G.MARGIN_L,
        0.55,
        headline,
        S["section"],
        F["head"],
        C["primary"],
        bold=True,
    )
    add_box(slide, G.MARGIN_L, G.RULE_Y, G.RIGHT_EDGE - G.MARGIN_L, 0.014, fill=C["muted"])


def bar_and_source(slide, text, source=SOURCE):
    bar = add_box(slide, G.MARGIN_L, G.BAR_Y, G.RIGHT_EDGE - G.MARGIN_L, G.BAR_H, fill=C["primary"])
    tf = bar.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name, run.font.bold = F["head"], True
    run.font.size = Pt(S["body"])
    run.font.color.rgb = C["bg"]
    add_text(
        slide,
        G.MARGIN_L,
        G.SOURCE_Y,
        G.RIGHT_EDGE - G.MARGIN_L,
        0.3,
        source,
        S["foot"],
        F["body"],
        C["muted"],
    )


def _panel_head(slide, x, y, w, title, sub, title_color, sub_color):
    add_text(slide, x, y, w, 0.3, title, S["head"], F["head"], title_color, bold=True)
    add_text(slide, x, y + 0.34, w, 0.26, sub, S["caption"], F["body"], sub_color)


def _item(slide, x, y, w, no, head, body, head_color, body_color, num_color):
    add_text(slide, x, y, 0.45, 0.26, f"0{no}", S["body"], F["head"], num_color, bold=True)
    add_text(slide, x + 0.42, y, w - 0.42, 0.26, head, S["body"], F["head"], head_color, bold=True)
    add_text(
        slide,
        x + 0.42,
        y + 0.3,
        w - 0.42,
        0.5,
        body,
        S["caption"],
        F["body"],
        body_color,
        line_spacing=1.2,
    )


def variant_a(prs):
    """A — 전환 대비 2패널: 좌 균열 4(회색) → 대형 화살표 → 우 구조 3(백)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    header(slide, "왜 리랭커를 버렸나 — 확률 신호에서 기준서의 구조로")
    add_text(
        slide,
        G.MARGIN_L,
        G.CONTENT_TOP,
        G.RIGHT_EDGE - G.MARGIN_L,
        0.3,
        "네 균열이 한 방향을 가리켰다 — 기준서의 명확한 구조를 확률·근거불명 신호로 뭉개고 있었다.",
        S["sub"],
        F["body"],
        C["text"],
    )
    # ── 좌 패널: 전환 전 (회색 톤) ──
    lx, lw = G.MARGIN_L, 5.5
    add_box(slide, lx, 2.2, lw, 4.15, fill=C["bg_alt"])
    px, pw = lx + 0.25, lw - 0.5
    _panel_head(
        slide,
        px,
        2.38,
        pw,
        "전환 전 — 유사도·가중치·리랭커 복합체",
        "벡터 + BM25 + RRF + 가중치 + 리랭커 — 지표는 좋아 보였다",
        C["primary"],
        C["muted"],
    )
    for i, (head, body) in enumerate(CRACKS):
        _item(slide, px, 3.1 + i * 0.8, pw, i + 1, head, body, C["primary"], C["muted"], C["muted"])
    # ── 중앙 전환 화살표 ──
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(6.25), Inches(4.05), Inches(0.72), Inches(0.5)
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = C["accent"]
    arrow.line.fill.background()
    arrow.shadow.inherit = False
    # ── 우 패널: 전환 후 (백 + accent 상단 룰) ──
    rx = 7.12
    rw = G.RIGHT_EDGE - rx  # 5.61
    add_box(slide, rx, 2.2, rw, 0.035, fill=C["accent"])
    add_box(slide, rx, 2.235, rw, 4.115, fill=C["bg"], line=C["bg_alt"], line_w=1.0)
    qx, qw = rx + 0.25, rw - 0.5
    _panel_head(
        slide,
        qx,
        2.38,
        qw,
        "전환 후 — 온톨로지 지식그래프",
        "확률로 추정할 필요가 없었다 — 관계는 이미 기준서에 있다",
        C["primary"],
        C["muted"],
    )
    for i, (head, body) in enumerate(STRUCTURES):
        _item(
            slide, qx, 3.2 + i * 1.02, qw, i + 1, head, body, C["primary"], C["text"], C["primary"]
        )
    bar_and_source(
        slide, "확률 신호 전량 폐기 — 임베딩·가중치·리랭커 없이 온톨로지 그래프로 동작한다"
    )
    return slide


SWAP_ROWS = [  # 7_JOURNEY §7.3 전환 표 — 6행 전량
    ("Cohere 리랭커 / reranker bypass", "그래프 위상 순서"),
    ("weight_score 가중치", "개념 계층·용어 변별력"),
    ("RRF / BM25 / HyDE 벡터검색", "개념 → 문단 → 상호참조 이웃 그래프 탐색"),
    ("핀포인트 직접 조회", "개념 → 관할 문단 관계"),
    ("3중 유사도 토픽 매칭", "용어사전 substring(결정적) + 토픽→개념 매핑"),
    ("26토픽 미검수 판단트리", "본문·부록B 원문 앵커 판단트리 41개"),
]


def _subhead(slide, x, y, w, text):
    add_text(slide, x, y, w, 0.3, text, S["head"], F["head"], C["primary"], bold=True)
    add_box(slide, x, y + 0.34, w, 0.012, fill=C["muted"])


def variant_b(prs):
    """B — 실측 인코딩: 105도트 그리드 + 균열 칩 3 + 버린 것→대체물 네이티브 표 6행."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    header(slide, "왜 리랭커를 버렸나 — 확률 신호에서 기준서의 구조로")
    add_text(
        slide,
        G.MARGIN_L,
        1.72,
        G.RIGHT_EDGE - G.MARGIN_L,
        0.28,
        "네 균열이 한 방향을 가리켰다 — 기준서의 명확한 구조를 확률·근거불명 신호로 뭉개고 있었다.",
        S["sub"],
        F["body"],
        C["text"],
    )
    # ── 상단 좌: 균열 ① 도트 그리드 (105건 실측 인코딩) ──
    lx, lw = G.MARGIN_L, 5.7
    _subhead(slide, lx, 2.14, lw, "균열 ① — 유사도 필터가 정답을 거른다")
    dot_d, pitch = 0.115, 0.155
    cols, total, survive = 21, 105, 2
    gx, gy = lx, 2.72
    for n in range(total):
        r, c = divmod(n, cols)
        alive = n >= total - survive  # 마지막 2개 = 통과분
        add_box(
            slide,
            gx + c * pitch,
            gy + r * pitch,
            dot_d,
            dot_d,
            fill=C["accent"] if alive else mix(C["bg"], C["muted"], 0.45),
            shape="oval",
        )
    grid_h = 5 * pitch  # 0.775
    add_text(
        slide,
        gx,
        gy + grid_h + 0.08,
        lw,
        0.5,
        [
            [
                ("전문가 큐레이션 문서 105건 중 리랭커 통과 ", {}),
                ("2건", {"bold": True, "color": C["accent"]}),
                (" — 탈락 98%", {"bold": True, "color": C["primary"]}),
            ],
            [("Cohere 리랭커는 query-document 표면 유사도로 정렬 — bypass로 우회 중이었다", {})],
        ],
        S["caption"],
        F["body"],
        C["muted"],
        line_spacing=1.3,
    )
    # ── 상단 우: 균열 ②③④ 칩 ──
    rx = 6.65
    rw = G.RIGHT_EDGE - rx  # 6.08
    add_box(slide, 6.45, 2.14, 0.012, 2.1, fill=C["bg_alt"])
    _subhead(slide, rx, 2.14, rw, "함께 드러난 균열 ②③④")
    chips = [
        ("점수에 근거가 없다", "“왜 이 문단이 3점인가”를 설명하지 못한다"),
        ("라우팅이 비결정적이다", "같은 질문에 토픽이 바뀌면 답도 바뀐다 — 3중 매칭 신호 충돌"),
        ("판단트리가 미검수다", "AI 교차생성 결과물을 기준서 원문 대조 없이 신뢰"),
    ]
    from pptx.enum.lang import MSO_LANGUAGE_ID

    for i, (head, body) in enumerate(chips):
        chip = add_box(slide, rx, 2.6 + i * 0.55, rw, 0.48, fill=C["bg_alt"], shape="round")
        tf = chip.text_frame
        tf.word_wrap = False
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.15)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p._p.get_or_add_pPr().set("eaLnBrk", "0")
        r1 = p.add_run()
        r1.text = head
        r1.font.name, r1.font.bold = F["head"], True
        r1.font.size = Pt(S["body"])
        r1.font.color.rgb = C["primary"]
        r2 = p.add_run()
        r2.text = f"   {body}"
        r2.font.name = F["body"]
        r2.font.size = Pt(S["caption"])
        r2.font.color.rgb = C["muted"]
        for r in (r1, r2):
            r.font.language_id = MSO_LANGUAGE_ID.KOREAN
    # ── 하단 풀폭: 전환 기록 네이티브 표 (버린 것 → 대체물, 6행 전량) ──
    _subhead(
        slide,
        G.MARGIN_L,
        4.42,
        G.RIGHT_EDGE - G.MARGIN_L,
        "전환 기록 — 버린 것과 대체물 (원문 표 전량)",
    )
    tbl_y = 4.72
    n_rows = len(SWAP_ROWS) + 1
    gf = slide.shapes.add_table(
        n_rows,
        3,
        Inches(G.MARGIN_L),
        Inches(tbl_y),
        Inches(G.RIGHT_EDGE - G.MARGIN_L),
        Inches(0.2 * n_rows),
    )
    tbl = gf.table
    for ci, cw in enumerate((5.2, 0.53, 6.4)):
        tbl.columns[ci].width = Inches(cw)
    data = [("버린 것", "", "대체물")] + [(a, "→", b) for a, b in SWAP_ROWS]
    for ri, row in enumerate(data):
        tbl.rows[ri].height = Inches(0.22 if ri else 0.24)
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = (
                C["primary"] if ri == 0 else (C["bg_alt"] if ri % 2 == 0 else C["bg"])
            )
            cell.margin_left = cell.margin_right = Inches(0.08)
            cell.margin_top = cell.margin_bottom = Inches(0.004)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p._p.get_or_add_pPr().set("eaLnBrk", "0")
            if ci == 1:
                p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = val
            r.font.name = F["head"] if ri == 0 else F["body"]
            r.font.size = Pt(S["caption"])
            r.font.bold = ri == 0
            r.font.color.rgb = C["bg"] if ri == 0 else (C["muted"] if ci == 0 else C["primary"])
            r.font.language_id = MSO_LANGUAGE_ID.KOREAN
    bar_and_source(
        slide, "확률 신호 전량 폐기 — 임베딩·가중치·리랭커 없이 온톨로지 그래프로 동작한다"
    )
    return slide


def _dash_arrow(slide, x1, y1, x2, y2):
    """점선 화살표 — 확률 신호(불확실) 표기용."""
    from pptx.enum.shapes import MSO_CONNECTOR
    from pptx.oxml.ns import qn

    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = C["muted"]
    conn.line.width = Pt(1.25)
    ln = conn.line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "sm", "len": "sm"}))
    return conn


def _solid_arrow(slide, x1, y1, x2, y2):
    from pptx.enum.shapes import MSO_CONNECTOR
    from pptx.oxml.ns import qn

    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = C["primary"]
    conn.line.width = Pt(1.75)
    ln = conn.line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    return conn


def _chip(slide, x, y, w, h, text, fill, color, border=None, bold=True, size=None):
    from pptx.enum.lang import MSO_LANGUAGE_ID

    box = add_box(slide, x, y, w, h, fill=fill, line=border, line_w=1.25, shape="round")
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.06)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p._p.get_or_add_pPr().set("eaLnBrk", "0")
    r = p.add_run()
    r.text = text
    r.font.name = F["head"] if bold else F["body"]
    r.font.bold = bold
    r.font.size = Pt(size or S["caption"])
    r.font.color.rgb = color
    r.font.language_id = MSO_LANGUAGE_ID.KOREAN
    return box


# 콘텐츠 기본값(골든 내용) — c=None이면 이 값, override 시 텍스트만 교체(좌표·색·도형 고정)
DEFAULT = {
    "headline": "같은 질문, 두 시스템 — 점수는 떨어뜨리고, 경로는 도달한다",
    "kicker": KICKER,
    "left_title": "리랭커 — 점수가 근거다, 그런데 점수엔 근거가 없다",
    "question": "“볼륨디스카운트 조항이 있는 계약은…”",
    "rank_labels": ["1위", "2위", "3위"],
    "surface_card": "표면 단어가\n비슷한 문단",
    "cut_label": "상위 컷 — 아래는 버려진다",
    "reject_mark": "✕",
    "reject_card": "전문가 큐레이션 문단 — 정답인데 표면 유사도가 낮다",
    "left_note": "실측 — 전문가 배정 큐레이션 문서가 표면 유사도가 낮아 105건 중 103건 탈락(98%). 이미 bypass로 우회 중이었다 — 유사도 필터가 정답을 거른다는 신호.",
    "right_title": "지식그래프 — 경로가 근거다",
    "steps": [
        ("용어사전 매칭 — 볼륨디스카운트 → 변동대가", "사람이 전수 검수한 색인 (등재 423)"),
        ("개념 노드 — 변동대가", "목차 위계의 그 자리 — 측정 > 거래가격 산정 > 변동대가"),
        ("관할 문단 50~54 — 변동대가 추정", "개념에 배정된 문단으로 — 간선 자체가 근거 ✓"),
    ],
    "right_note": "탈락시킬 점수가 없다 — 연결이 있으면 도달하고, 어떤 간선을 지났는지가 그대로 답변의 근거가 된다.",
    "bar": "확률 신호 전량 폐기 — 임베딩·가중치·리랭커 없이 온톨로지 그래프로 동작한다",
    "source": SOURCE,
}


def variant_c(prs, c=None):
    """C — 동일 질문 A/B 시뮬레이션: 좌 점수 산탄·컷라인 탈락 vs 우 결정 경로 도달.

    c: 텍스트 내용 override(None=골든 기본값). 좌표·색·도형은 고정.
    """
    c = {**DEFAULT, **(c or {})}
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    header(slide, c["headline"], c["kicker"])
    # ── 좌 판: 리랭커 (bg_alt) ──
    lx, lw = G.MARGIN_L, 5.75
    add_box(slide, lx, 1.85, lw, 4.5, fill=C["bg_alt"])
    px = lx + 0.25
    add_text(
        slide,
        px,
        2.0,
        lw - 0.5,
        0.28,
        c["left_title"],
        S["head"],
        F["head"],
        C["primary"],
        bold=True,
    )
    q_w = 3.7
    q_x = lx + (lw - q_w) / 2
    _chip(
        slide,
        q_x,
        2.42,
        q_w,
        0.42,
        c["question"],
        C["bg"],
        C["primary"],
        border=C["muted"],
        bold=False,
    )
    # 점선 산탄 3 + 점수 라벨(실측 범위 양끝 0.47·0.29, 가운데는 근거불명 ?)
    cards_y = 3.75
    card_w = 1.62
    card_xs = [px, px + 1.82, px + 3.64]
    labels = c["rank_labels"]
    qcx = lx + lw / 2
    for cx, lab in zip(card_xs, labels):
        ccx = cx + card_w / 2
        _dash_arrow(slide, qcx, 2.84, ccx, cards_y)
        add_text(
            slide,
            (qcx + ccx) / 2 - 0.25,
            (2.84 + cards_y) / 2 - 0.12,
            0.5,
            0.22,
            lab,
            S["caption"],
            F["head"],
            C["primary"],
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        _chip(
            slide,
            cx,
            cards_y,
            card_w,
            0.5,
            c["surface_card"],
            C["bg"],
            C["muted"],
            border=C["muted"],
            bold=False,
        )
    # 컷라인 (점선) + 탈락 카드
    cut_y = 4.62
    from pptx.enum.shapes import MSO_CONNECTOR
    from pptx.oxml.ns import qn

    cut = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(px), Inches(cut_y), Inches(lx + lw - 0.25), Inches(cut_y)
    )
    cut.line.color.rgb = C["primary"]
    cut.line.width = Pt(1.25)
    cut.line._get_or_add_ln().append(
        cut.line._get_or_add_ln().makeelement(qn("a:prstDash"), {"val": "dash"})
    )
    add_text(
        slide,
        lx + lw - 2.3,
        cut_y - 0.26,
        2.05,
        0.22,
        c["cut_label"],
        S["caption"],
        F["body"],
        C["muted"],
        align=PP_ALIGN.RIGHT,
    )
    add_text(
        slide,
        px + 0.35,
        4.86,
        0.4,
        0.34,
        c["reject_mark"],
        S["title"],
        F["head"],
        C["primary"],
        bold=True,
    )
    _chip(
        slide,
        px + 0.85,
        4.82,
        3.4,
        0.5,
        c["reject_card"],
        mix(C["bg_alt"], C["muted"], 0.25),
        C["primary"],
        bold=False,
    )
    add_text(
        slide,
        px,
        5.5,
        lw - 0.5,
        0.6,
        c["left_note"],
        S["caption"],
        F["body"],
        C["muted"],
        line_spacing=1.25,
    )
    # ── 우 판: 지식그래프 (백 + accent 상단 룰) ──
    rx = 6.75
    rw = G.RIGHT_EDGE - rx  # 5.98
    add_box(slide, rx, 1.85, rw, 0.035, fill=C["accent"])
    add_box(slide, rx, 1.885, rw, 4.465, fill=C["bg"], line=C["bg_alt"], line_w=1.0)
    qx = rx + 0.25
    add_text(
        slide,
        qx,
        2.0,
        rw - 0.5,
        0.28,
        c["right_title"],
        S["head"],
        F["head"],
        C["primary"],
        bold=True,
    )
    chain_w = 3.55
    ccx = rx + 0.35 + chain_w / 2  # 체인 축
    _chip(
        slide,
        rx + 0.35,
        2.42,
        chain_w,
        0.42,
        c["question"],
        C["bg"],
        C["primary"],
        border=C["muted"],
        bold=False,
    )
    # 스텝 스타일(구조·색 고정: fill, color, border), 텍스트는 c["steps"]에서
    step_style = [
        (C["bg"], C["primary"], C["muted"]),
        (C["primary"], C["bg"], None),
        (C["bg"], C["primary"], C["accent"]),
    ]
    step_y = [3.2, 3.98, 4.76]
    prev_bottom = 2.84
    for (text, note), (fill, color, border), sy in zip(c["steps"], step_style, step_y):
        _solid_arrow(slide, ccx, prev_bottom, ccx, sy)
        _chip(slide, rx + 0.35, sy, chain_w, 0.46, text, fill, color, border=border)
        add_text(
            slide,
            rx + 0.35 + chain_w + 0.15,
            sy + 0.1,
            rw - chain_w - 0.85,
            0.4,
            note,
            S["caption"],
            F["body"],
            C["muted"],
            line_spacing=1.15,
        )
        prev_bottom = sy + 0.46
    add_text(
        slide,
        qx,
        5.5,
        rw - 0.5,
        0.6,
        c["right_note"],
        S["caption"],
        F["body"],
        C["text"],
        line_spacing=1.25,
    )
    bar_and_source(slide, c["bar"], source=c["source"])
    return slide


def audit(prs):
    from pptx.oxml.ns import qn

    fails = []
    for si, sl in enumerate(prs.slides):
        accents = 0
        for sh in sl.shapes:
            t, le = sh.top / 914400, sh.left / 914400
            b, r = t + sh.height / 914400, le + sh.width / 914400
            full_bleed = sh.width / 914400 >= SLIDE_W - 0.05
            if t < 6.4 and b > 6.37 and not full_bleed:
                fails.append((si, "bottom", round(b, 2), sh.shape_id))
            if r > 12.75 and not full_bleed:
                fails.append((si, "right", round(r, 2), sh.shape_id))
            for el in sh._element.iter(qn("a:srgbClr")):
                if el.get("val") == "D66E3A":
                    if el.getparent().tag.endswith("}solidFill"):
                        accents += 1
        if accents > 4:
            fails.append((si, "accent>4", accents))
    assert not fails, f"AUDIT FAIL {fails}"
    print("audit pass")


def main():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    variant_a(prs)
    variant_b(prs)
    variant_c(prs)
    audit(prs)
    out = Path(__file__).parent / "variants" / "s14_variants.pptx"
    prs.save(str(out))
    print(f"saved {out}")


if __name__ == "__main__":
    main()
