"""brand-kit.yaml 로더 + 공용 도형 헬퍼 — **유일 소스**(2026-07-15 단일화로 golden/ 사본 삭제).

색·폰트·크기는 전부 brand-kit.yaml에서 온다. 이 모듈에 hex·pt 리터럴을 두지 않는다.
"""

from pathlib import Path

import yaml
from pptx.dml.color import RGBColor
from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[2]  # goldenfab -> scripts -> pptx-build
BRAND_KIT = ROOT / "assets/brand-kit.yaml"

SLIDE_W, SLIDE_H = 13.333, 7.5  # 16:9 inch


def load_kit() -> dict:
    kit = yaml.safe_load(BRAND_KIT.read_text(encoding="utf-8"))
    kit["rgb"] = {k: RGBColor.from_string(v) for k, v in kit["colors"].items()}
    return kit


def mix(c1: RGBColor, c2: RGBColor, t: float) -> RGBColor:
    """두 브랜드 색의 중간색 (t=0 → c1, t=1 → c2). 고스트 숫자 등 저대비 장식용."""
    return RGBColor(*(round(a + (b - a) * t) for a, b in zip(c1, c2)))


def fit_picture(slide, path, x, y, max_w, max_h, *, align="left", line=None, line_w=0.01):
    """이미지를 **박스 안에 맞춰** 넣는다. 실제 파일 비율로 축소하고 폭·높이 상한을 둘 다 지킨다.

    반환: (left, top, w, h) 실제 배치 좌표.

    왜 이 헬퍼가 필요한가(2026-07-15 버그):
        `add_picture(..., height=Inches(4.75))`처럼 **한 축만** 주면 다른 축은 실제 이미지 비율이
        정한다. s10은 비율을 `3200/2000`(1.60)으로 **가정**하고 폭 7.6"를 전제로 옆 칼럼을 8.6"에
        놨는데, 실제 파일은 1899×926(**2.05**)이라 폭이 **9.74"**로 나와 텍스트를 1.74" 침범했다.
        게다가 그 `img_w`는 **계산만 하고 안 쓰는 죽은 변수**였다 — 가짜 값이 레이아웃을 구동했다.
        공장 레이아웃은 **어떤 비율이 와도** 안 깨져야 한다. 그래서 비율을 파일에서 읽고 박스에 가둔다.
    """
    from PIL import Image

    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(max_w / iw, max_h / ih)
    w, h = iw * scale, ih * scale
    left = x if align == "left" else x + (max_w - w) / 2
    pic = slide.shapes.add_picture(str(path), Inches(left), Inches(y), Inches(w), Inches(h))
    if line is not None:
        pic.line.color.rgb = line
        pic.line.width = Inches(line_w)
    return left, y, w, h


def add_box(slide, x, y, w, h, fill=None, line=None, line_w=None, shape="rect"):
    """민무늬 도형. shape: rect|oval|chevron|round|pentagon. fill/line은 RGBColor 또는 None(투명)."""
    from pptx.enum.shapes import MSO_SHAPE

    kinds = {
        "rect": MSO_SHAPE.RECTANGLE,
        "oval": MSO_SHAPE.OVAL,
        "chevron": MSO_SHAPE.CHEVRON,
        "round": MSO_SHAPE.ROUNDED_RECTANGLE,
        "pentagon": MSO_SHAPE.PENTAGON,
        "diamond": MSO_SHAPE.DIAMOND,
        "l_brace": MSO_SHAPE.LEFT_BRACE,
        "r_brace": MSO_SHAPE.RIGHT_BRACE,
    }
    shp = slide.shapes.add_shape(kinds[shape], Inches(x), Inches(y), Inches(w), Inches(h))
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w or 1)
    return shp


def set_shape_text(shp, text, size, font, color, *, bold=False):
    """도형 중앙에 단일 런 텍스트를 채운다 (번호 원·필 배지용)."""
    p = shp.text_frame.paragraphs[0]
    p._p.get_or_add_pPr().set("eaLnBrk", "0")  # 한글 어절 단위 줄바꿈 강제
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.language_id = MSO_LANGUAGE_ID.KOREAN


def add_text(
    slide,
    x,
    y,
    w,
    h,
    runs,
    size,
    font,
    color,
    *,
    bold=False,
    align=PP_ALIGN.LEFT,
    space_after=None,
    line_spacing=None,
    anchor=None,
):
    """텍스트 상자. runs: str 또는 [(text, {bold/color/size 오버라이드}), ...] 또는 문단 리스트."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0)
    if anchor is not None:
        tf.vertical_anchor = anchor
    paragraphs = runs if isinstance(runs, list) and runs and isinstance(runs[0], list) else [runs]
    for i, para in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p._p.get_or_add_pPr().set(
            "eaLnBrk", "0"
        )  # 한글 어절 단위 줄바꿈 강제 (글자 단위 깨짐 금지 — 골든 §5)
        p.alignment = align
        if space_after is not None:
            p.space_after = Pt(space_after)
        if line_spacing is not None:
            p.line_spacing = line_spacing
        for run in para if isinstance(para, list) else [(para, {})]:
            text, ov = run if isinstance(run, tuple) else (run, {})
            r = p.add_run()
            r.text = text
            r.font.name = ov.get("font", font)
            r.font.size = Pt(ov.get("size", size))
            r.font.bold = ov.get("bold", bold)
            r.font.color.rgb = ov.get("color", color)
            r.font.language_id = MSO_LANGUAGE_ID.KOREAN  # eaLnBrk가 먹으려면 런 언어 태그 필요
    return tb
