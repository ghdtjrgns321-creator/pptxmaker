"""S11 시안 — 기술 3 판단트리 (신규 어휘: 조건-분기 도해. UI 패스).

콘텐츠 실물: k-ifrs-1115/data/ontology/judgment_trees.json 트리 1번
"기간에 걸쳐 vs 한 시점 인식" — 문단 32·35·36·37·38·B3~B5·B9·B11 (창작 0).
선택 기준·실측 균열: 00_factsheet.md §B·§C (미주입 17/18).
실행: 이 타입은 goldenfab 레지스트리 경유로만 렌더된다 — 골든 19장 확인은
      `uv run python golden/build_golden.py`(2026-07-15 단일화로 시안 개별 실행 경로 폐지).
"""

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from . import grid as G
from .kit import SLIDE_H, SLIDE_W, add_box, add_text, load_kit, set_shape_text

K = load_kit()
C, S, F = K["rgb"], K["sizes"], K["fonts"]

KICKER = "3. 기술 설명 — TECH 03 · Generate"
SOURCE = (
    "출처: data/ontology/judgment_trees.json(트리 41 · 사용자 검수 완료) · "
    "3_KNOWLEDGE-GRAPH.md §3.3 · 4_SEARCH-PIPELINE.md L72·78"
)

NARRATIVES = [  # §8d 서사 3칼럼 — 두괄식 볼드 헤드 + 완전 문장
    (
        "왜 필요한가",
        "하나의 회계 판단이 문단 32·35·36·37과 B9~B13에 흩어져 있다. "
        "문단을 조각으로 검색해서는 판단의 순서가 복원되지 않는다.",
    ),
    (
        "파이프라인에서의 역할",
        "Generate 직전에 주입되는 답변의 골격. 질문에 걸린 트리가 "
        "판단 순서를 제공하므로 AI가 순서를 창작하지 않는다.",
    ),
    (
        "어떻게 만들었나",
        "기준서 본문의 조건-분기만 원문 앵커(문단 번호)와 함께 옮겨 "
        "41개를 미리 조립 — AI 창작 아님, 전 트리 사용자 검수 완료.",
    ),
]

CRITERIA = [  # 트리 1번 실물 — 문단 35의 3기준
    "고객이 효익을 동시에\n얻고 소비한다 (B3~B4)",
    "고객이 통제하는 자산을\n만들거나 가치를 높인다 (B5)",
    "대체용도 없음 + 지급청구권\n둘 다 충족 (문단 36·37) ★",
]


def header(slide, headline):
    add_text(
        slide, G.MARGIN_L, 0.42, 8.0, 0.28, KICKER, S["caption"], F["head"], C["muted"], bold=True
    )
    add_text(
        slide,
        G.MARGIN_L,
        0.72,
        G.RIGHT_EDGE - G.MARGIN_L,
        0.55,
        headline,
        S["section"],
        F["head"],
        C["primary"],
        bold=True,
    )
    add_box(slide, G.MARGIN_L, G.RULE_Y, G.RIGHT_EDGE - G.MARGIN_L, 0.014, fill=C["muted"])


def _subhead(slide, x, y, w, text):
    add_text(slide, x, y, w, 0.3, text, S["head"], F["head"], C["primary"], bold=True)
    add_box(slide, x, y + 0.34, w, 0.012, fill=C["muted"])


def _arrow(slide, x1, y1, x2, y2):
    from pptx.oxml.ns import qn

    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = C["muted"]
    conn.line.width = Pt(1.5)
    ln = conn.line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    return conn


def bar_and_source(slide, text):
    bar = add_box(slide, G.MARGIN_L, G.BAR_Y, G.RIGHT_EDGE - G.MARGIN_L, G.BAR_H, fill=C["primary"])
    tf = bar.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name, run.font.bold = F["head"], True
    run.font.size = Pt(S["body"])
    run.font.color.rgb = C["bg"]
    add_text(
        slide,
        G.MARGIN_L,
        G.SOURCE_Y,
        G.RIGHT_EDGE - G.MARGIN_L,
        0.3,
        SOURCE,
        S["foot"],
        F["body"],
        C["muted"],
    )


def _two_line_box(slide, x, y, w, h, head, sub, line=None, line_w=1.0, fill=None):
    box = add_box(slide, x, y, w, h, fill=fill or C["bg"], line=line or C["muted"], line_w=line_w)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.08)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    p1._p.get_or_add_pPr().set("eaLnBrk", "0")
    r1 = p1.add_run()
    r1.text = head
    r1.font.name, r1.font.bold = F["head"], True
    r1.font.size = Pt(S["body"])
    r1.font.color.rgb = C["primary"]
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2._p.get_or_add_pPr().set("eaLnBrk", "0")
    r2 = p2.add_run()
    r2.text = sub
    r2.font.name = F["body"]
    r2.font.size = Pt(S["caption"])
    r2.font.color.rgb = C["muted"]
    from pptx.enum.lang import MSO_LANGUAGE_ID

    for r in (r1, r2):
        r.font.language_id = MSO_LANGUAGE_ID.KOREAN
    return box


def variant_a(prs):
    """A — 조건-분기 도해(트리 1번 실물) + 서사 3칼럼 + 선택 기준 패널."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    header(slide, "판단트리 41개 — 흩어진 조건-분기를 판단 순서로 미리 조립")
    # ── 상단: 서사 3칼럼 ──
    nar_w = (G.RIGHT_EDGE - G.MARGIN_L - 2 * 0.4) / 3
    for i, (head, body) in enumerate(NARRATIVES):
        nx = G.MARGIN_L + i * (nar_w + 0.4)
        add_text(
            slide,
            nx,
            G.CONTENT_TOP,
            nar_w,
            0.28,
            f"0{i + 1}  {head}",
            S["head"],
            F["head"],
            C["primary"],
            bold=True,
        )
        add_text(
            slide,
            nx,
            G.CONTENT_TOP + 0.36,
            nar_w,
            0.95,
            body,
            S["caption"],
            F["body"],
            C["muted"],
            line_spacing=1.25,
        )
        if i < 2:
            add_box(slide, nx + nar_w + 0.2, G.CONTENT_TOP, 0.012, 1.2, fill=C["bg_alt"])
    # ── 하단 좌: 조건-분기 도해 — 트리 1번 「기간에 걸쳐 vs 한 시점 인식」 ──
    dx, dw = G.MARGIN_L, 7.9
    _subhead(slide, dx, 3.05, dw, "실물 — 트리 1번 「기간에 걸쳐 vs 한 시점 인식」 (41개 중)")
    _two_line_box(
        slide, dx, 3.75, 2.2, 0.85, "수익을 언제 인식하나", "문단 32 · 계약 개시시점에 판단"
    )
    _arrow(slide, dx + 2.2, 4.175, 3.35, 4.175)
    dia = add_box(
        slide, 3.4, 3.75, 1.5, 0.85, fill=C["bg"], line=C["accent"], line_w=1.75, shape="diamond"
    )
    set_shape_text(
        dia, "3기준 중\n하나라도 충족?", S["caption"], F["head"], C["primary"], bold=True
    )
    add_text(
        slide,
        3.25,
        4.66,
        1.8,
        0.24,
        "문단 35",
        S["caption"],
        F["body"],
        C["muted"],
        align=PP_ALIGN.CENTER,
    )
    # 예 → 기간에 걸쳐 / 아니오 → 한 시점
    _arrow(slide, 4.9, 4.05, 6.0, 3.85)
    _arrow(slide, 4.9, 4.3, 6.0, 4.85)
    add_text(slide, 5.25, 3.72, 0.7, 0.22, "예", S["caption"], F["head"], C["primary"], bold=True)
    add_text(
        slide, 5.15, 4.62, 0.9, 0.22, "아니오", S["caption"], F["head"], C["primary"], bold=True
    )
    _two_line_box(
        slide,
        6.0,
        3.5,
        2.5,
        0.7,
        "기간에 걸쳐 인식",
        "진행기준 적용",
        line=C["accent"],
        line_w=1.75,
    )
    _two_line_box(slide, 6.0, 4.55, 2.5, 0.7, "한 시점에 인식", "문단 38 · 통제 이전 지표 5개 종합")
    # 3기준 밴드
    add_text(
        slide,
        dx,
        5.38,
        dw,
        0.24,
        "3기준 — 어느 하나라도 충족하면 기간에 걸쳐 인식한다",
        S["caption"],
        F["body"],
        C["muted"],
    )
    crit_w, gap = 2.48, 0.23
    for i, crit in enumerate(CRITERIA):
        cx = dx + i * (crit_w + gap)
        box = add_box(slide, cx, 5.68, crit_w, 0.64, fill=C["bg_alt"])
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = Inches(0.08)
        from pptx.enum.lang import MSO_LANGUAGE_ID

        for li, seg in enumerate(crit.split("\n")):
            p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            p._p.get_or_add_pPr().set("eaLnBrk", "0")
            r = p.add_run()
            r.text = seg
            r.font.name = F["body"]
            r.font.size = Pt(S["caption"])
            r.font.color.rgb = C["primary"]
            r.font.language_id = MSO_LANGUAGE_ID.KOREAN
    # ── 하단 우: 선택 기준 패널 ──
    px = 8.9
    pw = G.RIGHT_EDGE - px
    add_box(slide, 8.7, 3.05, 0.012, G.CONTENT_BOTTOM - 3.05, fill=C["bg_alt"])
    _subhead(slide, px, 3.05, pw, "언제 걸리나 — 트리 선택 기준")
    rules = [
        (
            "개념 겹침",
            "질문에서 진입한 개념 집합과 트리의 발동 개념이 가장 많이 겹치는 트리를 고른다.",
        ),
        (
            "주제 직속",
            "질문 주제에 직접 걸린 트리는 전부 이어붙여 주입한다 — 구 시스템의 미주입 균열 17/18건을 구조로 해소.",
        ),
    ]
    for i, (head, body) in enumerate(rules):
        ry = 3.6 + i * 0.92
        add_text(
            slide, px, ry, 0.45, 0.26, f"0{i + 1}", S["head"], F["head"], C["accent"], bold=True
        )
        add_text(
            slide,
            px + 0.42,
            ry,
            pw - 0.42,
            0.26,
            head,
            S["body"],
            F["head"],
            C["primary"],
            bold=True,
        )
        add_text(
            slide,
            px + 0.42,
            ry + 0.3,
            pw - 0.42,
            0.55,
            body,
            S["caption"],
            F["body"],
            C["muted"],
            line_spacing=1.2,
        )
    star = add_box(slide, px, 5.52, pw, 0.8, fill=C["bg_alt"])
    tf = star.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.1)
    p = tf.paragraphs[0]
    p._p.get_or_add_pPr().set("eaLnBrk", "0")
    r = p.add_run()
    r.text = (
        "★ 기준 ③의 핵심 분기 — 지급청구권은 '고객 종료권 유무'로 갈린다. "
        "있으면 문단 37·B9, 없으면 B11. 트리가 이 깊이까지 원문 앵커로 담는다."
    )
    r.font.name = F["body"]
    r.font.size = Pt(S["caption"])
    r.font.color.rgb = C["primary"]
    from pptx.enum.lang import MSO_LANGUAGE_ID

    r.font.language_id = MSO_LANGUAGE_ID.KOREAN
    bar_and_source(
        slide,
        "판단의 순서는 검색하지 않는다 — 기준서의 조건-분기 41개를 원문 앵커와 함께 미리 조립해 주입",
    )
    return slide


STEPS = [  # 진입 규칙 3단계 — 4_SEARCH-PIPELINE.md §4.2·§4.4 (After 현행)
    (
        "주제 지목",
        "질문 분석 AI가 35개 주제 목록에서 최대 3개를 지목한다 — 개념 80개를 직접 고르지 않는다.",
    ),
    (
        "직속 개념 매칭",
        "지목된 주제의 직속 개념과 트리의 발동 개념이 겹치면 그 트리가 걸린다 — 딸려온 배경 개념은 제외.",
    ),
    (
        "걸린 트리 전부 주입",
        "1개만 고르지 않고 걸린 판단 절차를 모두 이어붙여 넣는다. 직속 개념이 없으면 진입 개념 전체로 폴백.",
    ),
]


def _multi_caption(slide, x, y, w, h, lines, fill=None):
    """캡션 크기 여러 줄 텍스트 박스 — 어절 줄바꿈 보존."""
    from pptx.enum.lang import MSO_LANGUAGE_ID

    if fill is not None:
        box = add_box(slide, x, y, w, h, fill=fill)
        tf = box.text_frame
        tf.margin_left = tf.margin_right = Inches(0.1)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    else:
        box = add_text(slide, x, y, w, h, "", S["caption"], F["body"], C["muted"])
        tf = box.text_frame
    tf.word_wrap = True
    for li, (seg, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
        p._p.get_or_add_pPr().set("eaLnBrk", "0")
        p.line_spacing = 1.25
        r = p.add_run()
        r.text = seg
        r.font.name = F["body"]
        r.font.size = Pt(S["caption"])
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.language_id = MSO_LANGUAGE_ID.KOREAN
    return box


def variant_b(prs):
    """B — 진입 규칙이 주인공: 3단계 플로우 + Before/After 좌측 대형, 실물 트리 우측 컴팩트."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    header(slide, "판단트리 41개 — 흩어진 조건-분기를 판단 순서로 미리 조립")
    nar_w = (G.RIGHT_EDGE - G.MARGIN_L - 2 * 0.4) / 3
    for i, (head, body) in enumerate(NARRATIVES):
        nx = G.MARGIN_L + i * (nar_w + 0.4)
        add_text(
            slide,
            nx,
            G.CONTENT_TOP,
            nar_w,
            0.28,
            f"0{i + 1}  {head}",
            S["head"],
            F["head"],
            C["primary"],
            bold=True,
        )
        add_text(
            slide,
            nx,
            G.CONTENT_TOP + 0.36,
            nar_w,
            0.95,
            body,
            S["caption"],
            F["body"],
            C["muted"],
            line_spacing=1.25,
        )
        if i < 2:
            add_box(slide, nx + nar_w + 0.2, G.CONTENT_TOP, 0.012, 1.2, fill=C["bg_alt"])
    # ── 하단 좌: 진입 규칙 3단계 + Before/After ──
    lx, lw = G.MARGIN_L, 7.5
    _subhead(slide, lx, 3.05, lw, "트리는 이렇게 걸린다 — 주제기반 다중주입 (재설계 후 현행)")
    step_w, gap = 2.3, 0.3
    for i, (head, body) in enumerate(STEPS):
        sx = lx + i * (step_w + gap)
        box = add_box(slide, sx, 3.65, step_w, 1.5, fill=C["bg"], line=C["muted"], line_w=1.0)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.08)
        from pptx.enum.lang import MSO_LANGUAGE_ID

        p1 = box.text_frame.paragraphs[0]
        p1._p.get_or_add_pPr().set("eaLnBrk", "0")
        rn = p1.add_run()
        rn.text = f"0{i + 1}  "
        rn.font.name, rn.font.bold = F["head"], True
        rn.font.size = Pt(S["body"])
        rn.font.color.rgb = C["muted"]
        rh = p1.add_run()
        rh.text = head
        rh.font.name, rh.font.bold = F["head"], True
        rh.font.size = Pt(S["body"])
        rh.font.color.rgb = C["primary"]
        p2 = tf.add_paragraph()
        p2._p.get_or_add_pPr().set("eaLnBrk", "0")
        p2.line_spacing = 1.25
        rb = p2.add_run()
        rb.text = body
        rb.font.name = F["body"]
        rb.font.size = Pt(S["caption"])
        rb.font.color.rgb = C["muted"]
        for r in (rn, rh, rb):
            r.font.language_id = MSO_LANGUAGE_ID.KOREAN
        if i < 2:
            _arrow(slide, sx + step_w, 4.4, sx + step_w + gap, 4.4)
    # 적용 예시 — 실물 경로 (aliases → topic_concept_map → judgment_trees 실제 항목)
    add_text(
        slide,
        lx,
        5.35,
        lw,
        0.24,
        "적용 예시 — 실물 데이터로 이어지는 경로",
        S["caption"],
        F["head"],
        C["primary"],
        bold=True,
    )
    ex_labels = ["질문", "지목된 주제 (35개 목록 중)", "걸린 트리 (41개 중)"]
    ex_chips = ["“볼륨디스카운트 조항이 있는 계약은…”", "변동대가", "「변동대가 (추정)」 주입"]
    ex_x, ex_ws = [lx, lx + 3.05, lx + 5.15], [2.75, 1.8, 2.35]
    for i in range(3):
        add_text(
            slide,
            ex_x[i],
            5.66,
            ex_ws[i],
            0.22,
            ex_labels[i],
            S["caption"],
            F["body"],
            C["muted"],
            align=PP_ALIGN.CENTER,
        )
        chip = add_box(
            slide,
            ex_x[i],
            5.92,
            ex_ws[i],
            0.42,
            fill=C["bg"] if i == 0 else C["primary"],
            line=C["muted"] if i == 0 else None,
            line_w=1.0,
            shape="round",
        )
        set_shape_text(
            chip,
            ex_chips[i],
            S["caption"],
            F["body"] if i == 0 else F["head"],
            C["primary"] if i == 0 else C["bg"],
            bold=i > 0,
        )
        if i < 2:
            _arrow(slide, ex_x[i] + ex_ws[i], 6.13, ex_x[i + 1] - 0.05, 6.13)
    # ── 하단 우: 실물 트리 컴팩트 ──
    add_box(slide, 8.3, 3.05, 0.012, G.CONTENT_BOTTOM - 3.05, fill=C["bg_alt"])
    tx = 8.5
    tw = G.RIGHT_EDGE - tx  # 4.233
    _subhead(slide, tx, 3.05, tw, "걸린 트리의 실물 — 「기간에 걸쳐 vs 한 시점」")
    _two_line_box(slide, tx, 3.6, tw, 0.5, "수익을 언제 인식하나", "문단 32 · 계약 개시시점에 판단")
    dia_w, dia_h = 1.3, 0.75
    dia_x = tx + (tw - dia_w) / 2
    _arrow(slide, tx + tw / 2, 4.1, tx + tw / 2, 4.28)
    dia = add_box(
        slide,
        dia_x,
        4.3,
        dia_w,
        dia_h,
        fill=C["bg"],
        line=C["accent"],
        line_w=1.75,
        shape="diamond",
    )
    set_shape_text(dia, "3기준 충족?", S["caption"], F["head"], C["primary"], bold=True)
    add_text(
        slide,
        dia_x + dia_w + 0.08,
        4.55,
        1.0,
        0.24,
        "문단 35",
        S["caption"],
        F["body"],
        C["muted"],
    )
    _arrow(slide, dia_x + 0.2, 5.0, tx + 0.9, 5.25)
    _arrow(slide, dia_x + dia_w - 0.2, 5.0, tx + tw - 0.9, 5.25)
    add_text(
        slide, tx + 0.35, 5.02, 0.5, 0.2, "예", S["caption"], F["head"], C["primary"], bold=True
    )
    add_text(
        slide,
        tx + tw - 1.0,
        5.02,
        0.7,
        0.2,
        "아니오",
        S["caption"],
        F["head"],
        C["primary"],
        bold=True,
    )
    _two_line_box(
        slide, tx, 5.3, 2.0, 0.6, "기간에 걸쳐", "진행기준 적용", line=C["accent"], line_w=1.75
    )
    _two_line_box(slide, tx + tw - 2.0, 5.3, 2.0, 0.6, "한 시점", "문단 38 · 지표 5개")
    _multi_caption(
        slide,
        tx,
        6.0,
        tw,
        0.35,
        [
            (
                "3기준: ①동시 효익·소비(B3~B4) ②고객 통제 자산(B5) ③대체용도 없음+지급청구권"
                "(문단 36·37) — ★③은 고객 종료권 유무로 갈린다(B9/B11)",
                C["muted"],
                False,
            ),
        ],
    )
    bar_and_source(
        slide, "질문이 여러 판단에 걸치면 걸린 절차를 모두 넣는다 — 트리 41개, 전부 원문 앵커"
    )
    return slide


def variant_c(prs):
    """C — 적용 예시·트리 도해 통합: 우측 세로 흐름(질문→주제→트리→분기), 좌하 본문 원문 발췌."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    header(slide, "판단트리 41개 — 흩어진 조건-분기를 판단 순서로 미리 조립")
    nar_w = (G.RIGHT_EDGE - G.MARGIN_L - 2 * 0.4) / 3
    for i, (head, body) in enumerate(NARRATIVES):
        nx = G.MARGIN_L + i * (nar_w + 0.4)
        add_text(
            slide,
            nx,
            G.CONTENT_TOP,
            nar_w,
            0.28,
            f"0{i + 1}  {head}",
            S["head"],
            F["head"],
            C["primary"],
            bold=True,
        )
        add_text(
            slide,
            nx,
            G.CONTENT_TOP + 0.36,
            nar_w,
            0.95,
            body,
            S["caption"],
            F["body"],
            C["muted"],
            line_spacing=1.25,
        )
        if i < 2:
            add_box(slide, nx + nar_w + 0.2, G.CONTENT_TOP, 0.012, 1.2, fill=C["bg_alt"])
    # ── 하단 좌: 진입 규칙 3단계 (vB 동일) ──
    lx, lw = G.MARGIN_L, 7.5
    _subhead(slide, lx, 3.05, lw, "트리는 이렇게 걸린다 — 주제기반 다중주입")
    step_w, gap = 2.3, 0.3
    from pptx.enum.lang import MSO_LANGUAGE_ID

    for i, (head, body) in enumerate(STEPS):
        sx = lx + i * (step_w + gap)
        box = add_box(slide, sx, 3.65, step_w, 1.5, fill=C["bg"], line=C["muted"], line_w=1.0)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.08)
        p1 = tf.paragraphs[0]
        p1._p.get_or_add_pPr().set("eaLnBrk", "0")
        rn = p1.add_run()
        rn.text = f"0{i + 1}  "
        rn.font.name, rn.font.bold = F["head"], True
        rn.font.size = Pt(S["body"])
        rn.font.color.rgb = C["muted"]
        rh = p1.add_run()
        rh.text = head
        rh.font.name, rh.font.bold = F["head"], True
        rh.font.size = Pt(S["body"])
        rh.font.color.rgb = C["primary"]
        p2 = tf.add_paragraph()
        p2._p.get_or_add_pPr().set("eaLnBrk", "0")
        p2.line_spacing = 1.25
        rb = p2.add_run()
        rb.text = body
        rb.font.name = F["body"]
        rb.font.size = Pt(S["caption"])
        rb.font.color.rgb = C["muted"]
        for r in (rn, rh, rb):
            r.font.language_id = MSO_LANGUAGE_ID.KOREAN
        if i < 2:
            _arrow(slide, sx + step_w, 4.4, sx + step_w + gap, 4.4)
    # ── 하단 좌 아래: 트리 본문 원문 발췌 카드 (judgment_trees.json variable-consid) ──
    add_text(
        slide,
        lx,
        5.3,
        lw,
        0.24,
        "트리 본문 실물 — judgment_trees.json 「변동대가 (추정)」 발췌 (문단 번호가 원문 앵커)",
        S["caption"],
        F["head"],
        C["primary"],
        bold=True,
    )
    add_box(slide, lx, 5.6, 0.05, 0.75, fill=C["accent"])
    _multi_caption(
        slide,
        lx + 0.2,
        5.6,
        lw - 0.2,
        0.75,
        [
            (
                "“변동대가 추정 판단(문단 50): 약속한 대가에 변동금액이 포함되면(할인·리베이트·환불…, 문단 51), "
                "문단 53의 두 방법 중 받을 권리 대가를 더 잘 예측할 것으로 예상하는 방법으로 추정한다 — "
                "(1) 기댓값 (2) 가능성이 가장 높은 금액. 계약 전체에 하나의 방법을 일관 적용한다(문단 54).”",
                C["text"],
                False,
            ),
        ],
    )
    # ── 하단 우: 통합 세로 흐름 — 질문 → 주제 → 트리 → 분기 ──
    add_box(slide, 8.3, 3.05, 0.012, G.CONTENT_BOTTOM - 3.05, fill=C["bg_alt"])
    tx = 8.5
    tw = G.RIGHT_EDGE - tx
    cx = tx + tw / 2
    _subhead(slide, tx, 3.05, tw, "적용 실물 — 질문에서 트리 분기까지")
    q = add_box(slide, tx, 3.5, tw, 0.46, fill=C["bg"], line=C["muted"], line_w=1.0, shape="round")
    set_shape_text(q, "“볼륨디스카운트 조항이 있는 계약은…”", S["caption"], F["body"], C["primary"])
    _arrow(slide, cx, 3.96, cx, 4.14)
    add_text(
        slide,
        cx + 0.15,
        3.94,
        tw / 2 - 0.15,
        0.2,
        "주제 지목 — 변동대가 (35개 중)",
        S["caption"],
        F["body"],
        C["muted"],
    )
    chip = add_box(slide, tx + 0.35, 4.14, tw - 0.7, 0.4, fill=C["primary"], shape="round")
    set_shape_text(
        chip, "걸린 트리 「변동대가 (추정)」 · 41개 중", S["caption"], F["head"], C["bg"], bold=True
    )
    _arrow(slide, cx, 4.54, cx, 4.7)
    dia_w, dia_h = 1.5, 0.66
    dia = add_box(
        slide,
        cx - dia_w / 2,
        4.7,
        dia_w,
        dia_h,
        fill=C["bg"],
        line=C["accent"],
        line_w=1.75,
        shape="diamond",
    )
    set_shape_text(dia, "변동금액 포함?", S["caption"], F["head"], C["primary"], bold=True)
    add_text(
        slide,
        cx + dia_w / 2 + 0.08,
        4.92,
        1.1,
        0.22,
        "문단 51",
        S["caption"],
        F["body"],
        C["muted"],
    )
    _arrow(slide, cx - 0.35, 5.3, tx + 1.0, 5.52)
    _arrow(slide, cx + 0.35, 5.3, tx + tw - 1.0, 5.52)
    _two_line_box(slide, tx, 5.52, 2.0, 0.52, "기댓값", "확률 가중 합")
    _two_line_box(slide, tx + tw - 2.0, 5.52, 2.0, 0.52, "가능성 최고 금액", "단일 결과치")
    _multi_caption(
        slide,
        tx,
        6.08,
        tw,
        0.27,
        [
            (
                "문단 53 — 더 잘 예측하는 방법 하나를 일관 적용(문단 54), 이후 제약 적용(문단 56~58)",
                C["muted"],
                False,
            ),
        ],
    )
    bar_and_source(
        slide, "질문이 여러 판단에 걸치면 걸린 절차를 모두 넣는다 — 트리 41개, 전부 원문 앵커"
    )
    return slide


DEFAULT = {  # 텍스트 내용만(좌표·색·크기·도형은 코드 고정). c=None이면 골든 기본값.
    "kicker": KICKER,
    "headline": "판단트리 41개 — 흩어진 조건-분기를 판단 순서로 미리 조립",
    "narratives": NARRATIVES,
    "steps_head": "트리는 이렇게 걸린다 — 주제기반 다중주입",
    "steps": STEPS,
    "branch_label": "주입 뒤 — 답변은 질문 성격에 따라 세 갈래로 갈린다",
    "flow_head": "적용 실물 — 질문에서 트리 분기까지",
    "question": "“볼륨디스카운트 조항이 있는 계약은…”",
    "topic_label": "주제 지목 — 변동대가",
    "tree_chip": "「변동대가 (추정)」 걸림",
    "diamond": "변동금액 포함?",
    "diamond_anchor": "문단 51",
    "method1_head": "기댓값",
    "method1_sub": "확률 가중 합",
    "method2_head": "가능성 최고 금액",
    "method2_sub": "단일 결과치",
    "flow_foot": "문단 53 — 더 잘 예측하는 방법 하나를 일관 적용(문단 54)",
    "bar": "질문이 여러 판단에 걸치면 걸린 절차를 모두 넣는다 — 트리 41개, 전부 원문 앵커",
    "source": SOURCE,
}


def variant_d(prs, c=None):
    """D — vC 개정: 좌하 '주입 뒤 세 갈래' 도해(§4.4), 우측 칩 축소로 밀도 완화.

    c: 텍스트 내용 override(None=골든 기본값). 좌표·색·크기·도형은 고정.
    header()/bar_and_source()는 kicker·source를 c에서 받기 위해 인라인(출력 동일).
    """
    c = {**DEFAULT, **(c or {})}
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    add_text(
        slide,
        G.MARGIN_L,
        0.42,
        8.0,
        0.28,
        c["kicker"],
        S["caption"],
        F["head"],
        C["muted"],
        bold=True,
    )
    add_text(
        slide,
        G.MARGIN_L,
        0.72,
        G.RIGHT_EDGE - G.MARGIN_L,
        0.55,
        c["headline"],
        S["section"],
        F["head"],
        C["primary"],
        bold=True,
    )
    add_box(slide, G.MARGIN_L, G.RULE_Y, G.RIGHT_EDGE - G.MARGIN_L, 0.014, fill=C["muted"])
    nar_w = (G.RIGHT_EDGE - G.MARGIN_L - 2 * 0.4) / 3
    for i, (head, body) in enumerate(c["narratives"]):
        nx = G.MARGIN_L + i * (nar_w + 0.4)
        add_text(
            slide,
            nx,
            G.CONTENT_TOP,
            nar_w,
            0.28,
            [(f"0{i + 1}", {"color": C["accent"]}), (f"  {head}", {})],
            S["head"],
            F["head"],
            C["primary"],
            bold=True,
        )
        add_text(
            slide,
            nx,
            G.CONTENT_TOP + 0.36,
            nar_w,
            0.95,
            body,
            S["caption"],
            F["body"],
            C["muted"],
            line_spacing=1.25,
        )
        if i < 2:
            add_box(slide, nx + nar_w + 0.2, G.CONTENT_TOP, 0.012, 1.2, fill=C["bg_alt"])
    # ── 하단 좌: 진입 규칙 단계(G10 numbered_steps) — 폭은 단계 수에서 파생 ──
    lx, lw = G.MARGIN_L, 7.5
    _subhead(slide, lx, 3.05, lw, c["steps_head"])
    steps = c["steps"]
    gap = 0.3
    # 최소 폭 2.30 = 이 상자(h 1.1)가 실제로 수용하는 폭. 렌더 실측(2026-07-25): 폭 1.725"로
    # 좁히면 12pt 헤드가 2줄, 9pt 본문이 5줄이 되어 필요 높이 1.33" > 상자 1.1" — 텍스트가
    # 아래로 흘러 이웃을 덮는다. `check_text_overflow`는 글자폭을 0.62×pt로 근사해 **한글에서
    # 이 넘침을 못 잡는다**(미탐 실측). 그래서 기하 상한으로 막고, 4단 이상은 세로 나열이나
    # adapted 변형으로 내려보낸다 — 조용히 깨진 4단을 내놓지 않는다.
    step_w = G.track(len(steps), lx, lx + lw, gap, 2.30, what="진입 단계")
    from pptx.enum.lang import MSO_LANGUAGE_ID

    for i, (head, body) in enumerate(steps):
        sx = lx + i * (step_w + gap)
        box = add_box(slide, sx, 3.6, step_w, 1.1, fill=C["bg"], line=C["muted"], line_w=1.0)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.08)
        p1 = tf.paragraphs[0]
        p1._p.get_or_add_pPr().set("eaLnBrk", "0")
        rn = p1.add_run()
        rn.text = f"0{i + 1}  "
        rn.font.name, rn.font.bold = F["head"], True
        rn.font.size = Pt(S["body"])
        rn.font.color.rgb = C["muted"]
        rh = p1.add_run()
        rh.text = head
        rh.font.name, rh.font.bold = F["head"], True
        rh.font.size = Pt(S["body"])
        rh.font.color.rgb = C["primary"]
        p2 = tf.add_paragraph()
        p2._p.get_or_add_pPr().set("eaLnBrk", "0")
        p2.line_spacing = 1.25
        rb = p2.add_run()
        rb.text = body
        rb.font.name = F["body"]
        rb.font.size = Pt(S["caption"])
        rb.font.color.rgb = C["muted"]
        for r in (rn, rh, rb):
            r.font.language_id = MSO_LANGUAGE_ID.KOREAN
        if i < len(steps) - 1:
            _arrow(slide, sx + step_w, 4.15, sx + step_w + gap, 4.15)
    # ── 하단 좌 아래: 주입 뒤 — 답변 세 갈래, 1:3 flat tree (스냅 스펙 dy=+0.80) ──
    add_text(
        slide,
        lx,
        4.82,
        lw,
        0.22,
        c["branch_label"],
        S["caption"],
        F["head"],
        C["primary"],
        bold=True,
    )
    from .s11_branch_snap import audit as branch_audit
    from .s11_branch_snap import build_branch

    br_shapes, br_centers = build_branch(slide, dy=0.80)
    branch_audit(br_shapes, br_centers)
    chip1_top = br_shapes["chips"][0].top / 914400
    label_gap = chip1_top - (4.82 + 0.22)
    print(f"label-chip gap {label_gap:.3f}")
    assert label_gap >= 0.12, f"LABEL GAP FAIL {label_gap:.3f}"
    # ── 하단 우: 통합 세로 흐름 (vC 대비 칩 축소·여백 확대) ──
    add_box(slide, 8.3, 3.05, 0.012, G.CONTENT_BOTTOM - 3.05, fill=C["bg_alt"])
    tx = 8.5
    tw = G.RIGHT_EDGE - tx
    cx = tx + tw / 2
    _subhead(slide, tx, 3.05, tw, c["flow_head"])
    q = add_box(
        slide,
        tx + 0.3,
        3.55,
        tw - 0.6,
        0.4,
        fill=C["bg"],
        line=C["muted"],
        line_w=1.0,
        shape="round",
    )
    set_shape_text(q, c["question"], S["caption"], F["body"], C["primary"])
    _arrow(slide, cx, 3.95, cx, 4.22)
    add_text(
        slide,
        cx + 0.15,
        3.97,
        tw / 2 - 0.15,
        0.2,
        c["topic_label"],
        S["caption"],
        F["body"],
        C["muted"],
    )
    chip = add_box(slide, tx + 0.85, 4.22, tw - 1.7, 0.38, fill=C["primary"], shape="round")
    set_shape_text(chip, c["tree_chip"], S["caption"], F["head"], C["bg"], bold=True)
    _arrow(slide, cx, 4.6, cx, 4.8)
    dia_w, dia_h = 1.5, 0.62
    dia = add_box(
        slide,
        cx - dia_w / 2,
        4.8,
        dia_w,
        dia_h,
        fill=C["bg"],
        line=C["accent"],
        line_w=1.75,
        shape="diamond",
    )
    set_shape_text(dia, c["diamond"], S["caption"], F["head"], C["primary"], bold=True)
    add_text(
        slide,
        cx + dia_w / 2 + 0.08,
        5.0,
        1.1,
        0.22,
        c["diamond_anchor"],
        S["caption"],
        F["body"],
        C["muted"],
    )
    _arrow(slide, cx - 0.35, 5.38, tx + 1.0, 5.62)
    _arrow(slide, cx + 0.35, 5.38, tx + tw - 1.0, 5.62)
    _two_line_box(slide, tx + 0.1, 5.62, 1.85, 0.5, c["method1_head"], c["method1_sub"])
    _two_line_box(slide, tx + tw - 1.95, 5.62, 1.85, 0.5, c["method2_head"], c["method2_sub"])
    add_text(
        slide,
        tx,
        6.14,
        tw,
        0.21,
        c["flow_foot"],
        S["caption"],
        F["body"],
        C["muted"],
        align=PP_ALIGN.CENTER,
    )
    # bar_and_source() 인라인 — bar·source를 c에서 받음(도형·좌표·색 동일)
    bar = add_box(slide, G.MARGIN_L, G.BAR_Y, G.RIGHT_EDGE - G.MARGIN_L, G.BAR_H, fill=C["primary"])
    tfbar = bar.text_frame
    tfbar.vertical_anchor = MSO_ANCHOR.MIDDLE
    pbar = tfbar.paragraphs[0]
    pbar.alignment = PP_ALIGN.CENTER
    rbar = pbar.add_run()
    rbar.text = c["bar"]
    rbar.font.name, rbar.font.bold = F["head"], True
    rbar.font.size = Pt(S["body"])
    rbar.font.color.rgb = C["bg"]
    add_text(
        slide,
        G.MARGIN_L,
        G.SOURCE_Y,
        G.RIGHT_EDGE - G.MARGIN_L,
        0.3,
        c["source"],
        S["foot"],
        F["body"],
        C["muted"],
    )
    return slide


def audit(prs):
    from pptx.oxml.ns import qn

    fails = []
    for si, sl in enumerate(prs.slides):
        accents, diamonds = 0, 0
        for sh in sl.shapes:
            t, le = sh.top / 914400, sh.left / 914400
            b, r = t + sh.height / 914400, le + sh.width / 914400
            full_bleed = sh.width / 914400 >= SLIDE_W - 0.05  # 배경 박스는 제외
            if t < 6.4 and b > 6.37 and not full_bleed:  # 콘텐츠가 바닥선 침범
                fails.append((si, "bottom", round(b, 2), sh.shape_id))
            if r > 12.75 and not full_bleed:
                fails.append((si, "right", round(r, 2), sh.shape_id))
            for el in sh._element.iter(qn("a:prstGeom")):
                if el.get("prst") == "diamond":
                    diamonds += 1
            for el in sh._element.iter(qn("a:srgbClr")):
                if el.get("val") == "D66E3A":
                    parent = el.getparent().tag
                    if parent.endswith("}solidFill"):
                        accents += 1
        if accents > 4:
            fails.append((si, "accent>4", accents))
        if diamonds < 1:
            fails.append((si, "diamond<1", diamonds))
    assert not fails, f"AUDIT FAIL {fails}"
    print("audit pass")


def main():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    variant_a(prs)
    variant_b(prs)
    variant_c(prs)
    variant_d(prs)
    audit(prs)
    out = Path(__file__).parent / "variants" / "s11_variants.pptx"
    prs.save(str(out))
    print(f"saved {out}")


if __name__ == "__main__":
    main()
