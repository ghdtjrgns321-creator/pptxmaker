"""S12 시안 — 기술 4 구조화 출력 (신규 어휘: 스키마 카드 + 거부-재시도 루프. UI 패스).

콘텐츠 실물: k-ifrs-1115/app/agents.py ClarifyOutput(필드 7)·output_validator(빈 인용/분기
거부 → ModelRetry), 5_INTERFACE.md(Split View·답변 유형), 7_JOURNEY.md(형식 분산 문제).
실행: uv run python golden/s12_variants.py → golden/variants/s12_variants.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from . import grid as G
from .kit import SLIDE_H, SLIDE_W, add_box, add_text, load_kit, set_shape_text

K = load_kit()
C, S, F = K["rgb"], K["sizes"], K["fonts"]

KICKER = "3. 기술 설명 — TECH 04 · Generate"
SOURCE = "출처: app/agents.py (ClarifyOutput · output_validator) · 5_INTERFACE.md · 7_JOURNEY.md"

NARRATIVES = [
    (
        "왜 필요한가",
        "같은 질문에 매번 다른 형식으로 답하면 정량 평가가 불가능하다. "
        "자유형 텍스트를 폐기하고 답변의 칸을 고정해야 한다.",
    ),
    (
        "파이프라인에서의 역할",
        "Generate의 출력 계약. 모든 답변이 결론·인용·후속질문의 같은 칸으로 "
        "나오므로 화면 배치와 채점이 기계적으로 가능해진다.",
    ),
    (
        "어떻게 만들었나",
        "답변 형식을 코드(스키마)로 선언하고, 검증기가 빈 인용·빈 분기를 "
        "거부해 자동 재시도 — 형식 미달 답변은 화면에 도달하지 못한다.",
    ),
]

FIELDS = [  # app/agents.py ClarifyOutput 실물 필드 7개 (설명은 Field description 축약)
    ("selected_branches", "결론 가이드에서 고른 분기 — 조건부면 전부, 확정이면 1개"),
    ("answer", "답변 본문 (마크다운)"),
    ("cited_paragraphs", "인용한 기준서 문단 번호 — 비면 검증기가 거부"),
    ("cited_cases", "인용한 질의회신·감리지적사례 ID"),
    ("cited_ie", "인용한 적용사례(IE) 번호"),
    ("follow_up_questions", "조건을 좁힐 후속 질문 3개 — 결론 후엔 자동 클리어"),
    ("is_conclusion", "결론 포함 여부 — 항상 True 강제"),
]


def header(slide, headline):
    add_text(
        slide, G.MARGIN_L, 0.42, 8.0, 0.28, KICKER, S["caption"], F["head"], C["muted"], bold=True
    )
    add_text(
        slide,
        G.MARGIN_L,
        0.72,
        G.RIGHT_EDGE - G.MARGIN_L,
        0.55,
        headline,
        S["section"],
        F["head"],
        C["primary"],
        bold=True,
    )
    add_box(slide, G.MARGIN_L, G.RULE_Y, G.RIGHT_EDGE - G.MARGIN_L, 0.014, fill=C["muted"])


def _subhead(slide, x, y, w, text):
    add_text(slide, x, y, w, 0.3, text, S["head"], F["head"], C["primary"], bold=True)
    add_box(slide, x, y + 0.34, w, 0.012, fill=C["muted"])


def _arrow(slide, x1, y1, x2, y2, head=True):
    from pptx.oxml.ns import qn

    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = C["muted"]
    conn.line.width = Pt(1.5)
    if head:
        ln = conn.line._get_or_add_ln()
        ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    return conn


def bar_and_source(slide, text):
    bar = add_box(slide, G.MARGIN_L, G.BAR_Y, G.RIGHT_EDGE - G.MARGIN_L, G.BAR_H, fill=C["primary"])
    tf = bar.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name, run.font.bold = F["head"], True
    run.font.size = Pt(S["body"])
    run.font.color.rgb = C["bg"]
    add_text(
        slide,
        G.MARGIN_L,
        G.SOURCE_Y,
        G.RIGHT_EDGE - G.MARGIN_L,
        0.3,
        SOURCE,
        S["foot"],
        F["body"],
        C["muted"],
    )


def _flow_box(slide, x, y, w, h, head, sub=None, line=None, line_w=1.0, fill=None):
    from pptx.enum.lang import MSO_LANGUAGE_ID

    box = add_box(slide, x, y, w, h, fill=fill or C["bg"], line=line or C["muted"], line_w=line_w)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.08)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    p1._p.get_or_add_pPr().set("eaLnBrk", "0")
    r1 = p1.add_run()
    r1.text = head
    r1.font.name, r1.font.bold = F["head"], True
    r1.font.size = Pt(S["body"])
    r1.font.color.rgb = C["primary"]
    r1.font.language_id = MSO_LANGUAGE_ID.KOREAN
    if sub:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2._p.get_or_add_pPr().set("eaLnBrk", "0")
        r2 = p2.add_run()
        r2.text = sub
        r2.font.name = F["body"]
        r2.font.size = Pt(S["caption"])
        r2.font.color.rgb = C["muted"]
        r2.font.language_id = MSO_LANGUAGE_ID.KOREAN
    return box


def variant_a(prs):
    """A — 좌 스키마 코드 카드(실물 필드 7) + 우 거부-재시도 루프 도해."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    header(slide, "구조화 출력 — 답변의 형식을 코드가 강제한다")
    nar_w = (G.RIGHT_EDGE - G.MARGIN_L - 2 * 0.4) / 3
    for i, (head, body) in enumerate(NARRATIVES):
        nx = G.MARGIN_L + i * (nar_w + 0.4)
        add_text(
            slide,
            nx,
            G.CONTENT_TOP,
            nar_w,
            0.28,
            f"0{i + 1}  {head}",
            S["head"],
            F["head"],
            C["primary"],
            bold=True,
        )
        add_text(
            slide,
            nx,
            G.CONTENT_TOP + 0.36,
            nar_w,
            0.95,
            body,
            S["caption"],
            F["body"],
            C["muted"],
            line_spacing=1.25,
        )
        if i < 2:
            add_box(slide, nx + nar_w + 0.2, G.CONTENT_TOP, 0.012, 1.2, fill=C["bg_alt"])
    # ── 하단 좌: 스키마 코드 카드 (다크 모노) ──
    lx, lw = G.MARGIN_L, 6.1
    _subhead(slide, lx, 3.05, lw, "출력 스키마 실물 — ClarifyOutput (app/agents.py)")
    card = add_box(slide, lx, 3.6, lw, 2.72, fill=C["primary"], shape="round")
    card.adjustments[0] = 0.04
    tf = card.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.14)
    from pptx.enum.lang import MSO_LANGUAGE_ID

    for i, (name, gloss) in enumerate(FIELDS):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p._p.get_or_add_pPr().set("eaLnBrk", "0")
        p.line_spacing = 1.55
        r1 = p.add_run()
        r1.text = name
        r1.font.name = "Consolas"
        r1.font.size = Pt(S["caption"] + 1)
        r1.font.bold = True
        r1.font.color.rgb = C["bg"]
        r2 = p.add_run()
        r2.text = f"   {gloss}"
        r2.font.name = F["body"]
        r2.font.size = Pt(S["caption"])
        r2.font.color.rgb = C["bg_alt"]
        r2.font.language_id = MSO_LANGUAGE_ID.KOREAN
    # ── 하단 우: 거부-재시도 루프 도해 ──
    add_box(slide, 6.9, 3.05, 0.012, G.CONTENT_BOTTOM - 3.05, fill=C["bg_alt"])
    tx = 7.15
    tw = G.RIGHT_EDGE - tx  # 5.58
    _subhead(slide, tx, 3.05, tw, "형식을 어기면 — 거부와 자동 재시도")
    fcx = tx + 2.3  # 흐름 축 (루프백이 구분선 오른쪽에 머물도록)
    _flow_box(slide, fcx - 1.2, 3.62, 2.4, 0.5, "답변 생성", "판단트리·근거 주입")
    _arrow(slide, fcx, 4.12, fcx, 4.34)
    dia = add_box(
        slide,
        fcx - 0.85,
        4.34,
        1.7,
        0.72,
        fill=C["bg"],
        line=C["accent"],
        line_w=1.75,
        shape="diamond",
    )
    set_shape_text(dia, "스키마 검사", S["caption"], F["head"], C["primary"], bold=True)
    add_text(
        slide,
        fcx + 0.95,
        4.5,
        G.RIGHT_EDGE - (tx + 2.3 + 0.95),
        0.6,
        "인용 문단 0개 · 분기 0개 ·\n결론 없는 단정 → 거부",
        S["caption"],
        F["body"],
        C["muted"],
    )
    _arrow(slide, fcx, 5.06, fcx, 5.55)
    add_text(
        slide, fcx + 0.12, 5.16, 0.9, 0.24, "통과", S["caption"], F["head"], C["primary"], bold=True
    )
    _flow_box(
        slide,
        fcx - 1.55,
        5.55,
        3.1,
        0.66,
        "화면 표시",
        "근거 문단과 나란히 — 인용은 좌측에서 하이라이트",
        line=C["accent"],
        line_w=1.75,
    )
    # 거부 루프백: 다이아 좌측 → 좌로 → 위로 → 생성 박스 좌변
    loop_x = fcx - 2.05
    _arrow(slide, fcx - 0.85, 4.7, loop_x, 4.7, head=False)
    _arrow(slide, loop_x, 4.7, loop_x, 3.87, head=False)
    _arrow(slide, loop_x, 3.87, fcx - 1.2, 3.87)
    add_text(
        slide,
        loop_x - 0.2,
        4.78,
        1.7,
        0.24,
        "거부 — 자동 재시도",
        S["caption"],
        F["head"],
        C["accent"],
        bold=True,
    )
    bar_and_source(
        slide, "형식을 지키지 못한 답변은 화면에 도달하지 못한다 — 스키마 선언 + 검증기 자동 재시도"
    )
    return slide


IMG_ANSWER = str(Path(__file__).resolve().parents[5] / "golden" / "ref" / "s12_answer.png")  # 508×721 실물 답변 캡처
IMG_CROP = str(Path(__file__).resolve().parents[5] / "golden" / "ref" / "s12_answer_crop.png")  # 508×382 발췌


def variant_b(prs):
    """B — 좌 서사 스택 / 중 실물 답변 발췌(대형) / 우 스키마 카드 + 검증 루프."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    header(slide, "구조화 출력 — 답변의 형식을 코드가 강제한다")
    from pptx.enum.lang import MSO_LANGUAGE_ID

    # ── 좌: 서사 3단 세로 스택 ──
    nx, nw = G.MARGIN_L, 2.55
    for i, (head, body) in enumerate(NARRATIVES):
        ny = G.CONTENT_TOP + i * 1.55
        add_text(
            slide,
            nx,
            ny,
            nw,
            0.28,
            f"0{i + 1}  {head}",
            S["head"],
            F["head"],
            C["primary"],
            bold=True,
        )
        add_text(
            slide,
            nx,
            ny + 0.34,
            nw,
            1.15,
            body,
            S["caption"],
            F["body"],
            C["muted"],
            line_spacing=1.25,
        )
    add_box(slide, 3.35, G.CONTENT_TOP, 0.012, G.CONTENT_BOTTOM - G.CONTENT_TOP, fill=C["bg_alt"])
    # ── 중: 실물 답변 발췌 (주인공, w 4.8) ──
    mx, mw = 3.5, 4.8
    _subhead(slide, mx, G.CONTENT_TOP, mw, "실물 답변 화면 — 발췌")
    img_w = 4.8
    img_h = img_w * 382 / 508  # 3.61
    pic = slide.shapes.add_picture(IMG_CROP, Inches(mx), Inches(2.25), width=Inches(img_w))
    pic.line.color.rgb = C["muted"]
    pic.line.width = Inches(0.01)
    add_box(slide, mx, 6.02, 0.4, 0.03, fill=C["accent"])
    add_text(
        slide,
        mx + 0.55,
        5.94,
        mw - 0.55,
        0.4,
        "전체 답변 중 [확인 질문]·[조건부 결론 Case 1·2] 발췌 — 본인 vs 대리인 질의",
        S["caption"],
        F["body"],
        C["muted"],
        line_spacing=1.2,
    )
    # ── 우: 스키마 카드 + 검증 루프 ──
    add_box(slide, 8.45, G.CONTENT_TOP, 0.012, G.CONTENT_BOTTOM - G.CONTENT_TOP, fill=C["bg_alt"])
    tx = 8.6
    tw = G.RIGHT_EDGE - tx  # 4.13
    _subhead(slide, tx, G.CONTENT_TOP, tw, "출력 스키마 — ClarifyOutput")
    card = add_box(slide, tx, 2.22, tw, 2.0, fill=C["primary"], shape="round")
    card.adjustments[0] = 0.05
    tf = card.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.1)
    glosses = [
        ("selected_branches", "고른 결론 분기"),
        ("answer", "답변 본문 (마크다운)"),
        ("cited_paragraphs", "인용 문단 — 비면 거부"),
        ("cited_cases", "인용 질의회신·감리 ID"),
        ("cited_ie", "인용 적용사례 번호"),
        ("follow_up_questions", "확인 질문 3개"),
        ("is_conclusion", "결론 포함 — 항상 True"),
    ]
    for i, (name, gloss) in enumerate(glosses):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p._p.get_or_add_pPr().set("eaLnBrk", "0")
        p.line_spacing = 1.35
        r1 = p.add_run()
        r1.text = name
        r1.font.name = "Consolas"
        r1.font.size = Pt(S["caption"])
        r1.font.bold = True
        r1.font.color.rgb = C["bg"]
        r2 = p.add_run()
        r2.text = f"  {gloss}"
        r2.font.name = F["body"]
        r2.font.size = Pt(S["caption"])
        r2.font.color.rgb = C["bg_alt"]
        r2.font.language_id = MSO_LANGUAGE_ID.KOREAN
    _subhead(slide, tx, 4.5, tw, "형식 검증 — 거부와 자동 재시도")
    _flow_box(slide, tx, 5.0, 1.15, 0.58, "답변 생성")
    _arrow(slide, tx + 1.15, 5.29, tx + 1.35, 5.29)
    dia = add_box(
        slide,
        tx + 1.35,
        4.96,
        1.15,
        0.66,
        fill=C["bg"],
        line=C["accent"],
        line_w=1.75,
        shape="diamond",
    )
    set_shape_text(dia, "스키마\n검사", S["caption"], F["head"], C["primary"], bold=True)
    _arrow(slide, tx + 2.5, 5.29, tx + 2.7, 5.29)
    _flow_box(slide, tx + 2.7, 5.0, tw - 2.7, 0.58, "화면 표시", line=C["accent"], line_w=1.75)
    dcx = tx + 1.925
    _arrow(slide, dcx, 5.62, dcx, 5.88, head=False)
    _arrow(slide, dcx, 5.88, tx + 0.55, 5.88, head=False)
    _arrow(slide, tx + 0.55, 5.88, tx + 0.55, 5.6)
    add_text(
        slide,
        tx,
        5.98,
        tw,
        0.3,
        "거부 — 자동 재시도: 인용 0개 · 분기 0개 · 결론 없는 단정",
        S["caption"],
        F["body"],
        C["muted"],
        line_spacing=1.2,
    )
    bar_and_source(
        slide, "형식을 지키지 못한 답변은 화면에 도달하지 못한다 — 스키마 선언 + 검증기 자동 재시도"
    )
    return slide


def audit(prs):
    from pptx.oxml.ns import qn

    fails = []
    for si, sl in enumerate(prs.slides):
        accents, diamonds = 0, 0
        for sh in sl.shapes:
            t, le = sh.top / 914400, sh.left / 914400
            b, r = t + sh.height / 914400, le + sh.width / 914400
            full_bleed = sh.width / 914400 >= SLIDE_W - 0.05
            if t < 6.4 and b > 6.37 and not full_bleed:
                fails.append((si, "bottom", round(b, 2), sh.shape_id))
            if r > 12.75 and not full_bleed:
                fails.append((si, "right", round(r, 2), sh.shape_id))
            for el in sh._element.iter(qn("a:prstGeom")):
                if el.get("prst") == "diamond":
                    diamonds += 1
            for el in sh._element.iter(qn("a:srgbClr")):
                if el.get("val") == "D66E3A":
                    if el.getparent().tag.endswith("}solidFill"):
                        accents += 1
        if accents > 4:
            fails.append((si, "accent>4", accents))
        if diamonds < 1:
            fails.append((si, "diamond<1", diamonds))
    assert not fails, f"AUDIT FAIL {fails}"
    print("audit pass")


def main():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    variant_a(prs)
    variant_b(prs)
    audit(prs)
    out = Path(__file__).parent / "variants" / "s12_variants.pptx"
    prs.save(str(out))
    print(f"saved {out}")


if __name__ == "__main__":
    main()
