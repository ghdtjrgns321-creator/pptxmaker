# 시안 H — 프로세스 밴드 정식화 + 하단 y예산 + 텍스트 리듬 규칙
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from .common import content_header, source_line
from .kit import SLIDE_H, SLIDE_W, add_box, add_text, load_kit, mix

K = load_kit()
C, S, F = K["rgb"], K["sizes"], K["fonts"]
MARGIN = K["layout"]["margin"]

# ── y좌표 예산 (겹침 계약: 각 블록 하단 + 0.08 ≤ 다음 블록 상단) ──
Y = {
    "zone1_title": 1.9,  # 룰 2.22
    "lane1_label": 2.42,
    "lane1_band": 2.68,  # h 0.78 → 3.46
    "lane2_label": 3.62,
    "lane2_band": 3.88,  # h 0.78 → 4.66
    "zone2_title": 4.88,  # 룰 5.2
    "panels": 5.34,  # h 1.0 → 6.34
    "band": 6.44,  # h 0.34 → 6.78
    "source_rule": SLIDE_H - 0.62,  # 6.88
}
BAND_H = 0.78


def _zone_title(slide, x, y, w, text):
    add_text(slide, x, y, w, 0.28, text, S["body"], F["head"], C["accent"], bold=True)
    add_box(slide, x, y + 0.32, w, 0.012, fill=C["muted"])


def _shape_step(slide, x, y, w, h, kind, fill, head, desc, head_ink, desc_ink):
    """프로세스 밴드 단계 — pentagon(첫)/chevron(후속). 노치 보정 인셋 + 2행(구/단어 리듬)."""
    shp = add_box(slide, x, y, w, h, fill=fill, shape=kind)
    try:
        shp.adjustments[0] = 0.55  # 꼭짓점 깊이(높이 비례) — 완만하게
    except (IndexError, ValueError):
        pass
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    notch = h * 0.55 / 2
    tf.margin_top = tf.margin_bottom = Inches(0)
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(h * 0.55 / 4)  # 화살촉 보정(1/2은 과보정 — 라벨 2줄 유발)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = head
    r1.font.name, r1.font.bold = F["head"], True
    r1.font.size = Pt(S["body"])
    r1.font.color.rgb = head_ink
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = desc
    r2.font.name = F["body"]
    r2.font.size = Pt(S["foot"])
    r2.font.color.rgb = desc_ink
    return shp


def variant_h(prs):
    """H안(v9) — 정식 프로세스 밴드(펜타곤+맞물린 셰브런), 하단 y예산, 텍스트 리듬 규칙."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    content_header(
        slide,
        1,
        "문제 정의",
        "일반 LLM을 회계에 못 쓰는 이유는 틀려서가 아니라, 틀리는 방향 때문이다",
    )

    GAP = 0.25
    ax, aw = MARGIN, 2.95
    rx = ax + aw + 2 * GAP + 0.014
    rw = SLIDE_W - MARGIN - rx

    # ── 좌측 앵커 (문장 구역 — 완전한 문장) ──
    add_text(
        slide, ax, 1.9, aw, 0.26, "핵심 리스크", S["caption"], F["head"], C["muted"], bold=True
    )
    add_text(slide, ax, 2.2, aw, 0.7, "2종 오류", S["display"], F["head"], C["accent"], bold=True)
    add_text(
        slide, ax, 2.9, aw, 0.48, "= 허위 확정", S["title"], F["head"], C["primary"], bold=True
    )
    add_text(
        slide,
        ax,
        3.55,
        aw,
        0.8,
        "근거 없는 확신은 부실감사와 재무제표 왜곡이라는 실제 손해로 이어진다.",
        S["sub"],
        F["body"],
        C["text"],
        line_spacing=1.25,
    )
    nums = [
        (
            "01",
            "검증 불가",
            "그럴듯한 답을 즉시 내놓지만, 어느 문단에 기반했는지 확인할 방법이 없어 감사 조서에 인용할 수 없다.",
        ),
        (
            "02",
            "판단 창작",
            "기준서에 실재하지 않는 기준을 확신에 찬 어조로 만들어내며, 오류는 언제나 확정 방향으로 표출된다.",
        ),
    ]
    for i, (no, head, desc) in enumerate(nums):
        ny = 4.6 + i * 0.95
        add_text(slide, ax, ny, 0.5, 0.3, no, S["head"], F["head"], C["accent"], bold=True)
        add_text(
            slide, ax + 0.5, ny, aw - 0.5, 0.28, head, S["body"], F["head"], C["primary"], bold=True
        )
        add_text(
            slide,
            ax + 0.5,
            ny + 0.3,
            aw - 0.5,
            0.62,
            desc,
            S["caption"],
            F["body"],
            C["muted"],
            line_spacing=1.2,
        )
    add_box(slide, ax + aw + GAP, 1.9, 0.014, 4.55, fill=C["bg_alt"])

    # ── 구획 1: 프로세스 밴드 2레인 (도형 구역 — 단어·구만) ──
    _zone_title(slide, rx, Y["zone1_title"], rw, "같은 질문, 두 개의 결말")
    qw, qh = 1.5, 1.35
    band_top, band_bottom = Y["lane1_band"], Y["lane2_band"] + BAND_H
    qy = (band_top + band_bottom) / 2 - qh / 2
    q = add_box(
        slide, rx, qy, qw, qh, fill=C["bg_alt"], line=C["muted"], line_w=0.75, shape="round"
    )
    tf = q.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = '"이 계약, 수익을\n지금 인식해도\n됩니까?"'
    r.font.name, r.font.bold = F["head"], True
    r.font.size = Pt(S["caption"])
    r.font.color.rgb = C["primary"]

    fx = rx + qw + GAP + 0.1
    fzone = rw - qw - GAP - 0.1
    overlap = BAND_H * 0.55 / 2  # 셰브런 노치 겹침 — 맞물리는 밴드
    step_w = (fzone + 3 * overlap) / 4  # 마지막은 결말 펜타곤
    lanes = [
        (
            Y["lane1_label"],
            Y["lane1_band"],
            "일반 LLM",
            C["accent"],
            True,
            [
                ("즉시 답변", "근거 조문 없음"),
                ("확신에 찬 어조", "판단 창작"),
                ("검증 불가", "근거 없음"),
            ],
            "치명",
        ),
        (
            Y["lane2_label"],
            Y["lane2_band"],
            "본 시스템",
            C["primary"],
            False,
            [
                ("근거 먼저 검색", "문단·사례 인용"),
                ("확실하면 확정", "결론+인용"),
                ("애매하면 유보", "정보 되물음"),
            ],
            "안전",
        ),
    ]
    for lbl_y, ly, name, ink, fatal, steps, verdict in lanes:
        add_text(slide, fx, lbl_y, 3.0, 0.24, name, S["caption"], F["head"], ink, bold=True)
        for i, (head, desc) in enumerate(steps):
            x = fx + i * (step_w - overlap)
            kind = "pentagon" if i == 0 else "chevron"
            tint = mix(ink, C["bg"], 0.9 if i < 2 else 0.82)
            _shape_step(
                slide,
                x,
                ly,
                step_w,
                BAND_H,
                kind,
                tint,
                head,
                desc,
                C["primary"],
                mix(C["primary"], C["bg"], 0.35),
            )
        vx = fx + 3 * (step_w - overlap)
        pent = _shape_step(
            slide, vx, ly, step_w, BAND_H, "chevron", ink, verdict, "", C["bg"], C["bg"]
        )

    # ── 구획 2: 대비 패널 (문장 구역 — 완전한 문장 2줄) ──
    _zone_title(slide, rx, Y["zone2_title"], rw, "오류의 두 방향 — 왜 2종만 치명적인가")
    half = (rw - GAP) / 2
    ph = 1.0
    panels = [
        (
            rx,
            C["bg_alt"],
            C["muted"],
            "1종 · 놓침 — 안전한 실패",
            C["primary"],
            '근거가 있는데 못 찾아 "모른다"고 답하는 실패는 추가 검토로 해소된다. 실측된 유보 편향(되물음 9건)도 프롬프트 재균형으로 관리하고 있다.',
            C["text"],
        ),
        (
            rx + half + GAP,
            C["primary"],
            C["accent"],
            "2종 · 허위 확정 — 치명적 실패",
            C["bg"],
            "근거가 없거나 틀렸는데 확신에 찬 답을 내놓는 순간 부실감사로 직결된다. 그래서 설계 전체를 이 방향의 차단에 편향시켰다.",
            C["bg_alt"],
        ),
    ]
    for px, fill, spine, head, head_ink, desc, desc_ink in panels:
        add_box(slide, px, Y["panels"], half, ph, fill=fill, shape="round")
        add_box(slide, px + 0.14, Y["panels"] + 0.14, 0.05, ph - 0.28, fill=spine)
        add_text(
            slide,
            px + 0.34,
            Y["panels"] + 0.12,
            half - 0.55,
            0.28,
            head,
            S["body"],
            F["head"],
            head_ink,
            bold=True,
        )
        add_text(
            slide,
            px + 0.34,
            Y["panels"] + 0.44,
            half - 0.55,
            0.52,
            desc,
            S["caption"],
            F["body"],
            desc_ink,
            line_spacing=1.2,
        )

    band = add_box(slide, MARGIN, Y["band"], SLIDE_W - 2 * MARGIN, 0.34, fill=C["primary"])
    tfb = band.text_frame
    tfb.vertical_anchor = MSO_ANCHOR.MIDDLE
    pb = tfb.paragraphs[0]
    pb.alignment = PP_ALIGN.CENTER
    rb = pb.add_run()
    rb.text = "확실할 때만 확정하고, 애매하면 근거를 보여주며 유보한다"
    rb.font.name, rb.font.bold = F["head"], True
    rb.font.size = Pt(S["body"])
    rb.font.color.rgb = C["bg"]
    source_line(
        slide, "출처: 1_OVERVIEW.md · PROJECT_OVERVIEW.md · 5_INTERFACE.md (00_factsheet.md §A·§B)"
    )
