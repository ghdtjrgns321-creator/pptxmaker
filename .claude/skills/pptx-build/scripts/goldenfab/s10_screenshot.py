"""S10 시안 — 스크린샷 전용 장 (신규 어휘: 풀블리드 실물 화면. UI 패스).

다크 배경 + 3D 뷰어 캡처 대형 + 헤더·캡션만. 실물 화면이 주인공인 장.
실행: 이 타입은 goldenfab 레지스트리 경유로만 렌더된다 — 골든 19장 확인은
      `uv run python golden/build_golden.py`(2026-07-15 단일화로 시안 개별 실행 경로 폐지).
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from . import grid as G
from .kit import SLIDE_H, SLIDE_W, add_box, add_text, fit_picture, load_kit

K = load_kit()
C, S, F = K["rgb"], K["sizes"], K["fonts"]

IMG_3D = r"C:\Users\ghdtj\workspace\portfolio\k-ifrs-1115\images\knowledge_graph_3d.png"
EMU_IN = 914400

# ── 캡처 박스 (GRID 파생) ──
# **박스가 레이아웃의 단일 출처다.** 실제 그림은 이 안에 맞춰 축소되고(fit_picture),
# 텍스트 칼럼 x는 그림의 실제 폭이 아니라 **박스 폭**에서 파생한다 — 그래야 어떤 비율의
# 이미지가 와도 칼럼이 안 밀린다(2026-07-15 버그: 비율 가정이 틀려 1.74" 침범).
IMG_X, IMG_Y = G.MARGIN_L, 1.8
IMG_MAX_W, IMG_MAX_H = 7.6, 4.75
CAP_GAP = 0.18  # 그림 bottom → 캡션 top
COL_GAP = 0.4
OX = IMG_X + IMG_MAX_W + COL_GAP  # 8.6 — 관찰 칼럼 시작

# ── 관찰 칼럼 세로 배치 (항목 수에서 파생 — 리터럴 피치 금지) ──
# 구: `py = 2.4 + i * 1.05` — 피치가 리터럴이라 **항목 수와 무관**했다. 4개일 때 마지막 본문
# bottom이 6.47로 하한(6.35)을 0.12 넘었다(오딧 P2③ 검출). 항목이 5개면 더 넘친다.
# 하한에서 역산하면 3개든 5개든 안 넘친다.
PT_TOP = 2.40  # 첫 항목 top (룰 2.15+0.012에서 공기 0.238)
PT_HEAD_H = 0.28  # 번호·헤드 높이
PT_BODY_DY = 0.32  # 항목 top → 본문 top
PT_BODY_H = 0.60  # 본문 높이(2줄)
PT_BLOCK_H = PT_BODY_DY + PT_BODY_H  # 0.92 — 한 항목이 차지하는 실제 높이


def point_pitch(n):
    """항목 n개를 PT_TOP ~ CONTENT_BOTTOM 안에 균등 배치하는 피치.

    마지막 항목 bottom = PT_TOP + (n-1)*pitch + PT_BLOCK_H ≤ CONTENT_BOTTOM 을 만족하는
    최대 피치. n=1이면 나눗셈이 없으므로 블록 높이만 확인.
    """
    if n <= 1:
        return PT_BLOCK_H
    room = G.CONTENT_BOTTOM - PT_TOP - PT_BLOCK_H
    return min(1.05, room / (n - 1))  # 1.05 = 여유 있을 때의 기존 리듬(상한)


SHOT_DEFAULTS = {
    "kicker": "3. 기술 설명 — TECH 02 · Retrieve",
    "headline": "노드 929개가 하나의 구조로 — 지식그래프 3D 뷰어",
    "img": IMG_3D,
    "caption": "실제 운영 화면 — 지식그래프 3D 뷰어 (초기 V1 스냅샷, 수치 표기는 그래프 v14 기준)",
    # 골든 기본 4포인트. **DEFAULT 밖(함수 몸통)에 두지 말 것** — 2026-07-15 제3자 반증에서
    # `"points": None` + 함수 내 `c["points"] or [리터럴]` 조합이 두 게이트를 동시에 실명시켰다:
    # ① assert_content는 키가 있으니 통과 ② golden_strings가 DEFAULT만 걷어 리터럴을 못 봐
    # 오염 검사도 통과. 계약 content를 다 채워도 골든 원문이 나갔다(원사고와 동일 클래스).
    "points": [
        ("큰 파란 노드 = 개념", "기준서 소제목 80개 — 클수록·밝을수록 상위 계층이다."),
        (
            "작은 점 = 문단·사례",
            "회색은 문단 250, 색 점은 QNA·감리·IE 사례 188 — 전부 개념 아래 매달린다.",
        ),
        (
            "선 = 관계 간선",
            "회색 계층(목차) · 노랑 상호참조(E3) · 빨강 선행판단(E2) — 검색이 걷는 길이 그대로 보인다.",
        ),
        ("고립 노드 0", "모든 노드가 최소 하나의 경로로 본문에 닿는다 — 감사로 확인."),
    ],
}


def variant_a(prs, c=None):
    c = {**SHOT_DEFAULTS, **(c or {})}
    """A — 흰 배경(본문 무드 통일) 풀블리드: 캡처가 주인공, 텍스트는 헤더·캡션만."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    add_text(
        slide,
        G.MARGIN_L,
        0.42,
        8.0,
        0.28,
        c["kicker"],
        S["caption"],
        F["head"],
        C["muted"],
        bold=True,
    )
    add_text(
        slide,
        G.MARGIN_L,
        0.72,
        G.RIGHT_EDGE - G.MARGIN_L,
        0.55,
        c["headline"],
        S["section"],
        F["head"],
        C["primary"],
        bold=True,
    )
    add_box(
        slide,
        G.MARGIN_L,
        G.RULE_Y,
        G.RIGHT_EDGE - G.MARGIN_L,
        0.014,
        fill=C["muted"],
    )
    # 캡처는 **박스에 맞춘다** — 한 축만 주면 나머지는 실제 파일 비율이 정하고, 가정이 틀리면
    # 옆 칼럼을 침범한다(실측: 비율 가정 1.60 vs 실제 2.05 → 폭 7.6" 전제인데 9.74"로 나와 1.74" 침범).
    ix, iy, iw, ih = fit_picture(
        slide, c["img"], IMG_X, IMG_Y, IMG_MAX_W, IMG_MAX_H, line=C["muted"]
    )
    cap_y = iy + ih + CAP_GAP  # 캡션은 **실제** 그림 아래에 (그림 높이는 비율마다 달라진다)
    add_box(slide, G.MARGIN_L, cap_y + 0.09, 0.5, 0.035, fill=C["accent"])
    add_text(
        slide,
        G.MARGIN_L + 0.65,
        cap_y,
        IMG_MAX_W - 0.65,
        0.3,
        c["caption"],
        S["caption"],
        F["body"],
        C["muted"],
    )
    # ── 우측 관찰 칼럼 — 캡처 범례 실제 항목과 1:1 대응 ──
    # 칼럼 x는 **이미지 박스** 기준(실제 그림 폭이 아니라) — 그래야 어떤 비율이 와도 안 밀린다.
    ox = OX
    ow = G.RIGHT_EDGE - ox
    add_text(
        slide,
        ox,
        1.8,
        ow,
        0.3,
        "이 화면에서 보이는 것",
        S["head"],
        F["head"],
        C["primary"],
        bold=True,
    )
    add_box(slide, ox, 2.15, ow, 0.012, fill=C["muted"])
    points = c["points"]  # 폴백 금지 — 기본값은 SHOT_DEFAULTS에만 (위 주석 참조)
    pitch = point_pitch(len(points))  # 항목 수에서 파생 — 리터럴 피치는 넘친다(위 주석)
    for i, (head, body) in enumerate(points):
        py = PT_TOP + i * pitch
        add_text(
            slide,
            ox,
            py,
            0.5,
            PT_HEAD_H,
            f"{i + 1:02d}",
            S["head"],
            F["head"],
            C["accent"],
            bold=True,
        )
        add_text(
            slide,
            ox + 0.45,
            py + 0.02,
            ow - 0.45,
            0.26,
            head,
            S["body"],
            F["head"],
            C["primary"],
            bold=True,
        )
        add_text(
            slide,
            ox + 0.45,
            py + PT_BODY_DY,
            ow - 0.45,
            PT_BODY_H,
            body,
            S["caption"],
            F["body"],
            C["muted"],
            line_spacing=1.2,
        )
    return slide


def audit(prs):
    fails = []
    for si, sl in enumerate(prs.slides):
        pics = sum(1 for sh in sl.shapes if sh.shape_type == 13)
        if pics != 1:
            fails.append((si, "picture≠1", pics))
    assert not fails, f"AUDIT FAIL {fails}"
    print("audit pass")


def main():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    variant_a(prs)
    audit(prs)
    out = Path(__file__).parent / "variants" / "s10_screenshot.pptx"
    prs.save(str(out))
    print(f"saved {out}")


if __name__ == "__main__":
    main()
