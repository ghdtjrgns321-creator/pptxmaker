"""S21 시안 — 클로징 (물성 선언: 클로징 = 표지와의 수미상관 — 다크·같은 축·같은 액센트).

신규 콘텐츠 0 — 기확정 문장 재사용: 가치 제안(표지), 결론 바 3개(S14·S15·S18).
좌측 축은 표지 GRID(COVER_L 0.65) 그대로, 진행 도트는 6개 전부 점등(완주 표시).
실행: uv run python golden/s21_closing.py → golden/variants/s21_closing.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches

from .kit import SLIDE_H, SLIDE_W, add_box, add_text, load_kit, mix

K = load_kit()
C, S, F = K["rgb"], K["sizes"], K["fonts"]

L = 0.65  # 표지 GRID 좌측 축
VALUE_PROP = "확실할 때만 확정하고,\n애매하면 근거를 보여주며 유보한다"
SUMMARY = [  # 기확정 결론 바 문장의 재사용 (S14 · S15 · S18)
    "2종 오류를 구조로 차단 — 확률 신호 전량 폐기, 온톨로지 그래프로 동작한다",
    "정답 문서를 차단한 채 78/92 — 실패 14건까지 전수 해부된 수치다",
    "못 하는 것까지 계약 — 경계 안에서만, 결정적으로",
]


CLOSING_DEFAULTS = {
    "kicker": "CLOSING — K-IFRS 1115 온톨로지 지식그래프 RAG",
    "value_prop": VALUE_PROP,
}


def variant_a(prs, c=None):
    c = {**CLOSING_DEFAULTS, **(c or {})}
    """vC — 여백형 클로징: 킥커 / 가치 제안(광학 중앙) / 완주 도트 / 백색 밴드. 그 외 0."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["primary"])
    add_text(
        slide,
        L,
        0.55,
        9.0,
        0.3,
        c["kicker"],
        S["caption"],
        F["head"],
        C["muted"],
        bold=True,
    )
    add_box(slide, L, 0.92, SLIDE_W - 2 * L, 0.014, fill=mix(C["primary"], C["muted"], 0.5))
    # 타이틀 블록 — 광학 중앙 (블록 2.62~4.45, 중심 약 3.54)
    add_box(slide, L, 2.62, 0.3, 0.3, fill=C["accent"])
    heads = []
    for i, seg in enumerate(c["value_prop"].split(chr(10))):
        heads.append(
            add_text(
                slide,
                L,
                3.05 + i * 0.72,
                12.0,
                0.7,
                seg,
                S["display"],
                F["head"],
                C["bg"],
                bold=True,
            )
        )
    # 완주 도트 — 6부 점등, 타이틀 블록 아래 호흡
    for i in range(6):
        add_box(slide, L + i * 0.35, 5.05, 0.13, 0.13, fill=C["accent"], shape="oval")
    # 하단 백색 밴드 — 표지 다크 밴드의 반전
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt

    band_top = 6.62
    add_box(slide, 0, band_top, SLIDE_W, SLIDE_H - band_top, fill=C["bg"])
    add_text(
        slide,
        L,
        band_top + (SLIDE_H - band_top) / 2 - 0.15,
        5.0,
        0.3,
        K["brand"]["name"],
        S["head"],
        F["head"],
        C["primary"],
        bold=True,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    ph = add_box(
        slide, SLIDE_W - L - 2.2, band_top + 0.2, 2.2, 0.48, fill=None, line=C["muted"], line_w=1.0
    )
    tf = ph.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    pp2 = tf.paragraphs[0]
    pp2.alignment = PP_ALIGN.CENTER
    r = pp2.add_run()
    r.text = "회사 로고"
    r.font.name = F["body"]
    r.font.size = Pt(S["caption"])
    r.font.color.rgb = C["muted"]
    shapes = {"heads": heads}
    return slide, shapes


def audit(slide, shapes):
    EMU = 914400
    fails = []
    head_lefts = [sh.left / EMU for sh in shapes["heads"]]
    if any(abs(v - L) > 0.001 for v in head_lefts):
        fails.append(("헤드 축 편차", head_lefts))
    # 광학 중앙: 타이틀 블록(사각 top 2.62 ~ 헤드 bottom) 중심 3.4~3.8
    top = 2.62
    bottom = max((sh.top + sh.height) / EMU for sh in shapes["heads"])
    center = (top + bottom) / 2
    if not (3.4 <= center <= 3.8):
        fails.append(("블록 중심", round(center, 2)))
    # 여백 계약: 텍스트 프레임 수 <= 6 (킥커·헤드2·워드마크·로고 + 여유 1)
    n_text = sum(1 for sh in slide.shapes if sh.has_text_frame and sh.text_frame.text.strip())
    if n_text > 6:
        fails.append(("텍스트 과다", n_text))
    assert not fails, f"AUDIT FAIL {fails}"
    print(f"audit pass — 텍스트 프레임 {n_text}/6, 블록 중심 {center:.2f}")


def main():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    slide, shapes = variant_a(prs)
    audit(slide, shapes)
    out = Path(__file__).parent / "variants" / "s21_closing.pptx"
    prs.save(str(out))
    print(f"saved {out}")


if __name__ == "__main__":
    main()
