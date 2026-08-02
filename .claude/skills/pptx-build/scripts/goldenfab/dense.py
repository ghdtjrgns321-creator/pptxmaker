"""goldenfab.dense — s06 고밀도 디자인 언어 공유 헬퍼 (2026-07-20).

s06_mid에서 확정한 규격을 골든 전 장이 공유한다(handoff 문법 승격):
- compact_header : 키커 + 14pt 헤드라인 + 굵은 룰 (골든 28pt 헤더·리드 문단 폐지)
- hero_card      : 배지·제목·태그·중앙 배너·히어로 아이콘·2단 불릿·아이콘 칩 (s06 카드)
                   — **카드는 이 컴포넌트 하나만.** 장별 경량 변형 금지(design-rules §8).
- chip_row / two_col_bullets / icon / hrule / source_line

원칙: **하단 전폭 한줄평/재서술 바 금지(형태 자체)** — 색 불문. 결론은 콘텐츠에 흡수.
Ember accent 소량, 전 텍스트 8~15pt(헤드라인 14 예외). 자세한 차단 근거는 design-rules §8.
"""

from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from . import grid as G
from .kit import ROOT, add_box, add_text, load_kit, mix, set_shape_text

K = load_kit()
C, S, F = K["rgb"], K["sizes"], K["fonts"]
ICONS = ROOT / "assets/icons"

FULL_W = G.RIGHT_EDGE - G.MARGIN_L
RULE_Y = 0.66
SOURCE_Y = 7.14

line_soft = mix(C["primary"], C["bg"], 0.80)
line_mid = mix(C["primary"], C["bg"], 0.55)
ink_body = mix(C["muted"], C["primary"], 0.45)


def icon(slide, name, variant, x, y, size):
    p = ICONS / f"{name}_{variant}.png"
    return slide.shapes.add_picture(str(p), Inches(x), Inches(y), Inches(size), Inches(size))


def hrule(slide, x, y, w, color=None, weight=0.75):
    return add_box(slide, x, y, w, 0.001, fill=None, line=color or line_soft, line_w=weight)


def arrow(slide, x1, y1, x2, y2, color, w=1.5):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = color
    conn.line.width = Pt(w)
    ln = conn.line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    return conn


BAND_RULE_DY = 0.36


def band_head(slide, y, text, x=None, w=None):
    """밴드 소제목 + 룰 — 한 장 안의 구획을 여백이 아니라 장치로 나눈다.

    페이지 요소다(도해가 아니다). 2층 부품은 이걸 소유하지 않는다 — 부품은 box 안 그래픽
    하나만 그리고, 구획을 나눌지는 장이 정한다.

    `x`·`w`를 주면 **그 자리 폭만큼**만 그린다. 골든이 좌우 두 자리를 쓸 때 소제목도 자리마다
    따로다(S12 좌 x0.60 w5.10 · S11 우 x6.60). 기본값(전폭)은 골든 sparse 장 호환.
    """
    x = G.MARGIN_L if x is None else x
    w = (G.RIGHT_EDGE - G.MARGIN_L) if w is None else w
    add_text(slide, x, y, w, 0.32, text, S["head"], F["head"], C["primary"], bold=True)
    add_box(slide, x, y + BAND_RULE_DY, w, 0.012, fill=C["muted"])


def figure_caption(slide, x, y, w, text):
    """도해 바로 밑 한 줄 — 그림이 말하지 못하는 분해·단서를 붙인다(골든 s08 실측 y3.36).

    조립 장이 얕아 보이는 가장 큰 이유가 **도해와 카드 사이가 통째로 비는 것**이었다
    (2026-07-29 실측: 골든 최대 백지 0.3" vs 조립 1.1~1.3"). 골든은 그 자리를 이 줄로 잇는다.
    """
    add_text(slide, x, y, w, 0.18, text, S["foot"], F["body"], C["muted"])


DETAIL_HEAD_H = 0.28
DETAIL_BODY_DY = 0.33
DETAIL_BODY_H = 1.05
DETAIL_GAP = 0.35


def detail_columns(slide, x, y, w, items, *, gap=DETAIL_GAP, body_h=DETAIL_BODY_H):
    """설명 칼럼 N개 — 머리글 + (볼드 리드 + 평어 본문) + 칼럼 사이 세로 구분선.

    골든 s06이 얇은 레인 도해 아래를 이 층으로 채운다. **같은 도해로도 장이 차는 이유**가
    여기다 — 새 사실을 만드는 게 아니라 같은 내용의 설명 층을 얹는 것이다.
    골든과 조립이 이 함수 하나를 공유한다(관용구가 두 벌이 되면 곧 갈라진다).

    items = [(머리글, 볼드 리드, 평어 본문), ...]
    """
    n = len(items)
    cw = (w - (n - 1) * gap) / n
    for i, (head, lead, body) in enumerate(items):
        dx = x + i * (cw + gap)
        add_text(
            slide, dx, y, cw, DETAIL_HEAD_H, head, S["caption"], F["head"], C["primary"], bold=True
        )
        add_text(
            slide,
            dx,
            y + DETAIL_BODY_DY,
            cw,
            body_h,
            [[(lead, {"bold": True, "color": C["primary"]}), (body, {})]],
            S["caption"],
            F["body"],
            C["muted"],
            line_spacing=1.25,
        )
        if i < n - 1:
            add_box(slide, dx + cw + gap / 2, y, 0.012, body_h, fill=C["bg_alt"])
    return cw


def compact_header(slide, kicker, headline):
    """s06 압축 헤더 — 키커(8pt) + 헤드라인(14pt) + 굵은 룰. 상단 의식 ~0.7\"."""
    add_text(
        slide, G.MARGIN_L, 0.14, 8.0, 0.16, kicker, S["foot"], F["head"], C["muted"], bold=True
    )
    add_text(
        slide,
        G.MARGIN_L,
        0.32,
        FULL_W,
        0.28,
        headline,
        S["sub"],
        F["head"],
        C["primary"],
        bold=True,
    )
    hrule(slide, G.MARGIN_L, RULE_Y, FULL_W, color=C["primary"], weight=1.4)


def source_line(slide, text):
    add_text(slide, G.MARGIN_L, SOURCE_Y, FULL_W, 0.2, text, S["foot"], F["body"], C["muted"])


# [차단 2026-07-20] 하단 전폭 '한줄평/재서술 바'는 형태 자체가 금지다(handoff 결함 #2).
# 검은색이든 라이트든 accent 좌변이든 — 슬라이드 맨 밑 전폭 한 줄 결론은 만들지 않는다.
# 결론 문장은 콘텐츠(불릿·배너)에 흡수한다. takeaway_strip 헬퍼는 그래서 폐기했다.


def chip_row(slide, x, y, w, ic, em, label, *, h=0.23):
    """아이콘 칩 — 아이콘 + 강조(accent) + 라벨(muted)."""
    chip = add_box(slide, x, y, w, h, fill=C["bg_alt"])
    tf = chip.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.30)
    p = tf.paragraphs[0]
    p._p.get_or_add_pPr().set("eaLnBrk", "0")
    p.alignment = PP_ALIGN.LEFT
    for txt, color, bold in ((em, C["accent"], True), (label, C["muted"], False)):
        if not txt:
            continue
        r = p.add_run()
        r.text = txt
        r.font.name = F["head"]
        r.font.size = Pt(S["caption"])
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.language_id = MSO_LANGUAGE_ID.KOREAN
    icon(slide, ic, "primary", x + 0.08, y + (h - 0.17) / 2, 0.17)


def two_col_bullets(slide, x, y, w, items, *, step=0.29, h=0.28):
    """2단 불릿 — items: [(리드, 설명)]. 볼드 리드 + 설명(서술형 금지)."""
    for j, (lead, detail) in enumerate(items):
        add_text(
            slide,
            x,
            y + j * step,
            w,
            h,
            [
                [
                    ("· ", {"bold": True, "color": C["accent"]}),
                    (lead, {"bold": True, "color": C["primary"]}),
                    (detail, {}),
                ]
            ],
            S["caption"],
            F["body"],
            ink_body,
            anchor=MSO_ANCHOR.MIDDLE,
            line_spacing=1.0,
        )


# [차단 2026-07-20] narrative_card(경량 카드) 폐기. 카드는 hero_card 하나만 쓴다 —
# 장별 경량 변형이 s06과 어긋나 "쓰레기" 카드를 낳았다(사용자 반려). 카드 = hero_card 단일.


def hero_card(slide, x, y, w, h, num, title, tag, ic, banner, items, chips, *, hero=0.48):
    """s06 표준 히어로 카드 — 골든 전 장 공용 단일 컴포넌트.

    배지→제목→태그→배너→히어로 아이콘(상단 고정) · 룰+칩(하단 고정) · 불릿(가운데 채움).
    위치를 콘텐츠·h에서 **파생**한다 — 불릿 수·카드 높이가 달라도 s06 룩 유지 + 죽은 여백 0.
    이 함수만 카드를 그린다 — 장별 즉흥 카드 금지(design-rules §8).

    **좁은 자리에도 쓴다 — `hero`·칩 수·불릿 수가 전부 호출자 인자다.** 필요 높이는 고정값이
    아니라 그 셋에서 나온다:

        h ≥ 1.44(배지·제목·태그·배너) + hero + 0.24×칩수 + 한줄높이×불릿수

    예) 불릿2·칩1·hero 0.30 → 2.50" · 불릿2·칩0·hero 0.26 → 2.22". 자리가 2.2"면 겹을 빼서
    경량 카드를 만드는 게 아니라 **hero를 줄여** 이 함수를 그대로 쓴다.

    (2026-08-02) 이 문단은 `h≥3.2 권장`을 대체한 것이다. 그 권장값은 07-29에 아카이브된
    껍데기 `card_row`의 가드(`_archive/code/goldenfab/figures/card_row.py:86`)에서 숫자만
    남은 잔해였고, **무엇을 줄이면 들어가는지가 없어서** "좁아서 못 쓴다"는 오판을 낳았다
    (local-ai-assist 2026-08-02: 카드 14장 전부 즉흥 신설). 가드는 아래에 복원했다."""
    add_box(slide, x, y, w, h, fill=C["bg"], line=line_soft, line_w=1.0)
    badge = add_box(slide, x + 0.14, y + 0.13, 0.28, 0.28, fill=C["primary"], shape="oval")
    set_shape_text(badge, num, S["foot"], F["head"], C["bg"], bold=True)
    add_text(
        slide,
        x + 0.52,
        y + 0.11,
        w - 0.66,
        0.32,
        title,
        S["head"],
        F["head"],
        C["primary"],
        bold=True,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        x + 0.14,
        y + 0.50,
        w - 0.28,
        0.20,
        tag,
        S["caption"],
        F["head"],
        C["accent"],
        bold=True,
    )
    ban = add_box(slide, x + 0.14, y + 0.76, w - 0.28, 0.30, fill=C["bg_alt"])
    set_shape_text(ban, banner, S["caption"], F["head"], C["primary"], bold=True)
    ban.text_frame.word_wrap = False
    icon(slide, ic, "accent", x + (w - hero) / 2, y + 1.10, hero)
    # 하단 고정: 룰 + 칩 (카드 바닥에 정렬)
    rule_y = y + h - len(chips) * 0.24 - 0.14
    hrule(slide, x + 0.14, rule_y, w - 0.28)
    for k, (cic, em, label) in enumerate(chips):
        chip_row(slide, x + 0.14, rule_y + 0.09 + k * 0.24, w - 0.28, cic, em, label, h=0.23)
    # 가운데 채움: 히어로 밑 ~ 룰 위 구간을 불릿 수로 균등 분배
    b_top = y + 1.10 + hero + 0.14
    b_bottom = rule_y - 0.06
    step = (b_bottom - b_top) / max(len(items), 1)
    # 자리 부족 가드 — 아카이브된 card_row(:86)에 있던 것을 알맹이로 옮긴 것(2026-08-02 복원).
    # 없으면 step이 음수가 되어 **음수 높이 텍스트박스**가 든 pptx가 경고 없이 저장되고,
    # PowerPoint가 파일 자체를 못 연다(재현: h1.81·칩1·불릿2 → step -0.06). 조용히 깨지는
    # 대신 무엇을 줄이면 되는지를 수치로 낸다. 임계는 리터럴이 아니라 brand-kit 본문 크기 파생.
    line_h = S["caption"] / 72 * 1.6
    if items and step < line_h:
        need = 1.44 + hero + 0.24 * len(chips) + line_h * len(items)
        raise ValueError(
            f'hero_card: h={h:.2f}"에서 불릿 간격이 {step:.2f}"다(한 줄 {line_h:.2f}" 필요). '
            f'불릿 {len(items)}·칩 {len(chips)}·hero {hero:.2f} 기준 필요 h ≥ {need:.2f}". '
            "자리를 넓히거나 hero·칩·불릿을 줄일 것 — 겹을 빼서 경량 카드를 만들지 말 것(§8)."
        )
    two_col_bullets(slide, x + 0.15, b_top, w - 0.26, items, step=step, h=min(0.32, step))
