"""실전 덱 기계 오딧 러너 — 변형(adapted)·신규(novel)·content 주입 골든 장에 전역 규칙을 건다.

사용:
    uv run python .claude/skills/pptx-build/scripts/audit_deck.py <deck-spec.json> <deck.pptx>

왜 이 파일이 있나(2026-07-16 배선 재설계): 골든이 "복제 템플릿"에서 "변형 가능한 출발점"이
되면서 스냅샷 회귀(compare_golden)는 goldenfab 보호 전용으로 물러났다 — 실전 덱의 변형·신규
장은 스냅샷 대상이 아니다. 그 빈자리를 3겹이 메운다:
    ① 이 러너 — 레이아웃 무관 전역 규칙(goldenfab.audit.generic_checks) + 밀도 밴드
    ② 밀도 밴드 — 골든 스냅샷에서 파생한 본문 장 최저 밀도(리터럴 아님, 골든 따라 이동)
    ③ 병렬 채점 — 판단형 결함(형식-내용 적합·색 의미)은 design-rules P4가 진다

오딧 대상: spec type이 golden.*(정형 제외)·adapted.*·novel인 장. 레거시 타입(bullets·chart 등)은
consistency-qa의 기존 게이트가 본다.

슬라이드별 예외는 spec에 선언한다(코드 수정 없이):
    "audit": {"dup_allow": N, "known": {"규칙명": 기준선}}
dup_allow = 재등장이 메시지인 노드 수(실측을 박는다, design-rules §E-2).
known = 사용자 승인으로 확정된 기존 빚의 래칫 기준선(goldenfab.audit.report 참조).
"""

import json
import sys
from pathlib import Path

from pptx import Presentation

HERE = Path(__file__).resolve().parent
sys.path.insert(0, str(HERE))

from goldenfab import audit as A  # noqa: E402
from goldenfab.kit import load_kit  # noqa: E402

FIXED = {"cover", "toc", "part", "closing"}  # 정형 장 — 고정 템플릿이라 스냅샷·계약이 이미 지킨다


def target_slides(spec):
    """오딧 대상 (슬라이드 번호, spec 항목) 목록."""
    out = []
    for i, sl in enumerate(spec["slides"], start=1):
        t = sl.get("type", "")
        if t.startswith("golden."):
            if t.split(".", 1)[1] not in FIXED:
                out.append((i, sl))
        elif t.startswith("adapted.") or t == "novel":
            out.append((i, sl))
    return out


def main():
    if len(sys.argv) != 3:
        print(__doc__)
        return 2
    spec = json.loads(Path(sys.argv[1]).read_text(encoding="utf-8"))
    prs = Presentation(sys.argv[2])
    slides = list(prs.slides)
    accent = str(load_kit()["rgb"]["accent"])
    band = A.density_band()

    all_fails = []
    targets = target_slides(spec)
    print(
        f"오딧 대상 {len(targets)}/{len(slides)}장 · "
        f"밀도 밴드(골든 본문 {band['n_body']}장 파생): "
        f"도형 ≥{band['shapes_min']} · 텍스트 프레임 ≥{band['frames_min']}"
    )
    for no, sl in targets:
        if no > len(slides):
            all_fails.append(f"S{no}: 슬라이드 없음")
            continue
        shapes = list(slides[no - 1].shapes)
        opts = sl.get("audit", {})
        print(f"\n── S{no} {sl['type']}: 도형 {len(shapes)}")
        res = A.generic_checks(shapes, accent, band=band, dup_allow=opts.get("dup_allow", 0))
        fails = A.report(res, known=opts.get("known"))
        all_fails += [f"S{no}: {f}" for f in fails]

    print("\n" + ("DECK AUDIT PASS" if not all_fails else f"DECK AUDIT FAIL {len(all_fails)}건"))
    for f in all_fails:
        print(f"  - {f}")
    return 1 if all_fails else 0


if __name__ == "__main__":
    sys.exit(main())
