#!/usr/bin/env python
"""덱 기계 채점 — 완성 pptx에 전역 규칙(generic_checks)과 밀도 밴드를 건다.

    uv run python .claude/skills/pptx-build/scripts/score_deck.py <deck.pptx> [--body 4,5,7,9]

## 왜 이 파일이 있나 (2026-07-29 신설)

옛 러너 `audit_deck.py`는 `deck-spec.json`을 인자로 받았다 — spec의 장 타입(golden./adapted./
novel)을 보고 검사 대상을 골랐기 때문이다. deck-spec 경로가 폐기되면서(2026-07-29 하네스 정리)
그 입구가 사라졌고, 남은 `goldenfab.audit`을 부를 진입점이 없어졌다. 이 파일이 그 자리다.

**spec을 안 본다.** pptx만 열어서 검사한다. 대상 장은 인자로 받거나(--body), 없으면 전 장을
검사하되 정형 장(도형이 적은 표지·목차·간지)은 밀도 검사에서 자동으로 빠진다.

## 이 러너가 판정자가 아닌 이유 (정직)

게이트는 사용자가 반려한 조립 7장을 전부 통과시킨 전례가 있다(면적 82~97%). 그래서 이 수치는
"통과=좋다"가 아니라 **"기계가 보는 것 기준으로 어디가 튀나"**를 보는 용도다. 형식이 내용에
맞는지는 기계가 못 본다 — 그건 렌더 PNG 눈검증(⑥)이 진다.
"""

import argparse
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))

from pptx import Presentation  # noqa: E402

from goldenfab import audit as A  # noqa: E402
from goldenfab.kit import load_kit  # noqa: E402

FIXED_MAX_SHAPES = 12  # 이보다 도형이 적으면 정형 장(표지·목차·간지)으로 보고 밀도 면제


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx", type=Path)
    ap.add_argument("--body", help="본문 장 번호(1-index, 콤마 구분). 없으면 전 장")
    ap.add_argument("--profile", default="dense")
    args = ap.parse_args()

    kit = load_kit()
    accent = kit["colors"]["accent"]
    band = A.density_band()
    prs = Presentation(str(args.pptx))

    body = None
    if args.body:
        body = {int(n) for n in args.body.split(",")}

    total = 0
    checked = 0
    for i, slide in enumerate(prs.slides, 1):
        shapes = list(slide.shapes)
        if body is not None:
            if i not in body:
                continue
        elif len(shapes) < FIXED_MAX_SHAPES:
            continue  # 정형 장 — 밀도 대상 아님
        checked += 1
        res = A.generic_checks(shapes, accent, band, profile=args.profile)
        bad = [r for r in res if not r[1]]
        total += len(bad)
        mark = "PASS" if not bad else f"FAIL {len(bad)}"
        print(f"\n== s{i:02d}  도형 {len(shapes):3d}  {mark}")
        for name, ok, msg, _n in res:
            if not ok:
                print(f"    x {name}: {msg}")

    print(
        f"\n본문 {checked}장 총 FAIL {total}건  (밀도 밴드 {band['chars_min']}자 · 도형 바닥 {band['shapes_floor']})"
    )
    return 1 if total else 0


if __name__ == "__main__":
    raise SystemExit(main())
