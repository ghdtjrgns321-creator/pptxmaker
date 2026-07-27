"""배선 스모크 테스트 — 파이프라인 연결이 살아있는지 기계로 확인한다.

왜 있나(2026-07-20): "연결했다"고 선언한 게 실제로는 끊겨 있는 일이 반복됐다(dense가
registry에 안 물림·게이트가 빌드에 안 걸림·다양성 게이트가 golden 타입에 no-op). 손으로 매번
추적하는 건 못 미더우니, 연결 불변식(invariant)을 여기 박아 회귀를 자동으로 잡는다.

    uv run python .claude/skills/pptx-build/scripts/test_wiring.py

각 검사 = 파이프라인의 한 연결점. 종료코드 1 = 어딘가 끊김. green = 전 체인 연결됨.
"""

import ast
import importlib
import json
import re
import sys
import tempfile
from pathlib import Path

HERE = Path(__file__).resolve().parent
ROOT = HERE.parents[3]  # 프로젝트 루트
sys.path.insert(0, str(HERE))

RESULTS = []


def check(name, ok, detail=""):
    RESULTS.append((name, bool(ok), detail))


def _read(rel):
    p = ROOT / rel
    return p.read_text(encoding="utf-8") if p.exists() else ""


# ── ① dense 건강 + 골든 승격 상태 (dense가 실제 빌드에 나올 준비가 됐나) ──────────
# 현 상태(2026-07-20): dense는 골든 **승격 대기**. 승격(registry 재지정)은 골든 회귀 하네스가
# sparse variant 기하에 박제돼 있어 선행 작업이 필요하다 — audit.py를 dense 인지로(accent 런
# 제외·경계 6.9), compare_golden/audit_golden 장별 오라클을 dense 좌표로, s15·s16 완결. 그래서
# 이 테스트는 "승격됐나"가 아니라 **"dense가 건강하고 승격 준비됐나 + 반쯤 승격해 하네스 안
# 깨졌나"**를 본다(상시 적색 게이트는 무시당하므로). 승격이 끝나면 승격분이 registry에 뜬다.
# WIP dense — 하드코딩·content 미파라미터화·미승인. 아직 계약키(DEFAULT) 불필요, 빌드만 확인.
# 완결(파라미터화)되면 여기서 뺀다 → 그때부터 이 테스트가 DEFAULT 노출을 강제.
WIP_DENSE = (
    set()
)  # s16 hero_card 재건·§8 통과로 완결(2026-07-21). 남은 dense 미완은 exec_graph(S6) 신규 예정


def test_dense_ready():
    from goldenfab.registry import LAYOUTS

    reg = {fn.__module__.rsplit(".", 1)[-1] for fn in LAYOUTS.values()}
    dense = sorted(p.stem for p in (HERE / "goldenfab").glob("s*_dense.py"))
    unhealthy = []
    for m in dense:
        mod = importlib.import_module(f"goldenfab.{m}")
        n_default = sum(
            1 for n, v in vars(mod).items() if "DEFAULT" in n.upper() and isinstance(v, dict)
        )
        if not hasattr(mod, "build"):
            unhealthy.append(f"{m}(build 없음)")
        elif m not in WIP_DENSE and n_default != 1:  # 완결 dense = 계약키 dict 정확히 1개 필수
            unhealthy.append(f"{m}(DEFAULT={n_default})")
    ready = sorted(m for m in dense if m not in WIP_DENSE)
    promoted = sorted(m for m in dense if m in reg)
    check(
        "① dense 건강+승격준비 (완결=빌드+계약키1 · WIP=빌드 · 반쯤승격 없음)",
        not unhealthy,
        f"불건강 {unhealthy}"
        if unhealthy
        else f"승격준비 {len(ready)}장 · WIP {sorted(WIP_DENSE)} · 현 승격 "
        + (str(promoted) if promoted else "0(골든=variant, 하네스 재조정 후 승격 예정)"),
    )


# ── ② 밀도 게이트 → 빌드 (오케스트레이터 없이도 빌드가 스스로 채점하나) ──────────
def test_gate_mechanical():
    src = _read(".claude/skills/pptx-build/scripts/build_pptx.py")
    # 빌드 경로가 밀도 게이트를 호출하는지(메서드 존재 + build에서 호출)
    has_gate = "_gate_density" in src
    called = src.count("_gate_density") >= 2  # 정의 + 호출
    check(
        "② 밀도 게이트→빌드 (build_pptx가 스스로 밀도 채점)",
        has_gate and called,
        "build_pptx._gate_density 정의+호출 확인"
        if (has_gate and called)
        else "build_pptx에 빌드시 밀도 게이트 없음(consistency-qa 소프트 단계에만 의존)",
    )


# ── ③ 다양성 게이트 → golden 타입 인식 (golden 덱에 no-op 아닌가) ────────────────
def test_diversity_golden_aware():
    sys.path.insert(0, str(ROOT / ".claude/skills/consistency-qa/scripts"))
    ap = importlib.import_module("audit_pptx")
    # golden 프레임 장은 본문에서 제외돼야
    frame_excluded = all(
        ap.is_frame({"type": t})
        if hasattr(ap, "is_frame")
        else (t in getattr(ap, "FRAME_TYPES", set()))
        for t in ("golden.cover", "golden.toc", "golden.part", "golden.closing")
    )
    # golden 본문 장은 의미있는 시각 키를 받아야(그냥 유일 문자열 반환이면 no-op)
    k1 = ap.visual_key({"type": "golden.problem_grid"})
    k2 = ap.visual_key({"type": "golden.exec_graph"})
    meaningful = k1 != "golden.problem_grid" or hasattr(ap, "GOLDEN_KIND")
    check(
        "③ 다양성 게이트→golden 인식 (프레임 제외 + 본문 계수)",
        frame_excluded and meaningful,
        f"프레임제외={frame_excluded} 본문키={k1!r},{k2!r}",
    )


# ── ④ 즉흥 밋밋 카드 게이트 배선 (hero_card 미사용을 기계로 막나) ────────────────
def test_adhoc_card_gate_wired():
    from goldenfab import audit as A

    # 검사가 존재하고 dense 경로에만 배선됐나(존재≠검사 hollow 방지).
    has_fn = hasattr(A, "check_adhoc_card")
    dense_names = [r[0] for r in A.generic_checks([], "D66E3A", profile="dense")]
    sparse_names = [r[0] for r in A.generic_checks([], "D66E3A", profile="sparse")]
    legacy_names = [r[0] for r in A.generic_checks([], "D66E3A", profile="legacy")]
    # legacy 111장 전수 오탐 0 실측 후 2026-07-25 배선 — sparse 골든만 제외(카드 idiom 상이)
    wired = (
        "§8 즉흥 카드" in dense_names
        and "§8 즉흥 카드" in legacy_names
        and "§8 즉흥 카드" not in sparse_names
    )
    check(
        "④ 즉흥 카드 게이트→generic_checks(dense·legacy) (밋밋 카드 기계 차단)",
        has_fn and wired,
        f"check_adhoc_card 존재={has_fn} · dense·legacy 배선={wired} · sparse 제외 확인"
        if (has_fn and wired)
        else f"미배선(존재={has_fn}, dense={'§8 즉흥 카드' in dense_names}, "
        f"legacy={'§8 즉흥 카드' in legacy_names})",
    )


# ── ④.5 텍스트 겹침 게이트 배선 (배경없는 세로 파일업을 generic_checks 전 경로에서 막나) ──
def test_text_collision_gate_wired():
    """2026-07-24: 텍스트 요소끼리 겹침(카드 목업 3회 반려)이 오딧 사각지대였다 — 하강.
    check_text_overflow(제 상자)·check_picture_overlap(그림)의 형제. **공유 res**라 전 프로파일에
    물려야 한다(존재≠검사 hollow 방지 + 한쪽만 물리면 실전/골든/legacy 한쪽이 샌다)."""
    from goldenfab import audit as A

    has_fn = hasattr(A, "check_text_collision")
    names = {p: [r[0] for r in A.generic_checks([], "D66E3A", profile=p)] for p in A.PROFILES}
    wired_all = all("텍스트 겹침" in v for v in names.values())
    # 검출력 프로브(2026-07-25 추가) — 오탐을 줄이려 넣은 **글자 실점유 클램프**가 검출까지
    # 죽이지 않았나. 상자가 글자로 꽉 찬 진짜 결함(같은 폭 한글 2줄이 0.05~0.25" 포갬)은
    # 반드시 잡혀야 하고, 상자 여유가 큰 경우(글자는 안 겹침)는 안 잡혀야 한다.
    from pptx.util import Inches, Pt

    from goldenfab.reference import new_presentation

    def _inject(gap, box_h):
        prs = new_presentation()
        sl = prs.slides.add_slide(prs.slide_layouts[6])
        y = 2.0
        for k, t in enumerate(
            (
                "임베딩 유사도만으로는 근거의 성립 여부를 가릴 수 없다는 것이 실측의 결론이었다",
                "그래서 판정 경로를 사전과 그래프로 고정해 점수 경쟁 자체를 없애는 설계로 바꿨다",
            )
        ):
            tb = sl.shapes.add_textbox(Inches(1.0), Inches(y), Inches(4.0), Inches(box_h))
            tb.text_frame.word_wrap = True
            r = tb.text_frame.paragraphs[0].add_run()
            r.text = t
            r.font.size = Pt(12)
            if k == 0:
                y += box_h - gap
        return list(sl.shapes)

    detects = [A.check_text_collision(_inject(g, 0.40))[2] for g in (0.05, 0.10, 0.15, 0.25)]
    slack_ok = A.check_text_collision(_inject(0.15, 0.62))[0]  # 상자만 겹침 → 잡히면 오탐
    probe_ok = all(d >= 1 for d in detects) and slack_ok
    check(
        "④.5 텍스트 겹침 게이트→generic_checks(전 프로파일) + 검출력 프로브(주입 결함↔여유 상자)",
        has_fn and wired_all and probe_ok,
        f"존재={has_fn} · {'·'.join(A.PROFILES)} 전부 배선={wired_all} · "
        f'꽉 찬 상자 0.05~0.25" 포갬 전건 검출{detects} · 여유 상자 오탐 0'
        if (has_fn and wired_all and probe_ok)
        else f"미배선/검출력 저하(존재={has_fn}, 배선={wired_all}, 주입검출={detects}, "
        f"여유상자오탐={not slack_ok})",
    )


# ── ⑤ 산출물→덱 에이전트 게이트 (pptx 저장 감지 → deck-smith "완료 실행" 경유 강제) ──
def test_artifact_agent_gate_wired():
    """2026-07-24 전수 감사: 게이트가 경로(goldenfab 소스)만 봐서 스크래치 pptx 작업이
    무검증 통과하던 구멍. 훅 3점(matcher·arm 감지·게이트 검사)이 전부 물려 있어야 한다.
    페어링 교체(2026-07-24): 옛 게이트는 subagent_type "문자열 존재"만 봐서 **거부된**
    deck-smith 호출(tool use rejected)도 거짓통과했다(실세션 실증). 게이트가 tool_use↔
    tool_result 페어링으로 '실제 완료'를 보는지(거부·무응답 불인정) 마커로 못 박는다 —
    이 마커(is_error·tool_result·거부 문자열)가 사라지면 naive grep으로의 회귀."""
    sj = _read(".claude/settings.json")
    arm = _read(".claude/hooks/pptx_arm.py")
    gate = _read(".claude/hooks/pptx_render_gate.sh")
    matcher_ok = "Write|Edit|Bash|PowerShell" in sj
    arm_ok = "pptx_signal" in arm and ".pptx.tsv" in arm
    # 존재 마커 + 페어링(완료 실행) 판정 마커. 후자가 naive grep과의 차이를 강제한다.
    exists_ok = "subagent_type" in gate and "deck-smith" in gate and ".pptx.tsv" in gate
    pairing_ok = (
        "tool_result" in gate and "is_error" in gate and "The user doesn't want to proceed" in gate
    )
    gate_ok = exists_ok and pairing_ok
    check(
        "⑤ 산출물→덱 에이전트 게이트 (pptx 저장 감지 → deck-smith 완료 실행 경유, 거부 불인정)",
        matcher_ok and arm_ok and gate_ok,
        f"matcher={matcher_ok} arm감지={arm_ok} 게이트검사(존재={exists_ok}·페어링={pairing_ok})",
    )


# ── ⑥ 규칙 택소노미 양방향 (design-rules 강제 태그 ↔ 검사 실체 — 유령·고아·부채 래칫) ──
PROSE_DEBT_CAP = (
    4  # [산문·하강대기: 태그 수(§3·§7·§8·E-2). 증감은 P1.5 백로그 경유 + 사용자 승인 전용.
)


def test_rules_taxonomy():
    """2026-07-25: design-rules(570줄)가 목차는 있는데 **강제 상태 축**이 없어 "산문인데
    강제라고 착각"이 반복됐다("규칙이 있어도 러너에 안 물리면 없는 규칙"). 각 § 머리의
    `> **강제:**` 태그를 기계로 양방향 검증한다:
    ①전 § 태그 존재(영구 census) ②[기계: check_x]는 generic_checks에, [기계·골든: check_x]는
    audit_golden에 실재·배선(유령 태그 차단 — 골든 전용 검사를 전 덱 규칙에 달면 유령 강제)
    ③audit.py 전 검사가 문서에 등장(고아 검사 차단 — 새 검사 추가 시 문서 태그 강제 래칫)
    ④[산문·하강대기:] 총수 ≤ CAP(산문 부채 래칫 — 넘으려면 이 상수를 같은 커밋에서 올려야
    해서 부채 증가가 항상 가시화된다). 태그 전이의 정당성은 기계가 판정 못 함(사용자 승인 몫)."""
    doc = _read(".claude/skills/pptx-build/references/design-rules.md")
    audit_src = _read(".claude/skills/pptx-build/scripts/goldenfab/audit.py")
    golden_src = _read(".claude/skills/pptx-build/scripts/audit_golden.py")
    probs = []
    # ① 전 §(##~####) 헤딩 밑 3줄 내 강제 태그 줄
    lines = doc.split("\n")
    for i, ln in enumerate(lines):
        if re.match(r"^#{2,4} ", ln):
            win = lines[i + 1 : i + 4]
            if not any(w.startswith("> **강제:**") or w.startswith("> (구조") for w in win):
                probs.append(f"태그 없음: {ln.lstrip('# ')[:24]}")
    # ② 태그의 check_*가 맞는 러너에 실재·배선
    defined = set(re.findall(r"^def (check_[a-z_]+)", audit_src, re.M))
    gen_body = (
        audit_src.split("def generic_checks", 1)[1].split("\ndef ", 1)[0]
        if "def generic_checks" in audit_src
        else ""
    )
    for m in re.finditer(r"\[(기계·골든|기계)[^\]]*\]", doc):
        cls = m.group(1)
        for tok in re.findall(r"check_[a-z_]+", m.group(0)):
            if tok not in defined:
                probs.append(f"유령 태그: {tok}(audit.py에 없음)")
            elif cls == "기계" and f"{tok}(" not in gen_body:
                probs.append(f"오배선: [기계] {tok}가 generic_checks에 없음")
            elif cls == "기계·골든" and f"{tok}(" not in golden_src:
                probs.append(f"오배선: [기계·골든] {tok}가 audit_golden에 없음")
    # ③ 고아 검사(정의됐는데 문서 미등장) — 새 검사는 반드시 규칙 태그를 동반
    orphans = sorted(c for c in defined if c not in doc)
    if orphans:
        probs.append(f"고아 검사(문서 태그 없음): {orphans}")
    # ④ 산문 부채 래칫
    n_debt = doc.count("[산문·하강대기:")
    if n_debt > PROSE_DEBT_CAP:
        probs.append(
            f"산문 부채 증가: {n_debt} > 상한 {PROSE_DEBT_CAP}(백로그 경유+사용자 승인 필요)"
        )
    check(
        "⑥ 규칙 택소노미 양방향 (전 § 태그 · 유령/고아 0 · 부채 래칫)",
        not probs,
        "; ".join(probs)
        if probs
        else f"§ 태그 전수 · 검사 {len(defined)}종 전부 문서 등장 · 부채 {n_debt}/{PROSE_DEBT_CAP}",
    )


# ── ⑦ goldenfab 고아 모듈 census (장 모듈이 어디에도 안 물린 채 생존 못 하게) ──────
# 허용 = registry 도달(전이 import 폐포) ∪ dense census ∪ 아래 두 상수. 상수를 건드리지 않으면
# 새 고아가 들어올 수 없다(래칫) — EXEMPT는 "정당한 예외", DEBT는 "거취 미결·삭제 대기"로
# **분리**한다. 섞으면 부채가 정당으로 위장되고, 그게 census를 만든 이유를 부정한다.
ORPHAN_EXEMPT = {
    # 정당한 예외만. 각 항목에 "왜 영구히 안 물려도 되나" 한 줄 필수. 현재 해당 없음.
}
ORPHAN_DEBT = {
    # 거취가 아직 안 정해진 것만. 각 항목에 "무엇을 기다리나 + 해소 시 무엇을 하나" 한 줄 필수.
    "s06_mid": "제안[d] 승격 대기 — S6 dense 기준작인데 미배선. 승격 시 이 항목 삭제",
}
ORPHAN_DEBT_CAP = 1  # 감소 전용 래칫. 늘리려면 사용자 승인 + 같은 커밋에서 이 상수를 올려야 한다.


def _intra_imports(path):
    """모듈이 import하는 goldenfab 내부 모듈명 — ast로 정확히(정규식·docstring 오탐 0).

    함수 본문 안의 지연 import까지 잡아야 한다(s11_variants:872가 s11_branch_snap을 그렇게
    쓴다 — 정규식이면 놓치거나 docstring의 사용 예시를 실 import로 오인한다)."""
    out = set()
    for n in ast.walk(ast.parse(path.read_text(encoding="utf-8"))):
        if isinstance(n, ast.ImportFrom):
            if n.level == 1 and n.module:  # from .X import ...
                out.add(n.module.split(".")[0])
            elif n.level == 1:  # from . import X, Y
                out |= {a.name for a in n.names}
            elif n.module and n.module.startswith("goldenfab."):  # from goldenfab.X import ...
                out.add(n.module.split(".")[1])
        elif isinstance(n, ast.Import):
            out |= {a.name.split(".")[1] for a in n.names if a.name.startswith("goldenfab.")}
    return out


def test_no_orphan_modules():
    """2026-07-25: 고아 장 모듈 2개(s06_pilot·s06_proto_f — NAVY 리터럴 하드코딩, 아무 데도
    안 물림)가 **게이트 전면 green 속에서 생존**하다 수동 삭제로만 정리됐다. 원인은 분모의
    구멍이다: dense census(①)의 분모가 `s*_dense.py` glob이라 `sNN_pilot`·`_proto_f`·`_mid`
    같은 이름은 **어떤 census에도 안 걸렸다**. 안 물린 모듈은 다음 사람이 헛짚는 재료가 되고
    (거기서 좌표·색을 베낀다), 죽은 채로 규칙 부채가 된다 — 그래서 분모를 이름 규칙이 아니라
    `goldenfab/sNN_*.py`(장 모듈) **전수**로 놓고, 허용을 도달성으로 증명한다.

    도달성 = registry LAYOUTS 모듈에서 시작한 전이 import 폐포 ∪ dense census 분모. 어디에도
    없으면 FAIL — EXEMPT/DEBT에 올리려면 같은 커밋에서 상수를 건드려야 하므로 고아 유입이
    항상 가시화된다. 기계는 '안 물림'까지만 판정한다 — 그 모듈을 지울지 물릴지는 사용자 몫."""
    from goldenfab.registry import LAYOUTS

    pkg = HERE / "goldenfab"
    denom = sorted(p.stem for p in pkg.glob("s[0-9][0-9]_*.py"))  # 장 모듈만 — sNN_ 형태
    dense = {p.stem for p in pkg.glob("s*_dense.py")}  # ①과 같은 분모(중복 정의 아님, 참조)
    seeds = {fn.__module__.rsplit(".", 1)[-1] for fn in LAYOUTS.values()} | dense
    reach, queue = set(), sorted(seeds)
    while queue:  # 전이 폐포 — 지연 import 포함
        m = queue.pop()
        if m in reach or not (pkg / f"{m}.py").exists():
            continue
        reach.add(m)
        queue += sorted(_intra_imports(pkg / f"{m}.py") - reach)
    orphans = [m for m in denom if m not in reach | set(ORPHAN_EXEMPT) | set(ORPHAN_DEBT)]
    n_debt = len(ORPHAN_DEBT)
    probs = [f"고아 모듈(어디에도 안 물림): {orphans}"] if orphans else []
    if n_debt > ORPHAN_DEBT_CAP:
        probs.append(f"고아 부채 증가: {n_debt} > 상한 {ORPHAN_DEBT_CAP}(사용자 승인 필요)")
    stale = sorted(
        m for m in set(ORPHAN_EXEMPT) | set(ORPHAN_DEBT) if not (pkg / f"{m}.py").exists()
    )
    if stale:  # 해소됐는데 상수에 남은 항목 — 상수가 실물과 갈라지면 다음 고아를 덮는다
        probs.append(f"유령 면제(파일 없음): {stale}")
    check(
        "⑦ goldenfab 고아 모듈 census (분모=sNN_*.py 전수 · 허용=registry 폐포∪dense∪면제)",
        not probs,
        "; ".join(probs)
        if probs
        else f"sNN_*.py {len(denom)}개 = 도달 {len(reach & set(denom))}"
        f"(registry 폐포 {len((reach & set(denom)) - dense)} + dense {len(dense)})"
        f" + 면제 {len(ORPHAN_EXEMPT)} + 부채 {n_debt}/{ORPHAN_DEBT_CAP} · 고아 0",
    )


# ── 통합 A: 밀도 지표 단일 출처 (audit.py 하나, 경쟁 임계 없음) ──────────────────
def test_density_metric_single_source():
    pf = _read(".claude/skills/pptx-build/scripts/preflight_dense.py")
    uses_audit = "A.generic_checks" in pf or "A.check_density" in pf  # 정본 재사용
    no_rival = (
        "def dense_band" not in pf
        and "def check_accent_fills" not in pf  # 자체 accent 중복 제거됨(→audit dense 인지)
        and "def check_offslide" not in pf  # 자체 경계 중복 제거됨(→audit dense 인지)
        and "660" not in pf
    )
    check(
        "통합A 밀도/검사 지표 단일 출처 (preflight가 audit.generic_checks 재사용, 중복 없음)",
        uses_audit and no_rival,
        f"audit 재사용={uses_audit} 중복없음={no_rival}",
    )


# ── 통합 B: 결정표 단일 출처 (layout-matching, 세 어휘 crosswalk) ────────────────
def test_decision_table_single_source():
    lm = _read(".claude/skills/deck-compose/references/layout-matching.md")
    vs = _read(".claude/skills/pptx-visuals/references/visual-selection.md")
    ac = _read(".claude/skills/pptx-visuals/references/archetype-catalog.md")
    lm_ok = "물성" in lm and "FT Visual Vocabulary" in lm and "검색 플라이휠" in lm
    cross = "layout-matching" in vs and "layout-matching" in ac
    check(
        "통합B 결정표 단일 출처 (layout-matching + 세 어휘 crosswalk)",
        lm_ok and cross,
        f"통합표={lm_ok} crosswalk={cross}",
    )


# ── 가드: content_contract가 전 registry 레이아웃에서 해석되나 (①이 계약 안 깼나)
#          + **두 문 다** 배선됐나(2026-07-25: adapted./novel 문이 비어 있던 구멍) ─────
def test_content_contract_resolves():
    from goldenfab.content_contract import required_keys
    from goldenfab.registry import LAYOUTS

    fails = []
    for name, fn in LAYOUTS.items():
        try:
            required_keys(name, fn)
        except Exception as e:  # noqa: BLE001
            fails.append(f"{name}: {type(e).__name__}")
    # 문 2개 배선 — golden.* 경로와 adapted./novel 경로가 **각자** 계약을 물어야 한다.
    # `_render_scripted`에 계약이 없어 골든 글이 그대로 나갈 수 있었다(2026-07-25 실측).
    bp = _read(".claude/skills/pptx-build/scripts/build_pptx.py")
    scripted = bp.split("def _render_scripted", 1)[-1].split("\n    def ", 1)[0]
    golden = bp.split("def _render_golden", 1)[-1].split("\n    def ", 1)[0]
    if "assert_content(" not in golden:
        fails.append("_render_golden에 assert_content 없음")
    if "assert_scripted_content(" not in scripted:
        fails.append("_render_scripted에 assert_scripted_content 없음(골든 글 유출 문)")
    check(
        "가드 content_contract 전 레이아웃 해석 + 두 문(golden·scripted) 배선",
        not fails,
        f"실패 {fails}" if fails else f"{len(LAYOUTS)}개 전부 해석 · golden·scripted 두 문 배선",
    )


# ── 통합 C: 카탈로그 철칙 — 모든 L-ID는 렌더러가 실존 (hollow 카탈로그 방지) ─────
def _renderer_keys():
    """전 렌더러 키 집합을 소스에서 수집(matplotlib·pptx 무거운 import 회피)."""
    vsrc = _read(".claude/skills/pptx-visuals/scripts/visuals.py")
    msrc = _read(".claude/skills/pptx-visuals/scripts/mpl_exhibits.py")
    bsrc = _read(".claude/skills/pptx-build/scripts/build_pptx.py")
    # 차트 = CHART_TYPES 키 + MPL_TYPES + panels(chart 렌더러 분기, build_pptx.chart)
    ct = re.search(r"CHART_TYPES\s*=\s*\{(.*?)\}", vsrc, re.S)
    chart = set(re.findall(r'"([a-z0-9_]+)"\s*:', ct.group(1))) if ct else set()
    mt = re.search(r"MPL_TYPES\s*=\s*\{(.*?)\}", msrc, re.S)
    chart |= set(re.findall(r'"([a-z0-9_]+)"', mt.group(1))) if mt else set()
    chart.add("panels")
    # 다이어그램 = add_diagram의 layout 분기(== 및 in (...))
    diagram = set(re.findall(r'layout == "([a-z0-9_]+)"', vsrc))
    for grp in re.findall(r"layout in \(([^)]*)\)", vsrc):
        diagram |= set(re.findall(r'"([a-z0-9_]+)"', grp))
    # 표계열 = build_pptx.RENDERERS 키
    rb = re.search(r"RENDERERS\s*=\s*\{(.*?)\}", bsrc, re.S)
    table = set(re.findall(r'"([a-z0-9_]+)"\s*:', rb.group(1))) if rb else set()
    return chart, diagram, table


# 카탈로그 철칙이 적용되는 문서 전수(결정표·카탈로그·근거표) — 셋은 같은 축의 다른 표현.
CATALOG_DOCS = (
    ".claude/skills/pptx-visuals/references/visual-selection.md",
    ".claude/skills/pptx-visuals/references/archetype-catalog.md",
    ".claude/skills/deck-compose/references/layout-matching.md",
)
_FAM_TOK = re.compile(r"\b(chart|diagram)\s+([a-z][a-z0-9_]*(?:·[a-z][a-z0-9_]*)*)")
_CODE_SPAN = re.compile(r"`([^`]+)`")
_LID_REF = re.compile(r"\bL(\d{2})\b")


def _doc_visual_claims(rel):
    """문서에서 **렌더러를 요구하는 선언만** 뽑는다 — 오탐 0이 이 함수의 유일한 계약.

    추출 범위를 3중으로 좁힌다: ①`미보유`/`백로그` 섹션 제외(철칙상 렌더러 없음이 정상 상태)
    ②취소선(`~~폐기~~`) 제외 ③가족 접두(`chart `/`diagram `) + **표 행 또는 코드 스팬 안에서만**.
    ③이 산문 오탐의 차단선이다 — 출처 줄의 `"chart chooser"`(visual-selection:17)가 실측에서
    실제로 걸렸다. 맨 낱말(table·metrics·flow·mapping…)은 **일부러 안 잡는다**: 산문 빈출이라
    오탐 위험이 크고 그 절반은 LID_KEY가 이미 덮는다(미탐 < 오탐 — 오탐 게이트는 백스톱에
    음소거돼 죽는다). 반환: [(가족, 종류, 줄)], [(L-ID, 줄)]."""
    fam, lids = [], []
    skip = False
    for i, raw in enumerate(_read(rel).split("\n"), 1):
        if raw.startswith("#"):
            skip = "백로그" in raw or "미보유" in raw
        if skip:
            continue
        ln = re.sub(r"~~.*?~~", "", raw)
        scopes = [ln] if ln.lstrip().startswith("|") else _CODE_SPAN.findall(ln)
        for sc in scopes:
            for f, blob in _FAM_TOK.findall(sc):
                fam += [(f, k, i) for k in blob.split("·")]
        lids += [(f"L{n}", i) for n in _LID_REF.findall(ln)]
    return fam, lids


def test_catalog_renderers_exist():
    """카탈로그 철칙(archetype-catalog:8) "모든 행은 렌더러가 실존"의 기계 미러.
    정본 = recommend_archetypes.LID_KEY(L01~L38 → family:kind) + _renderer_keys() 하나.

    왜 문서까지 확장했나(2026-07-25): 옛 통합C는 **LID_KEY만** 봐서 결정표·카탈로그 문서는
    시야 밖이었다. 그래서 hollow 3행(`diagram process_grid`·`diagram category_spine_table`·
    `diagram mapping` — visual-selection 표 B)이 렌더러 없이 결정표에 앉아 **게이트 green 속에서
    생존**했고 2026-07-25 수동 삭제로만 정리됐다. 문서가 부품을 지시하는데 그 부품이 없으면
    deck-composer는 지시대로 배정하고 빌더가 즉흥한다 — 재발 차단이 이 확장의 목적.
    두 방향을 본다: ①문서의 `chart|diagram <종류>` 선언 → _renderer_keys()에 실존
    ②문서의 L-ID 참조 → LID_KEY에 실존(2홉: 문서→LID_KEY→렌더러). 추출 정밀도는
    _doc_visual_claims 참조(오탐 0이 전제 — 정밀 > 양)."""
    sys.path.insert(0, str(ROOT / ".claude/skills/pptx-visuals/scripts"))
    from recommend_archetypes import LID_KEY

    chart, diagram, table = _renderer_keys()
    missing = []
    for lid, key in LID_KEY.items():
        fam, _, kind = key.partition(":")
        if fam == "chart":
            ok = kind in chart
        elif fam == "diagram":
            ok = kind in diagram
        else:  # fam ∈ {table,metrics,two_column,bullets}; "table:matrix"는 base table
            ok = fam in table
        if not ok:
            missing.append(f"{lid}={key}")
    n_tok = n_lid = 0
    for rel in CATALOG_DOCS:
        fams, lids = _doc_visual_claims(rel)
        doc = rel.rsplit("/", 1)[-1]
        for f, kind, ln in fams:
            n_tok += 1
            if kind not in (chart if f == "chart" else diagram):
                missing.append(f"{doc}:{ln} `{f} {kind}`(렌더러 없음)")
        for lid, ln in lids:
            n_lid += 1
            if lid not in LID_KEY:
                missing.append(f"{doc}:{ln} {lid}(LID_KEY 밖)")
    check(
        "통합C 카탈로그 렌더러 실존 (LID_KEY + 결정표·카탈로그 문서 — hollow 행 금지)",
        not missing,
        f"미해석(렌더러 없는 행) {missing}"
        if missing
        else f"L-ID {len(LID_KEY)}종 전부 실존 (chart {len(chart)}·diagram {len(diagram)}·표 "
        f"{len(table)}) · 문서 {len(CATALOG_DOCS)}종 선언 {n_tok}토큰+L-ID참조 {n_lid} 전부 해석",
    )


# ── 통합 E: 골든 도해 문법 축 등재 검증 (A축이 hollow 아닌가 · 결정표↔카탈로그 양방향) ──
# 왜 있나(2026-07-25 어휘 정본 이전 2단계): 카탈로그·결정표가 (B) 범용 어휘만 가리켜 골든 도해가
# **자동 경로에서 후보로 오르지 못했다**(실전 덱이 폴백으로 샌 원인). (A)를 등재하면서 통합C와 같은
# hollow를 만들 위험이 생긴다 — (A)는 generic 렌더러가 없는 **장 함수**라 `add_diagram(layout=…)`으로
# 검증할 수 없다. 그래서 (A)의 철칙을 다르게 박는다: **선언한 `모듈:심볼`이 실존하고, 선언한
# `c:키`가 그 모듈에서 실제로 읽히고, `const:`가 실존하고, 항목 수 계약이 명시돼 있다.**
# 키 실존만 보면 "주면 바뀐다"는 거짓 주장을 못 잡는다(실측: dense 모듈이 주입 키를 무시하고
# 모듈 상수를 그리는 사례 6건 — 통합F).
GOLDEN_GRAMMAR_IDS = {  # 등재 확정 15종. census 분모 — 추가·삭제는 사용자 승인 + 같은 커밋.
    "G01", "G02", "G04", "G05", "G06", "G08", "G09", "G10",
    "G11", "G13", "G14", "G15", "G17", "G19", "G21",
}  # fmt: skip
# 제외 판정(사용자 확정)이 조용히 되살아나는 것 차단: G16(=L04 doughnut 동등·B 우위)·
# G18(=L35 icon_rows 흡수)·G20(미러 매트릭스 재탕 유인) + 부품 백로그 G03·G07·G12·G22·G23은
# 위 집합에 없으므로, 등재 표에 다시 나타나면 census 불일치로 FAIL 난다.
# 항목 수 파생 완료 = 좌표가 개수에서 나온다(리터럴 리스트 zip 절단 없음). **증가 전용** 래칫.
# 2026-07-25 G10·G11·G14 추가(3단계 선행 작업) — 자는 `grid.pitch`/`grid.track` 단일 출처이고
# 수용 한계 초과는 시끄럽게 죽는다. G19는 한계 카드만 파생(밖 1+1은 좌우 대칭 차단 = 문법 자체 ·
# 캔버스 1.9"에 칩 0.66 2단이 안 들어간다는 실측) → 여기 없음.
# 2026-07-26 G02·G04 추가 — 부품으로 꺼내면서 고정 계약이 파생이 됐다. G02는 항목 4 고정(zip
# 절단)이 pitch 파생으로, G04는 노드 7·태그 4 고정이 "폭은 종류에서·간격은 남는 자리에서"로.
ARITY_DERIVED = {
    "G02",
    "G04",
    "G05",
    "G08",
    "G10",
    "G11",
    "G14",
}
GRAMMAR_CATALOG = ".claude/skills/pptx-visuals/references/archetype-catalog.md"
GRAMMAR_TABLE = ".claude/skills/deck-compose/references/layout-matching.md"
# 등재가 아닌 언급(제외 판정·부품 백로그·창고)은 줄 단위로 제외 — 통합C의 취소선 제외와 같은 규율.
_GSKIP = ("백로그", "제외", "승격 금지", "폐기")
_GID = re.compile(r"\bG\d{2}\b")
# 하위 패키지 허용(2026-07-26 부품화) — 도해가 `goldenfab/figures/<부품>.py`로 나가기 시작했고
# 슬래시를 막아두면 부품으로 꺼낸 문법이 등재될 자리가 없어 카탈로그가 옛 자리를 가리킨 채 굳는다.
_GSRC = re.compile(r"`goldenfab/((?:[A-Za-z0-9_]+/)*[A-Za-z0-9_]+)\.py:([A-Za-z_][A-Za-z0-9_]*)`")
_GKEY = re.compile(r"`c:([a-z0-9_]+)`")
_GCONST = re.compile(r"`const:([A-Z][A-Z0-9_]*)`")
_GARITY = re.compile(r"`\[(파생|고정:[^\]`]+)\]`")
_GTYPE = re.compile(r"`(golden|adapted)\.([a-z_]+)`")


def _module_symbols(path):
    """모듈 최상위 심볼(def·대입·import) — 문법 출처가 실존하는 함수·상수인지 ast로 확인."""
    out = set()
    for n in ast.parse(path.read_text(encoding="utf-8")).body:
        if isinstance(n, (ast.FunctionDef, ast.AsyncFunctionDef, ast.ClassDef)):
            out.add(n.name)
        elif isinstance(n, ast.Assign):
            out |= {t.id for t in n.targets if isinstance(t, ast.Name)}
        elif isinstance(n, ast.AnnAssign) and isinstance(n.target, ast.Name):
            out.add(n.target.id)
        elif isinstance(n, (ast.Import, ast.ImportFrom)):
            out |= {(a.asname or a.name).split(".")[0] for a in n.names}
    return out


from goldenfab.content_contract import content_keys_read as _content_keys_read  # noqa: E402

# ↑ 2026-07-25: 이 정규식은 여기 사본으로 있었는데, 같은 판정을 **프로덕션 게이트**
# (`assert_scripted_content`)도 하게 되면서 정본을 content_contract로 옮기고 여기서 재사용한다.
# 사본을 남기면 한쪽만 고쳐져 "테스트는 통과하는데 게이트는 다르게 판정"이 된다.


def _grammar_lines(rel):
    """표 행(| …)만 · 등재 아닌 언급 줄 제외 → [(줄번호, 줄)]."""
    out = []
    for i, raw in enumerate(_read(rel).split("\n"), 1):
        ln = re.sub(r"~~.*?~~", "", raw)
        if ln.lstrip().startswith("|") and not any(t in ln for t in _GSKIP):
            out.append((i, ln))
    return out


def test_golden_grammar_axis():
    from goldenfab.registry import LAYOUTS

    pkg = HERE / "goldenfab"
    probs = []
    # ① 카탈로그 등재 행 파싱 + 행별 계약 검증
    seen = {}
    for lineno, ln in _grammar_lines(GRAMMAR_CATALOG):
        gids = _GID.findall(ln)
        srcs = _GSRC.findall(ln)
        if not gids or not srcs:  # 문법 행 = G-ID + 출처를 동시에 가진 행
            continue
        if len(gids) > 1:
            probs.append(f"카탈로그:{lineno} 한 행에 G-ID {gids} — 행=문법 1개 규약 위반")
            continue
        gid = gids[0]
        seen[gid] = lineno
        mods, syms = [], []
        for mod, sym in srcs:
            p = pkg / f"{mod}.py"
            if not p.exists():
                probs.append(f"{gid}: 출처 모듈 없음 goldenfab/{mod}.py")
                continue
            if sym not in _module_symbols(p):
                probs.append(f"{gid}: 심볼 없음 {mod}.py:{sym}")
            mods.append(mod)
            syms.append(sym)
        srcs_text = "\n".join((pkg / f"{m}.py").read_text(encoding="utf-8") for m in mods)
        read_keys = _content_keys_read(srcs_text)
        top_syms = set().union(*(_module_symbols(pkg / f"{m}.py") for m in mods)) if mods else set()
        for key in _GKEY.findall(ln):  # 선언한 계약 키가 실제로 읽히나(거짓 주장 차단)
            if key not in read_keys:
                probs.append(f"{gid}: 계약 키 `c:{key}`를 출처 {mods}가 안 읽는다")
        for const in _GCONST.findall(ln):
            if const not in top_syms:
                probs.append(f"{gid}: `const:{const}` 상수가 출처 {mods}에 없음")
        arity = _GARITY.findall(ln)
        if len(arity) != 1:
            probs.append(f"{gid}: 항목 수 토큰 {arity} — `[파생]`/`[고정:N]` 정확히 1개 필수")
        elif (arity[0] == "파생") != (gid in ARITY_DERIVED):
            probs.append(f"{gid}: 항목 수 `[{arity[0]}]`가 ARITY_DERIVED 래칫과 불일치")
    # ② census 양방향 — 등재 집합 == 분모 == 결정표 참조 집합(유령 문법·고아 어휘 0)
    if set(seen) != GOLDEN_GRAMMAR_IDS:
        probs.append(
            f"카탈로그 등재 불일치 — 누락 {sorted(GOLDEN_GRAMMAR_IDS - set(seen))} "
            f"· 미승인 등재 {sorted(set(seen) - GOLDEN_GRAMMAR_IDS)}"
        )
    lm_ref = {g for _, ln in _grammar_lines(GRAMMAR_TABLE) for g in _GID.findall(ln)}
    if lm_ref != GOLDEN_GRAMMAR_IDS:
        probs.append(
            f"결정표 참조 불일치 — 결정표에 없음 {sorted(GOLDEN_GRAMMAR_IDS - lm_ref)} "
            f"· 카탈로그에 없는 참조 {sorted(lm_ref - GOLDEN_GRAMMAR_IDS)}"
        )
    # ③ 두 문서의 장 타입 토큰이 registry에 실존(golden./adapted. 오타·유령 타입 차단)
    for rel in (GRAMMAR_CATALOG, GRAMMAR_TABLE):
        doc = rel.rsplit("/", 1)[-1]
        for i, raw in enumerate(_read(rel).split("\n"), 1):
            for pre, key in _GTYPE.findall(raw):
                if key not in LAYOUTS:
                    probs.append(f"{doc}:{i} `{pre}.{key}` — registry LAYOUTS에 없음")
    check(
        "통합E 골든 문법 축 등재 (출처 실존 · 계약 키 실제 읽힘 · 항목수 래칫 · 결정표 양방향)",
        not probs,
        "; ".join(probs)
        if probs
        else f"문법 {len(seen)}종 등재 = 분모 {len(GOLDEN_GRAMMAR_IDS)} = 결정표 참조 "
        f"· 파생 {sorted(ARITY_DERIVED)}·고정 {len(GOLDEN_GRAMMAR_IDS - ARITY_DERIVED)}종 "
        f"· 장 타입 전부 registry 실존",
    )


# ── 통합 I: 부품 선택 규칙 배선 (선언 ↔ 창고 ↔ 스킬 ↔ 회귀검증 양방향) ──────────
# 왜 있나(2026-07-27): 이 프로젝트는 "기계가 고른다"고 **선언**해놓고 기계는 다른 걸 보던
# 사고를 겪었다(결정표가 "(A) 먼저"라고 적었는데 판정 기계는 (B)만 후보로 올렸다). 그래서
# 선택 규칙은 네 지점이 동시에 물려 있어야 한다:
#   ① 창고 전수 = select.PARTS (부품 파일이 늘면 목록도 늘어야 — 손 목록이 뒤처지면 유령)
#   ② 부품마다 accepts 선언 (축 전수 · 개수 범위)
#   ③ 스킬이 러너를 부른다 (산문만 있으면 그 길로 안 간다 — 이 프로젝트의 반복 실패)
#   ④ 회귀 검증이 존재하고 통과한다 (골든이 손으로 고른 선택을 재현)
def test_selection_wired():
    from goldenfab import frames as FR
    from goldenfab import select as SEL

    pkg = HERE / "goldenfab" / "figures"
    on_disk = {
        p.stem for p in pkg.glob("*.py") if p.stem not in ("__init__", "elements")
    }
    ghost = sorted(set(SEL.PARTS) - on_disk)      # 목록에 있는데 파일 없음
    orphan = sorted(on_disk - set(SEL.PARTS))     # 파일 있는데 목록에 없음(선택에서 영영 빠진다)
    cat = SEL.catalog()
    no_acc = sorted(n for n, m in cat.items() if not m.get("accepts"))
    bad_axis = sorted(
        n for n, m in cat.items()
        if m.get("accepts") and set(SEL.AXES) - set(m["accepts"])
    )
    skill = _read(".claude/skills/deck-compose/SKILL.md")
    wired = "pick_parts.py" in skill and "verify_selection.py" in skill
    runner = (HERE / "pick_parts.py").exists() and (HERE / "verify_selection.py").exists()
    # 회귀 검증을 실제로 돌린다 — 존재만 확인하면 hollow다
    import subprocess
    rc = subprocess.run(
        [sys.executable, str(HERE / "verify_selection.py")], capture_output=True
    ).returncode
    frames_ok = bool(FR.FRAMES) and all("needs" in f for f in FR.FRAMES.values())
    ok = not (ghost or orphan or no_acc or bad_axis) and wired and runner and rc == 0 and frames_ok
    check(
        "통합I 부품 선택 규칙 배선 (창고 census · accepts 전수 · 스킬 호출 · 회귀 green)",
        ok,
        f"부품 {len(cat)}종 census 일치 · accepts 축 {len(SEL.AXES)}개 전수 · "
        f"틀 {len(FR.FRAMES)}종 needs 선언 · 스킬 배선={wired} · verify rc={rc}"
        if ok
        else f"유령 {ghost} · 고아 {orphan} · accepts 없음 {no_acc} · 축 누락 {bad_axis} · "
        f"스킬배선={wired} · 러너={runner} · verify rc={rc}",
    )


# ── 통합 F: dense 계약 구멍 래칫 (주입해도 무시되는 키 — 골든 글 유출 경계 고정) ──────
# 왜 있나(2026-07-25 통합E 작업 중 실측): registry에 물린 sparse 골든 모듈은 DEFAULT 키를 **전부**
# 읽는데(미사용 0/21), dense 기준작은 일부 키를 무시하고 **모듈 상수를 그린다**. 그래서 dense를
# adapted로 쓰면 그 자리에 골든(K-IFRS) 글이 남는다 — `_render_scripted`는 `assert_content`를
# 부르지 않으므로 content_contract도 못 잡는다. 고치는 건 dense 파라미터화(별 작업·사용자 결정)이고,
# 여기서는 **경계를 고정**한다: 새 구멍 유입 = FAIL, 줄어들면 통과(감소 전용 래칫).
# `bar`·`narratives`는 dense 개편의 의도된 폐지(§8 전폭 바 금지·서술행 제거)이고, 나머지가 실구멍이다.
DENSE_CONTRACT_HOLES = {
    # headline·kicker는 2026-07-26 dense 승격에서 메웠다(PARAM 오라클이 잡은 이식 누락)
    # converge_note·limits는 2026-07-26 fan_in 부품화로 메웠다 — 부품이 글을 소유하지 않으므로
    # 모듈 상수(LIMITS)가 사라지고 주입한 값이 실제로 그려진다.
    "s04_dense": {"bar", "source"},
    "s08_dense": {"bar", "narratives"},
    "s09_dense": {"bar", "narratives"},
    "s10_dense": {"caption", "points"},
    "s11_dense": {"bar", "flow_foot", "narratives", "steps"},
    "s12_dense": {
        "bar",
        "caption",
        "diamond",
        "flow_display",
        "flow_generate",
        "narratives",
        "reject_note",
        "validate_subhead",
    },  # fmt: skip
    "s14_dense": {"bar", "source"},  # headline·kicker 메움(2026-07-26)
    "s15_dense": set(),
    "s16_dense": set(),
}


def test_dense_contract_holes():
    pkg = HERE / "goldenfab"
    probs, shrunk = [], []
    for m in sorted(p.stem for p in pkg.glob("s*_dense.py")):
        mod = importlib.import_module(f"goldenfab.{m}")
        keys = set()
        for n, v in vars(mod).items():
            if "DEFAULT" in n.upper() and isinstance(v, dict):
                keys |= set(v)
        holes = keys - _content_keys_read((pkg / f"{m}.py").read_text(encoding="utf-8"))
        base = DENSE_CONTRACT_HOLES.get(m)
        if base is None:
            probs.append(f"{m}: 기준선 미등록(새 dense 모듈 — 구멍 {sorted(holes)})")
            continue
        if holes - base:
            probs.append(f"{m}: 신규 계약 구멍 {sorted(holes - base)}")
        if base - holes:
            shrunk.append(f"{m}-{sorted(base - holes)}")
    check(
        "통합F dense 계약 구멍 래칫 (주입해도 무시되는 키 — 증가 시 FAIL)",
        not probs,
        "; ".join(probs)
        if probs
        else f"구멍 {sum(len(v) for v in DENSE_CONTRACT_HOLES.values())}건 기준선 유지"
        + (f" · 개선(기준선 조일 것) {shrunk}" if shrunk else ""),
    )


# ── 가드: 오딧은 **읽기 전용**이어야 한다 (검사기가 산출물을 바꾸면 안 된다) ────────
def test_audit_is_readonly():
    """2026-07-25 실측 버그: `sh.line.color.type`·`run.font.color`는 python-pptx가 XML을 **생성**한다
    (`<a:ln><a:solidFill/></a:ln>`). PowerPoint는 빈 solidFill을 **검정 헤어라인**으로 그리고,
    빌드 게이트는 `prs.save()` **전에** 도니 그 테두리가 파일에 저장됐다 — legacy 배선 후 fixture
    렌더 눈검증에서 전 텍스트 상자 테두리로 발각(게이트는 전부 green이었다). 색 읽기를 lxml
    읽기 전용으로 바꿨고, 여기서 **전 검사 통과 후 XML 무변화**를 못 박는다."""
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    from goldenfab import audit as A
    from goldenfab.reference import new_presentation

    prs = new_presentation()
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    tb = sl.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(0.5))
    tb.text_frame.word_wrap = True
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = "판정 단어와 색을 읽는 대상"
    r.font.size = Pt(12)
    sq = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6), Inches(1), Inches(3), Inches(1.2))
    sq.fill.solid()
    sq.fill.fore_color.rgb = RGBColor(0xF2, 0xF3, 0xF5)
    ov = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.2), Inches(1.1), Inches(0.3), Inches(0.3))
    ov.line.fill.background()
    shapes = list(sl.shapes)
    before = [sh._element.xml for sh in shapes]
    for p in A.PROFILES:
        A.generic_checks(shapes, "D66E3A", profile=p)
    after = [sh._element.xml for sh in list(sl.shapes)]
    changed = [i for i, (a, b) in enumerate(zip(before, after)) if a != b]
    check(
        "가드 오딧 읽기 전용 (검사가 도형 XML을 변형하지 않음 — 빈 solidFill 테두리 오염 방지)",
        not changed,
        f"변형된 도형 인덱스 {changed} — 검사기가 산출물을 오염시킨다"
        if changed
        else f"{len(shapes)}도형 × 프로파일 {len(A.PROFILES)}종 전수 무변화",
    )


# ── 통합 G: legacy 경로 정본 게이트 배선 + 오탐 0 기준선 ──────────────────────────
# 왜 있나(2026-07-25 4단계 실측): `generic_checks`가 legacy 타입에 **한 번도 안 돌았다** —
# `_gate_density`·`audit_deck`는 golden./adapted./novel만 대상이고 `consistency-qa/audit_pptx`는
# goldenfab.audit을 **import조차 안 했다**. 그래서 남은 legacy 어휘(차트 12·다이어그램 9·표·문서형)가
# 무검증으로 나갔다. 배선했으니 **끊기지 않게 못 박고**, 동시에 "물린 검사 집합"을 고정한다:
# 오탐 3종(채움률·판정 대비·밀도 밴드)이 몰래 legacy로 들어오면 상시 적색이 되고 게이트가 음소거된다.
# fixture(assets/legacy-refs.json)는 생존 어휘 1장씩의 **통과작 기준선**이다 — 실제로 빌드해
# 채점하므로 임계를 잘못 조이면 여기서 즉시 FAIL 난다(오탐 0 래칫).
LEGACY_WIRED = {  # legacy 111장 전수 실측으로 승인된 배선 집합 (증감 = 사용자 승인)
    "§2 accent 상한",
    "P4③ 노드 재탕",
    "P4⑩ 노드 클래스",
    "§F 그림 침범",
    "텍스트 넘침",
    "텍스트 겹침",
    "P2③④ 경계",
    "§8 즉흥 카드",
}
LEGACY_UNWIRED = {  # 오탐 실측 때문에 **의도적으로** 빼둔 것 — 들어오면 상시 적색
    "P4④ 채움률",  # 풀폭 배너(12.13"×0.36 다크 스트립) 구조상 26% → 오탐 1/111
    "P4⑤ 판정 단어 대비",  # check_matrix 회색 ✓ 3.22·강조어 accent 3.41 = 어휘 idiom → 오탐 4/111
    "§6-D 밀도 밴드",  # 골든 458자 하한은 legacy 프레임에 부적합 → 93/111. 정본은 audit_pptx 60단어
}


def test_legacy_gate_wired():
    from goldenfab import audit as A

    probs = []
    names = {r[0] for r in A.generic_checks([], "D66E3A", profile="legacy")}
    if names != LEGACY_WIRED:
        probs.append(
            f"legacy 검사 집합 변동: 추가 {sorted(names - LEGACY_WIRED)} 누락 {sorted(LEGACY_WIRED - names)}"
        )
    if names & LEGACY_UNWIRED:
        probs.append(
            f"오탐 검사가 legacy에 유입: {sorted(names & LEGACY_UNWIRED)}(상시 적색 → 음소거)"
        )
    # 러너 2곳 배선(존재≠검사) — QA 러너가 정본을 import하고 legacy 프로파일로 부르나
    qa = _read(".claude/skills/consistency-qa/scripts/audit_pptx.py")
    if "from goldenfab import audit" not in qa or 'profile="legacy"' not in qa:
        probs.append("audit_pptx가 정본 게이트를 legacy 프로파일로 안 부름")
    if "legacy_gate(" not in qa or qa.count("legacy_gate") < 2:
        probs.append("audit_pptx.legacy_gate 정의/호출 누락")
    bp = _read(".claude/skills/pptx-build/scripts/build_pptx.py")
    if 'profile = "legacy"' not in bp:
        probs.append("build_pptx._gate_density가 legacy 장을 프로파일 배정 안 함")
    # fixture 어휘 census — LID_KEY 분모 전수 커버(새 legacy 어휘는 기준선 슬라이드를 동반해야 한다)
    sys.path.insert(0, str(ROOT / ".claude/skills/pptx-visuals/scripts"))
    from recommend_archetypes import LID_KEY

    fx_path = ROOT / ".claude/skills/pptx-build/assets/legacy-refs.json"
    fx = json.loads(fx_path.read_text(encoding="utf-8"))

    def vkey(s):
        t = s["type"]
        if t == "chart":
            return f"chart:{'panels' if s.get('panels') else s.get('chart_type', 'bar')}"
        if t == "diagram":
            return f"diagram:{s['layout']}"
        if t == "table":
            return "table:matrix" if s.get("style") == "matrix" else "table"
        return t

    covered = {vkey(s) for s in fx["slides"]}
    missing = set(LID_KEY.values()) - covered
    if missing:
        probs.append(f"fixture 미커버 어휘 {sorted(missing)}(기준선 슬라이드 필요)")
    # 실제 빌드 + 채점 — 기준선이 진짜 통과하나(hollow fixture 방지 · 임계 오탐 래칫)
    import build_pptx as B

    deck = B.Deck(B.load_brand(ROOT / ".claude/skills/pptx-build/assets/brand-kit.yaml"))
    deck.render_dir = Path(tempfile.mkdtemp(prefix="legacy_refs_")) / "render"
    deck._spec_dir = ROOT
    prs = deck.build(fx)
    viol = []
    for i, (sl, sp) in enumerate(zip(prs.slides, fx["slides"]), 1):
        for n, ok, d, _c in A.generic_checks(list(sl.shapes), "D66E3A", profile="legacy"):
            if not ok:
                viol.append(f"s{i}({vkey(sp)}) {n}: {d[:60]}")
    if viol:
        probs.append(f"기준선 위반 {len(viol)}건 — 임계가 통과작을 잡는다(오탐): {viol[:3]}")
    check(
        "통합G legacy 정본 게이트 배선 (검사 집합 고정 · 러너 2곳 · fixture 어휘 census + 오탐 0)",
        not probs,
        "; ".join(probs)
        if probs
        else f"검사 {len(names)}종 배선·오탐 3종 제외 · fixture {len(fx['slides'])}장 "
        f"(L-ID {len(LID_KEY)}종 전수 커버) 위반 0",
    )


# ── 통합 H: (A) 우선 판정이 **추천 기계**에 내려왔나 (결정표 1순위가 산문으로 남지 않게) ──
# 왜 있나(2026-07-25 4단계): 결정표는 "(A) 1순위"로 재작성됐는데 판정 기계(recommend_archetypes)는
# (B) L-ID만 랭킹해서 자동 경로를 타면 여전히 (B)만 후보로 올라왔다 — "고를 수 있게 만들었지만
# 그 길로 가는지는 산문"이던 자리. 세 축으로 못 박는다:
# ① GRAMMAR_KEY의 G-ID가 카탈로그 등재 집합 안에 있고 장 타입이 registry에 실재(유령 추천 차단)
# ② 추천 출력이 (A)를 **먼저** 낸다(순서가 곧 1순위)
# ③ check_candidates가 "(A) 후보 0개 + why_not_golden 없음"을 **FAIL로** 잡는다 + 그 반대는 통과
#    (민감도·오탐 양방향 프로브 — 의무가 vacuous하거나 과잉이면 여기서 걸린다)
def test_grammar_first_wired():
    sys.path.insert(0, str(ROOT / ".claude/skills/pptx-visuals/scripts"))
    import check_candidates as CC
    import recommend_archetypes as RA
    from goldenfab.registry import LAYOUTS

    probs = []
    unknown = set(RA.GRAMMAR_KEY) - GOLDEN_GRAMMAR_IDS
    if unknown:
        probs.append(f"카탈로그 미등재 G-ID 추천: {sorted(unknown)}")
    for gid, t in RA.GRAMMAR_KEY.items():
        if t.split(".", 1)[1] not in LAYOUTS:
            probs.append(f"{gid}: 추천 장 타입 {t}가 registry에 없음(유령 추천)")
    ranked = {g for v in RA.GRAMMAR_RANK.values() for g in v}
    if ranked - set(RA.GRAMMAR_KEY):
        probs.append(f"랭킹에만 있는 G-ID(키 없음): {sorted(ranked - set(RA.GRAMMAR_KEY))}")
    # ② 순서 — (A) 1순위 물성의 추천 첫 항목이 A축인가
    spec = {"slides": [{"type": "bullets", "title": "T", "bullets": ["가", "나", "다"]}]}
    r = RA.recommend(spec)[0]
    if not r["grammar_first"] or r["recommended"][0]["axis"] != "A":
        probs.append(f"정성 물성인데 추천 1순위가 (A)가 아님: {r['recommended'][:1]}")
    # ③ 의무 양방향 프로브
    cand_b = {
        "slides": [
            {
                "no": 1,
                "shape": "정성",
                "visual_candidates": [
                    {"type": "diagram", "layout": "icon_rows", "emphasis": "x", "why_not": "이유"},
                    {"type": "diagram", "layout": "venn", "emphasis": "x"},
                    {"type": "table", "emphasis": "x", "why_not": "이유"},
                ],
            }
        ]
    }
    cand_a = json.loads(json.dumps(cand_b))
    cand_a["slides"][0]["visual_candidates"][1] = {"type": "golden.problem_grid", "emphasis": "x"}

    def run(cand):
        import io
        from contextlib import redirect_stdout

        argv = sys.argv
        d = Path(tempfile.mkdtemp(prefix="cand_"))
        (d / "c.json").write_text(json.dumps(cand, ensure_ascii=False), encoding="utf-8")
        (d / "s.json").write_text(json.dumps(spec, ensure_ascii=False), encoding="utf-8")
        sys.argv = ["x", str(d / "c.json"), str(d / "s.json")]
        buf = io.StringIO()
        try:
            with redirect_stdout(buf):
                rc = CC.main()
        finally:
            sys.argv = argv
        return rc, buf.getvalue()

    rc_b, out_b = run(cand_b)
    rc_a, _out_a = run(cand_a)
    if rc_b == 0 or "(A) 골든 문법 1순위" not in out_b:
        probs.append("(B)만 낸 후보군을 통과시킴 — (A) 검토 의무가 vacuous")
    if rc_a != 0:
        probs.append("(A) 후보를 낸 후보군을 FAIL시킴 — 의무 오탐")
    check(
        "통합H (A) 우선 판정 기계화 (추천 1순위=A · 유령 추천 0 · (A)검토 의무 양방향)",
        not probs,
        "; ".join(probs)
        if probs
        else f"G-ID {len(RA.GRAMMAR_KEY)}종 전부 등재·registry 실재 · (A)먼저 물성 "
        f"{sorted(RA.GRAMMAR_RANK)} · 의무 검출 O·오탐 X",
    )


# ── 통합 D: brand-kit 밖 hex 하드코딩 래칫 (색 단일 출처 — design-rules §2) ────────
def test_no_new_hardcoded_hex():
    # 규율 "신규 hex 금지 — 색은 brand-kit 단일 출처"를 래칫으로 기계화. 기준선 = 승인된 예외
    # (mpl 데이터 베이스 회색 + 흰색 그라디언트 끝점 2건). 초과 = 새 하드코딩 유입 → FAIL.
    base = {
        ".claude/skills/pptx-visuals/scripts/visuals.py": 0,
        ".claude/skills/pptx-build/scripts/build_pptx.py": 0,
        ".claude/skills/pptx-visuals/scripts/mpl_exhibits.py": 2,
    }
    pat = re.compile(r"#[0-9a-fA-F]{6}\b|0x[0-9a-fA-F]{6}\b")
    over = []
    for rel, b in base.items():
        n = len(pat.findall(_read(rel)))
        if n > b:
            over.append(f"{rel.rsplit('/', 1)[-1]}: {n} > 기준선 {b}")
    check(
        "통합D brand-kit 밖 hex 하드코딩 래칫 (초과=신규 하드코딩)",
        not over,
        f"초과 {over}"
        if over
        else "visuals 0 · build_pptx 0 · mpl 2(승인 예외) — 신규 하드코딩 없음",
    )


def main():
    for t in (
        test_dense_ready,
        test_gate_mechanical,
        test_diversity_golden_aware,
        test_adhoc_card_gate_wired,
        test_text_collision_gate_wired,
        test_artifact_agent_gate_wired,
        test_rules_taxonomy,
        test_no_orphan_modules,
        test_density_metric_single_source,
        test_decision_table_single_source,
        test_content_contract_resolves,
        test_catalog_renderers_exist,
        test_golden_grammar_axis,
        test_dense_contract_holes,
        test_audit_is_readonly,
        test_legacy_gate_wired,
        test_grammar_first_wired,
        test_no_new_hardcoded_hex,
        test_selection_wired,
    ):
        try:
            t()
        except Exception as e:  # noqa: BLE001
            check(t.__name__, False, f"예외 — {type(e).__name__}: {e}")

    print("═" * 70)
    fails = 0
    for name, ok, detail in RESULTS:
        mark = "OK  " if ok else "✘FAIL"
        print(f"{mark} {name}\n       {detail}")
        fails += not ok
    print("═" * 70)
    print(
        "WIRING PASS — 전 체인 연결됨" if not fails else f"WIRING FAIL {fails}건 — 끊긴 연결 있음"
    )
    return 1 if fails else 0


if __name__ == "__main__":
    sys.exit(main())
