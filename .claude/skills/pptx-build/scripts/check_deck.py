"""덱 기계 게이트 — 완성된 pptx를 열어 **눈으로 잡기 어려운 좌표 결함**만 본다.

## 왜 두 개뿐인가 (2026-08-03)

사용자가 무하네스 산출물 19장을 훑고 결함 4건을 짚었다. 넷을 다 재보니 **둘만 기계로 갈렸다.**

| 결함 | 좋다고 한 10장 | 문제라 한 장 | 판정 |
| --- | --- | --- | --- |
| 우변 초과 | 전부 0 | 0.14 · 0.54 · 0.44 · 0.46" | **갈린다 → 검사** |
| 이미지 폭 | 2.45" 이상 | 1.25 · 1.71 · 1.95" | **갈린다 → 검사** |
| 밀도(덮은 면적) | 137~200% | 172% | 안 갈린다 → 안 함 |
| 블록 경계 정렬 | 46·53쌍 어긋남 | 15쌍 | 안 갈린다(반대) → 안 함 |

뒤 둘은 검사하지 않는다. 표 행 높이는 `blocks.compare_table`이 경고만 내고, **정렬은 조판이
좌표로 직접 잡는다** — 이걸 구조로 막겠다고 만든 `kit.Panel`은 실전 19장이 0회 부른 채
새 컨텍스트만 6/6이 불러 화면을 회색으로 덮었고, 2026-08-03에 삭제했다(근거는 kit.py 주석).

규칙을 늘리지 않는 이유는 이 저장소가 이미 겪었기 때문이다: "이 장 밋밋해"(한 장의 지적)
→ "카드 4장으로 채워라"(모든 장의 명령) → 36장이 동일한 4카드.

    uv run python .claude/skills/pptx-build/scripts/check_deck.py <경로.pptx>
"""

import re
import sys
from pathlib import Path

import yaml
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

sys.path.insert(0, str(Path(__file__).parent))
from deckkit import outline as O  # noqa: E402

_B = yaml.safe_load((Path(__file__).parent / "deckkit/brand.yaml").read_text("utf-8"))
FR, LIM = _B["frame"], _B["limits"]
ACC = _B["colors"]["accent"]
RIGHT = FR["slide_w"] - FR["margin"]
BOT = FR["bottom"]
MATERIAL = None  # --material=<README> 를 주면 재료 커버리지를 보고한다
OUTLINE = None  # --outline=<01.5_outline.md> 를 주면 헤드라인 일치를 검사한다


def _nums(text: str) -> set:
    """수치 토큰. 재료가 덱으로 얼마나 옮겨졌나를 세는 단위 — 반복해도 늘지 않는다."""
    out = set()
    for m in re.finditer(r"\d[\d,]*\.?\d*\s*%?", text):
        v = m.group().strip().replace(" ", "")
        if len(v.replace(",", "").replace(".", "").replace("%", "")) >= 2:
            out.add(v)
    return out


def coverage(deck: Path, material: Path) -> str:
    """재료 커버리지 보고 — **판정하지 않는다.**

    글자 수 하한은 텍스트로 채우게 만들어 시각을 밀어냈다(2026-08-03 실측). 커버리지는
    같은 절을 늘려 써도 안 오르고 새 사실을 가져와야 오르므로, 같은 왜곡을 만들지 않는다.
    그래도 하한이 아니라 **보고**로 둔다 — 19장이 재료 전량을 담을 수는 없고, 무엇을 담을지는
    계약이 정한 것이지 기계가 정할 것이 아니다.
    """
    src = _nums(material.read_text(encoding="utf-8"))
    txt = " ".join(
        sh.text_frame.text
        for sl in Presentation(str(deck)).slides
        for sh in sl.shapes
        if sh.has_text_frame
    )
    got = _nums(txt) & src
    r = len(got) / len(src) if src else 1.0
    mark = "" if r >= LIM["material_cover_warn"] else f"  ← {LIM['material_cover_warn']:.0%} 미만"
    return f"재료 커버리지 {len(got)}/{len(src)}종 ({r:.1%}){mark}"


def headlines(outline: Path) -> dict:
    """계약의 `헤드라인(확정)` 열을 읽는다. 열이 없으면 빈 dict — 구판 계약은 검사하지 않는다."""
    out, idx = {}, None
    for ln in outline.read_text(encoding="utf-8").splitlines():
        if not ln.strip().startswith("|"):
            continue
        cells = [c.strip() for c in ln.strip().strip("|").split("|")]
        if idx is None:
            for i, c in enumerate(cells):
                if c.startswith("헤드라인"):
                    idx = i
            continue
        m = re.match(r"^S?(\d+)$", cells[0])
        if m and idx is not None and idx < len(cells) and cells[idx] not in ("", "—"):
            out[int(m.group(1))] = cells[idx]
    return out


def formal(outline: Path) -> list[str]:
    """계약 `## 정형 장 문구` 절의 문구들 — 표지 4칸 + 부 제목·리드.

    2026-08-03 실측: 부 이름 5개가 **전부** 조판 단계 창작이었다(헤드라인은 12장 중 8장).
    정형 장이 오히려 더 많이 흔들렸다 — 계약에 칸이 없으면 지어내기 때문이다.
    """
    want, sec, idx, lead, first = [], False, None, None, True
    for ln in outline.read_text(encoding="utf-8").splitlines():
        if ln.startswith("## "):
            sec = ln.startswith("## 정형 장 문구")
            continue
        if not sec or not ln.strip().startswith("|"):
            idx, lead, first = None, None, True
            continue
        c = [x.strip() for x in ln.strip().strip("|").split("|")]
        if set("".join(c)) <= set("-: "):
            continue
        if first:
            first = False
            for i, v in enumerate(c):
                if v.startswith(("문구", "제목")):
                    idx = i
                if v.startswith("리드"):  # 부 표에만 있다. 표지 표의 '재료 근거'와 섞이면 안 된다
                    lead = i
            continue
        if idx is not None and idx < len(c) and c[idx] not in ("", "—"):
            want.append(c[idx])
            if lead is not None and lead < len(c) and c[lead] not in ("", "—"):
                want.append(c[lead])
    return want


def nav(prs, outline: Path) -> list[str]:
    """부 내비게이션이 계약과 같은가 — 문구·순서·**강조된 부**까지.

    본문 전 장의 상단에 5줄이 깔리므로, 틀리면 한 장이 아니라 덱 전체가 틀린다. 강조까지 보는
    이유는 문구만 맞고 강조가 엉뚱한 부에 붙으면 내비가 오히려 위치를 속이기 때문이다.
    """
    want = [f"{r} {t}" for r, t, _ in O.parts(outline)]
    of = O.part_of(outline)
    if not want or not of:
        return []
    bad = []
    for n, sl in enumerate(prs.slides, 1):
        lines = [sh for sh in sl.shapes if sh.name.startswith("부내비:")]
        if not lines and n in of:
            bad.append(f"s{n:02d} 부 내비가 없다 — 계약은 이 장을 「{of[n]}」부로 둔다")
            continue
        if not lines:
            continue
        got = [sh.text_frame.text.strip() for sh in lines]
        if got != want:
            bad.append(f"s{n:02d} 부 내비가 계약과 다르다 — 계약 {want} / 조판 {got}")
            continue
        cur = [
            sh.name.split(":", 1)[1]
            for sh in lines
            for p in sh.text_frame.paragraphs
            for r in p.runs
            if r.font.color and r.font.color.type is not None and str(r.font.color.rgb) == ACC
        ]
        if len(set(cur)) != 1:
            bad.append(f"s{n:02d} 부 내비 강조가 {len(set(cur))}개다 — 정확히 하나여야 한다")
        elif n in of and cur[0] != of[n]:
            bad.append(f"s{n:02d} 부 내비 강조 「{cur[0]}」 ≠ 계약 「{of[n]}」")
    return bad


def check(path: Path, body=()) -> int:
    prs = Presentation(str(path))
    bad = []
    for n, sl in enumerate(prs.slides, 1):
        for sh in sl.shapes:
            if sh.left is None or sh.width is None:
                continue
            x, w = Emu(sh.left).inches, Emu(sh.width).inches
            y, hh = Emu(sh.top).inches, Emu(sh.height).inches
            # 블리드는 결함이 아니다 — 좌변이 여백 밖이거나(전폭 배경), 우변이 슬라이드 끝에
            # 정확히 닿는 것(목차 띠처럼 한쪽만 흘리는 설계). 어중간하게 넘는 것만 잡는다.
            full_bleed = (
                x < FR["margin"] - LIM["right_tolerance"]
                or x + w >= FR["slide_w"] - LIM["right_tolerance"]
            )
            if not full_bleed and x + w > RIGHT + LIM["right_tolerance"]:
                txt = sh.text_frame.text[:18].replace("\n", " ") if sh.has_text_frame else ""
                bad.append(f's{n:02d} 우변 초과 {x + w - RIGHT:+.2f}"  {txt}')
            # 하변 초과 — 푸터를 뚫고 나간다. 2026-08-03: 이걸 안 봐서 캡처 아래 목록이
            # 푸터 밖으로 나갔는데 WARN 0 · 게이트 0이었고 렌더를 보고서야 알았다.
            if not full_bleed and y + hh > BOT + LIM["right_tolerance"] and y < BOT:
                bad.append(f's{n:02d} 하변 초과 {y + hh - BOT:+.2f}" (푸터 침범)')
            # 아이콘은 판독 하한 대상이 아니다 — 화면 캡처만 본다.
            if (
                sh.shape_type == MSO_SHAPE_TYPE.PICTURE
                and LIM["icon_max_w"] <= w < LIM["image_min_w"]
            ):
                bad.append(
                    f's{n:02d} 이미지 폭 {w:.2f}" < 하한 {LIM["image_min_w"]}"'
                    " — 개수를 줄이거나 크롭할 것"
                )
    if OUTLINE:
        # 조판이 계약 헤드라인을 그대로 썼나. 2026-08-03: 본문 12장 중 8장이 조판 단계에서
        # 바뀌었고 사용자는 그걸 볼 자리가 없었다. 제목은 계약에서 확정한다.
        want = headlines(OUTLINE)
        for n, sl in enumerate(prs.slides, 1):
            if n not in want:
                continue
            # 머리띠 = 헤더 밑줄 위. 좌표를 박지 않고 토큰에서 뽑는다 — 2026-08-03에 밑줄을
            # 1.18에서 올렸더니 `0.4 < top < 1.1`로 박아둔 창이 제목을 통째로 놓쳤다.
            got = [
                sh.text_frame.text.strip()
                for sh in sl.shapes
                if sh.has_text_frame
                and sh.top is not None
                and Emu(sh.top).inches < FR["rule"]
                and not sh.name.startswith("부내비:")
                and sh.text_frame.text.strip()
            ]
            if want[n] not in got:
                bad.append(
                    f"s{n:02d} 헤드라인이 계약과 다르다 — 계약 「{want[n]}」"
                    f" / 조판 「{got[0] if got else '없음'}」"
                )
    if OUTLINE:
        all_txt = " ".join(
            sh.text_frame.text for sl in prs.slides for sh in sl.shapes if sh.has_text_frame
        )
        for w in formal(OUTLINE):
            if w not in all_txt:
                bad.append(f"정형 장 문구가 덱에 없다 — 계약 「{w[:44]}」")
        bad += nav(prs, OUTLINE)
    print(f"{path.name} — {len(prs.slides._sldIdLst)}장 · 결함 {len(bad)}건")
    if MATERIAL:
        print("  ·", coverage(path, MATERIAL))
    for b in bad:
        print("  ★", b)
    return 1 if bad else 0


def main() -> int:
    if len(sys.argv) < 2:
        print(__doc__)
        return 2
    p = Path(sys.argv[1])
    if not p.exists():
        raise SystemExit(f"없는 파일: {p}")
    body = ()
    global MATERIAL, OUTLINE
    for a in sys.argv[2:]:
        if a.startswith("--outline="):
            OUTLINE = Path(a.split("=", 1)[1])
        if a.startswith("--material="):
            MATERIAL = Path(a.split("=", 1)[1])
        if a.startswith("--body="):
            body = tuple(int(v) for v in a.split("=", 1)[1].split(",") if v.strip())
    return check(p, body)


if __name__ == "__main__":
    sys.exit(main())
