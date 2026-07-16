"""check_contract — 아웃라인 계약(주문서) vs 빌드된 pptx 대조. **정답은 계약이다.**

왜 필요한가(2026-07-15 실증): 계약에 8장 전부 "README §번호"와 실제 이미지 4종이 "확정"으로
적혔는데 반영 0장이었다. 그런데 ②.5 목업 승인·④ QA·회귀 게이트가 전부 PASS를 냈다.
게이트들이 계약을 안 읽고 자기들끼리 일관성만 검사했기 때문이다 — 골든과 같으면 통과,
팩트시트에 있으면 통과. 골든은 이 프로젝트로 만들어져서 기본값이 팩트시트에 "대응"한다.
그래서 "계약을 0% 지킨 덱"이 만점을 받았다.

이 스크립트는 그 구멍을 메운다. 검사 대상은 골든도 팩트시트도 아니고 **사용자가 확정한 계약**.

사용:
    uv run python .claude/skills/consistency-qa/scripts/check_contract.py \
        _workspace/01.5_outline.md _workspace/02_deck-spec.json _workspace/deck.pptx
출력: _workspace/04_contract-check.md (행별 PASS/FAIL + 슬라이드 번호), 종료코드 1 = FAIL 존재
"""

import hashlib
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

ROOT = Path(__file__).resolve().parents[4]
sys.path.insert(0, str(ROOT / ".claude/skills/pptx-build/scripts"))

# 템플릿 플레이스홀더 — 계약이 요구한 적 없는데 나가면 사고.
# 값 = 고치는 법(되돌릴 대상). 어디서 오는지에 따라 대상이 다르므로 명시한다 —
# "pptx-builder로 되돌림"이라고만 쓰면 content로 못 고치는 항목에서 무한 왕복이 된다.
FORBIDDEN = {
    "My Company": "brand-kit.yaml `brand.name`을 실제 이름으로 (브랜드 단일 출처)",
    "회사 로고": (
        "**코드 수정 필요** — layouts.py(toc)·s21_closing.py에 하드코딩된 로고 슬롯 라벨. "
        "content·brand-kit으로 못 고친다. brand-kit 구동으로 바꾸되 golden/·goldenfab/ 양쪽을 "
        "함께 고쳐야 compare_golden이 산다"
    ),
    "Lorem ipsum": "deck-composer — 실제 본문으로 교체",
    "TODO": "deck-composer — 미완 표식 제거",
    "샘플 텍스트": "deck-composer — 실제 본문으로 교체",
}

# 장 제목이 역할 라벨인 행 — 리터럴 텍스트로 대조하지 않는다(제목이 아니라 역할명)
ROLE_LABELS = {"표지", "목차", "간지"}
ROLE_PREFIXES = ("마무리",)

# 앵커 수치: 3자리 이상 정수 · 분수(78/92) · 소수 퍼센트(59.1%)
ANCHOR_RE = re.compile(r"\d+\.\d+%|\d+/\d+|\d{3,}")

# 계약 제목에서 골라낼 수 없는 잡토큰
STOP_TOKENS = {"vs", "및", "그리고", "병합", "유지", "장"}


def parse_contract(path):
    """계약 md의 장 목록 표 → 행 리스트. 표 형식: | # | 부 | 장 제목 | 전하고 싶은 것 | 재료 출처 | 이미지 | 상태 |"""
    rows = []
    for line in path.read_text(encoding="utf-8").splitlines():
        if not line.strip().startswith("|"):
            continue
        cells = [c.strip() for c in line.strip().strip("|").split("|")]
        if len(cells) < 7 or not cells[0].isdigit():
            continue
        rows.append(
            {
                "no": int(cells[0]),
                "part": cells[1],
                "title": cells[2],
                "message": cells[3],
                "source": cells[4],
                "image": cells[5],
                "status": cells[6],
            }
        )
    return rows


def slide_text(slide):
    out = []
    for sh in slide.shapes:
        if sh.has_text_frame:
            out.append(sh.text_frame.text)
        if getattr(sh, "has_table", False):
            for r in sh.table.rows:
                for c in r.cells:
                    out.append(c.text)
    return "\n".join(out)


def slide_image_md5(slide):
    out = set()
    for sh in slide.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            out.add(hashlib.md5(sh.image.blob).hexdigest())
    return out


def golden_strings(layout_key):
    """레이아웃의 살아있는 골든 기본값에서 식별력 있는 문자열(len>=12) 수집.

    골든 기본값 JSON 스냅샷(_workspace/golden_defaults.json)은 스테일이다(S4 재설계 후
    14키로 남아있으나 코드는 9키). 그래서 코드에서 직접 읽는다 — 자동 최신.
    """
    from goldenfab.content_contract import golden_defaults
    from goldenfab.registry import LAYOUTS

    fn = LAYOUTS.get(layout_key)
    if fn is None:
        return []
    d = golden_defaults(fn)
    if d is None:
        return []
    out = []

    def walk(v):
        if isinstance(v, str):
            if len(v.strip()) >= 12:
                out.append(v.strip())
        elif isinstance(v, dict):
            for x in v.values():
                walk(x)
        elif isinstance(v, (list, tuple)):
            for x in v:
                walk(x)

    walk(d)
    return out


def layout_key(type_str):
    """spec type → 골든 기본값 대조용 레이아웃 키.

    `adapted.<key>`도 그 골든 레이아웃이 출발점이므로(2026-07-16 변형 출발점 배선)
    기본값 유출을 같은 잣대로 검사한다 — 변형 스크립트에 골든 글이 복사돼 남는 사고 방지.
    """
    return re.sub(r"^(golden|adapted)\.", "", type_str or "")


def title_tokens(title):
    """계약 장 제목 → 식별 토큰. 번호표식(①—+)·잡토큰 제거."""
    cleaned = re.sub(r"[①②③④⑤\[\]\-—+·]", " ", title)
    cleaned = re.sub(r"\*\*|\[병합\]|\[1장 유지\]", " ", cleaned)
    toks = [t for t in re.split(r"\s+", cleaned) if len(t) >= 2 and t not in STOP_TOKENS]
    return toks


def main():
    if len(sys.argv) != 4:
        print(__doc__)
        return 2
    contract_p, spec_p, pptx_p = (Path(a) for a in sys.argv[1:4])

    import json

    rows = parse_contract(contract_p)
    spec = json.loads(spec_p.read_text(encoding="utf-8"))
    prs = Presentation(str(pptx_p))
    slides = list(prs.slides)

    # 계약이 지정한 이미지 파일 → md5
    img_dir = ROOT / "_workspace/images"
    asset_md5 = {p.name: hashlib.md5(p.read_bytes()).hexdigest() for p in img_dir.glob("*.png")}

    out = [
        "# check_contract — 계약(주문서) vs 빌드 산출물 대조",
        "",
        f"- 계약: `{contract_p}` (행 {len(rows)}개)",
        f"- 스펙: `{spec_p}` · 산출물: `{pptx_p}` (슬라이드 {len(slides)}장)",
        "- **정답은 계약이다.** 골든 기본값·팩트시트 대응은 통과 근거가 아니다.",
        "",
    ]
    fails = []

    # ── 0. 장 수 ──
    if len(rows) != len(slides):
        fails.append(f"장 수 불일치: 계약 {len(rows)} vs pptx {len(slides)}")
        out.append(f"**FAIL — 장 수: 계약 {len(rows)} ≠ pptx {len(slides)}**\n")
    else:
        out.append(f"장 수 {len(rows)}/{len(rows)} 일치\n")

    out += [
        "| # | 장 제목 | 제목대조 | 골든오염 | 계약이미지 | 금지어 | 앵커수치 |",
        "| - | ------- | -------- | -------- | ---------- | ------ | -------- |",
    ]

    for row in rows:
        n = row["no"]
        if n > len(slides):
            fails.append(f"S{n}: 슬라이드 없음")
            continue
        sl = slides[n - 1]
        text = slide_text(sl)
        spec_sl = spec["slides"][n - 1] if n <= len(spec["slides"]) else {}
        layout = layout_key(spec_sl.get("type", ""))

        # ── 1. 제목 대조 (역할 라벨은 부 컬럼으로) ──
        expect = row["part"] if row["title"] in ROLE_LABELS and row["part"] != "—" else row["title"]
        is_role = row["title"] in ("표지", "목차") or row["title"].startswith(ROLE_PREFIXES)
        if is_role:
            title_v = "—(역할)"
        else:
            toks = title_tokens(expect)
            hit = [t for t in toks if t in text]
            if hit:
                title_v = f"OK({hit[0]})"
            else:
                title_v = f"**FAIL** {toks}"
                fails.append(f"S{n}: 계약 제목 토큰 {toks} 중 슬라이드에 등장 0개")

        # ── 2. 골든 기본값 오염 ── (이 검사가 이번 사고의 본체)
        leaked = [g for g in golden_strings(layout) if g in text]
        if leaked:
            golden_v = f"**FAIL** {len(leaked)}건"
            fails.append(
                f'S{n}({layout}): 골든 기본값 {len(leaked)}건 유출 — 예: "{leaked[0][:40]}…"'
            )
        else:
            golden_v = "OK"

        # ── 3. 계약 이미지 실존 (md5 대조 — 아무 그림이나가 아니라 그 그림) ──
        want = re.findall(r"[\w-]+\.png", row["image"])
        if not want:
            img_v = "—"
        else:
            have = slide_image_md5(sl)
            miss = [w for w in want if asset_md5.get(w) not in have]
            if miss:
                img_v = f"**FAIL** 누락 {miss}"
                fails.append(f"S{n}: 계약 확정 이미지 미배치 {miss} (PICTURE {len(have)}개)")
            else:
                img_v = f"OK {len(want)}종"

        # ── 4. 금지어 ──
        bad = [w for w in FORBIDDEN if w in text]
        if bad:
            forb_v = f"**FAIL** {bad}"
            for w in bad:
                fails.append(f'S{n}: 템플릿 플레이스홀더 "{w}" → {FORBIDDEN[w]}')
        else:
            forb_v = "OK"

        # ── 5. 앵커 수치 ──
        anchors = set(ANCHOR_RE.findall(row["message"]))
        if not anchors:
            anch_v = "—"
        else:
            miss_a = sorted(a for a in anchors if a not in text)
            if miss_a:
                anch_v = f"**FAIL** {miss_a}"
                fails.append(f"S{n}: 계약 수치 {miss_a} 슬라이드에 없음")
            else:
                anch_v = f"OK {len(anchors)}개"

        out.append(
            f"| {n} | {row['title'][:22]} | {title_v} | {golden_v} | {img_v} | {forb_v} | {anch_v} |"
        )

    out += ["", "## 판정", ""]
    if fails:
        out.append(f"**FAIL — 위반 {len(fails)}건**")
        out.append("")
        out += [f"- {f}" for f in fails]
        out += [
            "",
            "되돌릴 대상: **deck-composer** (골든오염·앵커수치·제목) / **pptx-builder** (이미지).",
            "플레이스홀더는 위 항목별 지시를 따를 것 — 출처가 달라 대상이 다르다"
            "(brand-kit / 레이아웃 코드).",
            "계약에 '확정'인 항목을 '사용자 선택 대기'로 격하하지 말 것 — 이미 지시받은 사항이다.",
        ]
    else:
        out.append(f"**PASS — 계약 {len(rows)}행 전수 이행, 위반 0건**")

    dst = ROOT / "_workspace/04_contract-check.md"
    dst.write_text("\n".join(out) + "\n", encoding="utf-8")
    print("\n".join(out[-min(len(out), 24) :]))
    print(f"\n-> {dst}")
    return 1 if fails else 0


if __name__ == "__main__":
    sys.exit(main())
