#!/usr/bin/env python
"""빌드된 pptx를 재파싱해 네이티브 객체·골격·폰트 일관성을 기계 검증한다.

deck-spec과 대조하여 표/차트가 이미지로 새지 않았는지, 골격이 지켜졌는지 확인한다.
판단(셀링 흐름 등)은 사람/에이전트 몫 — 여기선 기계로 셀 수 있는 것만 센다.

usage: python audit_pptx.py <out.pptx> <deck-spec.json> [--brand brand-kit.yaml]
exit code 0=PASS, 1=FAIL
"""

import argparse
import json
import sys
from pathlib import Path

import yaml
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# 이미지 경로 차트 유형은 pptx-visuals가 단일 출처
sys.path.insert(0, str(Path(__file__).resolve().parents[2] / "pptx-visuals" / "scripts"))
from mpl_exhibits import MPL_TYPES  # noqa: E402

# "박스+화살표" 어휘 — 다양성 게이트 3(30% 상한) 분모·분자 정의
BOX_LAYOUTS = {"flow", "layers", "cards", "branch", "from_to"}
# 프레임(정형) 레이아웃 — 본문 다양성 계수에서 제외. golden.<layout>·adapted.<layout>도
# 접두사를 벗겨 판정한다(golden.cover/toc/part/closing이 본문으로 오계수되던 버그, 2026-07-20).
FRAME_LAYOUTS = {"cover", "toc", "part", "section", "cta", "closing"}
FRAME_TYPES = {"cover", "toc", "part", "section", "cta"}  # 하위호환(레거시 타입)


def _base_layout(t):
    """golden.<layout>·adapted.<layout> → <layout>, 그 외는 타입 그대로."""
    if t.startswith("golden.") or t.startswith("adapted."):
        return t.split(".", 1)[1]
    return t


def is_frame(sl):
    """정형 장인가 — golden/adapted 접두사를 벗겨 판정(다양성 본문에서 제외)."""
    return _base_layout(sl["type"]) in FRAME_LAYOUTS


def visual_key(sl):
    """본문 슬라이드의 시각 어휘 키 — 다양성 게이트의 계수 단위."""
    t = sl["type"]
    if t == "chart":
        return f"chart:{'panels' if sl.get('panels') else sl.get('chart_type', 'bar')}"
    if t == "diagram":
        return f"diagram:{sl.get('layout', 'flow')}"
    if t.startswith("golden.") or t.startswith("adapted."):
        return f"golden:{_base_layout(t)}"  # golden 레이아웃도 정규화 계수(각기 고유 어휘)
    if t == "novel":
        return f"novel:{sl.get('script', '?')}"
    return t  # table·bullets·metrics·two_column


def diversity_checks(slides):
    """다양성 게이트 4종 — 아키타입 세트 제약(archetype-catalog.md)을 기계로 강제."""
    body = [sl for sl in slides if not is_frame(sl)]
    keys = [visual_key(sl) for sl in body]
    out = []
    # 게이트 1(쿨다운): 동일 시각 유형 간격 ≥3 — 연속(간격1)·한 장 건너(간격2) 모두 위반
    near = [
        f"본문{i + 1}·{j + 1}={keys[i]}"
        for i in range(len(keys))
        for j in range(i + 1, min(i + 3, len(keys)))
        if keys[i] == keys[j]
    ]
    out.append(("diversity_cooldown(동일 유형 간격 ≥3)", not near, f"위반 {len(near)}쌍 {near}"))
    # 게이트 2: 덱 전체 최소 5종 (본문 5장 이상일 때)
    kinds = len(set(keys))
    need = 5 if len(body) >= 5 else len(set(keys)) or 1
    out.append((f"diversity_min_kinds(≥{need}종)", kinds >= need, f"{kinds}종/{len(body)}장"))
    # 게이트 3: 박스 다이어그램 30% 이하
    boxes = sum(
        1 for sl in body if sl["type"] == "diagram" and sl.get("layout", "flow") in BOX_LAYOUTS
    )
    ratio = boxes / len(body) if body else 0
    out.append(("diversity_box_ratio(≤30%)", ratio <= 0.30, f"{boxes}/{len(body)}={ratio:.0%}"))
    # 게이트 4: 동일 유형 덱 전체 ≤2회
    from collections import Counter

    over = {k: n for k, n in Counter(keys).items() if n > 2}
    out.append(("diversity_max_repeat(동일 유형 ≤2회)", not over, f"초과 {over or '0건'}"))
    return out


def collect_fonts(prs):
    fonts = set()
    for slide in prs.slides:
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.name:
                        fonts.add(run.font.name)
    return fonts


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx")
    ap.add_argument("spec")
    ap.add_argument("--brand", default="")
    args = ap.parse_args()

    prs = Presentation(args.pptx)
    spec = json.loads(Path(args.spec).read_text(encoding="utf-8"))
    slides = spec["slides"]

    # 네이티브 표 생성 경로 3종: table 타입 + diagram band_table + chart sub_table
    want_tables = (
        sum(1 for s in slides if s["type"] == "table")
        + sum(1 for s in slides if s["type"] == "diagram" and s.get("layout") == "band_table")
        + sum(1 for s in slides if s["type"] == "chart" and s.get("sub_table"))
    )

    def is_image(s):  # 하이브리드: 확장 유형·render:image는 mpl PNG로 삽입된다
        return s.get("chart_type") in MPL_TYPES or s.get("render") == "image"

    # 멀티패널 chart는 패널 수만큼 네이티브 차트를 만든다
    want_charts = sum(
        len(s.get("panels") or [None]) for s in slides if s["type"] == "chart" and not is_image(s)
    )
    want_exhibit_pics = sum(1 for s in slides if s["type"] == "chart" and is_image(s))

    got_tables = got_charts = got_pics = 0
    for slide in prs.slides:
        for sh in slide.shapes:
            if sh.has_table:
                got_tables += 1
            elif sh.has_chart:
                got_charts += 1
            elif sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                got_pics += 1

    checks = []
    checks.append(
        (
            "slide_count",
            len(prs.slides) == len(slides),
            f"{len(prs.slides)}/{len(slides)}",
        )
    )
    checks.append(("native_tables", got_tables == want_tables, f"{got_tables}/{want_tables}"))
    checks.append(("native_charts", got_charts == want_charts, f"{got_charts}/{want_charts}"))
    # 이미지 익스히빗(mpl PNG)은 정확히 spec에 선언된 수만큼만 — 네이티브가 이미지로 새면 FAIL
    if want_exhibit_pics:
        checks.append(
            (
                "image_exhibits(선언 수 이상)",
                got_pics >= want_exhibit_pics,
                f"{got_pics}/{want_exhibit_pics}(아이콘 포함 가능)",
            )
        )
    # 골격: 첫 두 장 cover/toc, 마지막 cta
    types = [s["type"] for s in slides]
    skeleton_ok = types[:2] == ["cover", "toc"] and types[-1] == "cta"
    checks.append(("skeleton(cover→toc..cta)", skeleton_ok, str(types[:2] + ["..", types[-1]])))

    # 정보 밀도: 본문 슬라이드(표지·목차·섹션·CTA 제외) 단어수 하한 — 빈 슬라이드 차단
    MIN_BODY_WORDS = 60
    frame_types = {"cover", "toc", "part", "section", "cta"}
    thin = []
    for idx, (sl, spec_sl) in enumerate(zip(prs.slides, slides), 1):
        if spec_sl["type"] in frame_types:
            continue
        words = sum(len(sh.text_frame.text.split()) for sh in sl.shapes if sh.has_text_frame)
        if words < MIN_BODY_WORDS:
            thin.append(f"s{idx}({spec_sl['type']}:{words})")
    n_body = sum(1 for s in slides if s["type"] not in frame_types)
    checks.append(
        (
            f"density(본문 {MIN_BODY_WORDS}단어 미만 0장)",
            not thin,
            f"미달 {len(thin)}/{n_body}" + (f" {thin}" if thin else ""),
        )
    )

    # 다양성 게이트 3종(v4) — "다양하게"를 말이 아니라 기계 규칙으로 강제
    checks.extend(diversity_checks(slides))

    # 구성 규칙(deck-compose 계약을 기계로 강제 — 프롬프트 규칙은 요동, 게이트는 불변)
    n_bullets = types.count("bullets")
    checks.append(("bullets_slides(≤2)", n_bullets <= 2, f"{n_bullets}/2"))
    n_parts = types.count("part")
    part_ok = n_parts >= 3 or len(slides) < 12
    checks.append(
        ("part_navigation(≥3 when slides≥12)", part_ok, f"part={n_parts}, slides={len(slides)}")
    )

    # 폰트 일관성 (brand-kit 제공 시)
    if args.brand:
        brand = yaml.safe_load(Path(args.brand).read_text(encoding="utf-8"))
        allowed = {v for v in brand["fonts"].values()}
        used = collect_fonts(prs)
        stray = used - allowed
        checks.append(
            (
                "font_consistency",
                not stray,
                f"stray={sorted(stray)} allowed={sorted(allowed)}",
            )
        )

    print(f"pptx: {args.pptx}")
    all_ok = True
    for name, ok, detail in checks:
        all_ok &= ok
        print(f"  [{'PASS' if ok else 'FAIL'}] {name}: {detail}")
    if got_pics:
        print(f"  [info] pictures={got_pics} (표/차트가 그림으로 샜는지 확인 필요)")
    print("RESULT:", "PASS" if all_ok else "FAIL")
    return 0 if all_ok else 1


if __name__ == "__main__":
    sys.exit(main())
