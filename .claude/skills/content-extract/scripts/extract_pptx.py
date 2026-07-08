#!/usr/bin/env python
"""NotebookLM 등 외부 pptx에서 재료(텍스트·표·차트데이터·이미지)를 추출한다.

차트 추출 3단계 폴백(A안):
  1) 네이티브 차트 → categories/series 수치 직접 추출 (charts.json)
  2) 도형 뭉치 차트 → 텍스트로 남으므로 슬라이드 텍스트에서 수치 복원은 사람이/에이전트가 판단
  3) 이미지 차트 → images/ 에 저장하고 extracted.md에 "이미지 폴백" 표기

usage: python extract_pptx.py <out_dir> <in1.pptx> [in2.pptx ...]
출력: <out_dir>/extracted.md, charts.json, images/
"""

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def walk_shapes(shapes):
    """그룹 도형 내부까지 평탄화해 순회한다(외부 pptx는 그룹이 흔함)."""
    for sh in shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from walk_shapes(sh.shapes)
        else:
            yield sh


def extract_chart(sh):
    """네이티브 차트에서 수치를 뽑는다. 실패하면 None(폴백 대상)."""
    try:
        ch = sh.chart
        plot = ch.plots[0]
        return {
            "chart_type": str(ch.chart_type),
            "categories": [str(c) for c in plot.categories],
            "series": {
                s.name or f"series{i}": list(s.values) for i, s in enumerate(plot.series, 1)
            },
        }
    except Exception as e:  # 차트 XML 변형이 다양해 실패는 폴백으로 넘긴다
        return {"error": str(e)}


def extract_table(sh):
    tbl = sh.table
    return [[cell.text.strip() for cell in row.cells] for row in tbl.rows]


def extract_file(path, img_dir, charts, lines):
    prs = Presentation(path)
    stem = Path(path).stem
    lines.append(f"\n## 파일: {Path(path).name} ({len(prs.slides)}슬라이드)\n")
    for si, slide in enumerate(prs.slides, 1):
        lines.append(f"### 슬라이드 {si}")
        texts, n_img = [], 0
        for sh in walk_shapes(slide.shapes):
            if sh.has_text_frame and sh.text_frame.text.strip():
                texts.append(sh.text_frame.text.strip().replace("\n", " / "))
            if getattr(sh, "has_table", False) and sh.has_table:
                rows = extract_table(sh)
                lines.append(f"- 표({len(rows)}행): " + json.dumps(rows, ensure_ascii=False))
            if getattr(sh, "has_chart", False) and sh.has_chart:
                data = extract_chart(sh)
                if data and "error" not in data:
                    cid = len(charts)
                    charts.append({"file": Path(path).name, "slide": si, **data})
                    lines.append(f"- 차트 #{cid}: 수치 추출 성공 → charts.json")
                else:
                    lines.append(
                        f"- 차트: 수치 추출 실패({(data or {}).get('error', '?')}) — 텍스트/이미지 폴백 필요"
                    )
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img = sh.image
                    n_img += 1
                    fn = f"{stem}_s{si}_{n_img}.{img.ext}"
                    (img_dir / fn).write_bytes(img.blob)
                    lines.append(
                        f"- 이미지 폴백: images/{fn} (차트/인포그래픽이면 수치 복원 시도할 것)"
                    )
                except Exception:
                    lines.append("- 이미지: blob 추출 실패(외부 링크 그림일 수 있음)")
        for t in texts:
            lines.append(f"- 텍스트: {t}")
        lines.append("")
    return len(prs.slides)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("out_dir")
    ap.add_argument("inputs", nargs="+")
    args = ap.parse_args()

    out = Path(args.out_dir)
    img_dir = out / "images"
    img_dir.mkdir(parents=True, exist_ok=True)

    charts, lines = [], ["# NotebookLM pptx 추출 결과", ""]
    total = 0
    for p in args.inputs:
        total += extract_file(p, img_dir, charts, lines)

    (out / "extracted.md").write_text("\n".join(lines), encoding="utf-8")
    (out / "charts.json").write_text(
        json.dumps(charts, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    print(
        f"OK: {len(args.inputs)}개 파일 {total}슬라이드 → {out}/extracted.md, "
        f"charts.json({len(charts)}건), images/"
    )


if __name__ == "__main__":
    sys.exit(main())
