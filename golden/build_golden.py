"""골든 덱 빌더 — 견본 17장을 부 단위로 누적 구현한다. 현재: S1 표지, S2 목차, S3 간지, S4 문제정의(v16 확정 — _variant_k).

콘텐츠 단일 출처: golden/00_factsheet.md (수치 창작 금지)
디자인 모방 대상: golden/ref/s01_cover.png(BCG NYCHA p1), s02_toc.png(우석진 p2)
실행: uv run python golden/build_golden.py → golden/golden-deck.pptx
"""

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches

from kit import SLIDE_H, SLIDE_W, add_box, add_text, load_kit, mix, set_shape_text

K = load_kit()
C, S, F = K["rgb"], K["sizes"], K["fonts"]
MARGIN = K["layout"]["margin"]

TITLE = "K-IFRS 1115 온톨로지 지식그래프 RAG"
VALUE_PROP = "확실할 때만 확정하고, 애매하면 근거를 보여주며 유보한다 — 2종 오류를 구조로 차단하는 수익인식 특화 챗봇"
KICKER = "TECHNICAL PROPOSAL · K-IFRS 1115"
YEAR = "2026"

PART_LEADS = [  # 간지 리드 — 00_factsheet.md(실데이터 v2) 부별 요약
    "일반 LLM의 허위 확정 리스크와, 초기 임베딩 RAG 운영에서 실측된 4가지 균열",
    "질문에서 답까지 4개 노드 전부가 결정적으로 동작합니다 — 진입부 임베딩 유사도 0",
    "기준서의 명시적 관계를 그대로 옮긴 지식그래프와, 회계 항등식으로 전수 검증한 데이터입니다",
    "실측된 네 균열마다 구조적 대체물을 배치했습니다 — 확률 신호의 전량 폐기",
    "점수 경쟁이 사라진 결정적 경로, AI 개입을 색인 하나로 제한한 신빙성 설계입니다",
    "정답 문서를 차단한 채 사람 전문가의 결론을 78/92 재현했습니다 — 수치가 실제 능력을 반영합니다",
]

TOC = [  # (제목, 서브카피 ' · ' 구분, 시작 페이지) — 00_factsheet.md §A~§H, 페이지는 16장 구성 파생
    ("문제 정의", "치명적 오류는 2종(허위 확정) · 일반 LLM이 실패하는 4가지 이유", "03"),
    ("파이프라인", "Analyze → Retrieve → Generate → Format · 4개 결정 노드", "05"),
    ("기술 설명", "용어사전 · 지식그래프 · 판단트리 · 구조화 출력 — 노드 순서대로 4개 기술", "07"),
    ("문제와 해결", "확률 신호의 실패 실측(탈락 103/105) · 전·후 5개 지표 재현 72→78", "13"),
    ("차별점", "일반 임베딩 RAG 대비 7개 축 비교", "16"),
    ("성과·증거", "홀드아웃 78/92 · 에러 0건 · 미재현 원인 분해", "19"),
]


# ── S1 전용 GRID (사용자 확정 2026-07-10 — 임의 좌표 금지) ──
COVER_L = 0.65  # 좌측 기준선 — 좌측정렬 요소 전부 이 축
COVER_R = SLIDE_W - COVER_L  # 12.683
COVER_RULE_Y = 0.92
COVER_ACCENT_TOP = 2.75  # 오렌지 사각형 0.3×0.3 — 타이틀 블록(사각형·제목·부제)은 2.75~4.6 응집
COVER_TITLE_TOP = 3.15  # 사각형 bottom(3.05) + 0.1
COVER_SUB_TOP = 4.25  # 제목 bottom + 0.35 — 블록 세로중심 ≈ 3.68 (광학 중앙보다 살짝 아래)
COVER_YEAR_TOP = 6.35  # 연도는 푸터 성격 — 블록에서 분리, 하단 앵커
COVER_BAND_TOP = 6.95  # 풀폭 밴드. YEAR bottom < 6.95 - 0.3 유지


def build_s01_cover(prs):
    """표지 — NYCHA p1 구도. 좌측정렬 통일(택1=A): 워드마크도 left=COVER_L."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rule_w = COVER_R - COVER_L  # 12.033
    add_text(slide, COVER_L, 0.55, 9.0, 0.3, KICKER, S["caption"], F["head"], C["muted"], bold=True)
    add_box(slide, COVER_L, COVER_RULE_Y, rule_w, 0.014, fill=C["muted"])
    # 타이틀 블록 — 사각형·제목·부제·연도 left 편차 0
    add_box(slide, COVER_L, COVER_ACCENT_TOP, 0.3, 0.3, fill=C["accent"])
    add_text(
        slide,
        COVER_L,
        COVER_TITLE_TOP,
        rule_w,
        0.75,
        TITLE,
        S["display"],
        F["head"],
        C["primary"],
        bold=True,
    )
    add_text(slide, COVER_L, COVER_SUB_TOP, 11.0, 0.35, VALUE_PROP, S["sub"], F["body"], C["text"])
    add_text(slide, COVER_L, COVER_YEAR_TOP, 4.0, 0.28, YEAR, S["head"], F["body"], C["muted"])
    # 하단 다크 밴드 + 좌측정렬 워드마크 (밴드 내 수직중앙)
    band_h = SLIDE_H - COVER_BAND_TOP
    add_box(slide, 0, COVER_BAND_TOP, SLIDE_W, band_h, fill=C["primary"])
    add_text(
        slide,
        COVER_L,
        COVER_BAND_TOP + band_h / 2 - 0.15,
        6.0,
        0.3,
        K["brand"]["name"],
        S["head"],
        F["head"],
        C["bg"],
        bold=True,
        anchor=MSO_ANCHOR.MIDDLE,
    )


# ── S2 목차 GRID (사용자 확정 2026-07-10, vC5 — 타이포 리스트+다크 패널+로고 슬롯) ──
TOC_TITLE_TOP, TOC_RULE_Y = 0.5, 1.35
TOC_PANEL_X = 9.6  # 우측 다크 패널
TOC_LIST_R = TOC_PANEL_X - 0.5  # 리스트 우측 끝 9.1
TOC_NUM_X, TOC_TEXT_X = MARGIN, 1.35  # 번호 축 0.6 / 제목·서브 축 1.35
TOC_ROW1_Y, TOC_ROW_PITCH = 2.0, 0.82


def build_s02_toc(prs):
    """목차 — 타이포 리스트(번호·제목·서브·p.번호·헤어라인) + 다크 패널(로고 슬롯+덱 타이틀).

    장식 모티프 금지(그래프 성좌·거미줄 기각 이력) — 패널은 로고 슬롯과 덱 타이틀만.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    add_box(slide, TOC_PANEL_X, 0, SLIDE_W - TOC_PANEL_X, SLIDE_H, fill=C["primary"])
    ph_c = mix(C["primary"], C["muted"], 0.45)  # 로고 플레이스홀더 — 실제 로고로 교체하는 슬롯
    ph_x, ph_y, ph_w, ph_h = TOC_PANEL_X + 0.9, 0.9, (SLIDE_W - TOC_PANEL_X) - 1.8, 0.9
    add_box(slide, ph_x, ph_y, ph_w, ph_h, fill=None, line=ph_c, line_w=1.0)
    add_text(
        slide,
        ph_x,
        ph_y + ph_h / 2 - 0.15,
        ph_w,
        0.3,
        "회사 로고",
        S["caption"],
        F["body"],
        ph_c,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        TOC_PANEL_X + 0.35,
        SLIDE_H - 1.15,
        SLIDE_W - TOC_PANEL_X - 0.7,
        0.7,
        [[(TITLE, {"bold": True, "color": C["bg"]})], [("기술 제안서 · 2026", {})]],
        S["caption"],
        F["body"],
        C["bg_alt"],
        line_spacing=1.4,
    )
    add_text(
        slide,
        MARGIN,
        TOC_TITLE_TOP,
        4.0,
        0.55,
        "목차",
        S["title"],
        F["head"],
        C["primary"],
        bold=True,
    )
    add_box(slide, MARGIN, TOC_RULE_Y, TOC_LIST_R - MARGIN, 0.014, fill=C["muted"])
    for i, (title, subs, page) in enumerate(TOC):
        y = TOC_ROW1_Y + i * TOC_ROW_PITCH
        add_text(
            slide, TOC_NUM_X, y, 0.6, 0.3, f"0{i + 1}", S["head"], F["head"], C["accent"], bold=True
        )
        add_text(
            slide,
            TOC_TEXT_X,
            y - 0.02,
            6.0,
            0.32,
            title,
            S["head"],
            F["head"],
            C["primary"],
            bold=True,
        )
        add_text(
            slide,
            TOC_TEXT_X,
            y + 0.3,
            TOC_LIST_R - TOC_TEXT_X - 1.0,
            0.26,
            subs,
            S["caption"],
            F["body"],
            C["muted"],
        )
        add_text(
            slide,
            TOC_LIST_R - 0.9,
            y,
            0.9,
            0.3,
            f"p.{page}",
            S["body"],
            F["body"],
            C["muted"],
            align=PP_ALIGN.RIGHT,
        )
        if i < 5:
            add_box(
                slide,
                TOC_TEXT_X,
                y + TOC_ROW_PITCH - 0.17,
                TOC_LIST_R - TOC_TEXT_X,
                0.01,
                fill=C["bg_alt"],
            )


# ── S3 간지 GRID (사용자 확정 2026-07-10 — 좌측 축 0.65 한 축: 태그·제목·바·도트) ──
PART_L = 0.65
PART_TAG_TOP, PART_TAG_W, PART_TAG_H = 2.0, 1.3, 0.42
PART_TITLE_TOP = 2.55
PART_SUB_TOP = (
    4.4  # 액센트 바 top=서브 top, height=서브 블록 height (감싸기 — 텍스트보다 짧으면 위반)
)
PART_SUB_H = 0.35  # 리드 1줄 기준. 2줄이 되면 어절 경계에서만 끊고 바 높이도 함께 늘릴 것
PART_SUB_X = PART_L + 0.05 + 0.2  # 0.9 — 바 오른쪽 여백 0.2
PART_DOTS_TOP = 5.75
PART_DOT_D, PART_DOT_PITCH = 0.13, 0.35  # 첫 도트 left=0.65, 중심간 0.35 균등
PART_WM_RIGHT = 12.0  # 고스트 워터마크 right ≤ 12.0 (우측 잘림 방지)


def build_part(prs, no):
    """간지 — Dallas p34 구도. 태그·제목·바·도트 left=0.65 한 축. no: 1~6."""
    title = TOC[no - 1][0]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["primary"])
    # 고스트 대형 숫자 — 제목과 광학적 세로 정렬, right ≤ 12.0
    add_text(
        slide,
        PART_WM_RIGHT - 6.0,
        1.9,
        6.0,
        2.6,
        f"0{no}",
        S["display"] * 4,
        F["head"],
        mix(C["primary"], C["muted"], 0.22),
        bold=True,
        align=PP_ALIGN.RIGHT,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    pill = add_box(slide, PART_L, PART_TAG_TOP, PART_TAG_W, PART_TAG_H, fill=C["accent"])
    set_shape_text(pill, f"PART 0{no}", S["caption"], F["head"], C["bg"], bold=True)
    add_text(
        slide, PART_L, PART_TITLE_TOP, 7.5, 0.75, title, S["display"], F["head"], C["bg"], bold=True
    )
    add_box(slide, PART_L, PART_SUB_TOP, 0.05, PART_SUB_H, fill=C["accent"])
    add_text(
        slide,
        PART_SUB_X,
        PART_SUB_TOP,
        9.5,
        PART_SUB_H,
        PART_LEADS[no - 1],
        S["sub"],
        F["body"],
        C["bg_alt"],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    for i in range(6):  # 진행 도트 — 목차 6부 인디케이터, 현재 PART만 accent
        dot_fill = C["accent"] if i == no - 1 else C["muted"]
        add_box(
            slide,
            PART_L + i * PART_DOT_PITCH,
            PART_DOTS_TOP,
            PART_DOT_D,
            PART_DOT_D,
            fill=dot_fill,
            shape="oval",
        )


def content_header(slide, part_no, part_title, headline):
    """본문 장 공통 헤더 — 킥커 + 헤드라인 1줄 + 헤어라인 (BCG 문서형)."""
    add_text(
        slide,
        MARGIN,
        0.42,
        6.0,
        0.28,
        f"{part_no}. {part_title}",
        S["caption"],
        F["head"],
        C["muted"],
        bold=True,
    )
    add_text(
        slide,
        MARGIN,
        0.72,
        SLIDE_W - 2 * MARGIN,
        0.55,
        headline,
        S["section"],
        F["head"],
        C["primary"],
        bold=True,
    )
    add_box(slide, MARGIN, 1.55, SLIDE_W - 2 * MARGIN, 0.014, fill=C["muted"])


def source_line(slide, text):
    """하단 출처선 — 골든 문법 고정 요소."""
    add_box(slide, MARGIN, SLIDE_H - 0.62, SLIDE_W - 2 * MARGIN, 0.01, fill=C["bg_alt"])
    add_text(
        slide,
        MARGIN,
        SLIDE_H - 0.55,
        SLIDE_W - 2 * MARGIN,
        0.3,
        text,
        S["foot"],
        F["body"],
        C["muted"],
    )


def pct_bar(slide, x, y, w, h, label, pct, pct_label, note):
    """100% 비례 바 — 회색 전체 = 100%, accent 구간 = pct(0~1). 폭은 비율에서 산출."""
    add_text(slide, x, y - 0.34, w, 0.3, label, S["body"], F["body"], C["primary"], bold=True)
    add_box(slide, x, y, w, h, fill=C["bg_alt"])
    seg_w = max(w * pct, 0.035)  # 비례상 불가시 구간은 최소 가시폭으로 (주석으로 정직 표기)
    add_box(slide, x, y, seg_w, h, fill=C["accent"])
    add_text(
        slide,
        x + w + 0.15,
        y + h / 2 - 0.22,
        1.6,
        0.44,
        pct_label,
        S["title"],
        F["head"],
        C["accent"],
        bold=True,
    )
    add_text(slide, x, y + h + 0.06, w + 1.6, 0.3, note, S["caption"], F["body"], C["muted"])


def main():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    build_s01_cover(prs)
    build_s02_toc(prs)
    build_part(prs, 1)
    from _variant_k import variant_k as build_s04

    build_s04(prs)
    build_part(prs, 2)
    from s06_variants import variant_c as build_s06  # S6 확정(vC9) — 분기 포함 실행 그래프

    build_s06(prs)
    build_part(prs, 3)
    from s08_variants import variant_c as build_s08  # S8 확정(C3) — 서사 3칼럼+실물 증거

    build_s08(prs)
    from s09_variants import variant_a as build_s09  # S9 확정(v8) — 위계 트리+구축 방법 표

    build_s09(prs)
    from s10_screenshot import variant_a as build_s10  # S10 확정(vA3) — 실물 화면+관찰 칼럼

    build_s10(prs)
    from s11_variants import (
        variant_d as build_s11,
    )  # S11 확정(vD8) — 진입 3단계+flat tree+적용 실물

    build_s11(prs)
    from s12_variants import (
        variant_b as build_s12,
    )  # S12 확정(vB3) — 실물 답변 발췌+스키마+검증 루프

    build_s12(prs)
    build_part(prs, 4)
    from s14_variants import variant_c as build_s14  # S14 확정(vC3) — 동일 질문 A/B 시뮬레이션

    build_s14(prs)
    from s15_variants import variant_c as build_s15  # S15 확정(vC3) — 문서 분해도+기록+미재현 해부

    build_s15(prs)
    build_part(prs, 5)
    from s17_variants import variant_c as build_s17  # S17 확정(vC) — 풀폭 미러 매트릭스

    build_s17(prs)
    from s18_variants import variant_b as build_s18  # S18 확정(vB) — 경계 도해+잔여 한계 3섹션

    build_s18(prs)
    from s21_closing import variant_a as build_closing  # 클로징 확정(vC) — 여백형 수미상관

    build_closing(prs)
    from pathlib import Path

    out = Path(__file__).parent / "golden-deck.pptx"
    try:
        prs.save(str(out))
    except PermissionError:  # PowerPoint가 열어둔 경우 — 대체 파일로 저장하고 계속
        out = Path(__file__).parent / "golden-deck.build.pptx"
        prs.save(str(out))
    print(f"saved {out} ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
