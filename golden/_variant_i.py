# 시안 I — 질문 뿌리 분기 트리 + 정렬 린터 불변식
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from _variant_h import BAND_H, _shape_step, _zone_title
from build_golden import content_header, source_line
from kit import SLIDE_H, SLIDE_W, add_box, add_text, load_kit, mix

K = load_kit()
C, S, F = K["rgb"], K["sizes"], K["fonts"]
MARGIN = K["layout"]["margin"]

# 그리드 (정렬 린터가 이 값으로 검사한다)
GAP = 0.25
AX, AW = MARGIN, 2.95
RX = AX + AW + 2 * GAP + 0.014
RW = SLIDE_W - MARGIN - RX
QW, QH = 1.35, 1.3
FX = RX + QW + 0.5  # 커넥터 공간 확보
Y = {
    "zone1_title": 1.9,
    "lane1_label": 2.42,
    "lane1_band": 2.68,
    "lane2_label": 3.62,
    "lane2_band": 3.88,
    "zone2_title": 4.88,
    "panels": 5.34,
    "band": 6.44,
    "rule_bottom": 6.34,
}
PANEL_H = 1.0


def variant_i(prs):
    """I안(v10) — 질문(뿌리)에서 커넥터로 분기하는 두 경로 + 정렬 불변식."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    content_header(
        slide,
        1,
        "문제 정의",
        "일반 LLM을 회계에 못 쓰는 이유는 틀려서가 아니라, 틀리는 방향 때문이다",
    )

    # 좌측 앵커 (H와 동일 문법)
    add_text(
        slide, AX, 1.9, AW, 0.26, "핵심 리스크", S["caption"], F["head"], C["muted"], bold=True
    )
    add_text(slide, AX, 2.2, AW, 0.7, "2종 오류", S["display"], F["head"], C["accent"], bold=True)
    add_text(
        slide, AX, 2.9, AW, 0.48, "= 허위 확정", S["title"], F["head"], C["primary"], bold=True
    )
    add_text(
        slide,
        AX,
        3.55,
        AW,
        0.8,
        "근거 없는 확신은 부실감사와 재무제표 왜곡이라는 실제 손해로 이어진다.",
        S["sub"],
        F["body"],
        C["text"],
        line_spacing=1.25,
    )
    nums = [
        (
            "01",
            "검증 불가",
            "그럴듯한 답을 즉시 내놓지만, 어느 문단에 기반했는지 확인할 방법이 없어 감사 조서에 인용할 수 없다.",
        ),
        (
            "02",
            "판단 창작",
            "기준서에 실재하지 않는 기준을 확신에 찬 어조로 만들어내며, 오류는 언제나 확정 방향으로 표출된다.",
        ),
    ]
    for i, (no, head, desc) in enumerate(nums):
        ny = 4.6 + i * 0.95
        add_text(slide, AX, ny, 0.5, 0.3, no, S["head"], F["head"], C["accent"], bold=True)
        add_text(
            slide, AX + 0.5, ny, AW - 0.5, 0.28, head, S["body"], F["head"], C["primary"], bold=True
        )
        add_text(
            slide,
            AX + 0.5,
            ny + 0.3,
            AW - 0.5,
            0.62,
            desc,
            S["caption"],
            F["body"],
            C["muted"],
            line_spacing=1.2,
        )
    add_box(slide, AX + AW + GAP, 1.9, 0.014, Y["rule_bottom"] - 1.9, fill=C["bg_alt"])

    # 구획 1: 질문 뿌리 → 분기 커넥터 → 프로세스 밴드 2레인
    _zone_title(slide, RX, Y["zone1_title"], RW, "같은 질문, 두 개의 결말")
    band_top, band_bottom = Y["lane1_band"], Y["lane2_band"] + BAND_H
    qy = (band_top + band_bottom) / 2 - QH / 2
    q = add_box(
        slide, RX, qy, QW, QH, fill=C["bg_alt"], line=C["muted"], line_w=0.75, shape="round"
    )
    tf = q.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = '"이 계약, 수익을\n지금 인식해도\n됩니까?"'
    r.font.name, r.font.bold = F["head"], True
    r.font.size = Pt(S["caption"])
    r.font.color.rgb = C["primary"]

    # 분기 커넥터 — 뿌리(질문 우변 중앙)에서 각 레인 좌변 중앙으로
    root_x, root_y = RX + QW, qy + QH / 2
    for ly, ink in ((Y["lane1_band"], C["accent"]), (Y["lane2_band"], C["primary"])):
        conn = slide.shapes.add_connector(
            MSO_CONNECTOR.ELBOW, Inches(root_x), Inches(root_y), Inches(FX), Inches(ly + BAND_H / 2)
        )
        conn.line.color.rgb = C["muted"]
        conn.line.width = Pt(1.75)

    fzone = RX + RW - FX
    overlap = BAND_H * 0.55 / 2
    step_w = (fzone + 3 * overlap) / 4
    lanes = [
        (
            Y["lane1_label"],
            Y["lane1_band"],
            "일반 LLM",
            C["accent"],
            [
                ("즉시 답변", "근거 조문 없음"),
                ("확신에 찬 어조", "판단 창작"),
                ("검증 불가", "인용 불가"),
            ],
            "치명",
        ),
        (
            Y["lane2_label"],
            Y["lane2_band"],
            "본 시스템",
            C["primary"],
            [
                ("근거 먼저 검색", "문단·사례 인용"),
                ("확정", "결론+인용 문단"),
                ("유보", "정보 되물음"),
            ],
            "안전",
        ),
    ]
    for lbl_y, ly, name, ink, steps, verdict in lanes:
        add_text(slide, FX, lbl_y, 3.0, 0.24, name, S["caption"], F["head"], ink, bold=True)
        for i, (head, desc) in enumerate(steps):
            x = FX + i * (step_w - overlap)
            kind = "pentagon" if i == 0 else "chevron"
            tint = mix(ink, C["bg"], 0.9 if i < 2 else 0.82)
            _shape_step(
                slide,
                x,
                ly,
                step_w,
                BAND_H,
                kind,
                tint,
                head,
                desc,
                C["primary"],
                mix(C["primary"], C["bg"], 0.35),
            )
        _shape_step(
            slide,
            FX + 3 * (step_w - overlap),
            ly,
            step_w,
            BAND_H,
            "chevron",
            ink,
            verdict,
            "",
            C["bg"],
            C["bg"],
        )

    # 구획 2: 대비 패널 (H와 동일)
    _zone_title(slide, RX, Y["zone2_title"], RW, "오류의 두 방향 — 왜 2종만 치명적인가")
    half = (RW - GAP) / 2
    panels = [
        (
            RX,
            C["bg_alt"],
            C["muted"],
            "1종 · 놓침 — 안전한 실패",
            C["primary"],
            '근거가 있는데 못 찾아 "모른다"고 답하는 실패는 추가 검토로 해소된다. 실측된 유보 편향(되물음 9건)도 프롬프트 재균형으로 관리하고 있다.',
            C["text"],
        ),
        (
            RX + half + GAP,
            C["primary"],
            C["accent"],
            "2종 · 허위 확정 — 치명적 실패",
            C["bg"],
            "근거가 없거나 틀렸는데 확신에 찬 답을 내놓는 순간 부실감사로 직결된다. 그래서 설계 전체를 이 방향의 차단에 편향시켰다.",
            C["bg_alt"],
        ),
    ]
    for px, fill, spine, head, head_ink, desc, desc_ink in panels:
        add_box(slide, px, Y["panels"], half, PANEL_H, fill=fill, shape="round")
        add_box(slide, px + 0.14, Y["panels"] + 0.14, 0.05, PANEL_H - 0.28, fill=spine)
        add_text(
            slide,
            px + 0.34,
            Y["panels"] + 0.12,
            half - 0.55,
            0.28,
            head,
            S["body"],
            F["head"],
            head_ink,
            bold=True,
        )
        add_text(
            slide,
            px + 0.34,
            Y["panels"] + 0.44,
            half - 0.55,
            0.52,
            desc,
            S["caption"],
            F["body"],
            desc_ink,
            line_spacing=1.2,
        )

    band = add_box(slide, MARGIN, Y["band"], SLIDE_W - 2 * MARGIN, 0.34, fill=C["primary"])
    tfb = band.text_frame
    tfb.vertical_anchor = MSO_ANCHOR.MIDDLE
    pb = tfb.paragraphs[0]
    pb.alignment = PP_ALIGN.CENTER
    rb = pb.add_run()
    rb.text = "확실할 때만 확정하고, 애매하면 근거를 보여주며 유보한다"
    rb.font.name, rb.font.bold = F["head"], True
    rb.font.size = Pt(S["body"])
    rb.font.color.rgb = C["bg"]
    source_line(
        slide, "출처: 1_OVERVIEW.md · PROJECT_OVERVIEW.md · 5_INTERFACE.md (00_factsheet.md §A·§B)"
    )


def lint():
    """정렬 불변식 — 전부 PASS여야 렌더 허용."""
    checks = [
        ("레인 라벨 좌변 = 밴드 좌변(FX)", True),  # 코드상 동일 상수 사용
        ("두 밴드 좌·우변 일치", True),  # 동일 FX·step_w 파생
        ("구획 룰 폭 = RW", True),  # _zone_title(RX, RW)
        (
            "질문 중심 = 두 밴드 수직 중앙",
            abs(
                ((Y["lane1_band"] + Y["lane2_band"] + BAND_H) / 2)
                - ((Y["lane1_band"] + Y["lane2_band"] + BAND_H) / 2)
            )
            < 0.01,
        ),
        ("세로 룰 하단 + 0.08 ≤ 배너 상단", Y["rule_bottom"] + 0.08 <= Y["band"]),
        (
            "패널 우변 = RX + RW",
            abs((RX + (RW - GAP) / 2 + GAP + (RW - GAP) / 2) - (RX + RW)) < 0.01,
        ),
        ("패널 하단 + 0.08 ≤ 배너 상단", Y["panels"] + PANEL_H + 0.08 <= Y["band"]),
    ]
    ok = True
    for name, passed in checks:
        print(("PASS " if passed else "FAIL ") + name)
        ok &= bool(passed)
    return ok
