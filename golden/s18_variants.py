"""S18 시안 — 차별점: 정직한 한계 (물성 선언: 한계 = 경계 — 안은 결정 동작, 밖은 차단·유보).

콘텐츠 실물: 6_TEST §6.3(타 기준서 4건 IAS38·1002·1008·1037, 진입 누락 2건),
4_SEARCH L49(타 기준서 용어 감지 시 코드 레벨 강제 OUT — 예: 1116호 증분차입이자율),
00_factsheet.md §D·§H(경계 안 실측 스탯).
실행: uv run python golden/s18_variants.py → golden/variants/s18_variants.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

import grid as G
from kit import SLIDE_H, SLIDE_W, add_box, add_text, load_kit, mix

K = load_kit()
C, S, F = K["rgb"], K["sizes"], K["fonts"]

KICKER = "5. 차별점"
SOURCE = "출처: 6_TEST-DECISIONS.md §6.3 · 4_SEARCH-PIPELINE.md (라우팅 강제 OUT) · 00_factsheet.md §D·§H"


def header(slide, headline):
    add_text(
        slide, G.MARGIN_L, 0.42, 8.0, 0.28, KICKER, S["caption"], F["head"], C["muted"], bold=True
    )
    add_text(
        slide,
        G.MARGIN_L,
        0.72,
        G.RIGHT_EDGE - G.MARGIN_L,
        0.55,
        headline,
        S["section"],
        F["head"],
        C["primary"],
        bold=True,
    )
    add_box(slide, G.MARGIN_L, G.RULE_Y, G.RIGHT_EDGE - G.MARGIN_L, 0.014, fill=C["muted"])


def bar_and_source(slide, text):
    bar = add_box(slide, G.MARGIN_L, G.BAR_Y, G.RIGHT_EDGE - G.MARGIN_L, G.BAR_H, fill=C["primary"])
    tf = bar.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name, run.font.bold = F["head"], True
    run.font.size = Pt(S["body"])
    run.font.color.rgb = C["bg"]
    add_text(
        slide,
        G.MARGIN_L,
        G.SOURCE_Y,
        G.RIGHT_EDGE - G.MARGIN_L,
        0.3,
        SOURCE,
        S["foot"],
        F["body"],
        C["muted"],
    )


def _dash_arrow(slide, x1, y1, x2, y2):
    from pptx.oxml.ns import qn

    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = C["muted"]
    conn.line.width = Pt(1.5)
    ln = conn.line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "sm", "len": "sm"}))
    return conn


def _outside_chip(slide, x, y, w, head, sub):
    from pptx.enum.lang import MSO_LANGUAGE_ID

    chip = add_box(slide, x, y, w, 0.8, fill=C["bg"], line=C["muted"], line_w=1.0, shape="round")
    tf = chip.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.12)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    p1._p.get_or_add_pPr().set("eaLnBrk", "0")
    r1 = p1.add_run()
    r1.text = head
    r1.font.name, r1.font.bold = F["head"], True
    r1.font.size = Pt(S["caption"])
    r1.font.color.rgb = C["primary"]
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2._p.get_or_add_pPr().set("eaLnBrk", "0")
    r2 = p2.add_run()
    r2.text = sub
    r2.font.name = F["body"]
    r2.font.size = Pt(S["caption"])
    r2.font.color.rgb = C["muted"]
    for r in (r1, r2):
        r.font.language_id = MSO_LANGUAGE_ID.KOREAN
    return chip


def variant_a(prs):
    """A — 경계 도해: 안(1115 결정 동작·실측 스탯) / 밖(타 기준서 강제 OUT ✕ · 미등재 유보)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    header(slide, "정직한 한계 — 이 시스템이 서 있는 경계")
    from pptx.enum.lang import MSO_LANGUAGE_ID

    # 바깥 세계 (bg_alt 캔버스)
    add_box(slide, G.MARGIN_L, 1.95, G.RIGHT_EDGE - G.MARGIN_L, 4.05, fill=C["bg_alt"])
    add_text(
        slide,
        G.MARGIN_L + 0.2,
        2.05,
        3.0,
        0.24,
        "경계 밖",
        S["caption"],
        F["head"],
        C["muted"],
        bold=True,
    )
    # 안쪽 영역 — 결정적으로 동작하는 경계
    in_x, in_y, in_w, in_h = 4.35, 2.3, 4.6, 3.35
    region = add_box(
        slide, in_x, in_y, in_w, in_h, fill=C["bg"], line=C["primary"], line_w=2.5, shape="round"
    )
    region.adjustments[0] = 0.05
    add_text(
        slide,
        in_x,
        in_y + 0.22,
        in_w,
        0.32,
        "K-IFRS 1115 — 경계 안",
        S["head"],
        F["head"],
        C["primary"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        in_x,
        in_y + 0.56,
        in_w,
        0.26,
        "이 안에서는 결정적으로 동작한다",
        S["caption"],
        F["body"],
        C["muted"],
        align=PP_ALIGN.CENTER,
    )
    stats = [("노드", "929"), ("간선", "2,697"), ("판단트리", "41"), ("등재 용어", "423")]
    for i, (k, v) in enumerate(stats):
        col, row = divmod(i, 2)
        sx = in_x + 0.55 + col * 2.15
        sy = in_y + 1.0 + row * 0.8
        add_text(slide, sx, sy, 1.9, 0.38, v, S["title"], F["head"], C["primary"], bold=True)
        add_text(slide, sx, sy + 0.46, 1.9, 0.2, k, S["caption"], F["body"], C["muted"])
    add_box(slide, in_x + 0.5, in_y + 2.62, in_w - 1.0, 0.012, fill=C["bg_alt"])
    add_text(
        slide,
        in_x,
        in_y + 2.76,
        in_w,
        0.3,
        [
            [
                ("홀드아웃 ", {}),
                ("78 / 92", {"bold": True, "color": C["accent"]}),
                (" — 실행 에러 0건", {}),
            ]
        ],
        S["body"],
        F["body"],
        C["muted"],
        align=PP_ALIGN.CENTER,
    )
    # 밖 좌: 타 기준서 질문 — 강제 OUT
    _outside_chip(slide, 0.85, 3.15, 2.6, "타 기준서 질문", "IAS38 · 1002 · 1008 · 1037")
    _dash_arrow(slide, 3.45, 3.55, 4.33, 3.55)
    add_box(slide, 4.3, 3.33, 0.05, 0.44, fill=C["primary"])  # 차단 바
    add_text(slide, 3.98, 3.62, 0.3, 0.26, "✕", S["body"], F["head"], C["primary"], bold=True)
    add_text(
        slide,
        0.85,
        4.05,
        2.9,
        0.55,
        "라우팅이 코드 레벨 강제 OUT — 추측 답변 대신 거절. 실측 4건 (예: 1116호 '증분차입이자율' 감지)",
        S["caption"],
        F["body"],
        C["muted"],
        line_spacing=1.2,
    )
    # 밖 우: 등재되지 않은 표현 — 유보
    _outside_chip(slide, 9.85, 3.15, 2.6, "등재되지 않은 표현", "임베딩식 유사 확장 없음")
    _dash_arrow(slide, 9.85, 3.55, 8.97, 3.55)
    add_box(slide, 8.95, 3.33, 0.05, 0.44, fill=C["primary"])
    add_text(slide, 9.32, 3.62, 0.3, 0.26, "✕", S["body"], F["head"], C["primary"], bold=True)
    add_text(
        slide,
        9.55,
        4.05,
        2.9,
        0.55,
        "색인에 없으면 진입 실패 → '못 찾음'으로 유보 응답. 홀드아웃 실측 진입 누락 2건",
        S["caption"],
        F["body"],
        C["muted"],
        line_spacing=1.2,
    )
    add_text(
        slide,
        G.MARGIN_L,
        6.08,
        G.RIGHT_EDGE - G.MARGIN_L,
        0.24,
        "경계의 비용은 실측돼 있다 — 미재현 14건 중 타 기준서 4건·진입 누락 2건이 여기서 나온다. 나머지는 경계 안의 문제가 아니었다.",
        S["caption"],
        F["body"],
        C["muted"],
    )
    bar_and_source(
        slide, "경계를 넓히는 대신 경계 안을 결정적으로 — 밖의 질문은 추측하지 않고 거절한다"
    )
    return slide


OTHER_LIMITS = [  # 6_TEST §6.2·§6.3·L52 실물 — 경계 밖 잔여 한계 3종
    (
        "결론을 확정하지 못한 케이스",
        "헤지 2건 — 근거 문단을 다 찾고도 결론을 유보했다. 검색이 아니라 생성 계층의 남은 과제.",
    ),
    (
        "인용 완전성",
        "하드 인용 재현율 59.1% — 결론이 맞아도 근거 인용을 전부 담지는 못한다. 별도 지표로 추적.",
    ),
    (
        "검증 자체의 한계",
        "92건은 개발 중 전수 열람돼 순수 홀드아웃이 아니다 — 최종 성능은 안 본 질문으로 검증해야 한다.",
    ),
]


def variant_b(prs):
    """B — 상단 경계 도해(안=정책, 밖=차단·유보) + 하단 잔여 한계 3섹션."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    header(slide, "정직한 한계 — 이 시스템이 서 있는 경계")
    from pptx.enum.lang import MSO_LANGUAGE_ID

    # ── 상단: 경계 캔버스 ──
    add_box(slide, G.MARGIN_L, 1.95, G.RIGHT_EDGE - G.MARGIN_L, 2.6, fill=C["bg_alt"])
    add_text(
        slide,
        G.MARGIN_L + 0.2,
        2.05,
        3.0,
        0.24,
        "경계 밖",
        S["caption"],
        F["head"],
        C["muted"],
        bold=True,
    )
    in_x, in_y, in_w, in_h = 4.6, 2.15, 4.1, 2.2
    region = add_box(
        slide, in_x, in_y, in_w, in_h, fill=C["bg"], line=C["primary"], line_w=2.5, shape="round"
    )
    region.adjustments[0] = 0.07
    add_text(
        slide,
        in_x,
        in_y + 0.2,
        in_w,
        0.32,
        "K-IFRS 1115 — 경계 안",
        S["head"],
        F["head"],
        C["primary"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        in_x + 0.35,
        in_y + 0.62,
        in_w - 0.7,
        1.4,
        [
            [("답하는 범위를 계약한다", {"bold": True, "color": C["primary"]})],
            [
                (
                    "경계 안 질문에만 결정적으로 답하고, 밖이면 거절, 모르면 유보 — "
                    "아는 척이 구조적으로 불가능하게 만든 설계다.",
                    {},
                )
            ],
        ],
        S["caption"],
        F["body"],
        C["muted"],
        line_spacing=1.35,
    )
    # 밖 좌: 타 기준서 — 강제 OUT
    _outside_chip(slide, 0.85, 2.5, 2.5, "타 기준서 질문", "IAS38 · 1002 · 1008 · 1037")
    _dash_arrow(slide, 3.45, 2.9, 4.58, 2.9)
    add_box(slide, 4.55, 2.68, 0.05, 0.44, fill=C["primary"])
    add_text(slide, 4.18, 2.96, 0.3, 0.26, "✕", S["body"], F["head"], C["primary"], bold=True)
    add_text(
        slide,
        0.85,
        3.45,
        3.2,
        0.55,
        "라우팅이 코드 레벨 강제 OUT — 추측 대신 거절. 실측 4건 (예: 1116호 '증분차입이자율' 감지)",
        S["caption"],
        F["body"],
        C["muted"],
        line_spacing=1.2,
    )
    # 밖 우: 미등재 표현 — 유보
    _outside_chip(slide, 9.95, 2.5, 2.5, "등재되지 않은 표현", "임베딩식 유사 확장 없음")
    _dash_arrow(slide, 9.95, 2.9, 8.72, 2.9)
    add_box(slide, 8.7, 2.68, 0.05, 0.44, fill=C["primary"])
    add_text(slide, 9.07, 2.96, 0.3, 0.26, "✕", S["body"], F["head"], C["primary"], bold=True)
    add_text(
        slide,
        9.35,
        3.45,
        3.1,
        0.55,
        "색인에 없으면 진입 실패 → '못 찾음' 유보 응답. 홀드아웃 실측 진입 누락 2건",
        S["caption"],
        F["body"],
        C["muted"],
        line_spacing=1.2,
    )
    # ── 하단: 잔여 한계 3섹션 ──
    add_text(
        slide,
        G.MARGIN_L,
        4.78,
        G.RIGHT_EDGE - G.MARGIN_L,
        0.3,
        "경계 그 밖의 한계 — 실측으로 아는 것",
        S["head"],
        F["head"],
        C["primary"],
        bold=True,
    )
    add_box(slide, G.MARGIN_L, 5.12, G.RIGHT_EDGE - G.MARGIN_L, 0.012, fill=C["muted"])
    sec_w, gap = 3.86, 0.27
    for i, (head, body) in enumerate(OTHER_LIMITS):
        sx = G.MARGIN_L + i * (sec_w + gap)
        card = add_box(slide, sx, 5.28, sec_w, 1.06, fill=C["bg_alt"], shape="round")
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.14)
        tf.margin_top = Inches(0.09)
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.LEFT
        p1._p.get_or_add_pPr().set("eaLnBrk", "0")
        rn = p1.add_run()
        rn.text = f"0{i + 1}  "
        rn.font.name, rn.font.bold = F["head"], True
        rn.font.size = Pt(S["caption"])
        rn.font.color.rgb = C["muted"]
        rh = p1.add_run()
        rh.text = head
        rh.font.name, rh.font.bold = F["head"], True
        rh.font.size = Pt(S["caption"])
        rh.font.color.rgb = C["primary"]
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        p2._p.get_or_add_pPr().set("eaLnBrk", "0")
        p2.line_spacing = 1.25
        rb = p2.add_run()
        rb.text = body
        rb.font.name = F["body"]
        rb.font.size = Pt(S["caption"])
        rb.font.color.rgb = C["muted"]
        for r in (rn, rh, rb):
            r.font.language_id = MSO_LANGUAGE_ID.KOREAN
    bar_and_source(
        slide, "경계를 넓히는 대신 경계 안을 결정적으로 — 못 하는 것까지 실측으로 세어 두었다"
    )
    return slide


def audit(prs):
    from pptx.oxml.ns import qn

    fails = []
    for si, sl in enumerate(prs.slides):
        accents = 0
        for sh in sl.shapes:
            t, le = sh.top / 914400, sh.left / 914400
            b, r = t + sh.height / 914400, le + sh.width / 914400
            full_bleed = sh.width / 914400 >= SLIDE_W - 0.05
            if t < 6.4 and b > 6.37 and not full_bleed:
                fails.append((si, "bottom", round(b, 2), sh.shape_id))
            if r > 12.75 and not full_bleed:
                fails.append((si, "right", round(r, 2), sh.shape_id))
            for el in sh._element.iter(qn("a:srgbClr")):
                if el.get("val") == "D66E3A":
                    if el.getparent().tag.endswith("}solidFill"):
                        accents += 1
        if accents > 4:
            fails.append((si, "accent>4", accents))
    assert not fails, f"AUDIT FAIL {fails}"
    print("audit pass")


def main():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    variant_a(prs)
    variant_b(prs)
    audit(prs)
    out = Path(__file__).parent / "variants" / "s18_variants.pptx"
    prs.save(str(out))
    print(f"saved {out}")


if __name__ == "__main__":
    main()
