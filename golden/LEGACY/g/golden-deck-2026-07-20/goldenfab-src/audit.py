"""goldenfab.audit — 슬라이드 기계 오딧. **채점자가 캔 규칙을 코드로 제련해 두는 곳.**

## 왜 이 파일이 있나 (2026-07-15)

S4 한 장에 제3자 채점 2회 = **30분**이 들었다. 그런데 두 채점자가 낸 FAIL 9건을 분류하니
**7건이 순수 산수**였다 — 대비비·채움률·높이 종류·텍스트 중복·라벨 정렬·같은 fill 충돌·accent 계수.
15분짜리 모델이 필요한 일이 아니었다. 내 오딧이 그 규칙을 **안 갖고 있어서** 못 잡았을 뿐이다.

즉 **채점자는 검사기가 아니라 광산이다.** 한 번 캔 광석을 코드로 내려두면 다시 캘 필요가 없다.
design-rules의 철학("확정할 때마다 판단을 박제")과 같되, **산문이 아니라 실행되는 코드로** 박제한다
— 산문 규칙은 잊히지만 assert는 건너뛸 수 없다(P1 주석의 자기 인정).

## 운영 규칙

1. 레이아웃을 만들면 먼저 이 오딧을 통과시킨다(2초).
2. 제3자 채점은 **1회만**. 재채점 금지 — 오늘 실증: 6/10 → 5/10으로 **안 수렴**했고 두 채점자가
   **서로 다른 항목**을 찾았다. 주관적 채점에는 고정점이 없다.
3. 채점자가 새 결함을 찾으면 **반드시 여기에 규칙으로 내린다**. 산문으로 남기면 다음 장에서 또 캔다.
4. 회를 거듭할수록 채점자가 찾을 게 줄어든다. 그게 이 파일이 하는 일이다.

## 규칙 출처

각 규칙에 발견 경위를 적는다 — 왜 있는지 모르는 규칙은 다음 사람이 지운다.
"""

import re

from pptx.enum.shapes import MSO_SHAPE

EMU = 914400
AIR_MIN = 0.12  # §6 공기 하한
ACCENT_MAX = 4  # §2 accent 상한
CONTENT_BOTTOM = 6.35  # P2③
RIGHT_EDGE = 12.733  # P2④
CONTRAST_MIN = 4.5  # WCAG AA (일반 텍스트)


# ── 픽셀·색 산수 ──────────────────────────────────────────────


def _lum(hexstr):
    """상대 휘도 (WCAG)."""

    def ch(v):
        v /= 255
        return v / 12.92 if v <= 0.03928 else ((v + 0.055) / 1.055) ** 2.4

    r, g, b = (int(hexstr[i : i + 2], 16) for i in (0, 2, 4))
    return 0.2126 * ch(r) + 0.7152 * ch(g) + 0.0722 * ch(b)


def contrast(fg, bg="FFFFFF"):
    """대비비. 판정 단어가 산문보다 흐린지 재는 데 쓴다."""
    a, b = _lum(str(fg)), _lum(str(bg))
    hi, lo = max(a, b), min(a, b)
    return round((hi + 0.05) / (lo + 0.05), 2)


def shape_kind(sh):
    try:
        return sh.auto_shape_type
    except Exception:
        return None


def fill_hex(sh):
    try:
        if sh.fill.type is not None:
            return str(sh.fill.fore_color.rgb)
    except Exception:
        pass
    return None


def line_hex(sh):
    try:
        if sh.line.color.type is not None:
            return str(sh.line.color.rgb)
    except Exception:
        pass
    return None


def runs_of(sh):
    """(텍스트, hex색, pt, bold) 목록."""
    out = []
    if not sh.has_text_frame:
        return out
    for p in sh.text_frame.paragraphs:
        for r in p.runs:
            try:
                col = str(r.font.color.rgb)
            except Exception:
                col = None
            out.append((r.text, col, r.font.size.pt if r.font.size else None, r.font.bold))
    return out


def box(sh):
    return (
        sh.left / EMU,
        sh.top / EMU,
        sh.width / EMU,
        sh.height / EMU,
    )


# ── 규칙 ──────────────────────────────────────────────────────


def check_accent(shapes, accent, cap=ACCENT_MAX):
    """§2 accent 상한. **채움만 세면 안 된다 — 테두리·커넥터 선도 accent 예산이다.**

    출처: v4에서 채움·런만 세서 4로 통과시켰는데 제3자가 7군데를 셌다(선 미계수).
    """
    f = [s for s in shapes if fill_hex(s) == accent]
    ln = [s for s in shapes if line_hex(s) == accent]
    rn = [r for s in shapes for r in runs_of(s) if r[1] == accent]
    tot = len(f) + len(ln) + len(rn)
    ok = tot <= cap
    return ok, f"accent 채움 {len(f)} + 선 {len(ln)} + 런 {len(rn)} = {tot} (상한 {cap})", tot


def check_verdict_contrast(shapes, prose_pt):
    """**판정을 말하는 단어는 그것을 설명하는 산문보다 흐릴 수 없다.**

    출처: v5에서 결론 단어 '↩ 되돌아옴'(3.22) · '✕ 불가역'(3.41)이 그것을 설명하는
    산문(8.47)보다 2.5배 흐렸다. 단순 명찰은 16.16을 가져갔다. 제3자 2회 모두 지적.
    """
    bad = []
    for s in shapes:
        for text, col, pt, _bold in runs_of(s):
            t = text.strip()
            if not col or not t:
                continue
            # 판정 단어는 **짧은 라벨**이다. 기호로 시작하거나(✕ 불가역 · ↩ 되돌아옴), 짧으면서
            # 판정어로 끝난다. 길이 상한이 없으면 산문 속 "없다"까지 걸린다 — 2026-07-15
            # tech_tree 실측: '문단 뭉치에는 순서도 관계도 없다'(17자, 서사 본문)를 판정으로
            # 오탐했다. 서사는 muted가 정상이다(§4 위계 4단). 규칙을 느슨하게 한 게 아니라
            # **판정과 산문을 구분**하도록 정확하게 했다 — 기호 시작 케이스는 그대로 잡힌다.
            is_verdict = bool(re.match(r"^\s*[✕✓↩→]", t)) or (
                len(t) <= 12 and ("불가역" in t or t.endswith("없다"))
            )
            if is_verdict:
                cr = contrast(col)
                if cr < CONTRAST_MIN:
                    bad.append((t[:18], col, cr))
    return not bad, f"판정 단어 대비 < {CONTRAST_MIN}: {len(bad)} {bad[:3]}", len(bad)


def check_fill_ratio(shapes, min_ratio=0.30, min_w=1.5):
    """도형이 글자에 비해 과하게 넓지 않은가(죽은 회색).

    출처: v5의 한계 칩이 폭 3.30"에 채움률 14~22%였다 — 단어와 커넥터 사이 238px 공백.
    글자폭은 pt·글자수로 근사(정밀 측정 불가) — 한국어 1자 ≈ 0.8×pt.
    """
    bad = []
    for s in shapes:
        if not s.has_text_frame or fill_hex(s) is None:
            continue
        _x, _y, w, _h = box(s)
        if w < min_w:
            continue
        rs = runs_of(s)
        if not rs:
            continue
        txt = "".join(r[0] for r in rs)
        pt = max((r[2] or 9) for r in rs)
        est = len(txt.strip()) * pt * 0.8 / 72
        ratio = est / w
        if ratio < min_ratio:
            bad.append((txt.strip()[:14], round(w, 2), f"{ratio:.0%}"))
    return not bad, f"채움률 < {min_ratio:.0%}: {len(bad)} {bad[:3]}", len(bad)


def check_duplicate_nodes(shapes, allow=0):
    """**구별 장치 없는** 반복 = 재탕(P4③).

    출처: v4에서 '2종 · 허위 확정'이 두 밴드에 같은 글자·같은 fill·같은 테두리로 있었다.

    ⚠ 테두리(색·점선)가 다르면 재탕이 아니라 **대조**다 — v7의 두 ◇ '근거 충분?'은 실선(성립)
    vs 점선(성립 불가)으로 뒤가 앞을 부정한다. 제3자도 "반복이 아니라 대조"로 PASS를 줬는데
    초판 규칙이 (텍스트, 채움)만 봐서 오탐을 냈다. 구별 장치까지 서명에 넣는다.

    `allow`: 같은 노드의 재등장이 **그 자체로 메시지**인 레이아웃을 위한 기대값(기본 0).
    tech_tree 카탈로그가 그렇다 — `변동대가`가 hier·cp의 출발이자 term의 도착으로 3번 나오는
    건 "한 개념에 세 종류의 간선이 붙는다"를 보이는 것이다. 규칙이 잡으려는 건 **정보량 0인**
    반복이고 이건 정보량이 있다. 다만 예외를 여는 게 아니라 **수를 박는다** — 실측보다 늘면
    FAIL이라, 무심코 하나 더 그리면 잡힌다.
    """
    seen = set()
    dup = []
    for s in shapes:
        f = fill_hex(s)
        if not f or not s.has_text_frame:
            continue
        t = s.text_frame.text.strip()
        if len(t) < 3:
            continue
        try:
            dash = str(s.line.dash_style)
        except Exception:
            dash = None
        sig = (t, f, line_hex(s), dash)  # 구별 장치 = 테두리 색 + 점선 여부
        if sig in seen:
            dup.append(t[:20])
        seen.add(sig)
    return len(dup) <= allow, f"구별 장치 없는 반복: {len(dup)} (허용 {allow}) {dup[:3]}", len(dup)


def check_node_class(shapes, min_len=2):
    """같은 이름의 노드가 한 장 안에서 **다른 스타일**로 그려지면 FAIL (P4⑩의 쌍둥이).

    출처: 2026-07-15 제3자 채점이 s9에 2/10을 줬고 최대 결함이 이것이었다 — 좌 트리는
    `bg_alt` 채움을 **내부 노드**(개념·문단)에, 우 카탈로그는 같은 `bg_alt`를 **사례**에 썼다.
    독자가 왼쪽에서 배운 "회색 = 내부 노드"가 오른쪽에서 "회색 = 사례"로 뒤집힌다.
    `변동대가`는 한 장에서 회색 덩어리이자 흰 바탕 검정 테두리 볼드였다.

    ⚠ `check_duplicate_nodes`의 **정확한 반대**다. 저건 "같은 이름 + 같은 서명"의 무의미한
    반복을 잡고, 이건 "같은 이름 + **다른** 서명"의 모순을 잡는다. 둘 다 있어야 한다 —
    duplicate만 있으면 "스타일을 바꾸면 통과"라는 잘못된 탈출구가 열린다.

    ⚠ **채움만 본다.** 초판이 (채움, 테두리)를 서명으로 삼았더니 S4의 두 ◇ `근거 충분?`을
    오탐했다 — 그건 실선(성립) vs 점선(성립 불가)으로 **뒤가 앞을 부정하는 의도된 대조**이고,
    제3자도 "반복이 아니라 대조"로 PASS를 준 설계다. 그대로 뒀으면 duplicate 규칙("구별 장치를
    둬라")과 이 규칙("같은 스타일이어야")이 **서로 반대를 요구**해 S4가 어느 쪽으로도 못 빠져
    나갔다. 경계는 이렇다:
        **채움 = 클래스**(이게 무엇인가 — 개념·문단·사례) → 같은 이름이면 같아야 한다.
        **테두리·점선 = 상태**(성립/불성립·자동/위임) → 달라도 된다, 그게 대조다.
    s9의 실제 결함은 채움이 갈린 것(`bg_alt` 좌 = 내부 노드 / 우 = 사례)이라 이 경계로 잡힌다.

    원인은 늘 구조적이다: 노드를 그리는 함수가 둘이면 색이 갈린다. 단일 출처(NODE_STYLE)로
    합치는 게 근본 수정이고, 이 규칙은 그게 다시 갈라지는지 지킨다.
    """
    sigs = {}
    bad = []
    for s in shapes:
        f = fill_hex(s)
        if not f or not s.has_text_frame:
            continue
        t = s.text_frame.text.strip()
        if len(t) < min_len:
            continue
        if t in sigs and sigs[t] != f:
            bad.append((t[:18], sigs[t], f))
        else:
            sigs.setdefault(t, f)
    return not bad, f"같은 노드 다른 채움: {len(bad)} {[b[0] for b in bad][:3]}", len(bad)


def check_ink_collision(shapes, ink, allow=1):
    """같은 채움색이 서로 다른 뜻을 지지 않는가(P4⑩) — 근사: 그 색 채움 도형 수.

    출처: v4에서 '재무제표 반영'(파국)과 결론 바(우리 원칙)가 **바이트 동일**한
    primary+샤프+흰볼드였다. 훑는 사람은 검은 덩어리 둘을 같은 종류로 분류한다.
    """
    f = [s for s in shapes if fill_hex(s) == ink]
    txt = [s.text_frame.text[:16] for s in f if s.has_text_frame and s.text_frame.text.strip()]
    return len(f) <= allow, f"{ink} 채움 {len(f)} (허용 {allow}) {txt}", len(f)


def check_node_heights(shapes, top_max, left_min=1.0, h_min=0.3):
    """같은 계열(행 노드) 높이 균일(P4⑧).

    출처: v4에서 확정 61 / 2종 69 / 일반LLM 71 / band2 84px — 한 종류에 높이 4종.
    계열 제외: ◇(다른 도형) · 풀폭 바 · 뿌리 카드(흐름 시작점) · 작은 칩.
    """
    nodes = [
        s
        for s in shapes
        if fill_hex(s)
        and shape_kind(s) != MSO_SHAPE.DIAMOND
        and box(s)[2] < 11
        and box(s)[0] > left_min
        and box(s)[1] < top_max
        and box(s)[3] > h_min
    ]
    hs = sorted({round(box(s)[3], 3) for s in nodes})
    return len(hs) <= 1, f"행 노드 {len(nodes)}개 · 높이 종류 {hs}", len(hs)


def check_progress_shapes(shapes, allowed=0):
    """§5 셰브런·펜타곤은 **분기 없는 단순 진행**에만.

    출처: XOR을 셰브런 밴드로 그려 "확정한 다음 유보한다"는 거짓 인과를 만든 반려 4회.
    """
    prog = [s for s in shapes if shape_kind(s) in (MSO_SHAPE.CHEVRON, MSO_SHAPE.PENTAGON)]
    return len(prog) <= allowed, f"진행형 도형 {len(prog)} (허용 {allowed})", len(prog)


def check_picture_overlap(shapes):
    """그림이 다른 요소를 침범하지 않는가.

    출처(2026-07-15 사용자 버그 리포트): s10이 이미지 비율을 `3200/2000`(1.60)으로 **가정**하고
    폭 7.6"를 전제로 옆 칼럼을 8.6"에 놨는데, 실제 파일은 1899×926(**2.05**)이라 폭 9.74"로
    렌더돼 텍스트를 **1.74" 침범**했다. 그 `img_w`는 계산만 하고 안 쓰는 죽은 변수였다.
    → `kit.fit_picture`가 근본 수정(박스에 맞춤). 이 규칙은 재발 감시.

    ⚠ python-pptx는 접근할 때마다 **새 프록시**를 만들어 `s is p`로는 자기 자신도 못 거른다 —
    `_element` 동일성으로 비교해야 한다(이 함정에 한 번 빠져 "그림이 자기와 겹침"을 냈다).
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    pics = [s for s in shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    hits = []
    for p in pics:
        x, y, w, h = box(p)
        for s in shapes:
            if s._element is p._element:
                continue
            sx, sy, sw, sh = box(s)
            if sw > 13 and sh > 7:
                continue  # 풀블리드 배경
            if sx < x + w and x < sx + sw and sy < y + h and y < sy + sh:
                lbl = s.text_frame.text.strip()[:16] if s.has_text_frame else "(무텍스트)"
                hits.append(lbl)
    return not hits, f"PICTURE {len(pics)}개 · 침범 {len(hits)} {hits[:3]}", len(hits)


def check_bounds(shapes, bottom=CONTENT_BOTTOM, right=RIGHT_EDGE, skip_tops=()):
    """P2③ 콘텐츠 하한 · P2④ 풀폭 우변."""
    over, badr = [], []
    for s in shapes:
        x, y, w, h = box(s)
        if w > 13 and h > 7:
            continue  # 풀블리드 배경
        if any(abs(y - t) < 0.01 for t in skip_tops):
            continue
        if y + h > bottom + 0.01:
            over.append((round(y + h, 2), s.text_frame.text[:14] if s.has_text_frame else ""))
        if abs(x - 0.6) < 0.001 and w > 11 and abs(x + w - right) > 0.005:
            badr.append(round(x + w, 3))
    return (
        not (over or badr),
        f"bottom>{bottom}: {len(over)} {over[:2]} · 풀폭 right 위반: {badr}",
        len(over) + len(badr),
    )


def check_air(pairs, air_min=AIR_MIN):
    """§6 공기 — (이름, 위 요소 bottom, 아래 요소 top) 목록을 받아 간격 검사.

    전체 쌍 자동 검사는 오탐이 많다(헤더 킥커→헤드라인 0.02는 확정 설계) — 설계상
    붙으면 안 되는 지점만 명시적으로 넘긴다.
    """
    bad = [(n, round(t - b, 3)) for n, b, t in pairs if t - b < air_min]
    lines = [
        f"  {n:28} 공기 {t - b:+.3f} {'OK' if t - b >= air_min else '**FAIL**'}"
        for n, b, t in pairs
    ]
    return not bad, "\n".join(lines), len(bad)


def density_band(snapshot_path=None):
    """골든 스냅샷에서 본문 장의 밀도 하한(도형 수·텍스트 프레임 수)을 파생.

    출처(2026-07-16 배선 재설계): 골든이 "복제 템플릿"에서 "변형 가능한 출발점"이 되면서
    스냅샷 회귀가 실전 덱을 못 지킨다 — 변형·신규 장의 "밀집화"를 재는 기계 게이트가 이것.
    임계는 리터럴이 아니라 **골든의 가장 성긴 본문 장**에서 파생한다(§4 범용 해결) —
    골든이 바뀌어 스냅샷을 재생성하면 밴드도 따라 움직인다.
    """
    import json
    from pathlib import Path

    from .reference import SLIDE_ORDER

    fixed = {"cover", "toc", "part", "closing"}  # 정형 장 — 밀도 대상 아님
    if snapshot_path is None:
        snapshot_path = Path(__file__).resolve().parents[2] / "assets/golden-snapshot.json"
    snap = json.loads(Path(snapshot_path).read_text(encoding="utf-8"))
    body = [sl for (key, _n), sl in zip(SLIDE_ORDER, snap["slides"]) if key not in fixed]
    shape_counts = [len(sl) for sl in body]
    frame_counts = [sum(1 for sh in sl if (sh.get("text") or "").strip()) for sl in body]
    return {
        "shapes_min": min(shape_counts),
        "frames_min": min(frame_counts),
        "n_body": len(body),
    }


def check_density(shapes, band):
    """골든 밀도 밴드 — 변형(adapted)·신규(novel) 장이 골든 본문 최저 밀도보다 성기면 FAIL."""
    n_shapes = len(shapes)
    n_frames = sum(1 for s in shapes if s.has_text_frame and s.text_frame.text.strip())
    ok = n_shapes >= band["shapes_min"] and n_frames >= band["frames_min"]
    deficit = max(0, band["shapes_min"] - n_shapes) + max(0, band["frames_min"] - n_frames)
    return (
        ok,
        f"도형 {n_shapes}(하한 {band['shapes_min']}) · "
        f"텍스트 프레임 {n_frames}(하한 {band['frames_min']}) — 골든 본문 {band['n_body']}장 파생",
        deficit,
    )


def generic_checks(shapes, accent, band=None, dup_allow=0):
    """레이아웃 무관 전역 오딧 묶음 — 변형(adapted)·신규(novel)·content 주입 골든 장에 적용.

    레이아웃 특정 규칙(진행형 도형 허용 수·ink_allow·공기 쌍·노드 높이)은 여기 없다 —
    그건 장 스크립트가 자체 assert하거나 채점(P4)이 본다. audit_deck.py가 이 묶음의 러너.
    """
    res = [
        ("§2 accent 상한", *check_accent(shapes, accent)),
        ("P4⑤ 판정 단어 대비", *check_verdict_contrast(shapes, None)),
        ("P4④ 채움률", *check_fill_ratio(shapes)),
        ("P4③ 노드 재탕", *check_duplicate_nodes(shapes, allow=dup_allow)),
        ("P4⑩ 노드 클래스", *check_node_class(shapes)),
        ("§F 그림 침범", *check_picture_overlap(shapes)),
        ("P2③④ 경계", *check_bounds(shapes, skip_tops=(6.60, 7.15))),
    ]
    if band:
        res.append(("§6-D 밀도 밴드", *check_density(shapes, band)))
    return res


def report(results, known=None):
    """(이름, ok, 상세, 위반수) 목록 → 출력 + 실패 목록.

    **위반 수는 규칙이 직접 반환한다** — 상세 문자열에서 정규식으로 긁으면 규칙마다 형식이
    달라 조용히 None이 되고 래칫이 죽는다(초판이 `채움 2 (허용 0)`을 못 읽어 그랬다).

    known: {규칙명: 기준선 위반 수} — 이 오딧 **이전에** 사용자 승인으로 확정된 장의 기존 빚.
      규칙을 끄지 않고 기준선으로 박는다(래칫): 기준선 이하면 통과, **넘으면 FAIL**.
      · 규칙을 빼면 hollow-PASS가 된다(존재≠검사).
      · 그대로 빨간불이면 게이트가 상시 적색이라 무시당하고, 결국 우회한다 —
        이 프로젝트가 이미 겪은 실패 경로다.
      기준선을 내리는 건 해당 장 재설계 때. 좋아지면 알려서 기준선을 조인다.
    """
    known = known or {}
    fails = []
    for name, ok, detail, n in results:
        base = known.get(name)
        if not ok and base is not None and n is not None and n <= base:
            mark = "DEBT" if n == base else "GOOD"
            note = "기지 빚(기준선)" if n == base else f"개선 {base}→{n} · 기준선을 낮출 것"
            print(f"{mark} {name:22} {detail}  ← {note}")
            continue
        print(f"{'OK  ' if ok else 'FAIL'} {name:22} {detail}")
        if not ok:
            fails.append(name)
    print("\n" + ("AUDIT PASS" if not fails else f"AUDIT FAIL: {fails}"))
    return fails
