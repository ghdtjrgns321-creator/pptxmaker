"""S6 시안 2종 — 4노드 결정적 파이프라인 (00_factsheet.md §C, 모방: s07_flow.png = McKinsey p60).

시안 A(사이드바형): 매크로 셰브런 밴드 + 마이크로 입출력 플로우 + 우측 다크 사이드바(왜 결정적인가).
시안 B(풀와이드형): 풀폭 밴드 + 노드별 상세 4칼럼 + 하단 2패널(유사도 0 · 35토픽) + 결론 바.
실행: 이 타입은 goldenfab 레지스트리 경유로만 렌더된다 — 골든 19장 확인은
      `uv run python golden/build_golden.py`(2026-07-15 단일화로 시안 개별 실행 경로 폐지).
"""

from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from . import grid as G
from ._variant_h import _shape_step
from .kit import SLIDE_H, SLIDE_W, add_box, add_text, load_kit, mix, set_shape_text

K = load_kit()
C, S, F = K["rgb"], K["sizes"], K["fonts"]

KICKER = "2. 파이프라인"
SOURCE = "출처: 4_SEARCH-PIPELINE.md (00_factsheet.md §C)"
OVERLAP = G.FLOW_H * 0.55 / 2

# §C — 매크로 4노드 (도형 안=단어, 설명은 도형 밖 칼럼)
NODES = [
    ("Analyze", "용어사전 substring 진입 — 임베딩이 아니다"),
    ("Retrieve", "graph.traverse hops=1 + 상호참조 E3 1홉"),
    ("Generate", "판단트리 via_topic 다중 주입 · PydanticAI 구조화 출력"),
    ("Format", "감리 지적사례 경고와 꼬리질문 부가"),
]
# §C — 마이크로 입출력 체인: 산출물 5 + 공정 4
ARTIFACTS = ["질문 원문", "매칭 토픽", "문단·사례 묶음", "구조화 결론", "최종 답변"]
PROCESSES = ["용어사전 매칭", "그래프 1홉 탐색", "판단트리 주입 생성", "경고·꼬리질문 부가"]


def header(slide, headline, right):
    add_text(
        slide, G.MARGIN_L, 0.42, 6.0, 0.28, KICKER, S["caption"], F["head"], C["muted"], bold=True
    )
    add_text(
        slide,
        G.MARGIN_L,
        0.72,
        right - G.MARGIN_L,
        0.55,
        headline,
        S["section"],
        F["head"],
        C["primary"],
        bold=True,
    )
    add_box(slide, G.MARGIN_L, G.RULE_Y, right - G.MARGIN_L, 0.014, fill=C["muted"])


def _subhead(slide, x, y, w, text):
    add_text(slide, x, y, w, 0.32, text, S["head"], F["head"], C["primary"], bold=True)
    add_box(slide, x, y + 0.35, w, 0.012, fill=C["muted"])


def macro_band(slide, x0, x1, y):
    step_w = (x1 - x0 + 3 * OVERLAP) / 4
    for i, (head, _desc) in enumerate(NODES):
        _shape_step(
            slide,
            x0 + i * (step_w - OVERLAP),
            y,
            step_w,
            G.FLOW_H,
            "pentagon" if i == 0 else "chevron",
            C["bg_alt"],
            head,
            f"NODE {i + 1}",
            C["primary"],
            mix(C["primary"], C["bg"], 0.35),
        )
    return step_w


def micro_flow(slide, x0, x1, y_c, d=0.44):
    pitch = (x1 - x0 - d) / 4
    centers = [x0 + d / 2 + i * pitch for i in range(5)]
    for a, b in zip(centers, centers[1:]):
        ln = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Inches(a + d / 2), Inches(y_c), Inches(b - d / 2), Inches(y_c)
        )
        ln.line.color.rgb = C["muted"]
        ln.line.width = Pt(1.5)
    for i, cx in enumerate(centers):
        fill = C["accent"] if i == len(centers) - 1 else C["bg"]
        line = None if i == len(centers) - 1 else C["muted"]
        add_box(
            slide, cx - d / 2, y_c - d / 2, d, d, fill=fill, line=line, line_w=1.25, shape="oval"
        )
        add_text(
            slide,
            cx - pitch / 2,
            y_c + d / 2 + 0.08,
            pitch,
            0.45,
            ARTIFACTS[i],
            S["caption"],
            F["body"],
            C["primary"],
            bold=(i in (0, 4)),
            align=PP_ALIGN.CENTER,
        )
    for i, proc in enumerate(PROCESSES):
        mid = (centers[i] + centers[i + 1]) / 2
        add_text(
            slide,
            mid - pitch / 2,
            y_c - d / 2 - 0.38,
            pitch,
            0.3,
            proc,
            S["caption"],
            F["body"],
            C["muted"],
            align=PP_ALIGN.CENTER,
        )


def variant_a(prs):
    """A — mck p60 이식: 매크로 밴드 + 마이크로 플로우 + 우측 다크 사이드바."""
    panel_x = 10.2
    body_r = panel_x - 0.4  # 9.8
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    add_box(slide, panel_x, 0, SLIDE_W - panel_x, SLIDE_H, fill=C["primary"])
    header(slide, "질문에서 답까지 — 점수 경쟁 없는 4노드 결정 경로", body_r)
    add_text(
        slide,
        G.MARGIN_L,
        G.CONTENT_TOP,
        body_r - G.MARGIN_L,
        0.3,
        "네 노드 전부가 결정적으로 동작한다 — 같은 질문은 항상 같은 경로를 지난다.",
        S["sub"],
        F["body"],
        C["text"],
    )
    step_w = macro_band(slide, G.MARGIN_L, body_r, 2.35)
    for i, (_head, desc) in enumerate(NODES):  # 밴드 아래 설명 칼럼 (도형 밖 완전 문장)
        x = G.MARGIN_L + i * (step_w - OVERLAP)
        add_text(
            slide,
            x + 0.1,
            2.35 + G.FLOW_H + 0.12,
            step_w - 0.55,
            0.6,
            desc,
            S["caption"],
            F["body"],
            C["muted"],
            line_spacing=1.2,
        )
    _subhead(slide, G.MARGIN_L, 4.15, body_r - G.MARGIN_L, "입출력 체인 — 각 단계가 넘기는 산출물")
    micro_flow(slide, G.MARGIN_L + 0.2, body_r - 0.2, 5.45)
    # 사이드바 — 왜 결정적인가 (§C)
    px, pw = panel_x + 0.35, SLIDE_W - panel_x - 0.7
    add_text(slide, px, 0.9, pw, 0.3, "왜 결정적인가", S["head"], F["head"], C["bg"], bold=True)
    add_box(slide, px, 1.3, pw, 0.014, fill=mix(C["primary"], C["muted"], 0.45))
    add_text(
        slide, px, 1.75, pw, 1.0, "0", int(S["display"] * 1.5), F["head"], C["accent"], bold=True
    )
    add_text(
        slide,
        px,
        2.85,
        pw,
        0.6,
        "진입부 임베딩 유사도 — 점수 경쟁이 존재하지 않는다",
        S["caption"],
        F["body"],
        C["bg_alt"],
        line_spacing=1.25,
    )
    bullets = [
        "LLM이 고르는 것은 35토픽 목록뿐 — 자유 생성이 아니라 목록 중 선택",
        "그래프 탐색은 hops=1 고정 — 어디까지 볼지도 결정적",
        "출력은 PydanticAI 구조화 — 형식 강제",
    ]
    for i, b in enumerate(bullets):
        y = 3.85 + i * 0.85
        add_box(slide, px, y + 0.06, 0.1, 0.1, fill=C["bg_alt"], shape="oval")
        add_text(
            slide,
            px + 0.25,
            y,
            pw - 0.25,
            0.7,
            b,
            S["caption"],
            F["body"],
            C["bg_alt"],
            line_spacing=1.25,
        )
    bar = add_box(slide, G.MARGIN_L, G.BAR_Y, body_r - G.MARGIN_L, G.BAR_H, fill=C["primary"])
    tf = bar.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "질문에서 답까지 4개 노드 전부가 결정적으로 동작한다"
    r.font.name, r.font.bold = F["head"], True
    r.font.size = Pt(S["body"])
    r.font.color.rgb = C["bg"]
    add_text(
        slide,
        G.MARGIN_L,
        G.SOURCE_Y,
        body_r - G.MARGIN_L,
        0.3,
        SOURCE,
        S["foot"],
        F["body"],
        C["muted"],
    )
    return slide


def variant_b(prs):
    """B — 풀와이드: 밴드 + 노드 상세 4칼럼(입→출 명시) + 하단 2패널 + 결론 바."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    header(slide, "같은 질문은 항상 같은 경로 — 4노드 결정적 파이프라인", G.RIGHT_EDGE)
    add_text(
        slide,
        G.MARGIN_L,
        G.CONTENT_TOP,
        G.RIGHT_EDGE - G.MARGIN_L,
        0.3,
        "진입부터 출력까지 확률 점수가 개입하는 지점이 없다 — 진입부 임베딩 유사도 0.",
        S["sub"],
        F["body"],
        C["text"],
    )
    step_w = macro_band(slide, G.MARGIN_L, G.RIGHT_EDGE, 2.35)
    io_pairs = [  # 칼럼: 설명 + 입→출 (ARTIFACTS 체인 파생)
        ("질문 원문 → 매칭 토픽", NODES[0][1]),
        ("매칭 토픽 → 문단·사례 묶음", NODES[1][1]),
        ("문단·사례 묶음 → 구조화 결론", NODES[2][1]),
        ("구조화 결론 → 최종 답변", NODES[3][1]),
    ]
    for i, (io, desc) in enumerate(io_pairs):
        x = G.MARGIN_L + i * (step_w - OVERLAP)
        add_text(
            slide,
            x + 0.1,
            3.25,
            step_w - 0.55,
            0.28,
            io,
            S["caption"],
            F["head"],
            C["primary"],
            bold=True,
        )
        add_text(
            slide,
            x + 0.1,
            3.58,
            step_w - 0.45,
            0.6,
            desc,
            S["caption"],
            F["body"],
            C["muted"],
            line_spacing=1.2,
        )
    _subhead(
        slide,
        G.MARGIN_L,
        G.LOW_HEAD_Y - 0.25,
        G.RIGHT_EDGE - G.MARGIN_L,
        "왜 결정적인가 — 점수 경쟁이 사라진 두 지점",
    )
    panels = [
        (
            G.MARGIN_L,
            "0",
            "진입부 임베딩 유사도",
            "진입은 용어사전 substring 매칭 — 유사도 점수 경쟁이 존재하지 않는다.",
        ),
        (
            G.MARGIN_L + (G.RIGHT_EDGE - G.MARGIN_L + G.BOX_GAP) / 2,
            "35",
            "LLM의 선택지 = 토픽 목록",
            "자유 생성이 아니라 35토픽 목록 중 선택 — 이후는 그래프가 결정한다.",
        ),
    ]
    pw = (G.RIGHT_EDGE - G.MARGIN_L - G.BOX_GAP) / 2
    for px, big, head, desc in panels:
        add_box(slide, px, G.BOX_Y, pw, G.BOX_H, fill=C["bg_alt"], shape="round")
        add_box(slide, px, G.BOX_Y, 0.05, G.BOX_H, fill=C["accent"])
        add_text(
            slide,
            px + 0.25,
            G.BOX_Y + 0.15,
            1.3,
            0.75,
            big,
            S["display"],
            F["head"],
            C["accent"],
            bold=True,
        )
        add_text(
            slide,
            px + 1.7,
            G.BOX_Y + 0.16,
            pw - 1.95,
            0.28,
            head,
            S["body"],
            F["head"],
            C["primary"],
            bold=True,
        )
        add_text(
            slide,
            px + 1.7,
            G.BOX_Y + 0.48,
            pw - 1.95,
            0.5,
            desc,
            S["caption"],
            F["body"],
            C["muted"],
            line_spacing=1.2,
        )
    bar = add_box(slide, G.MARGIN_L, G.BAR_Y, G.RIGHT_EDGE - G.MARGIN_L, G.BAR_H, fill=C["primary"])
    tf = bar.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "질문에서 답까지 4개 노드 전부가 결정적으로 동작한다"
    r.font.name, r.font.bold = F["head"], True
    r.font.size = Pt(S["body"])
    r.font.color.rgb = C["bg"]
    add_text(
        slide,
        G.MARGIN_L,
        G.SOURCE_Y,
        G.RIGHT_EDGE - G.MARGIN_L,
        0.3,
        SOURCE,
        S["foot"],
        F["body"],
        C["muted"],
    )
    return slide


def _arrow(slide, x1, y1, x2, y2, elbow=False):
    """방향 화살표 커넥터 — tailEnd 삼각형 XML 부가."""
    from pptx.oxml.ns import qn

    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.ELBOW if elbow else MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    conn.line.color.rgb = C["muted"]
    conn.line.width = Pt(1.5)
    ln = conn.line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    return conn


DEFAULT_C = {
    "kicker": "2. 파이프라인",
    "headline": "질문에서 응답까지 — 분기까지 전부 결정적인 실행 그래프",
    "subtitle": "진입부 임베딩 유사도 0 — LLM이 고르는 것은 35토픽 목록뿐, 경로는 그래프가 결정한다.",
    "node_names": ["질문", "Analyze", "판정", "Retrieve", "Generate", "Format", "응답"],
    "node_tags": ["용어사전 매칭", "그래프 1홉 탐색", "판단트리 주입", "경고·꼬리질문"],
    "in_label": "IN",
    "reject_box": "거절 메시지",
    "out_label": "OUT",
    "reject_desc": "범위 밖 질문은 본선에 진입하지 못하고 즉시 거절된다",
    "detail_head": "핵심 기술 — 결정성을 어떻게 만들었나",
    "details": [
        (
            "01  용어사전 — 후보 진입점",
            "등재 423 · AI 신규 창작 0.",
            " 사람이 만든 자료 3종(질의 매핑 288 · 사례 제목 123 · 부록A 정의 9)에서 AI 초안 + 사람 전수 검수로 등재. 확정 라우터가 아니라 후보 진입점이다.",
        ),
        (
            "02  지식그래프 — 기계 생성 뼈대",
            "노드 929 · 간선 2,697.",
            " 기준서 공식 소제목이 그대로 개념 80, 문단 250 배정. 계층·참조 간선은 기준서에서 기계 생성 — AI가 만드는 것은 용어 색인 하나뿐이다.",
        ),
        (
            "03  판단트리 — 판단 순서의 사전 조립",
            "조건-분기 골격 41개, 원문 앵커 포함.",
            " 기준서에 흩어진 판단 절차를 미리 조립해 두고, 진입 개념과 트리거 개념의 최다 겹침 + 주제 직속 트리를 전부 프롬프트에 주입한다.",
        ),
        (
            "04  구조화 출력 — 코드 수준 강제",
            "PydanticAI 스키마 강제.",
            " 답변 형식을 코드로 강제하고, 위반하면 result_validator가 자동 재시도한다. 감리 경고·꼬리질문은 마지막에 부가되는 보조 장치다.",
        ),
    ],
    "rerank_note": "초기 설계에 있던 유사도 재정렬(rerank) 단계는 제거 — 그래프가 이미 결정적으로 선별하므로 재정렬할 것이 없다",
    "bar": "질문에서 답까지 4개 노드 전부가 결정적으로 동작한다",
    "source": "출처: 4_SEARCH-PIPELINE.md (00_factsheet.md §C)",
}


def variant_c(prs, c=None):
    """C — 원문(4_SEARCH-PIPELINE.md L8~19) 아스키 구조의 도형화: 분기 포함 실행 그래프.

    c: 텍스트 내용 override(None=골든 기본값). 좌표·색·도형 종류는 고정."""
    c = {**DEFAULT_C, **(c or {})}
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    # 헤더 — header() 인라인(kicker·headline만 c에서, 출력은 header 호출과 동일)
    add_text(
        slide,
        G.MARGIN_L,
        0.42,
        6.0,
        0.28,
        c["kicker"],
        S["caption"],
        F["head"],
        C["muted"],
        bold=True,
    )
    add_text(
        slide,
        G.MARGIN_L,
        0.72,
        G.RIGHT_EDGE - G.MARGIN_L,
        0.55,
        c["headline"],
        S["section"],
        F["head"],
        C["primary"],
        bold=True,
    )
    add_box(slide, G.MARGIN_L, G.RULE_Y, G.RIGHT_EDGE - G.MARGIN_L, 0.014, fill=C["muted"])
    add_text(
        slide,
        G.MARGIN_L,
        G.CONTENT_TOP,
        G.RIGHT_EDGE - G.MARGIN_L,
        0.3,
        c["subtitle"],
        S["sub"],
        F["body"],
        C["text"],
    )
    # ── 메인 레인 — 질문 ▶ Analyze ▶ ◇판정 ▶ Retrieve ▶ Generate ▶ Format ▶ 응답 ──
    LANE_Y, NODE_H = 2.35, 0.6
    mid = LANE_Y + NODE_H / 2
    widths = [1.0, 1.45, 1.15, 1.45, 1.45, 1.45, 1.0]
    gap = (G.RIGHT_EDGE - G.MARGIN_L - sum(widths)) / 6
    xs, x = [], G.MARGIN_L
    for w in widths:
        xs.append(x)
        x += w + gap
    # 노드 스타일(도형 종류·색 고정), 이름만 c["node_names"]에서
    node_style = [
        ("round", C["bg"], C["muted"], C["primary"]),
        ("round", C["bg_alt"], None, C["primary"]),
        ("diamond", C["bg"], C["accent"], C["primary"]),
        ("round", C["bg_alt"], None, C["primary"]),
        ("round", C["bg_alt"], None, C["primary"]),
        ("round", C["bg_alt"], None, C["primary"]),
        ("round", C["primary"], None, C["bg"]),
    ]
    for (kind, fill, line, ink), name, nx, w in zip(node_style, c["node_names"], xs, widths):
        h = NODE_H + 0.25 if kind == "diamond" else NODE_H
        y = mid - h / 2
        shp = add_box(
            slide, nx, y, w, h, fill=fill, line=line, line_w=1.25 if line else None, shape=kind
        )
        set_shape_text(
            shp, name, S["caption" if kind == "diamond" else "body"], F["head"], ink, bold=True
        )
    for i in range(6):
        _arrow(slide, xs[i] + widths[i], mid, xs[i + 1], mid)
    tag_indices = [1, 3, 4, 5]
    for i, tag in zip(tag_indices, c["node_tags"]):
        add_text(
            slide,
            xs[i] - 0.35,
            LANE_Y + NODE_H + 0.1,
            widths[i] + 0.7,
            0.24,
            tag,
            S["caption"],
            F["head"],
            C["primary"],
            bold=True,
            align=PP_ALIGN.CENTER,
        )
    add_text(
        slide,
        xs[2] + widths[2] + 0.03,
        mid - 0.32,
        0.5,
        0.22,
        c["in_label"],
        S["caption"],
        F["head"],
        C["accent"],
        bold=True,
    )
    # ── OUT 분기 — 다이아 아래 수직 낙하, 거절 박스 + 설명은 박스 우측 인라인 ──
    dia_cx = xs[2] + widths[2] / 2
    rej_w, rej_h = 2.0, 0.45
    rej_x, rej_y = dia_cx - rej_w / 2, 3.55
    rej = add_box(slide, rej_x, rej_y, rej_w, rej_h, fill=C["bg_alt"], shape="round")
    set_shape_text(rej, c["reject_box"], S["caption"], F["head"], C["primary"], bold=True)
    dia_bot = mid + (NODE_H + 0.25) / 2
    _arrow(slide, dia_cx, dia_bot, dia_cx, rej_y)
    add_text(
        slide,
        dia_cx + 0.12,
        dia_bot + 0.12,
        0.6,
        0.22,
        c["out_label"],
        S["caption"],
        F["head"],
        C["accent"],
        bold=True,
    )
    add_text(
        slide,
        rej_x + rej_w + 0.2,
        rej_y + 0.1,
        4.2,
        0.26,
        c["reject_desc"],
        S["caption"],
        F["body"],
        C["muted"],
    )
    # ── 하단: 노드별 핵심 기술 — 평어 상세 4칼럼 (원문 §C·L47·L53·L63~65) ──
    _subhead(slide, G.MARGIN_L, 4.35, G.RIGHT_EDGE - G.MARGIN_L, c["detail_head"])
    det_w = (G.RIGHT_EDGE - G.MARGIN_L - 3 * 0.35) / 4  # 2.771
    for i, (head, lead, body) in enumerate(c["details"]):
        dx = G.MARGIN_L + i * (det_w + 0.35)
        add_text(
            slide, dx, 4.85, det_w, 0.28, head, S["caption"], F["head"], C["primary"], bold=True
        )
        add_text(
            slide,
            dx,
            5.18,
            det_w,
            1.05,
            [[(lead, {"bold": True, "color": C["primary"]}), (body, {})]],
            S["caption"],
            F["body"],
            C["muted"],
            line_spacing=1.25,
        )
        if i < 3:
            add_box(slide, dx + det_w + 0.175, 4.85, 0.012, 1.05, fill=C["bg_alt"])
    add_text(
        slide,
        G.MARGIN_L,
        6.05,
        G.RIGHT_EDGE - G.MARGIN_L,
        0.3,
        c["rerank_note"],
        S["caption"],
        F["body"],
        C["muted"],
    )
    bar = add_box(slide, G.MARGIN_L, G.BAR_Y, G.RIGHT_EDGE - G.MARGIN_L, G.BAR_H, fill=C["primary"])
    tf = bar.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = c["bar"]
    r.font.name, r.font.bold = F["head"], True
    r.font.size = Pt(S["body"])
    r.font.color.rgb = C["bg"]
    add_text(
        slide,
        G.MARGIN_L,
        G.SOURCE_Y,
        G.RIGHT_EDGE - G.MARGIN_L,
        0.3,
        c["source"],
        S["foot"],
        F["body"],
        C["muted"],
    )
    return slide


def audit(prs):
    """오딧 — bottom>6.35 / 밴드·풀폭 right / 룰-채움 겹침 / accent ≤4."""
    EMU = 914400
    accent_hex = str(C["accent"])
    fails = []
    for si, sl in enumerate(prs.slides):
        rules, solids, acc = [], [], 0
        for sh in sl.shapes:
            L, T, W, H = sh.left / EMU, sh.top / EMU, sh.width / EMU, sh.height / EMU
            if (L == 0 and T == 0 and W > 13) or (W < 4 and H > 6):  # 배경·사이드바 제외
                continue
            if H <= 0.02 and W > 1:
                rules.append((L, T, W))
            elif H > 0.3 and W > 0.5 and sh.shape_type != 17:
                solids.append((L, T, W, H))
            if T + H > G.CONTENT_BOTTOM + 0.02 and T < G.BAR_Y - 0.02:
                fails.append((si, "bottom>6.35", round(T + H, 3)))
            if W > 11.5 and W < 13 and abs((L + W) - G.RIGHT_EDGE) > 0.02:
                fails.append((si, "풀폭 right≠12.733", round(L + W, 3)))
            try:
                if str(sh.fill.fore_color.rgb) == accent_hex:
                    acc += 1
            except Exception:
                pass
            if sh.has_text_frame:
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        if (
                            r.font.color
                            and r.font.color.rgb is not None
                            and str(r.font.color.rgb) == accent_hex
                        ):
                            acc += 1
        for rl, rt, rw in rules:
            for sl_, st, sw, shh in solids:
                if rl < sl_ + sw and rl + rw > sl_ and st - 0.005 <= rt <= st + shh - 0.005:
                    fails.append((si, "룰-채움 겹침", round(rt, 3)))
        if acc > 4:
            fails.append((si, "accent 초과", acc))
        print(f"slide{si}: accent {acc}곳")
    assert not fails, f"AUDIT FAIL {fails}"
    print("audit pass")


def main():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    variant_a(prs)
    variant_b(prs)
    variant_c(prs)
    audit(prs)
    from pathlib import Path

    out = Path(__file__).parent / "variants" / "s06_variants.pptx"
    prs.save(str(out))
    print(f"saved {out}")


if __name__ == "__main__":
    main()
