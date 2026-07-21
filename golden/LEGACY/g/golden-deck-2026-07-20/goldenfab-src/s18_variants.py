"""S18 시안 — 차별점: 정직한 한계 (물성 선언: 한계 = 경계 — 안은 결정 동작, 밖은 차단·유보).

콘텐츠 실물: 6_TEST §6.3(타 기준서 4건 IAS38·1002·1008·1037, 진입 누락 2건),
4_SEARCH L49(타 기준서 용어 감지 시 코드 레벨 강제 OUT — 예: 1116호 증분차입이자율),
00_factsheet.md §D·§H(경계 안 실측 스탯).
실행: 이 타입은 goldenfab 레지스트리 경유로만 렌더된다 — 골든 19장 확인은
      `uv run python golden/build_golden.py`(2026-07-15 단일화로 시안 개별 실행 경로 폐지).
"""

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from . import grid as G
from .kit import SLIDE_H, SLIDE_W, add_box, add_text, load_kit

K = load_kit()
C, S, F = K["rgb"], K["sizes"], K["fonts"]

KICKER = "4. 문제와 해결"  # 2026-07-16 5부 개편 — 한계 장은 검증 장 뒤 Part Ⅳ 소속
SOURCE = "출처: 6_TEST-DECISIONS.md §6.3 · 4_SEARCH-PIPELINE.md (라우팅 강제 OUT) · 00_factsheet.md §D·§H"


def header(slide, headline):
    add_text(
        slide, G.MARGIN_L, 0.42, 8.0, 0.28, KICKER, S["caption"], F["head"], C["muted"], bold=True
    )
    add_text(
        slide,
        G.MARGIN_L,
        0.72,
        G.RIGHT_EDGE - G.MARGIN_L,
        0.55,
        headline,
        S["section"],
        F["head"],
        C["primary"],
        bold=True,
    )
    add_box(slide, G.MARGIN_L, G.RULE_Y, G.RIGHT_EDGE - G.MARGIN_L, 0.014, fill=C["muted"])


def bar_and_source(slide, text):
    bar = add_box(slide, G.MARGIN_L, G.BAR_Y, G.RIGHT_EDGE - G.MARGIN_L, G.BAR_H, fill=C["primary"])
    tf = bar.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name, run.font.bold = F["head"], True
    run.font.size = Pt(S["body"])
    run.font.color.rgb = C["bg"]
    add_text(
        slide,
        G.MARGIN_L,
        G.SOURCE_Y,
        G.RIGHT_EDGE - G.MARGIN_L,
        0.3,
        SOURCE,
        S["foot"],
        F["body"],
        C["muted"],
    )


def _dash_arrow(slide, x1, y1, x2, y2):
    from pptx.oxml.ns import qn

    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = C["muted"]
    conn.line.width = Pt(1.5)
    ln = conn.line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "sm", "len": "sm"}))
    return conn


def _outside_chip(slide, x, y, w, head, sub):
    from pptx.enum.lang import MSO_LANGUAGE_ID

    chip = add_box(slide, x, y, w, 0.8, fill=C["bg"], line=C["muted"], line_w=1.0, shape="round")
    tf = chip.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.12)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    p1._p.get_or_add_pPr().set("eaLnBrk", "0")
    r1 = p1.add_run()
    r1.text = head
    r1.font.name, r1.font.bold = F["head"], True
    r1.font.size = Pt(S["caption"])
    r1.font.color.rgb = C["primary"]
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2._p.get_or_add_pPr().set("eaLnBrk", "0")
    r2 = p2.add_run()
    r2.text = sub
    r2.font.name = F["body"]
    r2.font.size = Pt(S["caption"])
    r2.font.color.rgb = C["muted"]
    for r in (r1, r2):
        r.font.language_id = MSO_LANGUAGE_ID.KOREAN
    return chip


def variant_a(prs):
    """A — 경계 도해: 안(1115 결정 동작·실측 스탯) / 밖(타 기준서 강제 OUT ✕ · 미등재 유보)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    header(slide, "정직한 한계 — 이 시스템이 서 있는 경계")
    from pptx.enum.lang import MSO_LANGUAGE_ID

    # 바깥 세계 (bg_alt 캔버스)
    add_box(slide, G.MARGIN_L, 1.95, G.RIGHT_EDGE - G.MARGIN_L, 4.05, fill=C["bg_alt"])
    add_text(
        slide,
        G.MARGIN_L + 0.2,
        2.05,
        3.0,
        0.24,
        "경계 밖",
        S["caption"],
        F["head"],
        C["muted"],
        bold=True,
    )
    # 안쪽 영역 — 결정적으로 동작하는 경계
    in_x, in_y, in_w, in_h = 4.35, 2.3, 4.6, 3.35
    region = add_box(
        slide, in_x, in_y, in_w, in_h, fill=C["bg"], line=C["primary"], line_w=2.5, shape="round"
    )
    region.adjustments[0] = 0.05
    add_text(
        slide,
        in_x,
        in_y + 0.22,
        in_w,
        0.32,
        "K-IFRS 1115 — 경계 안",
        S["head"],
        F["head"],
        C["primary"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        in_x,
        in_y + 0.56,
        in_w,
        0.26,
        "이 안에서는 결정적으로 동작한다",
        S["caption"],
        F["body"],
        C["muted"],
        align=PP_ALIGN.CENTER,
    )
    stats = [("노드", "929"), ("간선", "2,697"), ("판단트리", "41"), ("등재 용어", "423")]
    for i, (k, v) in enumerate(stats):
        col, row = divmod(i, 2)
        sx = in_x + 0.55 + col * 2.15
        sy = in_y + 1.0 + row * 0.8
        add_text(slide, sx, sy, 1.9, 0.38, v, S["title"], F["head"], C["primary"], bold=True)
        add_text(slide, sx, sy + 0.46, 1.9, 0.2, k, S["caption"], F["body"], C["muted"])
    add_box(slide, in_x + 0.5, in_y + 2.62, in_w - 1.0, 0.012, fill=C["bg_alt"])
    add_text(
        slide,
        in_x,
        in_y + 2.76,
        in_w,
        0.3,
        [
            [
                ("홀드아웃 ", {}),
                ("78 / 92", {"bold": True, "color": C["accent"]}),
                (" — 실행 에러 0건", {}),
            ]
        ],
        S["body"],
        F["body"],
        C["muted"],
        align=PP_ALIGN.CENTER,
    )
    # 밖 좌: 타 기준서 질문 — 강제 OUT
    _outside_chip(slide, 0.85, 3.15, 2.6, "타 기준서 질문", "IAS38 · 1002 · 1008 · 1037")
    _dash_arrow(slide, 3.45, 3.55, 4.33, 3.55)
    add_box(slide, 4.3, 3.33, 0.05, 0.44, fill=C["primary"])  # 차단 바
    add_text(slide, 3.98, 3.62, 0.3, 0.26, "✕", S["body"], F["head"], C["primary"], bold=True)
    add_text(
        slide,
        0.85,
        4.05,
        2.9,
        0.55,
        "라우팅이 코드 레벨 강제 OUT — 추측 답변 대신 거절. 실측 4건 (예: 1116호 '증분차입이자율' 감지)",
        S["caption"],
        F["body"],
        C["muted"],
        line_spacing=1.2,
    )
    # 밖 우: 등재되지 않은 표현 — 유보
    _outside_chip(slide, 9.85, 3.15, 2.6, "등재되지 않은 표현", "임베딩식 유사 확장 없음")
    _dash_arrow(slide, 9.85, 3.55, 8.97, 3.55)
    add_box(slide, 8.95, 3.33, 0.05, 0.44, fill=C["primary"])
    add_text(slide, 9.32, 3.62, 0.3, 0.26, "✕", S["body"], F["head"], C["primary"], bold=True)
    add_text(
        slide,
        9.55,
        4.05,
        2.9,
        0.55,
        "색인에 없으면 진입 실패 → '못 찾음'으로 유보 응답. 홀드아웃 실측 진입 누락 2건",
        S["caption"],
        F["body"],
        C["muted"],
        line_spacing=1.2,
    )
    add_text(
        slide,
        G.MARGIN_L,
        6.08,
        G.RIGHT_EDGE - G.MARGIN_L,
        0.24,
        "경계의 비용은 실측돼 있다 — 미재현 14건 중 타 기준서 4건·진입 누락 2건이 여기서 나온다. 나머지는 경계 안의 문제가 아니었다.",
        S["caption"],
        F["body"],
        C["muted"],
    )
    bar_and_source(
        slide, "경계를 넓히는 대신 경계 안을 결정적으로 — 밖의 질문은 추측하지 않고 거절한다"
    )
    return slide


# 콘텐츠 기본값(골든 내용) — c=None이면 이 값, override 시 텍스트만 교체(좌표·색·크기 고정)
DEFAULT = {
    "headline": "정직한 한계 — 이 시스템이 서 있는 경계",
    "outside_label": "경계 밖",
    "inside_title": "K-IFRS 1115 — 경계 안",
    "inside_policy": (
        "답하는 범위를 계약한다 — 안이면 결정적으로 답하고, 밖이면 거절, 모르면 유보. "
        "아는 척이 구조적으로 불가능한 설계다."
    ),
    "out_left_chip_head": "타 기준서 질문",
    "out_left_chip_sub": "IAS38 · 1002 · 1008 · 1037",
    "out_left_desc": "라우팅이 코드 레벨 강제 OUT — 추측 대신 거절. 실측 4건 (예: 1116호 '증분차입이자율' 감지)",
    "out_right_chip_head": "등재되지 않은 표현",
    "out_right_chip_sub": "임베딩식 유사 확장 없음",
    "out_right_desc": "색인에 없으면 진입 실패 → '못 찾음' 유보 응답. 홀드아웃 실측 진입 누락 2건",
    "bottom_head": "경계 그 밖의 한계 — 실측으로 아는 것",
    "limits": [  # 6_TEST §6.2·§6.3·L52 실물 — (한계, 실측, 남은 과제) 3종
        (
            "결론을 확정하지 못한 케이스",
            "헤지 2건 — 근거 문단을 다 찾고도 결론을 유보했다.",
            "검색이 아니라 생성 계층의 남은 과제",
        ),
        (
            "인용 완전성",
            "하드 인용 재현율 59.1% — 결론이 맞아도 근거 인용을 전부 담지는 못한다.",
            "인용률은 별도 지표로 추적",
        ),
        (
            "검증 자체의 한계",
            "92건은 개발 중 전수 열람돼 순수 홀드아웃이 아니다.",
            "최종 성능은 안 본 질문으로 검증",
        ),
    ],
    "bar": "경계를 넓히는 대신 경계 안을 결정적으로 — 못 하는 것까지 실측으로 세어 두었다",
}

# 경계 도해 좌표(2026-07-16 재재깎기) — 경계 안 상자는 **내용 크기로 압축**한다. 4유형 그리드를
# 넣었던 판은 사용자 기각: 확정/유보는 S4, 조건부·되물음은 S12 실물 발췌, IN/OUT 거절은 S6이
# 이미 보여준 재탕이었다. 빈 상자 문제의 근본 원인은 내용 부족이 아니라 **상자가 내용보다 큰 것**
# — 상자를 줄이고, 확보된 세로 공간은 이 장의 고유 재료(잔여 한계 3종)에 준다.
CANVAS_TOP, CANVAS_BOTTOM = 1.95, 3.95
IN_W, IN_H = 4.6, 1.25
IN_X = (SLIDE_W - IN_W) / 2  # 캔버스 정중앙
IN_Y = 2.15  # 영역 2.15~3.40 — 캔버스 안 광학 중앙
CHIP_Y, CHIP_H = 2.38, 0.8
LINE_Y = CHIP_Y + CHIP_H / 2  # 2.78 — 접근 화살표·차단 바·✕의 축
DESC_Y = 3.32  # 칩 bottom(3.18)과 공기 0.14
BHEAD_Y, BRULE_Y = 4.20, 4.54
CARD_Y, CARD_H = 4.70, 1.50  # 카드 bottom 6.20 — 하한(6.35)과 공기 0.27


def variant_b(prs, c=None):
    """B — 상단 경계 도해(안=범위 계약, 밖=차단·유보) + 하단 잔여 한계 3섹션(한계/실측/남은 과제).

    2026-07-16 재재깎기: 경계 안을 내용 크기로 압축(재탕 4유형 그리드 제거 — 사용자 기각),
    확보한 세로 공간을 이 장의 고유 재료인 잔여 한계 카드에 배분 — 카드가 3단(한계 12pt bold /
    실측 / → 남은 과제)으로 깊어졌다. c=텍스트 override(None=골든).
    """
    c = {**DEFAULT, **(c or {})}
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    header(slide, c["headline"])
    from pptx.enum.lang import MSO_LANGUAGE_ID

    # ── 상단: 경계 캔버스 ──
    add_box(
        slide,
        G.MARGIN_L,
        CANVAS_TOP,
        G.RIGHT_EDGE - G.MARGIN_L,
        CANVAS_BOTTOM - CANVAS_TOP,
        fill=C["bg_alt"],
    )
    add_text(
        slide,
        G.MARGIN_L + 0.2,
        CANVAS_TOP + 0.10,
        3.0,
        0.24,
        c["outside_label"],
        S["caption"],
        F["head"],
        C["muted"],
        bold=True,
    )
    region = add_box(
        slide, IN_X, IN_Y, IN_W, IN_H, fill=C["bg"], line=C["primary"], line_w=2.5, shape="round"
    )
    region.adjustments[0] = 0.10
    add_text(
        slide,
        IN_X,
        IN_Y + 0.15,
        IN_W,
        0.30,
        c["inside_title"],
        S["head"],
        F["head"],
        C["primary"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        IN_X + 0.3,
        IN_Y + 0.53,
        IN_W - 0.6,
        0.5,
        c["inside_policy"],
        S["caption"],
        F["body"],
        C["muted"],
        align=PP_ALIGN.CENTER,
        line_spacing=1.3,
    )
    # 밖 좌: 타 기준서 — 강제 OUT. 화살촉은 ✕ 앞에서 멈춰 보이게, ✕는 충돌점(선상·바 직전)에
    # 올린다 — 바 뒤에 숨은 화살촉·선 아래 떠 있는 ✕는 충돌 표식을 둘로 쪼갠다(채점 지적).
    _outside_chip(slide, 0.75, CHIP_Y, 2.5, c["out_left_chip_head"], c["out_left_chip_sub"])
    _dash_arrow(slide, 3.35, LINE_Y, IN_X - 0.42, LINE_Y)
    add_box(slide, IN_X - 0.05, LINE_Y - 0.22, 0.05, 0.44, fill=C["primary"])
    add_text(
        slide,
        IN_X - 0.40,
        LINE_Y - 0.13,
        0.3,
        0.26,
        "✕",
        S["body"],
        F["head"],
        C["primary"],
        bold=True,
    )
    add_text(
        slide,
        0.75,
        DESC_Y,
        3.35,
        0.55,
        c["out_left_desc"],
        S["caption"],
        F["body"],
        C["muted"],
        line_spacing=1.2,
    )
    # 밖 우: 미등재 표현 — 유보
    _outside_chip(slide, 10.03, CHIP_Y, 2.5, c["out_right_chip_head"], c["out_right_chip_sub"])
    _dash_arrow(slide, 10.03, LINE_Y, IN_X + IN_W + 0.47, LINE_Y)
    add_text(
        slide,
        IN_X + IN_W + 0.13,
        LINE_Y - 0.13,
        0.3,
        0.26,
        "✕",
        S["body"],
        F["head"],
        C["primary"],
        bold=True,
    )
    add_box(slide, IN_X + IN_W, LINE_Y - 0.22, 0.05, 0.44, fill=C["primary"])
    add_text(
        slide,
        9.38,
        DESC_Y,
        G.RIGHT_EDGE - 9.38 - 0.15,  # 캔버스 우변에서 0.15 안쪽 — 끝에 붙으면 잘려 보인다
        0.55,
        c["out_right_desc"],
        S["caption"],
        F["body"],
        C["muted"],
        line_spacing=1.2,
    )
    # ── 하단: 잔여 한계 3섹션 — 한계(12pt bold) / 실측 / → 남은 과제 ──
    add_text(
        slide,
        G.MARGIN_L,
        BHEAD_Y,
        G.RIGHT_EDGE - G.MARGIN_L,
        0.3,
        c["bottom_head"],
        S["head"],
        F["head"],
        C["primary"],
        bold=True,
    )
    add_box(slide, G.MARGIN_L, BRULE_Y, G.RIGHT_EDGE - G.MARGIN_L, 0.012, fill=C["muted"])
    sec_w, gap = 3.86, 0.27
    for i, (head, measured, task) in enumerate(c["limits"]):
        sx = G.MARGIN_L + i * (sec_w + gap)
        card = add_box(slide, sx, CARD_Y, sec_w, CARD_H, fill=C["bg_alt"], shape="round")
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.16)
        tf.margin_top = Inches(0.12)
        specs = [  # (런들, line_spacing) — 3단: 번호+한계 / 실측 / → 남은 과제
            (
                [
                    (f"0{i + 1}  ", C["accent"], True, S["body"]),
                    (head, C["primary"], True, S["body"]),
                ],
                None,
            ),
            ([(measured, C["muted"], False, S["caption"])], 1.3),
            ([(f"→ {task}", C["primary"], False, S["caption"])], 1.3),
        ]
        for li, (runs, spacing) in enumerate(specs):
            p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p._p.get_or_add_pPr().set("eaLnBrk", "0")
            if spacing:
                p.line_spacing = spacing
            if li > 0:
                p.space_before = Pt(4)
            for txt, col, bold, pt in runs:
                r = p.add_run()
                r.text = txt
                r.font.name = F["head"] if bold else F["body"]
                r.font.bold = bold
                r.font.size = Pt(pt)
                r.font.color.rgb = col
                r.font.language_id = MSO_LANGUAGE_ID.KOREAN
    bar_and_source(slide, c["bar"])
    return slide


def audit(prs):
    from pptx.oxml.ns import qn

    EMU = 914400
    fails = []
    for si, sl in enumerate(prs.slides):
        accents = 0
        for sh in sl.shapes:
            t, le = sh.top / EMU, sh.left / EMU
            b, r = t + sh.height / EMU, le + sh.width / EMU
            full_bleed = sh.width / EMU >= SLIDE_W - 0.05
            if t < 6.4 and b > 6.37 and not full_bleed:
                fails.append((si, "bottom", round(b, 2), sh.shape_id))
            if r > 12.75 and not full_bleed:
                fails.append((si, "right", round(r, 2), sh.shape_id))
            for el in sh._element.iter(qn("a:srgbClr")):
                if el.get("val") == "D66E3A":
                    if el.getparent().tag.endswith("}solidFill"):
                        accents += 1
        if accents > 4:
            fails.append((si, "accent>4", accents))
        if si == 1:  # variant_b 재재깎기 구조: 잔여 한계 카드 3(h=CARD_H) + 채움률·판정 대비
            from . import audit as GA

            cards = sum(
                1
                for sh in sl.shapes
                if abs(sh.height / EMU - CARD_H) < 0.01 and sh.width / EMU > 3.0
            )
            n_limits = len(DEFAULT["limits"])
            if cards != n_limits:
                fails.append((si, f"한계 카드≠{n_limits}", cards))
            shapes_b = list(sl.shapes)
            for rule_name, (ok, msg, _n) in (
                ("채움률", GA.check_fill_ratio(shapes_b)),
                ("판정 대비", GA.check_verdict_contrast(shapes_b, None)),
            ):
                if not ok:
                    fails.append((si, rule_name, msg))
    assert not fails, f"AUDIT FAIL {fails}"
    print("audit pass")


def main():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    variant_a(prs)
    variant_b(prs)
    audit(prs)
    out = Path(__file__).parent / "variants" / "s18_variants.pptx"
    prs.save(str(out))
    print(f"saved {out}")


if __name__ == "__main__":
    main()
