"""S6 장표형 프로토타입 F6 — 사용자 스케치(2026-07-18) 직접 구현. 골든 미편입.

사용자 레이아웃 스펙 (원문 스케치):
  [좌 레일 WHY — 네이비 밴드 "바꾸려는 것"]  AS-IS 카드 3행 → ▼(주황) → TO-BE 카드 2행
  [중앙 HOW — 네이비 밴드 "결정적 실행 그래프"]  칩 플로우(판정만 주황 솔리드) + 아이콘 카드 2×2
  [우 레일 IMPACT — 네이비 밴드 "보장되는 것"]  미니 지표 카드 3(숫자 크게)
  [하단 풀폭 네이비 푸터밴드]  핵심 장치: 후보 진입점 · 1홉 탐색 · 순서 사전 조립 · 코드 수준 강제

색: 네이비(밴드·배지·아이콘 — 레퍼런스 실측 1B2A6B) + 주황 accent(판정 칩·전환 화살표·핵심 수치).
제약 해제 세션 — 네이비는 브랜드킷 밖 리터럴이며 방향 확정 시 brand-kit 토큰으로 승격한다.
문구·수치: 사용자 스케치 + 00_factsheet.md §A'·§C (창작 0).

실행:
    uv run python -c "
    import sys; sys.path.insert(0, '.claude/skills/pptx-build/scripts')
    from goldenfab.reference import new_presentation
    from goldenfab.s06_proto_f import variant_f
    prs = new_presentation(); variant_f(prs); prs.save('golden/_proto_s06_f6.pptx')"
"""

from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches

from . import grid as G
from .kit import SLIDE_H, SLIDE_W, add_box, add_text, load_kit, mix, set_shape_text
from .s06_variants import DEFAULT_C, NODES

K = load_kit()
C, S, F = K["rgb"], K["sizes"], K["fonts"]

NAVY = RGBColor.from_string("1B2A6B")  # 레퍼런스 장표 실측 — 방향 확정 시 brand-kit 승격
NAVY_TINT = None  # 계산: mix(NAVY, bg, 0.92) — 칩 보더·연한 구분용

# ── 콘텐츠 (사용자 스케치 + 팩트시트) ──
ASIS_ROWS = ["임베딩 유사도 라우팅", "rerank 순서 변동", "LLM 자유 경로"]
TOBE_ROWS = ["경로는 그래프가 결정", "LLM은 토픽 선택만"]
FLOW = ["질문", "Analyze", "판정", "Retrieve", "Generate", "Format"]
HOW_CARDS = [  # (아이콘, 제목, 스탯, 라벨, 설명=NODES desc)
    ("book-open", "용어사전 진입", "423", "등재 · AI 창작 0"),
    ("link", "지식그래프", "929 · 2,697", "노드 · 간선"),
    ("git-branch", "판단트리 주입", "41", "골격 · 원문 앵커"),
    ("shield-check", "구조화 출력", "스키마 강제", "PydanticAI · 자동 재시도"),
]
IMPACT_CARDS = [
    ("0", "임베딩 의존 — 진입부 유사도"),
    ("제거", "재정렬(rerank) 단계"),
    ("즉시 거절", "범위 밖 질문"),
]
FOOTER = "핵심 장치 :   후보 진입점  ·  1홉 탐색  ·  순서 사전 조립  ·  코드 수준 강제"

# ── GRID ──
RAIL_TOP, RAIL_BOTTOM = 1.2, 6.35
L_X, L_W = G.MARGIN_L, 2.55
M_X, M_W = 3.37, 6.55
R_X = 10.14
BAND_H = 0.5


def _icon_chip(slide, x, y, icon, *, side=0.32, fill=None):
    from .kit import ROOT as _KIT_ROOT

    add_box(slide, x, y, side, side, fill=fill or NAVY)
    p = _KIT_ROOT / "assets/icons" / f"{icon}_white.png"
    pad = side * 0.18
    slide.shapes.add_picture(
        str(p), Inches(x + pad), Inches(y + pad), Inches(side - 2 * pad), Inches(side - 2 * pad)
    )


def _band(slide, x, y, w, en, ko):
    """네이비 밴드 — 소문자 영문 킥커 + 한글 라벨(흰 볼드) 센터."""
    band = add_box(slide, x, y, w, BAND_H, fill=NAVY)
    tf = band.text_frame
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    from pptx.enum.lang import MSO_LANGUAGE_ID
    from pptx.util import Pt

    p = tf.paragraphs[0]
    p._p.get_or_add_pPr().set("eaLnBrk", "0")
    p.alignment = PP_ALIGN.CENTER
    for txt, sz, bold, color in (
        (en + "  ", S["foot"], True, mix(NAVY, C["bg"], 0.55)),
        (ko, S["body"], True, C["bg"]),
    ):
        r = p.add_run()
        r.text = txt
        r.font.name = F["head"]
        r.font.size = Pt(sz)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.language_id = MSO_LANGUAGE_ID.KOREAN
    return band


def _card(slide, x, y, w, h, line):
    return add_box(slide, x, y, w, h, fill=C["bg"], line=line, line_w=1.0)


def variant_f(prs, c=None):
    """F6 — WHY | HOW | IMPACT 3레일 장표 (사용자 스케치 구현)."""
    c = {**DEFAULT_C, **(c or {})}
    navy_line = mix(NAVY, C["bg"], 0.75)
    body_ink = mix(C["muted"], C["primary"], 0.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    FULL_W = G.RIGHT_EDGE - G.MARGIN_L
    R_W = G.RIGHT_EDGE - R_X

    # 압축 헤더
    add_text(
        slide, G.MARGIN_L, 0.28, 6.0, 0.22, c["kicker"], S["foot"], F["head"], C["muted"], bold=True
    )
    add_text(
        slide,
        G.MARGIN_L,
        0.5,
        FULL_W,
        0.4,
        c["headline"],
        S["compact"],
        F["head"],
        C["primary"],
        bold=True,
    )

    # ── 좌 레일 WHY ──
    _band(slide, L_X, RAIL_TOP, L_W, "WHY", "바꾸려는 것")
    asis_y, asis_h = RAIL_TOP + 0.7, 1.66
    _card(slide, L_X, asis_y, L_W, asis_h, navy_line)
    pill = add_box(
        slide,
        L_X + 0.14,
        asis_y + 0.14,
        1.5,
        0.3,
        fill=C["bg"],
        line=NAVY,
        line_w=1.25,
        shape="round",
    )
    set_shape_text(pill, "AS-IS (기존 한계)", S["foot"], F["head"], NAVY, bold=True)
    for j, row in enumerate(ASIS_ROWS):
        add_text(
            slide,
            L_X + 0.2,
            asis_y + 0.56 + j * 0.34,
            L_W - 0.4,
            0.28,
            [[("·  ", {"bold": True, "color": NAVY}), (row, {})]],
            S["caption"],
            F["body"],
            body_ink,
        )
    arr_y = asis_y + asis_h + 0.1
    add_text(
        slide,
        L_X,
        arr_y - 0.06,
        L_W,
        0.4,
        "▼",
        S["head"],
        F["head"],
        C["accent"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    tobe_y, tobe_h = arr_y + 0.42, 1.32
    _card(slide, L_X, tobe_y, L_W, tobe_h, navy_line)
    pill = add_box(slide, L_X + 0.14, tobe_y + 0.14, 1.5, 0.3, fill=NAVY, shape="round")
    set_shape_text(pill, "TO-BE (전환)", S["foot"], F["head"], C["bg"], bold=True)
    for j, row in enumerate(TOBE_ROWS):
        add_text(
            slide,
            L_X + 0.2,
            tobe_y + 0.56 + j * 0.34,
            L_W - 0.4,
            0.28,
            [
                [
                    ("·  ", {"bold": True, "color": C["accent"]}),
                    (row, {"bold": True, "color": C["primary"]}),
                ]
            ],
            S["caption"],
            F["body"],
            C["primary"],
        )

    # ── 중앙 레일 HOW ──
    _band(slide, M_X, RAIL_TOP, M_W, "HOW", "결정적 실행 그래프")
    # 칩 플로우 — 판정만 주황 솔리드
    chip_y, chip_h = RAIL_TOP + 0.72, 0.42
    n = len(FLOW)
    chip_w, arr_w = 0.92, 0.21
    total = n * chip_w + (n - 1) * arr_w
    fx = M_X + (M_W - total) / 2
    for j, name in enumerate(FLOW):
        cxx = fx + j * (chip_w + arr_w)
        if name == "판정":
            chip = add_box(slide, cxx, chip_y, chip_w, chip_h, fill=C["accent"], shape="round")
            set_shape_text(chip, name, S["caption"], F["head"], C["bg"], bold=True)
        else:
            chip = add_box(
                slide,
                cxx,
                chip_y,
                chip_w,
                chip_h,
                fill=C["bg"],
                line=navy_line,
                line_w=1.25,
                shape="round",
            )
            set_shape_text(chip, name, S["caption"], F["head"], NAVY, bold=True)
        if j < n - 1:
            add_text(
                slide,
                cxx + chip_w - 0.02,
                chip_y,
                arr_w + 0.04,
                chip_h,
                "▸",
                S["caption"],
                F["head"],
                NAVY,
                bold=True,
                align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE,
            )
    # 아이콘 카드 2×2
    CGAP = 0.22
    card_w = (M_W - CGAP) / 2
    card_h = 1.72
    top = chip_y + chip_h + 0.28
    for i, (icon, title, num, label) in enumerate(HOW_CARDS):
        cx = M_X + (i % 2) * (card_w + CGAP)
        cy = top + (i // 2) * (card_h + 0.2)
        _card(slide, cx, cy, card_w, card_h, navy_line)
        _icon_chip(slide, cx + 0.14, cy + 0.14, icon)
        add_text(
            slide,
            cx + 0.58,
            cy + 0.14,
            card_w - 0.7,
            0.32,
            title,
            S["caption"],
            F["head"],
            NAVY,
            bold=True,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        add_text(
            slide,
            cx + 0.16,
            cy + 0.56,
            card_w - 0.32,
            0.38,
            [
                [
                    (num, {"bold": True, "color": C["accent"], "size": S["compact"]}),
                    (f"  {label}", {"color": C["muted"], "size": S["foot"]}),
                ]
            ],
            S["compact"],
            F["head"],
            C["accent"],
        )
        add_text(
            slide,
            cx + 0.16,
            cy + 1.02,
            card_w - 0.32,
            card_h - 1.14,
            NODES[i][1],
            S["foot"],
            F["body"],
            body_ink,
            line_spacing=1.2,
        )

    # ── 우 레일 IMPACT ──
    _band(slide, R_X, RAIL_TOP, R_W, "IMPACT", "보장되는 것")
    m_top = RAIL_TOP + 0.7
    m_h = (RAIL_BOTTOM - m_top - 2 * 0.2) / 3
    for j, (num, label) in enumerate(IMPACT_CARDS):
        my = m_top + j * (m_h + 0.2)
        _card(slide, R_X, my, R_W, m_h, navy_line)
        add_text(
            slide,
            R_X + 0.16,
            my + 0.12,
            R_W - 0.32,
            0.42,
            num,
            S["compact"],
            F["head"],
            C["accent"],
            bold=True,
        )
        add_text(
            slide,
            R_X + 0.16,
            my + 0.62,
            R_W - 0.32,
            m_h - 0.74,
            label,
            S["foot"],
            F["body"],
            body_ink,
            line_spacing=1.15,
        )

    # ── 하단 풀폭 네이비 푸터밴드 ──
    fb = add_box(slide, G.MARGIN_L, RAIL_BOTTOM + 0.2, FULL_W, 0.46, fill=NAVY)
    tf = fb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    from pptx.enum.lang import MSO_LANGUAGE_ID
    from pptx.util import Pt

    p = tf.paragraphs[0]
    p._p.get_or_add_pPr().set("eaLnBrk", "0")
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = FOOTER
    r.font.name = F["head"]
    r.font.size = Pt(S["caption"])
    r.font.bold = True
    r.font.color.rgb = C["bg"]
    r.font.language_id = MSO_LANGUAGE_ID.KOREAN

    add_text(slide, G.MARGIN_L, 7.12, FULL_W, 0.26, c["source"], S["foot"], F["body"], C["muted"])
    return slide
