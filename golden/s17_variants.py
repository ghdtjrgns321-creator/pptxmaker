"""S17 시안 — 차별점: 일반 임베딩 RAG 대비 7축 (아키타입: 경쟁 매트릭스 + 다크 요지 사이드바).

모방: golden/ref/s14_matrix.png (McKinsey p26 — 본문 매트릭스 + 우측 다크 패널).
콘텐츠 실물: 00_factsheet.md §G — 7축 전수 + 정직 조항(1115 한정·타 기준서 4건·유연성 없음).
실행: uv run python golden/s17_variants.py → golden/variants/s17_variants.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

import grid as G
from kit import SLIDE_H, SLIDE_W, add_box, add_text, load_kit, mix

K = load_kit()
C, S, F = K["rgb"], K["sizes"], K["fonts"]

KICKER = "5. 차별점"
SOURCE = "출처: 00_factsheet.md §G (6_TEST-DECISIONS.md·4_SEARCH-PIPELINE.md)"

AXES = [  # §G 7축 전수
    ("검색 진입", "임베딩 유사도 — 점수 경쟁", "용어사전 문자 매칭 — 결정적"),
    ("근거 설명", "유사도 점수 (근거 불명)", "그래프 경로 — 어떤 간선을 지났는지"),
    ("판단 절차", "없음 — 생성이 즉석 구성", "원문 앵커 판단트리 41개 사전 조립"),
    ("AI 개입 범위", "검색 전 구간", "용어 색인 1개로 한정"),
    ("오류 방향", "허위 확정 (2종 오류)", "못 찾으면 유보 — 1종 안전"),
    ("데이터 검증", "불가 — 벡터는 검산이 없다", "회계 항등식 전수 검증"),
    ("재현성", "비결정 — 같은 질문, 다른 답", "결정적 — 같은 질문, 같은 경로"),
]


def header(slide, headline):
    add_text(
        slide, G.MARGIN_L, 0.42, 8.0, 0.28, KICKER, S["caption"], F["head"], C["muted"], bold=True
    )
    add_text(
        slide,
        G.MARGIN_L,
        0.72,
        G.RIGHT_EDGE - G.MARGIN_L,
        0.55,
        headline,
        S["section"],
        F["head"],
        C["primary"],
        bold=True,
    )
    add_box(slide, G.MARGIN_L, G.RULE_Y, G.RIGHT_EDGE - G.MARGIN_L, 0.014, fill=C["muted"])


def variant_a(prs):
    """A — 좌 7축 네이티브 매트릭스 + 우 다크 요지 사이드바(McKinsey p26 구도)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    header(slide, "일반 임베딩 RAG과 무엇이 다른가 — 7개 축 전부 같은 방향")
    from pptx.enum.lang import MSO_LANGUAGE_ID

    # ── 좌: 매트릭스 ──
    tbl_x, tbl_y, tbl_w = G.MARGIN_L, 2.0, 8.6
    n_rows = len(AXES) + 1
    gf = slide.shapes.add_table(
        n_rows, 3, Inches(tbl_x), Inches(tbl_y), Inches(tbl_w), Inches(0.5 * n_rows)
    )
    tbl = gf.table
    for ci, cw in enumerate((1.65, 3.3, 3.65)):
        tbl.columns[ci].width = Inches(cw)
    data = [("비교 축", "일반 임베딩 RAG", "이 시스템")] + list(AXES)
    for ri, row in enumerate(data):
        tbl.rows[ri].height = Inches(0.42 if ri == 0 else 0.5)
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            if ri == 0:
                cell.fill.fore_color.rgb = C["primary"]
            elif ci == 2:
                cell.fill.fore_color.rgb = mix(C["bg"], C["accent"], 0.08)
            else:
                cell.fill.fore_color.rgb = C["bg_alt"] if ri % 2 == 0 else C["bg"]
            cell.margin_left = cell.margin_right = Inches(0.1)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p._p.get_or_add_pPr().set("eaLnBrk", "0")
            r = p.add_run()
            r.text = val
            r.font.name = F["head"] if ri == 0 or ci == 0 else F["body"]
            r.font.size = Pt(S["caption"])
            r.font.bold = ri == 0 or ci == 0 or ci == 2
            if ri == 0:
                r.font.color.rgb = C["bg"]
            elif ci == 1:
                r.font.color.rgb = C["muted"]
            else:
                r.font.color.rgb = C["primary"]
            r.font.language_id = MSO_LANGUAGE_ID.KOREAN
    # ── 우: 다크 요지 사이드바 ──
    px = 9.55
    pw = G.RIGHT_EDGE - px  # 3.18
    add_box(slide, px, 1.75, pw, G.CONTENT_BOTTOM - 1.75, fill=C["primary"])
    tx, tw = px + 0.3, pw - 0.6
    add_text(
        slide,
        tx,
        2.15,
        tw,
        1.3,
        [
            [("차이는 성능이 아니라 구조다", {"bold": True, "color": C["bg"]})],
            [
                (
                    "일곱 축 모두 같은 문장으로 요약된다 — 확률 신호가 하던 일을 기준서의 구조가 대신한다.",
                    {},
                )
            ],
        ],
        S["body"],
        F["body"],
        C["bg_alt"],
        line_spacing=1.35,
    )
    add_box(slide, tx, 3.85, tw, 0.014, fill=mix(C["primary"], C["muted"], 0.5))
    add_text(
        slide,
        tx,
        4.1,
        tw,
        1.9,
        [
            [("정직한 한계", {"bold": True, "color": C["bg"]})],
            [
                (
                    "코퍼스는 1115호 한정 — 타 기준서 질문 4건은 재현하지 못했다. "
                    "임베딩식의 유연한 유사 확장도 없다. 대신 그 경계 안에서는 결정적으로 동작한다.",
                    {},
                )
            ],
        ],
        S["caption"],
        F["body"],
        C["bg_alt"],
        line_spacing=1.35,
    )
    # 결론 바 + 출처
    bar = add_box(slide, G.MARGIN_L, G.BAR_Y, G.RIGHT_EDGE - G.MARGIN_L, G.BAR_H, fill=C["primary"])
    tf = bar.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "일반 RAG의 확률 신호 자리마다 기준서의 구조가 들어가 있다 — 7축 전부, 예외 없이"
    run.font.name, run.font.bold = F["head"], True
    run.font.size = Pt(S["body"])
    run.font.color.rgb = C["bg"]
    add_text(
        slide,
        G.MARGIN_L,
        G.SOURCE_Y,
        G.RIGHT_EDGE - G.MARGIN_L,
        0.3,
        SOURCE,
        S["foot"],
        F["body"],
        C["muted"],
    )
    return slide


MIRROR = [  # (축, RAG 키워드, 시스템 키워드) — §G 7축, 키워드 우선
    ("검색 진입", "임베딩 유사도\n점수 경쟁", "용어사전 문자 매칭\n결정적 진입"),
    ("근거 설명", "유사도 점수\n이유를 설명 못 함", "그래프 경로\n지나온 간선이 곧 근거"),
    ("판단 절차", "없음\n생성이 즉석에서 구성", "판단트리 41개\n원문 앵커로 사전 조립"),
    ("AI 개입", "검색 전 구간", "용어 색인 1개로 한정"),
    ("오류 방향", "허위 확정\n2종 오류", "못 찾으면 유보\n1종 안전"),
    ("데이터 검증", "불가\n벡터는 검산이 없다", "회계 항등식\n전수 검증"),
    ("재현성", "같은 질문, 다른 답", "같은 질문, 같은 경로"),
]


def _wing_card(slide, x, y, w, h, text, dashed, align_right):
    from pptx.enum.lang import MSO_LANGUAGE_ID
    from pptx.oxml.ns import qn

    card = add_box(slide, x, y, w, h, fill=C["bg"] if not dashed else C["bg_alt"])
    ln_el = card.line
    ln_el.color.rgb = C["muted"] if dashed else C["primary"]
    ln_el.width = Pt(1.0 if dashed else 1.25)
    if dashed:
        _l = ln_el._get_or_add_ln()
        _l.append(_l.makeelement(qn("a:prstDash"), {"val": "dash"}))
    tf = card.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.12)
    lines = text.split("\n")
    for li, seg in enumerate(lines):
        p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.RIGHT if align_right else PP_ALIGN.LEFT
        p._p.get_or_add_pPr().set("eaLnBrk", "0")
        if li == 0 and not dashed:
            rc = p.add_run()
            rc.text = "✓  "
            rc.font.name, rc.font.bold = F["head"], True
            rc.font.size = Pt(S["caption"])
            rc.font.color.rgb = C["primary"]
        r = p.add_run()
        r.text = seg
        r.font.name = F["head"] if li == 0 else F["body"]
        r.font.bold = li == 0
        r.font.size = Pt(S["caption"])
        r.font.color.rgb = (C["muted"] if dashed else C["primary"]) if li == 0 else C["muted"]
        r.font.language_id = MSO_LANGUAGE_ID.KOREAN
    return card


def variant_b(prs):
    """B — 미러 매트릭스: 중앙 축 기둥, 좌 점선 날개(RAG) vs 우 실선 날개(이 시스템)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    header(slide, "일반 임베딩 RAG과 무엇이 다른가 — 7개 축 전부 같은 방향")
    from kit import set_shape_text

    wing_l_x, wing_w = G.MARGIN_L, 3.15
    spine_x, spine_w = 3.95, 1.3
    wing_r_x = 5.45
    wing_r_w = 3.45
    row0, pitch, card_h = 2.42, 0.55, 0.47
    add_text(
        slide,
        wing_l_x,
        2.02,
        wing_w,
        0.3,
        "일반 임베딩 RAG",
        S["body"],
        F["head"],
        C["muted"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        wing_r_x,
        2.02,
        wing_r_w,
        0.3,
        "이 시스템",
        S["body"],
        F["head"],
        C["primary"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    for i, (axis, rag, ours) in enumerate(MIRROR):
        y = row0 + i * pitch
        _wing_card(slide, wing_l_x, y, wing_w, card_h, rag, dashed=True, align_right=True)
        pill = add_box(
            slide, spine_x, y + (card_h - 0.34) / 2, spine_w, 0.34, fill=C["primary"], shape="round"
        )
        set_shape_text(pill, axis, S["caption"], F["head"], C["bg"], bold=True)
        _wing_card(slide, wing_r_x, y, wing_r_w, card_h, ours, dashed=False, align_right=False)
    # ── 우: 다크 요지 사이드바 ──
    px = 9.2
    pw = G.RIGHT_EDGE - px
    add_box(slide, px, 1.75, pw, G.CONTENT_BOTTOM - 1.75, fill=C["primary"])
    tx, tw = px + 0.3, pw - 0.6
    add_text(
        slide,
        tx,
        2.15,
        tw,
        1.5,
        [
            [("차이는 성능이 아니라\n구조다", {"bold": True, "color": C["bg"]})],
            [
                (
                    "일곱 축 모두 같은 문장으로 요약된다 — 확률 신호가 하던 일을 기준서의 구조가 대신한다.",
                    {},
                )
            ],
        ],
        S["body"],
        F["body"],
        C["bg_alt"],
        line_spacing=1.35,
    )
    add_box(slide, tx, 4.0, tw, 0.014, fill=mix(C["primary"], C["muted"], 0.5))
    add_text(
        slide,
        tx,
        4.25,
        tw,
        1.9,
        [
            [("정직한 한계", {"bold": True, "color": C["bg"]})],
            [
                (
                    "코퍼스는 1115호 한정 — 타 기준서 질문 4건은 재현하지 못했다. "
                    "임베딩식의 유연한 유사 확장도 없다. 대신 그 경계 안에서는 결정적으로 동작한다.",
                    {},
                )
            ],
        ],
        S["caption"],
        F["body"],
        C["bg_alt"],
        line_spacing=1.35,
    )
    bar = add_box(slide, G.MARGIN_L, G.BAR_Y, G.RIGHT_EDGE - G.MARGIN_L, G.BAR_H, fill=C["primary"])
    tf = bar.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "일반 RAG의 확률 신호 자리마다 기준서의 구조가 들어가 있다 — 7축 전부, 예외 없이"
    run.font.name, run.font.bold = F["head"], True
    run.font.size = Pt(S["body"])
    run.font.color.rgb = C["bg"]
    add_text(
        slide,
        G.MARGIN_L,
        G.SOURCE_Y,
        G.RIGHT_EDGE - G.MARGIN_L,
        0.3,
        SOURCE,
        S["foot"],
        F["body"],
        C["muted"],
    )
    return slide


def variant_c(prs):
    """C — 풀폭 미러 매트릭스 (사이드바 제거 — 한계는 S18 전용 장으로 분리)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    header(slide, "일반 임베딩 RAG과 무엇이 다른가 — 7개 축 전부 같은 방향")
    from kit import set_shape_text

    wing_l_x, wing_w = G.MARGIN_L, 5.17
    spine_x, spine_w = 5.97, 1.4
    wing_r_x = 7.57
    wing_r_w = G.RIGHT_EDGE - 7.57
    row0, pitch, card_h = 2.45, 0.55, 0.48
    add_text(
        slide,
        wing_l_x,
        2.02,
        wing_w,
        0.3,
        "일반 임베딩 RAG",
        S["body"],
        F["head"],
        C["muted"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        wing_r_x,
        2.02,
        wing_r_w,
        0.3,
        "이 시스템",
        S["body"],
        F["head"],
        C["primary"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    for i, (axis, rag, ours) in enumerate(MIRROR):
        y = row0 + i * pitch
        _wing_card(slide, wing_l_x, y, wing_w, card_h, rag, dashed=True, align_right=True)
        pill = add_box(
            slide, spine_x, y + (card_h - 0.34) / 2, spine_w, 0.34, fill=C["primary"], shape="round"
        )
        set_shape_text(pill, axis, S["caption"], F["head"], C["bg"], bold=True)
        _wing_card(slide, wing_r_x, y, wing_r_w, card_h, ours, dashed=False, align_right=False)
    bar = add_box(slide, G.MARGIN_L, G.BAR_Y, G.RIGHT_EDGE - G.MARGIN_L, G.BAR_H, fill=C["primary"])
    tf = bar.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "일반 RAG의 확률 신호 자리마다 기준서의 구조가 들어가 있다 — 7축 전부, 예외 없이"
    run.font.name, run.font.bold = F["head"], True
    run.font.size = Pt(S["body"])
    run.font.color.rgb = C["bg"]
    add_text(
        slide,
        G.MARGIN_L,
        G.SOURCE_Y,
        G.RIGHT_EDGE - G.MARGIN_L,
        0.3,
        SOURCE,
        S["foot"],
        F["body"],
        C["muted"],
    )
    return slide


def audit(prs):
    from pptx.oxml.ns import qn

    fails = []
    for si, sl in enumerate(prs.slides):
        accents, tables = 0, 0
        for sh in sl.shapes:
            t, le = sh.top / 914400, sh.left / 914400
            b, r = t + sh.height / 914400, le + sh.width / 914400
            full_bleed = sh.width / 914400 >= SLIDE_W - 0.05
            if t < 6.4 and b > 6.37 and not full_bleed:
                fails.append((si, "bottom", round(b, 2), sh.shape_id))
            if r > 12.75 and not full_bleed:
                fails.append((si, "right", round(r, 2), sh.shape_id))
            if sh.has_table:
                tables += 1
                if len(sh.table.rows) != 8:
                    fails.append((si, "rows≠8", len(sh.table.rows)))
            for el in sh._element.iter(qn("a:srgbClr")):
                if el.get("val") == "D66E3A":
                    if el.getparent().tag.endswith("}solidFill"):
                        accents += 1
        if accents > 4:
            fails.append((si, "accent>4", accents))
        if si == 0 and tables != 1:
            fails.append((si, "table≠1", tables))
        if si in (1, 2):
            EMU = 914400
            wings = sum(
                1
                for sh in sl.shapes
                if abs(sh.height / EMU - 0.47) < 0.015 and sh.width / EMU > 3.0
            )
            pills = sum(
                1
                for sh in sl.shapes
                if abs(sh.height / EMU - 0.34) < 0.005 and 1.25 < sh.width / EMU < 1.45
            )
            if wings != 14 or pills != 7:
                fails.append((si, "mirror", wings, pills))
    assert not fails, f"AUDIT FAIL {fails}"
    print("audit pass")


def main():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    variant_a(prs)
    variant_b(prs)
    variant_c(prs)
    audit(prs)
    out = Path(__file__).parent / "variants" / "s17_variants.pptx"
    prs.save(str(out))
    print(f"saved {out}")


if __name__ == "__main__":
    main()
