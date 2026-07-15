# 시안 K — 2밴드 재설계(사용자 지시 2026-07-14 "섹션 2칸 반반")
# 물성 선언(design-rules P1): 상단=분기(뿌리→엘보→셰브런 밴드) · 하단=대비(동일 지오메트리 쌍, 명도만 차등)
# golden/ 동결 원본 — goldenfab/_variant_k(c=None)와 도형 동일해야 함(compare_golden 기준).
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

import grid as G
from _variant_h import _shape_step
from kit import SLIDE_H, SLIDE_W, add_box, add_text, load_kit, mix

K = load_kit()
C, S, F = K["rgb"], K["sizes"], K["fonts"]

B1_HEAD_Y = G.CONTENT_TOP  # 1.75 상단 밴드 소제목 (룰 2.11)
LANE_LABEL_DY = 0.35  # 라벨 top = row_y - 0.35, h 0.22 → 공기 0.13 (§6 ≥0.12)
LANE_LABEL_H = 0.22
ROW1_Y = 2.60
ROW2_Y = 3.72
FLOW_H = 0.62
QW = 1.45
CONN_W = 0.4
FX = G.MARGIN_L + QW + CONN_W  # 2.45
OVERLAP = FLOW_H * 0.55 / 2
STEP_W = (G.RIGHT_EDGE - FX + 3 * OVERLAP) / 4
Q_CY = (ROW1_Y + ROW2_Y + FLOW_H) / 2  # 3.47 — 뿌리 세로중심 = 두 레인 전체 중심(§5)
QH = 1.3
B2_HEAD_Y = 4.52  # 룰 4.88
PANEL_Y = 5.02
PANEL_H = G.CONTENT_BOTTOM - PANEL_Y  # 1.33
PANEL_GAP = 0.3
PANEL_W = (G.RIGHT_EDGE - G.MARGIN_L - PANEL_GAP) / 2

DEFAULT = {
    "kicker": "1. 문제 정의",
    "headline": "일반 LLM은 틀리는 게 아니라, 틀리는 방향이 문제다",
    "band1_head": "같은 질문, 두 개의 결말",
    "question": '"이 계약, 수익을\n지금 인식해도\n됩니까?"',
    "lanes": [
        {
            "name": "일반 LLM",
            "steps": [
                ("즉시 답변", "근거 조문 없음"),
                ("확신형 어조", "판단 창작"),
                ("검증 불가", "인용 불가"),
            ],
            "verdict": "치명",
        },
        {
            "name": "본 시스템",
            "steps": [
                ("근거 먼저 검색", "문단·사례 인용"),
                ("확정", "결론+인용"),
                ("유보", "정보 되물음"),
            ],
            "verdict": "안전",
        },
    ],
    "band2_head": "오류의 두 방향 — 왜 2종만 치명적인가",
    "panels": [
        {
            "head": "1종 · 놓침 (false negative)",
            "desc": '근거가 존재하는데도 찾지 못해 "모른다"고 물러선다. 사용자는 다른 방법을 찾을 뿐이라 손실은 시간뿐이고, 추가 검색으로 보완된다 — 위험도는 낮다. 실측된 유보 편향(되물음 9건)도 프롬프트 재균형으로 관리하고 있다.',
        },
        {
            "head": "2종 · 허위 확정 (false positive)",
            "desc": "근거가 없거나 틀렸는데 확신에 찬 결론을 준다. 그대로 재무제표에 반영되면 왜곡으로 직결된다 — 회계에서 가장 위험하다. 일반 LLM의 오류는 이 방향으로 표출되며, 이 시스템은 설계 전체를 이 방향의 차단에 편향시켰다.",
        },
    ],
    "bar": "확실할 때만 확정하고, 애매하면 근거를 보여주며 유보한다",
    "source": "출처: 1_OVERVIEW.md · PROJECT_OVERVIEW.md · 5_INTERFACE.md (00_factsheet.md §A·§B)",
}


def _band_head(slide, y, text):
    w = G.RIGHT_EDGE - G.MARGIN_L
    add_text(slide, G.MARGIN_L, y, w, 0.32, text, S["head"], F["head"], C["primary"], bold=True)
    add_box(slide, G.MARGIN_L, y + 0.36, w, 0.012, fill=C["muted"])


def _center(box, text, size, color, bold=True):
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p._p.get_or_add_pPr().set("eaLnBrk", "0")
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name, r.font.bold = F["head"], bold
    r.font.size = Pt(size)
    r.font.color.rgb = color


def variant_k(prs):
    """2밴드 — 상단: 분기(같은 질문 두 결말) · 하단: 대비(1종/2종). 골든 기준(고정 내용)."""
    c = DEFAULT
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    add_text(
        slide,
        G.MARGIN_L,
        0.42,
        6.0,
        0.28,
        c["kicker"],
        S["caption"],
        F["head"],
        C["muted"],
        bold=True,
    )
    add_text(
        slide,
        G.MARGIN_L,
        0.72,
        G.RIGHT_EDGE - G.MARGIN_L,
        0.55,
        c["headline"],
        S["section"],
        F["head"],
        C["primary"],
        bold=True,
    )
    add_box(slide, G.MARGIN_L, G.RULE_Y, G.RIGHT_EDGE - G.MARGIN_L, 0.014, fill=C["muted"])

    # 상단 밴드: 분기
    _band_head(slide, B1_HEAD_Y, c["band1_head"])
    q = add_box(
        slide,
        G.MARGIN_L,
        Q_CY - QH / 2,
        QW,
        QH,
        fill=C["bg_alt"],
        line=C["muted"],
        line_w=0.75,
        shape="round",
    )
    _center(q, c["question"], S["caption"], C["primary"])
    for row_y in (ROW1_Y, ROW2_Y):
        conn = slide.shapes.add_connector(
            MSO_CONNECTOR.ELBOW,
            Inches(G.MARGIN_L + QW),
            Inches(Q_CY),
            Inches(FX),
            Inches(row_y + FLOW_H / 2),
        )
        conn.line.color.rgb = C["muted"]
        conn.line.width = Pt(1.75)
    lane_style = [(ROW1_Y, C["accent"]), (ROW2_Y, C["primary"])]
    for (row_y, ink), lane in zip(lane_style, c["lanes"]):
        add_text(
            slide,
            FX,
            row_y - LANE_LABEL_DY,
            3.0,
            LANE_LABEL_H,
            lane["name"],
            S["caption"],
            F["head"],
            C["primary"],
            bold=True,
        )
        for i, (head, desc) in enumerate(lane["steps"]):
            _shape_step(
                slide,
                FX + i * (STEP_W - OVERLAP),
                row_y,
                STEP_W,
                FLOW_H,
                "pentagon" if i == 0 else "chevron",
                C["bg_alt"],
                head,
                desc,
                C["primary"],
                mix(C["primary"], C["bg"], 0.35),
            )
        _shape_step(
            slide,
            FX + 3 * (STEP_W - OVERLAP),
            row_y,
            STEP_W,
            FLOW_H,
            "chevron",
            ink,
            lane["verdict"],
            "",
            C["bg"],
            C["bg"],
        )

    # 하단 밴드: 대비 쌍
    _band_head(slide, B2_HEAD_Y, c["band2_head"])
    panel_style = [
        (C["bg_alt"], mix(C["muted"], C["bg"], 0.3), C["primary"], C["text"]),
        (C["primary"], C["accent"], C["bg"], C["bg_alt"]),
    ]
    for i, ((fill, spine, head_ink, desc_ink), panel) in enumerate(zip(panel_style, c["panels"])):
        px = G.MARGIN_L + i * (PANEL_W + PANEL_GAP)
        add_box(slide, px, PANEL_Y, PANEL_W, PANEL_H, fill=fill, shape="round")
        add_box(slide, px, PANEL_Y, 0.05, PANEL_H, fill=spine)
        add_text(
            slide,
            px + 0.2,
            PANEL_Y + 0.14,
            PANEL_W - 0.4,
            0.28,
            panel["head"],
            S["body"],
            F["head"],
            head_ink,
            bold=True,
        )
        add_text(
            slide,
            px + 0.2,
            PANEL_Y + 0.5,
            PANEL_W - 0.4,
            PANEL_H - 0.62,
            panel["desc"],
            S["caption"],
            F["body"],
            desc_ink,
            line_spacing=1.3,
        )

    bar = add_box(slide, G.MARGIN_L, G.BAR_Y, G.RIGHT_EDGE - G.MARGIN_L, G.BAR_H, fill=C["primary"])
    _center(bar, c["bar"], S["body"], C["bg"])
    add_text(
        slide,
        G.MARGIN_L,
        G.SOURCE_Y,
        G.RIGHT_EDGE - G.MARGIN_L,
        0.3,
        c["source"],
        S["foot"],
        F["body"],
        C["muted"],
    )
