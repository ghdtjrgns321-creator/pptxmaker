"""S06 파일럿 v3 — 초점 재설계 (2026-07-18). 골든 미편입.

v2 실패 진단: 77단위가 전부 같은 굵기로 경쟁 → 시선이 앉을 자리가 없었고,
골든 도해(팬인·경계)를 좁은 레일·스트립에 미니어처로 구겨 넣어 원형이 죽었다.

v3 원칙 — 주인공 하나, 조연 하나, 나머지는 속삭임:
- 주인공: 파이프라인 플로우(골든 s06 원형 스케일) — 큰 노드·노드 캡션·판정 다이아몬드·OUT 드랍
- 조연: 4노드 카드 2×2 (태그칩·수치칩 상자 제거 — 텍스트로 강등)
- 좌: 팬인 도해에 제 비례(긴 수렴선) 반환. 경계 다이어그램은 경계 장으로 반납(이 장에서 삭제)
- 어두운 존 바는 HOW 하나만. WHY·IMPACT는 연회색 밴드로 강등 — 위계로 초점을 만든다
- accent는 판정·IN/OUT·▼·핵심 수치에만

실행:
    uv run python -c "
    import sys; sys.path.insert(0, '.claude/skills/pptx-build/scripts')
    from goldenfab.reference import new_presentation
    from goldenfab.s06_pilot import build
    prs = new_presentation(); build(prs); prs.save('golden/_pilot_s06/pilot.pptx')"
"""

from pptx.enum.dml import MSO_LINE_DASH_STYLE as MSO_LINE
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from . import grid as G
from .kit import SLIDE_H, SLIDE_W, add_box, add_text, load_kit, mix, set_shape_text

K = load_kit()
C, S, F = K["rgb"], K["sizes"], K["fonts"]

# ── 콘텐츠 (units.md 재고 — 이 파일 밖 문구 금지) ──────────────────────────
ASIS_LABEL = "AS-IS · 점수 기반 선별 (폐기)"
ASIS_FAN = ["임베딩 유사도", "리랭커 점수", "weight_score", "슬롯 상한 3·2·3"]
ASIS_FAN_Q = "근거 선별?"
ASIS_FAN_X = "✕ 세울 수 없다"
ASIS_FAN_CAP = "실측 — 리랭커에 넣은 105건 중 103건 탈락"
TOBE_LABEL = "TO-BE · 결정적 진입 (현재)"
TOBE_CHAIN = [("용어 423", "진입"), ("개념 80", "1홉"), ("문단 250", None)]
TOBE_CAP = "경로는 그래프가 결정 — LLM은 주제 지목만 (via_topic)"
WHY_CAP = "검색 진입부 임베딩 호출 0"

FLOW = [  # (노드, 노드 캡션) — 골든 s06 원형: 캡션은 노드 아래 문서형
    ("질문", None),
    ("analyze", "용어사전 매칭"),
    ("판정", None),
    ("retrieve", "그래프 1홉 탐색"),
    ("generate", "판단트리 주입"),
    ("format", "감리 넛지"),
    ("응답", None),
]
OUT_BOX = "거절 메시지"
OUT_NOTE = "범위 밖 질문은 본선에 진입하지 못하고 즉시 거절된다"

NODE_CARDS = [  # (배지, 노드명, 역할, 처리 3, 수치 텍스트런[(강조, 라벨)])
    (
        "1",
        "analyze",
        "범위·주제 판정",
        [
            "라우팅 IN/OUT + scope guard",
            "멀티턴 → standalone 질의 재구성",
            "쿼리매핑 288 결정적 용어확장",
        ],
        [("임베딩 0", ""), ("288", " 사전"), ("42/42", " calc 판정")],
    ),
    (
        "2",
        "retrieve",
        "결정적 1홉 탐색",
        ["aliases 423 용어사전 진입", "관할 문단 → e3 이웃 → 역인덱스", "케이스·IE via_topic 직결"],
        [("423", " 용어"), ("69→23", " 케이스 정제"), ("3.7→0.6s", "")],
    ),
    (
        "3",
        "generate",
        "듀얼 LLM 생성",
        ["계산 gpt-4.1-mini · 서술 Gemini", "판단트리 41 다중 주입", "인용 교집합 필터 — 실재분만"],
        [("산술 100%", ""), ("0/31", " 트리 오선택"), ("temp 0.0", "")],
    ),
    (
        "4",
        "format",
        "감리 넛지 부착",
        [
            "수집 문서 중 감리사례 1건 선별",
            "상황형(clarify) 넛지 스킵",
            "MongoDB 원문 조회 — 인용분만",
        ],
        [("감리 22건", ""), ("임베딩 매칭 제거", ""), ("넛지 1건", "")],
    ),
]

IMPACT_CELLS = [  # (층, 실측)
    ("결정성", "같은 질문 → 같은 근거 문단"),
    ("재현율", "하드 49.1→59.1% · 소프트 78/92"),
    ("응답 속도", "retrieve 3.7→0.6s · 전체 −22%"),
    ("비용", "건당 약 14원 · 상한 12건"),
]
IMPACT_CAP = "근거 선정에\nLLM 개입 0"

KICKER = "2. 파이프라인"
HEADLINE = "질문에서 응답까지 — 분기까지 전부 결정적인 실행 그래프"
SOURCE = "출처: 2_PIPELINE.md §2.5·2.6 · 5_SEARCH-RUNTIME.md §5.2~5.13 (단위 재고: golden/_pilot_s06/units.md)"

# ── GRID ──────────────────────────────────────────────────────────────────
BAND_Y, BAND_H = 0.92, 0.40
L_X, L_W = G.MARGIN_L, 4.35
H_X = 5.11
H_W = G.RIGHT_EDGE - H_X
ZONE_BOT = 5.42
IMP_Y, IMP_H = 5.56, 1.06
FULL_W = G.RIGHT_EDGE - G.MARGIN_L

line_soft = mix(C["primary"], C["bg"], 0.80)
ink_body = mix(C["muted"], C["primary"], 0.45)


def _band(slide, x, y, w, en, ko, *, dark):
    """존 헤더 밴드 — 주인공 존만 dark, 나머지는 연회색으로 강등."""
    band = add_box(slide, x, y, w, BAND_H, fill=C["primary"] if dark else C["bg_alt"])
    tf = band.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.05)
    from pptx.enum.lang import MSO_LANGUAGE_ID

    p = tf.paragraphs[0]
    p._p.get_or_add_pPr().set("eaLnBrk", "0")
    p.alignment = PP_ALIGN.CENTER
    en_color = mix(C["primary"], C["bg"], 0.55) if dark else C["muted"]
    ko_color = C["bg"] if dark else C["primary"]
    for txt, sz, color in ((en + "   ", S["foot"], en_color), (ko, S["body"] - 1, ko_color)):
        r = p.add_run()
        r.text = txt
        r.font.name = F["head"]
        r.font.size = Pt(sz)
        r.font.bold = True
        r.font.color.rgb = color
        r.font.language_id = MSO_LANGUAGE_ID.KOREAN


def _cap_bar(slide, x, y, w, text):
    bar = add_box(slide, x, y, w, 0.36, fill=C["primary"])
    set_shape_text(bar, text, S["foot"], F["head"], C["bg"], bold=True)


def build(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])

    add_text(slide, G.MARGIN_L, 0.22, 4.0, 0.2, KICKER, S["foot"], F["head"], C["muted"], bold=True)
    add_text(
        slide,
        G.MARGIN_L,
        0.42,
        FULL_W,
        0.38,
        HEADLINE,
        S["compact"],
        F["head"],
        C["primary"],
        bold=True,
    )

    # ══ 좌 존. WHY — 걷어낸 것 → 남긴 것 (조용한 밴드 + 제 비례 팬인) ══
    _band(slide, L_X, BAND_Y, L_W, "WHY", "걷어낸 것 → 남긴 것", dark=False)
    add_box(slide, L_X, 1.44, L_W, 3.50, fill=C["bg"], line=line_soft, line_w=1.0)

    add_text(
        slide,
        L_X + 0.15,
        1.54,
        L_W - 0.3,
        0.2,
        ASIS_LABEL,
        S["foot"],
        F["head"],
        C["muted"],
        bold=True,
    )
    # 팬인 — 긴 수렴선이 본체 (골든 원형 비례)
    fan_ys = [1.86, 2.22, 2.58, 2.94]
    fan_w, fan_h = 1.45, 0.29
    dia_w, dia_h = 1.10, 0.74
    dia_x, dia_cy = 3.42, 2.55
    for fy, txt in zip(fan_ys, ASIS_FAN):
        fb = add_box(slide, L_X + 0.15, fy, fan_w, fan_h, fill=C["bg_alt"])
        set_shape_text(fb, txt, S["foot"], F["head"], ink_body)
        fb.text_frame.word_wrap = False
        conn = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(L_X + 0.15 + fan_w),
            Inches(fy + fan_h / 2),
            Inches(dia_x),
            Inches(dia_cy),
        )
        conn.line.color.rgb = line_soft
        conn.line.width = Pt(1.0)
    dia = add_box(
        slide,
        dia_x,
        dia_cy - dia_h / 2,
        dia_w,
        dia_h,
        fill=C["bg"],
        line=C["muted"],
        line_w=1.2,
        shape="diamond",
    )
    dia.line.dash_style = MSO_LINE.DASH
    set_shape_text(dia, ASIS_FAN_Q, S["foot"], F["head"], C["muted"], bold=True)
    dia.text_frame.word_wrap = False
    add_text(
        slide,
        dia_x - 0.20,
        dia_cy + dia_h / 2 + 0.05,
        dia_w + 0.4,
        0.22,
        ASIS_FAN_X,
        S["caption"],
        F["head"],
        C["primary"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide, L_X + 0.15, 3.36, L_W - 0.3, 0.18, ASIS_FAN_CAP, S["foot"] - 1, F["body"], C["muted"]
    )

    add_text(
        slide,
        L_X,
        3.56,
        L_W,
        0.2,
        "▼",
        S["caption"],
        F["head"],
        C["accent"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    add_text(
        slide,
        L_X + 0.15,
        3.80,
        L_W - 0.3,
        0.2,
        TOBE_LABEL,
        S["foot"],
        F["head"],
        C["primary"],
        bold=True,
    )
    ch_y, ch_w, ch_h = 4.08, 1.10, 0.34
    ch_xs = [L_X + 0.15, L_X + 1.62, L_X + 3.09]
    for (label, arrow_label), bx in zip(TOBE_CHAIN, ch_xs):
        cb = add_box(slide, bx, ch_y, ch_w, ch_h, fill=C["bg"], line=C["primary"], line_w=1.2)
        set_shape_text(cb, label, S["foot"], F["head"], C["primary"], bold=True)
        cb.text_frame.word_wrap = False
        if arrow_label:
            add_text(
                slide,
                bx + ch_w,
                ch_y,
                0.37,
                ch_h,
                "→",
                S["caption"],
                F["head"],
                C["accent"],
                bold=True,
                align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE,
            )
            add_text(
                slide,
                bx + ch_w - 0.05,
                ch_y + ch_h + 0.02,
                0.47,
                0.14,
                arrow_label,
                S["foot"] - 2,
                F["body"],
                C["muted"],
                align=PP_ALIGN.CENTER,
            )
    add_text(slide, L_X + 0.15, 4.64, L_W - 0.3, 0.2, TOBE_CAP, S["foot"], F["body"], ink_body)
    _cap_bar(slide, L_X, ZONE_BOT - 0.36, L_W, WHY_CAP)

    # ══ 우 존. HOW — 주인공: 골든 원형 스케일 플로우 ══
    _band(slide, H_X, BAND_Y, H_W, "HOW", "4노드 결정적 파이프라인", dark=True)

    node_w = [0.80, 1.00, 1.00, 1.00, 1.05, 0.95, 0.80]
    gap = (H_W - sum(node_w)) / (len(FLOW) - 1)
    flow_y, node_h = 1.46, 0.40
    nx = H_X
    dia_cx2 = 0.0
    for j, ((name, cap), w) in enumerate(zip(FLOW, node_w)):
        if name == "판정":
            dia_cx2 = nx + w / 2
            nd = add_box(
                slide,
                nx + (w - 0.92) / 2,
                flow_y - 0.13,
                0.92,
                node_h + 0.26,
                fill=C["bg"],
                line=C["accent"],
                line_w=1.6,
                shape="diamond",
            )
            set_shape_text(nd, name, S["caption"], F["head"], C["accent"], bold=True)
            nd.text_frame.word_wrap = False
        else:
            first, last = j == 0, j == len(FLOW) - 1
            nd = add_box(
                slide,
                nx,
                flow_y,
                w,
                node_h,
                fill=C["primary"] if last else C["bg"],
                line=None if last else (C["primary"] if first else line_soft),
                line_w=1.2,
                shape="round",
            )
            set_shape_text(
                nd, name, S["caption"], F["head"], C["bg"] if last else C["primary"], bold=True
            )
            nd.text_frame.word_wrap = False
        if cap:
            add_text(
                slide,
                nx - 0.15,
                flow_y + node_h + 0.05,
                w + 0.3,
                0.18,
                cap,
                S["foot"],
                F["body"],
                C["muted"],
                align=PP_ALIGN.CENTER,
            )
        if j < len(FLOW) - 1:
            after_dia = name == "판정"
            add_text(
                slide,
                nx + w - 0.02,
                flow_y,
                gap + 0.04,
                node_h,
                "→",
                S["caption"],
                F["head"],
                C["muted"],
                bold=True,
                align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE,
            )
            if after_dia:
                add_text(
                    slide,
                    nx + w - 0.02,
                    flow_y - 0.20,
                    gap + 0.04,
                    0.18,
                    "IN",
                    S["foot"],
                    F["head"],
                    C["accent"],
                    bold=True,
                    align=PP_ALIGN.CENTER,
                )
        nx += w + gap

    # OUT 드랍 — 다이아몬드 아래 거절 종점 + 문서형 주석 (골든 원형)
    drop_y0 = flow_y + node_h + 0.13
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(dia_cx2),
        Inches(drop_y0),
        Inches(dia_cx2),
        Inches(drop_y0 + 0.22),
    )
    conn.line.color.rgb = C["accent"]
    conn.line.width = Pt(1.4)
    add_text(
        slide,
        dia_cx2 + 0.05,
        drop_y0 - 0.02,
        0.5,
        0.18,
        "OUT",
        S["foot"],
        F["head"],
        C["accent"],
        bold=True,
    )
    ob_w = 1.30
    ob = add_box(slide, dia_cx2 - ob_w / 2, drop_y0 + 0.24, ob_w, 0.30, fill=C["bg_alt"])
    set_shape_text(ob, OUT_BOX, S["foot"], F["head"], C["primary"], bold=True)
    add_text(
        slide,
        dia_cx2 + ob_w / 2 + 0.25,
        drop_y0 + 0.24,
        H_X + H_W - dia_cx2 - ob_w / 2 - 0.4,
        0.30,
        OUT_NOTE,
        S["foot"],
        F["body"],
        C["muted"],
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # 조연 — 4노드 카드 2×2 (칩 상자 제거, 텍스트 강등)
    card_w = (H_W - 0.16) / 2
    card_top = 2.65
    card_h = (ZONE_BOT - card_top - 0.14) / 2
    for i, (num, name, role, procs, stats) in enumerate(NODE_CARDS):
        cx = H_X + (i % 2) * (card_w + 0.16)
        cy = card_top + (i // 2) * (card_h + 0.14)
        add_box(slide, cx, cy, card_w, card_h, fill=C["bg"], line=line_soft, line_w=1.0)
        badge = add_box(slide, cx + 0.12, cy + 0.09, 0.24, 0.24, fill=C["primary"], shape="oval")
        set_shape_text(badge, num, S["foot"], F["head"], C["bg"], bold=True)
        add_text(
            slide,
            cx + 0.46,
            cy + 0.07,
            1.6,
            0.28,
            name,
            S["caption"],
            F["head"],
            C["primary"],
            bold=True,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        add_text(
            slide,
            cx + card_w - 1.55,
            cy + 0.07,
            1.4,
            0.28,
            role,
            S["foot"],
            F["body"],
            C["muted"],
            align=PP_ALIGN.RIGHT,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        for j, txt in enumerate(procs):
            add_text(
                slide,
                cx + 0.16,
                cy + 0.40 + j * 0.20,
                card_w - 0.3,
                0.21,
                [[("·  ", {"bold": True, "color": C["accent"]}), (txt, {})]],
                S["foot"],
                F["body"],
                ink_body,
                anchor=MSO_ANCHOR.MIDDLE,
            )
        runs = []
        for k, (em, label) in enumerate(stats):
            if k:
                runs.append(("   ·   ", {"color": mix(C["muted"], C["bg"], 0.3)}))
            runs.append((em, {"bold": True, "color": C["accent"]}))
            if label:
                runs.append((label, {"color": C["muted"], "size": S["foot"] - 1}))
        add_text(
            slide,
            cx + 0.16,
            cy + card_h - 0.30,
            card_w - 0.3,
            0.24,
            [runs],
            S["foot"],
            F["head"],
            C["accent"],
            anchor=MSO_ANCHOR.MIDDLE,
        )

    # ══ 하단 존. IMPACT — 속삭임 스탯 밴드 ══
    add_box(slide, G.MARGIN_L, IMP_Y, FULL_W, IMP_H, fill=C["bg"], line=line_soft, line_w=1.0)
    lead = add_box(slide, G.MARGIN_L + 0.10, IMP_Y + 0.10, 1.50, IMP_H - 0.20, fill=C["bg_alt"])
    tfl = lead.text_frame
    tfl.vertical_anchor = MSO_ANCHOR.MIDDLE
    from pptx.enum.lang import MSO_LANGUAGE_ID

    for i, (txt, sz, color) in enumerate(
        (("IMPACT", S["foot"], C["muted"]), ("보장되는 것", S["caption"], C["primary"]))
    ):
        p = tfl.paragraphs[0] if i == 0 else tfl.add_paragraph()
        p._p.get_or_add_pPr().set("eaLnBrk", "0")
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = txt
        r.font.name = F["head"]
        r.font.size = Pt(sz)
        r.font.bold = True
        r.font.color.rgb = color
        r.font.language_id = MSO_LANGUAGE_ID.KOREAN

    st_x = G.MARGIN_L + 1.86
    st_w = 2.20
    for i, (tier, meas) in enumerate(IMPACT_CELLS):
        x0 = st_x + i * (st_w + 0.06)
        if i:
            add_box(
                slide,
                x0 - 0.03,
                IMP_Y + 0.18,
                0.001,
                IMP_H - 0.36,
                fill=None,
                line=line_soft,
                line_w=0.75,
            )
        add_text(
            slide,
            x0 + 0.12,
            IMP_Y + 0.18,
            st_w - 0.2,
            0.22,
            tier,
            S["caption"],
            F["head"],
            C["primary"],
            bold=True,
        )
        add_text(
            slide,
            x0 + 0.12,
            IMP_Y + 0.46,
            st_w - 0.2,
            IMP_H - 0.6,
            meas,
            S["foot"],
            F["body"],
            ink_body,
            line_spacing=1.15,
        )
    cap_x = st_x + 4 * (st_w + 0.06) + 0.04
    cap = add_box(
        slide, cap_x, IMP_Y + 0.10, G.RIGHT_EDGE - 0.10 - cap_x, IMP_H - 0.20, fill=C["primary"]
    )
    tfc = cap.text_frame
    tfc.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, txt in enumerate(IMPACT_CAP.split("\n")):
        p = tfc.paragraphs[0] if i == 0 else tfc.add_paragraph()
        p._p.get_or_add_pPr().set("eaLnBrk", "0")
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = txt
        r.font.name = F["head"]
        r.font.size = Pt(S["foot"])
        r.font.bold = True
        r.font.color.rgb = C["bg"]
        r.font.language_id = MSO_LANGUAGE_ID.KOREAN

    add_text(slide, G.MARGIN_L, 7.08, FULL_W, 0.22, SOURCE, S["foot"] - 1, F["body"], C["muted"])
    return slide
