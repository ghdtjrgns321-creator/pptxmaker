"""S2 목차 재설계 시안 — 3차(클린): 밴드·카드 제거, 타이포와 여백 중심.

시안 C(델로이트 TOC형): 좌측 타이포 리스트 + 헤어라인 구분 + 우측 다크 패널.
시안 D(BCG Altagamma형): 좌측 다크 컬럼 + 우측 대형 고스트 숫자 2열 타이포 그리드.
실행: uv run python golden/s02_variants.py → golden/variants/s02_variants.pptx
"""

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches

from kit import SLIDE_H, SLIDE_W, add_box, add_text, load_kit, mix

K = load_kit()
C, S, F = K["rgb"], K["sizes"], K["fonts"]
MARGIN = K["layout"]["margin"]
RIGHT_EDGE = SLIDE_W - MARGIN  # 12.733

DECK_TITLE = "K-IFRS 1115 온톨로지 지식그래프 RAG"

# (제목, 서브 1줄, 시작 페이지) — 서브 수치는 00_factsheet.md, 페이지는 05-plan 17장 구성 파생
TOC6 = [
    ("문제 정의", "치명적 오류는 2종(허위 확정) · 일반 LLM이 실패하는 4가지 이유", "03"),
    ("파이프라인", "Analyze → Retrieve → Generate → Format · 4개 결정 노드", "06"),
    ("기술 설명", "그래프 지식베이스 노드 929 · 간선 2,697 · 회계 항등식 전수 검증", "08"),
    ("문제와 해결", "확률 신호의 실패 실측(탈락 103/105) · 전·후 5개 지표 재현 72→78", "11"),
    ("차별점", "일반 임베딩 RAG 대비 7개 축 비교", "13"),
    ("성과·증거", "홀드아웃 78/92 · 에러 0건 · 미재현 원인 분해", "15"),
]

# ── 시안 C GRID ──
C_PANEL_X = 9.6  # 우측 다크 패널
C_LIST_R = C_PANEL_X - 0.5  # 리스트 우측 끝 9.1
C_NUM_X, C_TEXT_X = MARGIN, 1.35  # 번호 축 0.6 / 제목·서브 축 1.35
C_ROW1_Y, C_ROW_PITCH = 2.0, 0.82


def _logo_placeholder(slide):
    """패널 상단 회사 로고 자리 — 저대비 테두리 + 안내 캡션 (실제 로고로 교체하는 슬롯)."""
    ph_c = mix(C["primary"], C["muted"], 0.45)
    x, y, w, h = C_PANEL_X + 0.9, 0.9, (SLIDE_W - C_PANEL_X) - 1.8, 0.9
    add_box(slide, x, y, w, h, fill=None, line=ph_c, line_w=1.0)
    add_text(
        slide,
        x,
        y + h / 2 - 0.15,
        w,
        0.3,
        "회사 로고",
        S["caption"],
        F["body"],
        ph_c,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )


def variant_c(prs):
    """C — 타이포 리스트 + 헤어라인 + 우측 다크 패널 (Deloitte TOC)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    add_box(slide, C_PANEL_X, 0, SLIDE_W - C_PANEL_X, SLIDE_H, fill=C["primary"])
    _logo_placeholder(slide)
    add_text(
        slide,
        C_PANEL_X + 0.35,
        SLIDE_H - 1.15,
        SLIDE_W - C_PANEL_X - 0.7,
        0.7,
        [[(DECK_TITLE, {"bold": True, "color": C["bg"]})], [("기술 제안서 · 2026", {})]],
        S["caption"],
        F["body"],
        C["bg_alt"],
        line_spacing=1.4,
    )
    add_text(slide, MARGIN, 0.5, 4.0, 0.55, "목차", S["title"], F["head"], C["primary"], bold=True)
    add_box(slide, MARGIN, 1.35, C_LIST_R - MARGIN, 0.014, fill=C["muted"])
    for i, (title, subs, page) in enumerate(TOC6):
        y = C_ROW1_Y + i * C_ROW_PITCH
        add_text(
            slide, C_NUM_X, y, 0.6, 0.3, f"0{i + 1}", S["head"], F["head"], C["accent"], bold=True
        )
        add_text(
            slide,
            C_TEXT_X,
            y - 0.02,
            6.0,
            0.32,
            title,
            S["head"],
            F["head"],
            C["primary"],
            bold=True,
        )
        add_text(
            slide,
            C_TEXT_X,
            y + 0.3,
            C_LIST_R - C_TEXT_X - 1.0,
            0.26,
            subs,
            S["caption"],
            F["body"],
            C["muted"],
        )
        add_text(
            slide,
            C_LIST_R - 0.9,
            y,
            0.9,
            0.3,
            f"p.{page}",
            S["body"],
            F["body"],
            C["muted"],
            align=PP_ALIGN.RIGHT,
        )
        if i < 5:
            add_box(
                slide, C_TEXT_X, y + C_ROW_PITCH - 0.17, C_LIST_R - C_TEXT_X, 0.01, fill=C["bg_alt"]
            )
    return slide


# ── 시안 D GRID ──
D_COL_W = 4.2  # 좌측 다크 컬럼 0~4.2
D_GRID_X = 4.9
D_CELL_GAP = 0.5
D_CELL_W = (RIGHT_EDGE - D_GRID_X - D_CELL_GAP) / 2  # 3.667
D_ROW1_Y, D_ROW_PITCH = 1.35, 1.95


def variant_d(prs):
    """D — 좌측 다크 컬럼 + 대형 고스트 숫자 타이포 그리드 (BCG Altagamma)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    add_box(slide, 0, 0, D_COL_W, SLIDE_H, fill=C["primary"])
    add_text(
        slide, MARGIN, 1.35, 2.8, 0.3, "CONTENTS", S["caption"], F["head"], C["muted"], bold=True
    )
    add_text(slide, MARGIN, 1.75, 2.8, 0.7, "목차", S["title"], F["head"], C["bg"], bold=True)
    add_box(slide, MARGIN, 2.75, 0.5, 0.035, fill=C["accent"])
    add_text(
        slide,
        MARGIN,
        3.1,
        2.9,
        1.0,
        "문제 정의에서 성과·증거까지, 여섯 부가 하나의 논증으로 이어집니다.",
        S["sub"],
        F["body"],
        C["bg_alt"],
        line_spacing=1.4,
    )
    add_text(
        slide,
        MARGIN,
        SLIDE_H - 1.0,
        2.9,
        0.3,
        "기술 제안서 · 2026",
        S["caption"],
        F["body"],
        C["muted"],
    )
    ghost = mix(C["bg"], C["muted"], 0.45)
    for i, (title, _subs, page) in enumerate(TOC6):
        col, row = i % 2, i // 2
        x = D_GRID_X + col * (D_CELL_W + D_CELL_GAP)
        y = D_ROW1_Y + row * D_ROW_PITCH
        add_text(slide, x, y, 1.6, 0.85, f"0{i + 1}", S["display"], F["head"], ghost, bold=True)
        add_text(
            slide,
            x,
            y + 0.85,
            D_CELL_W,
            0.34,
            title,
            S["head"] + 2,
            F["head"],
            C["primary"],
            bold=True,
        )
        add_text(
            slide,
            x,
            y + 1.24,
            D_CELL_W,
            0.26,
            f"p.{page}",
            S["caption"],
            F["body"],
            C["muted"],
        )
    return slide


def audit(prs):
    inch = lambda e: round(Emu(e).inches, 3)
    counts = [len(list(sl.shapes)) for sl in prs.slides]
    # 시안 C: 번호 축 0.6 / 제목·서브 축 1.35 정렬 검증
    c = prs.slides[0]
    lefts = [
        inch(sh.left)
        for sh in c.shapes
        if sh.has_text_frame and sh.text_frame.text[:1].isdigit() and inch(sh.left) < 1.0
    ]
    assert all(x == MARGIN for x in lefts), f"C 번호 축 이탈 {lefts}"
    # 시안 D: 그리드 셀 좌축 2열 검증
    d = prs.slides[1]
    cell_lefts = {
        inch(sh.left) for sh in d.shapes if sh.has_text_frame and inch(sh.left) >= D_GRID_X - 0.001
    }
    expect = (D_GRID_X, D_GRID_X + D_CELL_W + D_CELL_GAP)
    stray = [x for x in cell_lefts if all(abs(x - e) > 0.002 for e in expect)]
    assert not stray, f"D 셀 축 이탈 {stray}"
    print(f"audit pass — shapes C={counts[0]}, D={counts[1]}")


def main():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    variant_c(prs)
    variant_d(prs)
    audit(prs)
    from pathlib import Path

    out = Path(__file__).parent / "variants" / "s02_variants.pptx"
    prs.save(str(out))
    print(f"saved {out}")


if __name__ == "__main__":
    main()
