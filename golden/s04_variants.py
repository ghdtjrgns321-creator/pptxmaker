"""S4 시안 경쟁 — A: 좌우 패널 대비(옛 덱 S4 개선) / B: 한 질문의 두 경로(스토리 flow).

재료: 00_factsheet.md §A(2종 오류 프레임)·§B(일반 LLM 결격 사유). 신규 수치 금지.
실행: uv run python s04_variants.py → scratchpad/s04_variants.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from build_golden import content_header, source_line
from kit import SLIDE_H, SLIDE_W, add_box, add_text, load_kit, set_shape_text

K = load_kit()
C, S, F = K["rgb"], K["sizes"], K["fonts"]
MARGIN = K["layout"]["margin"]

SCRATCH = Path(
    r"C:\Users\ghdtj\AppData\Local\Temp\claude\C--Users-ghdtj-workspace-portfolio-pptmaker\64942c94-5596-4938-84ac-0442d4d08c39\scratchpad"
)


def variant_a(prs):
    """A안 — 좌(1종·안전, 흰 패널) vs 우(2종·치명, 다크 패널). 옛 덱 S4 구도의 골든판."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    content_header(slide, 1, "문제 정의", "회계 챗봇의 치명적 오류는 1종이 아니라 2종이다")
    add_text(
        slide,
        MARGIN,
        1.72,
        SLIDE_W - 2 * MARGIN,
        0.35,
        "일반 LLM은 그럴듯한 답을 즉시 내놓지만 근거 조문을 확인할 수 없고, 기준서에 없는 판단을 지어낸다 — 오류가 치명적인 방향으로 표출된다",
        S["sub"],
        F["body"],
        C["text"],
    )

    pw = (SLIDE_W - 2 * MARGIN - 0.3) / 2
    py, ph = 2.35, 3.3
    # 좌: 1종 — 안전한 실패 (연회색 패널)
    add_box(slide, MARGIN, py, pw, ph, fill=C["bg_alt"])
    add_text(
        slide,
        MARGIN + 0.3,
        py + 0.22,
        pw - 0.6,
        0.4,
        "1종 오류 — 놓침 (수용 가능한 실패)",
        S["head"],
        F["head"],
        C["primary"],
        bold=True,
    )
    rows_a = [
        ("안전한 실패", '근거가 있는데 못 찾아 "모른다"고 답함 — 불편하지만 추가 검토로 해소'),
        (
            "측정·교정의 대상",
            "유보 편향(사실관계 충분한데 되물음 9건)도 방치가 아니라 프롬프트 재균형으로 관리",
        ),
    ]
    for i, (head, desc) in enumerate(rows_a):
        ry = py + 0.85 + i * 1.1
        add_box(slide, MARGIN + 0.3, ry + 0.05, 0.05, 0.8, fill=C["muted"])
        add_text(
            slide,
            MARGIN + 0.55,
            ry,
            pw - 0.9,
            0.3,
            f"{head}:",
            S["body"],
            F["head"],
            C["primary"],
            bold=True,
        )
        add_text(
            slide,
            MARGIN + 0.55,
            ry + 0.32,
            pw - 0.9,
            0.7,
            desc,
            S["body"],
            F["body"],
            C["text"],
            line_spacing=1.15,
        )
    # 우: 2종 — 치명 (다크 패널)
    rx = MARGIN + pw + 0.3
    add_box(slide, rx, py, pw, ph, fill=C["primary"])
    add_text(
        slide,
        rx + 0.3,
        py + 0.22,
        pw - 0.6,
        0.4,
        "2종 오류 — 허위 확정 (치명적 실패)",
        S["head"],
        F["head"],
        C["bg"],
        bold=True,
    )
    rows_b = [
        ("실제 손해", "근거가 없거나 틀렸는데 확신에 찬 답 — 부실감사·재무제표 왜곡으로 직결"),
        ("창작 리스크", "일반 LLM은 언어적 매끄러움에 집중해 실재하지 않는 기준을 확정적으로 제시"),
        ("정직한 목표", '"환각 0" 과장 대신, 근거 없는 확정을 구조적으로 어렵게 만드는 것'),
    ]
    for i, (head, desc) in enumerate(rows_b):
        ry = py + 0.8 + i * 0.82
        add_box(slide, rx + 0.3, ry + 0.05, 0.05, 0.6, fill=C["accent"])
        add_text(
            slide,
            rx + 0.55,
            ry,
            pw - 0.9,
            0.3,
            f"{head}:",
            S["body"],
            F["head"],
            C["bg"],
            bold=True,
        )
        add_text(
            slide,
            rx + 0.55,
            ry + 0.3,
            pw - 0.9,
            0.5,
            desc,
            S["caption"],
            F["body"],
            C["bg_alt"],
            line_spacing=1.15,
        )
    band = add_box(slide, MARGIN, py + ph + 0.25, SLIDE_W - 2 * MARGIN, 0.5, fill=C["accent"])
    set_shape_text(
        band,
        "확실할 때만 확정하고, 애매하면 근거를 보여주며 유보한다 — 설계 전체를 2종 차단 방향으로 편향",
        S["sub"],
        F["head"],
        C["bg"],
        bold=True,
    )
    source_line(slide, "출처: 1_OVERVIEW.md · PROJECT_OVERVIEW.md (00_factsheet.md §A·§B)")


def variant_b(prs):
    """B안 — 같은 질문 하나가 두 시스템에서 다른 결말로: 2레인 스토리 flow."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    content_header(
        slide, 1, "문제 정의", "같은 질문, 두 개의 결말 — 오류가 표출되는 방향이 운명을 가른다"
    )

    # 공통 질문 카드
    qw = 2.5
    qy = 3.35
    q = add_box(slide, MARGIN, qy, qw, 1.15, fill=C["bg_alt"], line=C["muted"], line_w=0.75)
    tf = q.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = '"이 계약, 수익을 지금 인식해도 됩니까?"'
    r.font.name, r.font.bold = F["head"], True
    from pptx.util import Pt

    r.font.size = Pt(S["body"])
    r.font.color.rgb = C["primary"]

    lanes = [
        (
            "일반 LLM",
            C["accent"],
            2.2,
            [
                ("즉시 답변", "근거 조문 표시 없음"),
                ("확신에 찬 어조", "기준서에 없는 판단도 창작"),
                ("검증 불가", "어느 문단에 기반했는지 확인 불가"),
            ],
            "허위 확정 → 부실감사·재무제표 왜곡",
            True,
        ),
        (
            "본 시스템",
            C["primary"],
            4.85,
            [
                ("근거 먼저 검색", "기준서 문단·사례 인용"),
                ("확실하면 확정", "단정 결론 + 인용 문단"),
                ("애매하면 유보", "부족한 정보를 근거와 함께 질문"),
            ],
            "안전한 실패 → 추가 검토로 해소 (되물음 9건도 재균형 관리)",
            False,
        ),
    ]
    sx, sw = MARGIN + qw + 0.45, 6.0
    step_w = (sw - 0.4) / 3
    for name, ink, ly, steps, outcome, fatal in lanes:
        add_text(slide, sx, ly - 0.38, 2.5, 0.3, name, S["head"], F["head"], ink, bold=True)
        for i, (head, desc) in enumerate(steps):
            x = sx + i * (step_w + 0.2)
            add_box(
                slide,
                x,
                ly,
                step_w,
                0.95,
                fill=C["bg_alt"] if not fatal else C["bg"],
                line=ink,
                line_w=1.0,
            )
            add_text(
                slide,
                x + 0.12,
                ly + 0.1,
                step_w - 0.24,
                0.3,
                head,
                S["body"],
                F["head"],
                C["primary"],
                bold=True,
            )
            add_text(
                slide,
                x + 0.12,
                ly + 0.42,
                step_w - 0.24,
                0.5,
                desc,
                S["caption"],
                F["body"],
                C["muted"],
                line_spacing=1.1,
            )
            if i < 2:
                add_text(
                    slide,
                    x + step_w + 0.01,
                    ly + 0.28,
                    0.2,
                    0.4,
                    "→",
                    S["head"],
                    F["head"],
                    C["muted"],
                    bold=True,
                )
        # 결말 칩
        chip = add_box(
            slide,
            sx + sw + 0.25,
            ly + 0.06,
            1.55,
            0.85,
            fill=C["accent"] if fatal else C["primary"],
        )
        set_shape_text(chip, "치명" if fatal else "안전", S["head"], F["head"], C["bg"], bold=True)
        add_text(
            slide,
            sx,
            ly + 1.05,
            sw + 1.8,
            0.3,
            outcome,
            S["caption"],
            F["body"],
            C["accent"] if fatal else C["muted"],
            bold=fatal,
        )
    # 질문 → 레인 연결선
    add_box(slide, MARGIN + qw, 2.7, 0.02, 2.4, fill=C["muted"])

    band = add_box(slide, MARGIN, 6.15, SLIDE_W - 2 * MARGIN, 0.5, fill=C["primary"])
    set_shape_text(
        band,
        "회계에서 1종(놓침)은 검토로 해소되지만, 2종(허위 확정)은 부실감사로 직결된다",
        S["sub"],
        F["head"],
        C["bg"],
        bold=True,
    )
    source_line(
        slide, "출처: 1_OVERVIEW.md · PROJECT_OVERVIEW.md · 5_INTERFACE.md (00_factsheet.md §A·§B)"
    )


def variant_c(prs):
    """C안 — B(스토리 flow 주 구획) + A(1종/2종 대비 보조 구획) 결합 고밀도판."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    content_header(
        slide, 1, "문제 정의", "같은 질문, 두 개의 결말 — 일반 LLM은 치명적인 방향으로 틀린다"
    )

    # ── 구획 1: 2레인 스토리 flow (주 시각) ──
    qw = 2.3
    q = add_box(slide, MARGIN, 2.55, qw, 1.0, fill=C["bg_alt"], line=C["muted"], line_w=0.75)
    set_shape_text(
        q, '"이 계약, 수익을 지금\n인식해도 됩니까?"', S["body"], F["head"], C["primary"], bold=True
    )
    add_box(slide, MARGIN + qw, 2.25, 0.02, 1.7, fill=C["muted"])

    lanes = [
        (
            "일반 LLM",
            C["accent"],
            2.02,
            [
                ("즉시 답변", "근거 조문 표시 없음"),
                ("확신에 찬 어조", "기준서 밖 판단도 창작"),
                ("검증 불가", "기반 문단 확인 불가"),
            ],
            "치명",
            True,
        ),
        (
            "본 시스템",
            C["primary"],
            3.42,
            [
                ("근거 먼저 검색", "기준서 문단·사례 인용"),
                ("확실하면 확정", "단정 결론 + 인용 문단"),
                ("애매하면 유보", "부족한 정보를 질문"),
            ],
            "안전",
            False,
        ),
    ]
    sx = MARGIN + qw + 0.4
    step_w = 2.15
    for name, ink, ly, steps, verdict, fatal in lanes:
        add_text(slide, sx, ly - 0.32, 2.5, 0.28, name, S["body"], F["head"], ink, bold=True)
        for i, (head, desc) in enumerate(steps):
            x = sx + i * (step_w + 0.22)
            add_box(
                slide,
                x,
                ly,
                step_w,
                0.78,
                fill=C["bg"] if fatal else C["bg_alt"],
                line=ink,
                line_w=1.0,
            )
            add_text(
                slide,
                x + 0.1,
                ly + 0.06,
                step_w - 0.2,
                0.28,
                head,
                S["body"],
                F["head"],
                C["primary"],
                bold=True,
            )
            add_text(
                slide,
                x + 0.1,
                ly + 0.38,
                step_w - 0.2,
                0.35,
                desc,
                S["caption"],
                F["body"],
                C["muted"],
            )
            if i < 2:
                add_text(
                    slide,
                    x + step_w - 0.01,
                    ly + 0.2,
                    0.25,
                    0.35,
                    "→",
                    S["head"],
                    F["head"],
                    C["muted"],
                    bold=True,
                    align=PP_ALIGN.CENTER,
                )
        chip = add_box(
            slide,
            sx + 3 * step_w + 2 * 0.22 + 0.15,
            ly,
            1.15,
            0.78,
            fill=C["accent"] if fatal else C["primary"],
        )
        set_shape_text(chip, verdict, S["head"], F["head"], C["bg"], bold=True)

    # ── 구획 2: 1종/2종 대비 미니 패널 (보조 구획) ──
    py, ph = 4.5, 1.3
    pw = (SLIDE_W - 2 * MARGIN - 0.3) / 2
    add_box(slide, MARGIN, py, pw, ph, fill=C["bg_alt"])
    add_text(
        slide,
        MARGIN + 0.25,
        py + 0.12,
        pw - 0.5,
        0.3,
        "1종 오류 — 놓침 (안전한 실패)",
        S["body"],
        F["head"],
        C["primary"],
        bold=True,
    )
    add_text(
        slide,
        MARGIN + 0.25,
        py + 0.48,
        pw - 0.5,
        0.75,
        '근거가 있는데 못 찾아 "모른다"고 답함 — 추가 검토로 해소. 유보 편향(되물음 9건 실측)도 프롬프트 재균형으로 관리한다.',
        S["caption"],
        F["body"],
        C["text"],
        line_spacing=1.2,
    )
    rx = MARGIN + pw + 0.3
    add_box(slide, rx, py, pw, ph, fill=C["primary"])
    add_text(
        slide,
        rx + 0.25,
        py + 0.12,
        pw - 0.5,
        0.3,
        "2종 오류 — 허위 확정 (치명적 실패)",
        S["body"],
        F["head"],
        C["bg"],
        bold=True,
    )
    add_text(
        slide,
        rx + 0.25,
        py + 0.48,
        pw - 0.5,
        0.75,
        "근거가 없거나 틀렸는데 확신에 찬 답 — 부실감사·재무제표 왜곡으로 직결. 일반 LLM의 오류는 이 방향으로 표출된다.",
        S["caption"],
        F["body"],
        C["bg_alt"],
        line_spacing=1.2,
    )

    band = add_box(slide, MARGIN, py + ph + 0.18, SLIDE_W - 2 * MARGIN, 0.48, fill=C["accent"])
    set_shape_text(
        band,
        "확실할 때만 확정하고, 애매하면 근거를 보여주며 유보한다 — 설계 전체를 2종 차단 방향으로 편향",
        S["sub"],
        F["head"],
        C["bg"],
        bold=True,
    )
    source_line(
        slide, "출처: 1_OVERVIEW.md · PROJECT_OVERVIEW.md · 5_INTERFACE.md (00_factsheet.md §A·§B)"
    )


def zone_title(slide, x, y, w, text):
    """구획 소제목 + 수평 룰 — mck_p128 문법."""
    add_text(slide, x, y, w, 0.28, text, S["body"], F["head"], C["accent"], bold=True)
    add_box(slide, x, y + 0.32, w, 0.012, fill=C["muted"])


def lead_item(slide, x, y, w, head, desc, ink_head, ink_desc, bar=None, desc_h=0.75):
    """색 세로바 + 볼드 리드인 + 완전 문장 — mck_p128 항목 문법."""
    add_box(slide, x, y + 0.04, 0.045, desc_h + 0.3, fill=bar or C["accent"])
    add_text(slide, x + 0.2, y, w - 0.2, 0.3, head, S["body"], F["head"], ink_head, bold=True)
    add_text(
        slide,
        x + 0.2,
        y + 0.32,
        w - 0.2,
        desc_h,
        desc,
        S["caption"],
        F["body"],
        ink_desc,
        line_spacing=1.2,
    )


def variant_d(prs):
    """D안(v5) — C 구도 + 레퍼런스 구획 문법: 소제목·룰·완전 문장·좌측 앵커·세로 헤어라인."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    content_header(
        slide, 1, "문제 정의", "같은 질문, 두 개의 결말 — 일반 LLM은 치명적인 방향으로 틀린다"
    )

    # ── 좌측 앵커 칼럼 ──
    ax, aw = MARGIN, 3.0
    q = add_box(slide, ax, 1.95, aw, 0.95, fill=C["bg_alt"], line=C["muted"], line_w=0.75)
    set_shape_text(
        q, '"이 계약, 수익을 지금\n인식해도 됩니까?"', S["body"], F["head"], C["primary"], bold=True
    )
    lead_item(
        slide,
        ax,
        3.15,
        aw,
        "왜 위험한 질문인가:",
        "수익인식(K-IFRS 1115)은 답이 틀리면 재무제표가 왜곡되는 영역이다 — 챗봇의 오답이 부실감사라는 실제 손해로 이어진다.",
        C["primary"],
        C["text"],
        desc_h=0.95,
    )
    lead_item(
        slide,
        ax,
        4.55,
        aw,
        "일반 LLM의 결격:",
        "그럴듯한 답을 즉시 내놓지만 어느 문단에 기반했는지 확인할 수 없고, 기준서에 없는 판단을 확신에 찬 어조로 지어낸다.",
        C["primary"],
        C["text"],
        desc_h=0.95,
    )
    add_box(slide, ax + aw + 0.35, 1.95, 0.014, 4.35, fill=C["bg_alt"])  # 세로 헤어라인

    # ── 우측: 경로 구획 2개 ──
    rx = ax + aw + 0.7
    rw = SLIDE_W - MARGIN - rx
    lanes = [
        (
            1.95,
            "일반 LLM의 응답 경로 — 오류가 치명적인 방향으로 표출된다",
            C["accent"],
            True,
            [
                ("즉시 답변", "근거 조문 없이 결론부터 제시한다"),
                ("확신에 찬 어조", "기준서에 없는 판단도 확정적으로 창작한다"),
                ("검증 불가", "감사 조서에 인용할 근거가 남지 않는다"),
            ],
            "치명",
        ),
        (
            3.75,
            "본 시스템의 응답 경로 — 틀려도 안전한 방향으로 무너진다",
            C["primary"],
            False,
            [
                ("근거 먼저 검색", "기준서 문단·사례를 결정적으로 수집한다"),
                ("확실하면 확정", "단정 결론에 인용 문단을 붙여 답한다"),
                ("애매하면 유보", "부족한 정보를 근거와 함께 되묻는다"),
            ],
            "안전",
        ),
    ]
    for zy, title, ink, fatal, steps, verdict in lanes:
        zone_title(slide, rx, zy, rw, title)
        step_w = (rw - 1.05 - 2 * 0.42) / 3
        ly = zy + 0.5
        for i, (head, desc) in enumerate(steps):
            x = rx + i * (step_w + 0.42)
            add_box(
                slide,
                x,
                ly,
                step_w,
                0.95,
                fill=C["bg"] if fatal else C["bg_alt"],
                line=ink,
                line_w=1.0,
            )
            add_text(
                slide,
                x + 0.12,
                ly + 0.08,
                step_w - 0.24,
                0.28,
                head,
                S["body"],
                F["head"],
                C["primary"],
                bold=True,
            )
            add_text(
                slide,
                x + 0.12,
                ly + 0.4,
                step_w - 0.24,
                0.5,
                desc,
                S["caption"],
                F["body"],
                C["muted"],
                line_spacing=1.1,
            )
            if i < 2:
                add_text(
                    slide,
                    x + step_w + 0.05,
                    ly + 0.3,
                    0.35,
                    0.35,
                    "→",
                    S["head"],
                    F["head"],
                    C["muted"],
                    bold=True,
                )
        chip = add_box(slide, rx + rw - 0.95, ly, 0.95, 0.95, fill=ink)
        set_shape_text(chip, verdict, S["head"], F["head"], C["bg"], bold=True)

    # ── 하단 구획: 오류의 두 방향 (정의) ──
    zy = 5.55
    zone_title(slide, rx, zy, rw, "오류의 두 방향 — 왜 2종만 치명적인가")
    half = (rw - 0.5) / 2
    lead_item(
        slide,
        rx,
        zy + 0.45,
        half,
        "1종 · 놓침 (안전한 실패):",
        '근거가 있는데 못 찾아 "모른다"고 답함 — 추가 검토로 해소되며, 유보 편향(되물음 9건 실측)도 재균형으로 관리한다.',
        C["primary"],
        C["text"],
        bar=C["muted"],
        desc_h=0.7,
    )
    lead_item(
        slide,
        rx + half + 0.5,
        zy + 0.45,
        half,
        "2종 · 허위 확정 (치명적 실패):",
        "근거가 없거나 틀렸는데 확신에 찬 답 — 부실감사·재무제표 왜곡으로 직결된다. 설계 전체를 이 방향 차단에 편향시켰다.",
        C["accent"],
        C["text"],
        desc_h=0.7,
    )

    source_line(
        slide, "출처: 1_OVERVIEW.md · PROJECT_OVERVIEW.md · 5_INTERFACE.md (00_factsheet.md §A·§B)"
    )


ICONS = Path(__file__).resolve().parents[1] / ".claude/skills/pptx-build/assets/icons"


def icon_item(slide, x, y, w, icon, head, desc, head_ink, desc_ink, desc_h=0.62):
    """아이콘 원 + 볼드 리드인 + 완전 문장 — mck_p128 항목 문법 (아이콘판)."""
    slide.shapes.add_picture(
        str(ICONS / icon), Inches(x), Inches(y + 0.02), Inches(0.4), Inches(0.4)
    )
    add_text(slide, x + 0.58, y, w - 0.58, 0.3, head, S["body"], F["head"], head_ink, bold=True)
    add_text(
        slide,
        x + 0.58,
        y + 0.32,
        w - 0.58,
        desc_h,
        desc,
        S["caption"],
        F["body"],
        desc_ink,
        line_spacing=1.2,
    )


def variant_e(prs):
    """E안(v6) — mck_p128 레이아웃 1:1 이식: 좌 대형 앵커 + 세로 룰 + 우 2구획(아이콘·문장)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    content_header(
        slide,
        1,
        "문제 정의",
        "일반 LLM을 회계에 못 쓰는 이유는 틀려서가 아니라, 틀리는 방향 때문이다",
    )

    # ── 좌측 앵커: 대형 용어 + 정의 + 볼드 문단 (mck 좌측 스탯 칼럼) ──
    ax, aw = MARGIN, 3.15
    add_text(slide, ax, 1.95, aw, 0.75, "2종 오류", S["display"], F["head"], C["accent"], bold=True)
    add_text(
        slide, ax, 2.7, aw, 0.55, "= 허위 확정", S["title"], F["head"], C["primary"], bold=True
    )
    add_text(
        slide,
        ax,
        3.35,
        aw,
        0.85,
        "근거가 없거나 틀렸는데 확신에 찬 답을 내놓는 오류 — 부실감사·재무제표 왜곡으로 직결된다.",
        S["body"],
        F["body"],
        C["text"],
        line_spacing=1.2,
    )
    add_text(
        slide,
        ax,
        4.45,
        aw,
        0.35,
        "1종(놓침)은 안전한 실패다",
        S["head"],
        F["head"],
        C["primary"],
        bold=True,
    )
    add_text(
        slide,
        ax,
        4.85,
        aw,
        1.1,
        '근거가 있는데 못 찾아 "모른다"고 답하는 것은 추가 검토로 해소된다. 유보 편향(되물음 9건 실측)도 재균형으로 관리한다 — 치명적인 것은 2종뿐이다.',
        S["body"],
        F["body"],
        C["text"],
        line_spacing=1.2,
    )
    add_box(slide, ax + aw + 0.3, 1.95, 0.014, 4.1, fill=C["muted"])  # 세로 룰

    # ── 우측 구획 1: 일반 LLM의 결격 (아이콘 2열) ──
    rx = ax + aw + 0.65
    rw = SLIDE_W - MARGIN - rx
    zone_title(slide, rx, 1.95, rw, "일반 LLM의 결격 — 왜 그대로 쓸 수 없는가")
    half = (rw - 0.4) / 2
    icon_item(
        slide,
        rx,
        2.5,
        half,
        "circle-x_accent.png",
        "근거 조문을 확인할 수 없다",
        "그럴듯한 답을 즉시 내놓지만 어느 문단에 기반했는지 검증 불가 — 감사 조서에 인용할 수 없다.",
        C["primary"],
        C["text"],
        desc_h=0.75,
    )
    icon_item(
        slide,
        rx + half + 0.4,
        2.5,
        half,
        "triangle-alert_accent.png",
        "기준서에 없는 판단을 창작한다",
        "언어적 매끄러움에 집중해 실재하지 않는 기준을 확신에 찬 어조로 제시한다.",
        C["primary"],
        C["text"],
        desc_h=0.75,
    )

    # ── 우측 구획 2: 이 시스템의 응답 원칙 (아이콘 2×2) ──
    zone_title(slide, rx, 3.75, rw, "이 시스템의 응답 원칙 — 확신의 수준이 구조로 드러난다")
    items = [
        (
            "circle-check-big_primary.png",
            "확정 (TYPE2)",
            "결정적 사실이 충분할 때만 — 단정 결론 + 인용 문단",
        ),
        ("git-branch_primary.png", "조건부 (TYPE1)", "사실관계가 갈릴 때 — Case 1/2 분기 결론"),
        (
            "search_primary.png",
            "되물음 (clarify)",
            "사실관계 부족 시 — 부족한 정보를 근거와 함께 질문",
        ),
        ("cpu_primary.png", "산술 (calc)", "계산 필요 판정 시 — 계산 전용 모델로 라우팅"),
    ]
    for i, (icon, head, desc) in enumerate(items):
        x = rx + (i % 2) * (half + 0.4)
        y = 4.3 + (i // 2) * 1.0
        icon_item(slide, x, y, half, icon, head, desc, C["primary"], C["text"], desc_h=0.55)

    band = add_box(slide, MARGIN, 6.35, SLIDE_W - 2 * MARGIN, 0.48, fill=C["primary"])
    set_shape_text(
        band,
        "확실할 때만 확정하고, 애매하면 근거를 보여주며 유보한다 — 설계 전체를 2종 차단 방향으로 편향",
        S["sub"],
        F["head"],
        C["bg"],
        bold=True,
    )
    source_line(
        slide,
        "출처: 1_OVERVIEW.md · PROJECT_OVERVIEW.md · 5_INTERFACE.md (00_factsheet.md §A·§B·§B')",
    )


def variant_f(prs):
    """F안(v7) — E의 구분 골격 + C의 2레인 flow + A의 1종/2종 대비 패널 합성."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    content_header(
        slide,
        1,
        "문제 정의",
        "일반 LLM을 회계에 못 쓰는 이유는 틀려서가 아니라, 틀리는 방향 때문이다",
    )

    # ── 좌측 앵커: 대형 용어 + 정의 (E 계승) ──
    ax, aw = MARGIN, 2.95
    add_text(slide, ax, 1.95, aw, 0.72, "2종 오류", S["display"], F["head"], C["accent"], bold=True)
    add_text(
        slide, ax, 2.67, aw, 0.5, "= 허위 확정", S["title"], F["head"], C["primary"], bold=True
    )
    add_text(
        slide,
        ax,
        3.3,
        aw,
        1.0,
        "근거가 없거나 틀렸는데 확신에 찬 답을 내놓는 오류 — 부실감사·재무제표 왜곡이라는 실제 손해로 직결된다.",
        S["body"],
        F["body"],
        C["text"],
        line_spacing=1.2,
    )
    add_text(
        slide, ax, 4.5, aw, 0.35, "일반 LLM의 결격:", S["head"], F["head"], C["primary"], bold=True
    )
    add_text(
        slide,
        ax,
        4.9,
        aw,
        1.2,
        "근거 조문을 확인할 수 없고, 기준서에 없는 판단을 확신에 찬 어조로 창작한다.",
        S["body"],
        F["body"],
        C["text"],
        line_spacing=1.2,
    )
    add_box(slide, ax + aw + 0.28, 1.95, 0.014, 4.25, fill=C["muted"])  # 세로 룰

    # ── 구획 1: 같은 질문, 두 개의 결말 (C의 flow) ──
    rx = ax + aw + 0.6
    rw = SLIDE_W - MARGIN - rx
    zone_title(slide, rx, 1.95, rw, "같은 질문, 두 개의 결말")
    qw = 1.75
    q = add_box(slide, rx, 2.95, qw, 1.55, fill=C["bg_alt"], line=C["muted"], line_w=0.75)
    set_shape_text(
        q,
        '"이 계약, 수익을\n지금 인식해도\n됩니까?"',
        S["caption"],
        F["head"],
        C["primary"],
        bold=True,
    )
    fx = rx + qw + 0.35
    fw = rw - qw - 0.35
    lanes = [
        (
            2.45,
            "일반 LLM",
            C["accent"],
            True,
            [
                ("즉시 답변", "근거 조문 없음"),
                ("확신에 찬 어조", "판단 창작"),
                ("검증 불가", "인용 근거 없음"),
            ],
            "치명",
        ),
        (
            3.6,
            "본 시스템",
            C["primary"],
            False,
            [
                ("근거 먼저 검색", "문단·사례 인용"),
                ("확실하면 확정", "결론+인용 문단"),
                ("애매하면 유보", "부족 정보 질문"),
            ],
            "안전",
        ),
    ]
    for ly, name, ink, fatal, steps, verdict in lanes:
        add_text(slide, fx, ly - 0.26, 2.0, 0.24, name, S["caption"], F["head"], ink, bold=True)
        step_w = (fw - 0.85 - 2 * 0.3) / 3
        for i, (head, desc) in enumerate(steps):
            x = fx + i * (step_w + 0.3)
            add_box(
                slide,
                x,
                ly,
                step_w,
                0.72,
                fill=C["bg"] if fatal else C["bg_alt"],
                line=ink,
                line_w=1.0,
            )
            add_text(
                slide,
                x + 0.1,
                ly + 0.06,
                step_w - 0.2,
                0.26,
                head,
                S["caption"],
                F["head"],
                C["primary"],
                bold=True,
            )
            add_text(
                slide, x + 0.1, ly + 0.36, step_w - 0.2, 0.3, desc, S["foot"], F["body"], C["muted"]
            )
            if i < 2:
                add_text(
                    slide,
                    x + step_w - 0.02,
                    ly + 0.2,
                    0.32,
                    0.3,
                    "→",
                    S["body"],
                    F["head"],
                    C["muted"],
                    bold=True,
                    align=PP_ALIGN.CENTER,
                )
        chip = add_box(slide, fx + fw - 0.8, ly, 0.8, 0.72, fill=ink)
        set_shape_text(chip, verdict, S["body"], F["head"], C["bg"], bold=True)

    # ── 구획 2: 오류의 두 방향 (A의 대비 패널) ──
    zy = 4.9
    zone_title(slide, rx, zy, rw, "오류의 두 방향 — 왜 2종만 치명적인가")
    half = (rw - 0.3) / 2
    ph = 1.1
    add_box(slide, rx, zy + 0.45, half, ph, fill=C["bg_alt"])
    slide.shapes.add_picture(
        str(ICONS / "circle-check-big_primary.png"),
        Inches(rx + 0.18),
        Inches(zy + 0.62),
        Inches(0.32),
        Inches(0.32),
    )
    add_text(
        slide,
        rx + 0.62,
        zy + 0.6,
        half - 0.8,
        0.3,
        "1종 · 놓침 — 안전한 실패",
        S["body"],
        F["head"],
        C["primary"],
        bold=True,
    )
    add_text(
        slide,
        rx + 0.2,
        zy + 0.98,
        half - 0.4,
        0.75,
        '못 찾으면 "모른다"고 답함 — 추가 검토로 해소. 되물음 9건도 재균형 관리.',
        S["caption"],
        F["body"],
        C["text"],
        line_spacing=1.2,
    )
    px = rx + half + 0.3
    add_box(slide, px, zy + 0.45, half, ph, fill=C["primary"])
    slide.shapes.add_picture(
        str(ICONS / "circle-x_white.png"),
        Inches(px + 0.18),
        Inches(zy + 0.62),
        Inches(0.32),
        Inches(0.32),
    )
    add_text(
        slide,
        px + 0.62,
        zy + 0.6,
        half - 0.8,
        0.3,
        "2종 · 허위 확정 — 치명적 실패",
        S["body"],
        F["head"],
        C["bg"],
        bold=True,
    )
    add_text(
        slide,
        px + 0.2,
        zy + 0.98,
        half - 0.4,
        0.75,
        "틀렸는데 확신에 찬 답 — 부실감사로 직결. 설계 전체를 이 방향 차단에 편향시켰다.",
        S["caption"],
        F["body"],
        C["bg_alt"],
        line_spacing=1.2,
    )

    band = add_box(slide, MARGIN, 6.42, SLIDE_W - 2 * MARGIN, 0.44, fill=C["primary"])
    set_shape_text(
        band,
        "확실할 때만 확정하고, 애매하면 근거를 보여주며 유보한다",
        S["sub"],
        F["head"],
        C["bg"],
        bold=True,
    )
    source_line(
        slide, "출처: 1_OVERVIEW.md · PROJECT_OVERVIEW.md · 5_INTERFACE.md (00_factsheet.md §A·§B)"
    )


def main():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    variant_a(prs)
    variant_b(prs)
    variant_c(prs)
    variant_d(prs)
    variant_e(prs)
    variant_f(prs)
    from _variant_g import variant_g
    variant_g(prs)
    from _variant_h import variant_h
    variant_h(prs)
    from _variant_i import lint, variant_i
    assert lint(), "정렬 린터 FAIL — 렌더 금지"
    variant_i(prs)
    from _variant_j import variant_j
    variant_j(prs)
    from _variant_k import variant_k
    variant_k(prs)
    out = SCRATCH / "s04_variants.pptx"
    prs.save(str(out))
    print(f"saved {out}")


if __name__ == "__main__":
    main()
