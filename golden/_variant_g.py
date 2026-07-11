# build_golden 통합 전 임시 모듈 — 시안 G (s04_variants.main에서 import)
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from build_golden import content_header, source_line
from kit import SLIDE_H, SLIDE_W, add_box, add_text, load_kit, mix, set_shape_text

K = load_kit()
C, S, F = K["rgb"], K["sizes"], K["fonts"]
MARGIN = K["layout"]["margin"]


def zone_title(slide, x, y, w, text):
    add_text(slide, x, y, w, 0.28, text, S["body"], F["head"], C["accent"], bold=True)
    add_box(slide, x, y + 0.32, w, 0.012, fill=C["muted"])


def chevron_step(slide, x, y, w, h, head, desc, tint):
    """셰브런 단계 도형 — 제목 + 설명 2문단."""
    cv = add_box(slide, x, y, w, h, fill=tint, shape="chevron")
    tf = cv.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = head
    r1.font.name, r1.font.bold = F["head"], True
    r1.font.size = Pt(S["caption"])
    r1.font.color.rgb = C["primary"]
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = desc
    r2.font.name = F["body"]
    r2.font.size = Pt(S["foot"])
    r2.font.color.rgb = C["muted"]


def variant_g(prs):
    """G안(v8) — F 구도 + 도형 어휘(셰브런·라운드·펜타곤) + 그리드 정렬 + 프로 텍스트."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C["bg"])
    content_header(
        slide,
        1,
        "문제 정의",
        "일반 LLM을 회계에 못 쓰는 이유는 틀려서가 아니라, 틀리는 방향 때문이다",
    )

    # 그리드 상수
    GAP = 0.25
    ax, aw = MARGIN, 2.95
    rx = ax + aw + 2 * GAP + 0.014
    rw = SLIDE_W - MARGIN - rx

    # 좌측 앵커: 캡션 라벨 + 대형 용어 + 리드 + 번호 리드인
    add_text(
        slide, ax, 1.95, aw, 0.26, "핵심 리스크", S["caption"], F["head"], C["muted"], bold=True
    )
    add_text(slide, ax, 2.25, aw, 0.7, "2종 오류", S["display"], F["head"], C["accent"], bold=True)
    add_text(
        slide, ax, 2.95, aw, 0.48, "= 허위 확정", S["title"], F["head"], C["primary"], bold=True
    )
    add_text(
        slide,
        ax,
        3.6,
        aw,
        0.75,
        "근거 없는 확신이 부실감사와 재무제표 왜곡으로 이어진다.",
        S["sub"],
        F["body"],
        C["text"],
        line_spacing=1.2,
    )
    nums = [
        ("01", "검증 불가", "그럴듯한 답을 내놓지만 어느 문단에 기반했는지 확인할 수 없다"),
        ("02", "판단 창작", "기준서에 없는 기준을 확신에 찬 어조로 지어낸다"),
    ]
    for i, (no, head, desc) in enumerate(nums):
        ny = 4.55 + i * 0.85
        add_text(slide, ax, ny, 0.5, 0.3, no, S["head"], F["head"], C["accent"], bold=True)
        add_text(
            slide, ax + 0.5, ny, aw - 0.5, 0.28, head, S["body"], F["head"], C["primary"], bold=True
        )
        add_text(
            slide,
            ax + 0.5,
            ny + 0.3,
            aw - 0.5,
            0.5,
            desc,
            S["caption"],
            F["body"],
            C["muted"],
            line_spacing=1.15,
        )
    add_box(slide, ax + aw + GAP, 1.95, 0.014, 4.25, fill=C["muted"])

    # 구획 1: 말풍선 질문 + 셰브런 flow 2레인 + 펜타곤 결말
    zone_title(slide, rx, 1.95, rw, "같은 질문, 두 개의 결말")
    lane1_y, lane2_y, card_h = 2.75, 3.95, 0.8
    qw, qh = 1.7, 1.3
    qy = (lane1_y + lane2_y + card_h) / 2 - qh / 2
    q = add_box(
        slide, rx, qy, qw, qh, fill=C["bg_alt"], line=C["muted"], line_w=0.75, shape="round"
    )
    set_shape_text(
        q,
        '"이 계약, 수익을\n지금 인식해도\n됩니까?"',
        S["caption"],
        F["head"],
        C["primary"],
        bold=True,
    )
    fx = rx + qw + GAP
    chev_zone = rw - qw - GAP - 0.95 - GAP
    chev_w = (chev_zone - 2 * 0.12) / 3
    lanes = [
        (
            lane1_y,
            "일반 LLM",
            C["accent"],
            True,
            [
                ("즉시 답변", "근거 조문 없음"),
                ("확신에 찬 어조", "판단 창작"),
                ("검증 불가", "인용 근거 없음"),
            ],
            "치명",
        ),
        (
            lane2_y,
            "본 시스템",
            C["primary"],
            False,
            [
                ("근거 먼저 검색", "문단·사례 인용"),
                ("확실하면 확정", "결론+인용 문단"),
                ("애매하면 유보", "부족 정보 질문"),
            ],
            "안전",
        ),
    ]
    for ly, name, ink, fatal, steps, verdict in lanes:
        add_text(
            slide, fx + 0.1, ly - 0.28, 2.0, 0.24, name, S["caption"], F["head"], ink, bold=True
        )
        for i, (head, desc) in enumerate(steps):
            x = fx + i * (chev_w + 0.12)
            tint = mix(ink, C["bg"], 0.92 if fatal else 0.88)
            chevron_step(slide, x, ly, chev_w, card_h, head, desc, tint)
        pent = add_box(slide, fx + chev_zone + GAP, ly, 0.95, card_h, fill=ink, shape="pentagon")
        set_shape_text(pent, verdict, S["body"], F["head"], C["bg"], bold=True)

    # 구획 2: 라운드 대비 카드 + 좌측 컬러 스파인
    zy = 5.05
    zone_title(slide, rx, zy, rw, "오류의 두 방향 — 왜 2종만 치명적인가")
    half = (rw - GAP) / 2
    ph = 1.0
    panels = [
        (
            rx,
            C["bg_alt"],
            C["muted"],
            "1종 · 놓침 — 안전한 실패",
            C["primary"],
            '못 찾으면 "모른다"고 답함 — 추가 검토로 해소. 되물음 9건도 재균형 관리.',
            C["text"],
        ),
        (
            rx + half + GAP,
            C["primary"],
            C["accent"],
            "2종 · 허위 확정 — 치명적 실패",
            C["bg"],
            "틀렸는데 확신에 찬 답 — 부실감사로 직결. 설계 전체를 이 방향 차단에 편향.",
            C["bg_alt"],
        ),
    ]
    for px, fill, spine, head, head_ink, desc, desc_ink in panels:
        add_box(slide, px, zy + 0.42, half, ph, fill=fill, shape="round")
        add_box(slide, px + 0.12, zy + 0.56, 0.05, ph - 0.28, fill=spine)
        add_text(
            slide,
            px + 0.32,
            zy + 0.52,
            half - 0.5,
            0.28,
            head,
            S["body"],
            F["head"],
            head_ink,
            bold=True,
        )
        add_text(
            slide,
            px + 0.32,
            zy + 0.84,
            half - 0.5,
            0.55,
            desc,
            S["caption"],
            F["body"],
            desc_ink,
            line_spacing=1.15,
        )

    band = add_box(slide, MARGIN, 6.42, SLIDE_W - 2 * MARGIN, 0.4, fill=C["primary"])
    set_shape_text(
        band,
        "확실할 때만 확정하고, 애매하면 근거를 보여주며 유보한다",
        S["sub"],
        F["head"],
        C["bg"],
        bold=True,
    )
    source_line(
        slide, "출처: 1_OVERVIEW.md · PROJECT_OVERVIEW.md · 5_INTERFACE.md (00_factsheet.md §A·§B)"
    )
