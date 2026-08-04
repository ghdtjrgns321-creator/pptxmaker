"""덱 진행 상태 — `_workspace/<프로젝트>/STATE.md`를 만들고, 실물과 대조한다.

왜 있나 (2026-08-03): 목차 계약(`01.5_outline.md`) 다음부터 산출 직전까지 **아무것도 파일로
남지 않았다.** 어떤 형태를 왜 골랐는지, 어디까지 왔는지, 무엇을 반려당했는지가 전부 대화
안에만 있다가 세션과 함께 사라졌다. 실측: 갤러리 단계에만 3시간 53분이 들고 전면 재시작이
2회 있었는데 그 왕복이 어디에도 없다. 반려도 마찬가지다 — 2026-08-02 14:25에 나온
"카드를 여러 개 만들어 골라 쓰자"가 13분 뒤 같은 장에서 재발하고 다음 날 또 재발했다.

**지도(MAP.md)와 역할이 다르다.** 지도는 코드에서 뽑으니 사람이 못 쓴다. 여기 담기는 것은
선택과 반려 — 사람만 아는 것이라 사람이 쓴다. 대신 기계가 **실물과 어긋나는지**를 본다.

    uv run python scripts/deck_state.py init    <프로젝트>   # 계약에서 뼈대 생성
    uv run python scripts/deck_state.py check   <프로젝트>   # 실물 대조 (어긋나면 exit 1)
    uv run python scripts/deck_state.py gallery <프로젝트>   # ③갤러리가 후보 자격을 갖췄나

## 검사 러너를 여기서 부른다 (2026-08-04)

`check_outline.py`·`check_deck.py`는 만들어만 두고 **아무도 부르지 않았다.** 절차 문서가
"돌려라"라고 적어도 이 저장소에서 산문 지시는 예외 없이 무시됐고, 실제로 두 러너는 배선 0회로
남아 있었다. 그래서 이미 게이트인 이 스크립트가 단계에 맞춰 직접 부른다 —
`init`이 ①계약을, `check`가 ⑦덱을 본다.

**import하지 않고 subprocess로 부른다.** 러너 셋이 서로를 import하지 않는 것은 지켜야 할
성질이다(한쪽 임계를 고치면 다른 쪽 판정이 조용히 따라 움직인다). 종료코드만 받는다.
"""

import importlib.util
import os
import re
import subprocess
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
WS = ROOT / "_workspace"
CHECK_OUTLINE = ROOT / "scripts" / "check_outline.py"
CHECK_DECK = ROOT / ".claude" / "skills" / "pptx-build" / "scripts" / "check_deck.py"
STEPS = ["①계약", "②형태후보", "③갤러리", "④선택", "⑤조판", "⑥눈검증", "⑦채점"]
FORMAL = ("표지", "목차", "간지")  # 정형 장 — 갤러리를 만들지 않는다
CORE_MIN = 0.5  # 후보 하나가 담아야 할 핵심 실체 비율 (v1 실측: 최저 3/5=60%, 불량 0/5)


def contract_slides(proj: Path) -> list[tuple[str, str]]:
    """아웃라인 계약의 장 목록 표에서 (장번호, 제목)을 뽑는다."""
    f = proj / "01.5_outline.md"
    if not f.exists():
        raise SystemExit(f"계약이 없다: {f.relative_to(ROOT)} — 계약 없이 진행하지 않는다(①)")
    rows = []
    for ln in f.read_text(encoding="utf-8").splitlines():
        m = re.match(r"^\|\s*(S\d+)\s*\|[^|]*\|\s*([^|]+?)\s*\|", ln)
        if m:
            rows.append((m.group(1), m.group(2)))
    if not rows:
        raise SystemExit(f"계약에서 장 목록 표를 못 찾았다: {f.relative_to(ROOT)}")
    return rows


def material(proj: Path) -> Path | None:
    """계약이 밝힌 재료 파일. 계약의 `- **재료**: \\`경로\\`` 줄에서 읽는다.

    경로를 계약에 적게 만드는 것 자체가 목적이다 — 안 적히면 `check_outline`이 영원히 안 돈다.
    재료는 이 저장소 밖(다른 프로젝트의 README)에 있는 게 정상이라 상위 폴더까지 훑는다.
    """
    f = proj / "01.5_outline.md"
    if not f.exists():
        return None
    m = re.search(r"^-\s*\*\*재료\*\*\s*:\s*`([^`]+)`", f.read_text(encoding="utf-8"), re.M)
    if not m:
        return None
    raw = m.group(1).strip()
    for base in (proj, ROOT, ROOT.parent):
        p = base / raw
        if p.is_file():
            return p
    return None


def run_check(script: Path, *args) -> tuple[int, str]:
    """검사 러너를 별도 프로세스로 돌리고 (종료코드, 출력)을 준다.

    한글 출력이 깨지지 않게 자식 쪽 표준출력 인코딩을 못 박는다 — 윈도우 기본은 cp949라
    안 박으면 검사 결과가 U+FFFD로 돌아온다.
    """
    if not script.exists():
        return 1, f"검사기가 없다: {script.relative_to(ROOT)}"
    env = {**os.environ, "PYTHONIOENCODING": "utf-8"}
    r = subprocess.run(
        [sys.executable, str(script), *(str(a) for a in args)],
        capture_output=True,
        text=True,
        encoding="utf-8",
        errors="replace",
        cwd=ROOT,
        env=env,
    )
    return r.returncode, (r.stdout + r.stderr).strip()


def build(proj: Path) -> str:
    slides = contract_slides(proj)
    L = [
        f"# 진행 상태 — {proj.name}",
        "",
        "> 뼈대는 `scripts/deck_state.py`가 계약에서 뽑았고, **내용은 사람이 채운다.**",
        f"> 실물과 어긋나는지는 `uv run python scripts/deck_state.py check {proj.name}` 가 본다.",
        "",
        "## 지금 어디",
        "",
        f"- **단계**: {STEPS[0]}   ({' · '.join(STEPS)})",
        "- **갱신**: (날짜)",
        "- **다음 할 일**: (한 줄)",
        "",
        "## 장별 — 무엇을 왜 골랐나",
        "",
        "형태를 고른 뒤 채운다. `왜`가 비면 다음 덱이 같은 선택을 재현하지 못한다.",
        "",
        f"`check`가 이 표로 **획일화**를 잰다 — 본문에서 같은 조합이 {SAME_COMBO_MAX}장을 넘거나 "
        f"한 부품이 {ONE_PART_SHARE:.0%}를 넘게 등장하면 경고. 표지·목차·간지(`layouts.*`)는 "
        "일부러 같아야 하므로 세지 않는다.",
        "",
        "| 장 | 제목 | 고른 형태 | 왜 이걸 골랐나 | 상태 |",
        "| --- | --- | --- | --- | --- |",
    ]
    for no, title in slides:
        L.append(f"| {no} | {title} | — | — | 대기 |")
    L += [
        "",
        "## 반려 기록 — 사용자가 무엇을 물렸나",
        "",
        "**이 표가 다음 덱으로 넘어가는 유일한 통로다.** 여기 안 적힌 반려는 사라진다.",
        "`하네스` 열이 `미정`인 줄이 남아 있으면 `check`가 경고한다 — 판단을 미룬 채 덮지 않기 위해서다.",
        "",
        "| # | 날짜 | 장 | 무엇을 물렸나 | 어떻게 고쳤나 | 하네스 |",
        "| --- | --- | --- | --- | --- | --- |",
        "| 1 | | | | | 미정 |",
        "",
        "`하네스` 열에 쓸 값: `불필요`(이 덱 한정) · `제안함`(사용자 승인 대기) · "
        "`반영`(코드·게이트로 내려감) · `미정`",
        "",
    ]
    return "\n".join(L)


# 획일화 임계 — local-ai-assist(본문 12장, 조합 12종)에서 오탐 0으로 맞춘 값.
# 잡으려는 실패는 2026-07-29 CHANGELOG에 박제된 것: "전 장이 같은 형식", "36장이 동일한 4카드".
SAME_COMBO_MAX = 2  # 완전히 같은 조합이 이 수를 넘으면 형태가 아니라 틀을 찍은 것이다
ONE_PART_SHARE = 0.60  # 한 부품이 본문 장의 이 비율을 넘으면 그 부품이 덱을 지배한다


def variety(txt: str) -> list[str]:
    """본문 장들이 서로 다른 형태를 쓰고 있나 — 부품을 늘려도 획일화가 자동으로 잡히게.

    정형 장(`layouts.*`)은 **일부러 같아야 하므로** 뺀다. 서식 일관성은 표지·목차·간지가
    지고, 형태 다양성은 본문이 진다 — 이 둘을 한 통에 세면 둘 다 못 본다.
    """
    combos: dict[str, list[str]] = {}
    part_hits: dict[str, set[str]] = {}
    for ln in txt.splitlines():
        m = re.match(r"^\|\s*(S\d+)\s*\|[^|]*\|([^|]*)\|", ln)
        if not m:
            continue
        # 사람이 쓰는 칸이라 표기가 흔들린다 — `a` + `b` 든 `a + b` 든 같게 읽는다
        cell = re.sub(r"[`*]", "", m.group(2)).strip()
        parts = sorted(
            {p.strip() for p in re.split(r"[+·,]", cell) if re.match(r"^[\w.]+$", p.strip())}
        )
        if not parts or any(p.startswith("layouts.") for p in parts):
            continue  # 미기록 장(—) · 정형 장(표지·목차·간지는 일부러 같다)
        combos.setdefault(" + ".join(parts), []).append(m.group(1))
        for p in parts:
            part_hits.setdefault(p, set()).add(m.group(1))

    n = sum(len(v) for v in combos.values())
    if n < 4:  # 표본이 적으면 판정하지 않는다
        return []
    out = []
    for combo, slides in sorted(combos.items(), key=lambda kv: -len(kv[1])):
        if len(slides) > SAME_COMBO_MAX:
            out.append(
                f"본문 {len(slides)}장이 같은 조합 `{combo}` ({', '.join(slides)})"
                " — 형태를 고른 게 아니라 틀을 찍었다"
            )
    for p, slides in sorted(part_hits.items(), key=lambda kv: -len(kv[1])):
        if len(slides) / n > ONE_PART_SHARE:
            out.append(
                f"`{p}`가 본문 {len(slides)}/{n}장({len(slides) / n:.0%})에 등장"
                " — 한 부품이 덱을 지배한다"
            )
    print(f"  형태: 본문 {n}장 · 서로 다른 조합 {len(combos)}종")
    return out


def gallery(proj: Path) -> int:
    """③갤러리 검사 — 후보가 "같은 장의 다른 배치"인지 기계가 본다.

    잡으려는 실패는 2026-08-03에 실제로 낸 것이다: 후보 ④에서 핵심 실체(5축)를 통째로
    빼고 깔때기만 그렸더니 사용자에게 **"주제가 다른 것 중에 고르라는 거냐"**로 읽혔다.
    SKILL.md ②가 "고르는 건 무엇을이 아니라 어떻게다"라고 이미 적어 뒀는데도 어겼다 —
    산문으로 있는 규칙은 이 저장소에서 예외 없이 무시됐으므로 기계로 내린다.

    갤러리 스크립트가 `CORE = {"S4": ["검증 범위", ...], ...}`로 장별 핵심 실체를 선언하고,
    이 검사가 후보 4개 각각에 그 절반 이상이 있는지 본다. 라벨 접기(… 외 N개)는 허용된다
    — 접혀도 실체 이름은 남기 때문이다(v1 실측: 최저 후보가 3/5).
    """
    from pptx import Presentation
    from pptx.util import Emu

    pptxs = sorted(proj.glob("gallery*.pptx"))
    pys = sorted(proj.glob("gallery*.py"))
    if not pptxs or not pys:
        print(f"★ 갤러리가 없다 (pptx {len(pptxs)} · py {len(pys)}) — ③을 아직 안 했다")
        return 1

    body = [(n, t) for n, t in contract_slides(proj) if t not in FORMAL]
    warn = []

    core: dict[str, list[str]] = {}
    for py in pys:
        spec = importlib.util.spec_from_file_location(py.stem, py)
        if spec is None or spec.loader is None:
            continue
        mod = importlib.util.module_from_spec(spec)
        sys.path.insert(0, str(proj))
        try:
            spec.loader.exec_module(mod)
            core.update(getattr(mod, "CORE", {}) or {})
        except Exception as e:
            warn.append(f"`{py.name}`을 못 읽는다: {type(e).__name__} {e}")
        finally:
            sys.path.remove(str(proj))
    if not core:
        warn.append(
            "갤러리 스크립트에 `CORE`가 없다 — 장별 핵심 실체를 선언하지 않으면"
            " 후보가 같은 내용인지 기계가 볼 수 없다"
        )

    seen = 0
    for pptx in pptxs:
        prs = Presentation(str(pptx))
        sw = Emu(prs.slide_width).inches
        for sl in prs.slides:
            seen += 1
            head = " ".join(
                sh.text_frame.text
                for sh in sl.shapes
                if sh.has_text_frame and sh.top is not None and Emu(sh.top).inches < 1.0
            )
            m = re.search(r"\bS(\d+)\b", head)
            key = f"S{m.group(1)}" if m else None
            if key not in core:
                if core:
                    warn.append(f"`{pptx.name}`의 '{head.strip()[:24]}' 장에 `CORE` 선언이 없다")
                continue
            words = core[key]
            need = max(1, round(len(words) * CORE_MIN))
            quads: list[list] = [[], [], [], []]
            for sh in sl.shapes:
                if sh.left is None or sh.top is None:
                    continue
                cx = Emu(sh.left).inches + Emu(sh.width).inches / 2
                cy = Emu(sh.top).inches + Emu(sh.height).inches / 2
                if not (1.0 <= cy <= 6.85):
                    continue
                quads[(0 if cy < 3.7 else 2) + (0 if cx < sw / 2 else 1)].append(sh)
            for i, q in enumerate(quads, 1):
                txt = " ".join(sh.text_frame.text for sh in q if sh.has_text_frame)
                hit = [w for w in words if w in txt]
                if len(hit) < need:
                    warn.append(
                        f"{key} 후보{i}: 핵심 실체 {len(hit)}/{len(words)} (하한 {need})"
                        " — 다른 배치가 아니라 다른 내용이다"
                    )

    if seen != len(body):
        warn.append(f"갤러리 {seen}장인데 본문은 {len(body)}장 — 빠진 장은 고를 수가 없다")

    print(f"갤러리 대조: {seen}장 / 본문 {len(body)}장 — 경고 {len(warn)}건")
    for w in warn:
        print("  ★", w)
    return 1 if warn else 0


def check(proj: Path) -> int:
    state = proj / "STATE.md"
    if not state.exists():
        print(f"★ 상태 파일이 없다: {state.relative_to(ROOT)} — `init`부터")
        return 1
    txt = state.read_text(encoding="utf-8")
    warn = []

    want = contract_slides(proj)
    got = re.findall(r"^\|\s*(S\d+)\s*\|", txt, re.M)
    if [n for n, _ in want] != got:
        warn.append(
            f"장 목록이 계약과 다르다 — 계약 {len(want)}장 / 상태 {len(got)}장"
            f" (계약에만: {sorted(set(n for n, _ in want) - set(got))} ·"
            f" 상태에만: {sorted(set(got) - set(n for n, _ in want))})"
        )

    deck = proj / "deck.pptx"
    if deck.exists():
        try:
            from pptx import Presentation

            n = len(Presentation(str(deck)).slides._sldIdLst)
            if n != len(want):
                warn.append(f"deck.pptx {n}장인데 계약은 {len(want)}장")
        except Exception as e:
            warn.append(f"deck.pptx를 못 읽는다: {type(e).__name__}")

        # ⑦ 규격 대조 — 여기가 `check_deck.py`의 유일한 호출처다(위 docstring 참조).
        args = [deck, f"--outline={proj / '01.5_outline.md'}"]
        mat = material(proj)
        if mat:
            args.append(f"--material={mat}")
        rc, out = run_check(CHECK_DECK, *args)
        for ln in out.splitlines():
            print("  │", ln)
        if rc:
            warn.append(f"`check_deck.py`가 규격 결함을 잡았다 (exit {rc}) — 위 `│` 줄")
        elif not mat:
            warn.append("계약에 `- **재료**: `경로`` 가 없어 재료 커버리지를 못 쟀다")

    rend = proj / "deck-render"
    if rend.exists():
        pngs = len(list(rend.glob("s*.png")))
        if deck.exists() and pngs and pngs != len(want):
            warn.append(f"렌더 PNG {pngs}장인데 계약은 {len(want)}장 — 렌더가 낡았다")

    step = re.search(r"\*\*단계\*\*:\s*(\S+)", txt)
    if not step:
        warn.append("`단계`가 비었다 — 지금 어디인지 안 적혀 있다")
    elif step.group(1) not in STEPS:
        warn.append(f"단계 `{step.group(1)}`는 정의에 없다 (허용 {' · '.join(STEPS)})")

    body = txt.split("## 장별")[-1]
    if body.count("| — | — | 대기 |") == len(want) and step and step.group(1) not in STEPS[:3]:
        warn.append(f"단계는 {step.group(1)}인데 고른 형태가 한 장도 안 적혔다")

    warn += variety(txt)

    # 반려 표: | # | 날짜 | 장 | 무엇을 | 어떻게 | 하네스 |
    # 마지막 칸이 '미정'이면 그 반려는 다음 덱으로 갈지 안 갈지 정해지지 않은 것이다.
    # 굵게·꼬리말이 붙어도 잡히도록 마크업을 벗기고 앞머리만 본다.
    undecided = []
    for ln in txt.splitlines():
        if not re.match(r"^\|\s*\d+\s*\|", ln):
            continue
        cells = [c.strip() for c in ln.strip().strip("|").split("|")]
        if len(cells) < 6 or not cells[3]:  # 내용이 안 적힌 뼈대 줄은 넘긴다
            continue
        verdict = re.sub(r"[*`_]", "", cells[5]).strip()
        if verdict.startswith("미정"):
            undecided.append(cells[0])
    if undecided:
        warn.append(
            f"반려 {len(undecided)}건이 `하네스: 미정` (#{', #'.join(undecided)})"
            " — 다음 덱으로 내릴지 안 정했다. 미정은 재발한다"
        )

    print(f"대조: {state.relative_to(ROOT)} — 경고 {len(warn)}건")
    for w in warn:
        print("  ★", w)
    return 1 if warn else 0


def main() -> int:
    if len(sys.argv) != 3 or sys.argv[1] not in ("init", "check", "gallery"):
        print(__doc__)
        return 2
    cmd, name = sys.argv[1], sys.argv[2]
    proj = WS / name
    if not proj.is_dir():
        raise SystemExit(f"작업 폴더가 없다: {proj.relative_to(ROOT)}")
    if cmd == "gallery":
        return gallery(proj)
    if cmd == "init":
        out = proj / "STATE.md"
        if out.exists():
            raise SystemExit(f"이미 있다: {out.relative_to(ROOT)} — 덮지 않는다. 직접 고칠 것")
        # ① 계약 대조 — 여기가 `check_outline.py`의 유일한 호출처다(위 docstring 참조).
        # 지어낸 제목 위에 STATE를 세우면 그 창작이 뼈대가 되어 끝까지 간다. 계약을 먼저 고친다.
        mat = material(proj)
        if not mat:
            raise SystemExit(
                f"계약에 재료가 없다: {(proj / '01.5_outline.md').relative_to(ROOT)} —"
                " `- **재료**: `경로`` 한 줄을 적어야 제목이 자료에서 나왔는지 대조할 수 있다"
            )
        rc, msg = run_check(CHECK_OUTLINE, proj / "01.5_outline.md", mat)
        for ln in msg.splitlines():
            print("  │", ln)
        if rc:
            raise SystemExit("★ 계약에 자료 밖 제목이 있다 — 계약을 고친 뒤 다시 `init`")
        out.write_text(build(proj), encoding="utf-8")
        print(f"생성: {out.relative_to(ROOT)} — 장 {len(contract_slides(proj))}개")
        return 0
    return check(proj)


if __name__ == "__main__":
    sys.exit(main())
